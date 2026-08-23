"""Generate the deterministic public image benchmark corpus."""

from __future__ import annotations

import argparse
import hashlib
import json
import zlib
from datetime import datetime
from io import BytesIO
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

import lxml.etree
import numpy as np
import PIL
import pptx
import reportlab
from PIL import Image, ImageDraw, ImageFilter, ImageFont, ImageOps, features
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
SCHEMA_VERSION = 1
LANDSCAPE = (1600, 900)
PORTRAIT = (900, 1600)

IMAGE_CASES = [
    ("01_zh_courseware", "image", "01-zh-courseware.png", 1),
    ("02_typography", "image", "02-typography.png", 1),
    ("03_flowchart", "image", "03-flowchart.png", 1),
    ("04_table_chart", "image", "04-table-chart.png", 1),
    ("05_photo_overlay", "image", "05-photo-overlay.png", 1),
    ("06_transparency_shadow", "image", "06-transparency-shadow.png", 1),
    ("07_compressed", "image", "07-compressed.jpg", 1),
    ("08_portrait", "image", "08-portrait.png", 1),
]
DOCUMENT_CASES = [
    ("09_document", "pdf", "09-document.pdf", 3),
    ("10_mixed", "pptx", "10-mixed.pptx", 3),
]
CASES = IMAGE_CASES + DOCUMENT_CASES
ROUTES = [
    {"id": "images", "cases": [case[0] for case in IMAGE_CASES], "pages": 8},
    {"id": "pdf", "cases": ["09_document"], "pages": 3},
    {"id": "mixed_pptx", "cases": ["10_mixed"], "pages": 3},
]
FIXED_TIME = (1980, 1, 1, 0, 0, 0)
CANONICAL_EXTERNAL_ATTR = 0o600 << 16
PDF_PAGE_SOURCES = (
    "01-zh-courseware.png",
    "04-table-chart.png",
    "08-portrait.png",
)
EXTENDED_PROPERTIES_NAMESPACE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
)
PRESENTATION_NAMESPACE = (
    "http://schemas.openxmlformats.org/presentationml/2006/main"
)

ENCODER_ERROR = (
    "deterministic benchmark encoder unavailable: requires Pillow 10.4.0, "
    "libjpeg_turbo 3.0.3, zlib 1.3.1, freetype2 2.13.2"
)
DOCUMENT_ENVIRONMENT_ERROR = (
    "deterministic benchmark document environment unavailable: requires "
    "ReportLab 5.0.0, python-pptx 1.0.2, stdlib zlib 1.3.1, lxml 6.1.1"
)


def _require_deterministic_encoder() -> None:
    fingerprint = (
        PIL.__version__,
        features.version_feature("libjpeg_turbo"),
        features.version_codec("zlib"),
        features.version_module("freetype2"),
    )
    if fingerprint != ("10.4.0", "3.0.3", "1.3.1", "2.13.2"):
        raise RuntimeError(ENCODER_ERROR)


def _require_deterministic_document_environment() -> None:
    fingerprint = (
        reportlab.Version,
        pptx.__version__,
        zlib.ZLIB_RUNTIME_VERSION,
        lxml.etree.LXML_VERSION,
    )
    if fingerprint != ("5.0.0", "1.0.2", "1.3.1", (6, 1, 1, 0)):
        raise RuntimeError(DOCUMENT_ENVIRONMENT_ERROR)


def _write_png(image: Image.Image, path: Path) -> None:
    image.save(path, format="PNG", optimize=False, compress_level=9)


def _record(
    root: Path, case_id: str, kind: str, name: str, pages: int
) -> dict[str, object]:
    path = root / name
    payload = path.read_bytes()
    return {
        "id": case_id,
        "kind": kind,
        "path": name,
        "pages": pages,
        "bytes": len(payload),
        "sha256": hashlib.sha256(payload).hexdigest(),
    }


def _canonical_sha256(value: object) -> str:
    payload = json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        allow_nan=False,
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _write_pdf(output: Path) -> None:
    document = canvas.Canvas(str(output / "09-document.pdf"), pagesize=LANDSCAPE, invariant=1)
    for name in PDF_PAGE_SOURCES:
        image = ImageReader(BytesIO((output / name).read_bytes()))
        document.drawImage(
            image,
            0,
            0,
            width=LANDSCAPE[0],
            height=LANDSCAPE[1],
            preserveAspectRatio=True,
            anchor="c",
        )
        document.showPage()
    document.save()


def _required_xml_element(root, namespace: str, name: str):
    elements = root.findall(f"{{{namespace}}}{name}")
    if len(elements) != 1:
        raise ValueError(f"expected exactly one {name} element")
    return elements[0]


def _canonicalize_pptx_xml(filename: str, payload: bytes) -> bytes:
    if filename == "docProps/app.xml":
        root = lxml.etree.fromstring(payload)
        _required_xml_element(
            root, EXTENDED_PROPERTIES_NAMESPACE, "PresentationFormat"
        ).text = "On-screen Show (16:9)"
        _required_xml_element(
            root, EXTENDED_PROPERTIES_NAMESPACE, "Slides"
        ).text = "3"
    elif filename == "ppt/presentation.xml":
        root = lxml.etree.fromstring(payload)
        _required_xml_element(root, PRESENTATION_NAMESPACE, "sldSz").set(
            "type", "screen16x9"
        )
    else:
        return payload
    return lxml.etree.tostring(
        root,
        encoding="UTF-8",
        xml_declaration=True,
        standalone=True,
    )


def _normalize_pptx(payload: bytes) -> bytes:
    normalized = BytesIO()
    with ZipFile(BytesIO(payload), "r") as source, ZipFile(
        normalized, "w", compression=ZIP_DEFLATED
    ) as target:
        for source_info in sorted(source.infolist(), key=lambda info: info.filename):
            info = ZipInfo(source_info.filename, FIXED_TIME)
            info.compress_type = ZIP_DEFLATED
            info.create_system = 0
            info.external_attr = CANONICAL_EXTERNAL_ATTR
            target.writestr(
                info,
                _canonicalize_pptx_xml(
                    source_info.filename, source.read(source_info.filename)
                ),
            )
    return normalized.getvalue()


def _add_text(slide, text: str, left: float, top: float, width: float, height: float, size: int, color: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    frame = box.text_frame
    frame.clear()
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.word_wrap = True
    run = frame.paragraphs[0].add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color)


def _write_normalized_pptx(output: Path) -> None:
    deck = Presentation()
    deck.slide_width = 12_192_000
    deck.slide_height = 6_858_000
    properties = deck.core_properties
    fixed_date = datetime(2020, 1, 1)
    properties.author = "image2editable"
    properties.created = fixed_date
    properties.last_modified_by = "image2editable"
    properties.modified = fixed_date
    properties.revision = 1
    properties.subject = "Editable conversion benchmark"
    properties.title = "Mixed editable benchmark"

    blank = deck.slide_layouts[6]
    first = deck.slides.add_slide(blank)
    first.background.fill.solid()
    first.background.fill.fore_color.rgb = RGBColor(245, 247, 250)
    accent = first.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.0), Inches(0.18), Inches(5.0)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(31, 78, 121)
    accent.line.fill.background()
    _add_text(first, "NATIVE CONTENT", 1.35, 1.35, 10.5, 0.7, 38, "1F4E79")
    _add_text(
        first,
        "Editable text and vector shapes preserve structure for downstream conversion.",
        1.35,
        2.45,
        9.8,
        1.2,
        22,
        "34495E",
    )
    rule = first.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.35), Inches(4.45), Inches(4.0), Inches(0.12)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = RGBColor(225, 126, 74)
    rule.line.fill.background()

    second = deck.slides.add_slide(blank)
    second.shapes.add_picture(
        str(output / "04-table-chart.png"),
        0,
        0,
        width=deck.slide_width,
        height=deck.slide_height,
    )

    third = deck.slides.add_slide(blank)
    third.background.fill.solid()
    third.background.fill.fore_color.rgb = RGBColor(16, 42, 67)
    _add_text(third, "MIXED CONTENT", 0.8, 1.15, 4.9, 0.7, 36, "F6AE2D")
    _add_text(
        third,
        "Native text remains editable while the image region retains visual detail.",
        0.8,
        2.25,
        4.6,
        1.55,
        20,
        "F0F4F8",
    )
    third.shapes.add_picture(
        str(output / "05-photo-overlay.png"),
        Inches(6.0),
        Inches(1.78),
        width=Inches(6.55),
        height=Inches(3.684375),
    )

    raw = BytesIO()
    deck.save(raw)
    (output / "10-mixed.pptx").write_bytes(_normalize_pptx(raw.getvalue()))


def generate_corpus(output: Path, source_root: Path = ROOT) -> None:
    _require_deterministic_encoder()
    _require_deterministic_document_environment()
    output = Path(output)
    source_root = Path(source_root)
    if output.exists() and (not output.is_dir() or any(output.iterdir())):
        raise FileExistsError(f"output must not exist or must be empty: {output}")
    output.mkdir(parents=True, exist_ok=True)

    rng = np.random.default_rng(20260815)
    regular = ImageFont.load_default(size=34)
    small = ImageFont.load_default(size=25)
    heading = ImageFont.load_default(size=60)

    with Image.open(source_root / "docs/images/demo-source-1.png") as source:
        courseware = ImageOps.fit(source.convert("RGB"), LANDSCAPE, Image.Resampling.LANCZOS)
    _write_png(courseware, output / "01-zh-courseware.png")

    typography = Image.new("RGB", LANDSCAPE, "#F4F0E8")
    draw = ImageDraw.Draw(typography)
    draw.rectangle((0, 0, 120, 900), fill="#17324D")
    draw.text((190, 115), "TYPE SYSTEM", font=small, fill="#D05A3A")
    draw.text((190, 180), "Hierarchy makes", font=heading, fill="#17324D")
    draw.text((190, 255), "meaning visible.", font=heading, fill="#17324D")
    draw.line((190, 370, 1410, 370), fill="#C9BFB0", width=3)
    draw.text((190, 430), "01  DISPLAY", font=regular, fill="#17324D")
    draw.text((620, 430), "Large type creates a clear entry point.", font=regular, fill="#4E5C68")
    draw.text((190, 540), "02  BODY", font=small, fill="#17324D")
    draw.multiline_text(
        (620, 535),
        "Supporting copy uses a quieter scale\nand a comfortable reading rhythm.",
        font=small,
        fill="#4E5C68",
        spacing=18,
    )
    draw.text((190, 730), "1600 / 900", font=small, fill="#D05A3A")
    _write_png(typography, output / "02-typography.png")

    flowchart = Image.new("RGB", LANDSCAPE, "#F7FAFC")
    draw = ImageDraw.Draw(flowchart)
    draw.text((110, 80), "SERVICE FLOW", font=heading, fill="#183B56")
    nodes = [
        ((120, 350, 400, 540), "INPUT", "#D6EAF8"),
        ((660, 350, 940, 540), "PROCESS", "#D5F5E3"),
        ((1200, 350, 1480, 540), "OUTPUT", "#FADBD8"),
    ]
    for box, label, color in nodes:
        draw.rounded_rectangle(box, radius=30, fill=color, outline="#183B56", width=4)
        draw.ellipse((box[0] + 95, box[1] + 30, box[0] + 185, box[1] + 120), outline="#183B56", width=6)
        draw.text((box[0] + 70, box[1] + 135), label, font=small, fill="#183B56")
    for start, end in [((400, 445), (660, 445)), ((940, 445), (1200, 445))]:
        draw.line((*start, *end), fill="#E67E22", width=12)
        draw.polygon([(end[0], end[1]), (end[0] - 35, end[1] - 24), (end[0] - 35, end[1] + 24)], fill="#E67E22")
    draw.text((120, 690), "Collect", font=regular, fill="#486581")
    draw.text((660, 690), "Transform", font=regular, fill="#486581")
    draw.text((1200, 690), "Deliver", font=regular, fill="#486581")
    _write_png(flowchart, output / "03-flowchart.png")

    table_chart = Image.new("RGB", LANDSCAPE, "#FFFFFF")
    draw = ImageDraw.Draw(table_chart)
    draw.text((90, 65), "QUARTERLY MIX", font=heading, fill="#243B53")
    left, top, row_height = 90, 225, 95
    widths = [260, 190, 190]
    headers = ["CHANNEL", "Q1", "Q2"]
    rows = [("Direct", "42", "55"), ("Partner", "34", "39"), ("Organic", "28", "45")]
    x_positions = [left, left + widths[0], left + widths[0] + widths[1], left + sum(widths)]
    for row_index in range(5):
        y = top + row_index * row_height
        draw.line((left, y, left + sum(widths), y), fill="#BCCCDC", width=2)
    for x in x_positions:
        draw.line((x, top, x, top + 4 * row_height), fill="#BCCCDC", width=2)
    for col, text_value in enumerate(headers):
        draw.text((x_positions[col] + 20, top + 25), text_value, font=small, fill="#243B53")
    for row_index, row in enumerate(rows, start=1):
        for col, text_value in enumerate(row):
            draw.text((x_positions[col] + 20, top + row_index * row_height + 25), text_value, font=small, fill="#486581")
    axis_left, axis_bottom, axis_top = 880, 700, 240
    draw.line((axis_left, axis_top, axis_left, axis_bottom, 1490, axis_bottom), fill="#334E68", width=4)
    values = [220, 330, 410]
    colors = ["#2E86AB", "#F6AE2D", "#D1495B"]
    for index, (height, color) in enumerate(zip(values, colors)):
        x = axis_left + 100 + index * 180
        draw.rectangle((x, axis_bottom - height, x + 100, axis_bottom), fill=color)
        draw.rectangle((930 + index * 190, 775, 960 + index * 190, 805), fill=color)
        draw.text((970 + index * 190, 775), rows[index][0], font=small, fill="#486581")
    _write_png(table_chart, output / "04-table-chart.png")

    y = np.linspace(0, 1, LANDSCAPE[1], dtype=np.float32)[:, None, None]
    x = np.linspace(0, 1, LANDSCAPE[0], dtype=np.float32)[None, :, None]
    base = np.array([35, 91, 110], dtype=np.float32)
    light = np.array([226, 168, 104], dtype=np.float32)
    texture = base * (1 - x) + light * x
    texture = texture * (0.78 + 0.22 * y)
    texture += rng.normal(0, 12, (LANDSCAPE[1], LANDSCAPE[0], 1))
    photo = Image.fromarray(np.clip(texture, 0, 255).astype(np.uint8), "RGB")
    draw = ImageDraw.Draw(photo, "RGBA")
    draw.rectangle((0, 0, 1600, 900), fill=(6, 20, 32, 55))
    draw.rounded_rectangle((110, 470, 980, 790), radius=30, fill=(7, 24, 39, 185))
    draw.text((170, 525), "COASTAL FIELD NOTES", font=heading, fill="white")
    draw.multiline_text((170, 625), "Texture, light and atmosphere\nwith a readable overlay.", font=regular, fill="#F8E9D2", spacing=18)
    draw.rounded_rectangle((1220, 90, 1480, 155), radius=28, fill=(255, 255, 255, 210))
    draw.text((1260, 105), "FEATURE", font=small, fill="#17324D")
    _write_png(photo, output / "05-photo-overlay.png")

    gradient_left = np.array([238, 242, 247], dtype=np.float32)
    gradient_right = np.array([207, 226, 243], dtype=np.float32)
    gradient = np.linspace(gradient_left, gradient_right, LANDSCAPE[0], dtype=np.uint8)
    gradient = np.broadcast_to(gradient, (LANDSCAPE[1], LANDSCAPE[0], 3)).copy()
    transparency = Image.fromarray(gradient, "RGB")
    shadow = Image.new("RGBA", LANDSCAPE, (0, 0, 0, 0))
    shadow_draw = ImageDraw.Draw(shadow)
    shadow_draw.rounded_rectangle((270, 240, 930, 730), radius=70, fill=(22, 34, 52, 105))
    shadow = shadow.filter(ImageFilter.GaussianBlur(32))
    transparency = Image.alpha_composite(transparency.convert("RGBA"), shadow)
    layers = Image.new("RGBA", LANDSCAPE, (0, 0, 0, 0))
    draw = ImageDraw.Draw(layers)
    draw.rounded_rectangle((220, 180, 880, 670), radius=70, fill=(41, 128, 185, 185))
    draw.ellipse((610, 250, 1210, 850), fill=(241, 196, 15, 145))
    draw.rounded_rectangle((910, 120, 1460, 620), radius=70, fill=(192, 57, 43, 140))
    draw.text((100, 60), "ALPHA / SHADOW / OVERLAP", font=heading, fill="#243B53")
    draw.text((580, 420), "LAYERED", font=heading, fill="white")
    transparency = Image.alpha_composite(transparency, layers).convert("RGB")
    _write_png(transparency, output / "06-transparency-shadow.png")

    noise = rng.normal(0, 21, (LANDSCAPE[1], LANDSCAPE[0], 3))
    xx = np.linspace(20, 225, LANDSCAPE[0], dtype=np.float32)[None, :, None]
    yy = np.linspace(15, 90, LANDSCAPE[1], dtype=np.float32)[:, None, None]
    compressed_array = np.empty((LANDSCAPE[1], LANDSCAPE[0], 3), dtype=np.float32)
    compressed_array[..., 0:1] = xx
    compressed_array[..., 1:2] = 80 + yy
    compressed_array[..., 2:3] = 210 - xx * 0.45
    compressed = Image.fromarray(np.clip(compressed_array + noise, 0, 255).astype(np.uint8), "RGB")
    draw = ImageDraw.Draw(compressed)
    draw.rectangle((120, 120, 700, 450), fill="#EAD7B7", outline="#473B2B", width=8)
    draw.text((175, 220), "LOW QUALITY", font=heading, fill="#473B2B")
    compressed.save(
        output / "07-compressed.jpg",
        format="JPEG",
        quality=28,
        subsampling=2,
        optimize=False,
        progressive=False,
    )

    portrait = Image.new("RGB", PORTRAIT, "#102A43")
    draw = ImageDraw.Draw(portrait)
    draw.rectangle((70, 70, 830, 1530), outline="#9FB3C8", width=4)
    draw.text((115, 145), "PORTRAIT", font=small, fill="#F6AE2D")
    draw.multiline_text((115, 240), "VERTICAL\nSTORIES", font=heading, fill="#F0F4F8", spacing=24)
    draw.ellipse((170, 570, 730, 1130), fill="#2E86AB", outline="#F6AE2D", width=14)
    draw.line((450, 610, 450, 1090), fill="#F0F4F8", width=10)
    draw.line((210, 850, 690, 850), fill="#F0F4F8", width=10)
    draw.text((115, 1290), "900 x 1600", font=regular, fill="#9FB3C8")
    draw.text((115, 1390), "A fixed vertical benchmark", font=small, fill="#F0F4F8")
    _write_png(portrait, output / "08-portrait.png")

    _write_pdf(output)
    _write_normalized_pptx(output)

    content = {
        "schema_version": SCHEMA_VERSION,
        "cases": [_record(output, *case) for case in CASES],
        "routes": ROUTES,
    }
    manifest = {**content, "corpus_sha256": _canonical_sha256(content)}
    (output / "manifest.json").write_bytes(
        (json.dumps(manifest, ensure_ascii=False, indent=2) + "\n").encode("utf-8")
    )


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", type=Path, required=True)
    args = parser.parse_args()
    generate_corpus(args.output)
