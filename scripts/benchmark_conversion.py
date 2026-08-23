import argparse
import hashlib
import importlib
import io
import json
import math
import os
import platform as platform_module
import re
import stat
import subprocess
import sys
import time
import warnings
import zipfile
from dataclasses import dataclass
from fractions import Fraction
from pathlib import Path, PurePosixPath, PureWindowsPath


_MANIFEST_LIMIT = 64 * 1024
_ASSET_LIMIT = 12 * 1024 * 1024
_SUMMARY_LIMIT = 4 * 1024 * 1024
_PAGE_RESULT_LIMIT = 1024 * 1024
_OUTPUT_LIMIT = 256 * 1024 * 1024
_PPTX_MEMBER_LIMIT = 2048
_PPTX_MEMBER_SIZE_LIMIT = 64 * 1024 * 1024
_PPTX_TOTAL_SIZE_LIMIT = 256 * 1024 * 1024
_PPTX_COMPRESSION_RATIO_LIMIT = 200
_MIN_VISIBLE_AREA_EMU2 = 9144 * 9144
_MIN_VISIBLE_AREA_RATIO_DENOMINATOR = 10_000
_PICTURE_BLOB_LIMIT = 16 * 1024 * 1024
_PICTURE_PIXEL_LIMIT = 25_000_000
_PERFORMANCE_MAX_INTEGER = 1_000_000_000_000
_SHA256_PATTERN = re.compile(r"[0-9a-f]{64}")
_IDENTIFIER_PATTERN = re.compile(r"[A-Za-z0-9][A-Za-z0-9_.:-]{0,127}")
_TOP_FIELDS = {"schema_version", "cases", "routes", "corpus_sha256"}
_CASE_FIELDS = {"id", "kind", "path", "pages", "bytes", "sha256"}
_ROUTE_FIELDS = {"id", "cases", "pages"}
_ROUTE_IDS = ("images", "pdf", "mixed_pptx")
_ROUTE_KINDS = {"images": "image", "pdf": "pdf", "mixed_pptx": "pptx"}
_ROUTE_INPUT_COUNTS = {"images": 8, "pdf": 1, "mixed_pptx": 1}
_PERFORMANCE_MAPS = {
    "model_loads",
    "stage_runs",
    "stage_duration_ms",
    "worker_runs",
    "worker_duration_ms",
    "inference_runs",
    "inference_operations",
    "inference_duration_ms",
}
_PERFORMANCE_SCALARS = {
    "agent_runs",
    "agent_image_count",
    "agent_total_bytes",
    "agent_duration_ms",
}
_WINDOWS_RESERVED_CHARS = frozenset('<>:"/\\|?*')
_WINDOWS_DEVICE_NAMES = {"con", "prn", "aux", "nul"} | {
    f"{prefix}{number}" for prefix in ("com", "lpt") for number in range(1, 10)
} | {
    f"{prefix}{number}" for prefix in ("com", "lpt") for number in ("¹", "²", "³")
}


class BenchmarkError(RuntimeError):
    pass


@dataclass(frozen=True)
class Route:
    identifier: str
    sources: tuple[Path, ...]
    expected_pages: int


@dataclass(frozen=True)
class CorpusManifest:
    root: Path
    corpus_sha256: str
    cases: tuple[dict[str, object], ...]
    routes: tuple[Route, ...]


@dataclass(frozen=True)
class RouteExecution:
    route: Route
    run_root: Path
    output_path: Path
    returncode: int
    stdout: str
    duration_seconds: float


def _identity(metadata: os.stat_result) -> tuple[int, int, int, int, int]:
    return (
        metadata.st_dev,
        metadata.st_ino,
        metadata.st_mode,
        metadata.st_size,
        metadata.st_nlink,
    )


def _read_regular_file(path: Path, limit: int, *, require_single_link: bool) -> bytes:
    before = path.lstat()
    if (
        not stat.S_ISREG(before.st_mode)
        or before.st_size > limit
        or (require_single_link and before.st_nlink != 1)
    ):
        raise ValueError
    with path.open("rb") as handle:
        opened = os.fstat(handle.fileno())
        if _identity(opened) != _identity(before):
            raise ValueError
        payload = handle.read(limit + 1)
    after = path.lstat()
    if len(payload) > limit or _identity(after) != _identity(before):
        raise ValueError
    return payload


def _object_without_duplicate_keys(pairs: list[tuple[str, object]]) -> dict[str, object]:
    result: dict[str, object] = {}
    for key, value in pairs:
        if key in result:
            raise ValueError
        result[key] = value
    return result


def _reject_json_constant(_: str) -> object:
    raise ValueError


def _strict_positive_int(value: object) -> bool:
    return type(value) is int and value > 0


def _safe_filename(value: object) -> bool:
    if not isinstance(value, str) or not value or value in {".", ".."}:
        return False
    windows_path = PureWindowsPath(value)
    device_name = value.split(".", 1)[0].casefold()
    return (
        not any(
            character in _WINDOWS_RESERVED_CHARS
            or ord(character) < 32
            or ord(character) == 127
            for character in value
        )
        and not value.endswith((".", " "))
        and device_name not in _WINDOWS_DEVICE_NAMES
        and not PurePosixPath(value).is_absolute()
        and not windows_path.is_absolute()
        and not windows_path.drive
        and not windows_path.is_reserved()
    )


def _windows_filename_key(value: str) -> str:
    return value.casefold()


def _canonical_sha256(value: object) -> str:
    payload = json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        allow_nan=False,
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _validate_assets(
    root: Path,
    cases: list[dict[str, object]] | tuple[dict[str, object], ...],
) -> None:
    for case in cases:
        source = root / case["path"]
        asset = _read_regular_file(source, _ASSET_LIMIT, require_single_link=True)
        if (
            len(asset) != case["bytes"]
            or hashlib.sha256(asset).hexdigest() != case["sha256"]
        ):
            raise ValueError


def _load_manifest(path: Path) -> CorpusManifest:
    payload = _read_regular_file(path, _MANIFEST_LIMIT, require_single_link=False)
    manifest = json.loads(payload, object_pairs_hook=_object_without_duplicate_keys)
    if not isinstance(manifest, dict) or set(manifest) != _TOP_FIELDS:
        raise ValueError
    if type(manifest["schema_version"]) is not int or manifest["schema_version"] != 1:
        raise ValueError
    cases = manifest["cases"]
    routes = manifest["routes"]
    corpus_sha256 = manifest["corpus_sha256"]
    if not isinstance(cases, list) or not isinstance(routes, list):
        raise ValueError
    if not isinstance(corpus_sha256, str) or not _SHA256_PATTERN.fullmatch(
        corpus_sha256
    ):
        raise ValueError

    canonical_value = {
        "schema_version": manifest["schema_version"],
        "cases": cases,
        "routes": routes,
    }
    if _canonical_sha256(canonical_value) != corpus_sha256:
        raise ValueError

    root = path.resolve(strict=True).parent
    case_by_id: dict[str, dict[str, object]] = {}
    path_keys: set[str] = set()
    declared_bytes = 0
    for case in cases:
        if not isinstance(case, dict) or set(case) != _CASE_FIELDS:
            raise ValueError
        identifier = case["id"]
        kind = case["kind"]
        filename = case["path"]
        pages = case["pages"]
        size = case["bytes"]
        digest = case["sha256"]
        if not isinstance(identifier, str) or not identifier or identifier in case_by_id:
            raise ValueError
        if kind not in {"image", "pdf", "pptx"}:
            raise ValueError
        if not _safe_filename(filename):
            raise ValueError
        filename_key = _windows_filename_key(filename)
        if filename_key in path_keys:
            raise ValueError
        if not _strict_positive_int(pages) or not _strict_positive_int(size):
            raise ValueError
        if not isinstance(digest, str) or not _SHA256_PATTERN.fullmatch(digest):
            raise ValueError
        case_by_id[identifier] = case
        path_keys.add(filename_key)
        declared_bytes += size
    if declared_bytes > _ASSET_LIMIT:
        raise ValueError

    _validate_assets(root, cases)

    if len(routes) != 3:
        raise ValueError
    built_routes = []
    used_case_ids: list[str] = []
    expected_route_shapes = (("image", 8, 8), ("pdf", 1, 3), ("pptx", 1, 3))
    for route, expected_id, shape in zip(routes, _ROUTE_IDS, expected_route_shapes):
        if not isinstance(route, dict) or set(route) != _ROUTE_FIELDS:
            raise ValueError
        route_cases = route["cases"]
        pages = route["pages"]
        if route["id"] != expected_id or not isinstance(route_cases, list):
            raise ValueError
        if not _strict_positive_int(pages):
            raise ValueError
        if any(not isinstance(identifier, str) for identifier in route_cases):
            raise ValueError
        try:
            selected = [case_by_id[identifier] for identifier in route_cases]
        except KeyError as error:
            raise ValueError from error
        expected_kind, expected_count, expected_pages = shape
        if (
            len(selected) != expected_count
            or pages != expected_pages
            or any(case["kind"] != expected_kind for case in selected)
            or sum(case["pages"] for case in selected) != expected_pages
        ):
            raise ValueError
        used_case_ids.extend(route_cases)
        built_routes.append(
            Route(
                identifier=expected_id,
                sources=tuple(root / case["path"] for case in selected),
                expected_pages=pages,
            )
        )
    if len(used_case_ids) != len(set(used_case_ids)) or set(used_case_ids) != set(
        case_by_id
    ):
        raise ValueError

    return CorpusManifest(
        root=root,
        corpus_sha256=corpus_sha256,
        cases=tuple(cases),
        routes=tuple(built_routes),
    )


def load_manifest(path: Path) -> CorpusManifest:
    try:
        if path.is_dir():
            path = path / "manifest.json"
        return _load_manifest(path)
    except (KeyboardInterrupt, SystemExit):
        raise
    except Exception:
        raise BenchmarkError("invalid_corpus") from None


def _base_command() -> list[str]:
    return [sys.executable, "-I", "-B", "-m", "image2editable"]


def _clean_environment() -> dict[str, str]:
    environment = os.environ.copy()
    environment.pop("PYTHONPATH", None)
    return environment


def _run_cli(
    arguments: list[str], *, cwd: Path, timeout: float | None = None
) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [*_base_command(), *arguments],
        cwd=cwd,
        env=_clean_environment(),
        text=True,
        capture_output=True,
        timeout=timeout,
        check=False,
    )


def _require_ready_doctor(cwd: Path) -> None:
    try:
        completed = _run_cli(["doctor", "--agent-local"], cwd=cwd, timeout=180.0)
        status = json.loads(
            completed.stdout,
            object_pairs_hook=_object_without_duplicate_keys,
            parse_constant=_reject_json_constant,
        )
        ready = (
            completed.returncode == 0
            and isinstance(status, dict)
            and set(status) == {"ready", "checks"}
            and status.get("ready") is True
            and isinstance(status.get("checks"), dict)
        )
    except (KeyboardInterrupt, SystemExit):
        raise
    except Exception:
        ready = False
    if not ready:
        raise BenchmarkError("doctor_not_ready")


def execute_routes(
    manifest: CorpusManifest, result_root: Path
) -> tuple[RouteExecution, ...]:
    try:
        _validate_assets(manifest.root, manifest.cases)
    except (KeyboardInterrupt, SystemExit):
        raise
    except Exception:
        raise BenchmarkError("invalid_corpus") from None

    try:
        result_root.lstat()
    except FileNotFoundError:
        pass
    except OSError:
        raise BenchmarkError("invalid_output") from None
    else:
        raise BenchmarkError("invalid_output")
    absolute_result_root = result_root.resolve()
    try:
        absolute_result_root.mkdir()
    except OSError:
        raise BenchmarkError("invalid_output") from None

    _require_ready_doctor(absolute_result_root)
    executions = []
    for route in manifest.routes:
        run_root = absolute_result_root / route.identifier / "run"
        output_path = absolute_result_root / route.identifier / "output.pptx"
        arguments = [
            "convert",
            *(str(source) for source in route.sources),
            "--run-dir",
            str(run_root),
            "--output",
            str(output_path),
            "--slide-size",
            "16:9",
            "--agent-provider",
            "local",
        ]
        started = time.perf_counter()
        try:
            completed = _run_cli(arguments, cwd=absolute_result_root, timeout=None)
            returncode = completed.returncode
            stdout = completed.stdout
        except Exception:
            returncode = -1
            stdout = ""
        duration = time.perf_counter() - started
        if not math.isfinite(duration) or duration < 0:
            duration = 0.0
        executions.append(
            RouteExecution(
                route=route,
                run_root=run_root,
                output_path=output_path,
                returncode=returncode,
                stdout=stdout,
                duration_seconds=duration,
            )
        )
    return tuple(executions)


def _strict_json(payload: str | bytes, limit: int) -> dict[str, object]:
    encoded = payload.encode("utf-8") if isinstance(payload, str) else payload
    if len(encoded) > limit:
        raise ValueError
    value = json.loads(
        encoded,
        object_pairs_hook=_object_without_duplicate_keys,
        parse_constant=_reject_json_constant,
    )
    if not isinstance(value, dict):
        raise ValueError
    return value


def _validate_performance(
    performance: object, expected_page_ids: list[str]
) -> dict[str, object]:
    fields = _PERFORMANCE_MAPS | _PERFORMANCE_SCALARS
    if not isinstance(performance, dict) or set(performance) != {"pages"}:
        raise ValueError
    pages = performance["pages"]
    if not isinstance(pages, dict) or set(pages) != set(expected_page_ids):
        raise ValueError
    safe_pages = {}
    for page_id, page in pages.items():
        if not _IDENTIFIER_PATTERN.fullmatch(page_id):
            raise ValueError
        if not isinstance(page, dict) or set(page) != fields:
            raise ValueError
        safe_page = {}
        for field in _PERFORMANCE_MAPS:
            metrics = page[field]
            if not isinstance(metrics, dict):
                raise ValueError
            for name, value in metrics.items():
                if (
                    not isinstance(name, str)
                    or not _IDENTIFIER_PATTERN.fullmatch(name)
                    or type(value) is not int
                    or value < 0
                    or value > _PERFORMANCE_MAX_INTEGER
                ):
                    raise ValueError
            total = sum(metrics.values())
            if total > _PERFORMANCE_MAX_INTEGER:
                raise ValueError
            safe_page[field] = total
        if any(
            type(page[field]) is not int
            or page[field] < 0
            or page[field] > _PERFORMANCE_MAX_INTEGER
            for field in _PERFORMANCE_SCALARS
        ):
            raise ValueError
        safe_page.update({field: page[field] for field in _PERFORMANCE_SCALARS})
        safe_pages[page_id] = safe_page
    return {"pages": safe_pages}


def _warning_present(value: dict[str, object]) -> bool:
    warning = value.get("warning")
    warnings = value.get("warnings", [])
    if warning is not None and warning != "":
        return True
    if not isinstance(warnings, list):
        raise ValueError
    return bool(warnings)


def _validate_runtime_results(
    execution: RouteExecution, summary: dict[str, object]
) -> tuple[int, set[int]]:
    page_ids = [f"page_{index:03d}" for index in range(1, execution.route.expected_pages + 1)]
    if execution.route.identifier != "mixed_pptx":
        warning_pages = 0
        for page_id in page_ids:
            payload = _read_regular_file(
                execution.run_root / "pages" / page_id / "page_result.json",
                _PAGE_RESULT_LIMIT,
                require_single_link=True,
            )
            result = _strict_json(payload, _PAGE_RESULT_LIMIT)
            if (
                type(result.get("schema_version")) is not int
                or result["schema_version"] != 1
                or result.get("page_id") != page_id
            ):
                raise ValueError
            if result.get("status") != "validated" or _warning_present(result):
                warning_pages += 1
        return warning_pages, set(range(execution.route.expected_pages))

    results = summary.get("page_results")
    if not isinstance(results, list) or len(results) != len(page_ids):
        raise ValueError
    statuses = []
    warning_indexes = set()
    for index, (page_id, result) in enumerate(zip(page_ids, results)):
        if (
            not isinstance(result, dict)
            or type(result.get("schema_version")) is not int
            or result["schema_version"] != 1
            or result.get("page_id") != page_id
        ):
            raise ValueError
        status = result.get("status")
        if status not in {"preserved", "replaced", "preserved_with_warning"}:
            raise ValueError
        statuses.append(status)
        if status == "preserved_with_warning" or _warning_present(result):
            warning_indexes.add(index)
    summary_warnings = summary.get("warnings")
    warning_count = summary.get("preserved_with_warning_pages")
    status_warning_count = statuses.count("preserved_with_warning")
    if (
        not isinstance(summary_warnings, list)
        or type(warning_count) is not int
        or warning_count < 0
        or warning_count > len(page_ids)
        or warning_count != status_warning_count
    ):
        raise ValueError
    for index in (1, 2):
        if statuses[index] != "replaced":
            warning_indexes.add(index)
    warning_pages = max(
        len(warning_indexes), min(len(summary_warnings), len(page_ids))
    )
    reconstructed = {index for index, status in enumerate(statuses) if status == "replaced"}
    return warning_pages, reconstructed


def _iter_shapes(
    shapes: object,
    transform: tuple[Fraction, Fraction, Fraction, Fraction] | None = None,
):
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    scale_x, scale_y, offset_x, offset_y = transform or (
        Fraction(1),
        Fraction(1),
        Fraction(0),
        Fraction(0),
    )
    for shape in shapes:
        if _shape_hidden(shape):
            continue
        start_x = offset_x + shape.left * scale_x
        end_x = offset_x + (shape.left + shape.width) * scale_x
        start_y = offset_y + shape.top * scale_y
        end_y = offset_y + (shape.top + shape.height) * scale_y
        left = min(start_x, end_x)
        top = min(start_y, end_y)
        shape_width = abs(end_x - start_x)
        shape_height = abs(end_y - start_y)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            xfrm = shape._element.grpSpPr.xfrm
            if (
                shape_width <= 0
                or shape_height <= 0
                or xfrm.chExt.cx <= 0
                or xfrm.chExt.cy <= 0
            ):
                continue
            child_scale_x = (end_x - start_x) / xfrm.chExt.cx
            child_scale_y = (end_y - start_y) / xfrm.chExt.cy
            child_offset_x = start_x - xfrm.chOff.x * child_scale_x
            child_offset_y = start_y - xfrm.chOff.y * child_scale_y
            if xfrm.flipH:
                child_scale_x = -child_scale_x
                child_offset_x = end_x - xfrm.chOff.x * child_scale_x
            if xfrm.flipV:
                child_scale_y = -child_scale_y
                child_offset_y = end_y - xfrm.chOff.y * child_scale_y
            yield from _iter_shapes(
                shape.shapes,
                (
                    child_scale_x,
                    child_scale_y,
                    child_offset_x,
                    child_offset_y,
                ),
            )
        else:
            yield shape, (left, top, shape_width, shape_height)


def _shape_hidden(shape: object) -> bool:
    properties = next(
        (
            node
            for node in shape._element.iter()
            if node.tag.endswith("}cNvPr")
        ),
        None,
    )
    return properties is not None and properties.get("hidden", "").casefold() in {
        "1",
        "true",
    }


def _visible_area(
    geometry: tuple[Fraction, Fraction, Fraction, Fraction],
    width: int,
    height: int,
) -> Fraction:
    left, top, shape_width, shape_height = geometry
    intersection_width = min(left + shape_width, width) - max(left, 0)
    intersection_height = min(top + shape_height, height) - max(top, 0)
    if intersection_width <= 0 or intersection_height <= 0:
        return Fraction(0)
    return intersection_width * intersection_height


def _area_is_visible(area: Fraction, width: int, height: int) -> bool:
    return (
        area >= _MIN_VISIBLE_AREA_EMU2
        and area * _MIN_VISIBLE_AREA_RATIO_DENOMINATOR >= width * height
    )


def _alpha_is_zero(element: object) -> bool:
    for node in element.iter():
        if node.tag.rsplit("}", 1)[-1] not in {"alpha", "alphaMod", "alphaModFix"}:
            continue
        value = node.get("val")
        if value is None:
            value = node.get("amt")
        if (
            not isinstance(value, str)
            or not value
            or not value.isascii()
            or not value.isdecimal()
            or int(value) > 100_000
        ):
            return True
        if int(value) == 0:
            return True
    return False


def _visible_text(shape: object) -> bool:
    for paragraph in shape.text_frame.paragraphs:
        paragraph_properties = paragraph._p.pPr
        paragraph_transparent = (
            paragraph_properties is not None
            and _alpha_is_zero(paragraph_properties)
        )
        for run in paragraph.runs:
            if (
                run.text.strip()
                and not paragraph_transparent
                and not _alpha_is_zero(run._r)
            ):
                return True
    return False


def _picture_alpha_visible(blip: object) -> bool:
    supported = {"alphaModFix": "amt", "alphaRepl": "a"}
    for effect in blip:
        effect_type = effect.tag.rsplit("}", 1)[-1]
        if not effect_type.startswith("alpha"):
            continue
        attribute = supported.get(effect_type)
        if attribute is None or set(effect.attrib) != {attribute}:
            return False
        value = effect.get(attribute)
        if (
            not value
            or not value.isascii()
            or not value.isdecimal()
            or int(value) > 100_000
            or int(value) == 0
        ):
            return False
    return True


def _visible_picture(
    shape: object, area: Fraction, width: int, height: int
) -> bool:
    from PIL import Image

    if not _picture_alpha_visible(shape._element.blipFill.blip):
        return False
    blob = shape.image.blob
    if len(blob) > _PICTURE_BLOB_LIMIT:
        raise ValueError
    with warnings.catch_warnings():
        warnings.simplefilter("error", Image.DecompressionBombWarning)
        with Image.open(io.BytesIO(blob)) as image:
            pixel_count = image.width * image.height
            if pixel_count <= 0 or pixel_count > _PICTURE_PIXEL_LIMIT:
                raise ValueError
            image.load()
            if "A" not in image.getbands() and "transparency" not in image.info:
                return True
            alpha = image.convert("RGBA").getchannel("A")
            visible_pixels = sum(alpha.histogram()[1:])
            return _area_is_visible(
                area * visible_pixels / pixel_count, width, height
            )


def _direct_fill(element: object) -> object | None:
    fill_names = {"blipFill", "gradFill", "grpFill", "noFill", "pattFill", "solidFill"}
    return next(
        (
            child
            for child in element
            if child.tag.rsplit("}", 1)[-1] in fill_names
        ),
        None,
    )


def _fill_is_transparent(fill: object | None) -> bool:
    if fill is None:
        return False
    fill_type = fill.tag.rsplit("}", 1)[-1]
    if fill_type == "noFill":
        return True
    if fill_type == "solidFill":
        return _alpha_is_zero(fill)
    if fill_type == "gradFill":
        stops = [
            node for node in fill.iter() if node.tag.rsplit("}", 1)[-1] == "gs"
        ]
        return bool(stops) and all(_alpha_is_zero(stop) for stop in stops)
    if fill_type == "pattFill":
        colors = [
            child
            for child in fill
            if child.tag.rsplit("}", 1)[-1] in {"fgClr", "bgClr"}
        ]
        return len(colors) == 2 and all(_alpha_is_zero(color) for color in colors)
    return False


def _native_shape_visible(shape: object) -> bool:
    from pptx.enum.dml import MSO_FILL_TYPE

    properties = shape._element.spPr
    line = next(
        (child for child in properties if child.tag.endswith("}ln")), None
    )
    shape_fill = _direct_fill(properties)
    line_fill = _direct_fill(line) if line is not None else None
    fill_transparent = (
        shape.fill.type == MSO_FILL_TYPE.BACKGROUND
        or _fill_is_transparent(shape_fill)
    )
    line_transparent = (
        shape.line.fill.type == MSO_FILL_TYPE.BACKGROUND
        or _fill_is_transparent(line_fill)
    )
    return not (fill_transparent and line_transparent)


def _validate_pptx_archive(payload: bytes) -> None:
    try:
        with zipfile.ZipFile(io.BytesIO(payload)) as archive:
            members = archive.infolist()
    except Exception:
        raise ValueError from None
    if len(members) > _PPTX_MEMBER_LIMIT:
        raise ValueError
    names = set()
    total_size = 0
    for member in members:
        name = member.orig_filename
        parts = name.split("/")
        windows_path = PureWindowsPath(name)
        name_key = name.casefold()
        if (
            not name
            or "\\" in name
            or "\x00" in name
            or name.startswith("/")
            or windows_path.is_absolute()
            or windows_path.drive
            or any(not part or part in {".", ".."} for part in parts)
            or name_key in names
            or member.flag_bits & 1
            or member.file_size > _PPTX_MEMBER_SIZE_LIMIT
            or (
                member.file_size > 0
                and (
                    member.compress_size == 0
                    or member.file_size
                    > member.compress_size * _PPTX_COMPRESSION_RATIO_LIMIT
                )
            )
        ):
            raise ValueError
        names.add(name_key)
        total_size += member.file_size
        if total_size > _PPTX_TOTAL_SIZE_LIMIT:
            raise ValueError


def _editable_slide(slide: object, width: int, height: int) -> bool:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    pictures = 0
    native = False
    for shape, geometry in _iter_shapes(slide.shapes):
        left, top, shape_width, shape_height = geometry
        area = _visible_area(geometry, width, height)
        visible = _area_is_visible(area, width, height)
        if (
            visible
            and getattr(shape, "has_text_frame", False)
            and _visible_text(shape)
        ):
            return True
        if shape.shape_type in {
            MSO_SHAPE_TYPE.PICTURE,
            MSO_SHAPE_TYPE.LINKED_PICTURE,
        }:
            if visible and _visible_picture(shape, area, width, height):
                pictures += 1
            continue
        if shape.shape_type in {
            MSO_SHAPE_TYPE.TEXT_BOX,
            MSO_SHAPE_TYPE.PLACEHOLDER,
            MSO_SHAPE_TYPE.COMMENT,
            MSO_SHAPE_TYPE.SCRIPT_ANCHOR,
        }:
            continue
        if shape.shape_type in {
                MSO_SHAPE_TYPE.AUTO_SHAPE,
                MSO_SHAPE_TYPE.CALLOUT,
                MSO_SHAPE_TYPE.FREEFORM,
                MSO_SHAPE_TYPE.LINE,
        } and not _native_shape_visible(shape):
            continue
        full_page = (
            left <= width // 100
            and top <= height // 100
            and left + shape_width >= width * 99 // 100
            and top + shape_height >= height * 99 // 100
        )
        if visible and not full_page:
            native = True
    return native or pictures >= 2


def _route_result(execution: RouteExecution) -> dict[str, object]:
    return {
        "id": execution.route.identifier,
        "kind": _ROUTE_KINDS[execution.route.identifier],
        "input_count": _ROUTE_INPUT_COUNTS[execution.route.identifier],
        "pages": execution.route.expected_pages,
        "duration_ms": max(0, int(execution.duration_seconds * 1000)),
        "status": "failed",
        "error_type": None,
        "warning_pages": 0,
        "output_sha256": None,
        "performance": None,
    }


def evaluate_execution(execution: RouteExecution) -> dict[str, object]:
    result = _route_result(execution)
    if execution.returncode != 0:
        result["error_type"] = "conversion_failed"
        return result
    try:
        summary = _strict_json(execution.stdout, _SUMMARY_LIMIT)
        page_ids = [
            f"page_{index:03d}"
            for index in range(1, execution.route.expected_pages + 1)
        ]
        if (
            type(summary.get("schema_version")) is not int
            or summary["schema_version"] != 1
            or summary.get("status") != "completed"
            or type(summary.get("pages")) is not int
            or summary["pages"] != execution.route.expected_pages
            or summary.get("outputs")
            != {
                "pptx" if execution.route.identifier == "mixed_pptx" else "16:9": str(
                    execution.output_path
                )
            }
        ):
            raise ValueError
        performance = _validate_performance(summary.get("performance"), page_ids)
    except Exception:
        result["error_type"] = "invalid_summary"
        return result

    try:
        payload = _read_regular_file(
            execution.output_path, _OUTPUT_LIMIT, require_single_link=True
        )
        digest = hashlib.sha256(payload).hexdigest()
        expected_hash = summary.get("output_sha256")
        hash_matches = (
            isinstance(expected_hash, str)
            and _SHA256_PATTERN.fullmatch(expected_hash)
            and expected_hash == digest
        )
        if (
            execution.route.identifier == "mixed_pptx"
            or "output_sha256" in summary
        ) and not hash_matches:
            raise ValueError
        _validate_pptx_archive(payload)
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(payload))
        if len(presentation.slides) != execution.route.expected_pages:
            raise ValueError
        if (
            execution.route.identifier != "mixed_pptx"
            and presentation.slide_width * 9 != presentation.slide_height * 16
        ):
            raise ValueError
    except Exception:
        result["error_type"] = "output_invalid"
        return result

    try:
        warning_pages, reconstructed = _validate_runtime_results(execution, summary)
    except Exception:
        result["error_type"] = "invalid_summary"
        return result
    result["warning_pages"] = warning_pages
    if warning_pages:
        result["error_type"] = "warning_fallback"
        return result

    try:
        if any(
            not _editable_slide(
                presentation.slides[index],
                presentation.slide_width,
                presentation.slide_height,
            )
            for index in reconstructed
        ):
            result["error_type"] = "flattened_output"
            return result
    except Exception:
        result["error_type"] = "output_invalid"
        return result

    result.update(
        status="passed",
        output_sha256=digest,
        performance=performance,
    )
    return result


def _device_interface() -> str:
    try:
        torch = importlib.import_module("torch")
        if getattr(getattr(torch, "version", None), "hip", None):
            return "rocm"
        if torch.cuda.is_available():
            return "cuda"
        mps = getattr(getattr(torch, "backends", None), "mps", None)
        if mps is not None and mps.is_available():
            return "mps"
        return "cpu"
    except Exception:
        return "unknown"


def _environment() -> dict[str, str]:
    supported = {"win32", "linux", "darwin"}
    return {
        "python": platform_module.python_version(),
        "platform": sys.platform if sys.platform in supported else "other",
        "device_interface": _device_interface(),
    }


def _report_from_routes(
    manifest: CorpusManifest, routes: list[dict[str, object]]
) -> dict[str, object]:
    return {
        "schema_version": 1,
        "status": "passed" if all(route["status"] == "passed" for route in routes) else "failed",
        "corpus_sha256": manifest.corpus_sha256,
        "environment": _environment(),
        "routes": routes,
        "totals": {
            "routes": len(manifest.routes),
            "inputs": sum(_ROUTE_INPUT_COUNTS[route.identifier] for route in manifest.routes),
            "pages": sum(route.expected_pages for route in manifest.routes),
            "duration_ms": sum(route["duration_ms"] for route in routes),
            "failed_routes": sum(route["status"] == "failed" for route in routes),
            "warning_pages": sum(route["warning_pages"] for route in routes),
        },
    }


def build_report(
    manifest: CorpusManifest, executions: tuple[RouteExecution, ...]
) -> dict[str, object]:
    return _report_from_routes(
        manifest, [evaluate_execution(execution) for execution in executions]
    )


def _failure_report(
    manifest: CorpusManifest, result_root: Path, error_type: str
) -> dict[str, object]:
    routes = []
    for route in manifest.routes:
        execution = RouteExecution(
            route,
            result_root / route.identifier / "run",
            result_root / route.identifier / "output.pptx",
            1,
            "",
            0.0,
        )
        result = _route_result(execution)
        result["error_type"] = error_type
        routes.append(result)
    return _report_from_routes(manifest, routes)


def write_report(path: Path, report: dict[str, object]) -> None:
    payload = (
        json.dumps(
            report,
            ensure_ascii=False,
            indent=2,
            sort_keys=True,
            allow_nan=False,
        )
        + "\n"
    ).encode("utf-8")
    flags = os.O_WRONLY | os.O_CREAT | os.O_EXCL | getattr(os, "O_BINARY", 0)
    descriptor = os.open(path, flags, 0o600)
    with os.fdopen(descriptor, "wb") as handle:
        handle.write(payload)


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--corpus", required=True, type=Path)
    parser.add_argument("--output-dir", required=True, type=Path)
    arguments = parser.parse_args(argv)
    try:
        manifest = load_manifest(arguments.corpus)
    except BenchmarkError:
        print("invalid_corpus", file=sys.stderr)
        return 1

    try:
        executions = execute_routes(manifest, arguments.output_dir)
        report = build_report(manifest, executions)
    except BenchmarkError as error:
        if str(error) != "doctor_not_ready":
            print("benchmark_failed", file=sys.stderr)
            return 1
        report = _failure_report(manifest, arguments.output_dir, "doctor_not_ready")
    except Exception:
        report = _failure_report(manifest, arguments.output_dir, "conversion_failed")

    try:
        write_report(arguments.output_dir / "benchmark-report.json", report)
    except Exception:
        print("benchmark_failed", file=sys.stderr)
        return 1
    return 0 if report["status"] == "passed" else 1


if __name__ == "__main__":
    raise SystemExit(main())
