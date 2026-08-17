"""Build the fixed, project-generated release benchmark corpus."""

from __future__ import annotations

import argparse
from datetime import datetime
from io import BytesIO
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile, ZipInfo

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pypdf import PdfReader, PdfWriter
from reportlab.lib.pagesizes import A4, LETTER, landscape
from reportlab.lib.utils import ImageReader
from reportlab.pdfgen import canvas


ROOT = Path(__file__).resolve().parents[1]
FONT_PATH = ROOT / "benchmarks" / "release" / "fonts" / "NotoSansSC[wght].ttf"
FONT_VARIATION = "Regular"
WIDE = (1600, 900)
PORTRAIT = (1000, 1400)
PPTX_WIDTH = 12_192_000
PPTX_HEIGHT = 6_858_000
ZIP_TIME = (1980, 1, 1, 0, 0, 0)
FIXED_DATE = datetime(2020, 1, 1)


def _font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont:
    if not FONT_PATH.is_file():
        raise FileNotFoundError(f"bundled release font not found: {FONT_PATH}")
    font = ImageFont.truetype(str(FONT_PATH), size)
    font.set_variation_by_name("Bold" if bold else FONT_VARIATION)
    return font


def _save(image: Image.Image, path: Path) -> None:
    image.convert("RGB").save(path, format="PNG", optimize=False, compress_level=9)


def _title(draw: ImageDraw.ImageDraw, kicker: str, title: str, color: str) -> None:
    draw.text((86, 58), kicker, font=_font(24, bold=True), fill=color)
    draw.text((86, 102), title, font=_font(50, bold=True), fill=color)


def _image_bilingual_dashboard(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#F4F7FB")
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 255, 900), fill="#17324D")
    draw.text((45, 52), "北辰 / POLAR", font=_font(28, bold=True), fill="white")
    for index, label in enumerate(("总览 OVERVIEW", "销售 SALES", "用户 USERS", "设置 SETTINGS")):
        y = 170 + index * 82
        if index == 0:
            draw.rounded_rectangle((28, y - 14, 225, y + 48), radius=16, fill="#2E86AB")
        draw.text((48, y), label, font=_font(19), fill="#E8F1F8")
    draw.text((305, 58), "运营仪表盘 / Operations", font=_font(42, bold=True), fill="#183B56")
    cards = (("今日订单", "1,284", "+18%"), ("活跃用户", "8,902", "+7%"), ("转化率", "24.6%", "+3%"))
    for index, (label, value, delta) in enumerate(cards):
        x = 305 + index * 400
        draw.rounded_rectangle((x, 145, x + 350, 300), radius=24, fill="white", outline="#D7E2EC", width=2)
        draw.text((x + 28, 170), label, font=_font(20), fill="#627D98")
        draw.text((x + 28, 210), value, font=_font(40, bold=True), fill="#17324D")
        draw.text((x + 245, 220), delta, font=_font(19, bold=True), fill="#1F9D77")
    draw.rounded_rectangle((305, 345, 1115, 805), radius=24, fill="white", outline="#D7E2EC", width=2)
    draw.text((338, 375), "每周收入 / Weekly revenue", font=_font(24, bold=True), fill="#183B56")
    points = [(360, 705), (460, 640), (560, 670), (660, 545), (760, 585), (860, 475), (960, 510), (1060, 415)]
    for y in (470, 570, 670):
        draw.line((360, y, 1060, y), fill="#E6EEF4", width=2)
    draw.line(points, fill="#2E86AB", width=8, joint="curve")
    for x, y in points:
        draw.ellipse((x - 8, y - 8, x + 8, y + 8), fill="#F6AE2D")
    draw.rounded_rectangle((1155, 345, 1515, 805), radius=24, fill="#102A43")
    draw.text((1190, 380), "区域 / REGION", font=_font(20, bold=True), fill="#9FB3C8")
    for index, (label, value) in enumerate((("华东 East", 76), ("华南 South", 61), ("华北 North", 48), ("海外 Global", 35))):
        y = 455 + index * 77
        draw.text((1190, y), label, font=_font(18), fill="white")
        draw.rounded_rectangle((1190, y + 31, 1440, y + 43), radius=6, fill="#334E68")
        draw.rounded_rectangle((1190, y + 31, 1190 + value * 3, y + 43), radius=6, fill="#F6AE2D")
    _save(image, path)


def _image_dense_comparison(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#FFFDF8")
    draw = ImageDraw.Draw(image)
    _title(draw, "SYSTEM MATRIX / 参数矩阵", "模型参数对比", "#243B53")
    left, top, row_h = 80, 205, 72
    columns = (360, 270, 270, 270, 270)
    headers = ("PARAMETER", "ALPHA", "BETA", "GAMMA", "DELTA")
    rows = (
        ("Context window", "32K", "64K", "128K", "64K"),
        ("Latency p95", "860 ms", "720 ms", "940 ms", "690 ms"),
        ("Accuracy", "91.2", "93.8", "95.1", "92.6"),
        ("Languages", "18", "28", "42", "24"),
        ("Tool calls", "Basic", "Full", "Full", "Basic"),
        ("Cost index", "1.0x", "1.3x", "1.8x", "0.9x"),
        ("Release", "Stable", "Stable", "Preview", "Stable"),
        ("Recommended", "Mobile", "General", "Research", "Batch"),
    )
    x_positions = [left]
    for width in columns:
        x_positions.append(x_positions[-1] + width)
    draw.rounded_rectangle((left, top, x_positions[-1], top + row_h), radius=18, fill="#17324D")
    for index, text in enumerate(headers):
        draw.text((x_positions[index] + 20, top + 22), text, font=_font(18, bold=True), fill="white")
    for row_index, row in enumerate(rows, start=1):
        y = top + row_index * row_h
        fill = "#F2F6FA" if row_index % 2 else "#FFFFFF"
        draw.rectangle((left, y, x_positions[-1], y + row_h), fill=fill)
        for column_index, text in enumerate(row):
            color = "#17324D" if column_index == 0 else "#486581"
            draw.text((x_positions[column_index] + 20, y + 22), text, font=_font(17, bold=column_index == 0), fill=color)
        draw.line((left, y + row_h, x_positions[-1], y + row_h), fill="#D9E2EC", width=2)
    for x in x_positions:
        draw.line((x, top, x, top + row_h * 9), fill="#D9E2EC", width=2)
    _save(image, path)


def _image_profiles(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#EEF3F8")
    draw = ImageDraw.Draw(image)
    _title(draw, "PEOPLE / TEAM", "把专业能力变成协作势能", "#17324D")
    profiles = (
        ("林岚 LIN LAN", "Product Strategy", "#2E86AB", "LL"),
        ("MAYA CHEN", "Data & Insight", "#D1495B", "MC"),
        ("NOAH WU", "Design Systems", "#F6AE2D", "NW"),
    )
    for index, (name, role, accent, initials) in enumerate(profiles):
        x = 90 + index * 500
        draw.rounded_rectangle((x, 230, x + 430, 790), radius=32, fill="white", outline="#D5E1EB", width=2)
        draw.ellipse((x + 105, 280, x + 325, 500), fill=accent)
        draw.text((x + 150, 352), initials, font=_font(42, bold=True), fill="white")
        draw.text((x + 38, 545), name, font=_font(25, bold=True), fill="#17324D")
        draw.text((x + 38, 592), role, font=_font(20), fill="#627D98")
        draw.line((x + 38, 645, x + 392, 645), fill="#D9E2EC", width=2)
        for tag_index, tag in enumerate(("DISCOVER", "BUILD", "SHIP")):
            tx = x + 38 + tag_index * 118
            draw.rounded_rectangle((tx, 690, tx + 104, 730), radius=18, fill="#EAF2F8")
            draw.text((tx + 12, 701), tag, font=_font(13, bold=True), fill="#2E86AB")
    _save(image, path)


def _image_timeline(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#FAF7F2")
    draw = ImageDraw.Draw(image)
    _title(draw, "ROADMAP 2026", "从信号到规模", "#243B53")
    draw.line((150, 475, 1450, 475), fill="#BCCCDC", width=8)
    stages = (
        (210, "01", "DISCOVER", "明确机会\n梳理约束", "#2E86AB"),
        (590, "02", "PROTOTYPE", "验证路径\n建立基线", "#F6AE2D"),
        (970, "03", "PILOT", "真实试点\n量化反馈", "#D1495B"),
        (1350, "04", "SCALE", "稳定发布\n持续改进", "#1F9D77"),
    )
    for index, (x, number, label, detail, accent) in enumerate(stages):
        draw.ellipse((x - 42, 433, x + 42, 517), fill=accent, outline="white", width=8)
        draw.text((x - 21, 450), number, font=_font(21, bold=True), fill="white")
        top = 280 if index % 2 == 0 else 570
        draw.rounded_rectangle((x - 135, top, x + 135, top + 145), radius=22, fill="white", outline=accent, width=3)
        draw.text((x - 105, top + 24), label, font=_font(20, bold=True), fill=accent)
        draw.multiline_text((x - 105, top + 62), detail, font=_font(18), fill="#486581", spacing=9)
    _save(image, path)


def _image_combo_chart(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#F7FAFC")
    draw = ImageDraw.Draw(image)
    _title(draw, "GROWTH SIGNAL", "渠道收入与转化趋势", "#17324D")
    left, right, bottom = 150, 1460, 770
    for index in range(6):
        y = bottom - index * 90
        draw.line((left, y, right, y), fill="#D9E2EC", width=2)
        draw.text((95, y - 13), str(index * 20), font=_font(16), fill="#829AB1")
    values = (42, 56, 68, 74, 91, 104)
    labels = ("JAN", "FEB", "MAR", "APR", "MAY", "JUN")
    points = []
    for index, value in enumerate(values):
        x = 220 + index * 210
        height = value * 4
        draw.rounded_rectangle((x, bottom - height, x + 105, bottom), radius=12, fill="#2E86AB")
        draw.text((x + 26, bottom + 22), labels[index], font=_font(16, bold=True), fill="#486581")
        points.append((x + 52, bottom - (value + 14) * 4))
    draw.line(points, fill="#D1495B", width=9, joint="curve")
    for x, y in points:
        draw.ellipse((x - 10, y - 10, x + 10, y + 10), fill="#F6AE2D", outline="#D1495B", width=4)
    draw.rectangle((1160, 92, 1190, 116), fill="#2E86AB")
    draw.text((1205, 91), "Revenue", font=_font(17), fill="#486581")
    draw.line((1320, 104, 1365, 104), fill="#D1495B", width=7)
    draw.text((1380, 91), "Conversion", font=_font(17), fill="#486581")
    _save(image, path)


def _image_flowchart(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#FFFFFF")
    draw = ImageDraw.Draw(image)
    _title(draw, "DECISION FLOW", "可编辑重建路径", "#17324D")
    nodes = (
        ((90, 340, 330, 485), "INPUT\nPNG / PDF", "#D6EAF8"),
        ((455, 340, 695, 485), "ANALYZE\nLayout + OCR", "#D5F5E3"),
        ((820, 220, 1060, 365), "NATIVE\nText + Shape", "#FCE7C1"),
        ((820, 520, 1060, 665), "RASTER\nScreenshot", "#F8D7DA"),
        ((1190, 340, 1510, 485), "VALIDATE\nEditable PPTX", "#E4D7F5"),
    )
    links = (((330, 412), (455, 412)), ((695, 412), (820, 292)), ((695, 412), (820, 592)), ((1060, 292), (1190, 412)), ((1060, 592), (1190, 412)))
    for start, end in links:
        draw.line((*start, *end), fill="#829AB1", width=8)
        draw.polygon(((end[0], end[1]), (end[0] - 24, end[1] - 17), (end[0] - 24, end[1] + 17)), fill="#829AB1")
    for box, label, fill in nodes:
        draw.rounded_rectangle(box, radius=25, fill=fill, outline="#334E68", width=3)
        draw.multiline_text((box[0] + 28, box[1] + 35), label, font=_font(21, bold=True), fill="#17324D", spacing=10)
    draw.text((748, 190), "high confidence", font=_font(16, bold=True), fill="#B7791F")
    draw.text((748, 690), "visual candidate", font=_font(16, bold=True), fill="#C53030")
    _save(image, path)


def _image_icon_matrix(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#102A43")
    draw = ImageDraw.Draw(image)
    _title(draw, "CAPABILITY MAP", "十二项能力矩阵", "#F0F4F8")
    labels = ("SEARCH", "DESIGN", "BUILD", "TEST", "SHIP", "LEARN", "PLAN", "WRITE", "CHART", "SYNC", "SECURE", "SCALE")
    colors = ("#2E86AB", "#F6AE2D", "#D1495B", "#1F9D77")
    for index, label in enumerate(labels):
        row, col = divmod(index, 4)
        x, y = 100 + col * 370, 230 + row * 190
        draw.rounded_rectangle((x, y, x + 310, y + 145), radius=25, fill="#17324D", outline="#486581", width=2)
        accent = colors[index % 4]
        draw.ellipse((x + 25, y + 34, x + 101, y + 110), fill=accent)
        if index % 3 == 0:
            draw.line((x + 46, y + 73, x + 80, y + 73), fill="white", width=7)
        elif index % 3 == 1:
            draw.rectangle((x + 49, y + 51, x + 78, y + 92), outline="white", width=5)
        else:
            draw.polygon(((x + 63, y + 47), (x + 84, y + 91), (x + 42, y + 91)), outline="white")
        draw.text((x + 125, y + 55), label, font=_font(20, bold=True), fill="#F0F4F8")
    _save(image, path)


def _image_gradient(path: Path) -> None:
    image = Image.new("RGB", WIDE)
    pixels = image.load()
    left, right = (44, 62, 116), (120, 74, 142)
    for x in range(WIDE[0]):
        ratio = x / (WIDE[0] - 1)
        color = tuple(round(a + (b - a) * ratio) for a, b in zip(left, right, strict=True))
        for y in range(WIDE[1]):
            pixels[x, y] = color
    draw = ImageDraw.Draw(image, "RGBA")
    draw.ellipse((980, -260, 1700, 460), fill=(246, 174, 45, 85))
    draw.ellipse((-240, 540, 440, 1220), fill=(46, 134, 171, 110))
    draw.text((125, 115), "LIGHT / 08", font=_font(22, bold=True), fill="#F7EDEA")
    draw.multiline_text((125, 245), "Clarity lives\nin the contrast.", font=_font(72, bold=True), fill="white", spacing=12)
    draw.text((130, 655), "浅色文字 × 固定渐变 × 大留白", font=_font(28), fill="#F8E9F0")
    draw.line((130, 735, 550, 735), fill=(255, 255, 255, 180), width=3)
    _save(image, path)


def _image_network(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#F9FBFD")
    draw = ImageDraw.Draw(image)
    _title(draw, "CONNECTED SYSTEM", "细线网络与多层关系", "#17324D")
    points = ((170, 460), (370, 290), (450, 620), (650, 420), (830, 250), (900, 650), (1080, 430), (1280, 270), (1410, 590), (1180, 700), (620, 720), (250, 730))
    edges = ((0, 1), (0, 2), (1, 3), (2, 3), (2, 11), (3, 4), (3, 5), (3, 10), (4, 6), (4, 7), (5, 6), (5, 9), (6, 7), (6, 8), (6, 9), (8, 9), (9, 10), (10, 11))
    for start, end in edges:
        draw.line((*points[start], *points[end]), fill="#9FB3C8", width=2)
    for index, (x, y) in enumerate(points):
        radius = 24 if index in (3, 6) else 14
        fill = "#D1495B" if index in (3, 6) else "#2E86AB"
        draw.ellipse((x - radius, y - radius, x + radius, y + radius), fill=fill, outline="white", width=4)
    draw.text((585, 375), "CORE A", font=_font(18, bold=True), fill="#D1495B")
    draw.text((1015, 380), "CORE B", font=_font(18, bold=True), fill="#D1495B")
    draw.text((1120, 800), "12 nodes / 18 edges / 2 hubs", font=_font(19), fill="#627D98")
    _save(image, path)


def _image_tiny_table(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#FFFFFF")
    draw = ImageDraw.Draw(image)
    _title(draw, "OPERATIONS LEDGER", "小元素密集表格", "#243B53")
    left, top, row_h = 70, 205, 48
    widths = (120, 330, 175, 170, 170, 180, 290)
    headers = ("ID", "WORKSTREAM", "OWNER", "START", "DUE", "STATUS", "SIGNAL")
    rows = tuple((f"R-{index:03d}", f"Release track {index}", ("Lan", "Maya", "Noah")[index % 3], f"08/{index + 1:02d}", f"09/{index + 4:02d}", ("Ready", "Review", "Blocked")[index % 3], ("Stable", "Watch", "Risk")[index % 3]) for index in range(12))
    xs = [left]
    for width in widths:
        xs.append(xs[-1] + width)
    draw.rectangle((left, top, xs[-1], top + row_h), fill="#17324D")
    for column, text in enumerate(headers):
        draw.text((xs[column] + 10, top + 15), text, font=_font(14, bold=True), fill="white")
    for row_index, row in enumerate(rows, start=1):
        y = top + row_index * row_h
        draw.rectangle((left, y, xs[-1], y + row_h), fill="#F4F7FA" if row_index % 2 else "white")
        for column, text in enumerate(row):
            color = "#C53030" if text in ("Blocked", "Risk") else "#486581"
            draw.text((xs[column] + 10, y + 15), text, font=_font(13, bold=column in (0, 5)), fill=color)
        draw.line((left, y + row_h, xs[-1], y + row_h), fill="#D9E2EC")
    for x in xs:
        draw.line((x, top, x, top + row_h * 13), fill="#D9E2EC")
    _save(image, path)


def _image_dark_poster(path: Path) -> None:
    image = Image.new("RGB", WIDE, "#080D16")
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 38, 900), fill="#F6AE2D")
    draw.ellipse((1050, 120, 1550, 620), outline="#2E86AB", width=18)
    draw.ellipse((1160, 230, 1440, 510), outline="#D1495B", width=12)
    draw.line((1020, 700, 1510, 210), fill="#F6AE2D", width=5)
    draw.text((110, 85), "RELEASE / 020", font=_font(24, bold=True), fill="#F6AE2D")
    draw.multiline_text((110, 220), "MAKE IT\nEDITABLE.", font=_font(92, bold=True), fill="#F0F4F8", spacing=4)
    draw.text((118, 690), "Fixed corpus for cross-format reliability", font=_font(27), fill="#9FB3C8")
    draw.text((118, 752), "30 PAGES  /  18 INPUTS  /  CC0", font=_font(18, bold=True), fill="#F0F4F8")
    _save(image, path)


def _image_portrait_infographic(path: Path) -> None:
    image = Image.new("RGB", PORTRAIT, "#F3F6F8")
    draw = ImageDraw.Draw(image)
    draw.rectangle((0, 0, 1000, 250), fill="#17324D")
    draw.text((70, 55), "RELEASE MAP", font=_font(24, bold=True), fill="#F6AE2D")
    draw.multiline_text((70, 102), "30 pages.\nOne strict gate.", font=_font(45, bold=True), fill="white", spacing=4)
    sections = (("12", "IMAGES", "Raster composition diversity", "#2E86AB"), ("06", "PDF PAGES", "Size, rotation, high DPI", "#D1495B"), ("12", "PPTX PAGES", "Image-only and mixed objects", "#F6AE2D"))
    for index, (number, label, detail, accent) in enumerate(sections):
        y = 315 + index * 285
        draw.rounded_rectangle((60, y, 940, y + 225), radius=30, fill="white", outline="#D5E1EB", width=2)
        draw.ellipse((95, y + 45, 245, y + 195), fill=accent)
        draw.text((135, y + 89), number, font=_font(38, bold=True), fill="white")
        draw.text((295, y + 48), label, font=_font(26, bold=True), fill="#17324D")
        draw.text((295, y + 102), detail, font=_font(20), fill="#627D98")
        for step in range(4):
            x = 295 + step * 125
            draw.rounded_rectangle((x, y + 160, x + 95, y + 178), radius=9, fill=accent if step <= index else "#D9E2EC")
    draw.line((70, 1230, 930, 1230), fill="#BCCCDC", width=3)
    draw.text((70, 1275), "PROJECT-GENERATED  /  CC0-1.0", font=_font(20, bold=True), fill="#486581")
    draw.text((70, 1320), "1000 × 1400 portrait fixture", font=_font(17), fill="#829AB1")
    _save(image, path)


def _pdf_canvas(path: Path, page_size: tuple[float, float]) -> canvas.Canvas:
    document = canvas.Canvas(str(path), pagesize=page_size, invariant=1, pageCompression=1)
    document.setAuthor("image2editable")
    document.setCreator("scripts/build_release_corpus.py")
    document.setSubject("CC0 release benchmark fixture")
    document.setTitle(path.stem)
    return document


def _draw_pdf_page(document: canvas.Canvas, size: tuple[float, float], title: str, accent: tuple[float, float, float]) -> None:
    width, height = size
    document.setFillColorRGB(0.96, 0.97, 0.98)
    document.rect(0, 0, width, height, stroke=0, fill=1)
    document.setFillColorRGB(*accent)
    document.roundRect(36, height - 118, width - 72, 70, 14, stroke=0, fill=1)
    document.setFillColorRGB(1, 1, 1)
    document.setFont("Helvetica-Bold", 24)
    document.drawString(58, height - 88, title)
    for index in range(3):
        x = 42 + index * ((width - 96) / 3)
        document.setFillColorRGB(1, 1, 1)
        document.roundRect(x, height - 300, (width - 126) / 3, 130, 10, stroke=0, fill=1)
        document.setFillColorRGB(0.10, 0.20, 0.30)
        document.setFont("Helvetica-Bold", 15)
        document.drawString(x + 18, height - 210, f"SECTION {index + 1}")
        document.setFont("Helvetica", 10)
        document.drawString(x + 18, height - 238, "Fixed geometry and metadata")
        document.drawString(x + 18, height - 258, "for release verification")
    document.setFillColorRGB(0.20, 0.30, 0.40)
    document.setFont("Helvetica", 11)
    document.drawString(42, 38, f"{round(width)} x {round(height)} pt / project-generated / CC0-1.0")


def _write_pdfs(output_root: Path) -> None:
    mixed = _pdf_canvas(output_root / "13-mixed-page-sizes.pdf", landscape(A4))
    _draw_pdf_page(mixed, landscape(A4), "A4 LANDSCAPE / MIXED SIZE", (0.18, 0.53, 0.67))
    mixed.showPage()
    mixed.setPageSize(LETTER)
    _draw_pdf_page(mixed, LETTER, "LETTER PORTRAIT / MIXED SIZE", (0.82, 0.29, 0.36))
    mixed.showPage()
    mixed.save()

    unrotated = output_root / "14-rotated-source.pdf"
    rotated_canvas = _pdf_canvas(unrotated, LETTER)
    _draw_pdf_page(rotated_canvas, LETTER, "PAGE 1 / STANDARD", (0.18, 0.53, 0.67))
    rotated_canvas.showPage()
    _draw_pdf_page(rotated_canvas, LETTER, "PAGE 2 / ROTATE 90", (0.96, 0.68, 0.18))
    rotated_canvas.showPage()
    rotated_canvas.save()
    reader = PdfReader(unrotated)
    writer = PdfWriter()
    writer.add_page(reader.pages[0])
    writer.add_page(reader.pages[1].rotate(90))
    writer.add_metadata({"/Author": "image2editable", "/Creator": "scripts/build_release_corpus.py", "/Title": "14-rotated-page", "/Subject": "CC0 release benchmark fixture"})
    with (output_root / "14-rotated-page.pdf").open("wb") as stream:
        writer.write(stream)
    unrotated.unlink()

    high_dpi = _pdf_canvas(output_root / "15-high-dpi.pdf", landscape(LETTER))
    source_names = ("01-bilingual-dashboard.png", "11-dark-poster.png")
    for page_number, name in enumerate(source_names, start=1):
        with Image.open(output_root / name) as source:
            enlarged = source.resize((2400, 1350), Image.Resampling.LANCZOS)
            buffer = BytesIO()
            enlarged.save(buffer, format="PNG", compress_level=9)
        buffer.seek(0)
        high_dpi.drawImage(ImageReader(buffer), 0, 0, width=landscape(LETTER)[0], height=landscape(LETTER)[1])
        high_dpi.setFillColorRGB(1, 1, 1)
        high_dpi.setFont("Helvetica-Bold", 13)
        high_dpi.drawRightString(755, 24, f"HIGH DPI {page_number} / 2400 x 1350")
        high_dpi.showPage()
    high_dpi.save()


def _set_pptx_properties(deck: Presentation, title: str) -> None:
    deck.slide_width = PPTX_WIDTH
    deck.slide_height = PPTX_HEIGHT
    properties = deck.core_properties
    properties.author = "image2editable"
    properties.created = FIXED_DATE
    properties.last_modified_by = "image2editable"
    properties.modified = FIXED_DATE
    properties.revision = 1
    properties.subject = "CC0 release benchmark fixture"
    properties.title = title


def _set_text(shape, text: str, size: int, color: str, *, bold: bool = False) -> None:
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = False
    frame.margin_left = Inches(0.14)
    frame.margin_right = Inches(0.14)
    frame.margin_top = Inches(0.08)
    frame.margin_bottom = Inches(0.08)
    paragraph = frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor.from_string(color)


def _add_picture(slide, path: Path, left: float, top: float, width: float, height: float) -> None:
    slide.shapes.add_picture(str(path), Inches(left), Inches(top), Inches(width), Inches(height))


def _canonicalize_pptx(path: Path) -> None:
    normalized = BytesIO()
    with ZipFile(path, "r") as source, ZipFile(normalized, "w", compression=ZIP_DEFLATED) as target:
        for source_info in sorted(source.infolist(), key=lambda item: item.filename):
            info = ZipInfo(source_info.filename, ZIP_TIME)
            info.compress_type = ZIP_DEFLATED
            info.create_system = 0
            info.external_attr = 0o600 << 16
            target.writestr(info, source.read(source_info.filename))
    path.write_bytes(normalized.getvalue())


def _save_pptx(deck: Presentation, path: Path) -> None:
    deck.save(path)
    _canonicalize_pptx(path)


def _write_image_only_pptx(output_root: Path) -> None:
    deck = Presentation()
    _set_pptx_properties(deck, "Image-only release fixture")
    blank = deck.slide_layouts[6]
    for name in ("01-bilingual-dashboard.png", "04-four-stage-timeline.png", "05-combo-chart.png", "11-dark-poster.png"):
        slide = deck.slides.add_slide(blank)
        _add_picture(slide, output_root / name, 0, 0, 13.333333, 7.5)
    _save_pptx(deck, output_root / "16-image-only.pptx")


def _write_mixed_native_pptx(output_root: Path) -> None:
    deck = Presentation()
    _set_pptx_properties(deck, "Mixed native release fixture")
    blank = deck.slide_layouts[6]
    content = (("NATIVE SIGNALS", "01-bilingual-dashboard.png"), ("STRUCTURED TEAMS", "03-profile-cards.png"), ("DECISION PATH", "06-flowchart.png"), ("RELEASE MATRIX", "10-tiny-element-table.png"))
    for page_number, (title, image_name) in enumerate(content, start=1):
        slide = deck.slides.add_slide(blank)
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(244, 247, 251)
        title_box = slide.shapes.add_textbox(Inches(0.65), Inches(0.38), Inches(6.0), Inches(0.62))
        _set_text(title_box, title, 30, "17324D", bold=True)
        _add_picture(slide, output_root / image_name, 6.75, 1.22, 5.9, 3.31875)
        for index, (label, value, color) in enumerate((("PAGE", f"0{page_number}", "2E86AB"), ("MODE", "NATIVE", "D1495B"), ("STATE", "EDITABLE", "1F9D77"))):
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.65), Inches(1.45 + index * 1.32), Inches(5.45), Inches(0.92))
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor.from_string("FFFFFF")
            box.line.color.rgb = RGBColor.from_string(color)
            _set_text(box, f"{label}   {value}", 20, color, bold=True)
        footer = slide.shapes.add_textbox(Inches(6.75), Inches(5.15), Inches(5.9), Inches(0.8))
        _set_text(footer, "Native text + vector shapes + raster candidate", 17, "486581")
    _save_pptx(deck, output_root / "17-mixed-native.pptx")


def _write_screenshot_candidates_pptx(output_root: Path) -> None:
    deck = Presentation()
    _set_pptx_properties(deck, "Mixed screenshot candidate fixture")
    blank = deck.slide_layouts[6]
    names = ("02-dense-parameter-comparison.png", "07-icon-matrix.png", "08-light-text-gradient.png", "12-non-16-9-infographic.png")
    for page_number, name in enumerate(names, start=1):
        slide = deck.slides.add_slide(blank)
        if name == "12-non-16-9-infographic.png":
            _add_picture(slide, output_root / name, 3.99, 0, 5.36, 7.5)
        else:
            _add_picture(slide, output_root / name, 0, 0, 13.333333, 7.5)
        banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.45), Inches(0.48), Inches(1.3), Inches(0.82))
        banner.fill.solid()
        banner.fill.fore_color.rgb = RGBColor(8, 13, 22)
        banner.fill.transparency = 12
        banner.line.fill.background()
        _set_text(banner, f"0{page_number}", 18, "FFFFFF", bold=True)
    _save_pptx(deck, output_root / "18-mixed-screenshot-candidates.pptx")


def build(output_root: Path) -> None:
    """Create all 18 inputs in a new directory."""
    output_root = Path(output_root)
    output_root.mkdir(parents=True, exist_ok=False)
    builders = (
        ("01-bilingual-dashboard.png", _image_bilingual_dashboard),
        ("02-dense-parameter-comparison.png", _image_dense_comparison),
        ("03-profile-cards.png", _image_profiles),
        ("04-four-stage-timeline.png", _image_timeline),
        ("05-combo-chart.png", _image_combo_chart),
        ("06-flowchart.png", _image_flowchart),
        ("07-icon-matrix.png", _image_icon_matrix),
        ("08-light-text-gradient.png", _image_gradient),
        ("09-thin-line-network.png", _image_network),
        ("10-tiny-element-table.png", _image_tiny_table),
        ("11-dark-poster.png", _image_dark_poster),
        ("12-non-16-9-infographic.png", _image_portrait_infographic),
    )
    for name, builder in builders:
        builder(output_root / name)
    _write_pdfs(output_root)
    _write_image_only_pptx(output_root)
    _write_mixed_native_pptx(output_root)
    _write_screenshot_candidates_pptx(output_root)


def main(argv: list[str] | None = None) -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("output_root", type=Path)
    args = parser.parse_args(argv)
    build(args.output_root)


if __name__ == "__main__":
    main()
