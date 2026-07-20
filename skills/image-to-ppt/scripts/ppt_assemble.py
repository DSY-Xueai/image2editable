#!/usr/bin/env python3
"""PPTX assembly module — compose background, foreground components, and text.

Builds a PowerPoint presentation with layered structure:
  - Bottom: repaired background image (full slide)
  - Middle: independent transparent foreground component images
  - Top: editable text boxes with font styling

Usage:
    from ppt_assemble import assemble_pptx
    assemble_pptx(bg_path, components, text_items, img_w, img_h, output_path)
"""

from __future__ import annotations

import logging
from dataclasses import dataclass
from numbers import Integral
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# Slide width in inches (standard widescreen reference)
SLIDE_WIDTH_INCHES = 40 / 3
SLIDE_HEIGHT_INCHES = 7.5


@dataclass(frozen=True)
class ContainTransform:
    slide_width: float
    slide_height: float
    content_width: float
    content_height: float
    offset_x: float
    offset_y: float


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------


def compute_slide_transform(
    img_width: int,
    img_height: int,
    slide_size: str,
) -> ContainTransform:
    """Map an image to its original aspect ratio or a fixed widescreen slide."""
    if slide_size not in {"original", "16:9"}:
        raise ValueError("slide_size must be 'original' or '16:9'")

    scale = min(
        SLIDE_WIDTH_INCHES / img_width,
        SLIDE_HEIGHT_INCHES / img_height,
    )
    content_width = img_width * scale
    content_height = img_height * scale
    if slide_size == "original":
        if min(content_width, content_height) < 1.0:
            size_scale = 1.0 / min(content_width, content_height)
            content_width *= size_scale
            content_height *= size_scale
        if max(content_width, content_height) > 56.0 + 1e-9:
            raise ValueError(
                "original slide aspect ratio exceeds PowerPoint's 1-56 inch range"
            )
        return ContainTransform(
            slide_width=content_width,
            slide_height=content_height,
            content_width=content_width,
            content_height=content_height,
            offset_x=0,
            offset_y=0,
        )

    return ContainTransform(
        slide_width=SLIDE_WIDTH_INCHES,
        slide_height=SLIDE_HEIGHT_INCHES,
        content_width=content_width,
        content_height=content_height,
        offset_x=(SLIDE_WIDTH_INCHES - content_width) / 2,
        offset_y=(SLIDE_HEIGHT_INCHES - content_height) / 2,
    )


def compute_contain_transform(img_width: int, img_height: int) -> ContainTransform:
    """Map an image into the center of a fixed widescreen slide."""
    return compute_slide_transform(img_width, img_height, "16:9")


def assemble_pptx(
    background_path: str | Path,
    components: list[dict],
    text_items: list[dict],
    img_width: int,
    img_height: int,
    output_path: str | Path,
    add_reference_slide: bool = False,
    original_image_path: str | Path | None = None,
    slide_size: str = "16:9",
    canvas_width: int | None = None,
    canvas_height: int | None = None,
    content_offset_x: int = 0,
    content_offset_y: int = 0,
) -> str:
    """Assemble a PPTX from background, foreground components, and text.

    Args:
        background_path: Path to the clean background PNG.
        components: List of component dicts (path, x, y, w, h, area).
        text_items: List of text dicts (box, text, font_size, color, bold, font, align).
        img_width: Original image width in pixels.
        img_height: Original image height in pixels.
        output_path: Where to save the PPTX.
        add_reference_slide: If True, add a second slide with the original image.
        original_image_path: Path to original image (for reference slide).
        canvas_width: Optional exact-16:9 canvas width for widescreen placement.
        canvas_height: Optional exact-16:9 canvas height for widescreen placement.
        content_offset_x: Source-image horizontal offset inside the canvas.
        content_offset_y: Source-image vertical offset inside the canvas.

    Returns:
        The output path as string.
    """
    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    use_canvas = False
    if slide_size == "16:9":
        use_canvas = _validate_canvas(
            img_width,
            img_height,
            canvas_width,
            canvas_height,
            content_offset_x,
            content_offset_y,
        )
    if use_canvas:
        transform = compute_slide_transform(canvas_width, canvas_height, slide_size)
    else:
        transform = compute_slide_transform(img_width, img_height, slide_size)

    prs.slide_width = Inches(transform.slide_width)
    prs.slide_height = Inches(transform.slide_height)
    prs._element.sldSz.set(
        "type",
        "screen16x9" if slide_size == "16:9" else "custom",
    )

    # Use blank layout
    blank_layout = prs.slide_layouts[6]

    # --- Main slide ---
    slide = prs.slides.add_slide(blank_layout)

    # Layer 1 (bottom): Background image
    slide.shapes.add_picture(
        str(background_path), 0, 0,
        Inches(transform.slide_width), Inches(transform.slide_height)
    )
    logger.info("Added background layer.")

    # Layer 2 (middle): Foreground components
    for comp in components:
        _add_component(
            slide,
            comp,
            img_width,
            img_height,
            transform,
            canvas_width if use_canvas else None,
            canvas_height if use_canvas else None,
            content_offset_x if use_canvas else 0,
            content_offset_y if use_canvas else 0,
        )

    logger.info("Added %d foreground components.", len(components))

    # Layer 3 (top): Editable text boxes
    for item in text_items:
        _add_textbox(
            slide,
            item,
            img_width,
            img_height,
            transform,
            canvas_width if use_canvas else None,
            canvas_height if use_canvas else None,
            content_offset_x if use_canvas else 0,
            content_offset_y if use_canvas else 0,
        )

    logger.info("Added %d text boxes.", len(text_items))

    # --- Reference slide (optional) ---
    if add_reference_slide and original_image_path:
        original_image_path = Path(original_image_path)
        if original_image_path.exists():
            ref_slide = prs.slides.add_slide(blank_layout)
            ref_slide.shapes.add_picture(
                str(background_path), 0, 0,
                Inches(transform.slide_width), Inches(transform.slide_height)
            )
            left, top, width, height = _map_bbox(
                0,
                0,
                img_width,
                img_height,
                img_width,
                img_height,
                transform,
                canvas_width if use_canvas else None,
                canvas_height if use_canvas else None,
                content_offset_x if use_canvas else 0,
                content_offset_y if use_canvas else 0,
            )
            ref_slide.shapes.add_picture(
                str(original_image_path),
                Inches(left), Inches(top), Inches(width), Inches(height)
            )
            logger.info("Added reference slide with original image.")

    prs.save(str(output_path))
    logger.info("Saved PPTX: %s", output_path)

    return str(output_path)


def assemble_pptx_multi(
    slides_data: list[dict],
    output_path: str | Path,
    add_reference: bool = False,
) -> str:
    """Assemble a multi-slide PPTX from multiple images' data.

    Each entry in slides_data is a dict with keys:
        background_path, components, text_items, img_width, img_height,
        original_image_path (optional), canvas_width, canvas_height,
        content_offset_x, content_offset_y (all canvas fields optional).

    Args:
        slides_data: List of per-image data dicts.
        output_path: Where to save the PPTX.
        add_reference: If True, add a reference slide after each content slide.

    Returns:
        The output path as string.
    """
    if not slides_data:
        raise ValueError("slides_data must not be empty")

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation()

    prs.slide_width = Inches(SLIDE_WIDTH_INCHES)
    prs.slide_height = Inches(SLIDE_HEIGHT_INCHES)
    prs._element.sldSz.set("type", "screen16x9")

    blank_layout = prs.slide_layouts[6]

    for idx, data in enumerate(slides_data):
        img_w = data["img_width"]
        img_h = data["img_height"]
        canvas_width = data.get("canvas_width")
        canvas_height = data.get("canvas_height")
        content_offset_x = data.get("content_offset_x", 0)
        content_offset_y = data.get("content_offset_y", 0)
        use_canvas = _validate_canvas(
            img_w,
            img_h,
            canvas_width,
            canvas_height,
            content_offset_x,
            content_offset_y,
        )
        if use_canvas:
            transform = compute_contain_transform(canvas_width, canvas_height)
        else:
            transform = compute_contain_transform(img_w, img_h)

        # --- Content slide ---
        slide = prs.slides.add_slide(blank_layout)

        # Layer 1: Background
        slide.shapes.add_picture(
            str(data["background_path"]), 0, 0,
            Inches(transform.slide_width), Inches(transform.slide_height)
        )

        # Layer 2: Foreground components
        for comp in data["components"]:
            _add_component(
                slide,
                comp,
                img_w,
                img_h,
                transform,
                canvas_width if use_canvas else None,
                canvas_height if use_canvas else None,
                content_offset_x if use_canvas else 0,
                content_offset_y if use_canvas else 0,
            )

        # Layer 3: Text boxes
        for item in data["text_items"]:
            _add_textbox(
                slide,
                item,
                img_w,
                img_h,
                transform,
                canvas_width if use_canvas else None,
                canvas_height if use_canvas else None,
                content_offset_x if use_canvas else 0,
                content_offset_y if use_canvas else 0,
            )

        logger.info("Slide %d: bg + %d components + %d text boxes.",
                    idx + 1, len(data["components"]), len(data["text_items"]))

        # --- Reference slide (optional) ---
        if add_reference and data.get("original_image_path"):
            orig = Path(data["original_image_path"])
            if orig.exists():
                ref_slide = prs.slides.add_slide(blank_layout)
                ref_slide.shapes.add_picture(
                    str(data["background_path"]), 0, 0,
                    Inches(transform.slide_width), Inches(transform.slide_height)
                )
                left, top, width, height = _map_bbox(
                    0,
                    0,
                    img_w,
                    img_h,
                    img_w,
                    img_h,
                    transform,
                    canvas_width if use_canvas else None,
                    canvas_height if use_canvas else None,
                    content_offset_x if use_canvas else 0,
                    content_offset_y if use_canvas else 0,
                )
                ref_slide.shapes.add_picture(
                    str(orig),
                    Inches(left), Inches(top), Inches(width), Inches(height)
                )

    prs.save(str(output_path))
    logger.info("Saved multi-slide PPTX (%d slides): %s", len(slides_data), output_path)

    return str(output_path)


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _hex_to_rgb(hex_color: str) -> tuple[int, int, int]:
    """Convert '#rrggbb' to (r, g, b) tuple."""
    hex_color = hex_color.lstrip("#")
    if len(hex_color) != 6:
        return (0, 0, 0)
    return (
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _validate_canvas(
    img_width: int,
    img_height: int,
    canvas_width: int | None,
    canvas_height: int | None,
    content_offset_x: int,
    content_offset_y: int,
) -> bool:
    """Validate an optional widescreen canvas and its source-image rectangle."""
    if not all(
        isinstance(value, Integral) and not isinstance(value, bool)
        for value in (content_offset_x, content_offset_y)
    ):
        raise ValueError("canvas offsets must be integers")
    if canvas_width is None and canvas_height is None:
        if content_offset_x != 0 or content_offset_y != 0:
            raise ValueError("canvas offsets require canvas dimensions")
        return False
    if canvas_width is None or canvas_height is None:
        raise ValueError("canvas_width and canvas_height must be provided together")
    if not all(
        isinstance(value, Integral) and not isinstance(value, bool)
        for value in (canvas_width, canvas_height)
    ):
        raise ValueError("canvas dimensions must be integers")
    if canvas_width <= 0 or canvas_height <= 0:
        raise ValueError("canvas dimensions must be positive")
    if canvas_width * 9 != canvas_height * 16:
        raise ValueError("canvas dimensions must have an exact 16:9 ratio")
    if (
        content_offset_x < 0
        or content_offset_y < 0
        or content_offset_x + img_width > canvas_width
        or content_offset_y + img_height > canvas_height
    ):
        raise ValueError("source image must fit completely inside the canvas")
    return True


def _map_bbox(
    x: int,
    y: int,
    width: int,
    height: int,
    img_w: int,
    img_h: int,
    transform: ContainTransform,
    canvas_width: int | None = None,
    canvas_height: int | None = None,
    content_offset_x: int = 0,
    content_offset_y: int = 0,
) -> tuple[float, float, float, float]:
    """Map a source-image bounding box to slide inches."""
    if canvas_width is not None and canvas_height is not None:
        return (
            (x + content_offset_x) / canvas_width * transform.slide_width,
            (y + content_offset_y) / canvas_height * transform.slide_height,
            width / canvas_width * transform.slide_width,
            height / canvas_height * transform.slide_height,
        )

    return (
        transform.offset_x + x / img_w * transform.content_width,
        transform.offset_y + y / img_h * transform.content_height,
        width / img_w * transform.content_width,
        height / img_h * transform.content_height,
    )


def _add_component(
    slide,
    component: dict,
    img_w: int,
    img_h: int,
    transform: ContainTransform,
    canvas_width: int | None = None,
    canvas_height: int | None = None,
    content_offset_x: int = 0,
    content_offset_y: int = 0,
) -> None:
    left, top, width, height = _map_bbox(
        component["x"],
        component["y"],
        component["w"],
        component["h"],
        img_w,
        img_h,
        transform,
        canvas_width,
        canvas_height,
        content_offset_x,
        content_offset_y,
    )
    slide.shapes.add_picture(
        component["path"],
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )


def _add_textbox(
    slide,
    item: dict,
    img_w: int,
    img_h: int,
    transform: ContainTransform,
    canvas_width: int | None = None,
    canvas_height: int | None = None,
    content_offset_x: int = 0,
    content_offset_y: int = 0,
) -> None:
    """Add an editable text box to the slide.

    Alignment is applied inside the OCR-detected bounds so nearby text keeps
    its original horizontal position.
    """
    x, y, w, h = item["box"]

    left, top, width, height = _map_bbox(
        x,
        y,
        w,
        h,
        img_w,
        img_h,
        transform,
        canvas_width,
        canvas_height,
        content_offset_x,
        content_offset_y,
    )

    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    from pptx.enum.text import MSO_VERTICAL_ANCHOR
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    tf.clear()

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = item.get("text", "")

    font = run.font
    _set_run_font(run, item.get("font", "Microsoft YaHei"))
    if canvas_width is not None:
        font_scale = img_w / canvas_width
    else:
        font_scale = transform.content_width / SLIDE_WIDTH_INCHES
    font_size = item.get("font_size", 12) * font_scale
    font.size = Pt(font_size)
    font.bold = item.get("bold", False)

    color = _hex_to_rgb(item.get("color", "#000000"))
    font.color.rgb = RGBColor(*color)

    # Alignment: 0=left, 1=center, 2=right
    from pptx.enum.text import PP_ALIGN
    align_map = {0: PP_ALIGN.LEFT, 1: PP_ALIGN.CENTER, 2: PP_ALIGN.RIGHT}
    p.alignment = align_map.get(item.get("align", 1), PP_ALIGN.CENTER)


def _set_run_font(run, font_name: str) -> None:
    """Set both Latin and East Asian font names for PowerPoint."""
    run.font.name = font_name
    rpr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        node = rpr.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            rpr.append(node)
        node.set("typeface", font_name)
