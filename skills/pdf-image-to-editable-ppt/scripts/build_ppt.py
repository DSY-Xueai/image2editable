from __future__ import annotations

from pathlib import Path
from zipfile import ZipFile, ZIP_DEFLATED


def _build_minimal_pptx(output_path: Path, slide_count: int) -> None:
    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/ppt/presentation.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.presentation.main+xml"/>
  <Override PartName="/ppt/slideMasters/slideMaster1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideMaster+xml"/>
  <Override PartName="/ppt/slideLayouts/slideLayout1.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slideLayout+xml"/>
  <Override PartName="/ppt/theme/theme1.xml" ContentType="application/vnd.openxmlformats-officedocument.theme+xml"/>
  <Override PartName="/docProps/core.xml" ContentType="application/vnd.openxmlformats-package.core-properties+xml"/>
  <Override PartName="/docProps/app.xml" ContentType="application/vnd.openxmlformats-officedocument.extended-properties+xml"/>
  {slides}
</Types>"""
    slide_overrides = "\n  ".join(
        f'<Override PartName="/ppt/slides/slide{i}.xml" ContentType="application/vnd.openxmlformats-officedocument.presentationml.slide+xml"/>'
        for i in range(1, slide_count + 1)
    )
    rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="ppt/presentation.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/package/2006/relationships/metadata/core-properties" Target="docProps/core.xml"/>
  <Relationship Id="rId3" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/extended-properties" Target="docProps/app.xml"/>
</Relationships>"""
    presentation = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:presentation xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:sldMasterIdLst><p:sldMasterId id="2147483648" r:id="rId1"/></p:sldMasterIdLst>
  <p:sldIdLst>{slides}</p:sldIdLst>
  <p:sldSz cx="9144000" cy="6858000"/>
  <p:notesSz cx="6858000" cy="9144000"/>
</p:presentation>"""
    presentation_slides = "".join(
        f'<p:sldId id="{255 + i}" r:id="rId{i + 1}"/>' for i in range(1, slide_count + 1)
    )
    presentation_rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideMaster" Target="slideMasters/slideMaster1.xml"/>
  {slides}
</Relationships>"""
    presentation_rel_slides = "\n  ".join(
        f'<Relationship Id="rId{i + 1}" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide" Target="slides/slide{i}.xml"/>'
        for i in range(1, slide_count + 1)
    )
    slide_master = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldMaster xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld name="Blank"><p:bg><p:bgRef idx="1001"><a:schemeClr val="bg1"/></p:bgRef></p:bg></p:cSld>
  <p:clrMap accent1="accent1" accent2="accent2" accent3="accent3" accent4="accent4" accent5="accent5" accent6="accent6" bg1="lt1" bg2="lt2" folHlink="folHlink" hlink="hlink" tx1="dk1" tx2="dk2"/>
  <p:sldLayoutIdLst><p:sldLayoutId id="1" r:id="rId1"/></p:sldLayoutIdLst>
  <p:txStyles><p:titleStyle/><p:bodyStyle/><p:otherStyle/></p:txStyles>
</p:sldMaster>"""
    slide_master_rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout" Target="../slideLayouts/slideLayout1.xml"/>
  <Relationship Id="rId2" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme" Target="../theme/theme1.xml"/>
</Relationships>"""
    slide_layout = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sldLayout xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" type="blank" preserve="1">
  <p:cSld name="Blank"/>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sldLayout>"""
    theme = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<a:theme xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" name="Office Theme">
  <a:themeElements>
    <a:clrScheme name="Office"><a:dk1><a:sysClr val="windowText" lastClr="000000"/></a:dk1><a:lt1><a:sysClr val="window" lastClr="FFFFFF"/></a:lt1><a:dk2><a:srgbClr val="1F497D"/></a:dk2><a:lt2><a:srgbClr val="EEECE1"/></a:lt2><a:accent1><a:srgbClr val="4F81BD"/></a:accent1><a:accent2><a:srgbClr val="C0504D"/></a:accent2><a:accent3><a:srgbClr val="9BBB59"/></a:accent3><a:accent4><a:srgbClr val="8064A2"/></a:accent4><a:accent5><a:srgbClr val="4BACC6"/></a:accent5><a:accent6><a:srgbClr val="F79646"/></a:accent6><a:hlink><a:srgbClr val="0000FF"/></a:hlink><a:folHlink><a:srgbClr val="800080"/></a:folHlink></a:clrScheme>
    <a:fontScheme name="Office"><a:majorFont><a:latin typeface="Arial"/></a:majorFont><a:minorFont><a:latin typeface="Arial"/></a:minorFont></a:fontScheme>
    <a:fmtScheme name="Office"><a:fillStyleLst/><a:lnStyleLst/><a:effectStyleLst/><a:bgFillStyleLst/></a:fmtScheme>
  </a:themeElements>
</a:theme>"""
    core = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<cp:coreProperties xmlns:cp="http://schemas.openxmlformats.org/package/2006/metadata/core-properties" xmlns:dc="http://purl.org/dc/elements/1.1/" xmlns:dcterms="http://purl.org/dc/terms/" xmlns:dcmitype="http://purl.org/dc/dcmitype/" xmlns:xsi="http://www.w3.org/2001/XMLSchema-instance">
  <dc:title>Generated Presentation</dc:title>
</cp:coreProperties>"""
    app = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Properties xmlns="http://schemas.openxmlformats.org/officeDocument/2006/extended-properties" xmlns:vt="http://schemas.openxmlformats.org/officeDocument/2006/docPropsVTypes">
  <Application>OpenAI Codex</Application>
</Properties>"""
    blank_slide = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<p:sld xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships" xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
  <p:cSld><p:spTree><p:nvGrpSpPr><p:cNvPr id="1" name=""/><p:cNvGrpSpPr/><p:nvPr/></p:nvGrpSpPr><p:grpSpPr/></p:spTree></p:cSld>
  <p:clrMapOvr><a:masterClrMapping/></p:clrMapOvr>
</p:sld>"""

    with ZipFile(output_path, "w", ZIP_DEFLATED) as archive:
        archive.writestr("[Content_Types].xml", content_types.format(slides=slide_overrides))
        archive.writestr("_rels/.rels", rels)
        archive.writestr(
            "ppt/presentation.xml", presentation.format(slides=presentation_slides)
        )
        archive.writestr(
            "ppt/_rels/presentation.xml.rels",
            presentation_rels.format(slides=presentation_rel_slides),
        )
        archive.writestr("ppt/slideMasters/slideMaster1.xml", slide_master)
        archive.writestr(
            "ppt/slideMasters/_rels/slideMaster1.xml.rels", slide_master_rels
        )
        archive.writestr("ppt/slideLayouts/slideLayout1.xml", slide_layout)
        archive.writestr("ppt/theme/theme1.xml", theme)
        archive.writestr("docProps/core.xml", core)
        archive.writestr("docProps/app.xml", app)
        for index in range(1, slide_count + 1):
            archive.writestr(f"ppt/slides/slide{index}.xml", blank_slide)


def build_presentation(page_plans, output_path: Path) -> None:
    slide_count = max(len(page_plans), 1)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    for page in page_plans:
        if not Path(page.background_path).exists():
            raise FileNotFoundError(page.background_path)

    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Emu, Pt
    except ImportError:
        _build_minimal_pptx(output_path, slide_count)
        return

    presentation = Presentation()
    while presentation.slides:
        slide_id = presentation.slides._sldIdLst[0]
        presentation.slides._sldIdLst.remove(slide_id)

    base_page = page_plans[0]
    slide_width = 9144000
    slide_height = int(slide_width * (base_page.height_px / max(base_page.width_px, 1)))
    presentation.slide_width = slide_width
    presentation.slide_height = slide_height

    alignment_map = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
        "justify": PP_ALIGN.JUSTIFY,
    }

    for page in page_plans:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        slide.shapes.add_picture(
            str(page.background_path),
            0,
            0,
            width=presentation.slide_width,
            height=presentation.slide_height,
        )

        scale_x = presentation.slide_width / max(page.width_px, 1)
        scale_y = presentation.slide_height / max(page.height_px, 1)

        for block in page.text_blocks:
            textbox = slide.shapes.add_textbox(
                Emu(int(block.left * scale_x)),
                Emu(int(block.top * scale_y)),
                Emu(int(block.width * scale_x)),
                Emu(int(block.height * scale_y)),
            )
            paragraph = textbox.text_frame.paragraphs[0]
            paragraph.text = block.text
            paragraph.alignment = alignment_map.get(block.alignment, PP_ALIGN.LEFT)
            run = paragraph.runs[0]
            run.font.size = Pt(block.font_size)
            run.font.color.rgb = RGBColor.from_string(block.color.lstrip("#"))
            if getattr(block, "font_name", None):
                run.font.name = block.font_name

        for block in page.image_blocks:
            if not Path(block.path).exists():
                continue
            slide.shapes.add_picture(
                block.path,
                Emu(int(block.left * scale_x)),
                Emu(int(block.top * scale_y)),
                width=Emu(int(block.width * scale_x)),
                height=Emu(int(block.height * scale_y)),
            )

        for _effect in getattr(page, "effect_blocks", []):
            # Stage 2 uses fail-closed effect handling. Unsupported effects stay in the background.
            continue

    presentation.save(output_path)
