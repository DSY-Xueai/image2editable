from __future__ import annotations

import hashlib
import io
import json
import math
import os
import stat
import uuid
from pathlib import Path
from pathlib import PurePosixPath

import image_to_ppt
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Pt


def build_reconstruction_donor(
    image_path: str | Path,
    output_pptx: str | Path,
    work_root: str | Path,
    *,
    decision: dict,
    lang: str = "ch",
) -> dict:
    """Run the isolated CV pipeline for one Agent-approved screenshot."""
    if (
        decision.get("runtime_action") != "shadow_run"
        or decision.get("eligible_for_shadow_run") is not True
        or not isinstance(decision.get("source_shape_id"), str)
        or not decision["source_shape_id"]
    ):
        raise ValueError(
            "reconstruction requires an Agent-approved shadow_run decision"
        )

    image = Path(image_path).resolve()
    output = Path(output_pptx).resolve()
    root = Path(work_root).resolve()
    if not image.is_file():
        raise FileNotFoundError(image)
    if output.exists():
        raise FileExistsError(output)
    root.mkdir(parents=True, exist_ok=True)

    slide_data, assets = image_to_ppt._prepare_single_image(
        image,
        lang,
        _work_root=root,
        _resource_isolation=True,
    )
    image_to_ppt._assemble_prepared_slide(
        slide_data,
        output,
        False,
        "original",
    )
    if not output.is_file():
        raise RuntimeError("reconstruction pipeline did not create donor PPTX")

    manifest = {
        "source_shape_id": decision["source_shape_id"],
        "candidate_image": str(image),
        "candidate_sha256": _sha256(image),
        "donor_pptx": str(output),
        "assets": str(assets),
        "components": len(slide_data["components"]),
        "text_boxes": len(slide_data["text_items"]),
        "quality": slide_data.get("quality"),
        "background_residual": slide_data.get("background_residual"),
    }
    manifest_path = root / "reconstruction_manifest.json"
    if manifest_path.exists():
        raise FileExistsError(manifest_path)
    temporary = manifest_path.with_suffix(".json.tmp")
    temporary.write_text(
        json.dumps(manifest, ensure_ascii=False, indent=2) + "\n",
        encoding="utf-8",
    )
    temporary.replace(manifest_path)
    return manifest


def _sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as file:
        while chunk := file.read(1024 * 1024):
            digest.update(chunk)
    return digest.hexdigest()


def build_reconstruction_donor_from_result(
    result_path: str | Path,
    output_pptx: str | Path,
    work_root: str | Path,
    *,
    source_screenshot_sha256: str,
    provider: str,
    initial_component_count: int = 0,
    expected_result_sha256: str | None = None,
    run_root: str | Path | None = None,
) -> dict:
    """Assemble a donor strictly from an accepted component result.

    This path deliberately does not invoke OCR/CV/model code; the result is
    treated as the already validated boundary artifact from Task 9.
    """
    result_file = Path(result_path).resolve()
    native_root = Path(run_root).resolve() if run_root is not None else None
    if native_root is not None:
        try:
            result_relative = result_file.relative_to(native_root).as_posix()
        except ValueError as error:
            raise ValueError("component result is outside the Run") from error
        result_file, result_payload = _read_run_ref(
            native_root,
            {"path": result_relative,
             "sha256": expected_result_sha256 or _sha256(result_file)},
            "component result",
        )
        if (
            expected_result_sha256 is not None
            and _sha256(result_file) != expected_result_sha256
        ):
            raise ValueError("component result hash mismatch")
        data = json.loads(result_payload.decode("utf-8"))
    else:
        if expected_result_sha256 is not None and _sha256(result_file) != expected_result_sha256:
            raise ValueError("component result hash mismatch")
        data = json.loads(result_file.read_text(encoding="utf-8"))
    native_graph = None
    native_graph_path = None
    native_background_payload = None
    native_reconstructed_payload = None
    native_text_mask_payload = None
    native_mask_payloads: dict[str, bytes] = {}
    if "accepted_asset_refs" in data and "components" not in data:
        if native_root is None:
            raise ValueError("component result run root is required")
        graph_ref = data.get("graph_ref")
        if not _is_run_ref(graph_ref):
            raise ValueError("component result graph reference is missing")
        graph_path, graph_payload = _read_run_ref(
            native_root, graph_ref, "component result graph"
        )
        graph = json.loads(graph_payload.decode("utf-8"))
        from image2editable.component_contracts import validate_component_graph
        graph = validate_component_graph(graph)
        accepted_refs = data.get("accepted_asset_refs", {})
        if not isinstance(accepted_refs, dict):
            raise ValueError("component result accepted assets are invalid")
        accepted_payloads = {}
        for name, ref in accepted_refs.items():
            if not _is_run_ref(ref):
                raise ValueError(f"component result {name} reference is invalid")
            _, accepted_payloads[name] = _read_run_ref(
                native_root, ref, f"component result {name}"
            )
        source_ref = accepted_refs.get("source")
        if not _is_run_ref(source_ref):
            raise ValueError("component result source reference is missing")
        expected_source_sha = data.get(
            "source_sha256", source_screenshot_sha256
        )
        if source_ref["sha256"] != expected_source_sha:
            raise ValueError("component result source hash mismatch")
        for name in ("background", "reconstructed", "text_mask"):
            ref = accepted_refs.get(name)
            if not _is_run_ref(ref):
                raise ValueError("component result accepted assets are invalid")
            if name == "background":
                native_background_payload = accepted_payloads[name]
            elif name == "reconstructed":
                native_reconstructed_payload = accepted_payloads[name]
            else:
                native_text_mask_payload = accepted_payloads[name]
        native_graph = graph
        native_graph_path = graph_path
        final_ids = data.get("final_component_ids", [])
        if (
            not isinstance(final_ids, list)
            or final_ids != sorted(set(final_ids))
            or any(type(value) is not str or not value for value in final_ids)
        ):
            raise ValueError("component result final IDs are invalid")
        by_id = {node["id"]: node for node in graph.get("nodes", [])}
        components = []
        assets = []
        for component_id in final_ids:
            node = by_id.get(component_id)
            if not isinstance(node, dict):
                raise ValueError("component result final ID is invalid")
            if node["state"] != "frozen":
                raise ValueError("component result final component is not frozen")
            mask_relative = (
                PurePosixPath(graph_ref["path"]).parent
                / PurePosixPath(node["mask"])
            ).as_posix()
            mask_path, mask_payload = _read_run_ref(
                native_root,
                {"path": mask_relative,
                 "sha256": node["mask_sha256"]},
                f"component result mask {component_id}",
            )
            native_mask_payloads[component_id] = mask_payload
            components.append({
                "id": component_id, "kind": node["kind"],
                "parent_id": node.get("parent_id"), "state": "active",
                "path": str(mask_path), "sha256": node["mask_sha256"],
            })
            assets.append({"path": str(mask_path), "sha256": node["mask_sha256"]})
        frozen_visual_ids = {
            node["id"] for node in graph["nodes"]
            if node["state"] == "frozen" and node["kind"] != "text"
        }
        if set(final_ids) - set(by_id) != set():
            raise ValueError("component result final IDs are invalid")
        if frozen_visual_ids - {
            node["id"] for node in graph["nodes"]
            if node["id"] in final_ids and node["kind"] != "text"
        }:
            raise ValueError("component result final IDs omit active components")
        text_items = data.get("text_items", [])
        if not isinstance(text_items, list):
            raise ValueError("component result text_items are invalid")
        data = {
            **data,
            "source_screenshot_sha256": expected_source_sha,
            "components": components,
            "text_items": text_items,
            "assets": assets,
        }
    if data.get("provider") != provider:
        raise ValueError("component result provider mismatch")
    if data.get("source_screenshot_sha256") != source_screenshot_sha256:
        raise ValueError("source screenshot SHA mismatch")
    components = data.get("components")
    if not isinstance(components, list):
        raise ValueError("component result components must be a list")
    active = [c for c in components if c.get("state", "active") == "active"]
    active_visuals = [c for c in active if c.get("kind") != "text"]
    if initial_component_count > 0 and not active_visuals:
        raise ValueError("non-zero initial components cannot produce zero-component donor")
    parents = {c.get("id") for c in active if c.get("kind") == "parent"}
    children = {c.get("parent_id") for c in active if c.get("kind") == "child"}
    if parents & children:
        raise ValueError("parent and child components cannot both be active")
    for asset in data.get("assets", []):
        if not isinstance(asset, dict) or not isinstance(asset.get("path"), str):
            raise ValueError("asset reference is invalid")
        path = Path(asset["path"])
        if not path.is_file() or _sha256(path) != asset.get("sha256"):
            raise ValueError("asset hash mismatch")
    output = Path(output_pptx).resolve()
    output.parent.mkdir(parents=True, exist_ok=True)
    work_root_path = Path(work_root).resolve()
    work_root_path.mkdir(parents=True, exist_ok=True)
    if output.exists():
        raise FileExistsError(output)
    # Keep a single slide so the OOXML patcher can consume the donor even
    # when this deterministic boundary carries no model-generated shapes.
    donor_presentation = Presentation()
    donor_slide = donor_presentation.slides.add_slide(
        donor_presentation.slide_layouts[6]
    )
    active_components = active_visuals
    raster_text_preserved = False
    if native_root is not None and native_graph is not None:
        if (
            native_background_payload is None
            or native_reconstructed_payload is None
            or native_text_mask_payload is None
        ):
            raise ValueError("component result reconstruction assets are missing")
        background = Image.open(
            io.BytesIO(native_background_payload)
        ).convert("RGB")
        reconstructed = Image.open(
            io.BytesIO(native_reconstructed_payload)
        ).convert("RGBA")
        text_mask = Image.open(io.BytesIO(native_text_mask_payload)).convert("L")
        if background.size != reconstructed.size or reconstructed.size != text_mask.size:
            raise ValueError("component result asset dimensions differ")
        width, height = reconstructed.size
        donor_slide.shapes.add_picture(
            io.BytesIO(native_background_payload),
            0,
            0,
            donor_presentation.slide_width,
            donor_presentation.slide_height,
        )
        active_nodes = [
            node for node in native_graph["nodes"]
            if node["state"] == "frozen" and node["kind"] != "text"
        ]
        text_items = _validated_text_items(data.get("text_items", []), width, height)
        raster_text_preserved = False
        component_manifest = []
        component_assets = []
        published_route = None
        if isinstance(data.get("page_id"), str):
            from image2editable.route_execution import load_published_route
            from image2editable.store import RunStore

            published_route = load_published_route(
                RunStore(native_root), result_file, page_id=data["page_id"]
            )
        if published_route is not None:
            from image2editable.route_execution import route_visual_elements
            from scripts.ppt_assemble import ContainTransform, _add_visual_element

            route_store = RunStore(native_root)
            elements = route_visual_elements(
                route_store, published_route["ir"], published_route["plan"]
            )
            slide_width = donor_presentation.slide_width / 914400
            slide_height = donor_presentation.slide_height / 914400
            transform = ContainTransform(
                slide_width,
                slide_height,
                slide_width,
                slide_height,
                0,
                0,
            )
            for element in elements:
                _add_visual_element(
                    donor_slide, element, width, height, transform
                )
                entry = {
                    "id": element["object_id"],
                    "kind": element["route"],
                }
                if element["route"] == "raster_component":
                    component_path = Path(element["component"]["path"])
                    entry.update(
                        path=str(component_path), sha256=_sha256(component_path)
                    )
                    component_assets.append(
                        {"path": entry["path"], "sha256": entry["sha256"]}
                    )
                component_manifest.append(entry)
        else:
            for node in active_nodes:
                if native_graph_path is None:
                    raise ValueError("component result graph path is missing")
                mask_payload = native_mask_payloads.get(node["id"])
                if mask_payload is None:
                    raise ValueError("component result mask snapshot is missing")
                mask = Image.open(io.BytesIO(mask_payload)).convert("L")
                if mask.size != reconstructed.size:
                    raise ValueError("component result mask dimensions differ")
                alpha = mask
                bbox = alpha.getbbox()
                if bbox is None:
                    raise ValueError("accepted component mask is empty")
                cropped = reconstructed.crop(bbox)
                cropped.putalpha(alpha.crop(bbox))
                component_path = work_root_path / f"component-{node['id']}.png"
                component_stream = io.BytesIO()
                cropped.save(component_stream, format="PNG")
                _publish_bytes_no_clobber(
                    component_path,
                    component_stream.getvalue(),
                    reuse_identical=True,
                )
                left, top, right, bottom = bbox
                donor_slide.shapes.add_picture(
                    str(component_path),
                    int(left / width * donor_presentation.slide_width),
                    int(top / height * donor_presentation.slide_height),
                    int((right - left) / width * donor_presentation.slide_width),
                    int((bottom - top) / height * donor_presentation.slide_height),
                )
                component_manifest.append({
                    "id": node["id"], "kind": node["kind"],
                    "path": str(component_path), "sha256": _sha256(component_path),
                })
            component_assets = [
                {"path": item["path"], "sha256": item["sha256"]}
                for item in component_manifest
            ]
        for item in text_items:
            left, top, right, bottom = item["box"]
            left_emu = int(left / width * donor_presentation.slide_width)
            top_emu = int(top / height * donor_presentation.slide_height)
            width_emu = int((right - left) / width * donor_presentation.slide_width)
            height_emu = int((bottom - top) / height * donor_presentation.slide_height)
            rotation = item.get("rotation", 0)
            if rotation in {90, 270}:
                center_x = left_emu + width_emu / 2
                center_y = top_emu + height_emu / 2
                width_emu, height_emu = height_emu, width_emu
                left_emu = int(center_x - width_emu / 2)
                top_emu = int(center_y - height_emu / 2)
            text_shape = donor_slide.shapes.add_textbox(
                left_emu,
                top_emu,
                width_emu,
                height_emu,
            )
            text_shape.rotation = rotation
            _style_reconstruction_textbox(text_shape, item)
        data["components"] = component_manifest
        data["text_items"] = text_items
        data["assets"] = list(data.get("assets", [])) + component_assets
        active_components = component_manifest
    donor_stream = io.BytesIO()
    donor_presentation.save(donor_stream)
    manifest = {
        "provider": provider,
        "source_screenshot_sha256": source_screenshot_sha256,
        "donor_pptx": str(output),
        "components": len(active_components),
        "component_ids": [item.get("id") for item in active_components],
        "component_kinds": [item.get("kind") for item in active_components],
        "text_boxes": len(data.get("text_items", [])),
        "assets": data.get("assets", []),
        "raster_text_preserved": raster_text_preserved,
        "warning": None,
    }
    root = work_root_path
    manifest_path = root / "reconstruction_manifest.json"
    manifest_payload = (json.dumps(
        manifest, ensure_ascii=False, indent=2, sort_keys=True
    ) + "\n").encode("utf-8")
    _publish_bytes_no_clobber(
        manifest_path, manifest_payload, reuse_identical=True
    )
    _publish_bytes_no_clobber(
        output, donor_stream.getvalue(), reuse_identical=False
    )
    return manifest


def _style_reconstruction_textbox(text_shape, item: dict) -> None:
    text_frame = text_shape.text_frame
    text_frame.word_wrap = False
    text_frame.margin_left = 0
    text_frame.margin_right = 0
    text_frame.margin_top = 0
    text_frame.margin_bottom = 0
    text_frame.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = item.get("text", "")
    font_name = item.get("font") or "Microsoft YaHei"
    run.font.name = font_name
    run_properties = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea"):
        node = run_properties.find(qn(tag))
        if node is None:
            node = OxmlElement(tag)
            run_properties.append(node)
        node.set("typeface", font_name)
    run.font.size = Pt(item.get("font_size", 12))
    run.font.bold = item.get("bold", False)
    color = str(item.get("color", "#000000")).lstrip("#")
    if len(color) != 6:
        color = "000000"
    run.font.color.rgb = RGBColor.from_string(color)
    paragraph.alignment = {
        0: PP_ALIGN.LEFT,
        1: PP_ALIGN.CENTER,
        2: PP_ALIGN.RIGHT,
    }.get(item.get("align", 1), PP_ALIGN.CENTER)


def _is_run_ref(value: object) -> bool:
    if not isinstance(value, dict) or set(value) != {"path", "sha256"}:
        return False
    path = value["path"]
    return (
        isinstance(path, str)
        and bool(path)
        and not Path(path).is_absolute()
        and "\\" not in path
        and ".." not in PurePosixPath(path).parts
        and isinstance(value["sha256"], str)
        and len(value["sha256"]) == 64
        and all(character in "0123456789abcdef" for character in value["sha256"])
    )


def _read_run_ref(root: Path, ref: dict, label: str) -> tuple[Path, bytes]:
    if not _is_run_ref(ref):
        raise ValueError(f"{label} reference is invalid")
    path = root
    parts = PurePosixPath(ref["path"]).parts
    before = None
    reparse_flag = getattr(stat, "FILE_ATTRIBUTE_REPARSE_POINT", 0)
    for index, part in enumerate(parts):
        path /= part
        try:
            before = path.lstat()
        except FileNotFoundError as error:
            raise ValueError(f"{label} is missing") from error
        if (
            stat.S_ISLNK(before.st_mode)
            or bool(getattr(before, "st_file_attributes", 0) & reparse_flag)
        ):
            raise ValueError(f"{label} contains a link or reparse point")
        if index < len(parts) - 1 and not stat.S_ISDIR(before.st_mode):
            raise ValueError(f"{label} parent is not a directory")
    if before is None:
        raise ValueError(f"{label} reference is invalid")
    if (
        not stat.S_ISREG(before.st_mode)
        or before.st_nlink != 1
    ):
        raise ValueError(f"{label} must be a regular file")
    with path.open("rb") as handle:
        payload = handle.read()
        digest = hashlib.sha256(payload).hexdigest()
    after = path.lstat()
    if (
        before.st_dev != after.st_dev
        or before.st_ino != after.st_ino
        or before.st_size != after.st_size
        or before.st_mtime_ns != after.st_mtime_ns
    ):
        raise ValueError(f"{label} changed during verification")
    if digest != ref["sha256"]:
        raise ValueError(f"{label} hash mismatch")
    return path, payload


def _publish_bytes_no_clobber(
    path: Path, payload: bytes, *, reuse_identical: bool
) -> None:
    if path.exists() or path.is_symlink():
        if reuse_identical and path.is_file() and path.read_bytes() == payload:
            return
        raise FileExistsError(path)
    temporary = path.with_name(f".{path.name}.{uuid.uuid4().hex}.tmp")
    try:
        temporary.write_bytes(payload)
        try:
            os.link(temporary, path)
        except FileExistsError:
            if reuse_identical and path.is_file() and path.read_bytes() == payload:
                return
            raise
    finally:
        temporary.unlink(missing_ok=True)


def _validated_text_items(items: object, width: int, height: int) -> list[dict]:
    if not isinstance(items, list):
        raise ValueError("component result text_items are invalid")
    validated = []
    for item in items:
        if not isinstance(item, dict) or not isinstance(item.get("text"), str):
            raise ValueError("component result text item is invalid")
        box = item.get("box")
        if (
            not isinstance(box, list)
            or len(box) != 4
            or any(type(value) not in {int, float} or not math.isfinite(value) for value in box)
            or box[0] < 0 or box[1] < 0
            or box[2] <= 0 or box[3] <= 0
            or box[0] + box[2] > width or box[1] + box[3] > height
        ):
            raise ValueError("component result text item box is invalid")
        x, y, w, h = box
        rotation = item.get("rotation", 0)
        if type(rotation) is not int or rotation not in {0, 90, 180, 270}:
            raise ValueError("component result text rotation is invalid")
        validated.append({**item, "box": [x, y, x + w, y + h]})
    return validated
