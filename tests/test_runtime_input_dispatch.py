from __future__ import annotations

from pathlib import Path
import subprocess
import sys
import textwrap

import pytest
from PIL import Image

from image2editable import inputs, runtime


def _write_image(path: Path) -> None:
    Image.new("RGB", (1, 1)).save(path)


def test_classify_inputs_returns_a_resolved_pdf(tmp_path: Path) -> None:
    source = tmp_path / "document.pdf"
    source.write_bytes(b"pdf")

    assert inputs.classify_inputs(source) == ("pdf", [source.resolve()])


def test_classify_inputs_returns_a_resolved_pptx(tmp_path: Path) -> None:
    source = tmp_path / "deck.pptx"
    source.write_bytes(b"pptx")

    assert inputs.classify_inputs(source) == ("pptx", [source.resolve()])


def test_classify_inputs_keeps_image_directory_behavior(tmp_path: Path) -> None:
    folder = tmp_path / "images"
    folder.mkdir()
    _write_image(folder / "b.png")
    _write_image(folder / "a.png")
    extra = tmp_path / "extra.png"
    _write_image(extra)

    input_type, paths = inputs.classify_inputs([folder, extra])

    assert input_type == "images"
    assert [path.name for path in paths] == ["a.png", "b.png", "extra.png"]


@pytest.mark.parametrize(
    "names",
    [
        ("first.pdf", "second.pdf"),
        ("first.pptx", "second.pptx"),
        ("image.png", "document.pdf"),
        ("document.pdf", "deck.pptx"),
    ],
)
def test_classify_inputs_rejects_multiple_inputs_with_documents(
    tmp_path: Path, names: tuple[str, str]
) -> None:
    paths = []
    for name in names:
        path = tmp_path / name
        if path.suffix == ".png":
            _write_image(path)
        else:
            path.write_bytes(b"document")
        paths.append(path)

    with pytest.raises(ValueError, match="one PDF or one PPTX"):
        inputs.classify_inputs(paths)


def test_classify_inputs_rejects_empty_inputs() -> None:
    with pytest.raises(ValueError, match="^No inputs provided$"):
        inputs.classify_inputs([])


def test_classify_inputs_reports_missing_paths(tmp_path: Path) -> None:
    missing = tmp_path / "missing.pdf"

    with pytest.raises(FileNotFoundError, match="Input path does not exist"):
        inputs.classify_inputs(missing)


def test_new_job_id_has_existing_format() -> None:
    assert inputs.new_job_id().count("-") == 1


@pytest.mark.parametrize(
    ("name", "input_type", "handler_name"),
    [
        ("source.png", "images", "prepare_image_job"),
        ("source.pdf", "pdf", "prepare_pdf_job"),
        ("source.pptx", "pptx", "prepare_pptx_job"),
    ],
)
def test_prepare_job_classifies_and_dispatches_all_supported_inputs(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    name: str,
    input_type: str,
    handler_name: str,
) -> None:
    source = tmp_path / name
    if source.suffix == ".png":
        _write_image(source)
    else:
        source.write_bytes(b"document")
    run_dir = tmp_path / "run"
    output = tmp_path / "output.pptx"
    calls = []

    def handler(*args, **kwargs):
        calls.append((args, kwargs))
        return run_dir

    monkeypatch.setattr(runtime, handler_name, handler)

    result = runtime.prepare_job(
        source,
        run_dir=run_dir,
        output_path=output,
        slide_size="original",
        lang="en",
    )

    expected_source = (
        [source.resolve()] if input_type == "images" else source.resolve()
    )
    assert result == run_dir
    assert calls == [
        (
            (expected_source,),
            {
                "run_dir": run_dir,
                "output_path": output,
                    "slide_size": "original",
                    "lang": "en",
                    "agent_provider": "host",
            },
        )
    ]


def test_package_help_and_pptx_runtime_work_without_pdfium(
    tmp_path: Path,
) -> None:
    script = textwrap.dedent(
        """
        import json
        from pathlib import Path
        import sys

        sys.modules["pypdfium2"] = None
        import image2editable
        from image2editable import cli, runtime
        from image2editable.store import RunStore
        from pptx import Presentation

        root = Path(sys.argv[1])
        source = root / "source.pptx"
        presentation = Presentation()
        presentation.slides.add_slide(presentation.slide_layouts[6])
        presentation.save(source)

        assert cli.build_parser().parse_args(["prepare", str(source)]).sources == [
            str(source)
        ]
        prepared = runtime.prepare_job(source, run_dir=root / "prepared")
        assert RunStore.open(prepared).read_json("run_state.json")["status"] == "prepared"
        summary = runtime.convert(
            source,
            run_dir=root / "converted",
            output_path=root / "output.pptx",
        )
        assert summary["status"] == "completed"
        assert (root / "output.pptx").read_bytes() == source.read_bytes()
        assert callable(image2editable.rerender_pdf_page)
        """
    )

    result = subprocess.run(
        [sys.executable, "-c", script, str(tmp_path)],
        cwd=Path(__file__).resolve().parents[1],
        capture_output=True,
        text=True,
        check=False,
    )

    assert result.returncode == 0, result.stderr


def test_pdf_runtime_reports_missing_pdfium_without_breaking_package(
    tmp_path: Path,
) -> None:
    script = textwrap.dedent(
        """
        from pathlib import Path
        import sys

        sys.modules["pypdfium2"] = None
        import image2editable

        source = Path(sys.argv[1]) / "source.pdf"
        source.write_bytes(b"%PDF")
        try:
            image2editable.prepare_job(source, run_dir=source.parent / "run")
        except ModuleNotFoundError as error:
            assert "pypdfium2" in str(error)
        else:
            raise AssertionError("PDF preparation unexpectedly succeeded")
        """
    )

    result = subprocess.run(
        [sys.executable, "-c", script, str(tmp_path)],
        cwd=Path(__file__).resolve().parents[1],
        capture_output=True,
        text=True,
        check=False,
    )

    assert result.returncode == 0, result.stderr
