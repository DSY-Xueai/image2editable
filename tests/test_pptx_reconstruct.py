from __future__ import annotations

import json
from pathlib import Path

import pytest

import image_to_ppt
from image2editable.pptx_reconstruct import build_reconstruction_donor


def test_build_reconstruction_donor_uses_isolated_original_slide_pipeline(
    tmp_path,
    monkeypatch,
):
    image = tmp_path / "candidate.png"
    image.write_bytes(b"candidate")
    output = tmp_path / "donor.pptx"
    work_root = tmp_path / "work"
    calls = []
    slide_data = {
        "background_original_path": str(tmp_path / "background.png"),
        "background_widescreen_path": str(tmp_path / "background-16x9.png"),
        "components": [{"path": str(tmp_path / "component.png")}],
        "text_items": [{"text": "editable"}],
        "img_width": 100,
        "img_height": 50,
        "original_image_path": str(image),
        "quality": {"p99": 1.0},
        "background_residual": {"retained_edge_ratio": 0.0},
    }

    def fake_prepare(path, lang, **kwargs):
        calls.append(("prepare", Path(path), lang, kwargs))
        return slide_data, work_root / "page_001"

    def fake_assemble(data, path, add_reference, slide_size):
        calls.append(("assemble", data, Path(path), add_reference, slide_size))
        Path(path).write_bytes(b"pptx")
        return str(path)

    monkeypatch.setattr(image_to_ppt, "_prepare_single_image", fake_prepare)
    monkeypatch.setattr(image_to_ppt, "_assemble_prepared_slide", fake_assemble)

    result = build_reconstruction_donor(
        image,
        output,
        work_root,
        decision={
            "runtime_action": "shadow_run",
            "eligible_for_shadow_run": True,
            "source_shape_id": "background",
        },
        lang="ch",
    )

    assert calls == [
        (
            "prepare",
            image,
            "ch",
            {
                "_work_root": work_root,
                "_resource_isolation": True,
            },
        ),
        (
            "assemble",
            slide_data,
            output,
            False,
            "original",
        ),
    ]
    assert result["source_shape_id"] == "background"
    assert result["components"] == 1
    assert result["text_boxes"] == 1
    assert (
        json.loads((work_root / "reconstruction_manifest.json").read_text("utf-8"))
        == result
    )


def test_build_reconstruction_donor_rejects_unapproved_action(tmp_path):
    image = tmp_path / "candidate.png"
    image.write_bytes(b"candidate")

    with pytest.raises(ValueError, match="Agent-approved shadow_run"):
        build_reconstruction_donor(
            image,
            tmp_path / "donor.pptx",
            tmp_path / "work",
            decision={
                "runtime_action": "preserve",
                "eligible_for_shadow_run": False,
                "source_shape_id": "background",
            },
        )


def test_donor_reuses_component_result_without_cv(monkeypatch, tmp_path):
    import json
    from image2editable.pptx_reconstruct import build_reconstruction_donor_from_result

    (tmp_path / "asset.png").write_bytes(b"asset")
    result = tmp_path / "component_result.json"
    import hashlib

    digest = hashlib.sha256(b"asset").hexdigest()
    result.write_text(
        json.dumps(
            {
                "provider": "host",
                "source_screenshot_sha256": "abc",
                "components": [
                    {
                        "id": "c1",
                        "state": "active",
                        "path": "asset.png",
                        "sha256": digest,
                    }
                ],
                "text_items": [],
                "assets": [{"path": str(tmp_path / "asset.png"), "sha256": digest}],
            }
        ),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        "image_to_ppt.prepare_component_layers",
        lambda *a, **k: (_ for _ in ()).throw(AssertionError("CV rerun")),
        raising=False,
    )
    manifest = build_reconstruction_donor_from_result(
        result,
        tmp_path / "donor.pptx",
        tmp_path / "work",
        source_screenshot_sha256="abc",
        provider="host",
    )
    assert manifest["components"] == 1


def test_component_result_donor_reuses_identical_manifest_without_clobber(
    tmp_path,
):
    from image2editable.pptx_reconstruct import build_reconstruction_donor_from_result
    import hashlib

    asset = tmp_path / "asset.png"
    asset.write_bytes(b"asset")
    digest = hashlib.sha256(asset.read_bytes()).hexdigest()
    result = tmp_path / "component_result.json"
    document = {
        "provider": "host",
        "source_screenshot_sha256": "abc",
        "components": [{"id": "c1", "state": "active"}],
        "text_items": [],
        "assets": [{"path": str(asset), "sha256": digest}],
    }
    result.write_text(json.dumps(document), encoding="utf-8")
    donor = tmp_path / "donor.pptx"
    work = tmp_path / "work"
    first = build_reconstruction_donor_from_result(
        result, donor, work, source_screenshot_sha256="abc", provider="host"
    )
    manifest_before = (work / "reconstruction_manifest.json").read_bytes()
    donor.unlink()

    assert (
        build_reconstruction_donor_from_result(
            result, donor, work, source_screenshot_sha256="abc", provider="host"
        )
        == first
    )

    donor.unlink()
    document["components"][0]["id"] = "c2"
    result.write_text(json.dumps(document), encoding="utf-8")
    with pytest.raises(FileExistsError):
        build_reconstruction_donor_from_result(
            result,
            donor,
            work,
            source_screenshot_sha256="abc",
            provider="host",
        )
    assert (work / "reconstruction_manifest.json").read_bytes() == manifest_before
    assert not donor.exists()


def test_nonzero_component_result_rejects_zero_donor(tmp_path):
    import json
    from image2editable.pptx_reconstruct import build_reconstruction_donor_from_result

    result = tmp_path / "component_result.json"
    result.write_text(
        json.dumps(
            {
                "provider": "host",
                "source_screenshot_sha256": "abc",
                "components": [],
                "text_items": [],
                "assets": [],
            }
        ),
        encoding="utf-8",
    )
    import pytest

    with pytest.raises(ValueError, match="non-zero"):
        build_reconstruction_donor_from_result(
            result,
            tmp_path / "donor.pptx",
            tmp_path / "work",
            source_screenshot_sha256="abc",
            provider="host",
            initial_component_count=1,
        )


def test_active_parent_and_child_are_rejected(tmp_path):
    import json
    import pytest
    from image2editable.pptx_reconstruct import build_reconstruction_donor_from_result

    p = tmp_path / "r.json"
    p.write_text(
        json.dumps(
            {
                "provider": "host",
                "source_screenshot_sha256": "s",
                "components": [
                    {"id": "p", "kind": "parent", "state": "active"},
                    {"id": "c", "kind": "child", "parent_id": "p", "state": "active"},
                ],
                "assets": [],
            }
        ),
        encoding="utf-8",
    )
    with pytest.raises(ValueError, match="parent and child"):
        build_reconstruction_donor_from_result(
            p,
            tmp_path / "d.pptx",
            tmp_path / "w",
            source_screenshot_sha256="s",
            provider="host",
        )


def test_native_component_result_hash_and_graph_escape_fail_closed(tmp_path):
    from image2editable.pptx_reconstruct import build_reconstruction_donor_from_result
    import hashlib

    run = tmp_path / "run"
    result_dir = run / "pages" / "page_001" / "reconstruction"
    result_dir.mkdir(parents=True)
    graph = result_dir / "graph.json"
    graph.write_text(json.dumps({"nodes": []}), encoding="utf-8")
    result = result_dir / "component_result.json"
    payload = {
        "provider": "host",
        "source_sha256": "a" * 64,
        "final_component_ids": [],
        "graph_ref": {
            "path": graph.relative_to(run).as_posix(),
            "sha256": hashlib.sha256(graph.read_bytes()).hexdigest(),
        },
        "accepted_asset_refs": {},
    }
    result.write_text(json.dumps(payload), encoding="utf-8")
    with pytest.raises(ValueError, match="result hash"):
        build_reconstruction_donor_from_result(
            result,
            tmp_path / "donor.pptx",
            result_dir,
            source_screenshot_sha256="a" * 64,
            provider="host",
            expected_result_sha256="0" * 64,
            run_root=run,
        )
    linked = result_dir / "component_result-linked.json"
    linked.hardlink_to(result)
    try:
        with pytest.raises(ValueError, match="regular file"):
            build_reconstruction_donor_from_result(
                result,
                tmp_path / "donor-hardlink.pptx",
                result_dir,
                source_screenshot_sha256="a" * 64,
                provider="host",
                expected_result_sha256=hashlib.sha256(result.read_bytes()).hexdigest(),
                run_root=run,
            )
    finally:
        linked.unlink()
    payload["graph_ref"]["path"] = "../../outside.json"
    result.write_text(json.dumps(payload), encoding="utf-8")
    with pytest.raises(ValueError, match="graph reference"):
        build_reconstruction_donor_from_result(
            result,
            tmp_path / "donor-escape.pptx",
            result_dir,
            source_screenshot_sha256="a" * 64,
            provider="host",
            expected_result_sha256=hashlib.sha256(result.read_bytes()).hexdigest(),
            run_root=run,
        )


def test_native_donor_preserves_component_fill_and_editable_text_style(tmp_path):
    import hashlib

    from PIL import Image
    from pptx import Presentation
    from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN

    from image2editable.pptx_reconstruct import build_reconstruction_donor_from_result

    run = tmp_path / "run"
    accepted = run / "pages/page_001/reconstruction/accepted"
    masks = accepted / "masks"
    masks.mkdir(parents=True)
    source = accepted / "source.png"
    background = accepted / "background.png"
    reconstructed = accepted / "reconstructed.png"
    text_mask = accepted / "text-mask.png"
    native = accepted / "native.json"
    mask = masks / "component_0001.png"
    Image.new("RGB", (4, 4), "white").save(source)
    Image.new("RGB", (4, 4), "white").save(background)
    Image.new("RGB", (4, 4), "red").save(reconstructed)
    Image.new("L", (4, 4), 0).save(text_mask)
    with Image.open(text_mask) as image:
        image.putpixel((1, 1), 255)
        image.save(text_mask)
    Image.new("L", (4, 4), 255).save(mask)
    native.write_text("{}", encoding="utf-8")

    def ref(path: Path) -> dict:
        return {
            "path": path.relative_to(run).as_posix(),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }

    graph = {
        "nodes": [{
            "id": "component_0001", "kind": "parent", "parent_id": None,
            "state": "frozen", "mask": "masks/component_0001.png",
            "mask_sha256": ref(mask)["sha256"], "bbox": [0, 0, 4, 4],
            "z_index": 0, "text_ids": [],
        }]
    }
    graph_path = accepted / "component-graph.json"
    graph_path.write_text(json.dumps(graph), encoding="utf-8")
    source_sha = ref(source)["sha256"]
    result = {
        "schema_version": 1, "page_id": "page_001",
        "status": "ready_for_assembly",
        "provider": "host", "source_sha256": source_sha,
        "final_component_ids": ["component_0001"],
        "graph_ref": ref(graph_path),
        "accepted_asset_refs": {
            "source": ref(source), "background": ref(background),
            "reconstructed": ref(reconstructed), "text_mask": ref(text_mask),
            "native_check": ref(native),
        },
        "text_items": [{
            "text": "优点", "box": [1, 1, 2, 1], "font_size": 14.5,
            "font": "Microsoft YaHei", "bold": True,
            "color": "#F5FAF6", "align": 1,
        }],
    }
    result_path = run / "pages/page_001/reconstruction/component_result.json"
    result_path.write_text(json.dumps(result), encoding="utf-8")
    donor = tmp_path / "donor.pptx"
    work = tmp_path / "work"

    build_reconstruction_donor_from_result(
        result_path, donor, work,
        source_screenshot_sha256=source_sha,
        provider="host",
        expected_result_sha256=hashlib.sha256(result_path.read_bytes()).hexdigest(),
        run_root=run,
    )

    with Image.open(work / "component-component_0001.png") as component:
        assert component.getchannel("A").getpixel((1, 1)) == 255
    presentation = Presentation(donor)
    text_frame = presentation.slides[0].shapes[-1].text_frame
    paragraph = text_frame.paragraphs[0]
    run_text = paragraph.runs[0]
    assert text_frame.word_wrap is False
    assert text_frame.margin_left == text_frame.margin_right == 0
    assert text_frame.margin_top == text_frame.margin_bottom == 0
    assert text_frame.vertical_anchor == MSO_VERTICAL_ANCHOR.MIDDLE
    assert paragraph.alignment == PP_ALIGN.CENTER
    assert run_text.text == "优点"
    assert run_text.font.name == "Microsoft YaHei"
    assert run_text.font.size.pt == 14.5
    assert run_text.font.bold is True
    assert str(run_text.font.color.rgb) == "F5FAF6"

    assert presentation.slides[0].shapes[0].shape_type == 13
    route = result_path.parent / "route"
    route.mkdir()
    ir = {
        "schema_version": 1,
        "page_id": "page_001",
        "canvas": {"width": 4, "height": 4},
        "objects": [{
            "id": "component_0001",
            "bbox": [0, 0, 4, 4],
            "z_index": 0,
            "source_refs": [ref(source)],
            "mask_ref": ref(mask),
            "relations": [],
            "candidate_representations": [
                {
                    "kind": "raster_component",
                    "confidence": 1.0,
                    "payload": {"asset_ref": ref(source)},
                    "evidence_refs": [],
                    "required_qa_checks": [],
                },
                {
                    "kind": "native_shape",
                    "confidence": 1.0,
                    "payload": {
                        "shape_type": "rectangle",
                        "fill_rgb": [255, 0, 0],
                    },
                    "evidence_refs": [],
                    "required_qa_checks": ["render_difference"],
                },
            ],
        }],
    }
    ir_path = route / "reconstruction-ir.json"
    ir_path.write_text(json.dumps(ir), encoding="utf-8")
    plan = {
        "schema_version": 1,
        "page_id": "page_001",
        "ir_sha256": ref(ir_path)["sha256"],
        "adapter": "pptx",
        "routes": [{
            "object_id": "component_0001",
            "selected_route": "native_shape",
            "fallback_route": "raster_component",
            "candidate_confidence": 1.0,
            "evidence_refs": [],
            "qa_requirements": ["render_difference"],
        }],
    }
    plan_path = route / "reconstruction-plan.json"
    plan_path.write_text(json.dumps(plan), encoding="utf-8")
    route_result_path = route / "route_result.json"
    route_result_path.write_text(json.dumps({
        "schema_version": 1,
        "page_id": "page_001",
        "status": "native_accepted",
        "component_result_sha256": ref(result_path)["sha256"],
        "ir_ref": ref(ir_path),
        "plan_ref": ref(plan_path),
        "qa_ref": None,
        "reason": None,
    }), encoding="utf-8")

    native_donor = tmp_path / "native-donor.pptx"
    build_reconstruction_donor_from_result(
        result_path, native_donor, tmp_path / "native-work",
        source_screenshot_sha256=source_sha,
        provider="host",
        expected_result_sha256=ref(result_path)["sha256"],
        run_root=run,
    )

    native_presentation = Presentation(native_donor)
    native_shape = native_presentation.slides[0].shapes[0]
    assert native_shape.shape_type == 1
    assert native_shape.name == "image2editable:component_0001"
