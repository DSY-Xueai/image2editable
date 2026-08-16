import hashlib
import json
import re
import subprocess
import sys
from io import BytesIO
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZIP_STORED, ZipFile, ZipInfo

import PIL
import pytest
from lxml import etree
from PIL import Image, ImageOps, ImageStat, features
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader
from pypdf.generic import ContentStream

import scripts.generate_benchmark_corpus as corpus_module

from scripts.generate_benchmark_corpus import (
    LANDSCAPE,
    PORTRAIT,
    SCHEMA_VERSION,
    generate_corpus,
)


EXPECTED_IMAGE_CASES = [
    ("01_zh_courseware", "image", "01-zh-courseware.png", 1),
    ("02_typography", "image", "02-typography.png", 1),
    ("03_flowchart", "image", "03-flowchart.png", 1),
    ("04_table_chart", "image", "04-table-chart.png", 1),
    ("05_photo_overlay", "image", "05-photo-overlay.png", 1),
    ("06_transparency_shadow", "image", "06-transparency-shadow.png", 1),
    ("07_compressed", "image", "07-compressed.jpg", 1),
    ("08_portrait", "image", "08-portrait.png", 1),
]

EXPECTED_CASES = EXPECTED_IMAGE_CASES + [
    ("09_document", "pdf", "09-document.pdf", 3),
    ("10_mixed", "pptx", "10-mixed.pptx", 3),
]

EXPECTED_ROUTES = [
    {
        "id": "images",
        "cases": [case[0] for case in EXPECTED_IMAGE_CASES],
        "pages": 8,
    },
    {"id": "pdf", "cases": ["09_document"], "pages": 3},
    {"id": "mixed_pptx", "cases": ["10_mixed"], "pages": 3},
]

GOLDEN_SHA256 = {
    "01-zh-courseware.png": "b8a4df5baf6978e0e353494a9d3ff3c991d1105daaf084b52dff36fabbe2c2f8",
    "02-typography.png": "e649cadfd47ee8dc270eb35aa36d6af018e710461e5458686decaaf83006aaca",
    "03-flowchart.png": "5d6c1ff4700a0fb80c566dc7b790f8517a5bd71d5637f4e1ab23e1ea544eb998",
    "04-table-chart.png": "5a5ed76087edbfe029f0f68c94a5080880e0d33b6bac3b6d0edb503109d409ec",
    "05-photo-overlay.png": "d9c1c57ff0e200139cd575f860ed9a5c893c2ae59c57b11d642f6d0bc067c181",
    "06-transparency-shadow.png": "590a5cf4796d575b6e1b71543e82d5624c64be12eac88b1e3c304e587671c1ba",
    "07-compressed.jpg": "9f055966531e0fa4bbb403e368224db1e5b73d685fc48bc6d30faede65fabad1",
    "08-portrait.png": "4891659393a0c3bfdf8cfa5be58b46c9de7cc12830f7f0fec0e7dd4dc9d4a41d",
    "09-document.pdf": "3518c9250d8bc00b66fb4648ddd7398658de60debe2374bdd01cd051ba669658",
    "10-mixed.pptx": "a5432c5f9ea30ed3c8626c03b80f855113d2665ba7ab142f7c361d624fd7ddf5",
    "manifest.json": "139e7031b40afbfda7f257e22fbcd47051db20a9bd6c77c49d73d83155bc7ca3",
}


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _canonical_sha256(value: object) -> str:
    payload = json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        allow_nan=False,
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def test_all_golden_sha256_values_are_approved() -> None:
    assert all(
        re.fullmatch(r"[0-9a-f]{64}", value) for value in GOLDEN_SHA256.values()
    )


def test_corpus_git_attributes_disable_text_conversion() -> None:
    root = Path(__file__).resolve().parents[1]
    expected_lines = {
        "benchmark/corpus/*.png binary",
        "benchmark/corpus/*.jpg binary",
        "benchmark/corpus/*.pdf binary",
        "benchmark/corpus/*.pptx binary",
        "benchmark/corpus/manifest.json text eol=lf",
    }
    attribute_lines = {
        line.strip()
        for line in (root / ".gitattributes").read_text(encoding="utf-8").splitlines()
        if line.strip() and not line.lstrip().startswith("#")
    }
    assert expected_lines <= attribute_lines

    for relative_path in (
        "benchmark/corpus/01-zh-courseware.png",
        "benchmark/corpus/07-compressed.jpg",
        "benchmark/corpus/09-document.pdf",
        "benchmark/corpus/10-mixed.pptx",
    ):
        result = subprocess.run(
            ["git", "check-attr", "diff", "merge", "text", "--", relative_path],
            cwd=root,
            capture_output=True,
            text=True,
            check=False,
        )
        assert result.returncode == 0, result.stderr
        values = {
            line.split(": ", 2)[1]: line.split(": ", 2)[2]
            for line in result.stdout.splitlines()
        }
        assert values == {"diff": "unset", "merge": "unset", "text": "unset"}

    result = subprocess.run(
        [
            "git",
            "check-attr",
            "text",
            "eol",
            "--",
            "benchmark/corpus/manifest.json",
        ],
        cwd=root,
        capture_output=True,
        text=True,
        check=False,
    )
    assert result.returncode == 0, result.stderr
    values = {
        line.split(": ", 2)[1]: line.split(": ", 2)[2]
        for line in result.stdout.splitlines()
    }
    assert values == {"text": "set", "eol": "lf"}


def test_benchmark_outputs_and_private_corpus_are_ignored_but_public_corpus_is_not() -> None:
    root = Path(__file__).resolve().parents[1]
    ignore_lines = (root / ".gitignore").read_text(encoding="utf-8").splitlines()
    assert ignore_lines[-2:] == ["/benchmark-results/", "/benchmark/private/"]

    ignored = subprocess.run(
        [
            "git",
            "check-ignore",
            "--no-index",
            "--",
            "benchmark-results/report.json",
            "benchmark/private/manifest.json",
        ],
        cwd=root,
        capture_output=True,
        text=True,
        check=False,
    )
    assert ignored.returncode == 0, ignored.stderr
    assert ignored.stdout.splitlines() == [
        "benchmark-results/report.json",
        "benchmark/private/manifest.json",
    ]

    public = subprocess.run(
        [
            "git",
            "check-ignore",
            "--no-index",
            "--",
            "benchmark/corpus/manifest.json",
        ],
        cwd=root,
        capture_output=True,
        text=True,
        check=False,
    )
    assert public.returncode == 1, public.stderr
    assert public.stdout == ""


def test_generate_corpus_rejects_unapproved_encoder_before_writing(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    message = (
        "deterministic benchmark encoder unavailable: requires Pillow 10.4.0, "
        "libjpeg_turbo 3.0.3, zlib 1.3.1, freetype2 2.13.2"
    )

    wrong_freetype = tmp_path / "wrong-freetype"
    with monkeypatch.context() as patch:
        patch.setattr(
            features,
            "version_module",
            lambda name: "2.13.3" if name == "freetype2" else None,
        )
        with pytest.raises(RuntimeError, match=f"^{re.escape(message)}$"):
            generate_corpus(wrong_freetype)
    assert not wrong_freetype.exists()

    wrong_pillow = tmp_path / "wrong-pillow"
    with monkeypatch.context() as patch:
        patch.setattr(PIL, "__version__", "11.1.0")
        with pytest.raises(RuntimeError, match=f"^{re.escape(message)}$"):
            generate_corpus(wrong_pillow)
    assert not wrong_pillow.exists()

    wrong_jpeg = tmp_path / "wrong-jpeg"
    wrong_jpeg.mkdir()
    with monkeypatch.context() as patch:
        patch.setattr(
            features,
            "version_feature",
            lambda name: "3.0.2" if name == "libjpeg_turbo" else None,
        )
        with pytest.raises(RuntimeError, match=f"^{re.escape(message)}$"):
            generate_corpus(wrong_jpeg)
    assert not any(wrong_jpeg.iterdir())

    wrong_zlib = tmp_path / "wrong-zlib"
    with monkeypatch.context() as patch:
        patch.setattr(
            features,
            "version_codec",
            lambda name: "1.3.0" if name == "zlib" else None,
        )
        with pytest.raises(RuntimeError, match=f"^{re.escape(message)}$"):
            generate_corpus(wrong_zlib)
    assert not wrong_zlib.exists()


@pytest.mark.parametrize(
    ("label", "patch_target", "wrong_version"),
    [
        ("reportlab", "reportlab.Version", "5.0.1"),
        ("python-pptx", "pptx.__version__", "1.0.3"),
        ("stdlib-zlib", "zlib.ZLIB_RUNTIME_VERSION", "1.3.0"),
        ("lxml", "lxml.etree.LXML_VERSION", (6, 1, 2, 0)),
    ],
)
def test_generate_corpus_rejects_unapproved_document_environment_before_writing(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    label: str,
    patch_target: str,
    wrong_version: object,
) -> None:
    output = tmp_path / label
    message = (
        "deterministic benchmark document environment unavailable: requires "
        "ReportLab 5.0.0, python-pptx 1.0.2, stdlib zlib 1.3.1, lxml 6.1.1"
    )

    monkeypatch.setattr(patch_target, wrong_version)
    with pytest.raises(RuntimeError, match=f"^{re.escape(message)}$"):
        generate_corpus(output)

    assert not output.exists()


def test_generate_corpus_writes_fixed_manifest_and_nonblank_images(
    tmp_path: Path,
) -> None:
    output = tmp_path / "corpus"

    generate_corpus(output)

    manifest_bytes = (output / "manifest.json").read_bytes()
    assert manifest_bytes.endswith(b"\n")
    assert b"\r" not in manifest_bytes
    manifest = json.loads(manifest_bytes)
    assert SCHEMA_VERSION == 1
    assert LANDSCAPE == (1600, 900)
    assert PORTRAIT == (900, 1600)
    assert list(manifest) == ["schema_version", "cases", "routes", "corpus_sha256"]
    assert set(manifest) == {"schema_version", "cases", "routes", "corpus_sha256"}
    assert manifest["schema_version"] == 1
    assert [
        (case["id"], case["kind"], case["path"], case["pages"])
        for case in manifest["cases"]
    ] == EXPECTED_CASES
    assert manifest["routes"] == EXPECTED_ROUTES
    assert manifest["corpus_sha256"] == _canonical_sha256(
        {
            "schema_version": manifest["schema_version"],
            "cases": manifest["cases"],
            "routes": manifest["routes"],
        }
    )

    for case in manifest["cases"]:
        assert list(case) == ["id", "kind", "path", "pages", "bytes", "sha256"]
        path = output / case["path"]
        assert case["bytes"] == path.stat().st_size
        assert case["sha256"] == _sha256(path)

    for index, case in enumerate(manifest["cases"][:8], start=1):
        path = output / case["path"]
        with Image.open(path) as image:
            assert image.format == ("JPEG" if index == 7 else "PNG")
            assert image.size == (PORTRAIT if index == 8 else LANDSCAPE)
            extrema = ImageStat.Stat(image.convert("RGB")).extrema
            assert any(low != high for low, high in extrema)
            if index == 6:
                assert image.getpixel((0, 899)) != image.getpixel((1599, 899))

    source_path = Path(__file__).resolve().parents[1] / "docs/images/demo-source-1.png"
    with (
        Image.open(source_path) as source,
        Image.open(output / "01-zh-courseware.png") as generated,
    ):
        expected = ImageOps.fit(
            source.convert("RGB"), LANDSCAPE, Image.Resampling.LANCZOS
        )
        assert generated.convert("RGB").tobytes() == expected.tobytes()


def test_generate_corpus_writes_valid_three_page_pdf(tmp_path: Path) -> None:
    output = tmp_path / "corpus"
    generate_corpus(output)

    reader = PdfReader(output / "09-document.pdf")

    assert len(reader.pages) == 3
    assert [
        (float(page.mediabox.width), float(page.mediabox.height))
        for page in reader.pages
    ] == [(1600.0, 900.0)] * 3


def test_pdf_pages_preserve_source_pixels_and_centered_placement(
    tmp_path: Path,
) -> None:
    output = tmp_path / "corpus"
    generate_corpus(output)
    reader = PdfReader(output / "09-document.pdf")
    expected = [
        ("01-zh-courseware.png", [1600.0, 0.0, 0.0, 900.0, 0.0, 0.0]),
        ("04-table-chart.png", [1600.0, 0.0, 0.0, 900.0, 0.0, 0.0]),
        ("08-portrait.png", [506.25, 0.0, 0.0, 900.0, 546.875, 0.0]),
    ]

    assert corpus_module.PDF_PAGE_SOURCES == tuple(name for name, _ in expected)
    for page, (source_name, expected_matrix) in zip(reader.pages, expected):
        images = list(page.images)
        assert len(images) == 1
        with Image.open(output / source_name) as source:
            decoded = images[0].image.convert("RGB")
            expected_image = source.convert("RGB")
            assert decoded.size == expected_image.size
            assert decoded.tobytes() == expected_image.tobytes()

        drawing_operations = [
            (operands, operator)
            for operands, operator in ContentStream(
                page.get_contents(), reader
            ).operations
            if operator in (b"cm", b"Do")
        ]
        assert [operator for _, operator in drawing_operations] == [
            b"cm",
            b"cm",
            b"Do",
        ]
        assert [float(value) for value in drawing_operations[0][0]] == pytest.approx(
            [1.0, 0.0, 0.0, 1.0, 0.0, 0.0], abs=1e-9
        )
        assert [float(value) for value in drawing_operations[1][0]] == pytest.approx(
            expected_matrix, abs=1e-9
        )
        image_name = images[0].name.removesuffix(".png")
        assert str(drawing_operations[2][0][0]).lstrip("/") == image_name


def test_generate_corpus_writes_expected_mixed_pptx_structure(tmp_path: Path) -> None:
    output = tmp_path / "corpus"
    generate_corpus(output)

    deck = Presentation(output / "10-mixed.pptx")
    assert len(deck.slides) == 3
    assert deck.slide_width * 9 == deck.slide_height * 16

    first, second, third = deck.slides
    first_pictures = [
        shape for shape in first.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    first_text = [
        shape for shape in first.shapes if shape.has_text_frame and shape.text.strip()
    ]
    first_sizes = [
        run.font.size.pt
        for shape in first_text
        for paragraph in shape.text_frame.paragraphs
        for run in paragraph.runs
        if run.text.strip() and run.font.size is not None
    ]
    assert first_text
    assert not first_pictures
    assert min(first_sizes) >= 16
    assert max(first_sizes) >= 35

    second_pictures = [
        shape for shape in second.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    assert len(second.shapes) == 1
    assert len(second_pictures) == 1
    picture = second_pictures[0]
    assert (picture.left, picture.top) == (0, 0)
    assert (picture.width, picture.height) == (deck.slide_width, deck.slide_height)

    third_pictures = [
        shape for shape in third.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    ]
    third_text = [
        shape for shape in third.shapes if shape.has_text_frame and shape.text.strip()
    ]
    assert third_pictures
    assert third_text
    body = next(shape for shape in third_text if shape.text.startswith("Native text"))
    assert body.text_frame.word_wrap is True
    assert any(
        shape.shape_type != MSO_SHAPE_TYPE.PICTURE for shape in third.shapes
    )


def test_normalize_pptx_canonicalizes_cross_platform_zip_metadata() -> None:
    def source_zip(create_system: int, external_attr: int) -> bytes:
        payload = BytesIO()
        with ZipFile(payload, "w") as archive:
            for filename in ("z-last.txt", "a-first.txt"):
                info = ZipInfo(filename, (2026, 8, 15, 12, 30, 0))
                info.compress_type = ZIP_STORED
                info.create_system = create_system
                info.external_attr = external_attr
                archive.writestr(info, filename.encode("ascii"))
        return payload.getvalue()

    windows_source = source_zip(0, 0x20)
    posix_source = source_zip(3, 0o100644 << 16)
    windows = corpus_module._normalize_pptx(windows_source)
    posix = corpus_module._normalize_pptx(posix_source)

    assert windows == posix
    with ZipFile(BytesIO(windows)) as archive:
        infos = archive.infolist()
        assert [info.filename for info in infos] == ["a-first.txt", "z-last.txt"]
        for info in infos:
            assert info.date_time == (1980, 1, 1, 0, 0, 0)
            assert info.compress_type == ZIP_DEFLATED
            assert info.create_system == 0
            assert info.external_attr == 0o600 << 16


def test_generated_pptx_has_canonical_zip_and_fixed_format_metadata(
    tmp_path: Path,
) -> None:
    output = tmp_path / "corpus"
    generate_corpus(output)

    with ZipFile(output / "10-mixed.pptx") as archive:
        infos = archive.infolist()
        assert [info.filename for info in infos] == sorted(
            info.filename for info in infos
        )
        for info in infos:
            assert info.date_time == (1980, 1, 1, 0, 0, 0)
            assert info.compress_type == ZIP_DEFLATED
            assert info.create_system == 0
            assert info.external_attr == 0o600 << 16

        app = etree.fromstring(archive.read("docProps/app.xml"))
        app_namespace = "http://schemas.openxmlformats.org/officeDocument/2006/extended-properties"
        assert app.findtext(f"{{{app_namespace}}}Slides") == "3"
        assert (
            app.findtext(f"{{{app_namespace}}}PresentationFormat")
            == "On-screen Show (16:9)"
        )

        presentation = etree.fromstring(archive.read("ppt/presentation.xml"))
        presentation_namespace = (
            "http://schemas.openxmlformats.org/presentationml/2006/main"
        )
        slide_size = presentation.find(f"{{{presentation_namespace}}}sldSz")
        assert slide_size is not None
        assert slide_size.get("type") == "screen16x9"
        assert int(slide_size.get("cx")) * 9 == int(slide_size.get("cy")) * 16


def test_generate_corpus_refuses_nonempty_output_before_writing(
    tmp_path: Path,
) -> None:
    output = tmp_path / "corpus"
    output.mkdir()
    conflict = output / "01-zh-courseware.png"
    conflict.write_bytes(b"keep me")

    with pytest.raises(FileExistsError):
        generate_corpus(output)

    assert {path.name for path in output.iterdir()} == {conflict.name}
    assert conflict.read_bytes() == b"keep me"


def test_generate_corpus_is_byte_reproducible_in_two_empty_directories(
    tmp_path: Path,
) -> None:
    first = tmp_path / "first"
    second = tmp_path / "second"
    first.mkdir()
    second.mkdir()

    generate_corpus(first)
    generate_corpus(second)

    expected_names = {case[2] for case in EXPECTED_CASES} | {"manifest.json"}
    assert len(expected_names) == 11
    assert {path.name for path in first.iterdir()} == expected_names
    assert {path.name for path in second.iterdir()} == expected_names
    assert {
        name: (first / name).read_bytes() for name in expected_names
    } == {name: (second / name).read_bytes() for name in expected_names}
    assert {name: _sha256(first / name) for name in expected_names} == GOLDEN_SHA256
    assert sum((first / name).stat().st_size for name in expected_names) <= 12 * 1024 * 1024


def test_tracked_corpus_matches_fresh_generation(tmp_path: Path) -> None:
    root = Path(__file__).resolve().parents[1]
    tracked = root / "benchmark/corpus"
    generated = tmp_path / "generated"
    generate_corpus(generated)

    assert tracked.is_dir()
    assert sorted(path.name for path in tracked.iterdir()) == sorted(
        path.name for path in generated.iterdir()
    )
    assert {
        path.name: path.read_bytes() for path in tracked.iterdir()
    } == {path.name: path.read_bytes() for path in generated.iterdir()}


def test_cli_generates_manifest_from_checkout(tmp_path: Path) -> None:
    root = Path(__file__).resolve().parents[1]
    output = tmp_path / "cli-corpus"

    result = subprocess.run(
        [
            sys.executable,
            "scripts/generate_benchmark_corpus.py",
            "--output",
            str(output),
        ],
        cwd=root,
        capture_output=True,
        text=True,
        check=False,
    )

    assert result.returncode == 0, result.stderr
    manifest = json.loads((output / "manifest.json").read_text(encoding="utf-8"))
    assert manifest["schema_version"] == 1
    assert len(manifest["cases"]) == 10
    assert sum(case["pages"] for case in manifest["cases"]) == 14
