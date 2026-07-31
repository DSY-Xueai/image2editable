from __future__ import annotations

import hashlib
import errno
import re
import zipfile
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from xml.etree import ElementTree as ET

import pytest
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR_TYPE, MSO_SHAPE
from pptx.util import Inches
from PIL import Image

from image2editable import pptx_input
from image2editable.pptx_input import scan_pptx
from image2editable.resources import safe_default_policy
from image2editable.store import RunStore


def test_picture_slide_coverage_uses_clipped_slide_intersection():
    picture = {
        "x": -10,
        "y": -10,
        "cx": 100,
        "cy": 120,
        "transform_reliable": True,
    }

    assert pptx_input.picture_slide_coverage(picture, 100, 100) == pytest.approx(
        0.9
    )


@pytest.mark.parametrize(("width", "height"), [(0, 100), (100, -1), (True, 100)])
def test_picture_slide_coverage_rejects_invalid_slide_size(width, height):
    picture = {
        "x": 0,
        "y": 0,
        "cx": 100,
        "cy": 100,
        "transform_reliable": True,
    }

    with pytest.raises(ValueError, match="positive"):
        pptx_input.picture_slide_coverage(picture, width, height)


def test_picture_slide_coverage_rejects_unreliable_transform():
    with pytest.raises(ValueError, match="reliable"):
        pptx_input.picture_slide_coverage(
            {
                "x": 0,
                "y": 0,
                "cx": 100,
                "cy": 100,
                "transform_reliable": False,
            },
            100,
            100,
        )


def test_picture_slide_coverage_rejects_unbounded_ooxml_integer():
    with pytest.raises(ValueError, match="reliable"):
        pptx_input.picture_slide_coverage(
            {
                "x": 0,
                "y": 0,
                "cx": 10**1000,
                "cy": 100,
                "transform_reliable": True,
            },
            100,
            100,
        )


@pytest.mark.parametrize(
    ("picture", "slide_width", "slide_height"),
    [
        ({"x": 0, "y": 0, "cx": 100, "cy": 100}, 100.0, 100),
        ({"x": 0.0, "y": 0, "cx": 100, "cy": 100}, 100, 100),
        ({"x": 0, "y": 0, "cx": 100.0, "cy": 100}, 100, 100),
        ({"x": 0, "y": 0, "cx": 100, "cy": 100}, 5e-324, 5e-324),
    ],
)
def test_picture_slide_coverage_rejects_non_integer_geometry(
    picture, slide_width, slide_height
):
    picture["transform_reliable"] = True

    with pytest.raises(ValueError):
        pptx_input.picture_slide_coverage(
            picture, slide_width, slide_height
        )


@pytest.mark.parametrize(
    ("value", "expected"),
    [
        ("0", 0),
        ("+9223372036854775807", 9223372036854775807),
        ("-9223372036854775808", -9223372036854775808),
        ("9223372036854775808", None),
        ("-9223372036854775809", None),
        ("000000000000000000000", None),
        ("１２", None),
        (" 1", None),
        ("+", None),
        ("", None),
    ],
)
def test_parse_ooxml_integer_is_ascii_bounded_and_python310_safe(
    value, expected
):
    assert pptx_input._parse_ooxml_integer(value) == expected


def test_publish_pptx_no_clobber_preserves_existing_target(tmp_path):
    temporary = tmp_path / "temporary.pptx"
    target = tmp_path / "output.pptx"
    temporary.write_bytes(b"new")
    target.write_bytes(b"existing")

    with pytest.raises(FileExistsError, match="already exists"):
        pptx_input._publish_pptx_no_clobber(temporary, target)

    assert target.read_bytes() == b"existing"
    assert temporary.read_bytes() == b"new"


def test_publish_pptx_no_clobber_returns_fixed_output_identity(tmp_path):
    temporary = tmp_path / "temporary.pptx"
    output = tmp_path / "output.pptx"
    temporary.write_bytes(b"published")

    token = pptx_input._publish_pptx_no_clobber(temporary, output)
    status = output.lstat()

    assert token == {
        "version": 1,
        "path": str(output),
        "dev": status.st_dev,
        "ino": status.st_ino,
        "mode": status.st_mode,
        "size": status.st_size,
        "mtime_ns": status.st_mtime_ns,
        "sha256": hashlib.sha256(b"published").hexdigest(),
    }


def test_execute_pptx_preserve_falls_back_when_hardlinks_are_unsupported(
    tmp_path, monkeypatch
):
    source = _large_picture_file(tmp_path, "fallback.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")

    def unsupported(_source, _target):
        raise OSError(errno.EPERM, "hardlinks unsupported")

    monkeypatch.setattr(pptx_input.os, "link", unsupported)

    summary = pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    output = run_dir / "final" / "output.pptx"
    assert output.read_bytes() == source.read_bytes()
    assert summary["output_sha256"] == hashlib.sha256(output.read_bytes()).hexdigest()
    status = output.lstat()
    assert summary["_output_identity"] == {
        "version": 1,
        "path": str(output),
        "dev": status.st_dev,
        "ino": status.st_ino,
        "mode": status.st_mode,
        "size": status.st_size,
        "mtime_ns": status.st_mtime_ns,
        "sha256": summary["output_sha256"],
    }
    assert not list(output.parent.glob(".*.tmp"))


def test_publish_fallback_never_overwrites_existing_target(
    tmp_path, monkeypatch
):
    temporary = tmp_path / "temporary.pptx"
    output = tmp_path / "output.pptx"
    temporary.write_bytes(b"new")
    output.write_bytes(b"existing")
    monkeypatch.setattr(
        pptx_input.os,
        "link",
        lambda _source, _target, **_kwargs: (_ for _ in ()).throw(
            OSError(errno.EPERM, "hardlinks unsupported")
        ),
    )

    with pytest.raises(FileExistsError, match="already exists"):
        pptx_input._publish_pptx_no_clobber(temporary, output)

    assert output.read_bytes() == b"existing"


def test_publish_fallback_allows_only_one_concurrent_winner(
    tmp_path, monkeypatch
):
    first = tmp_path / "first.pptx"
    second = tmp_path / "second.pptx"
    output = tmp_path / "output.pptx"
    first.write_bytes(b"first")
    second.write_bytes(b"second")
    monkeypatch.setattr(
        pptx_input.os,
        "link",
        lambda _source, _target, **_kwargs: (_ for _ in ()).throw(
            OSError(errno.EPERM, "hardlinks unsupported")
        ),
    )

    def publish(source):
        try:
            pptx_input._publish_pptx_no_clobber(source, output)
            return "success"
        except FileExistsError:
            return "exists"

    with ThreadPoolExecutor(max_workers=2) as executor:
        results = list(executor.map(publish, (first, second)))

    assert sorted(results) == ["exists", "success"]
    assert output.read_bytes() in {b"first", b"second"}


def test_publish_fallback_removes_only_its_partial_output_on_failure(
    tmp_path, monkeypatch
):
    temporary = tmp_path / "temporary.pptx"
    output = tmp_path / "output.pptx"
    temporary.write_bytes(b"new")
    monkeypatch.setattr(
        pptx_input.os,
        "link",
        lambda _source, _target: (_ for _ in ()).throw(
            OSError(errno.EPERM, "hardlinks unsupported")
        ),
    )
    monkeypatch.setattr(
        pptx_input,
        "_copy_pptx_to_fd",
        lambda _source, _fd: (_ for _ in ()).throw(
            RuntimeError("copy failed")
        ),
        raising=False,
    )

    with pytest.raises(RuntimeError, match="copy failed"):
        pptx_input._publish_pptx_no_clobber(temporary, output)

    assert not output.exists()
    assert temporary.read_bytes() == b"new"


def test_publish_fallback_failure_preserves_concurrent_replacement(
    tmp_path, monkeypatch
):
    temporary = tmp_path / "temporary.pptx"
    output = tmp_path / "output.pptx"
    replacement = tmp_path / "replacement.pptx"
    temporary.write_bytes(b"new")
    replacement.write_bytes(b"user replacement")
    monkeypatch.setattr(
        pptx_input.os,
        "link",
        lambda _source, _target, **_kwargs: (_ for _ in ()).throw(
            OSError(errno.EPERM, "hardlinks unsupported")
        ),
    )
    original_sha256_file = pptx_input.sha256_file
    replaced = False

    def replace_then_fail(path):
        nonlocal replaced
        if Path(path) == output and not replaced:
            replaced = True
            pptx_input.os.replace(replacement, output)
            raise RuntimeError("verification failed")
        return original_sha256_file(path)

    monkeypatch.setattr(pptx_input, "sha256_file", replace_then_fail)

    with pytest.raises(RuntimeError, match="verification failed"):
        pptx_input._publish_pptx_no_clobber(temporary, output)

    assert output.read_bytes() == b"user replacement"
    assert not list(tmp_path.glob(".output.pptx.recovery-*.tmp"))


def test_restore_isolated_symlink_never_follows_target(tmp_path, monkeypatch):
    target = tmp_path / "user.pptx"
    isolated = tmp_path / "isolated.pptx"
    output = tmp_path / "output.pptx"
    target.write_bytes(b"user")
    try:
        isolated.symlink_to(target)
    except OSError as error:
        pytest.skip(f"symlink unavailable: {error}")
    calls = []

    def unsupported(_source, _target, **kwargs):
        calls.append(kwargs)
        raise NotImplementedError("no no-follow hardlinks")

    monkeypatch.setattr(pptx_input.os, "link", unsupported)

    with pytest.raises(NotImplementedError, match="no-follow"):
        pptx_input._restore_isolated_file(isolated, output)

    assert calls == [{"follow_symlinks": False}]
    assert isolated.is_symlink()
    assert target.read_bytes() == b"user"
    assert not output.exists()


def test_scan_pptx_inventories_slide_background_image_before_shapes(tmp_path):
    source = _slide_background_picture_file(tmp_path, "background.pptx")
    _move_picture_to_slide_background(source)

    result = scan_pptx(source)
    objects = result["slides"][0]["objects"]

    assert [item["name"] for item in objects] == [
        "Slide Background Image",
        "Overlay",
    ]
    background = objects[0]
    assert {
        key: background[key]
        for key in (
            "shape_id",
            "name",
            "type",
            "z_order",
            "group_path",
            "inside_group",
            "slide_part",
            "x",
            "y",
            "cx",
            "cy",
            "rotation",
            "rotation_degrees",
            "flip_h",
            "flip_v",
            "crop_left",
            "crop_top",
            "crop_right",
            "crop_bottom",
            "has_extension",
            "has_timing_reference",
            "slide_coverage",
            "pixel_width",
            "pixel_height",
            "action",
            "safety_reasons",
        )
    } == {
        "shape_id": "background",
        "name": "Slide Background Image",
        "type": "slide_background_image",
        "z_order": -1,
        "group_path": [],
        "inside_group": False,
        "slide_part": "ppt/slides/slide1.xml",
        "x": 0,
        "y": 0,
        "cx": result["slide_width"],
        "cy": result["slide_height"],
        "rotation": 0,
        "rotation_degrees": 0.0,
        "flip_h": False,
        "flip_v": False,
        "crop_left": 0,
        "crop_top": 0,
        "crop_right": 0,
        "crop_bottom": 0,
        "has_extension": False,
        "has_timing_reference": False,
        "slide_coverage": 1.0,
        "pixel_width": 40,
        "pixel_height": 20,
        "action": "candidate",
        "safety_reasons": [],
    }
    assert len(background["relationships"]) == 1
    assert background["primary_relationship"] == background["relationships"][0]
    assert background["blip_sources"] == [
        {"kind": "embed", "id": background["relationships"][0]["id"]}
    ]
    assert background["media_valid"] is True
    assert re.fullmatch(r"[0-9a-f]{64}", background["xml_c14n_sha256"])
    assert "_transform_reliable" not in background
    assert objects[1]["z_order"] == 0


def test_scan_pptx_preserves_cropped_slide_background_image(tmp_path):
    source = _slide_background_picture_file(tmp_path, "cropped-background.pptx")
    _move_picture_to_slide_background(source, crop_left=10000)

    background = scan_pptx(source)["slides"][0]["objects"][0]

    assert background["crop_left"] == 10000
    assert background["action"] == "preserve"
    assert background["safety_reasons"] == ["nonzero_crop"]


@pytest.mark.parametrize("kind", ["solid", "gradient", "theme"])
def test_scan_pptx_ignores_slide_background_without_blip_fill(tmp_path, kind):
    source = _slide_background_picture_file(tmp_path, f"{kind}.pptx")
    _move_picture_to_slide_background(source, kind=kind)

    objects = scan_pptx(source)["slides"][0]["objects"]

    assert [(item["name"], item["z_order"]) for item in objects] == [
        ("Overlay", 0)
    ]


@pytest.mark.parametrize(
    ("mutation", "reason"),
    [
        ("external", "external_relationship"),
        ("missing", "missing_media"),
        ("invalid", "invalid_media"),
        ("extension", "unsupported_extension"),
    ],
)
def test_scan_pptx_preserves_unsafe_slide_background_image(
    tmp_path, mutation, reason
):
    source = _slide_background_picture_file(tmp_path, f"{mutation}.pptx")
    if mutation == "external":
        _make_picture_relationship_external(source)
    elif mutation == "missing":
        _remove_member(source, "ppt/media/image1.png")
    elif mutation == "invalid":
        _replace_member(source, "ppt/media/image1.png", lambda _: b"not an image")
    _move_picture_to_slide_background(
        source,
        extension=mutation == "extension",
    )

    background = scan_pptx(source)["slides"][0]["objects"][0]

    assert background["type"] == "slide_background_image"
    assert background["action"] == "preserve"
    assert background["safety_reasons"] == [reason]


def test_scan_pptx_marks_only_large_structurally_safe_picture_candidate(
    tmp_path,
):
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_path = tmp_path / "candidate.png"
    Image.new("RGB", (40, 20), "navy").save(image_path)
    picture = slide.shapes.add_picture(
        str(image_path),
        Inches(0.5),
        Inches(0.25),
        Inches(9),
        Inches(4.5),
    )
    picture.name = "Candidate"
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1))
    source = tmp_path / "candidate.pptx"
    presentation.save(source)

    objects = scan_pptx(source)["slides"][0]["objects"]

    assert objects[0]["action"] == "candidate"
    assert objects[0]["safety_reasons"] == []
    assert objects[0]["slide_coverage"] == pytest.approx(0.81)
    assert objects[1]["action"] == "preserve"
    assert objects[1]["safety_reasons"] == []


def test_scan_pptx_records_stable_picture_safety_reasons(tmp_path):
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_path = tmp_path / "rules.png"
    Image.new("RGB", (40, 20), "navy").save(image_path)

    logo = slide.shapes.add_picture(
        str(image_path), Inches(0), Inches(0), Inches(1), Inches(0.5)
    )
    logo.name = "Logo"
    rotated = slide.shapes.add_picture(
        str(image_path), Inches(0.5), Inches(0.25), Inches(9), Inches(4.5)
    )
    rotated.name = "Rotated"
    rotated.rotation = 10
    cropped = slide.shapes.add_picture(
        str(image_path), Inches(0.5), Inches(0.25), Inches(9), Inches(4.5)
    )
    cropped.name = "Cropped"
    cropped.crop_left = 0.1
    group = slide.shapes.add_group_shape()
    group.name = "Group"
    grouped = group.shapes.add_picture(
        str(image_path), Inches(0.5), Inches(0.25), Inches(9), Inches(4.5)
    )
    grouped.name = "Grouped"
    flipped = slide.shapes.add_picture(
        str(image_path), Inches(0.5), Inches(0.25), Inches(9), Inches(4.5)
    )
    flipped.name = "Flipped"
    extended = slide.shapes.add_picture(
        str(image_path), Inches(0.5), Inches(0.25), Inches(9), Inches(4.5)
    )
    extended.name = "Extended"
    unreliable = slide.shapes.add_picture(
        str(image_path), Inches(0.5), Inches(0.25), Inches(9), Inches(4.5)
    )
    unreliable.name = "Unreliable"
    source = tmp_path / "rules.pptx"
    presentation.save(source)
    _mutate_picture_safety(source, "Flipped", flip_h=True)
    _mutate_picture_safety(source, "Extended", extension=True, timing=True)
    _mutate_picture_safety(source, "Unreliable", remove_transform=True)

    by_name = {
        item["name"]: item for item in scan_pptx(source)["slides"][0]["objects"]
    }

    assert by_name["Logo"]["safety_reasons"] == ["coverage_below_threshold"]
    assert by_name["Rotated"]["safety_reasons"] == ["rotation_or_flip"]
    assert by_name["Cropped"]["safety_reasons"] == ["nonzero_crop"]
    assert by_name["Grouped"]["safety_reasons"] == ["inside_group"]
    assert by_name["Flipped"]["safety_reasons"] == ["rotation_or_flip"]
    assert by_name["Extended"]["safety_reasons"] == ["unsupported_extension"]
    assert by_name["Unreliable"]["safety_reasons"] == ["unreliable_transform"]
    assert all(
        by_name[name]["action"] == "preserve"
        for name in (
            "Logo",
            "Rotated",
            "Cropped",
            "Grouped",
            "Flipped",
            "Extended",
            "Unreliable",
        )
    )


def test_scan_pptx_marks_unbounded_transform_unreliable(tmp_path):
    source = _large_picture_file(tmp_path, "unbounded-transform.pptx")
    _mutate_picture_safety(source, "Picture 1", huge_extent=True)

    picture = scan_pptx(source)["slides"][0]["objects"][0]

    assert picture["action"] == "preserve"
    assert picture["safety_reasons"] == ["unreliable_transform"]


@pytest.mark.parametrize(
    ("mutation", "reason"),
    [
        ("external", "external_relationship"),
        ("missing", "missing_media"),
        ("non_image", "missing_media"),
    ],
)
def test_scan_pptx_preserves_picture_without_internal_image_media(
    tmp_path, mutation, reason
):
    source = _large_picture_file(tmp_path, f"{mutation}.pptx")
    if mutation == "external":
        _make_picture_relationship_external(source)
    elif mutation == "missing":
        _remove_member(source, "ppt/media/image1.png")
    else:
        _set_content_type(source, "/ppt/media/image1.png", "application/xml")

    picture = scan_pptx(source)["slides"][0]["objects"][0]

    assert picture["action"] == "preserve"
    assert picture["safety_reasons"] == [reason]


@pytest.mark.parametrize("mutation", ["truncated", "format_mismatch"])
def test_scan_pptx_preserves_invalid_picture_media(tmp_path, mutation):
    source = _large_picture_file(tmp_path, f"{mutation}.pptx")
    if mutation == "truncated":
        _replace_member(
            source,
            "ppt/media/image1.png",
            lambda contents: contents[:45],
        )
    else:
        replacement = tmp_path / "replacement.gif"
        Image.new("RGB", (40, 20), "navy").save(replacement, format="GIF")
        replacement_bytes = replacement.read_bytes()
        _replace_member(
            source,
            "ppt/media/image1.png",
            lambda _: replacement_bytes,
        )

    picture = scan_pptx(source)["slides"][0]["objects"][0]

    assert picture["action"] == "preserve"
    assert picture["safety_reasons"] == ["invalid_media"]
    assert picture["media_valid"] is False


def test_scan_pptx_preserves_pillow_decompression_bomb_warning(
    tmp_path, monkeypatch
):
    source = _large_picture_file(tmp_path, "bomb-warning.pptx")
    monkeypatch.setattr(Image, "MAX_IMAGE_PIXELS", 500)

    picture = scan_pptx(source)["slides"][0]["objects"][0]

    assert picture["action"] == "preserve"
    assert picture["safety_reasons"] == ["invalid_media"]
    assert picture["media_valid"] is False


def test_scan_pptx_rejects_oversized_picture_before_decoding(
    tmp_path, monkeypatch
):
    source = _large_picture_file(tmp_path, "oversized-picture.pptx")
    counts = _count_pillow_decoding(monkeypatch)
    monkeypatch.setattr(pptx_input, "PIXEL_COUNT_CEILING", 100)

    picture = scan_pptx(source)["slides"][0]["objects"][0]

    assert picture["action"] == "preserve"
    assert picture["safety_reasons"] == ["invalid_media"]
    assert counts == {"open": 1, "verify": 0, "load": 0}


def test_scan_pptx_decodes_shared_media_once_per_scan(tmp_path, monkeypatch):
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_path = tmp_path / "shared.png"
    Image.new("RGB", (40, 20), "navy").save(image_path)
    for _ in range(50):
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            presentation.slide_width,
            presentation.slide_height,
        )
    source = tmp_path / "shared-media.pptx"
    presentation.save(source)
    counts = _count_pillow_decoding(monkeypatch)

    pictures = scan_pptx(source)["slides"][0]["objects"]

    assert len(pictures) == 50
    assert len(
        {
            picture["primary_relationship"]["target"]
            for picture in pictures
        }
    ) == 1
    assert all(picture["media_valid"] is True for picture in pictures)
    assert counts == {"open": 2, "verify": 1, "load": 1}


@pytest.mark.parametrize("mutation", ["dual_attribute", "second_blip"])
def test_scan_pptx_preserves_ambiguous_blip_sources(tmp_path, mutation):
    source = _large_picture_file(tmp_path, f"{mutation}.pptx")
    _add_ambiguous_blip_source(source, mutation)

    picture = scan_pptx(source)["slides"][0]["objects"][0]

    assert picture["action"] == "preserve"
    assert picture["safety_reasons"] == ["ambiguous_media_source"]
    assert len(picture["blip_sources"]) == 2
    assert len(picture["relationships"]) == 1


def test_scan_pptx_preserves_picture_with_extra_external_relationship(tmp_path):
    source = _large_picture_file(tmp_path, "extra-external.pptx")
    _add_picture_hyperlink(source)

    picture = scan_pptx(source)["slides"][0]["objects"][0]

    assert picture["action"] == "preserve"
    assert picture["safety_reasons"] == [
        "external_relationship",
        "ambiguous_media_source",
    ]
    assert len(picture["relationships"]) == 2
    assert {
        relationship["target_mode"]
        for relationship in picture["relationships"]
    } == {"Internal", "External"}


def test_prepare_pptx_job_writes_inventory_manifest_and_analyzed_pages(tmp_path):
    source = _large_picture_file(tmp_path, "prepare.pptx")
    before = source.read_bytes()
    run_dir = tmp_path / "run"

    prepared = pptx_input.prepare_pptx_job(
        source,
        run_dir=run_dir,
        slide_size="original",
        lang="en",
    )

    assert prepared == run_dir.resolve()
    assert source.read_bytes() == before
    store = RunStore.open(prepared)
    manifest = store.read_json("job_manifest.json")
    native_path = run_dir / "pages/page_001/native_objects.json"
    candidates_path = run_dir / "pages/page_001/screenshot_candidates.json"
    native = store.read_json("pages/page_001/native_objects.json")
    candidates = store.read_json(
        "pages/page_001/screenshot_candidates.json"
    )
    assert manifest["input"] == {
        "type": "pptx",
        "original_path": str(source.resolve()),
        "source": "input/original.pptx",
        "sha256": hashlib.sha256(before).hexdigest(),
        "slide_count": 1,
        "object_count": 1,
        "candidate_count": 1,
        "slide_width": 9144000,
        "slide_height": 4572000,
        "inventories": [
            {
                "page_id": "page_001",
                "slide_index": 1,
                "slide_part": "ppt/slides/slide1.xml",
                "slide_width": 9144000,
                "slide_height": 4572000,
                "native_objects_sha256": hashlib.sha256(
                    native_path.read_bytes()
                ).hexdigest(),
                "screenshot_candidates_sha256": hashlib.sha256(
                    candidates_path.read_bytes()
                ).hexdigest(),
            }
        ],
    }
    assert manifest["pages"] == ["page_001"]
    assert manifest["options"] == {
        "agent_provider": "host",
        "lang": "en",
        "slide_size": "original",
        "output_path": None,
        "resource_policy": safe_default_policy(),
    }
    assert native["page_id"] == "page_001"
    assert len(native["objects"]) == 1
    assert candidates["candidates"] == native["objects"]
    assert store.read_json("run_state.json")["status"] == "prepared"
    assert (
        store.read_json("page_jobs.json")["pages"]["page_001"]["status"]
        == "analyzed"
    )


@pytest.mark.parametrize("agent_provider", ["host", "local"])
def test_prepare_pptx_job_freezes_agent_provider(tmp_path, agent_provider):
    source = _large_picture_file(tmp_path, "provider.pptx")

    run = pptx_input.prepare_pptx_job(
        source,
        run_dir=tmp_path / agent_provider,
        agent_provider=agent_provider,
    )

    assert RunStore.open(run).read_json("job_manifest.json")["options"][
        "agent_provider"
    ] == agent_provider


def test_prepare_pptx_job_cleans_failed_run(tmp_path):
    source = tmp_path / "broken.pptx"
    source.write_bytes(b"not a zip")
    run_dir = tmp_path / "failed-run"

    with pytest.raises(ValueError, match="Cannot open PPTX ZIP"):
        pptx_input.prepare_pptx_job(source, run_dir=run_dir)

    assert run_dir.is_dir()
    assert list(run_dir.iterdir()) == []


def test_execute_pptx_preserve_copies_bytes_and_reports_candidates(tmp_path):
    source = _large_picture_file(tmp_path, "execute.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    state_before = store.read_json("run_state.json")
    pages_before = store.read_json("page_jobs.json")

    summary = pptx_input.execute_pptx_preserve(store)

    output = run_dir / "final" / "output.pptx"
    assert output.read_bytes() == source.read_bytes()
    token = summary.pop("_output_identity")
    assert summary == {
        "schema_version": 1,
        "status": "completed",
        "pages": 1,
        "preserved_objects": 1,
        "pending_candidates": 1,
        "warnings": [
            "P1 preserved screenshot candidates without replacement"
        ],
        "outputs": {"pptx": str(output.resolve())},
        "input_sha256": hashlib.sha256(source.read_bytes()).hexdigest(),
        "output_sha256": hashlib.sha256(output.read_bytes()).hexdigest(),
    }
    status = output.lstat()
    assert token == {
        "version": 1,
        "path": str(output.resolve()),
        "dev": status.st_dev,
        "ino": status.st_ino,
        "mode": status.st_mode,
        "size": status.st_size,
        "mtime_ns": status.st_mtime_ns,
        "sha256": summary["output_sha256"],
    }
    assert store.read_json("run_state.json") == state_before
    assert store.read_json("page_jobs.json") == pages_before


def test_execute_pptx_shadow_reports_replaced_pages(
    tmp_path,
    monkeypatch,
):
    source = _large_picture_file(tmp_path, "shadow.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    plans = [
        {
            "page_id": "page_001",
            "slide_part": "ppt/slides/slide1.xml",
            "image_path": "candidate.png",
            "work_root": "reconstruction",
            "decision": {
                "runtime_action": "shadow_run",
                "eligible_for_shadow_run": True,
                "source_shape_id": "2",
            },
        }
    ]

    def fake_run(source_path, output_path, received, *, run_root, lang):
        assert Path(source_path) == run_dir / "input/original.pptx"
        assert received == plans
        assert Path(run_root) == run_dir
        assert lang == "ch"
        Path(output_path).parent.mkdir(parents=True, exist_ok=True)
        Path(output_path).write_bytes(b"shadow")
        digest = hashlib.sha256(b"shadow").hexdigest()
        return {
            "page_results": [
                {
                    "schema_version": 1,
                    "page_id": "page_001",
                    "status": "replaced",
                }
            ],
            "output_sha256": digest,
            "_output_identity": pptx_input._published_output_identity(
                Path(output_path),
                Path(output_path).lstat(),
                digest,
            ),
        }

    monkeypatch.setattr(
        pptx_input,
        "run_shadow_replacements",
        fake_run,
        raising=False,
    )

    summary = pptx_input.execute_pptx_shadow(store, plans)

    assert summary["preserved_objects"] == 0
    assert summary["pending_candidates"] == 0
    assert summary["replaced_pages"] == 1
    assert summary["preserved_with_warning_pages"] == 0
    assert summary["page_results"][0]["status"] == "replaced"
    assert summary["output_sha256"] == hashlib.sha256(b"shadow").hexdigest()


def test_execute_pptx_preserve_honors_explicit_output_and_refuses_existing(
    tmp_path,
):
    source = _large_picture_file(tmp_path, "explicit.pptx")
    output = tmp_path / "exports" / "preserved.pptx"
    run_dir = pptx_input.prepare_pptx_job(
        source,
        run_dir=tmp_path / "run",
        output_path=output,
    )
    output.parent.mkdir()
    output.write_bytes(b"existing")

    with pytest.raises(FileExistsError, match="already exists"):
        pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    assert output.read_bytes() == b"existing"
    output.unlink()
    summary = pptx_input.execute_pptx_preserve(RunStore.open(run_dir))
    assert summary["outputs"]["pptx"] == str(output.resolve())
    assert output.read_bytes() == source.read_bytes()


def test_execute_pptx_preserve_rejects_manifest_source_as_output(tmp_path):
    source = _large_picture_file(tmp_path, "same-path.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    copied_source = run_dir / manifest["input"]["source"]
    before = copied_source.read_bytes()
    manifest["options"]["output_path"] = str(copied_source)
    store.write_json("job_manifest.json", manifest)

    with pytest.raises(ValueError, match="overwrites source"):
        pptx_input.execute_pptx_preserve(store)

    assert copied_source.read_bytes() == before


def test_execute_pptx_preserve_validates_counts_before_creating_output(tmp_path):
    source = _large_picture_file(tmp_path, "bad-count.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    manifest["input"]["candidate_count"] = -1
    store.write_json("job_manifest.json", manifest)

    with pytest.raises(RuntimeError, match="candidate_count"):
        pptx_input.execute_pptx_preserve(store)

    assert not (run_dir / "final" / "output.pptx").exists()


def test_validate_pptx_inventories_accepts_bound_inventory_documents(tmp_path):
    source = _large_picture_file(tmp_path, "validated-inventories.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")

    assert pptx_input.validate_pptx_inventories(RunStore.open(run_dir)) == (
        1,
        1,
    )


def test_execute_pptx_preserve_removes_corrupt_temporary_copy(
    tmp_path, monkeypatch
):
    source = _large_picture_file(tmp_path, "corrupt-copy.pptx")
    before = source.read_bytes()
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")

    def corrupt_copy(_source, target):
        Path(target).write_bytes(b"corrupt")
        return str(target)

    monkeypatch.setattr(pptx_input.shutil, "copyfile", corrupt_copy)

    with pytest.raises(RuntimeError, match="hash mismatch"):
        pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    final = run_dir / "final"
    assert not (final / "output.pptx").exists()
    assert not list(final.glob(".*.tmp"))
    assert source.read_bytes() == before


def test_execute_pptx_preserve_keeps_primary_error_when_tmp_cleanup_fails(
    tmp_path, monkeypatch
):
    source = _large_picture_file(tmp_path, "cleanup-error.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")

    def corrupt_copy(_source, target):
        Path(target).write_bytes(b"corrupt")
        return str(target)

    original_unlink = Path.unlink

    def fail_temporary_cleanup(path, *args, **kwargs):
        if path.name.startswith(".output.pptx.") and path.suffix == ".tmp":
            raise OSError("cleanup failed")
        return original_unlink(path, *args, **kwargs)

    monkeypatch.setattr(pptx_input.shutil, "copyfile", corrupt_copy)
    monkeypatch.setattr(Path, "unlink", fail_temporary_cleanup)

    with pytest.raises(RuntimeError, match="hash mismatch") as caught:
        pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    assert isinstance(caught.value.__cause__, OSError)
    assert str(caught.value.__cause__) == "cleanup failed"


def test_execute_pptx_preserve_uses_no_clobber_publisher(
    tmp_path, monkeypatch
):
    source = _large_picture_file(tmp_path, "publish-race.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    expected_output = run_dir / "final" / "output.pptx"

    def race(_temporary, target):
        Path(target).write_bytes(b"winner")
        raise FileExistsError(f"PPTX output already exists: {target}")

    monkeypatch.setattr(pptx_input, "_publish_pptx_no_clobber", race)

    with pytest.raises(FileExistsError, match="already exists"):
        pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    assert expected_output.read_bytes() == b"winner"


def test_execute_pptx_preserve_rejects_dangling_output_symlink(tmp_path):
    source = _large_picture_file(tmp_path, "symlink.pptx")
    output = tmp_path / "dangling-output.pptx"
    try:
        output.symlink_to(tmp_path / "missing-target.pptx")
    except OSError as error:
        pytest.skip(f"symlink unavailable: {error}")
    with pytest.raises(ValueError, match="symlink"):
        pptx_input.prepare_pptx_job(
            source,
            run_dir=tmp_path / "run",
            output_path=output,
        )

    assert not (tmp_path / "missing-target.pptx").exists()


def test_execute_pptx_preserve_default_output_cannot_escape_symlink(tmp_path):
    source = _large_picture_file(tmp_path, "symlink-parent.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    external = tmp_path / "external"
    external.mkdir()
    try:
        (run_dir / "final").symlink_to(external, target_is_directory=True)
    except OSError as error:
        pytest.skip(f"symlink unavailable: {error}")

    with pytest.raises(ValueError, match="symlink"):
        pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    assert not (external / "output.pptx").exists()


@pytest.mark.parametrize(
    "mutate",
    [
        lambda manifest: manifest.update(schema_version=2),
        lambda manifest: manifest.update(pages=[]),
        lambda manifest: manifest["input"].update(candidate_count=2),
        lambda manifest: manifest["input"].update(slide_count=10**1000),
    ],
)
def test_execute_pptx_preserve_rejects_inconsistent_manifest_before_output(
    tmp_path, mutate
):
    source = _large_picture_file(tmp_path, "bad-manifest.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    mutate(manifest)
    store.write_json("job_manifest.json", manifest)

    with pytest.raises((ValueError, RuntimeError)):
        pptx_input.execute_pptx_preserve(store)

    assert not (run_dir / "final" / "output.pptx").exists()


@pytest.mark.parametrize(
    "mutation",
    [
        "native_schema",
        "native_page_id",
        "objects_type",
        "object_type",
        "object_action",
        "candidates_type",
        "candidate_action",
        "candidate_mismatch",
        "manifest_object_count",
        "manifest_candidate_count",
    ],
)
def test_execute_pptx_preserve_rejects_corrupt_authoritative_inventory(
    tmp_path, mutation
):
    source = _large_picture_file(tmp_path, f"inventory-{mutation}.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    native_path = "pages/page_001/native_objects.json"
    candidates_path = "pages/page_001/screenshot_candidates.json"
    native = store.read_json(native_path)
    candidates = store.read_json(candidates_path)
    manifest = store.read_json("job_manifest.json")

    if mutation == "native_schema":
        native["schema_version"] = 2
    elif mutation == "native_page_id":
        native["page_id"] = "page_999"
    elif mutation == "objects_type":
        native["objects"] = {}
    elif mutation == "object_type":
        native["objects"] = ["not-an-object"]
    elif mutation == "object_action":
        native["objects"][0]["action"] = "replace"
    elif mutation == "candidates_type":
        candidates["candidates"] = {}
    elif mutation == "candidate_action":
        candidates["candidates"][0]["action"] = "preserve"
    elif mutation == "candidate_mismatch":
        candidates["candidates"] = []
    elif mutation == "manifest_object_count":
        manifest["input"]["object_count"] = 2
    else:
        manifest["input"]["candidate_count"] = 0

    store.write_json(native_path, native)
    store.write_json(candidates_path, candidates)
    store.write_json("job_manifest.json", manifest)

    with pytest.raises((ValueError, RuntimeError)):
        pptx_input.execute_pptx_preserve(store)

    assert not (run_dir / "final" / "output.pptx").exists()


def test_execute_pptx_preserve_rejects_oversized_inventory(tmp_path):
    source = _large_picture_file(tmp_path, "oversized-inventory.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    native_path = run_dir / "pages" / "page_001" / "native_objects.json"
    native_path.write_text(
        '{"padding":"' + ("x" * (16 * 1024 * 1024)) + '"}',
        encoding="utf-8",
    )

    with pytest.raises(RuntimeError, match="too large"):
        pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    assert not (run_dir / "final" / "output.pptx").exists()


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("name", "tampered-name"),
        ("type", "tampered-type"),
        ("xml_c14n_sha256", "0" * 64),
    ],
)
def test_execute_pptx_preserve_rejects_consistently_tampered_object_fields(
    tmp_path, field, value
):
    source = _large_picture_file(tmp_path, f"tampered-{field}.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    native_path = "pages/page_001/native_objects.json"
    candidates_path = "pages/page_001/screenshot_candidates.json"
    native = store.read_json(native_path)
    candidates = store.read_json(candidates_path)
    native["objects"][0][field] = value
    candidates["candidates"][0][field] = value
    store.write_json(native_path, native)
    store.write_json(candidates_path, candidates)

    with pytest.raises(RuntimeError, match="hash"):
        pptx_input.execute_pptx_preserve(store)

    assert not (run_dir / "final" / "output.pptx").exists()


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("slide_index", 99),
        ("slide_part", "ppt/slides/slide99.xml"),
        ("slide_width", 1),
        ("slide_height", 1),
    ],
)
def test_execute_pptx_preserve_rejects_consistently_tampered_inventory_metadata(
    tmp_path, field, value
):
    source = _large_picture_file(tmp_path, f"metadata-{field}.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    native_path = "pages/page_001/native_objects.json"
    candidates_path = "pages/page_001/screenshot_candidates.json"
    native = store.read_json(native_path)
    candidates = store.read_json(candidates_path)
    native[field] = value
    candidates[field] = value
    store.write_json(native_path, native)
    store.write_json(candidates_path, candidates)

    with pytest.raises(RuntimeError, match="hash"):
        pptx_input.execute_pptx_preserve(store)

    assert not (run_dir / "final" / "output.pptx").exists()


def test_execute_pptx_preserve_rejects_semantically_unchanged_candidate_file(
    tmp_path,
):
    source = _large_picture_file(tmp_path, "candidate-whitespace.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    candidates_path = (
        run_dir / "pages/page_001/screenshot_candidates.json"
    )
    candidates_path.write_bytes(candidates_path.read_bytes() + b"\n")

    with pytest.raises(RuntimeError, match="hash"):
        pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    assert not (run_dir / "final" / "output.pptx").exists()


@pytest.mark.parametrize(
    "mutate",
    [
        lambda manifest: manifest["input"].pop("inventories"),
        lambda manifest: manifest["input"].update(inventories={}),
        lambda manifest: manifest["input"]["inventories"].append(
            dict(manifest["input"]["inventories"][0])
        ),
        lambda manifest: manifest["input"]["inventories"][0].update(
            native_objects_sha256="A" * 64
        ),
        lambda manifest: manifest["input"]["inventories"][0].pop(
            "screenshot_candidates_sha256"
        ),
        lambda manifest: manifest["input"]["inventories"][0].update(
            unexpected=True
        ),
    ],
)
def test_execute_pptx_preserve_rejects_invalid_manifest_inventory_bindings(
    tmp_path, mutate
):
    source = _large_picture_file(tmp_path, "bad-bindings.pptx")
    run_dir = pptx_input.prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    if "inventories" not in manifest["input"]:
        native_path = run_dir / "pages/page_001/native_objects.json"
        candidates_path = (
            run_dir / "pages/page_001/screenshot_candidates.json"
        )
        native = store.read_json(native_path.relative_to(run_dir))
        manifest["input"]["inventories"] = [
            {
                key: native[key]
                for key in (
                    "page_id",
                    "slide_index",
                    "slide_part",
                    "slide_width",
                    "slide_height",
                )
            }
            | {
                "native_objects_sha256": hashlib.sha256(
                    native_path.read_bytes()
                ).hexdigest(),
                "screenshot_candidates_sha256": hashlib.sha256(
                    candidates_path.read_bytes()
                ).hexdigest(),
            }
        ]
    mutate(manifest)
    store.write_json("job_manifest.json", manifest)

    with pytest.raises(RuntimeError, match="inventor"):
        pptx_input.execute_pptx_preserve(store)

    assert not (run_dir / "final" / "output.pptx").exists()


def test_scan_pptx_reports_source_slide_size_and_textbox(tmp_path):
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(3), Inches(1))
    shape.name = "Title"
    shape.text = "hello"
    source = tmp_path / "sample.pptx"
    presentation.save(source)

    result = scan_pptx(source)

    assert result["schema_version"] == 1
    assert result["source"] == str(source.resolve())
    assert len(result["source_sha256"]) == 64
    assert result["slide_count"] == 1
    assert result["slide_width"] == 9144000
    assert result["slide_height"] == 4572000
    assert result["slides"][0]["objects"] == [
        {
            "shape_id": str(shape.shape_id),
            "name": "Title",
            "type": "text",
            "z_order": 0,
            "group_path": [],
            "inside_group": False,
            "slide_part": "ppt/slides/slide1.xml",
            "x": 914400,
            "y": 1828800,
            "cx": 2743200,
            "cy": 914400,
            "rotation": 0,
            "rotation_degrees": 0.0,
            "flip_h": False,
            "flip_v": False,
            "crop_left": 0,
            "crop_top": 0,
            "crop_right": 0,
            "crop_bottom": 0,
            "has_extension": False,
            "has_timing_reference": False,
            "relationships": [],
            "xml_c14n_sha256": result["slides"][0]["objects"][0]["xml_c14n_sha256"],
            "action": "preserve",
            "safety_reasons": [],
        }
    ]


def test_scan_pptx_scans_grouped_objects_relationships_notes_and_picture_data(tmp_path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    text.name = "Text"
    text.text = "text"
    plain = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(2), Inches(1))
    plain.name = "Plain"
    table = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(2), Inches(1))
    table.name = "Table"
    image_path = tmp_path / "image.png"
    Image.new("RGB", (7, 5), "red").save(image_path)
    picture = slide.shapes.add_picture(str(image_path), Inches(4), Inches(1), Inches(2), Inches(1))
    picture.name = "Picture"
    picture.crop_left = 0.1
    picture.crop_top = 0.2
    group = slide.shapes.add_group_shape()
    group.name = "Group"
    grouped = group.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1), Inches(1), Inches(1))
    grouped.name = "Grouped"
    connector = slide.shapes.add_connector(MSO_CONNECTOR_TYPE.STRAIGHT, Inches(6), Inches(1), Inches(7), Inches(2))
    connector.name = "Connector"
    slide.notes_slide.notes_text_frame.text = "note"
    source = tmp_path / "objects.pptx"
    presentation.save(source)
    _remove_text_bodies(source, {"Plain", "Grouped"})
    _add_shape_metadata(source, "Text")
    before = source.read_bytes()

    result = scan_pptx(source)

    assert source.read_bytes() == before
    objects = result["slides"][0]["objects"]
    by_name = {item["name"]: item for item in objects}
    assert [item["name"] for item in objects] == ["Text", "Plain", "Table", "Picture", "Group", "Grouped", "Connector"]
    assert [by_name[name]["type"] for name in ("Text", "Plain", "Table", "Picture", "Group", "Grouped", "Connector")] == [
        "text", "autoshape", "table", "picture", "group", "autoshape", "connector"
    ]
    assert by_name["Grouped"]["group_path"] == [by_name["Group"]["shape_id"]]
    assert by_name["Grouped"]["inside_group"] is True
    assert by_name["Grouped"]["z_order"] == 0
    assert (by_name["Text"]["rotation"], by_name["Text"]["rotation_degrees"]) == (540000, 9.0)
    assert by_name["Text"]["flip_h"] is True
    assert by_name["Text"]["flip_v"] is True
    assert by_name["Text"]["has_extension"] is True
    assert by_name["Text"]["has_timing_reference"] is True
    assert (by_name["Table"]["x"], by_name["Table"]["y"], by_name["Table"]["cx"], by_name["Table"]["cy"]) == (
        914400, 2743200, 1828800, 914400
    )
    assert by_name["Picture"]["crop_left"] == 10000
    assert by_name["Picture"]["crop_top"] == 20000
    assert by_name["Picture"]["media_sha256"] == result["parts"][by_name["Picture"]["primary_relationship"]["target"]]
    assert (by_name["Picture"]["pixel_width"], by_name["Picture"]["pixel_height"]) == (7, 5)
    assert by_name["Picture"]["primary_relationship"]["content_type"] == "image/png"
    assert result["slides"][0]["notes_part"] == "ppt/notesSlides/notesSlide1.xml"
    assert result["slides"][0]["notes_sha256"] == result["parts"]["ppt/notesSlides/notesSlide1.xml"]
    assert scan_pptx(source) == result


def test_scan_pptx_accepts_strict_namespaces_and_relationships(tmp_path):
    source = _strict_presentation_file(tmp_path)
    source_bytes = source.read_bytes()

    first = scan_pptx(source)
    second = scan_pptx(source)

    assert first == second
    assert source.read_bytes() == source_bytes
    assert first["source_sha256"] == hashlib.sha256(source_bytes).hexdigest()
    assert (first["slide_width"], first["slide_height"]) == (
        9144000,
        4572000,
    )
    assert [
        slide["slide_part"] for slide in first["slides"]
    ] == [
        "ppt/slides/slide1.xml",
        "ppt/slides/slide2.xml",
    ]
    assert [
        item["name"] for item in first["slides"][0]["objects"]
    ] == ["Backdrop", "Title"]
    assert [
        item["type"] for item in first["slides"][0]["objects"]
    ] == ["picture", "text"]
    picture = first["slides"][0]["objects"][0]
    assert picture["action"] == "candidate"
    assert picture["primary_relationship"]["type"] == (
        "http://purl.oclc.org/ooxml/officeDocument/relationships/image"
    )
    assert picture["primary_relationship"]["target"] == "ppt/media/image1.png"
    assert picture["media_sha256"] == first["parts"]["ppt/media/image1.png"]
    assert [
        item["name"] for item in first["slides"][1]["objects"]
    ] == ["Second", "Shape"]
    assert [
        item["type"] for item in first["slides"][1]["objects"]
    ] == ["text", "autoshape"]
    assert first["slides"][0]["notes_part"] == (
        "ppt/notesSlides/notesSlide1.xml"
    )
    assert first["slides"][0]["notes_sha256"] == first["parts"][
        "ppt/notesSlides/notesSlide1.xml"
    ]
    with zipfile.ZipFile(source) as archive:
        slide_bytes = archive.read("ppt/slides/slide1.xml")
    assert first["parts"]["ppt/slides/slide1.xml"] == hashlib.sha256(
        slide_bytes
    ).hexdigest()
    assert [
        item["xml_c14n_sha256"]
        for slide in first["slides"]
        for item in slide["objects"]
    ] == [
        item["xml_c14n_sha256"]
        for slide in second["slides"]
        for item in slide["objects"]
    ]


def test_prepare_and_execute_preserve_strict_pptx_bytes(tmp_path):
    source = _strict_presentation_file(tmp_path)
    source_bytes = source.read_bytes()
    run_dir = pptx_input.prepare_pptx_job(
        source,
        run_dir=tmp_path / "strict-run",
    )

    summary = pptx_input.execute_pptx_preserve(RunStore.open(run_dir))

    output = run_dir / "final" / "output.pptx"
    assert output.read_bytes() == source_bytes
    assert source.read_bytes() == source_bytes
    assert summary["pages"] == 2
    assert summary["pending_candidates"] == 1
    assert summary["input_sha256"] == hashlib.sha256(source_bytes).hexdigest()
    assert summary["output_sha256"] == summary["input_sha256"]


def test_scan_pptx_selects_mc_choice_requiring_strict_p_and_a(tmp_path):
    source = _presentation_file(tmp_path)
    _replace_shape_tree_with_alternate_content(
        source,
        choice_requires="p a",
    )
    _strictify_pptx(source)

    objects = scan_pptx(source)["slides"][0]["objects"]

    assert [item["name"] for item in objects] == ["Choice"]


def test_scan_pptx_does_not_accept_attacker_namespace_suffix(tmp_path):
    source = _presentation_file(tmp_path)
    _replace_member(
        source,
        "ppt/presentation.xml",
        lambda xml: _rename_xml_namespaces(
            xml,
            {
                "http://schemas.openxmlformats.org/presentationml/2006/main": (
                    "https://attacker.invalid/ooxml/presentationml/main"
                )
            },
        ),
    )

    with pytest.raises(ValueError, match="Missing slide size"):
        scan_pptx(source)


def test_scan_pptx_uses_blip_relationship_as_picture_primary(tmp_path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_path = tmp_path / "image.png"
    Image.new("RGB", (11, 13), "blue").save(image_path)
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(1))
    source = tmp_path / "hyperlinked-picture.pptx"
    presentation.save(source)
    _add_picture_hyperlink(source)

    picture = scan_pptx(source)["slides"][0]["objects"][0]

    assert picture["relationships"][0]["target_mode"] == "External"
    assert picture["primary_relationship"]["target"] == "ppt/media/image1.png"
    assert (picture["pixel_width"], picture["pixel_height"]) == (11, 13)


@pytest.mark.parametrize(
    ("name", "contents", "message"),
    [
        ("not-pptx.txt", b"x", "existing .pptx"),
        ("bad.pptx", b"not a zip", "Cannot open PPTX ZIP"),
    ],
)
def test_scan_pptx_rejects_invalid_source(tmp_path, name, contents, message):
    source = tmp_path / name
    source.write_bytes(contents)

    with pytest.raises(ValueError, match=message):
        scan_pptx(source)


def test_scan_pptx_rejects_duplicate_zip_member(tmp_path):
    source = _presentation_file(tmp_path)
    with pytest.warns(UserWarning, match="Duplicate name"):
        with zipfile.ZipFile(source, "a") as archive:
            archive.writestr("ppt/presentation.xml", archive.read("ppt/presentation.xml"))

    with pytest.raises(ValueError, match="Duplicate ZIP member: ppt/presentation.xml"):
        scan_pptx(source)


def test_scan_pptx_accepts_empty_canonical_directory_members(tmp_path):
    source = _presentation_file(tmp_path)
    with zipfile.ZipFile(source, "a") as archive:
        archive.writestr("_rels/", b"")
        archive.writestr("ppt/", b"")
        archive.writestr("ppt/slides/", b"")

    inventory = scan_pptx(source)

    assert inventory["slide_count"] == 1
    assert "_rels/" not in inventory["parts"]
    assert "ppt/" not in inventory["parts"]
    assert "ppt/slides/" not in inventory["parts"]
    assert "ppt/presentation.xml" in inventory["parts"]


@pytest.mark.parametrize(
    "name",
    [
        "",
        "/absolute.xml",
        "C:/absolute.xml",
        "C:relative.xml",
        "C:",
        "../outside.xml",
        "ppt/../outside.xml",
        "ppt/./slide.xml",
        "ppt//slide.xml",
        r"ppt\slide.xml",
    ],
)
def test_normalized_member_name_rejects_unsafe_or_noncanonical_names(name):
    with pytest.raises(ValueError, match="Unsafe ZIP member name"):
        pptx_input._normalized_member_name(name)


def test_scan_pptx_rejects_nul_in_original_zip_member_name(tmp_path):
    source = _presentation_file(tmp_path)
    safe_name = b"nulXhidden.xml"
    unsafe_name = b"nul\x00hidden.xml"
    with zipfile.ZipFile(source, "a") as archive:
        archive.writestr(safe_name.decode(), b"unused")
    contents = source.read_bytes()
    assert contents.count(safe_name) == 2
    source.write_bytes(contents.replace(safe_name, unsafe_name))
    with zipfile.ZipFile(source) as archive:
        info = archive.getinfo("nul")
        assert info.filename == "nul"
        assert info.orig_filename == unsafe_name.decode()

    with pytest.raises(ValueError, match="Unsafe ZIP member name"):
        scan_pptx(source)


@pytest.mark.parametrize(
    ("contents", "compression"),
    [
        (b"x", zipfile.ZIP_STORED),
        (b"", zipfile.ZIP_DEFLATED),
    ],
)
def test_scan_pptx_rejects_directory_members_with_payload(
    tmp_path, contents, compression
):
    source = _presentation_file(tmp_path)
    with zipfile.ZipFile(source, "a") as archive:
        info = zipfile.ZipInfo("payload/")
        info.compress_type = compression
        archive.writestr(info, contents)

    with pytest.raises(ValueError, match="Directory ZIP member has payload"):
        scan_pptx(source)


def test_scan_pptx_rejects_directory_file_identity_conflict(tmp_path):
    source = _presentation_file(tmp_path)
    with zipfile.ZipFile(source, "a") as archive:
        archive.writestr("ambiguous", b"")
        archive.writestr("ambiguous/", b"")

    with pytest.raises(ValueError, match="Duplicate normalized ZIP member"):
        scan_pptx(source)


def test_scan_pptx_rejects_encrypted_zip_member(tmp_path):
    source = _presentation_file(tmp_path)
    contents = bytearray(source.read_bytes())
    central_header = contents.index(b"PK\x01\x02")
    contents[central_header + 8] |= 1
    source.write_bytes(contents)

    with pytest.raises(ValueError, match="Encrypted ZIP member"):
        scan_pptx(source)


def test_scan_pptx_rejects_encrypted_directory_member(tmp_path):
    source = _presentation_file(tmp_path)
    with zipfile.ZipFile(source, "a") as archive:
        archive.writestr("encrypted/", b"")
    contents = bytearray(source.read_bytes())
    central_header = contents.rfind(b"PK\x01\x02")
    contents[central_header + 8] |= 1
    source.write_bytes(contents)
    with zipfile.ZipFile(source) as archive:
        info = archive.getinfo("encrypted/")
        assert info.is_dir()
        assert info.flag_bits & 1

    with pytest.raises(ValueError, match="Encrypted ZIP member"):
        scan_pptx(source)


def test_scan_pptx_wraps_invalid_utf8_central_member_name(tmp_path):
    source = _presentation_file(tmp_path)
    name = b"centralXinvalid.xml"
    with zipfile.ZipFile(source, "a") as archive:
        archive.writestr(name.decode(), b"unused")
    contents = bytearray(source.read_bytes())
    central_header = contents.rfind(b"PK\x01\x02")
    assert contents[central_header + 46 : central_header + 46 + len(name)] == name
    contents[central_header + 9] |= 0x08
    contents[central_header + 46 + name.index(b"X")] = 0xFF
    source.write_bytes(contents)

    with pytest.raises(ValueError, match="invalid member name encoding") as caught:
        scan_pptx(source)

    assert type(caught.value) is ValueError
    assert isinstance(caught.value.__cause__, UnicodeDecodeError)


def test_scan_pptx_wraps_invalid_utf8_local_member_name(tmp_path):
    source = _presentation_file(tmp_path)
    name = b"localXinvalid.xml"
    with zipfile.ZipFile(source, "a") as archive:
        archive.writestr(name.decode(), b"unused")
    contents = bytearray(source.read_bytes())
    local_header = contents.rfind(b"PK\x03\x04")
    assert contents[local_header + 30 : local_header + 30 + len(name)] == name
    contents[local_header + 7] |= 0x08
    contents[local_header + 30 + name.index(b"X")] = 0xFF
    source.write_bytes(contents)

    with pytest.raises(ValueError, match="invalid member name encoding") as caught:
        scan_pptx(source)

    assert type(caught.value) is ValueError
    assert isinstance(caught.value.__cause__, UnicodeDecodeError)


def test_scan_pptx_translates_corrupt_member_to_value_error(tmp_path):
    source = _presentation_file(tmp_path)
    contents = bytearray(source.read_bytes())
    central_header = contents.index(b"PK\x01\x02")
    contents[central_header + 16] ^= 1
    source.write_bytes(contents)

    with pytest.raises(ValueError, match="Cannot read PPTX ZIP"):
        scan_pptx(source)


def test_scan_pptx_rejects_high_compression_ratio_member_before_read(tmp_path):
    source = _presentation_file(tmp_path)
    with zipfile.ZipFile(source, "a", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("ppt/media/bomb.bin", b"0" * (1024 * 1024))

    with pytest.raises(ValueError, match="compression ratio"):
        scan_pptx(source)


def test_scan_pptx_hashes_source_without_path_read_bytes(tmp_path, monkeypatch):
    source = _presentation_file(tmp_path)
    expected = hashlib.sha256(source.read_bytes()).hexdigest()

    def fail_read_bytes(_):
        raise AssertionError("source must be hashed as a stream")

    monkeypatch.setattr(Path, "read_bytes", fail_read_bytes)

    assert scan_pptx(source)["source_sha256"] == expected


def test_scan_pptx_accepts_absolute_internal_relationship_target(tmp_path):
    source = _presentation_file(tmp_path)
    _replace_member(
        source,
        "ppt/_rels/presentation.xml.rels",
        lambda xml: re.sub(
            br'Target="slides/slide1.xml"',
            b'Target="/ppt/slides/slide1.xml"',
            xml,
        ),
    )

    result = scan_pptx(source)

    assert result["slides"][0]["slide_part"] == "ppt/slides/slide1.xml"


def test_scan_pptx_rejects_wrong_presentation_slide_relationship_type(tmp_path):
    source = _presentation_file(tmp_path)
    _set_presentation_slide_relationship_type(
        source,
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
    )

    with pytest.raises(ValueError, match="slide relationship type"):
        scan_pptx(source)


def test_scan_pptx_rejects_attacker_uri_ending_in_slide(tmp_path):
    source = _presentation_file(tmp_path)
    _set_presentation_slide_relationship_type(
        source,
        "https://attacker.invalid/slide",
    )

    with pytest.raises(ValueError, match="slide relationship type"):
        scan_pptx(source)


def test_scan_pptx_rejects_wrong_slide_content_type(tmp_path):
    source = _presentation_file(tmp_path)
    _set_content_type(
        source,
        "/ppt/slides/slide1.xml",
        "application/xml",
    )

    with pytest.raises(ValueError, match="content type.*ppt/slides/slide1.xml"):
        scan_pptx(source)


def test_scan_pptx_rejects_duplicate_content_type_default(tmp_path):
    source = _presentation_file(tmp_path)

    def duplicate_default(xml):
        root = ET.fromstring(xml)
        ct = "{http://schemas.openxmlformats.org/package/2006/content-types}"
        ET.SubElement(
            root,
            f"{ct}Default",
            Extension="XML",
            ContentType="application/xml",
        )
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "[Content_Types].xml", duplicate_default)

    with pytest.raises(ValueError, match="Duplicate content type Default"):
        scan_pptx(source)


def test_scan_pptx_rejects_duplicate_relationship_id(tmp_path):
    source = _presentation_file(tmp_path)

    def duplicate_id(xml):
        root = ET.fromstring(xml)
        relationship = next(iter(root))
        root.append(ET.fromstring(ET.tostring(relationship)))
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(
        source,
        "ppt/_rels/presentation.xml.rels",
        duplicate_id,
    )

    with pytest.raises(ValueError, match="Duplicate relationship Id"):
        scan_pptx(source)


def test_scan_pptx_rejects_wrong_notes_content_type(tmp_path):
    source = _notes_presentation_file(tmp_path)
    _set_content_type(
        source,
        "/ppt/notesSlides/notesSlide1.xml",
        "application/xml",
    )

    with pytest.raises(ValueError, match="notes content type"):
        scan_pptx(source)


def test_scan_pptx_rejects_external_notes_relationship(tmp_path):
    source = _notes_presentation_file(tmp_path)
    _mutate_notes_relationship(
        source,
        target_mode="External",
        target="https://attacker.invalid/notes.xml",
    )

    with pytest.raises(ValueError, match="notes relationship TargetMode"):
        scan_pptx(source)


def test_scan_pptx_rejects_attacker_notes_relationship_type(tmp_path):
    source = _notes_presentation_file(tmp_path)
    _mutate_notes_relationship(
        source,
        relationship_type="https://attacker.invalid/notesSlide",
    )

    with pytest.raises(ValueError, match="notes relationship type"):
        scan_pptx(source)


def test_scan_pptx_canonical_hash_ignores_registered_namespace_prefix(tmp_path):
    source = _presentation_file(tmp_path)
    presentation_namespace = (
        "http://schemas.openxmlformats.org/presentationml/2006/main"
    )
    ET.register_namespace("p", presentation_namespace)
    first = scan_pptx(source)["slides"][0]["objects"][0][
        "xml_c14n_sha256"
    ]
    try:
        ET.register_namespace("foo", presentation_namespace)
        second = scan_pptx(source)["slides"][0]["objects"][0][
            "xml_c14n_sha256"
        ]
    finally:
        ET.register_namespace("p", presentation_namespace)

    assert second == first


def test_scan_pptx_expands_alternate_content_without_extlst_object(tmp_path):
    source = _presentation_file(tmp_path)
    _replace_shape_tree_with_alternate_content(source)

    objects = scan_pptx(source)["slides"][0]["objects"]

    assert [(item["name"], item["type"], item["z_order"]) for item in objects] == [
        ("Wrapped", "autoshape", 0)
    ]
    assert all(item["shape_id"] for item in objects)


def test_scan_pptx_records_complete_unknown_object_but_skips_metadata(tmp_path):
    source = _presentation_file(tmp_path)
    _replace_shape_tree_with_unknown_object(source)

    objects = scan_pptx(source)["slides"][0]["objects"]

    assert [(item["shape_id"], item["name"], item["type"]) for item in objects] == [
        ("77", "Mystery", "unknown")
    ]


def test_scan_pptx_prefers_first_supported_alternate_content_choice(tmp_path):
    source = _presentation_file(tmp_path)
    ET.register_namespace(
        "p",
        "http://schemas.openxmlformats.org/presentationml/2006/main",
    )
    _replace_shape_tree_with_alternate_content(
        source,
        choice_requires="p",
    )

    objects = scan_pptx(source)["slides"][0]["objects"]

    assert [item["name"] for item in objects] == ["Choice"]


def test_scan_pptx_resolves_renamed_requires_prefix_to_supported_uri(tmp_path):
    source = _presentation_file(tmp_path)
    presentation_namespace = (
        "http://schemas.openxmlformats.org/presentationml/2006/main"
    )
    ET.register_namespace("foo", presentation_namespace)
    try:
        _replace_shape_tree_with_alternate_content(
            source,
            choice_requires="foo",
        )
    finally:
        ET.register_namespace("p", presentation_namespace)

    objects = scan_pptx(source)["slides"][0]["objects"]

    assert [item["name"] for item in objects] == ["Choice"]


def test_scan_pptx_falls_back_for_requires_bound_to_unknown_uri(tmp_path):
    source = _presentation_file(tmp_path)
    ET.register_namespace("future", "urn:example:future")
    _replace_shape_tree_with_alternate_content(
        source,
        choice_requires="future",
        choice_namespace_uri="urn:example:future",
    )

    objects = scan_pptx(source)["slides"][0]["objects"]

    assert [item["name"] for item in objects] == ["Wrapped"]


@pytest.mark.parametrize(
    "contents",
    [
        b"<!DOCTYPE notes><notes/>",
        b"<notes>",
    ],
)
def test_scan_pptx_rejects_unsafe_notes_xml(tmp_path, contents):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "note"
    source = tmp_path / "unsafe-notes.pptx"
    presentation.save(source)
    _replace_member(
        source,
        "ppt/notesSlides/notesSlide1.xml",
        lambda _: contents,
    )

    with pytest.raises(ValueError, match="ppt/notesSlides/notesSlide1.xml"):
        scan_pptx(source)


def test_scan_pptx_rejects_utf16_doctype_before_parsing(tmp_path):
    source = _presentation_file(tmp_path)

    def inject_utf16_doctype(xml):
        text = xml.decode("utf-8").replace("UTF-8", "UTF-16")
        declaration_end = text.index("?>") + 2
        text = text[:declaration_end] + "<!DOCTYPE p:sld>" + text[declaration_end:]
        return text.encode("utf-16")

    _replace_member(
        source,
        "ppt/slides/slide1.xml",
        inject_utf16_doctype,
    )

    with pytest.raises(
        ValueError,
        match="Unsafe XML in ppt/slides/slide1.xml",
    ):
        scan_pptx(source)


@pytest.mark.parametrize(
    ("part", "replace", "message"),
    [
        ("ppt/slides/slide1.xml", lambda xml: b"<!DOCTYPE x>" + xml, "Unsafe XML in ppt/slides/slide1.xml"),
        ("ppt/_rels/presentation.xml.rels", lambda xml: re.sub(br'Target="slides/slide1.xml"', b'Target="../../escape.xml"', xml), "escapes package"),
    ],
)
def test_scan_pptx_rejects_unsafe_xml_and_relationship_targets(tmp_path, part, replace, message):
    source = _presentation_file(tmp_path)
    _replace_member(source, part, replace)

    with pytest.raises(ValueError, match=message):
        scan_pptx(source)


@pytest.mark.parametrize(
    ("uri", "expected"),
    [
        ("http://schemas.openxmlformats.org/drawingml/2006/chart", "chart"),
        ("http://schemas.openxmlformats.org/drawingml/2006/diagram", "smartart"),
        ("urn:example:unknown", "unknown"),
    ],
)
def test_scan_pptx_classifies_graphic_frame_uri(tmp_path, uri, expected):
    source = _presentation_file(tmp_path)
    _append_graphic_frame(source, uri)

    object_ = scan_pptx(source)["slides"][0]["objects"][-1]

    assert object_["name"] == "Injected"
    assert object_["type"] == expected
    assert (object_["x"], object_["y"], object_["cx"], object_["cy"]) == (1, 2, 3, 4)


def _presentation_file(tmp_path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(1), Inches(1)).text = "x"
    source = tmp_path / "sample.pptx"
    presentation.save(source)
    return source


def _slide_background_picture_file(tmp_path, name):
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_path = tmp_path / f"{Path(name).stem}.png"
    Image.new("RGB", (40, 20), "navy").save(image_path)
    slide.shapes.add_picture(str(image_path), 0, 0, Inches(10), Inches(5))
    overlay = slide.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    overlay.name = "Overlay"
    overlay.text = "overlay"
    source = tmp_path / name
    presentation.save(source)
    return source


def _move_picture_to_slide_background(
    source,
    *,
    kind="picture",
    extension=False,
    crop_left=0,
):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    def transform(xml):
        root = ET.fromstring(xml)
        common_slide_data = root.find(f"{p}cSld")
        shape_tree = common_slide_data.find(f"{p}spTree")
        picture = shape_tree.find(f"{p}pic")
        shape_tree.remove(picture)
        background = ET.Element(f"{p}bg")
        if kind == "theme":
            ET.SubElement(background, f"{p}bgRef", idx="1001")
        else:
            properties = ET.SubElement(background, f"{p}bgPr")
            if kind == "picture":
                fill = picture.find(f"{p}blipFill")
                fill.tag = f"{a}blipFill"
                if crop_left:
                    ET.SubElement(fill, f"{a}srcRect", l=str(crop_left))
                properties.append(fill)
            elif kind == "solid":
                ET.SubElement(properties, f"{a}solidFill")
            elif kind == "gradient":
                ET.SubElement(properties, f"{a}gradFill")
            if extension:
                ET.SubElement(properties, f"{p}extLst")
        common_slide_data.insert(0, background)
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", transform)


def _large_picture_file(tmp_path, name):
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_path = tmp_path / f"{Path(name).stem}.png"
    Image.new("RGB", (40, 20), "navy").save(image_path)
    slide.shapes.add_picture(
        str(image_path),
        Inches(0.5),
        Inches(0.25),
        Inches(9),
        Inches(4.5),
    )
    source = tmp_path / name
    presentation.save(source)
    return source


def _strict_presentation_file(tmp_path):
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    first = presentation.slides.add_slide(presentation.slide_layouts[6])
    image_path = tmp_path / "strict.png"
    Image.new("RGB", (40, 20), "navy").save(image_path)
    picture = first.shapes.add_picture(
        str(image_path),
        Inches(0.5),
        Inches(0.25),
        Inches(9),
        Inches(4.5),
    )
    picture.name = "Backdrop"
    title = first.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    title.name = "Title"
    title.text = "strict"
    first.notes_slide.notes_text_frame.text = "strict note"
    second = presentation.slides.add_slide(presentation.slide_layouts[6])
    second_text = second.shapes.add_textbox(
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(1),
    )
    second_text.name = "Second"
    second_text.text = "second"
    shape = second.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4),
        Inches(1),
        Inches(1),
        Inches(1),
    )
    shape.name = "Shape"
    source = tmp_path / "strict.pptx"
    presentation.save(source)
    _remove_text_bodies(
        source,
        {"Shape"},
        "ppt/slides/slide2.xml",
    )
    _strictify_pptx(source)
    return source


def _strictify_pptx(source):
    namespace_mapping = {
        "http://schemas.openxmlformats.org/presentationml/2006/main": (
            "http://purl.oclc.org/ooxml/presentationml/main"
        ),
        "http://schemas.openxmlformats.org/drawingml/2006/main": (
            "http://purl.oclc.org/ooxml/drawingml/main"
        ),
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships": (
            "http://purl.oclc.org/ooxml/officeDocument/relationships"
        ),
    }
    relationship_types = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slide": (
            "http://purl.oclc.org/ooxml/officeDocument/relationships/slide"
        ),
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide": (
            "http://purl.oclc.org/ooxml/officeDocument/relationships/notesSlide"
        ),
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image": (
            "http://purl.oclc.org/ooxml/officeDocument/relationships/image"
        ),
    }
    strict_parts = {
        "ppt/presentation.xml",
        "ppt/slides/slide1.xml",
        "ppt/slides/slide2.xml",
        "ppt/notesSlides/notesSlide1.xml",
    }
    temporary = source.with_suffix(".strict.tmp")
    package_relationship = (
        "{http://schemas.openxmlformats.org/package/2006/relationships}"
        "Relationship"
    )
    with zipfile.ZipFile(source) as old, zipfile.ZipFile(temporary, "w") as new:
        for info in old.infolist():
            contents = old.read(info.filename)
            if info.filename in strict_parts:
                contents = _rename_xml_namespaces(
                    contents,
                    namespace_mapping,
                )
            elif info.filename.endswith(".rels"):
                root = ET.fromstring(contents)
                for relationship in root.findall(package_relationship):
                    relationship_type = relationship.get("Type")
                    if relationship_type in relationship_types:
                        relationship.set(
                            "Type",
                            relationship_types[relationship_type],
                        )
                contents = ET.tostring(
                    root,
                    encoding="utf-8",
                    xml_declaration=True,
                )
            new.writestr(info, contents)
    temporary.replace(source)


def _rename_xml_namespaces(xml, namespace_mapping):
    root = ET.fromstring(xml)
    for node in root.iter():
        if node.tag.startswith("{"):
            namespace, separator, local = node.tag[1:].partition("}")
            if separator and namespace in namespace_mapping:
                node.tag = f"{{{namespace_mapping[namespace]}}}{local}"
        attributes = {}
        for name, value in node.attrib.items():
            if name.startswith("{"):
                namespace, separator, local = name[1:].partition("}")
                if separator and namespace in namespace_mapping:
                    name = f"{{{namespace_mapping[namespace]}}}{local}"
            attributes[name] = value
        node.attrib.clear()
        node.attrib.update(attributes)
    strict = (
        "http://purl.oclc.org/ooxml/presentationml/main",
        "http://purl.oclc.org/ooxml/drawingml/main",
        "http://purl.oclc.org/ooxml/officeDocument/relationships",
    )
    for prefix, namespace in zip(("p", "a", "r"), strict):
        ET.register_namespace(prefix, namespace)
    try:
        return ET.tostring(
            root,
            encoding="utf-8",
            xml_declaration=True,
        )
    finally:
        transitional = (
            "http://schemas.openxmlformats.org/presentationml/2006/main",
            "http://schemas.openxmlformats.org/drawingml/2006/main",
            "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
        )
        for prefix, namespace in zip(("p", "a", "r"), transitional):
            ET.register_namespace(prefix, namespace)


def _count_pillow_decoding(monkeypatch):
    original_open = pptx_input.Image.open
    counts = {"open": 0, "verify": 0, "load": 0}

    class CountedImage:
        def __init__(self, image):
            self.image = image

        def __enter__(self):
            return self

        def __exit__(self, *_):
            self.image.close()

        def __getattr__(self, name):
            return getattr(self.image, name)

        def verify(self):
            counts["verify"] += 1
            return self.image.verify()

        def load(self):
            counts["load"] += 1
            return self.image.load()

    def counted_open(*args, **kwargs):
        counts["open"] += 1
        return CountedImage(original_open(*args, **kwargs))

    monkeypatch.setattr(pptx_input.Image, "open", counted_open)
    return counts


def _notes_presentation_file(tmp_path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.notes_slide.notes_text_frame.text = "note"
    source = tmp_path / "notes.pptx"
    presentation.save(source)
    return source


def _replace_member(source, name, transform):
    temporary = source.with_suffix(".tmp")
    with zipfile.ZipFile(source) as old, zipfile.ZipFile(temporary, "w") as new:
        for info in old.infolist():
            contents = old.read(info.filename)
            new.writestr(info, transform(contents) if info.filename == name else contents)
    temporary.replace(source)


def _remove_member(source, name):
    temporary = source.with_suffix(".tmp")
    with zipfile.ZipFile(source) as old, zipfile.ZipFile(temporary, "w") as new:
        for info in old.infolist():
            if info.filename != name:
                new.writestr(info, old.read(info.filename))
    temporary.replace(source)


def _mutate_picture_safety(
    source,
    name,
    *,
    flip_h=False,
    extension=False,
    timing=False,
    remove_transform=False,
    huge_extent=False,
):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    def transform(xml):
        root = ET.fromstring(xml)
        for picture in root.findall(f".//{p}pic"):
            non_visual = picture.find(f"{p}nvPicPr/{p}cNvPr")
            if non_visual is None or non_visual.get("name") != name:
                continue
            properties = picture.find(f"{p}spPr")
            xfrm = properties.find(f"{a}xfrm")
            if flip_h:
                xfrm.set("flipH", "1")
            if extension:
                ET.SubElement(properties, f"{a}extLst")
            if timing:
                timing_node = root.find(f"{p}timing")
                if timing_node is None:
                    timing_node = ET.SubElement(root, f"{p}timing")
                ET.SubElement(
                    timing_node,
                    f"{p}spTgt",
                    spid=non_visual.get("id"),
                )
            if remove_transform:
                properties.remove(xfrm)
            if huge_extent:
                xfrm.find(f"{a}ext").set("cx", str(10**1000))
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", transform)


def _make_picture_relationship_external(source):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    pr = "{http://schemas.openxmlformats.org/package/2006/relationships}"

    def change_blip(xml):
        root = ET.fromstring(xml)
        blip = root.find(f".//{p}pic/{p}blipFill/{a}blip")
        rel_id = blip.attrib.pop(f"{r}embed")
        blip.set(f"{r}link", rel_id)
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    def change_relationship(xml):
        root = ET.fromstring(xml)
        for relationship in root.findall(f"{pr}Relationship"):
            if relationship.get("Type", "").endswith("/image"):
                relationship.set("Target", "https://example.invalid/image.png")
                relationship.set("TargetMode", "External")
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", change_blip)
    _replace_member(
        source,
        "ppt/slides/_rels/slide1.xml.rels",
        change_relationship,
    )


def _add_ambiguous_blip_source(source, mutation):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"

    def transform(xml):
        root = ET.fromstring(xml)
        blip_fill = root.find(f".//{p}pic/{p}blipFill")
        blip = blip_fill.find(f"{a}blip")
        rel_id = blip.get(f"{r}embed")
        if mutation == "dual_attribute":
            blip.set(f"{r}link", rel_id)
        else:
            ET.SubElement(blip_fill, f"{a}blip", {f"{r}embed": rel_id})
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", transform)


def _remove_text_bodies(
    source,
    names,
    slide_part="ppt/slides/slide1.xml",
):
    namespace = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    def transform(xml):
        root = ET.fromstring(xml)
        for shape in root.findall(f".//{namespace}sp"):
            non_visual = shape.find(f".//{namespace}cNvPr")
            if non_visual is not None and non_visual.get("name") in names:
                shape.remove(shape.find(f"{namespace}txBody"))
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, slide_part, transform)


def _append_graphic_frame(source, uri):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    def transform(xml):
        root = ET.fromstring(xml)
        tree = root.find(f".//{p}spTree")
        frame = ET.SubElement(tree, f"{p}graphicFrame")
        non_visual = ET.SubElement(frame, f"{p}nvGraphicFramePr")
        ET.SubElement(non_visual, f"{p}cNvPr", id="99", name="Injected")
        ET.SubElement(non_visual, f"{p}cNvGraphicFramePr")
        ET.SubElement(non_visual, f"{p}nvPr")
        xfrm = ET.SubElement(frame, f"{p}xfrm")
        ET.SubElement(xfrm, f"{a}off", x="1", y="2")
        ET.SubElement(xfrm, f"{a}ext", cx="3", cy="4")
        graphic = ET.SubElement(frame, f"{a}graphic")
        graphic_data = ET.SubElement(graphic, f"{a}graphicData", uri=uri)
        preview = ET.SubElement(graphic_data, f"{a}xfrm")
        ET.SubElement(preview, f"{a}off", x="101", y="102")
        ET.SubElement(preview, f"{a}ext", cx="103", cy="104")
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", transform)


def _add_shape_metadata(source, name):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"

    def transform(xml):
        root = ET.fromstring(xml)
        for shape in root.findall(f".//{p}sp"):
            non_visual = shape.find(f".//{p}cNvPr")
            if non_visual is not None and non_visual.get("name") == name:
                xfrm = shape.find(f".//{a}xfrm")
                xfrm.set("rot", "540000")
                xfrm.set("flipH", "1")
                xfrm.set("flipV", "true")
                ET.SubElement(shape.find(f"{p}spPr"), f"{a}extLst")
                timing = ET.SubElement(root, f"{p}timing")
                ET.SubElement(timing, f"{p}spTgt", spid=non_visual.get("id"))
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", transform)


def _add_picture_hyperlink(source):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    r = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}"
    pr = "{http://schemas.openxmlformats.org/package/2006/relationships}"

    def add_click(xml):
        root = ET.fromstring(xml)
        non_visual = root.find(f".//{p}pic/{p}nvPicPr/{p}cNvPr")
        ET.SubElement(non_visual, f"{a}hlinkClick", {f"{r}id": "rIdHyperlink"})
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    def add_relationship(xml):
        root = ET.fromstring(xml)
        ET.SubElement(
            root,
            f"{pr}Relationship",
            Id="rIdHyperlink",
            Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink",
            Target="https://example.com/",
            TargetMode="External",
        )
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", add_click)
    _replace_member(source, "ppt/slides/_rels/slide1.xml.rels", add_relationship)


def _set_presentation_slide_relationship_type(source, relationship_type):
    pr = "{http://schemas.openxmlformats.org/package/2006/relationships}"

    def transform(xml):
        root = ET.fromstring(xml)
        for relationship in root.findall(f"{pr}Relationship"):
            if relationship.get("Target") in {
                "slides/slide1.xml",
                "/ppt/slides/slide1.xml",
            }:
                relationship.set("Type", relationship_type)
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/_rels/presentation.xml.rels", transform)


def _mutate_notes_relationship(
    source,
    *,
    relationship_type=None,
    target_mode=None,
    target=None,
):
    pr = "{http://schemas.openxmlformats.org/package/2006/relationships}"

    def transform(xml):
        root = ET.fromstring(xml)
        for relationship in root.findall(f"{pr}Relationship"):
            if relationship.get("Type", "").endswith("/notesSlide"):
                if relationship_type is not None:
                    relationship.set("Type", relationship_type)
                if target_mode is not None:
                    relationship.set("TargetMode", target_mode)
                if target is not None:
                    relationship.set("Target", target)
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/_rels/slide1.xml.rels", transform)


def _set_content_type(source, part_name, content_type):
    ct = "{http://schemas.openxmlformats.org/package/2006/content-types}"

    def transform(xml):
        root = ET.fromstring(xml)
        changed = False
        for override in root.findall(f"{ct}Override"):
            if override.get("PartName") == part_name:
                override.set("ContentType", content_type)
                changed = True
        if not changed:
            extension = Path(part_name).suffix.lstrip(".")
            for default in root.findall(f"{ct}Default"):
                if default.get("Extension", "").casefold() == extension.casefold():
                    default.set("ContentType", content_type)
                    break
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "[Content_Types].xml", transform)


def _replace_shape_tree_with_alternate_content(
    source,
    *,
    choice_requires="future",
    choice_namespace_uri=None,
):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"
    a = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    mc = "{http://schemas.openxmlformats.org/markup-compatibility/2006}"

    def add_shape(parent, shape_id, name):
        shape = ET.SubElement(parent, f"{p}sp")
        non_visual = ET.SubElement(shape, f"{p}nvSpPr")
        ET.SubElement(non_visual, f"{p}cNvPr", id=shape_id, name=name)
        ET.SubElement(non_visual, f"{p}cNvSpPr")
        ET.SubElement(non_visual, f"{p}nvPr")
        properties = ET.SubElement(shape, f"{p}spPr")
        transform = ET.SubElement(properties, f"{a}xfrm")
        ET.SubElement(transform, f"{a}off", x="1", y="2")
        ET.SubElement(transform, f"{a}ext", cx="3", cy="4")

    def transform(xml):
        root = ET.fromstring(xml)
        tree = root.find(f".//{p}spTree")
        for child in list(tree):
            if child.tag not in {f"{p}nvGrpSpPr", f"{p}grpSpPr"}:
                tree.remove(child)
        ET.SubElement(tree, f"{p}extLst")
        alternate = ET.SubElement(tree, f"{mc}AlternateContent")
        choice = ET.SubElement(
            alternate,
            f"{mc}Choice",
            Requires=choice_requires,
        )
        if choice_namespace_uri is not None:
            choice.set(f"{{{choice_namespace_uri}}}marker", "1")
        add_shape(choice, "98", "Choice")
        fallback = ET.SubElement(alternate, f"{mc}Fallback")
        add_shape(fallback, "99", "Wrapped")
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", transform)


def _replace_shape_tree_with_unknown_object(source):
    p = "{http://schemas.openxmlformats.org/presentationml/2006/main}"

    def transform(xml):
        root = ET.fromstring(xml)
        tree = root.find(f".//{p}spTree")
        for child in list(tree):
            if child.tag not in {f"{p}nvGrpSpPr", f"{p}grpSpPr"}:
                tree.remove(child)
        ET.SubElement(tree, f"{p}extLst")
        unknown = ET.SubElement(tree, f"{p}mystery")
        non_visual = ET.SubElement(unknown, f"{p}nvMysteryPr")
        ET.SubElement(non_visual, f"{p}cNvPr", id="77", name="Mystery")
        ET.SubElement(non_visual, f"{p}cNvMysteryPr")
        ET.SubElement(non_visual, f"{p}nvPr")
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", transform)
