from __future__ import annotations

import errno
import hashlib
import json
from numbers import Integral
import os
import stat
import subprocess
import sys
import types
import weakref
from pathlib import Path

import pytest
from PIL import Image

import image_to_ppt
from scripts import (
    bg_model,
    fg_extract,
    lama_inpaint,
    sam_worker,
    text_detect,
    visual_worker,
    visual_segment,
)


def test_visual_worker_rejects_source_hash_mismatch_before_loading_pipeline(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (8, 6), "red").save(source)
    mask = tmp_path / "mask.png"
    Image.new("L", (8, 6), 0).save(mask)
    mask_content = mask.read_bytes()
    text_clean = tmp_path / "text-clean.png"
    Image.new("RGB", (8, 6), "white").save(text_clean)
    text_clean_content = text_clean.read_bytes()
    request = tmp_path / "request.json"
    request_content = json.dumps({
        "text_analysis": {
            "items": [{"text": "x", "box": [1, 1, 2, 2]}],
            "mask_path": str(mask),
            "text_clean_path": str(text_clean),
        },
        "text_mask_sha256": hashlib.sha256(mask_content).hexdigest(),
        "text_mask_size": len(mask_content),
        "text_clean_sha256": hashlib.sha256(text_clean_content).hexdigest(),
        "text_clean_size": len(text_clean_content),
    }).encode("utf-8")
    request.write_bytes(request_content)
    monkeypatch.setattr(
        visual_worker,
        "_load_process_image",
        lambda: pytest.fail("visual pipeline must not load before source validation"),
    )
    monkeypatch.setattr(sys, "argv", [
        "visual_worker.py", "--image", str(source), "--work-dir", str(tmp_path),
        "--lang", "en", "--request", str(request),
        "--request-sha256", hashlib.sha256(request_content).hexdigest(),
        "--request-size", str(len(request_content)),
        "--source-sha256", "0" * 64, "--source-size", str(source.stat().st_size),
        "--result", str(tmp_path / "result.json"),
    ])

    with pytest.raises(ValueError, match="sha256 mismatch"):
        visual_worker.main()


def test_visual_worker_processes_verified_in_memory_source_snapshot(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (8, 6), "red").save(source)
    mask = tmp_path / "mask.png"
    Image.new("L", (8, 6), 0).save(mask)
    mask_content = mask.read_bytes()
    text_clean = tmp_path / "text-clean.png"
    Image.new("RGB", (8, 6), "white").save(text_clean)
    text_clean_content = text_clean.read_bytes()
    request = tmp_path / "request.json"
    request_content = json.dumps({
        "text_analysis": {
            "items": [{"text": "x", "box": [1, 1, 2, 2]}],
            "mask_path": str(mask),
            "text_clean_path": str(text_clean),
        },
        "text_mask_sha256": hashlib.sha256(mask_content).hexdigest(),
        "text_mask_size": len(mask_content),
        "text_clean_sha256": hashlib.sha256(text_clean_content).hexdigest(),
        "text_clean_size": len(text_clean_content),
    }).encode("utf-8")
    request.write_bytes(request_content)
    source_content = source.read_bytes()

    def fake_process(path, work_dir, *args, **kwargs):
        Image.new("RGB", (8, 6), "black").save(path)
        snapshot = kwargs["_source_image"]
        assert snapshot.shape == (6, 8, 3)
        assert tuple(snapshot[0, 0]) == (255, 0, 0)
        assert kwargs["_text_mask"].shape == (6, 8)
        assert tuple(kwargs["_text_clean_image"][0, 0]) == (255, 255, 255)
        return {"snapshot": "verified"}

    monkeypatch.setattr(visual_worker, "_load_process_image", lambda: fake_process)
    monkeypatch.setattr(sys, "argv", [
        "visual_worker.py", "--image", str(source), "--work-dir", str(tmp_path),
        "--lang", "en", "--request", str(request),
        "--request-sha256", hashlib.sha256(request_content).hexdigest(),
        "--request-size", str(len(request_content)),
        "--source-sha256", hashlib.sha256(source_content).hexdigest(),
        "--source-size", str(len(source_content)),
        "--result", str(tmp_path / "result.json"),
    ])

    assert visual_worker.main() == 0
    assert json.loads((tmp_path / "result.json").read_text(encoding="utf-8")) == {
        "snapshot": "verified"
    }


def test_isolated_visual_request_binds_source_snapshot(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (8, 6), "red").save(source)
    Image.new("L", (8, 6), 0).save(tmp_path / "mask.png")
    captured = {}

    def fake_worker(command, **kwargs):
        request_path = Path(command[command.index("--request") + 1])
        captured["request"] = json.loads(request_path.read_text(encoding="utf-8"))
        captured["source_sha256"] = command[command.index("--source-sha256") + 1]
        captured["source_size"] = command[command.index("--source-size") + 1]
        captured["request_sha256"] = command[command.index("--request-sha256") + 1]
        return types.SimpleNamespace(returncode=1, stderr="stop", stdout="")

    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_worker)

    with pytest.raises(RuntimeError, match="stop"):
        image_to_ppt._process_image_isolated(
            source,
            tmp_path,
            "en",
            {"items": [], "mask_path": str(tmp_path / "mask.png")},
        )

    source_content = source.read_bytes()
    assert captured["source_sha256"] == hashlib.sha256(source_content).hexdigest()
    assert captured["source_size"] == str(len(source_content))
    assert captured["request_sha256"] == hashlib.sha256(
        (tmp_path / "visual-worker-request.json").read_bytes()
    ).hexdigest()
    assert set(captured["request"]) == {
        "text_analysis", "text_mask_sha256", "text_mask_size",
        "text_clean_sha256", "text_clean_size",
    }


def test_visual_worker_does_not_trust_replaced_request_source_binding(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (8, 6), "red").save(source)
    original = source.read_bytes()
    request = tmp_path / "request.json"
    mask = tmp_path / "mask.png"
    Image.new("L", (8, 6), 0).save(mask)
    mask_content = mask.read_bytes()
    request_payload = {
        "text_analysis": {"items": [], "mask_path": str(mask)},
        "text_mask_sha256": hashlib.sha256(mask_content).hexdigest(),
        "text_mask_size": len(mask_content),
        "text_clean_sha256": None,
        "text_clean_size": None,
        "source_sha256": hashlib.sha256(original).hexdigest(),
        "source_size": len(original),
    }
    request_content = json.dumps(request_payload).encode("utf-8")
    request.write_bytes(request_content)
    Image.new("RGB", (8, 6), "black").save(source)
    request_payload["source_sha256"] = hashlib.sha256(source.read_bytes()).hexdigest()
    request_payload["source_size"] = source.stat().st_size
    request.write_bytes(json.dumps(request_payload).encode("utf-8"))
    monkeypatch.setattr(
        visual_worker,
        "_load_process_image",
        lambda: pytest.fail("replaced source must fail before pipeline load"),
    )
    monkeypatch.setattr(sys, "argv", [
        "visual_worker.py", "--image", str(source), "--work-dir", str(tmp_path),
        "--lang", "en", "--request", str(request),
        "--request-sha256", hashlib.sha256(request_content).hexdigest(),
        "--request-size", str(len(request_content)),
        "--source-sha256", hashlib.sha256(original).hexdigest(),
        "--source-size", str(len(original)),
        "--result", str(tmp_path / "result.json"),
    ])

    with pytest.raises(ValueError, match="mismatch"):
        visual_worker.main()


def test_visual_worker_rejects_replaced_bound_text_mask_before_pipeline_load(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (8, 6), "red").save(source)
    mask = tmp_path / "mask.png"
    Image.new("L", (8, 6), 0).save(mask)
    source_content = source.read_bytes()
    mask_content = mask.read_bytes()
    request_content = json.dumps({
        "text_analysis": {"items": [], "mask_path": str(mask)},
        "text_mask_sha256": hashlib.sha256(mask_content).hexdigest(),
        "text_mask_size": len(mask_content),
        "text_clean_sha256": None,
        "text_clean_size": None,
    }).encode("utf-8")
    request = tmp_path / "request.json"
    request.write_bytes(request_content)
    Image.new("L", (8, 6), 255).save(mask)
    monkeypatch.setattr(
        visual_worker,
        "_load_process_image",
        lambda: pytest.fail("replaced text mask must fail before pipeline load"),
    )
    monkeypatch.setattr(sys, "argv", [
        "visual_worker.py", "--image", str(source), "--work-dir", str(tmp_path),
        "--lang", "en", "--request", str(request),
        "--request-sha256", hashlib.sha256(request_content).hexdigest(),
        "--request-size", str(len(request_content)),
        "--source-sha256", hashlib.sha256(source_content).hexdigest(),
        "--source-size", str(len(source_content)),
        "--result", str(tmp_path / "result.json"),
    ])

    with pytest.raises(ValueError, match="text mask.*mismatch"):
        visual_worker.main()

from image_to_ppt import _parse_reference_option
from image_to_ppt import _merge_foreground_masks
from image_to_ppt import _load_rgb
from image_to_psd import _resolve_output_paths
from scripts.bg_model import _build_inpaint_mask
from scripts.bg_model import _original_based_background
from scripts.bg_model import _build_component_repair_mask
from scripts.bg_model import _fill_text_regions
from scripts.bg_model import _should_use_fg_hint
from scripts.bg_model import build_background
from scripts.fg_extract import extract_foreground_mask
from scripts.fg_extract import _build_text_ink_mask
from scripts.fg_extract import _keep_detector_mask
from scripts.fg_extract import _limit_combined_mask
from scripts.fg_extract import split_components
from scripts.ppt_assemble import assemble_pptx
from scripts.ppt_assemble import assemble_pptx_multi
from scripts.ppt_assemble import _set_run_font
from scripts.psd_assemble import assemble_psd
from scripts.psd_assemble import AsposePsdLicenseError
from scripts.text_detect import (
    _adjust_font_size,
    _filter_noise,
    _select_font,
    _should_force_regular_weight,
    _try_tesseract,
)
from scripts.visual_compare_qa import write_visual_compare
from pptx import Presentation


def test_multi_slide_original_uses_one_custom_canvas_for_matching_ratios(
    tmp_path: Path,
) -> None:
    backgrounds = []
    originals = []
    for index, size in enumerate(((400, 200), (800, 400)), start=1):
        background = tmp_path / f"background_{index}.png"
        original = tmp_path / f"original_{index}.png"
        Image.new("RGB", size, "white").save(background)
        Image.new("RGB", size, "black").save(original)
        backgrounds.append(background)
        originals.append(original)

    out_path = tmp_path / "original.pptx"
    assemble_pptx_multi(
        [
            {
                "background_original_path": str(background),
                "components": [],
                "text_items": [],
                "img_width": width,
                "img_height": height,
                "original_image_path": str(original),
            }
            for background, original, (width, height) in zip(
                backgrounds, originals, ((400, 200), (800, 400))
            )
        ],
        out_path,
        slide_size="original",
    )

    presentation = Presentation(out_path)
    assert len(presentation.slides) == 2
    assert abs(presentation.slide_width / presentation.slide_height - 2.0) < 1e-6
    assert presentation.element.sldSz.get("type") == "custom"


def test_multi_slide_original_contains_near_ratio_pages_without_stretching(
    tmp_path: Path,
) -> None:
    component_path = tmp_path / "component.png"
    Image.new("RGBA", (10, 10), "red").save(component_path)
    slides_data = []
    for index, (width, height) in enumerate(((10000, 5000), (9999, 5000)), start=1):
        background = tmp_path / f"background_{index}.png"
        original = tmp_path / f"original_{index}.png"
        Image.new("RGB", (20, 10), "white").save(background)
        Image.new("RGB", (20, 10), "black").save(original)
        slides_data.append(
            {
                "background_original_path": str(background),
                "components": [
                    {
                        "path": str(component_path),
                        "x": 1000,
                        "y": 1000,
                        "w": 1000,
                        "h": 1000,
                    }
                ],
                "text_items": [
                    {
                        "box": [3000, 1000, 1000, 1000],
                        "text": "Square",
                        "font_size": 20,
                    }
                ],
                "img_width": width,
                "img_height": height,
                "original_image_path": str(original),
            }
        )

    output = tmp_path / "near-ratio.pptx"
    assemble_pptx_multi(
        slides_data,
        output,
        add_reference=True,
        slide_size="original",
    )

    presentation = Presentation(output)
    content_slide = presentation.slides[2]
    reference_slide = presentation.slides[3]
    background = content_slide.shapes[0]
    component = content_slide.shapes[1]
    text = content_slide.shapes[2]
    reference_background = reference_slide.shapes[0]
    reference_image = reference_slide.shapes[1]
    page_ratio = 9999 / 5000

    assert len(presentation.slides) == 4
    assert background.left > 0
    assert abs(background.width / background.height - page_ratio) < 1e-5
    assert abs(component.width / component.height - 1.0) < 1e-5
    assert abs(text.width / text.height - 1.0) < 1e-5
    assert reference_background.left == background.left
    assert reference_background.width == background.width
    assert abs(reference_image.width / reference_image.height - page_ratio) < 1e-5


def test_batch_original_uses_pdf_physical_ratio_for_quantized_renders(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    background = tmp_path / "background.png"
    component = tmp_path / "component.png"
    Image.new("RGB", (20, 10), "white").save(background)
    Image.new("RGBA", (10, 10), "red").save(component)
    slides_data = [
        {
            "original_image_path": str(tmp_path / f"page-{index}.png"),
            "background_original_path": str(background),
            "background_path": str(background),
            "components": [
                {"path": str(component), "x": 30, "y": 30, "w": 30, "h": 30}
            ],
            "text_items": [],
            "img_width": width,
            "img_height": height,
        }
        for index, (width, height) in enumerate(((300, 150), (305, 153)), start=1)
    ]
    monkeypatch.setattr(
        image_to_ppt,
        "_prepare_multiple_images",
        lambda image_paths, lang, _work_root=None: slides_data,
    )

    result = image_to_ppt.convert_batch_variants(
        ["page-1.png", "page-2.png"],
        output_path=tmp_path / "physical.pptx",
        include_widescreen=False,
        combine_original=True,
        original_aspect_ratio=2.0,
    )

    presentation = Presentation(result["original"])
    second_component = presentation.slides[1].shapes[1]
    assert len(presentation.slides) == 2
    assert abs(presentation.slide_width / presentation.slide_height - 2.0) < 1e-6
    assert abs(second_component.width / second_component.height - 1.0) < 1e-5


def test_multi_slide_original_physical_ratio_rejects_clear_pixel_mismatch(
    tmp_path: Path,
) -> None:
    output = tmp_path / "mismatch.pptx"

    with pytest.raises(ValueError, match="same aspect ratio"):
        assemble_pptx_multi(
            [
                {
                    "background_original_path": "first.png",
                    "components": [],
                    "text_items": [],
                    "img_width": 300,
                    "img_height": 150,
                },
                {
                    "background_original_path": "second.png",
                    "components": [],
                    "text_items": [],
                    "img_width": 300,
                    "img_height": 160,
                },
            ],
            output,
            slide_size="original",
            original_aspect_ratio=2.0,
        )

    assert not output.exists()


def test_multi_slide_original_rejects_mixed_ratios_before_writing(
    tmp_path: Path,
) -> None:
    output = tmp_path / "mixed.pptx"

    with pytest.raises(ValueError, match="same aspect ratio"):
        assemble_pptx_multi(
            [
                {
                    "background_original_path": "first.png",
                    "components": [],
                    "text_items": [],
                    "img_width": 400,
                    "img_height": 200,
                },
                {
                    "background_original_path": "second.png",
                    "components": [],
                    "text_items": [],
                    "img_width": 400,
                    "img_height": 300,
                },
            ],
            output,
            slide_size="original",
        )

    assert not output.exists()


def test_batch_variants_combines_original_slides_without_repreparing(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    prepared = 0
    calls = []
    slides_data = [
        {
            "original_image_path": str(tmp_path / "first.png"),
            "background_original_path": "first-original.png",
            "background_path": "first-wide.png",
            "components": [],
            "text_items": [],
            "img_width": 400,
            "img_height": 200,
        },
        {
            "original_image_path": str(tmp_path / "second.png"),
            "background_original_path": "second-original.png",
            "background_path": "second-wide.png",
            "components": [],
            "text_items": [],
            "img_width": 800,
            "img_height": 400,
        },
    ]

    def fake_prepare(image_paths, lang, _work_root=None):
        nonlocal prepared
        prepared += 1
        return slides_data

    def fake_assemble(**kwargs):
        calls.append(kwargs)
        return str(kwargs["output_path"])

    monkeypatch.setattr(image_to_ppt, "_prepare_multiple_images", fake_prepare)
    monkeypatch.setattr(image_to_ppt, "assemble_pptx_multi", fake_assemble)

    result = image_to_ppt.convert_batch_variants(
        ["first.png", "second.png"],
        output_path=tmp_path / "deck.pptx",
        combine_original=True,
    )

    assert prepared == 1
    assert result == {
        "16:9": str((tmp_path / "deck_16x9.pptx").resolve()),
        "original": str((tmp_path / "deck_original.pptx").resolve()),
    }
    assert [call.get("slide_size", "16:9") for call in calls] == ["16:9", "original"]


def test_batch_variants_keeps_original_per_image_outputs_by_default(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    slides_data = [
        {
            "original_image_path": str(tmp_path / "first.png"),
            "background_original_path": "first-original.png",
            "background_path": "first-wide.png",
            "components": [],
            "text_items": [],
            "img_width": 400,
            "img_height": 200,
        }
    ]

    monkeypatch.setattr(
        image_to_ppt,
        "_prepare_multiple_images",
        lambda image_paths, lang, _work_root=None: slides_data,
    )
    monkeypatch.setattr(
        image_to_ppt, "assemble_pptx_multi", lambda **kwargs: str(kwargs["output_path"])
    )
    monkeypatch.setattr(
        image_to_ppt, "_assemble_prepared_slide", lambda data, output, ref, size: str(output)
    )

    result = image_to_ppt.convert_batch_variants(
        ["first.png"], output_path=tmp_path / "deck.pptx", include_widescreen=False
    )

    assert result == {
        "16:9": None,
        "original": [
            str((tmp_path / "deck_original" / "first_original.pptx").resolve())
        ],
    }


def test_batch_variants_keeps_legacy_sixth_positional_bg_period(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    slide_data = {
        "original_image_path": str(tmp_path / "first.png"),
        "background_original_path": "first-original.png",
        "background_path": "first-wide.png",
        "components": [],
        "text_items": [],
        "img_width": 400,
        "img_height": 200,
    }
    multi_calls = []

    monkeypatch.setattr(
        image_to_ppt,
        "_prepare_multiple_images",
        lambda image_paths, lang, _work_root=None: [slide_data],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "assemble_pptx_multi",
        lambda **kwargs: multi_calls.append(kwargs) or str(kwargs["output_path"]),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_assemble_prepared_slide",
        lambda data, output, ref, size: str(output),
    )

    result = image_to_ppt.convert_batch_variants(
        ["first.png"],
        tmp_path / "deck.pptx",
        "en",
        False,
        False,
        64,
        21.0,
        30,
    )

    assert multi_calls == []
    assert result["original"] == [
        str((tmp_path / "deck_original" / "first_original.pptx").resolve())
    ]


def test_image_to_psd_skill_cli_help_starts() -> None:
    script = (
        Path(__file__).resolve().parents[1]
        / "skills"
        / "image-to-psd"
        / "scripts"
        / "image_to_psd.py"
    )
    result = subprocess.run(
        [sys.executable, str(script), "--help"],
        capture_output=True,
        text=True,
        check=False,
    )

    assert result.returncode == 0, result.stderr


def test_image_to_psd_skill_uses_shared_text_reconstruction() -> None:
    scripts_dir = (
        Path(__file__).resolve().parents[1]
        / "skills"
        / "image-to-psd"
        / "scripts"
    )
    launcher = (scripts_dir / "image_to_psd.py").read_text(encoding="utf-8")

    assert "image2editable.cli" in launcher
    assert not (scripts_dir / "text_detect.py").exists()


def test_filter_noise_keeps_meaningful_all_caps_text() -> None:
    boxes = [
        {"text": "PROJECT ROADMAP", "box": (0, 0, 100, 20), "confidence": 0.95},
        {"text": "MCOULE ST:SETMP", "box": (0, 30, 100, 20), "confidence": 0.95},
    ]

    filtered = _filter_noise(boxes)

    assert [item["text"] for item in filtered] == ["PROJECT ROADMAP"]


def test_filter_noise_removes_vertical_decorative_ocr_fragments() -> None:
    boxes = [
        {"text": "目录", "box": (65, 49, 245, 134), "confidence": 0.95},
        {"text": "01", "box": (353, 254, 59, 45), "confidence": 0.95},
        {"text": "革", "box": (1236, 402, 167, 178), "confidence": 0.95},
        {"text": "命", "box": (1239, 547, 226, 180), "confidence": 0.95},
        {"text": "19111", "box": (1219, 445, 23, 115), "confidence": 0.95},
        {"text": "WUCCHANG", "box": (1202, 522, 14, 99), "confidence": 0.95},
        {"text": "武昌", "box": (1208, 640, 27, 52), "confidence": 0.95},
    ]

    filtered = _filter_noise(boxes)

    assert [item["text"] for item in filtered] == ["目录", "01"]


def test_tesseract_fallback_uses_requested_language(
    tmp_path: Path, monkeypatch
) -> None:
    image_path = tmp_path / "sample.png"
    Image.new("RGB", (60, 20), "white").save(image_path)

    captured = {}

    fake_tesseract = types.SimpleNamespace(tesseract_cmd="")

    def image_to_data(img, lang, output_type):
        captured["lang"] = lang
        return {
            "text": ["Hello"],
            "conf": ["95"],
            "left": [1],
            "top": [2],
            "width": [30],
            "height": [10],
            "block_num": [1],
            "par_num": [1],
            "line_num": [1],
        }

    fake_module = types.SimpleNamespace(
        Output=types.SimpleNamespace(DICT="DICT"),
        pytesseract=fake_tesseract,
        image_to_data=image_to_data,
    )
    monkeypatch.setitem(sys.modules, "pytesseract", fake_module)

    boxes = _try_tesseract(image_path, 0.7, lang="eng")

    assert captured["lang"] == "eng"
    assert boxes and boxes[0]["text"] == "Hello"


def test_paddleocr_session_is_reused_per_language(monkeypatch) -> None:
    created = []

    class FakeOCR:
        def predict(self, image_path):
            return []

    def fake_create(lang):
        created.append(lang)
        return FakeOCR()

    text_detect.close_ocr_engines()
    monkeypatch.setattr(text_detect, "_create_paddleocr", fake_create)
    try:
        assert text_detect._try_paddleocr(Path("one.png"), "ch", 0.7) == []
        assert text_detect._try_paddleocr(Path("two.png"), "ch", 0.7) == []
        assert text_detect._try_paddleocr(Path("three.png"), "en", 0.7) == []
        assert created == ["ch", "en"]

        text_detect.close_ocr_engines()
        assert text_detect._try_paddleocr(Path("four.png"), "ch", 0.7) == []
        assert created == ["ch", "en", "ch"]
    finally:
        text_detect.close_ocr_engines()


@pytest.mark.parametrize("failure", ["construction", "predict", "conversion"])
def test_paddleocr_falls_back_on_engine_errors(monkeypatch, failure) -> None:
    class FakeOCR:
        def predict(self, image_path):
            if failure == "predict":
                raise RuntimeError("predict failed")
            return [{
                "rec_texts": ["Hello"],
                "rec_scores": [0.95],
                "dt_polys": [],
            }]

    def fake_create(lang):
        if failure == "construction":
            raise RuntimeError("construction failed")
        return FakeOCR()

    text_detect.close_ocr_engines()
    monkeypatch.setattr(text_detect, "_create_paddleocr", fake_create)
    try:
        assert text_detect._try_paddleocr(Path("sample.png"), "en", 0.7) is None
    finally:
        text_detect.close_ocr_engines()


def _create_fake_sam_generator(
    monkeypatch,
    cuda_available,
    resource_safe=False,
):
    captured = {}
    fake_torch = types.SimpleNamespace(
        cuda=types.SimpleNamespace(is_available=lambda: cuda_available)
    )
    fake_build_sam = types.SimpleNamespace(
        build_sam2=lambda *args, **kwargs: object()
    )

    class FakeGenerator:
        def __init__(self, model, **kwargs):
            captured.update(kwargs)
            captured["generator"] = self
            self.min_mask_region_area = kwargs["min_mask_region_area"]

    fake_modules = {
        "torch": fake_torch,
        "sam2.build_sam": fake_build_sam,
        "sam2.automatic_mask_generator": types.SimpleNamespace(
            SAM2AutomaticMaskGenerator=FakeGenerator
        ),
    }
    monkeypatch.setattr(
        visual_segment.importlib,
        "import_module",
        lambda name: fake_modules[name],
    )
    if resource_safe:
        monkeypatch.setattr(
            visual_segment,
            "_build_resource_safe_sam_model",
            lambda *args: object(),
        )
    visual_segment.create_sam_generator(
        Path("sam.pt"),
        resource_safe=resource_safe,
    )
    return captured


def test_sam_generator_uses_four_points_per_batch_on_cuda(monkeypatch) -> None:
    captured = _create_fake_sam_generator(monkeypatch, cuda_available=True)

    assert captured["points_per_side"] == 16
    assert captured["points_per_batch"] == 4


def test_sam_generator_uses_four_points_per_batch_on_cpu(monkeypatch) -> None:
    captured = _create_fake_sam_generator(monkeypatch, cuda_available=False)

    assert captured["points_per_side"] == 16
    assert captured["points_per_batch"] == 4


def test_sam_generator_uses_opencv_region_cleanup_without_cuda_extension(
    monkeypatch,
) -> None:
    captured = _create_fake_sam_generator(monkeypatch, cuda_available=False)

    assert captured["min_mask_region_area"] == 0
    assert captured["generator"].min_mask_region_area == 20


def test_sam_resource_safe_generator_uses_single_point_and_rle(
    monkeypatch,
) -> None:
    captured = _create_fake_sam_generator(
        monkeypatch,
        cuda_available=True,
        resource_safe=True,
    )

    assert captured["points_per_side"] == 16
    assert captured["points_per_batch"] == 1
    assert captured["output_mode"] == "uncompressed_rle"


def test_resource_safe_sam_build_assigns_mmap_checkpoint_on_target_device(
    monkeypatch,
) -> None:
    events = []

    class EmptyWeights:
        def __enter__(self):
            events.append("empty-enter")

        def __exit__(self, *args):
            events.append("empty-exit")

    class Model:
        def load_state_dict(self, state, assign=False):
            events.append(("assign", state, assign))
            return [], []

        def eval(self):
            events.append("eval")
            return self

    model = Model()
    build_sam = types.SimpleNamespace(
        compose=lambda **kwargs: (
            events.append(("compose", kwargs)) or {"model": "config"}
        ),
        OmegaConf=types.SimpleNamespace(
            resolve=lambda config: events.append(("resolve", config))
        ),
        instantiate=lambda config, **kwargs: (
            events.append(("instantiate", config, kwargs)) or model
        ),
    )
    torch = types.SimpleNamespace(
        load=lambda path, **kwargs: (
            events.append(("load", path, kwargs)) or {"model": "weights"}
        )
    )
    monkeypatch.setattr(
        visual_segment.importlib,
        "import_module",
        lambda name: types.SimpleNamespace(
            init_empty_weights=lambda: EmptyWeights()
        ),
    )

    actual = visual_segment._build_resource_safe_sam_model(
        build_sam,
        torch,
        Path("sam.pt"),
        "cuda",
    )

    assert actual is model
    assert ("load", Path("sam.pt"), {
        "map_location": "cuda",
        "weights_only": True,
        "mmap": True,
    }) in events
    assert ("assign", "weights", True) in events
    assert events[-1] == "eval"


def test_isolated_lama_invokes_worker(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image_path = tmp_path / "source.png"
    mask_path = tmp_path / "mask.png"
    output_path = tmp_path / "clean.png"
    Image.new("RGB", (20, 10), "white").save(image_path)
    Image.new("L", (20, 10), 0).save(mask_path)
    calls = []

    def fake_run(command, **kwargs):
        calls.append((command, kwargs))
        Image.open(image_path).save(output_path)
        return types.SimpleNamespace(returncode=0, stdout="", stderr="")

    monkeypatch.setattr(
        lama_inpaint,
        "run_isolated_worker",
        fake_run,
        raising=False,
    )
    monkeypatch.setattr(
        subprocess,
        "run",
        lambda *args, **kwargs: pytest.fail("LaMa must use the shared runner"),
        raising=False,
    )

    lama_inpaint.inpaint_large_mask_isolated(
        image_path,
        mask_path,
        output_path,
    )

    assert output_path.is_file()
    assert len(calls) == 1
    assert calls[0][0][0] == sys.executable
    assert Path(calls[0][0][1]).name == "lama_worker.py"
    assert calls[0][1] == {"capture_output": True, "text": True}


def test_sam_cuda_inference_uses_bfloat16_autocast(monkeypatch) -> None:
    import numpy as np

    events = []

    class FakeContext:
        def __init__(self, name):
            self.name = name

        def __enter__(self):
            events.append(f"enter:{self.name}")

        def __exit__(self, exc_type, exc, traceback):
            events.append(f"exit:{self.name}")

    fake_torch = types.SimpleNamespace(
        bfloat16="bfloat16",
        inference_mode=lambda: FakeContext("inference"),
        autocast=lambda **kwargs: FakeContext(
            f"autocast:{kwargs['device_type']}:{kwargs['dtype']}"
        ),
    )
    monkeypatch.setattr(
        visual_segment.importlib,
        "import_module",
        lambda name: fake_torch,
    )

    class FakeGenerator:
        _image2editable_device = "cuda"

        def generate(self, image):
            events.append("generate")
            return []

    visual_segment.generate_mask_candidates(
        np.zeros((10, 10, 3), dtype=np.uint8),
        FakeGenerator(),
        include_geometry=False,
    )

    assert events == [
        "enter:inference",
        "enter:autocast:cuda:bfloat16",
        "generate",
        "exit:autocast:cuda:bfloat16",
        "exit:inference",
    ]


def _candidate_batch_signature(candidate) -> tuple:
    return (
        candidate.mask.tolist(),
        candidate.score,
        candidate.source,
        candidate.crop_box,
        candidate.touches_crop_edge,
        candidate.label,
        candidate.role,
        candidate.object_box,
    )


def _candidate_batch_record(shape=(6, 8)) -> dict:
    mask = image_to_ppt.np.zeros(shape, dtype=bool)
    mask[1:-1, 2:-2] = True
    return {
        "mask": image_to_ppt.base64.b64encode(
            image_to_ppt.np.packbits(mask, axis=None)
        ).decode("ascii"),
        "mask_shape": list(shape),
        "score": 0.94,
        "source": "sam",
        "crop_box": [0, 0, shape[1], shape[0]],
        "touches_crop_edge": False,
        "label": "",
        "role": "",
        "object_box": None,
    }


def _candidate_batch_payload(
    prompted: list[dict] | None = None,
    automatic: list[dict] | None = None,
) -> dict:
    return {
        "schema_version": 1,
        "operations": [
            {
                "id": "prompted",
                "kind": "prompted",
                "candidates": prompted or [],
            },
            {
                "id": "automatic",
                "kind": "automatic",
                "candidates": automatic or [],
            },
        ],
    }


def _candidate_batch_worker_fixture(root: Path) -> dict:
    root.mkdir(parents=True, exist_ok=True)
    image_path = root / "image.png"
    text_mask_path = root / "text-mask.png"
    proposals_path = root / "proposals.json"
    request_path = root / "request.json"
    Image.new("RGB", (8, 6), "white").save(image_path)
    Image.new("L", (8, 6), 0).save(text_mask_path)
    proposals = [{
        "box_xyxy": [1.0, 1.0, 7.0, 5.0],
        "score": 0.91,
        "label": "badge",
        "role": "object",
        "source": "full",
        "crop_box": [0, 0, 8, 6],
        "touches_crop_edge": False,
    }]
    proposals_path.write_text(json.dumps(proposals), encoding="utf-8")
    request = {
        "schema_version": 1,
        "operations": [
            {
                "id": "prompted",
                "kind": "prompted",
                "image": image_path.name,
                "text_mask": text_mask_path.name,
                "proposals": proposals_path.name,
            },
            {
                "id": "automatic",
                "kind": "automatic",
                "image": image_path.name,
            },
        ],
    }
    request_path.write_text(json.dumps(request), encoding="utf-8")
    return {
        "root": root,
        "request": request_path,
        "image": image_path,
        "text_mask": text_mask_path,
        "proposals": proposals_path,
        "proposal_records": proposals,
    }


def _run_candidate_batch_worker_main(
    monkeypatch,
    request_path: Path,
    result_path: Path,
) -> None:
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "sam_worker.py",
            "--mode",
            "batch",
            "--request",
            str(request_path),
            "--result",
            str(result_path),
        ],
    )
    sam_worker.main()


def test_candidate_batch_worker_loads_sam_once_and_preserves_candidate_semantics(
    tmp_path: Path,
    monkeypatch,
) -> None:
    batch_helper = getattr(
        image_to_ppt,
        "_generate_sam_candidate_batch_isolated",
        None,
    )
    assert batch_helper is not None
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    image[:, :, 1] = 37
    text_mask = image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8)
    text_mask[0, 0] = 255
    proposal = image_to_ppt.ObjectProposal(
        box_xyxy=(1.0, 1.0, 7.0, 5.0),
        score=0.91,
        label="badge",
        role="object",
        source="full",
        crop_box=(0, 0, 8, 6),
        touches_crop_edge=True,
    )
    prompted_mask = image_to_ppt.np.zeros((6, 8), dtype=bool)
    prompted_mask[1:5, 1:7] = True
    automatic_mask = image_to_ppt.np.zeros((6, 8), dtype=bool)
    automatic_mask[2:4, 2:6] = True
    generator = object()
    events = []
    calls = []

    def create_generator(checkpoint, *, resource_safe):
        events.append(("create", checkpoint, resource_safe))
        return generator

    def generate_prompted(actual_image, proposals, actual_generator, actual_text_mask):
        events.append("prompted")
        assert actual_generator is generator
        assert image_to_ppt.np.array_equal(actual_image, image)
        assert image_to_ppt.np.array_equal(actual_text_mask, text_mask)
        assert proposals == [proposal]
        return [
            image_to_ppt.MaskCandidate(
                mask=prompted_mask.copy(),
                score=0.89,
                source="grounded:full:object",
                crop_box=(0, 0, 8, 6),
                touches_crop_edge=True,
                label="badge",
                role="object",
                object_box=(1.0, 1.0, 7.0, 5.0),
            )
        ]

    def generate_automatic(actual_image, actual_generator, **kwargs):
        events.append(("automatic", kwargs))
        assert actual_generator is generator
        assert image_to_ppt.np.array_equal(actual_image, image)
        return [
            image_to_ppt.MaskCandidate(
                mask=automatic_mask.copy(),
                score=0.97,
                source="sam",
                crop_box=(0, 0, 8, 6),
                touches_crop_edge=False,
                label="shape",
                role="decoration",
                object_box=None,
            )
        ]

    monkeypatch.setattr(
        sam_worker,
        "_load_tools",
        lambda: (
            image_to_ppt.ObjectProposal,
            create_generator,
            generate_automatic,
            generate_prompted,
            lambda: Path("sam2.1-large.pt"),
            visual_segment.VisualElement,
            lambda *args: None,
        ),
    )

    def fake_run(command, **kwargs):
        calls.append((command, kwargs))
        request_path = Path(command[command.index("--request") + 1])
        request = json.loads(request_path.read_text(encoding="utf-8"))
        assert set(request) == {"schema_version", "operations"}
        assert request["schema_version"] == 1
        assert request["operations"] == [
            {
                "id": "prompted",
                "kind": "prompted",
                "image": "image.png",
                "text_mask": "text-mask.png",
                "proposals": "proposals.json",
            },
            {"id": "automatic", "kind": "automatic", "image": "image.png"},
        ]
        previous_argv = sys.argv
        try:
            sys.argv = command[1:]
            returncode = sam_worker.main()
        finally:
            sys.argv = previous_argv
        return subprocess.CompletedProcess(command, returncode, "", "")

    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_run)

    prompted, automatic = batch_helper(image, text_mask, [proposal], tmp_path)

    expected_prompted = generate_prompted(image, [proposal], generator, text_mask)
    expected_automatic = generate_automatic(
        image,
        generator,
        crop_size=max(image.shape[:2]),
        include_geometry=False,
        min_score=0.90,
    )
    assert [_candidate_batch_signature(item) for item in prompted] == [
        _candidate_batch_signature(item) for item in expected_prompted
    ]
    assert [_candidate_batch_signature(item) for item in automatic] == [
        _candidate_batch_signature(item) for item in expected_automatic
    ]
    assert len(calls) == 1
    command, kwargs = calls[0]
    assert command[command.index("--mode") + 1] == "batch"
    assert "--image" not in command
    assert kwargs == {"capture_output": True, "text": True}
    assert events[:3] == [
        ("create", Path("sam2.1-large.pt"), True),
        "prompted",
        (
            "automatic",
            {"crop_size": 8, "include_geometry": False, "min_score": 0.90},
        ),
    ]


def test_candidate_batch_main_matches_legacy_candidate_records(
    tmp_path: Path,
    monkeypatch,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    prompted_mask = image_to_ppt.np.zeros((6, 8), dtype=bool)
    prompted_mask[1:5, 1:7] = True
    automatic_mask = image_to_ppt.np.zeros((6, 8), dtype=bool)
    automatic_mask[2:4, 2:6] = True

    def generate_prompted(*args):
        return [
            image_to_ppt.MaskCandidate(
                prompted_mask.copy(),
                -0.25,
                "grounded:full:object",
                (0, 0, 8, 6),
                True,
                "badge",
                "object",
                (1.0, 1.0, 7.0, 5.0),
            )
        ]

    def generate_automatic(*args, **kwargs):
        return [
            image_to_ppt.MaskCandidate(
                automatic_mask.copy(),
                0.97,
                "sam",
                (0, 0, 8, 6),
            )
        ]

    monkeypatch.setattr(
        sam_worker,
        "_load_tools",
        lambda: (
            image_to_ppt.ObjectProposal,
            lambda *args, **kwargs: object(),
            generate_automatic,
            generate_prompted,
            lambda: Path("sam2.1-large.pt"),
            visual_segment.VisualElement,
            lambda *args: None,
        ),
    )

    def run_main(arguments):
        monkeypatch.setattr(sys, "argv", ["sam_worker.py", *arguments])
        assert sam_worker.main() == 0

    legacy_prompted = fixture["root"] / "legacy-prompted.json"
    run_main(
        [
            "--mode", "prompted",
            "--image", str(fixture["image"]),
            "--text-mask", str(fixture["text_mask"]),
            "--proposals", str(fixture["proposals"]),
            "--result", str(legacy_prompted),
        ]
    )
    legacy_automatic = fixture["root"] / "legacy-automatic.json"
    run_main(
        [
            "--mode", "automatic",
            "--image", str(fixture["image"]),
            "--result", str(legacy_automatic),
        ]
    )
    batch_result = fixture["root"] / "batch-result.json"
    run_main(
        [
            "--mode", "batch",
            "--request", str(fixture["request"]),
            "--result", str(batch_result),
        ]
    )

    batch_payload = json.loads(batch_result.read_text(encoding="utf-8"))
    assert batch_payload["operations"][0]["candidates"] == json.loads(
        legacy_prompted.read_text(encoding="utf-8")
    )
    assert batch_payload["operations"][1]["candidates"] == json.loads(
        legacy_automatic.read_text(encoding="utf-8")
    )


def test_candidate_batch_legacy_main_still_requires_image(
    tmp_path: Path,
    monkeypatch,
) -> None:
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "sam_worker.py",
            "--mode", "automatic",
            "--result", str(tmp_path / "result.json"),
        ],
    )

    with pytest.raises(SystemExit) as error:
        sam_worker.main()

    assert error.value.code == 2
    assert loads == []


@pytest.mark.parametrize(
    "malformation",
    [
        "empty",
        "missing",
        "duplicate-id",
        "reordered",
        "unknown-kind",
        "wrong-mask-shape",
    ],
)
def test_candidate_batch_caller_rejects_the_entire_malformed_result(
    tmp_path: Path,
    monkeypatch,
    malformation: str,
) -> None:
    batch_helper = getattr(
        image_to_ppt,
        "_generate_sam_candidate_batch_isolated",
        None,
    )
    assert batch_helper is not None
    payload = {
        "schema_version": 1,
        "operations": [
            {
                "id": "prompted",
                "kind": "prompted",
                "candidates": [_candidate_batch_record()],
            },
            {
                "id": "automatic",
                "kind": "automatic",
                "candidates": [_candidate_batch_record()],
            },
        ],
    }
    if malformation == "empty":
        payload = {}
    elif malformation == "missing":
        payload["operations"].pop()
    elif malformation == "duplicate-id":
        payload["operations"][1]["id"] = "prompted"
    elif malformation == "reordered":
        payload["operations"].reverse()
    elif malformation == "unknown-kind":
        payload["operations"][1]["kind"] = "mystery"
    else:
        payload["operations"][1]["candidates"][0]["mask_shape"] = [3, 8]

    result_path = None

    def fake_run(command, **kwargs):
        nonlocal result_path
        result_path = Path(command[command.index("--result") + 1])
        result_path.write_text(json.dumps(payload), encoding="utf-8")
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_run)

    with pytest.raises(RuntimeError, match="SAM candidate batch"):
        batch_helper(
            image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8),
            image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8),
            [],
            tmp_path,
        )

    assert result_path is not None
    assert not result_path.exists()


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("score", float("nan")),
        ("source", 7),
        ("crop_box", "0,0,8,6"),
        ("touches_crop_edge", 1),
        ("label", "bad\nlabel"),
        ("source", "bad\u202esource"),
        ("score", True),
        ("object_box", [9.0, 1.0, 12.0, 4.0]),
    ],
)
def test_candidate_batch_caller_rejects_invalid_candidate_metadata(
    tmp_path: Path,
    monkeypatch,
    field: str,
    value,
) -> None:
    record = _candidate_batch_record()
    record[field] = value
    payload = _candidate_batch_payload(prompted=[record])

    def fake_run(command, **kwargs):
        result_path = Path(command[command.index("--result") + 1])
        result_path.write_text(json.dumps(payload), encoding="utf-8")
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_run)

    with pytest.raises(RuntimeError, match="SAM candidate batch"):
        image_to_ppt._generate_sam_candidate_batch_isolated(
            image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8),
            image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8),
            [],
            tmp_path,
        )


def test_candidate_batch_caller_accepts_negative_finite_sam_score(
    tmp_path: Path,
    monkeypatch,
) -> None:
    record = _candidate_batch_record()
    record["score"] = -0.25
    record["object_box"] = [-2.5, -1.0, 5.0, 4.0]
    payload = _candidate_batch_payload(prompted=[record])

    def fake_run(command, **kwargs):
        result_path = Path(command[command.index("--result") + 1])
        result_path.write_text(json.dumps(payload), encoding="utf-8")
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_run)
    proposal = image_to_ppt.ObjectProposal(
        box_xyxy=(1.0, 1.0, 7.0, 5.0),
        score=0.91,
        label="badge",
        role="object",
        source="full",
        crop_box=(0, 0, 8, 6),
    )

    prompted, automatic = image_to_ppt._generate_sam_candidate_batch_isolated(
        image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8),
        image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8),
        [proposal],
        tmp_path,
    )

    assert [candidate.score for candidate in prompted] == [-0.25]
    assert prompted[0].object_box == (-2.5, -1.0, 5.0, 4.0)
    assert automatic == []


def test_candidate_batch_caller_validates_all_records_before_constructing_any(
    tmp_path: Path,
    monkeypatch,
) -> None:
    good = _candidate_batch_record()
    bad = _candidate_batch_record()
    bad["crop_box"] = 7
    payload = _candidate_batch_payload(prompted=[good, bad])
    constructions = []
    original_candidate = image_to_ppt.MaskCandidate

    def tracking_candidate(*args, **kwargs):
        constructions.append((args, kwargs))
        return original_candidate(*args, **kwargs)

    def fake_run(command, **kwargs):
        result_path = Path(command[command.index("--result") + 1])
        result_path.write_text(json.dumps(payload), encoding="utf-8")
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(image_to_ppt, "MaskCandidate", tracking_candidate)
    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_run)

    with pytest.raises(RuntimeError, match="SAM candidate batch"):
        image_to_ppt._generate_sam_candidate_batch_isolated(
            image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8),
            image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8),
            [],
            tmp_path,
        )

    assert constructions == []


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("score", float("nan")),
        ("source", 7),
        ("crop_box", "0,0,8,6"),
        ("touches_crop_edge", 1),
        ("source", "bad\u202esource"),
        ("score", True),
    ],
)
def test_candidate_batch_worker_rejects_invalid_proposal_before_loading_model(
    tmp_path: Path,
    monkeypatch,
    field: str,
    value,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    proposal = fixture["proposal_records"][0].copy()
    proposal[field] = value
    fixture["proposals"].write_text(json.dumps([proposal]), encoding="utf-8")
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            fixture["root"] / "result.json",
        )

    assert loads == []


def test_candidate_batch_worker_accepts_intersecting_out_of_bounds_dino_box() -> None:
    record = {
        "box_xyxy": [-2.5, -1.0, 5.0, 4.0],
        "score": 0.91,
        "label": "badge",
        "role": "object",
        "source": "full",
        "crop_box": [0, 0, 8, 6],
        "touches_crop_edge": False,
    }

    validated = sam_worker._validate_batch_proposals([record], (6, 8))

    assert validated[0]["box_xyxy"] == (-2.5, -1.0, 5.0, 4.0)


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("box_xyxy", [9.0, 1.0, 12.0, 4.0]),
        ("box_xyxy", [0.0, 1.0, float("inf"), 4.0]),
        ("crop_box", [0.0, 0, 8, 6]),
    ],
)
def test_candidate_batch_worker_rejects_invalid_box_domains(
    field: str,
    value,
) -> None:
    record = {
        "box_xyxy": [1.0, 1.0, 7.0, 5.0],
        "score": 0.91,
        "label": "badge",
        "role": "object",
        "source": "full",
        "crop_box": [0, 0, 8, 6],
        "touches_crop_edge": False,
    }
    record[field] = value

    with pytest.raises(ValueError):
        sam_worker._validate_batch_proposals([record], (6, 8))


@pytest.mark.parametrize(
    "operation_change",
    ["missing", "extra", "reordered"],
)
def test_candidate_batch_worker_requires_exact_candidate_operations(
    tmp_path: Path,
    monkeypatch,
    operation_change: str,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    request = json.loads(fixture["request"].read_text(encoding="utf-8"))
    if operation_change == "missing":
        request["operations"].pop()
    elif operation_change == "extra":
        request["operations"].append(request["operations"][1].copy())
        request["operations"][2]["id"] = "extra"
    else:
        request["operations"].reverse()
    fixture["request"].write_text(json.dumps(request), encoding="utf-8")
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            fixture["root"] / "result.json",
        )

    assert loads == []


def test_candidate_batch_worker_requires_one_shared_image_before_loading_model(
    tmp_path: Path,
    monkeypatch,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    other_image = fixture["root"] / "other-image.png"
    Image.new("RGB", (16, 12), "black").save(other_image)
    request = json.loads(fixture["request"].read_text(encoding="utf-8"))
    request["operations"][1]["image"] = other_image.name
    fixture["request"].write_text(json.dumps(request), encoding="utf-8")
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError, match="same image"):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            fixture["root"] / "result.json",
        )

    assert loads == []


def test_candidate_batch_worker_bounds_request_before_loading_model(
    tmp_path: Path,
    monkeypatch,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    request_bytes = fixture["request"].read_bytes()
    monkeypatch.setattr(
        sam_worker,
        "_BATCH_MAX_REQUEST_BYTES",
        len(request_bytes) - 1,
        raising=False,
    )
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError, match="request"):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            fixture["root"] / "result.json",
        )

    assert loads == []


def test_candidate_batch_worker_bounds_proposals_before_loading_model(
    tmp_path: Path,
    monkeypatch,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    fixture["proposals"].write_text(
        json.dumps(fixture["proposal_records"] * 2),
        encoding="utf-8",
    )
    monkeypatch.setattr(
        sam_worker,
        "sam_candidate_batch_max_proposals",
        lambda image_shape: 1,
    )
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError, match="proposal"):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            fixture["root"] / "result.json",
        )

    assert loads == []


@pytest.mark.parametrize(
    ("fixture_name", "limit_kind"),
    [
        ("image", "image"),
        ("proposals", "proposals"),
    ],
)
def test_candidate_batch_worker_bounds_input_files_before_loading_model(
    tmp_path: Path,
    monkeypatch,
    fixture_name: str,
    limit_kind: str,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    limit = fixture[fixture_name].stat().st_size - 1
    if limit_kind == "image":
        monkeypatch.setattr(sam_worker, "_BATCH_MAX_INPUT_BYTES", limit)
    else:
        monkeypatch.setattr(
            sam_worker,
            "sam_candidate_batch_proposals_max_bytes",
            lambda image_shape: limit,
        )
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError, match="size limit"):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            fixture["root"] / "result.json",
        )

    assert loads == []


def test_candidate_batch_worker_bounds_candidates_before_decoding_masks(
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    payload = _candidate_batch_payload(
        automatic=[_candidate_batch_record(), _candidate_batch_record()]
    )
    monkeypatch.setattr(
        sam_worker,
        "sam_candidate_batch_max_automatic_candidates",
        lambda: 1,
    )
    decode_calls = []
    monkeypatch.setattr(
        sam_worker.base64,
        "b64decode",
        lambda *args, **kwargs: decode_calls.append((args, kwargs)),
    )

    with pytest.raises(RuntimeError, match="candidate"):
        sam_worker._validate_batch_output(payload, operations)

    assert decode_calls == []


def test_candidate_batch_worker_rejects_generated_candidate_overflow_before_encoding(
    tmp_path: Path,
    monkeypatch,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    candidate = object()
    maximum = sam_worker.sam_candidate_batch_max_automatic_candidates()
    monkeypatch.setattr(
        sam_worker,
        "_load_tools",
        lambda: (
            image_to_ppt.ObjectProposal,
            lambda *args, **kwargs: object(),
            lambda *args, **kwargs: [candidate] * (maximum + 1),
            lambda *args, **kwargs: [],
            lambda: Path("sam2.1-large.pt"),
            visual_segment.VisualElement,
            lambda *args: None,
        ),
    )
    monkeypatch.setattr(
        sam_worker,
        "_candidate_record",
        lambda candidate: pytest.fail("overflow candidates must not be encoded"),
    )

    with pytest.raises(RuntimeError, match="candidate count"):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            fixture["root"] / "result.json",
        )

    assert not (fixture["root"] / "result.json").exists()


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("score", float("nan")),
        ("source", 7),
        ("crop_box", "0,0,8,6"),
        ("touches_crop_edge", 1),
        ("source", "bad\u202esource"),
        ("score", True),
        ("object_box", [9.0, 1.0, 12.0, 4.0]),
    ],
)
def test_candidate_batch_worker_rejects_invalid_candidate_metadata(
    field: str,
    value,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    record = _candidate_batch_record()
    record[field] = value

    with pytest.raises(RuntimeError, match="candidate"):
        sam_worker._validate_batch_output(
            _candidate_batch_payload(automatic=[record]),
            operations,
        )


def test_candidate_batch_caller_bounds_candidates_before_decoding_masks(
    tmp_path: Path,
    monkeypatch,
) -> None:
    payload = _candidate_batch_payload(
        prompted=[_candidate_batch_record()],
        automatic=[_candidate_batch_record(), _candidate_batch_record()]
    )
    monkeypatch.setattr(
        image_to_ppt,
        "sam_candidate_batch_max_automatic_candidates",
        lambda: 1,
    )
    decode_calls = []

    def fake_run(command, **kwargs):
        result_path = Path(command[command.index("--result") + 1])
        result_path.write_text(json.dumps(payload), encoding="utf-8")
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_run)
    monkeypatch.setattr(
        image_to_ppt.base64,
        "b64decode",
        lambda *args, **kwargs: decode_calls.append((args, kwargs)),
    )

    with pytest.raises(RuntimeError, match="SAM candidate batch"):
        image_to_ppt._generate_sam_candidate_batch_isolated(
            image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8),
            image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8),
            [],
            tmp_path,
        )

    assert decode_calls == []


def test_candidate_batch_caller_bounds_result_before_parsing(
    tmp_path: Path,
    monkeypatch,
) -> None:
    payload_bytes = json.dumps(_candidate_batch_payload()).encode("utf-8")
    monkeypatch.setattr(
        image_to_ppt,
        "sam_candidate_batch_result_max_bytes",
        lambda image_shape, proposal_count: len(payload_bytes) - 1,
    )

    def fake_run(command, **kwargs):
        result_path = Path(command[command.index("--result") + 1])
        result_path.write_bytes(payload_bytes)
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_run)

    with pytest.raises(RuntimeError, match="SAM candidate batch"):
        image_to_ppt._generate_sam_candidate_batch_isolated(
            image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8),
            image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8),
            [],
            tmp_path,
        )


def test_candidate_batch_caller_rejects_oversized_mask_before_decoding(
    tmp_path: Path,
    monkeypatch,
) -> None:
    record = _candidate_batch_record()
    record["mask"] += "AAAA"
    payload = _candidate_batch_payload(automatic=[record])
    decode_calls = []

    def fake_run(command, **kwargs):
        result_path = Path(command[command.index("--result") + 1])
        result_path.write_text(json.dumps(payload), encoding="utf-8")
        return subprocess.CompletedProcess(command, 0, "", "")

    monkeypatch.setattr(image_to_ppt, "run_isolated_worker", fake_run)
    monkeypatch.setattr(
        image_to_ppt.base64,
        "b64decode",
        lambda *args, **kwargs: decode_calls.append((args, kwargs)),
    )

    with pytest.raises(RuntimeError, match="SAM candidate batch"):
        image_to_ppt._generate_sam_candidate_batch_isolated(
            image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8),
            image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8),
            [],
            tmp_path,
        )

    assert decode_calls == []


def test_candidate_batch_limits_cover_current_sam_generation_bounds() -> None:
    image_shape = (4096, 4096)
    assert sam_worker.sam_candidate_batch_max_proposals(image_shape) == 50 * 900
    assert sam_worker.sam_candidate_batch_max_automatic_candidates() == 16 * 16 * 3
    assert (
        sam_worker.sam_candidate_batch_max_prompted_candidates(50 * 900)
        == 2 * 50 * 900
    )


def test_candidate_batch_multitile_proposal_boundary_is_dynamic() -> None:
    image_shape = (4096, 4096)
    maximum = sam_worker.sam_candidate_batch_max_proposals(image_shape)
    record = {
        "box_xyxy": [1.0, 1.0, 7.0, 5.0],
        "score": 0.91,
        "label": "badge",
        "role": "object",
        "source": "full",
        "crop_box": [0, 0, 4096, 4096],
        "touches_crop_edge": False,
    }

    assert len(
        sam_worker._validate_batch_proposals([record] * maximum, image_shape)
    ) == maximum
    with pytest.raises(ValueError, match="proposal count"):
        sam_worker._validate_batch_proposals(
            [record] * (maximum + 1),
            image_shape,
        )


def test_candidate_batch_4k_result_budget_covers_exact_automatic_masks() -> None:
    image_shape = (2160, 3840)
    packed_bytes = (image_shape[0] * image_shape[1] + 7) // 8
    encoded_bytes = ((packed_bytes + 2) // 3) * 4
    metadata_record = _candidate_batch_record((1, 8))
    metadata_record["mask"] = ""
    metadata_record["score"] = -(10**4299)
    metadata_record["source"] = "😀" * 256
    metadata_record["label"] = "😀" * 256
    metadata_record["role"] = "😀" * 256
    metadata_record["object_box"] = [-(10**4299), -1, 8, 1]
    metadata_bytes = len(
        json.dumps(metadata_record, ensure_ascii=False).encode("utf-8")
    )
    automatic_count = 16 * 16 * 3
    minimum = automatic_count * (encoded_bytes + metadata_bytes)

    automatic_budget = sam_worker.sam_candidate_batch_result_max_bytes(
        image_shape,
        proposal_count=0,
    )
    maximum_proposals = sam_worker.sam_candidate_batch_max_proposals(image_shape)
    worker_budget = sam_worker.sam_candidate_batch_result_max_bytes(
        image_shape,
        proposal_count=maximum_proposals,
    )
    caller_budget = image_to_ppt.sam_candidate_batch_result_max_bytes(
        image_shape,
        proposal_count=maximum_proposals,
    )
    assert automatic_budget >= minimum
    assert caller_budget == worker_budget
    assert worker_budget > 1024**3


@pytest.mark.parametrize(
    "invalid_request",
    ["unknown-kind", "duplicate-id", "escape"],
)
def test_candidate_batch_worker_rejects_invalid_schema_before_loading_sam(
    tmp_path: Path,
    monkeypatch,
    invalid_request: str,
) -> None:
    Image.new("RGB", (8, 6), "white").save(tmp_path / "image.png")
    Image.new("L", (8, 6), 0).save(tmp_path / "text-mask.png")
    (tmp_path / "proposals.json").write_text("[]", encoding="utf-8")
    operations = [
        {
            "id": "prompted",
            "kind": "prompted",
            "image": "image.png",
            "text_mask": "text-mask.png",
            "proposals": "proposals.json",
        },
        {"id": "automatic", "kind": "automatic", "image": "image.png"},
    ]
    if invalid_request == "unknown-kind":
        operations[1]["kind"] = "mystery"
    elif invalid_request == "duplicate-id":
        operations[1]["id"] = "prompted"
    elif invalid_request == "escape":
        operations[1]["image"] = "../image.png"
    request_path = tmp_path / "request.json"
    request_path.write_text(
        json.dumps({"schema_version": 1, "operations": operations}),
        encoding="utf-8",
    )
    result_path = tmp_path / "result.json"
    load_events = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: load_events.append("load"))
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "sam_worker.py",
            "--mode",
            "batch",
            "--request",
            str(request_path),
            "--result",
            str(result_path),
        ],
    )

    with pytest.raises(ValueError):
        sam_worker.main()

    assert load_events == []
    assert not result_path.exists()
    assert not list(tmp_path.glob(".result.json.*.tmp"))


@pytest.mark.parametrize(
    "result_name",
    ["request", "image", "text_mask", "proposals"],
)
def test_candidate_batch_worker_rejects_result_alias_without_deleting_input(
    tmp_path: Path,
    monkeypatch,
    result_name: str,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    before = {
        name: fixture[name].read_bytes()
        for name in ("request", "image", "text_mask", "proposals")
    }
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            fixture[result_name],
        )

    assert loads == []
    assert {
        name: fixture[name].read_bytes()
        for name in ("request", "image", "text_mask", "proposals")
    } == before


def test_candidate_batch_worker_rejects_existing_result_without_deleting_it(
    tmp_path: Path,
    monkeypatch,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    result_path = fixture["root"] / "result.json"
    result_path.write_bytes(b"existing result")
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError, match="already exists"):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            result_path,
        )

    assert loads == []
    assert result_path.read_bytes() == b"existing result"


def test_candidate_batch_worker_rejects_result_hardlink_to_input(
    tmp_path: Path,
    monkeypatch,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    result_path = fixture["root"] / "result.json"
    os.link(fixture["image"], result_path)
    before = fixture["image"].read_bytes()
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError, match="already exists"):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            result_path,
        )

    assert loads == []
    assert fixture["image"].read_bytes() == before
    assert result_path.read_bytes() == before


@pytest.mark.parametrize("target_name", ["image", "external"])
def test_candidate_batch_worker_rejects_result_symlink_without_touching_target(
    tmp_path: Path,
    monkeypatch,
    target_name: str,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "batch")
    if target_name == "image":
        target = fixture["image"]
    else:
        target = tmp_path / "external.txt"
        target.write_bytes(b"external content")
    before = target.read_bytes()
    result_path = fixture["root"] / "result.json"
    try:
        result_path.symlink_to(target)
    except OSError as exc:
        pytest.skip(f"symlink creation is unavailable: {exc}")
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError):
        _run_candidate_batch_worker_main(
            monkeypatch,
            fixture["request"],
            result_path,
        )

    assert loads == []
    assert target.read_bytes() == before
    assert result_path.is_symlink()


def test_candidate_batch_worker_rejects_linked_parent_before_loading_model(
    tmp_path: Path,
    monkeypatch,
) -> None:
    fixture = _candidate_batch_worker_fixture(tmp_path / "actual")
    linked_root = tmp_path / "linked"
    try:
        linked_root.symlink_to(fixture["root"], target_is_directory=True)
    except OSError as exc:
        if os.name != "nt":
            pytest.skip(f"directory symlink creation is unavailable: {exc}")
        junction = subprocess.run(
            ["cmd", "/c", "mklink", "/J", str(linked_root), str(fixture["root"])],
            capture_output=True,
            text=True,
        )
        if junction.returncode != 0:
            pytest.skip(
                "directory symlink and junction creation are unavailable: "
                f"{junction.stderr.strip()}"
            )
    loads = []
    monkeypatch.setattr(sam_worker, "_load_tools", lambda: loads.append("load"))

    with pytest.raises(ValueError, match="directory"):
        _run_candidate_batch_worker_main(
            monkeypatch,
            linked_root / "request.json",
            linked_root / "result.json",
        )

    assert loads == []
    assert fixture["request"].is_file()
    assert not (fixture["root"] / "result.json").exists()


def test_candidate_batch_worker_rejects_resolved_path_escape(
    tmp_path: Path,
    monkeypatch,
) -> None:
    batch_root = tmp_path / "batch"
    batch_root.mkdir()
    linked_image = batch_root / "linked-image.png"
    linked_image.write_bytes(b"inside placeholder")
    outside_image = tmp_path / "outside.png"
    outside_image.write_bytes(b"outside target")
    resolved_root = batch_root.resolve()
    resolved_outside = outside_image.resolve()
    original_resolve = Path.resolve

    def fake_resolve(path, *args, **kwargs):
        if path == linked_image:
            return resolved_outside
        return original_resolve(path, *args, **kwargs)

    monkeypatch.setattr(Path, "resolve", fake_resolve)

    with pytest.raises(ValueError, match="stay inside"):
        sam_worker._batch_file(resolved_root, linked_image.name, "image")


def test_candidate_batch_worker_validates_output_before_replacing_result(
    tmp_path: Path,
    monkeypatch,
) -> None:
    Image.new("RGB", (8, 6), "white").save(tmp_path / "image.png")
    Image.new("L", (8, 6), 0).save(tmp_path / "text-mask.png")
    (tmp_path / "proposals.json").write_text("[]", encoding="utf-8")
    request_path = tmp_path / "request.json"
    request_path.write_text(
        json.dumps(
            {
                "schema_version": 1,
                "operations": [
                    {
                        "id": "prompted",
                        "kind": "prompted",
                        "image": "image.png",
                        "text_mask": "text-mask.png",
                        "proposals": "proposals.json",
                    },
                    {"id": "automatic", "kind": "automatic", "image": "image.png"},
                ],
            }
        ),
        encoding="utf-8",
    )
    result_path = tmp_path / "result.json"
    wrong_mask = image_to_ppt.np.ones((2, 2), dtype=bool)
    monkeypatch.setattr(
        sam_worker,
        "_load_tools",
        lambda: (
            image_to_ppt.ObjectProposal,
            lambda *args, **kwargs: object(),
            lambda *args, **kwargs: [
                image_to_ppt.MaskCandidate(wrong_mask, 0.99, "sam")
            ],
            lambda *args, **kwargs: [],
            lambda: Path("sam.pt"),
            visual_segment.VisualElement,
            lambda *args: None,
        ),
    )
    monkeypatch.setattr(
        sys,
        "argv",
        [
            "sam_worker.py",
            "--mode",
            "batch",
            "--request",
            str(request_path),
            "--result",
            str(result_path),
        ],
    )

    with pytest.raises(RuntimeError, match="mask shape"):
        sam_worker.main()

    assert not result_path.exists()
    assert not list(tmp_path.glob(".result.json.*.tmp"))


def test_candidate_batch_worker_uses_exclusive_random_result_temp(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    payload = _candidate_batch_payload(automatic=[_candidate_batch_record()])
    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    result_binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    created_files = []
    actual_create = sam_worker._create_batch_result_file

    def tracked_create(binding, parent_handle):
        created = actual_create(binding, parent_handle)
        created_files.append(created[1:])
        return created

    monkeypatch.setattr(sam_worker, "_create_batch_result_file", tracked_create)
    monkeypatch.setattr(
        sam_worker.json,
        "dumps",
        lambda *args, **kwargs: pytest.fail("batch writer must stream JSON"),
    )
    monkeypatch.setattr(
        sam_worker.json.JSONEncoder,
        "encode",
        lambda *args, **kwargs: pytest.fail("batch writer must use iterencode"),
    )

    sam_worker._write_batch_result(result_binding, payload, operations)

    assert len(created_files) == 1
    created_path, direct_result = created_files[0]
    if created_path is not None and not direct_result:
        assert created_path.parent == tmp_path
        assert created_path.name.startswith(".result.json.")
        assert created_path.name.endswith(".tmp")
    assert json.loads(result_path.read_text(encoding="utf-8")) == payload


def test_candidate_batch_worker_publishes_verified_temp_inode_when_path_is_replaced(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    payload = _candidate_batch_payload()
    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    result_binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    attack = b"unverified replacement"
    actual_verify = sam_worker._verify_batch_result_binding
    verify_calls = 0
    replaced_path = None
    replacement_name = None

    def replace_temp_after_verification(binding):
        nonlocal verify_calls, replaced_path, replacement_name
        actual_verify(binding)
        verify_calls += 1
        if verify_calls == 2:
            temporary_paths = list(tmp_path.glob(".result.json.*.tmp"))
            if not temporary_paths:
                return
            temporary_path = temporary_paths[0]
            replacement_name = temporary_path.name
            replaced_path = tmp_path / "validated-original"
            temporary_path.replace(replaced_path)
            temporary_path.write_bytes(attack)

    monkeypatch.setattr(
        sam_worker,
        "_verify_batch_result_binding",
        replace_temp_after_verification,
    )

    sam_worker._write_batch_result(result_binding, payload, operations)

    assert json.loads(result_path.read_text(encoding="utf-8")) == payload
    if replaced_path is not None:
        assert (tmp_path / replacement_name).read_bytes() == attack


def test_candidate_batch_worker_rejects_parent_replacement_during_publish(
    tmp_path: Path,
    monkeypatch,
) -> None:
    owned_parent = tmp_path / "owned"
    owned_parent.mkdir()
    moved_parent = tmp_path / "moved"
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    payload = _candidate_batch_payload()
    result_path = owned_parent / "result.json"
    root_status = owned_parent.lstat()
    result_binding = {
        "path": result_path,
        "parent": owned_parent,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    attack = b"replacement-parent temp"
    actual_verify = sam_worker._verify_batch_result_binding
    verify_calls = 0

    def replace_parent_after_verification(binding):
        nonlocal verify_calls
        actual_verify(binding)
        verify_calls += 1
        if verify_calls == 2:
            temporary_paths = list(owned_parent.glob(".result.json.*.tmp"))
            temporary_name = (
                temporary_paths[0].name
                if temporary_paths
                else ".result.json.replacement.tmp"
            )
            try:
                owned_parent.replace(moved_parent)
            except PermissionError as exc:
                pytest.skip(f"open directory handles prevent rename: {exc}")
            owned_parent.mkdir()
            (owned_parent / temporary_name).write_bytes(attack)

    monkeypatch.setattr(
        sam_worker,
        "_verify_batch_result_binding",
        replace_parent_after_verification,
    )

    with pytest.raises(RuntimeError, match="directory changed"):
        sam_worker._write_batch_result(result_binding, payload, operations)

    assert not result_path.exists()
    assert not (moved_parent / "result.json").exists()
    assert next(owned_parent.glob(".result.json.*.tmp")).read_bytes() == attack


def test_candidate_batch_worker_cleans_partial_stream_when_limit_is_exceeded(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    payload = _candidate_batch_payload(automatic=[_candidate_batch_record()])
    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    result_binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    monkeypatch.setattr(
        sam_worker,
        "sam_candidate_batch_result_max_bytes",
        lambda image_shape, proposal_count: 32,
    )

    with pytest.raises(RuntimeError, match="size limit"):
        sam_worker._write_batch_result(result_binding, payload, operations)

    assert not result_path.exists()
    assert not list(tmp_path.glob(".result.json.*.tmp"))


def test_candidate_batch_linux_publish_falls_back_to_proc_fd_without_capability(
    monkeypatch,
) -> None:
    import errno

    calls = []
    outcomes = iter([errno.ENOENT, 0])

    def fake_linkat(old_fd, old_path, new_fd, new_path, flags):
        calls.append((old_fd, old_path, new_fd, new_path, flags))
        return next(outcomes)

    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(
        sam_worker,
        "_linkat_batch_result",
        fake_linkat,
        raising=False,
    )

    sam_worker._publish_batch_result(
        17,
        19,
        {"path": Path("result.json")},
    )

    assert calls == [
        (17, b"", 19, b"result.json", 0x1000),
        (-100, b"/proc/self/fd/17", 19, b"result.json", 0x400),
    ]


def test_candidate_batch_linux_publish_reports_both_link_mechanisms_unsupported(
    monkeypatch,
) -> None:
    outcomes = iter([errno.ENOENT, errno.ENOENT])
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(
        sam_worker,
        "_linkat_batch_result",
        lambda *args: next(outcomes),
    )

    with pytest.raises(sam_worker._BatchResultPublishingUnsupported):
        sam_worker._publish_batch_result(
            17,
            19,
            {"path": Path("capability-result")},
        )


def test_candidate_batch_other_posix_never_creates_partial_final(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setattr(sam_worker.sys, "platform", "freebsd14")
    monkeypatch.setattr(
        sam_worker,
        "_open_batch_result_parent",
        lambda binding: 7,
    )
    monkeypatch.setattr(
        sam_worker,
        "_close_batch_result_parent",
        lambda handle: None,
    )
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    with pytest.raises(RuntimeError, match="unsupported platform"):
        sam_worker._write_batch_result(
            binding,
            _candidate_batch_payload(automatic=[_candidate_batch_record()]),
            operations,
        )

    assert not result_path.exists()
    assert not hasattr(sam_worker, "_unlink_owned_posix_result")


@pytest.mark.parametrize(
    ("clone_error", "message"),
    [(0, None), (errno.EEXIST, "already exists"), (errno.ENOTSUP, "not supported")],
)
def test_candidate_batch_darwin_fclone_is_descriptor_bound_and_no_clobber(
    monkeypatch,
    clone_error: int,
    message: str | None,
) -> None:
    calls = []
    monkeypatch.setattr(sam_worker.sys, "platform", "darwin")
    monkeypatch.setattr(
        sam_worker,
        "_fclonefileat_batch_result",
        lambda source_fd, parent_fd, name, flags: (
            calls.append((source_fd, parent_fd, name, flags)) or clone_error
        ),
        raising=False,
    )

    if message is None:
        sam_worker._publish_batch_result(17, 19, {"path": Path("result.json")})
    else:
        with pytest.raises(RuntimeError, match=message):
            sam_worker._publish_batch_result(
                17,
                19,
                {"path": Path("result.json")},
            )

    assert calls == [(17, 19, b"result.json", 0)]


def test_candidate_batch_darwin_without_anonymous_source_fails_before_path_creation(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setattr(sam_worker.sys, "platform", "darwin")
    monkeypatch.setattr(
        sam_worker.os,
        "open",
        lambda *args, **kwargs: pytest.fail("Darwin must not create a named temp"),
    )
    monkeypatch.setattr(
        sam_worker.os,
        "unlink",
        lambda *args, **kwargs: pytest.fail("Darwin must not path-unlink cleanup"),
    )

    with pytest.raises(RuntimeError, match="anonymous.*not supported"):
        sam_worker._create_batch_result_file(
            {"path": tmp_path / "result.json"},
            19,
        )

    assert not (tmp_path / "result.json").exists()


def test_candidate_batch_linux_tmpfile_failure_does_not_create_final(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    monkeypatch.setattr(
        sam_worker.os,
        "open",
        lambda *args, **kwargs: (_ for _ in ()).throw(OSError("unsupported")),
    )
    result_path = tmp_path / "result.json"

    with pytest.raises(RuntimeError, match="unsupported"):
        sam_worker._create_batch_result_file({"path": result_path}, 19)

    assert not result_path.exists()


@pytest.mark.skipif(sys.platform != "win32", reason="Windows sharing contract")
def test_candidate_batch_windows_temp_denies_external_writers(tmp_path: Path) -> None:
    import ctypes
    from ctypes import wintypes

    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    result_binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    parent_handle = sam_worker._open_batch_result_parent(result_binding)
    file_descriptor = None
    try:
        file_descriptor, temporary_path, published = (
            sam_worker._create_batch_result_file(result_binding, parent_handle)
        )
        assert temporary_path is not None
        assert not published
        kernel32 = ctypes.WinDLL("kernel32", use_last_error=True)
        kernel32.CreateFileW.argtypes = [
            wintypes.LPCWSTR,
            wintypes.DWORD,
            wintypes.DWORD,
            ctypes.c_void_p,
            wintypes.DWORD,
            wintypes.DWORD,
            wintypes.HANDLE,
        ]
        kernel32.CreateFileW.restype = wintypes.HANDLE
        attacker_handle = kernel32.CreateFileW(
            str(temporary_path),
            0x40000000,
            0x00000001 | 0x00000002 | 0x00000004,
            None,
            3,
            0x00000080,
            None,
        )
        if attacker_handle != ctypes.c_void_p(-1).value:
            kernel32.CloseHandle(attacker_handle)
            pytest.fail("batch result temp accepted an external writer")
        assert ctypes.get_last_error() == 32
    finally:
        if file_descriptor is not None:
            sam_worker._delete_windows_batch_result(file_descriptor)
            os.close(file_descriptor)
        sam_worker._close_batch_result_parent(parent_handle)

    assert not list(tmp_path.glob(".result.json.*.tmp"))


def test_candidate_batch_cleanup_failure_preserves_primary_error_and_closes_handles(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    file_descriptor, temporary_name = sam_worker.tempfile.mkstemp(dir=tmp_path)
    parent_handle = object()
    closed = []
    actual_close = sam_worker.os.close
    monkeypatch.setattr(sam_worker.sys, "platform", "win32")
    monkeypatch.setattr(
        sam_worker,
        "_open_batch_result_parent",
        lambda actual_binding: parent_handle,
    )
    monkeypatch.setattr(
        sam_worker,
        "_create_batch_result_file",
        lambda actual_binding, actual_parent: (
            file_descriptor,
            Path(temporary_name),
            False,
        ),
    )
    monkeypatch.setattr(
        sam_worker,
        "_delete_windows_batch_result",
        lambda descriptor: (_ for _ in ()).throw(RuntimeError("cleanup failed")),
    )
    monkeypatch.setattr(
        sam_worker.os,
        "close",
        lambda descriptor: closed.append(("file", descriptor)),
    )
    monkeypatch.setattr(
        sam_worker,
        "_close_batch_result_parent",
        lambda handle: closed.append(("parent", handle)),
    )
    monkeypatch.setattr(
        sam_worker,
        "sam_candidate_batch_result_max_bytes",
        lambda image_shape, proposal_count: 32,
    )

    with pytest.raises(RuntimeError, match="size limit") as raised:
        sam_worker._write_batch_result(
            binding,
            _candidate_batch_payload(automatic=[_candidate_batch_record()]),
            operations,
        )

    assert closed == [("file", file_descriptor), ("parent", parent_handle)]
    assert any("cleanup failed" in note for note in raised.value.__notes__)
    actual_close(file_descriptor)
    Path(temporary_name).unlink()


def test_candidate_batch_close_failure_does_not_override_primary_error(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    actual_close = sam_worker.os.close
    close_calls = []

    def failing_close(descriptor):
        close_calls.append(descriptor)
        actual_close(descriptor)
        raise OSError("close failed")

    monkeypatch.setattr(sam_worker.os, "close", failing_close)
    monkeypatch.setattr(
        sam_worker,
        "sam_candidate_batch_result_max_bytes",
        lambda image_shape, proposal_count: 32,
    )

    with pytest.raises(RuntimeError, match="size limit") as raised:
        sam_worker._write_batch_result(
            binding,
            _candidate_batch_payload(automatic=[_candidate_batch_record()]),
            operations,
        )

    assert close_calls
    assert any("close failed" in note for note in raised.value.__notes__)


@pytest.mark.skipif(sys.platform != "win32", reason="Windows handle publish")
def test_candidate_batch_windows_close_failure_after_publish_keeps_success(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    payload = _candidate_batch_payload()
    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    actual_close = sam_worker.os.close

    def failing_close(descriptor):
        actual_close(descriptor)
        raise OSError("close failed after publish")

    monkeypatch.setattr(sam_worker.os, "close", failing_close)

    sam_worker._write_batch_result(binding, payload, operations)

    assert json.loads(result_path.read_text(encoding="utf-8")) == payload


def test_candidate_batch_worker_does_not_replace_result_created_during_publish(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    operations = [
        {"id": "prompted", "kind": "prompted", "image": image},
        {"id": "automatic", "kind": "automatic", "image": image},
    ]
    payload = _candidate_batch_payload()
    result_path = tmp_path / "result.json"
    root_status = tmp_path.lstat()
    result_binding = {
        "path": result_path,
        "parent": tmp_path,
        "parent_identity": (root_status.st_dev, root_status.st_ino),
    }
    intruder = b"created during publish"
    actual_verify = sam_worker._verify_batch_result_binding
    verify_calls = 0

    def create_result_after_verification(binding):
        nonlocal verify_calls
        actual_verify(binding)
        verify_calls += 1
        if verify_calls == 2:
            result_path.write_bytes(intruder)

    monkeypatch.setattr(
        sam_worker,
        "_verify_batch_result_binding",
        create_result_after_verification,
    )

    with pytest.raises(RuntimeError, match="already exists"):
        sam_worker._write_batch_result(result_binding, payload, operations)

    assert result_path.read_bytes() == intruder
    assert not list(tmp_path.glob(".result.json.*.tmp"))


def test_candidate_batch_stage_supported_uses_one_batch_worker(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    text_mask = image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8)
    proposal = image_to_ppt.ObjectProposal(
        box_xyxy=(1.0, 1.0, 7.0, 5.0),
        score=0.9,
        label="object",
        role="foreground",
        source="dino",
        crop_box=(0, 0, 8, 6),
    )
    prompted = [
        image_to_ppt.MaskCandidate(
            image_to_ppt.np.ones((6, 8), dtype=bool), 0.8, "sam"
        )
    ]
    automatic = [
        image_to_ppt.MaskCandidate(
            image_to_ppt.np.eye(6, 8, dtype=bool), 0.7, "sam"
        )
    ]
    batch_calls = []
    monkeypatch.setattr(
        image_to_ppt,
        "sam_candidate_batch_output_supported",
        lambda work_dir: True,
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidate_batch_isolated",
        lambda *args: batch_calls.append(args) or (prompted, automatic),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidates_isolated",
        lambda *args, **kwargs: pytest.fail("legacy workers must not run"),
    )

    actual = image_to_ppt._generate_sam_candidate_stage_isolated(
        image,
        text_mask,
        [proposal],
        tmp_path,
    )

    assert actual == (prompted, automatic)
    assert batch_calls == [(image, text_mask, [proposal], tmp_path)]


def test_candidate_batch_stage_unsupported_uses_two_legacy_workers_in_order(
    tmp_path: Path,
    monkeypatch,
) -> None:
    image = image_to_ppt.np.zeros((6, 8, 3), dtype=image_to_ppt.np.uint8)
    text_mask = image_to_ppt.np.zeros((6, 8), dtype=image_to_ppt.np.uint8)
    proposals = []
    prompted = [
        image_to_ppt.MaskCandidate(
            image_to_ppt.np.ones((6, 8), dtype=bool), 0.8, "sam"
        )
    ]
    automatic = [
        image_to_ppt.MaskCandidate(
            image_to_ppt.np.eye(6, 8, dtype=bool), 0.7, "sam"
        )
    ]
    legacy_calls = []
    monkeypatch.setattr(
        image_to_ppt,
        "sam_candidate_batch_output_supported",
        lambda work_dir: False,
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidate_batch_isolated",
        lambda *args: pytest.fail("batch worker must not run"),
    )

    def fake_legacy(actual_image, actual_mask, actual_proposals, work_dir, *, mode):
        legacy_calls.append(
            (actual_image, actual_mask, actual_proposals, work_dir, mode)
        )
        return prompted if mode == "prompted" else automatic

    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidates_isolated",
        fake_legacy,
    )

    actual = image_to_ppt._generate_sam_candidate_stage_isolated(
        image,
        text_mask,
        proposals,
        tmp_path,
    )

    assert actual == (prompted, automatic)
    assert legacy_calls == [
        (image, text_mask, proposals, tmp_path, "prompted"),
        (image, None, None, tmp_path, "automatic"),
    ]


def test_candidate_batch_stage_does_not_fallback_after_batch_failure(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setattr(
        image_to_ppt,
        "sam_candidate_batch_output_supported",
        lambda work_dir: True,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidate_batch_isolated",
        lambda *args: (_ for _ in ()).throw(RuntimeError("batch failed")),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidates_isolated",
        lambda *args, **kwargs: pytest.fail("failed batch must not be retried"),
    )

    with pytest.raises(RuntimeError, match="batch failed"):
        image_to_ppt._generate_sam_candidate_stage_isolated(
            image_to_ppt.np.zeros((2, 2, 3), dtype=image_to_ppt.np.uint8),
            image_to_ppt.np.zeros((2, 2), dtype=image_to_ppt.np.uint8),
            [],
            tmp_path,
        )


@pytest.mark.parametrize("platform", ["darwin", "freebsd14"])
def test_candidate_batch_output_capability_is_false_without_path_probe(
    tmp_path: Path,
    monkeypatch,
    platform: str,
) -> None:
    monkeypatch.setattr(sam_worker.sys, "platform", platform)
    monkeypatch.setattr(
        sam_worker.os,
        "open",
        lambda *args, **kwargs: pytest.fail("unsupported platforms need no probe file"),
    )

    assert sam_worker.sam_candidate_batch_output_supported(tmp_path) is False


def _mock_linux_candidate_batch_probe_directory(tmp_path: Path, monkeypatch) -> Path:
    probe_dir = tmp_path / "private-probe"
    probe_dir.mkdir(mode=0o700)
    identity = (probe_dir.stat().st_dev, probe_dir.stat().st_ino)
    monkeypatch.setattr(
        sam_worker,
        "_create_linux_batch_capability_directory",
        lambda root: (probe_dir, identity),
    )
    monkeypatch.setattr(
        sam_worker,
        "_remove_linux_batch_capability_directory",
        lambda path, expected: path.rmdir(),
    )
    return probe_dir


@pytest.mark.parametrize(
    "error_number",
    [errno.EOPNOTSUPP, errno.EINVAL, errno.EISDIR, errno.ENOENT, errno.ENOSYS],
)
def test_candidate_batch_linux_known_unsupported_filesystem_returns_false(
    tmp_path: Path,
    monkeypatch,
    error_number: int,
) -> None:
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    _mock_linux_candidate_batch_probe_directory(tmp_path, monkeypatch)
    monkeypatch.setattr(sam_worker, "_open_batch_result_parent", lambda binding: 90)
    monkeypatch.setattr(sam_worker, "_close_batch_result_parent", lambda handle: None)
    monkeypatch.setattr(
        sam_worker.os,
        "open",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            OSError(error_number, "unsupported")
        ),
    )

    assert sam_worker.sam_candidate_batch_output_supported(tmp_path) is False


def test_candidate_batch_linux_indeterminate_probe_error_prevents_inference(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    _mock_linux_candidate_batch_probe_directory(tmp_path, monkeypatch)
    monkeypatch.setattr(sam_worker, "_open_batch_result_parent", lambda binding: 90)
    monkeypatch.setattr(sam_worker, "_close_batch_result_parent", lambda handle: None)
    monkeypatch.setattr(
        sam_worker.os,
        "open",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            PermissionError(errno.EACCES, "denied")
        ),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "sam_candidate_batch_output_supported",
        sam_worker.sam_candidate_batch_output_supported,
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidate_batch_isolated",
        lambda *args: pytest.fail("batch inference must not run"),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidates_isolated",
        lambda *args, **kwargs: pytest.fail("legacy inference must not run"),
    )

    with pytest.raises(PermissionError, match="denied"):
        image_to_ppt._generate_sam_candidate_stage_isolated(
            image_to_ppt.np.zeros((2, 2, 3), dtype=image_to_ppt.np.uint8),
            image_to_ppt.np.zeros((2, 2), dtype=image_to_ppt.np.uint8),
            [],
            tmp_path,
        )


def test_candidate_batch_linux_missing_work_dir_is_not_treated_as_unsupported(
    tmp_path: Path,
    monkeypatch,
) -> None:
    missing = tmp_path / "missing"
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    monkeypatch.setattr(
        sam_worker.os,
        "open",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            FileNotFoundError(errno.ENOENT, "missing")
        ),
    )

    with pytest.raises(ValueError, match="does not exist"):
        sam_worker.sam_candidate_batch_output_supported(missing)


def test_candidate_batch_linux_supported_probe_closes_anonymous_descriptor(
    tmp_path: Path,
    monkeypatch,
) -> None:
    closed = []
    published = []
    cleaned = []
    probe_dir = tmp_path / "private-probe"
    probe_dir.mkdir(mode=0o700)
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    monkeypatch.setattr(
        sam_worker,
        "_create_linux_batch_capability_directory",
        lambda root: (probe_dir, probe_dir.stat()),
        raising=False,
    )
    monkeypatch.setattr(
        sam_worker,
        "_remove_linux_batch_capability_directory",
        lambda path, identity: path.rmdir(),
        raising=False,
    )
    monkeypatch.setattr(sam_worker, "_open_batch_result_parent", lambda binding: 90)
    monkeypatch.setattr(
        sam_worker,
        "_close_batch_result_parent",
        lambda handle: closed.append(handle),
    )
    monkeypatch.setattr(sam_worker.os, "open", lambda *args, **kwargs: 91)
    monkeypatch.setattr(sam_worker.os, "close", closed.append)
    monkeypatch.setattr(
        sam_worker,
        "_publish_batch_result",
        lambda descriptor, parent, binding: published.append(
            (descriptor, parent, binding["path"])
        ),
    )
    monkeypatch.setattr(
        sam_worker,
        "_unlink_linux_batch_capability_result",
        lambda descriptor, parent, name: cleaned.append((descriptor, parent, name)),
        raising=False,
    )

    assert sam_worker.sam_candidate_batch_output_supported(tmp_path) is True
    assert closed == [91, 90]
    assert len(published) == 1
    assert published[0][:2] == (91, 90)
    assert published[0][2].parent == probe_dir
    assert cleaned == [(91, 90, published[0][2].name)]
    assert not list(tmp_path.iterdir())


@pytest.mark.skipif(os.name != "posix", reason="POSIX permission bits")
def test_candidate_batch_linux_private_probe_directory_is_mode_700(
    tmp_path: Path,
) -> None:
    probe_dir, identity = sam_worker._create_linux_batch_capability_directory(
        tmp_path
    )

    assert stat.S_IMODE(probe_dir.stat().st_mode) == 0o700
    sam_worker._remove_linux_batch_capability_directory(probe_dir, identity)
    assert not list(tmp_path.iterdir())


@pytest.mark.skipif(os.name == "posix", reason="simulates invalid permission bits")
def test_candidate_batch_linux_unsafe_private_probe_directory_is_removed(
    tmp_path: Path,
) -> None:
    with pytest.raises(RuntimeError, match="unsafe"):
        sam_worker._create_linux_batch_capability_directory(tmp_path)

    assert not list(tmp_path.iterdir())


def test_candidate_batch_linux_publish_unsupported_returns_false_and_cleans(
    tmp_path: Path,
    monkeypatch,
) -> None:
    probe_dir = tmp_path / "private-probe"
    probe_dir.mkdir(mode=0o700)
    cleaned = []
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    monkeypatch.setattr(
        sam_worker,
        "_create_linux_batch_capability_directory",
        lambda root: (probe_dir, probe_dir.stat()),
        raising=False,
    )
    monkeypatch.setattr(
        sam_worker,
        "_remove_linux_batch_capability_directory",
        lambda path, identity: path.rmdir(),
        raising=False,
    )
    monkeypatch.setattr(sam_worker, "_open_batch_result_parent", lambda binding: 90)
    monkeypatch.setattr(sam_worker, "_close_batch_result_parent", lambda handle: None)
    monkeypatch.setattr(sam_worker.os, "open", lambda *args, **kwargs: 91)
    monkeypatch.setattr(sam_worker.os, "close", lambda descriptor: None)
    monkeypatch.setattr(
        sam_worker,
        "_publish_batch_result",
        lambda *args: (_ for _ in ()).throw(
            sam_worker._BatchResultPublishingUnsupported("unsupported")
        ),
    )
    monkeypatch.setattr(
        sam_worker,
        "_unlink_linux_batch_capability_result",
        lambda *args: cleaned.append(args),
        raising=False,
    )

    assert sam_worker.sam_candidate_batch_output_supported(tmp_path) is False
    assert cleaned == []
    assert not list(tmp_path.iterdir())


def test_candidate_batch_linux_probe_cleanup_failure_is_not_fallback(
    tmp_path: Path,
    monkeypatch,
) -> None:
    probe_dir = tmp_path / "private-probe"
    probe_dir.mkdir(mode=0o700)
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    monkeypatch.setattr(
        sam_worker,
        "_create_linux_batch_capability_directory",
        lambda root: (probe_dir, probe_dir.stat()),
        raising=False,
    )
    monkeypatch.setattr(
        sam_worker,
        "_remove_linux_batch_capability_directory",
        lambda path, identity: path.rmdir(),
        raising=False,
    )
    monkeypatch.setattr(sam_worker, "_open_batch_result_parent", lambda binding: 90)
    monkeypatch.setattr(sam_worker, "_close_batch_result_parent", lambda handle: None)
    monkeypatch.setattr(sam_worker.os, "open", lambda *args, **kwargs: 91)
    monkeypatch.setattr(sam_worker.os, "close", lambda descriptor: None)
    monkeypatch.setattr(sam_worker, "_publish_batch_result", lambda *args: None)
    monkeypatch.setattr(
        sam_worker,
        "_unlink_linux_batch_capability_result",
        lambda *args: (_ for _ in ()).throw(RuntimeError("cleanup failed")),
        raising=False,
    )

    with pytest.raises(RuntimeError, match="cleanup failed"):
        sam_worker.sam_candidate_batch_output_supported(tmp_path)


def test_candidate_batch_linux_probe_parent_open_failure_removes_private_dir(
    tmp_path: Path,
    monkeypatch,
) -> None:
    probe_dir = tmp_path / "private-probe"
    probe_dir.mkdir(mode=0o700)
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    monkeypatch.setattr(
        sam_worker,
        "_create_linux_batch_capability_directory",
        lambda root: (probe_dir, (probe_dir.stat().st_dev, probe_dir.stat().st_ino)),
    )
    monkeypatch.setattr(
        sam_worker,
        "_open_batch_result_parent",
        lambda binding: (_ for _ in ()).throw(RuntimeError("open failed")),
    )

    with pytest.raises(RuntimeError, match="open failed"):
        sam_worker.sam_candidate_batch_output_supported(tmp_path)

    assert not list(tmp_path.iterdir())


def test_candidate_batch_linux_probe_target_occupancy_is_not_overwritten(
    tmp_path: Path,
    monkeypatch,
) -> None:
    probe_dir = tmp_path / "private-probe"
    probe_dir.mkdir(mode=0o700)
    intruder = b"occupied"
    occupied = []
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    monkeypatch.setattr(
        sam_worker,
        "_create_linux_batch_capability_directory",
        lambda root: (probe_dir, (probe_dir.stat().st_dev, probe_dir.stat().st_ino)),
    )
    monkeypatch.setattr(sam_worker, "_open_batch_result_parent", lambda binding: 90)
    monkeypatch.setattr(sam_worker, "_close_batch_result_parent", lambda handle: None)
    monkeypatch.setattr(sam_worker.os, "open", lambda *args, **kwargs: 91)
    monkeypatch.setattr(sam_worker.os, "close", lambda descriptor: None)

    def occupy_target(descriptor, parent, binding):
        binding["path"].write_bytes(intruder)
        occupied.append(binding["path"])
        raise RuntimeError("already exists")

    monkeypatch.setattr(sam_worker, "_publish_batch_result", occupy_target)

    with pytest.raises(RuntimeError, match="already exists"):
        sam_worker.sam_candidate_batch_output_supported(tmp_path)

    assert occupied[0].read_bytes() == intruder
    occupied[0].unlink()
    probe_dir.rmdir()


@pytest.mark.skipif(sys.platform != "win32", reason="Windows capability probe")
def test_candidate_batch_windows_capability_probe_uses_real_publish_api(
    tmp_path: Path,
) -> None:
    assert sam_worker.sam_candidate_batch_output_supported(tmp_path) is True
    assert not list(tmp_path.iterdir())


@pytest.mark.skipif(sys.platform != "win32", reason="Windows capability probe")
def test_candidate_batch_windows_unsupported_publish_returns_false(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setattr(
        sam_worker,
        "_publish_windows_batch_result",
        lambda *args: (_ for _ in ()).throw(
            sam_worker._BatchResultPublishingUnsupported("unsupported")
        ),
    )

    assert sam_worker.sam_candidate_batch_output_supported(tmp_path) is False
    assert not list(tmp_path.iterdir())


@pytest.mark.skipif(sys.platform != "win32", reason="Windows capability probe")
def test_candidate_batch_windows_probe_retries_delete_before_close(
    tmp_path: Path,
    monkeypatch,
) -> None:
    original_delete = sam_worker._delete_windows_batch_result
    delete_calls = []

    def fail_first_delete(file_descriptor):
        delete_calls.append(file_descriptor)
        if len(delete_calls) == 1:
            raise RuntimeError("first delete failed")
        return original_delete(file_descriptor)

    monkeypatch.setattr(
        sam_worker,
        "_delete_windows_batch_result",
        fail_first_delete,
    )

    with pytest.raises(RuntimeError, match="first delete failed"):
        sam_worker.sam_candidate_batch_output_supported(tmp_path)

    assert len(delete_calls) == 2
    assert not list(tmp_path.iterdir())


@pytest.mark.skipif(sys.platform != "win32", reason="Windows capability probe")
def test_candidate_batch_windows_unsupported_does_not_hide_close_failure(
    tmp_path: Path,
    monkeypatch,
) -> None:
    original_close = sam_worker.os.close
    monkeypatch.setattr(
        sam_worker,
        "_publish_windows_batch_result",
        lambda *args: (_ for _ in ()).throw(
            sam_worker._BatchResultPublishingUnsupported("unsupported")
        ),
    )

    def close_then_fail(descriptor):
        original_close(descriptor)
        raise RuntimeError("close failed")

    monkeypatch.setattr(sam_worker.os, "close", close_then_fail)

    with pytest.raises(RuntimeError, match="close failed"):
        sam_worker.sam_candidate_batch_output_supported(tmp_path)

    assert not list(tmp_path.iterdir())


@pytest.mark.skipif(sys.platform != "win32", reason="Windows capability probe")
def test_candidate_batch_windows_close_failure_does_not_hide_primary_error(
    tmp_path: Path,
    monkeypatch,
) -> None:
    original_close = sam_worker.os.close
    monkeypatch.setattr(
        sam_worker,
        "_publish_windows_batch_result",
        lambda *args: (_ for _ in ()).throw(RuntimeError("publish failed")),
    )

    def close_then_fail(descriptor):
        original_close(descriptor)
        raise RuntimeError("close failed")

    monkeypatch.setattr(sam_worker.os, "close", close_then_fail)

    with pytest.raises(RuntimeError, match="publish failed") as error:
        sam_worker.sam_candidate_batch_output_supported(tmp_path)

    assert any("close failed" in note for note in error.value.__notes__)
    assert not list(tmp_path.iterdir())


def test_candidate_batch_linux_close_failure_does_not_hide_primary_error(
    tmp_path: Path,
    monkeypatch,
) -> None:
    monkeypatch.setattr(sam_worker.sys, "platform", "linux")
    monkeypatch.setattr(sam_worker.os, "O_TMPFILE", 0x410000, raising=False)
    _mock_linux_candidate_batch_probe_directory(tmp_path, monkeypatch)
    monkeypatch.setattr(sam_worker, "_open_batch_result_parent", lambda binding: 90)
    monkeypatch.setattr(
        sam_worker,
        "_close_batch_result_parent",
        lambda handle: (_ for _ in ()).throw(RuntimeError("close failed")),
    )
    monkeypatch.setattr(
        sam_worker.os,
        "open",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            PermissionError(errno.EACCES, "denied")
        ),
    )

    with pytest.raises(PermissionError, match="denied") as error:
        sam_worker.sam_candidate_batch_output_supported(tmp_path)

    assert any("close failed" in note for note in error.value.__notes__)


@pytest.mark.parametrize("batch_supported", [True, False])
def test_candidate_batch_process_routes_initial_and_residual_stage(
    tmp_path: Path,
    monkeypatch,
    batch_supported: bool,
) -> None:
    image = image_to_ppt.np.full((10, 20, 3), 40, dtype=image_to_ppt.np.uint8)
    image_path = tmp_path / "source.png"
    Image.fromarray(image).save(image_path)
    work_dir = tmp_path / "work"
    work_dir.mkdir()
    text_mask_path = work_dir / "source-text-mask.png"
    Image.fromarray(image_to_ppt.np.zeros((10, 20), dtype=image_to_ppt.np.uint8)).save(
        text_mask_path
    )
    mask = image_to_ppt.np.zeros((10, 20), dtype=bool)
    mask[2:8, 4:16] = True
    candidate = image_to_ppt.MaskCandidate(mask.copy(), 0.95, "sam")
    element = visual_segment.VisualElement(mask.copy(), 0, 0.95, "sam")
    batch_calls = []
    legacy_calls = []

    monkeypatch.setattr(
        image_to_ppt,
        "sam_candidate_batch_output_supported",
        lambda target: batch_supported,
    )

    monkeypatch.setattr(
        image_to_ppt,
        "_generate_filtered_object_proposals_isolated",
        lambda *args: [],
    )

    def fake_batch(actual_image, actual_text_mask, proposals, target):
        batch_calls.append((actual_image.copy(), actual_text_mask.copy(), proposals, target))
        return ([candidate], []) if len(batch_calls) == 1 else ([], [])

    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidate_batch_isolated",
        fake_batch,
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_sam_candidates_isolated",
        lambda actual_image, actual_mask, proposals, target, *, mode: (
            legacy_calls.append(
                (actual_image.copy(), actual_mask, proposals, target, mode)
            )
            or (
                [candidate]
                if mode == "prompted"
                and sum(call[4] == "prompted" for call in legacy_calls) == 1
                else []
            )
        ),
    )
    monkeypatch.setattr(image_to_ppt, "filter_prompt_free_candidates", lambda *args: [])
    monkeypatch.setattr(image_to_ppt, "generate_geometry_candidates", lambda *args: [])
    monkeypatch.setattr(
        image_to_ppt,
        "combine_residual_candidates",
        lambda **kwargs: ([], 0),
    )
    monkeypatch.setattr(image_to_ppt, "resolve_visual_elements", lambda candidates: [element])
    monkeypatch.setattr(image_to_ppt, "validate_visual_masks", lambda masks: None)
    monkeypatch.setattr(
        image_to_ppt,
        "build_clean_background",
        lambda actual_image, *args, **kwargs: actual_image.copy(),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_recheck_visual_element_holes_isolated",
        lambda *args: None,
    )
    monkeypatch.setattr(image_to_ppt, "_release_visual_resources", lambda: None)
    monkeypatch.setattr(image_to_ppt, "export_visual_components", lambda *args, **kwargs: [])
    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        lambda background, **kwargs: (background, 0, 0, "identity"),
    )

    image_to_ppt._process_image(
        image_path,
        work_dir,
        object_detector=None,
        mask_generator=None,
        lang="en",
        text_analysis={"items": [], "mask_path": str(text_mask_path)},
        defer_quality=True,
        _resource_isolation=True,
    )

    if batch_supported:
        assert len(batch_calls) == 2
        assert legacy_calls == []
        assert image_to_ppt.np.array_equal(batch_calls[0][0], image)
        assert image_to_ppt.np.array_equal(batch_calls[1][0], image)
        assert all(call[2] == [] and call[3] == work_dir for call in batch_calls)
    else:
        assert batch_calls == []
        assert [call[4] for call in legacy_calls] == [
            "prompted",
            "automatic",
            "prompted",
            "automatic",
        ]
        assert all(
            call[2] == [] and call[3] == work_dir
            for call in (legacy_calls[0], legacy_calls[2])
        )
        assert all(
            call[1] is None and call[2] is None and call[3] == work_dir
            for call in (legacy_calls[1], legacy_calls[3])
        )
        assert all(image_to_ppt.np.array_equal(call[0], image) for call in legacy_calls)


def test_sam_rle_mask_is_decoded_without_full_frame_copy(monkeypatch) -> None:
    import numpy as np

    decoded = np.zeros((10, 10), dtype=bool)
    decoded[2:5, 3:7] = True
    monkeypatch.setattr(
        visual_segment.importlib,
        "import_module",
        lambda name: types.SimpleNamespace(
            rle_to_mask=lambda rle: decoded,
        ),
    )

    class FakeGenerator:
        def generate(self, image):
            return [{
                "segmentation": {"size": [10, 10], "counts": [100]},
                "predicted_iou": 0.95,
                "stability_score": 0.96,
            }]

    candidates = visual_segment.generate_mask_candidates(
        np.zeros((10, 10, 3), dtype=np.uint8),
        FakeGenerator(),
        include_geometry=False,
    )

    assert len(candidates) == 1
    assert candidates[0].mask is decoded


def test_low_score_sam_rle_is_skipped_before_decode(monkeypatch) -> None:
    import numpy as np

    def fail_decode(rle):
        raise AssertionError("low-score RLE must not be decoded")

    monkeypatch.setattr(
        visual_segment.importlib,
        "import_module",
        lambda name: types.SimpleNamespace(rle_to_mask=fail_decode),
    )

    class FakeGenerator:
        def generate(self, image):
            return [{
                "segmentation": {"size": [10, 10], "counts": [100]},
                "predicted_iou": 0.89,
                "stability_score": 0.96,
            }]

    candidates = visual_segment.generate_mask_candidates(
        np.zeros((10, 10, 3), dtype=np.uint8),
        FakeGenerator(),
        include_geometry=False,
        min_score=0.90,
    )

    assert candidates == []


def test_prompt_free_filter_reuses_candidate_storage() -> None:
    import numpy as np

    mask = np.zeros((20, 20), dtype=bool)
    mask[2:12, 3:13] = True
    candidate = visual_segment.MaskCandidate(
        mask=mask,
        score=0.95,
        source="sam",
    )

    retained = visual_segment.filter_prompt_free_candidates(
        [candidate],
        [],
        np.zeros((20, 20), dtype=np.uint8),
    )

    assert len(retained) == 1
    assert retained[0] is candidate


def test_combine_residual_candidates_keeps_prompt_free_object() -> None:
    import numpy as np

    source = np.full((80, 120, 3), 240, dtype=np.uint8)
    right = np.zeros((80, 120), dtype=bool)
    right[24:56, 72:104] = True
    source[right] = 20
    clean_background = np.full_like(source, 240)
    clean_background[right] = source[right]
    automatic = [visual_segment.MaskCandidate(right, 0.96, "sam")]

    residual, attached = visual_segment.combine_residual_candidates(
        source=source,
        clean_background=clean_background,
        prompted=[],
        prompt_free=automatic,
        existing=[],
        text_mask=np.zeros(right.shape, dtype=np.uint8),
    )

    assert attached == 0
    assert len(residual) == 1
    assert np.array_equal(residual[0].mask, right)


def test_combine_residual_candidates_keeps_geometry_when_sam_is_empty() -> None:
    import numpy as np

    source = np.full((80, 120, 3), 240, dtype=np.uint8)
    shape = np.zeros((80, 120), dtype=bool)
    shape[20:60, 30:90] = True
    source[shape] = 20
    geometry = visual_segment.MaskCandidate(shape, 0.70, "geometry")

    residual, attached = visual_segment.combine_residual_candidates(
        source=source,
        clean_background=source.copy(),
        prompted=[],
        prompt_free=[geometry],
        existing=[],
        text_mask=np.zeros(shape.shape, dtype=np.uint8),
    )

    assert attached == 0
    assert len(residual) == 1
    assert residual[0] is geometry


def test_resolve_visual_elements_reuses_candidate_as_semantic_support() -> None:
    import numpy as np

    mask = np.zeros((20, 20), dtype=bool)
    mask[2:12, 3:13] = True
    candidate = visual_segment.MaskCandidate(mask, 0.95, "sam")

    elements = visual_segment.resolve_visual_elements([candidate])

    assert len(elements) == 1
    assert elements[0].semantic_mask is candidate.mask


def test_reference_option_is_explicit() -> None:
    assert _parse_reference_option(reference=False, no_reference=False) is False
    assert _parse_reference_option(reference=True, no_reference=False) is True
    assert _parse_reference_option(reference=True, no_reference=True) is False


def test_huge_foreground_hint_is_rejected_for_background_refinement() -> None:
    assert not _should_use_fg_hint(nonzero_pixels=600, total_pixels=1000)
    assert _should_use_fg_hint(nonzero_pixels=200, total_pixels=1000)


def test_huge_detector_mask_is_rejected() -> None:
    assert not _keep_detector_mask(nonzero_pixels=600, total_pixels=1000)
    assert _keep_detector_mask(nonzero_pixels=200, total_pixels=1000)


def test_combined_foreground_mask_is_rejected_after_union() -> None:
    import numpy as np

    mask = np.zeros((10, 10), dtype=bool)
    mask[:, :6] = True
    fallback = np.zeros((10, 10), dtype=bool)
    fallback[0, 0] = True

    limited = _limit_combined_mask(mask, fallback)

    assert int(np.count_nonzero(limited)) == 1


def test_oversized_edge_fallback_is_rejected() -> None:
    import numpy as np

    mask = np.ones((10, 10), dtype=bool)
    fallback = np.zeros((10, 10), dtype=bool)

    limited = _limit_combined_mask(mask, fallback)

    assert int(np.count_nonzero(limited)) == 0


def test_merge_foreground_masks_restores_lost_compact_component() -> None:
    import numpy as np

    initial = np.zeros((100, 120), dtype=np.uint8)
    initial[30:50, 40:65] = 255
    refined = np.zeros((100, 120), dtype=np.uint8)

    merged = _merge_foreground_masks(initial, refined)

    assert np.count_nonzero(merged[30:50, 40:65]) == 20 * 25


def test_merge_foreground_masks_does_not_restore_sparse_full_slide_shell() -> None:
    import numpy as np

    initial = np.zeros((100, 120), dtype=np.uint8)
    initial[5:95, 5:8] = 255
    initial[5:95, 112:115] = 255
    initial[5:8, 5:115] = 255
    initial[92:95, 5:115] = 255
    refined = np.zeros((100, 120), dtype=np.uint8)

    merged = _merge_foreground_masks(initial, refined)

    assert np.count_nonzero(merged) == 0


def test_sparse_full_slide_layout_shell_is_rejected() -> None:
    import cv2
    import numpy as np

    img = np.full((100, 180, 3), 245, dtype=np.uint8)
    bg = img.copy()
    text_mask = np.zeros((100, 180), dtype=np.uint8)

    img[5:95, 5:8] = [150, 110, 60]
    img[5:95, 172:175] = [150, 110, 60]
    img[5:8, 5:175] = [150, 110, 60]
    img[92:95, 5:175] = [150, 110, 60]
    for y in (25, 40, 55, 70):
        cv2.line(img, (5, y), (175, y), (150, 110, 60), 2)

    mask = extract_foreground_mask(img, bg, text_mask, diff_threshold=20)

    assert np.count_nonzero(mask) == 0


def test_text_mask_does_not_cut_holes_in_underlying_graphic() -> None:
    import numpy as np

    img = np.full((120, 160, 3), 255, dtype=np.uint8)
    img[30:90, 40:120] = [220, 20, 20]

    text_mask = np.zeros((120, 160), dtype=np.uint8)
    text_mask[50:70, 55:105] = 255

    bg = build_background(img, text_mask=text_mask)
    fg_mask = extract_foreground_mask(img, bg, text_mask, diff_threshold=20)

    covered_graphic = fg_mask[50:70, 55:105]

    assert np.count_nonzero(covered_graphic) >= covered_graphic.size * 0.9


def test_hole_recheck_uses_sam_inference_context(monkeypatch) -> None:
    from contextlib import contextmanager

    import numpy as np

    mask = np.zeros((24, 24), dtype=bool)
    mask[2:22, 2:22] = True
    mask[8:16, 8:16] = False
    filled = np.zeros_like(mask)
    filled[2:22, 2:22] = True
    events = []

    class Predictor:
        def set_image(self, image):
            events.append("set-image")

        def predict(self, **kwargs):
            events.append("predict")
            return np.asarray([filled]), np.asarray([1.0]), None

    generator = types.SimpleNamespace(
        predictor=Predictor(),
        _image2editable_device="cuda",
    )
    element = visual_segment.VisualElement(
        mask=mask,
        z_index=0,
        score=0.95,
        source="sam",
        object_box=(2.0, 2.0, 22.0, 22.0),
    )

    @contextmanager
    def recording_context(actual_generator):
        assert actual_generator is generator
        events.append("enter")
        yield
        events.append("exit")

    monkeypatch.setattr(
        visual_segment,
        "_sam_inference_context",
        recording_context,
    )

    visual_segment.recheck_visual_element_holes(
        np.zeros((24, 24, 3), dtype=np.uint8),
        [element],
        generator,
    )

    assert events == [
        "enter",
        "set-image",
        "exit",
        "enter",
        "predict",
        "exit",
    ]


def test_thin_foreground_line_is_preserved() -> None:
    import numpy as np

    img = np.full((80, 120, 3), 255, dtype=np.uint8)
    img[40:41, 20:100] = [0, 0, 0]
    text_mask = np.zeros((80, 120), dtype=np.uint8)

    bg = build_background(img, text_mask=text_mask)
    fg_mask = extract_foreground_mask(img, bg, text_mask, diff_threshold=20)

    assert np.count_nonzero(fg_mask[40:41, 20:100]) >= 70


def test_connected_shape_is_not_split_into_separate_components(tmp_path: Path) -> None:
    import cv2
    import numpy as np

    img = np.full((200, 300, 3), 255, dtype=np.uint8)
    mask = np.zeros((200, 300), dtype=np.uint8)
    cv2.circle(mask, (65, 100), 20, 255, -1)
    cv2.circle(mask, (135, 100), 20, 255, -1)
    mask[98:102, 65:135] = 255
    img[mask > 0] = [0, 120, 200]

    components = split_components(img, mask, tmp_path, min_area=20)

    assert len(components) == 1


def test_inpaint_mask_does_not_expand_foreground_repair_area() -> None:
    import numpy as np

    exclude_mask = np.zeros((20, 20), dtype=np.uint8)
    exclude_mask[8:12, 8:12] = 255

    inpaint_mask = _build_inpaint_mask(exclude_mask)

    assert np.array_equal(inpaint_mask, exclude_mask)


def test_assemble_pptx_defaults_to_no_reference_slide(tmp_path: Path) -> None:
    bg_path = tmp_path / "bg.png"
    Image.new("RGB", (20, 10), "white").save(bg_path)

    out_path = tmp_path / "out.pptx"
    assemble_pptx(
        background_path=bg_path,
        components=[],
        text_items=[],
        img_width=20,
        img_height=10,
        output_path=out_path,
        original_image_path=bg_path,
    )

    prs = Presentation(out_path)

    assert len(prs.slides) == 1


def test_centered_textbox_keeps_detected_horizontal_bounds(tmp_path: Path) -> None:
    bg_path = tmp_path / "bg.png"
    Image.new("RGB", (100, 50), "white").save(bg_path)
    out_path = tmp_path / "out.pptx"

    assemble_pptx(
        background_path=bg_path,
        components=[],
        text_items=[{
            "box": [20, 10, 40, 20],
            "text": "Centered",
            "font_size": 18,
            "color": "#000000",
            "bold": False,
            "font": "Arial",
            "align": 1,
        }],
        img_width=100,
        img_height=50,
        output_path=out_path,
    )

    prs = Presentation(out_path)
    text_box = prs.slides[0].shapes[1]

    assert abs(text_box.left / 914400 - 2.6666) < 0.01
    assert abs(text_box.width / 914400 - 5.3332) < 0.01


def test_widescreen_canvas_places_component_without_aspect_change(
    tmp_path: Path,
) -> None:
    bg_path = tmp_path / "bg.png"
    component_path = tmp_path / "component.png"
    Image.new("RGB", (368, 207), "white").save(bg_path)
    Image.new("RGBA", (20, 40), "red").save(component_path)
    out_path = tmp_path / "out.pptx"

    assemble_pptx(
        background_path=bg_path,
        components=[{
            "path": str(component_path),
            "x": 10,
            "y": 20,
            "w": 20,
            "h": 40,
        }],
        text_items=[],
        img_width=100,
        img_height=200,
        output_path=out_path,
        canvas_width=368,
        canvas_height=207,
        content_offset_x=134,
        content_offset_y=3,
    )

    component = Presentation(out_path).slides[0].shapes[1]
    assert abs(component.width / component.height - 0.5) < 1e-4
    assert abs(component.left / 914400 - 144 / 368 * 13.333333) < 0.01
    assert abs(component.top / 914400 - 23 / 207 * 7.5) < 0.01


def test_widescreen_canvas_rejects_invalid_geometry(tmp_path: Path) -> None:
    import pytest

    bg_path = tmp_path / "bg.png"
    Image.new("RGB", (100, 100), "white").save(bg_path)
    base_args = {
        "background_path": bg_path,
        "components": [],
        "text_items": [],
        "img_width": 100,
        "img_height": 100,
        "output_path": tmp_path / "out.pptx",
    }

    for canvas_args in (
        {"canvas_width": 368},
        {"canvas_height": 207},
        {"canvas_width": 0, "canvas_height": 0},
        {"canvas_width": 368.0, "canvas_height": 207.0},
        {"canvas_width": 367, "canvas_height": 207},
        {"canvas_width": 160, "canvas_height": 90, "content_offset_x": 61},
        {"canvas_width": 160, "canvas_height": 90, "content_offset_x": -1},
        {"canvas_width": 160, "canvas_height": 90, "content_offset_x": 1.0},
        {"canvas_width": True, "canvas_height": 90},
        {"canvas_width": 160, "canvas_height": 90, "content_offset_y": True},
        {"content_offset_x": 1},
    ):
        with pytest.raises(ValueError):
            assemble_pptx(**base_args, **canvas_args)


def test_widescreen_canvas_accepts_numpy_integral_geometry(tmp_path: Path) -> None:
    import numpy as np

    bg_path = tmp_path / "bg.png"
    Image.new("RGB", (160, 90), "white").save(bg_path)

    assemble_pptx(
        background_path=bg_path,
        components=[],
        text_items=[],
        img_width=100,
        img_height=80,
        output_path=tmp_path / "out.pptx",
        canvas_width=np.int64(160),
        canvas_height=np.int64(90),
        content_offset_x=np.int64(30),
        content_offset_y=np.int64(5),
    )


def test_multi_slide_rejects_invalid_canvas_geometry(tmp_path: Path) -> None:
    import pytest

    bg_path = tmp_path / "bg.png"
    Image.new("RGB", (160, 90), "white").save(bg_path)
    base_data = {
        "background_path": str(bg_path),
        "components": [],
        "text_items": [],
        "img_width": 100,
        "img_height": 80,
    }

    for canvas_args in (
        {"canvas_width": 160.0, "canvas_height": 90.0},
        {"canvas_height": 90},
    ):
        with pytest.raises(ValueError):
            assemble_pptx_multi(
                [{**base_data, **canvas_args}],
                tmp_path / "out.pptx",
            )


def test_widescreen_canvas_is_exact_ratio_and_contains_source() -> None:
    for source_width, source_height in (
        (1600, 900),
        (1122, 1402),
        (2400, 900),
        (37, 901),
    ):
        canvas_width, canvas_height, offset_x, offset_y = (
            bg_model.compute_widescreen_canvas(source_width, source_height)
        )

        assert all(
            isinstance(value, Integral)
            for value in (canvas_width, canvas_height, offset_x, offset_y)
        )
        assert canvas_width * 9 == canvas_height * 16
        assert offset_x >= 0
        assert offset_y >= 0
        assert offset_x + source_width <= canvas_width
        assert offset_y + source_height <= canvas_height
        assert abs(offset_x - (canvas_width - offset_x - source_width)) <= 1
        assert abs(offset_y - (canvas_height - offset_y - source_height)) <= 1


def test_exact_widescreen_canvas_is_identity() -> None:
    assert bg_model.compute_widescreen_canvas(1920, 1080) == (1920, 1080, 0, 0)


def test_widescreen_background_preserves_complete_source_region() -> None:
    import numpy as np

    from scripts.lama_inpaint import LargeMaskInpaintError

    source = np.arange(20 * 30 * 3, dtype=np.uint8).reshape(20, 30, 3)
    was_called = False

    def fail_inpaint(*_args):
        nonlocal was_called
        was_called = True
        raise LargeMaskInpaintError("expected test fallback")

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source,
        large_inpainter=fail_inpaint,
    )

    height, width = background.shape[:2]
    assert width * 9 == height * 16
    assert np.array_equal(
        background[offset_y:offset_y + 20, offset_x:offset_x + 30], source
    )
    assert was_called
    assert method == "ambient"


def test_widescreen_ambient_preserves_first_extension_ring() -> None:
    import numpy as np

    from scripts.lama_inpaint import LargeMaskInpaintError

    source = np.full((20, 30, 3), 100, dtype=np.uint8)

    def fail_inpaint(*_args):
        raise LargeMaskInpaintError("expected ambient fallback")

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source,
        large_inpainter=fail_inpaint,
    )
    x2 = offset_x + source.shape[1]
    y2 = offset_y + source.shape[0]

    assert method == "ambient"
    assert np.array_equal(background[offset_y - 1, offset_x:x2], source[0])
    assert np.array_equal(background[y2, offset_x:x2], source[-1])
    assert np.array_equal(background[offset_y:y2, offset_x - 1], source[:, 0])
    assert np.array_equal(background[offset_y:y2, x2], source[:, -1])


def test_widescreen_ambient_does_not_repeat_reflected_high_frequency() -> None:
    import cv2
    import numpy as np

    yy, xx = np.indices((113, 63))
    checker = ((xx + yy) % 2 * 255).astype(np.uint8)
    source = np.stack((checker, 255 - checker, checker), axis=2)
    canvas_width, canvas_height, offset_x, offset_y = (
        bg_model.compute_widescreen_canvas(source.shape[1], source.shape[0])
    )

    background = bg_model._build_ambient_backdrop(
        source, canvas_width, canvas_height, offset_x, offset_y
    )
    legacy_reflection = cv2.copyMakeBorder(
        source,
        offset_y,
        canvas_height - offset_y - source.shape[0],
        offset_x,
        canvas_width - offset_x - source.shape[1],
        cv2.BORDER_REFLECT,
    )
    left = background[offset_y:offset_y + source.shape[0], :offset_x - 4]
    legacy_left = legacy_reflection[
        offset_y:offset_y + source.shape[0], :offset_x - 4
    ]
    source_gray = cv2.cvtColor(source, cv2.COLOR_RGB2GRAY)
    left_gray = cv2.cvtColor(left, cv2.COLOR_RGB2GRAY)
    legacy_gray = cv2.cvtColor(legacy_left, cv2.COLOR_RGB2GRAY)
    correlation = np.corrcoef(left_gray.ravel(), legacy_gray.ravel())[0, 1]

    assert correlation < 0.6
    assert cv2.Laplacian(left_gray, cv2.CV_32F).var() < (
        cv2.Laplacian(source_gray, cv2.CV_32F).var() * 0.15
    )


def test_widescreen_ambient_is_deterministic_contiguous_and_lossless() -> None:
    import numpy as np

    source = np.random.default_rng(12).integers(
        0, 256, size=(79, 37, 3), dtype=np.uint8
    )
    canvas_width, canvas_height, offset_x, offset_y = (
        bg_model.compute_widescreen_canvas(source.shape[1], source.shape[0])
    )

    first = bg_model._build_ambient_backdrop(
        source, canvas_width, canvas_height, offset_x, offset_y
    )
    second = bg_model._build_ambient_backdrop(
        source, canvas_width, canvas_height, offset_x, offset_y
    )

    assert first.dtype == np.uint8
    assert first.flags.c_contiguous
    assert np.array_equal(first, second)
    assert np.array_equal(
        first[offset_y:offset_y + source.shape[0],
              offset_x:offset_x + source.shape[1]],
        source,
    )


def test_widescreen_ambient_handles_portrait_wide_odd_and_tiny_sources() -> None:
    import numpy as np

    rng = np.random.default_rng(21)
    for height, width in ((23, 11), (11, 23), (19, 35), (3, 7), (7, 3)):
        source = rng.integers(0, 256, size=(height, width, 3), dtype=np.uint8)
        canvas_width, canvas_height, offset_x, offset_y = (
            bg_model.compute_widescreen_canvas(width, height)
        )
        background = bg_model._build_ambient_backdrop(
            source, canvas_width, canvas_height, offset_x, offset_y
        )
        x2 = offset_x + width
        y2 = offset_y + height

        assert background.shape == (canvas_height, canvas_width, 3)
        assert np.array_equal(background[offset_y:y2, offset_x:x2], source)
        if offset_x:
            assert np.array_equal(background[offset_y:y2, offset_x - 1], source[:, 0])
        if x2 < canvas_width:
            assert np.array_equal(background[offset_y:y2, x2], source[:, -1])
        if offset_y:
            assert np.array_equal(background[offset_y - 1, offset_x:x2], source[0])
        if y2 < canvas_height:
            assert np.array_equal(background[y2, offset_x:x2], source[-1])


def test_widescreen_ambient_tiny_cover_keeps_uniform_geometry() -> None:
    import numpy as np

    source = np.zeros((3, 7, 3), dtype=np.uint8)
    source[:, 2:5] = 255

    cover = bg_model._resize_cover(source, 160, 90)
    bright_y, bright_x = np.where(cover[:, :, 0] >= 128)

    bright_width = bright_x.max() - bright_x.min() + 1
    bright_height = bright_y.max() - bright_y.min() + 1
    assert abs(bright_width - bright_height) <= 1


def test_widescreen_ambient_cover_aligns_pixel_centers() -> None:
    import numpy as np

    source = np.zeros((5, 9, 3), dtype=np.uint8)
    source[1:4, 3:6] = 255

    cover = bg_model._resize_cover(source, 160, 90)
    bright_y, bright_x = np.where(cover[:, :, 0] >= 128)
    bright_center_x = (bright_x.min() + bright_x.max()) / 2
    bright_center_y = (bright_y.min() + bright_y.max()) / 2

    assert abs(bright_center_x - (cover.shape[1] - 1) / 2) <= 1
    assert abs(bright_center_y - (cover.shape[0] - 1) / 2) <= 1


def test_widescreen_ambient_gradient_seam_is_low_frequency() -> None:
    import cv2
    import numpy as np

    y = np.linspace(30, 220, 91, dtype=np.uint8)[:, None]
    source = np.stack(
        (
            np.broadcast_to(y, (91, 41)),
            np.broadcast_to(255 - y, (91, 41)),
            np.full((91, 41), 96, dtype=np.uint8),
        ),
        axis=2,
    )
    canvas_width, canvas_height, offset_x, offset_y = (
        bg_model.compute_widescreen_canvas(source.shape[1], source.shape[0])
    )
    background = bg_model._build_ambient_backdrop(
        source, canvas_width, canvas_height, offset_x, offset_y
    )
    adjacent = background[
        offset_y:offset_y + source.shape[0], max(0, offset_x - 2)
    ]
    low_source_edge = cv2.GaussianBlur(
        source[:, :1], (0, 0), sigmaX=1, sigmaY=3
    )[:, 0]

    assert np.abs(adjacent.astype(np.int16) - low_source_edge.astype(np.int16)).mean() < 32
    adjacent_gray = cv2.cvtColor(adjacent[:, None, :], cv2.COLOR_RGB2GRAY)[:, 0]
    assert cv2.Laplacian(adjacent_gray, cv2.CV_32F).var() < 20


def test_widescreen_ambient_keeps_multicolor_extension_decorative() -> None:
    import numpy as np

    source = np.zeros((81, 39, 3), dtype=np.uint8)
    source[:40, :19] = (240, 30, 30)
    source[:40, 19:] = (30, 220, 60)
    source[40:, :19] = (30, 60, 230)
    source[40:, 19:] = (230, 210, 30)
    canvas_width, canvas_height, offset_x, offset_y = (
        bg_model.compute_widescreen_canvas(source.shape[1], source.shape[0])
    )
    background = bg_model._build_ambient_backdrop(
        source, canvas_width, canvas_height, offset_x, offset_y
    )
    extension_mask = np.ones(background.shape[:2], dtype=bool)
    extension_mask[
        offset_y:offset_y + source.shape[0],
        offset_x:offset_x + source.shape[1],
    ] = False

    assert np.unique(background[extension_mask], axis=0).shape[0] > 8


def test_widescreen_background_rejects_bad_outpaint() -> None:
    import numpy as np

    x = np.arange(80, dtype=np.uint8)[None, :]
    y = np.arange(40, dtype=np.uint8)[:, None]
    source = np.stack((np.broadcast_to(x, (40, 80)),
                       np.broadcast_to(y, (40, 80)),
                       np.broadcast_to(x + y, (40, 80))), axis=2)

    def black_outpaint(image, mask):
        result = image.copy()
        result[mask > 0] = 0
        return result

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source,
        large_inpainter=black_outpaint,
    )
    new_region = np.ones(background.shape[:2], dtype=bool)
    new_region[offset_y:offset_y + 40, offset_x:offset_x + 80] = False

    assert method == "ambient"
    assert np.array_equal(
        background[offset_y:offset_y + 40, offset_x:offset_x + 80], source
    )
    assert background[new_region].size > 0
    assert np.var(background[new_region]) > 0


def test_widescreen_background_rejects_invalid_inpainter_arrays() -> None:
    import numpy as np

    source = np.full((80, 90, 3), 64, dtype=np.uint8)
    results = []
    for dtype, value in ((np.float32, np.nan), (np.int32, 64)):
        def invalid_outpaint(image, _mask, *, dtype=dtype, value=value):
            return np.full(image.shape, value, dtype=dtype)

        results.append(bg_model.build_widescreen_background(
            source,
            large_inpainter=invalid_outpaint,
        ))

    for background, offset_x, offset_y, method in results:
        assert method == "ambient"
        assert background.dtype == np.uint8
        assert np.isfinite(background).all()
        assert np.array_equal(
            background[offset_y:offset_y + 80, offset_x:offset_x + 90], source
        )


def test_widescreen_background_normalizes_valid_inpainter_memory_order() -> None:
    import numpy as np

    source = np.full((80, 90, 3), 64, dtype=np.uint8)

    def fortran_outpaint(image, mask):
        result = image.copy()
        result[mask > 0] = 64
        return np.asfortranarray(result)

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source,
        large_inpainter=fortran_outpaint,
    )

    assert method == "outpaint"
    assert background.flags.c_contiguous
    assert np.array_equal(
        background[offset_y:offset_y + 80, offset_x:offset_x + 90], source
    )


def test_widescreen_background_rejects_bad_intermediate_outpaint_seam() -> None:
    import numpy as np

    source = np.full((80, 90, 3), 64, dtype=np.uint8)
    calls = 0

    def broken_staged_outpaint(image, mask):
        nonlocal calls
        calls += 1
        result = image.copy()
        result[mask > 0] = 64 if calls == 1 else 128
        return result

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source,
        large_inpainter=broken_staged_outpaint,
    )

    assert calls >= 2
    assert method == "ambient"
    assert np.array_equal(
        background[offset_y:offset_y + 80, offset_x:offset_x + 90], source
    )


def test_widescreen_background_avoids_oversized_tiny_axis_stage() -> None:
    import numpy as np

    source = np.full((1, 16, 3), 64, dtype=np.uint8)
    was_called = False

    def record_outpaint(image, _mask):
        nonlocal was_called
        was_called = True
        return image

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source,
        large_inpainter=record_outpaint,
    )

    assert not was_called
    assert method == "ambient"
    assert np.array_equal(
        background[offset_y:offset_y + 1, offset_x:offset_x + 16], source
    )


def test_widescreen_background_accepts_quality_checked_outpaint() -> None:
    import numpy as np

    source = np.full((80, 90, 3), 64, dtype=np.uint8)

    def constant_outpaint(image, mask):
        result = image.copy()
        result[mask > 0] = 64
        return result

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source,
        large_inpainter=constant_outpaint,
    )

    assert method == "outpaint"
    assert np.array_equal(
        background[offset_y:offset_y + 80, offset_x:offset_x + 90], source
    )


def test_widescreen_background_stages_use_centered_border_masks() -> None:
    import numpy as np

    source = np.full((80, 90, 3), 64, dtype=np.uint8)
    calls = []

    def record_outpaint(image, mask):
        calls.append((image.shape, mask.copy()))
        result = image.copy()
        result[mask > 0] = 64
        return result

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source,
        large_inpainter=record_outpaint,
    )
    target_width, target_height, _, _ = bg_model.compute_widescreen_canvas(90, 80)
    previous_width, previous_height = 90, 80
    previous_source_x = previous_source_y = 0

    for shape, mask in calls:
        height, width = shape[:2]
        assert width - previous_width <= previous_width // 4
        assert height - previous_height <= previous_height // 4
        if previous_width == target_width:
            assert width == previous_width
        if previous_height == target_height:
            assert height == previous_height
        source_x = (width - 90) // 2
        source_y = (height - 80) // 2
        left = source_x - previous_source_x
        top = source_y - previous_source_y
        expected = np.full((height, width), 255, dtype=np.uint8)
        expected[top:top + previous_height, left:left + previous_width] = 0
        assert np.array_equal(mask, expected)
        assert np.count_nonzero(mask == 0) == previous_width * previous_height
        previous_width, previous_height = width, height
        previous_source_x, previous_source_y = source_x, source_y

    assert len(calls) >= 2
    assert (previous_width, previous_height) == (target_width, target_height)
    assert method == "outpaint"
    assert np.array_equal(
        background[offset_y:offset_y + 80, offset_x:offset_x + 90], source
    )


def test_exact_widescreen_background_returns_unchanged_pixels() -> None:
    import numpy as np

    source = np.random.default_rng(0).integers(
        0, 256, size=(90, 160, 3), dtype=np.uint8
    )

    background, offset_x, offset_y, method = bg_model.build_widescreen_background(
        source
    )

    assert np.array_equal(background, source)
    assert (offset_x, offset_y, method) == (0, 0, "identity")


def test_canvas_reference_slide_preserves_source_rect(tmp_path: Path) -> None:
    bg_path = tmp_path / "bg.png"
    original_path = tmp_path / "original.png"
    Image.new("RGB", (368, 207), "white").save(bg_path)
    Image.new("RGB", (100, 200), "red").save(original_path)
    out_path = tmp_path / "out.pptx"

    assemble_pptx(
        background_path=bg_path,
        components=[],
        text_items=[],
        img_width=100,
        img_height=200,
        output_path=out_path,
        add_reference_slide=True,
        original_image_path=original_path,
        canvas_width=368,
        canvas_height=207,
        content_offset_x=134,
        content_offset_y=3,
    )

    prs = Presentation(out_path)
    reference = prs.slides[1].shapes[1]
    assert len(prs.slides) == 2
    assert abs(reference.left / 914400 - 134 / 368 * 13.333333) < 0.01
    assert abs(reference.top / 914400 - 3 / 207 * 7.5) < 0.01
    assert abs(reference.width / 914400 - 100 / 368 * 13.333333) < 0.01
    assert abs(reference.height / 914400 - 200 / 207 * 7.5) < 0.01
    assert abs(reference.width / reference.height - 0.5) < 1e-4


def test_widescreen_canvas_places_text_with_uniform_geometry(
    tmp_path: Path,
) -> None:
    bg_path = tmp_path / "bg.png"
    Image.new("RGB", (368, 207), "white").save(bg_path)
    out_path = tmp_path / "out.pptx"

    assemble_pptx(
        background_path=bg_path,
        components=[],
        text_items=[{
            "box": [10, 20, 40, 20],
            "text": "Mapped",
            "font_size": 20,
            "color": "#000000",
            "bold": False,
            "font": "Arial",
            "align": 1,
        }],
        img_width=100,
        img_height=200,
        output_path=out_path,
        canvas_width=368,
        canvas_height=207,
        content_offset_x=134,
        content_offset_y=3,
    )

    text_box = Presentation(out_path).slides[0].shapes[1]
    x_scale = text_box.width / 40
    y_scale = text_box.height / 20
    font_size = text_box.text_frame.paragraphs[0].runs[0].font.size.pt
    assert abs(text_box.left / 914400 - 144 / 368 * 13.333333) < 0.01
    assert abs(text_box.top / 914400 - 23 / 207 * 7.5) < 0.01
    assert abs(x_scale - y_scale) < 1
    assert abs(font_size - 20 * 100 / 368) < 0.01


def test_original_slide_ignores_widescreen_canvas(tmp_path: Path) -> None:
    bg_path = tmp_path / "bg.png"
    component_path = tmp_path / "component.png"
    Image.new("RGB", (100, 100), "white").save(bg_path)
    Image.new("RGBA", (20, 20), "red").save(component_path)
    out_path = tmp_path / "out.pptx"

    assemble_pptx(
        background_path=bg_path,
        components=[{
            "path": str(component_path),
            "x": 10,
            "y": 20,
            "w": 20,
            "h": 20,
        }],
        text_items=[],
        img_width=100,
        img_height=100,
        output_path=out_path,
        slide_size="original",
        canvas_width=368,
        canvas_height=207,
        content_offset_x=134,
        content_offset_y=3,
    )

    prs = Presentation(out_path)
    component = prs.slides[0].shapes[1]
    assert abs(prs.slide_width / 914400 - 7.5) < 0.01
    assert abs(prs.slide_height / 914400 - 7.5) < 0.01
    assert abs(component.left / 914400 - 0.75) < 0.01
    assert abs(component.top / 914400 - 1.5) < 0.01
    assert abs(component.width / 914400 - 1.5) < 0.01
    assert abs(component.height / 914400 - 1.5) < 0.01


def test_multi_slide_uses_canvas_and_keeps_one_image_per_page(
    tmp_path: Path,
) -> None:
    bg_path_1 = tmp_path / "bg_1.png"
    bg_path_2 = tmp_path / "bg_2.png"
    component_path_1 = tmp_path / "component_1.png"
    component_path_2 = tmp_path / "component_2.png"
    original_path_1 = tmp_path / "original_1.png"
    original_path_2 = tmp_path / "original_2.png"
    Image.new("RGB", (368, 207), "white").save(bg_path_1)
    Image.new("RGB", (240, 135), "white").save(bg_path_2)
    Image.new("RGBA", (20, 40), "red").save(component_path_1)
    Image.new("RGBA", (60, 30), "blue").save(component_path_2)
    Image.new("RGB", (100, 200), "red").save(original_path_1)
    Image.new("RGB", (240, 90), "blue").save(original_path_2)
    out_path = tmp_path / "out.pptx"
    slide_data_1 = {
        "background_path": str(bg_path_1),
        "components": [{
            "path": str(component_path_1),
            "x": 10,
            "y": 20,
            "w": 20,
            "h": 40,
        }],
        "text_items": [{
            "box": [5, 10, 30, 20],
            "text": "First",
            "font_size": 18,
        }],
        "img_width": 100,
        "img_height": 200,
        "canvas_width": 368,
        "canvas_height": 207,
        "content_offset_x": 134,
        "content_offset_y": 3,
        "original_image_path": str(original_path_1),
    }
    slide_data_2 = {
        "background_path": str(bg_path_2),
        "components": [{
            "path": str(component_path_2),
            "x": 120,
            "y": 40,
            "w": 60,
            "h": 30,
        }],
        "text_items": [{
            "box": [20, 5, 80, 15],
            "text": "Second",
            "font_size": 24,
        }],
        "img_width": 240,
        "img_height": 90,
        "canvas_width": 240,
        "canvas_height": 135,
        "content_offset_x": 0,
        "content_offset_y": 22,
        "original_image_path": str(original_path_2),
    }

    assemble_pptx_multi(
        [slide_data_1, slide_data_2],
        out_path,
        add_reference=True,
    )

    prs = Presentation(out_path)
    component_1 = prs.slides[0].shapes[1]
    text_1 = prs.slides[0].shapes[2]
    reference_1 = prs.slides[1].shapes[1]
    component_2 = prs.slides[2].shapes[1]
    text_2 = prs.slides[2].shapes[2]
    reference_2 = prs.slides[3].shapes[1]

    assert len(prs.slides) == 4
    assert abs(component_1.left / 914400 - 144 / 368 * 13.333333) < 0.01
    assert abs(component_1.top / 914400 - 23 / 207 * 7.5) < 0.01
    assert abs(component_1.width / component_1.height - 0.5) < 1e-4
    assert abs(text_1.left / 914400 - 139 / 368 * 13.333333) < 0.01
    assert abs(text_1.top / 914400 - 13 / 207 * 7.5) < 0.01
    assert abs(text_1.text_frame.paragraphs[0].runs[0].font.size.pt - 18 * 100 / 368) < 0.01
    assert abs(reference_1.left / 914400 - 134 / 368 * 13.333333) < 0.01
    assert abs(reference_1.top / 914400 - 3 / 207 * 7.5) < 0.01
    assert abs(reference_1.width / reference_1.height - 0.5) < 1e-4

    assert abs(component_2.left / 914400 - 120 / 240 * 13.333333) < 0.01
    assert abs(component_2.top / 914400 - 62 / 135 * 7.5) < 0.01
    assert abs(component_2.width / component_2.height - 2.0) < 1e-4
    assert abs(text_2.left / 914400 - 20 / 240 * 13.333333) < 0.01
    assert abs(text_2.top / 914400 - 27 / 135 * 7.5) < 0.01
    assert abs(text_2.text_frame.paragraphs[0].runs[0].font.size.pt - 24.0) < 0.01
    assert abs(reference_2.left / 914400) < 0.01
    assert abs(reference_2.top / 914400 - 22 / 135 * 7.5) < 0.01
    assert abs(reference_2.width / reference_2.height - 240 / 90) < 1e-4

    for slide in prs.slides:
        for shape in slide.shapes:
            assert shape.left >= 0
            assert shape.top >= 0
            assert shape.left + shape.width <= prs.slide_width
            assert shape.top + shape.height <= prs.slide_height


def test_prepared_widescreen_slide_passes_canvas_geometry(
    monkeypatch,
    tmp_path: Path,
) -> None:
    captured = {}

    def fake_assemble_pptx(**kwargs):
        captured.update(kwargs)
        return str(kwargs["output_path"])

    monkeypatch.setattr(image_to_ppt, "assemble_pptx", fake_assemble_pptx)
    image_to_ppt._assemble_prepared_slide(
        {
            "background_original_path": "original.png",
            "background_widescreen_path": "widescreen.png",
            "components": [],
            "text_items": [],
            "img_width": 100,
            "img_height": 200,
            "canvas_width": 368,
            "canvas_height": 207,
            "content_offset_x": 134,
            "content_offset_y": 3,
            "original_image_path": "source.png",
        },
        tmp_path / "out.pptx",
        False,
        "16:9",
    )

    assert captured["canvas_width"] == 368
    assert captured["canvas_height"] == 207
    assert captured["content_offset_x"] == 134
    assert captured["content_offset_y"] == 3


def test_prepared_original_slide_ignores_widescreen_canvas_geometry(
    monkeypatch,
    tmp_path: Path,
) -> None:
    captured = {}

    def fake_assemble_pptx(**kwargs):
        captured.update(kwargs)
        return str(kwargs["output_path"])

    monkeypatch.setattr(image_to_ppt, "assemble_pptx", fake_assemble_pptx)
    image_to_ppt._assemble_prepared_slide(
        {
            "background_original_path": "original.png",
            "background_widescreen_path": "widescreen.png",
            "components": [],
            "text_items": [],
            "img_width": 100,
            "img_height": 200,
            "canvas_width": 368,
            "canvas_height": 207,
            "content_offset_x": 134,
            "content_offset_y": 3,
            "original_image_path": "source.png",
        },
        tmp_path / "out.pptx",
        False,
        "original",
    )

    assert captured["background_path"] == "original.png"
    assert captured["canvas_width"] is None
    assert captured["canvas_height"] is None
    assert captured["content_offset_x"] == 0
    assert captured["content_offset_y"] == 0


def test_process_image_records_widescreen_canvas_without_mutating_components(
    monkeypatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((200, 100, 3), 40, dtype=np.uint8)
    image_path = tmp_path / "source.png"
    Image.fromarray(source, mode="RGB").save(image_path)
    component_path = tmp_path / "component.png"
    Image.new("RGBA", (4, 3), "red").save(component_path)
    component = {
        "path": str(component_path),
        "x": 10,
        "y": 20,
        "w": 4,
        "h": 3,
        "area": 12,
        "z_index": 0,
    }
    expected_component = component.copy()
    empty_mask = np.zeros(source.shape[:2], dtype=np.uint8)
    export_kwargs = {}
    sam_generation_calls = []

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], empty_mask.copy()),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_generate_filtered_object_proposals",
        lambda *args, **kwargs: [],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "generate_prompted_mask_candidates",
        lambda *args: [],
    )
    def generate_masks(*args, **kwargs):
        sam_generation_calls.append(kwargs)
        return []

    monkeypatch.setattr(image_to_ppt, "generate_mask_candidates", generate_masks)
    monkeypatch.setattr(
        image_to_ppt,
        "filter_prompt_free_candidates",
        lambda *args: [],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "combine_residual_candidates",
        lambda **kwargs: ([], 0),
    )
    monkeypatch.setattr(image_to_ppt, "recheck_visual_element_holes", lambda *args: None)
    monkeypatch.setattr(
        image_to_ppt,
        "build_clean_background",
        lambda *args, **kwargs: source.copy(),
    )

    def fake_export(*args, **kwargs):
        export_kwargs.update(kwargs)
        return [component]

    canvas = np.zeros((207, 368, 3), dtype=np.uint8)
    canvas[3:203, 134:234] = source
    backgrounds = iter([
        (canvas.copy(), 134, 3, "ambient"),
        (source.copy(), 0, 0, "identity"),
    ])
    monkeypatch.setattr(image_to_ppt, "export_visual_components", fake_export)
    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        lambda background: next(backgrounds),
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "visual_difference",
        lambda *args: {"mae": 0.0, "p95": 0.0},
    )
    monkeypatch.setattr(image_to_ppt, "write_segmentation_diagnostics", lambda *args, **kwargs: None)
    monkeypatch.setattr(image_to_ppt, "require_visual_quality", lambda *args: None)

    slide_data = image_to_ppt._process_image(
        image_path,
        tmp_path / "work",
        object_detector=None,
        mask_generator=None,
        lang="en",
    )

    assert export_kwargs == {"semantic_masks": []}
    assert component == expected_component
    assert slide_data["components"] == [expected_component]
    assert slide_data["canvas_width"] == 368
    assert slide_data["canvas_height"] == 207
    assert slide_data["content_offset_x"] == 134
    assert slide_data["content_offset_y"] == 3
    assert slide_data["widescreen_background_method"] == "ambient"
    assert slide_data["background_path"] == slide_data["background_widescreen_path"]
    with Image.open(slide_data["background_widescreen_path"]) as background:
        assert background.size == (368, 207)

    identity_data = image_to_ppt._process_image(
        image_path,
        tmp_path / "identity-work",
        object_detector=None,
        mask_generator=None,
        lang="en",
    )
    assert identity_data["widescreen_background_method"] == "identity"
    assert (
        identity_data["background_widescreen_path"]
        == identity_data["background_original_path"]
    )
    assert not (tmp_path / "identity-work" / "background-16x9.png").exists()
    assert [call["include_geometry"] for call in sam_generation_calls] == [
        False, True, False, True,
    ]
    assert all(call["min_score"] == 0.90 for call in sam_generation_calls)


def test_text_only_background_removes_detected_box_without_changing_outside() -> None:
    import cv2
    import numpy as np

    img = np.full((80, 160, 3), [5, 20, 45], dtype=np.uint8)
    cv2.putText(
        img,
        "ABC",
        (35, 50),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.0,
        (255, 255, 255),
        2,
        cv2.LINE_AA,
    )
    original = img.copy()
    text_items = [{"box": [30, 20, 80, 40]}]

    assert hasattr(bg_model, "build_text_only_background")
    cleaned = bg_model.build_text_only_background(img, text_items)
    original_text = np.all(original > 220, axis=2)

    assert np.count_nonzero(np.all(cleaned[original_text] > 220, axis=1)) == 0
    assert np.array_equal(cleaned[:15], original[:15])


def test_adjacent_same_style_text_is_merged_but_different_color_is_not() -> None:
    items = [
        {"box": [10, 22, 50, 20], "text": "Kylian", "font_size": 20,
         "color": "#ffffff", "bold": True, "font": "Arial", "confidence": 0.99},
        {"box": [62, 20, 60, 20], "text": "Mbappé", "font_size": 20,
         "color": "#ffffff", "bold": True, "font": "Arial", "confidence": 0.98},
        {"box": [130, 20, 50, 20], "text": "TOP10", "font_size": 20,
         "color": "#d6a84a", "bold": True, "font": "Arial", "confidence": 0.99},
    ]

    assert hasattr(text_detect, "_merge_adjacent_text_items")
    merged = text_detect._merge_adjacent_text_items(items)

    assert [item["text"] for item in merged] == ["Kylian Mbappé", "TOP10"]


def test_merged_chinese_label_keeps_space_before_date() -> None:
    items = [
        {"box": [10, 20, 40, 20], "text": "截至", "font_size": 20,
         "color": "#ffffff", "bold": True, "font": "Microsoft YaHei",
         "confidence": 0.99},
        {"box": [52, 20, 90, 20], "text": "2026.07.03", "font_size": 20,
         "color": "#ffffff", "bold": True, "font": "Arial", "confidence": 0.99},
    ]

    merged = text_detect._merge_adjacent_text_items(items)

    assert [item["text"] for item in merged] == ["截至 2026.07.03"]


def test_overlapping_number_and_chinese_suffix_are_merged() -> None:
    items = [
        {"box": [10, 20, 28, 40], "text": "1", "font_size": 20,
         "color": "#ffffff", "bold": True, "font": "Arial", "confidence": 0.99},
        {"box": [27, 20, 62, 40], "text": "助攻", "font_size": 20,
         "color": "#ffffff", "bold": True, "font": "Microsoft YaHei",
         "confidence": 0.99},
    ]

    merged = text_detect._merge_adjacent_text_items(items)

    assert [item["text"] for item in merged] == ["1助攻"]


def test_merged_chinese_text_reselects_east_asian_font() -> None:
    items = [
        {"box": [10, 20, 28, 40], "text": "6", "font_size": 20,
         "color": "#ffffff", "bold": True, "font": "Arial", "confidence": 0.99},
        {"box": [40, 20, 40, 40], "text": "球", "font_size": 20,
         "color": "#ffffff", "bold": True, "font": "Microsoft YaHei",
         "confidence": 0.99},
    ]

    merged = text_detect._merge_adjacent_text_items(items)

    assert merged[0]["font"] == "Microsoft YaHei"


def test_large_chinese_title_uses_sans_serif_bold() -> None:
    assert _select_font("世界杯射手榜", 60) == "Microsoft YaHei"
    assert not _should_force_regular_weight("世界杯射手榜", 60)


def test_text_ink_mask_covers_antialiased_glyph_edges() -> None:
    import numpy as np

    img = np.full((40, 80, 3), 255, dtype=np.uint8)
    img[12:28, 22:58] = [160, 160, 160]
    img[15:25, 28:52] = [0, 0, 0]
    text_mask = np.zeros((40, 80), dtype=np.uint8)
    text_mask[10:30, 20:60] = 255

    ink_mask = _build_text_ink_mask(img, text_mask)

    assert np.all(ink_mask[12:28, 22:58] == 255)


def test_detected_text_is_removed_from_foreground_after_cleanup() -> None:
    import cv2
    import numpy as np

    img = np.full((80, 160, 3), 245, dtype=np.uint8)
    img[20:60, 20:140] = [160, 20, 20]
    cv2.putText(
        img,
        "ABC",
        (52, 47),
        cv2.FONT_HERSHEY_SIMPLEX,
        0.9,
        (255, 255, 255),
        2,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((80, 160), dtype=np.uint8)
    text_mask[24:52, 48:112] = 255

    bg = build_background(img, text_mask=text_mask)
    fg_mask = extract_foreground_mask(img, bg, text_mask, diff_threshold=20)
    ink_mask = _build_text_ink_mask(img, text_mask)

    assert np.count_nonzero(fg_mask[ink_mask > 0]) == 0
    assert np.count_nonzero(fg_mask[30:55, 25:45]) > 200


def test_text_holes_do_not_split_one_graphic_into_many_components(tmp_path: Path) -> None:
    import cv2
    import numpy as np

    img = np.full((80, 160, 3), 245, dtype=np.uint8)
    img[20:60, 20:140] = [160, 20, 20]
    cv2.putText(
        img,
        "ABC",
        (52, 47),
        cv2.FONT_HERSHEY_SIMPLEX,
        0.9,
        (255, 255, 255),
        2,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((80, 160), dtype=np.uint8)
    text_mask[24:52, 48:112] = 255

    bg = build_background(img, text_mask=text_mask)
    fg_mask = extract_foreground_mask(img, bg, text_mask, diff_threshold=20)
    components = split_components(img, fg_mask, tmp_path, min_area=20, text_mask=text_mask)

    assert len(components) == 1


def test_component_rgb_cleans_text_pixels_over_graphic(tmp_path: Path) -> None:
    import cv2
    import numpy as np

    img = np.full((80, 160, 3), 245, dtype=np.uint8)
    img[20:60, 20:140] = [160, 20, 20]
    cv2.putText(
        img,
        "ABC",
        (52, 47),
        cv2.FONT_HERSHEY_SIMPLEX,
        0.9,
        (255, 255, 255),
        2,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((80, 160), dtype=np.uint8)
    text_mask[24:52, 48:112] = 255

    bg = build_background(img, text_mask=text_mask)
    fg_mask = extract_foreground_mask(img, bg, text_mask, diff_threshold=20)
    ink_mask = _build_text_ink_mask(img, text_mask)
    components = split_components(img, fg_mask, tmp_path, min_area=20, text_mask=text_mask)

    comp = components[0]
    rgba = Image.open(comp["path"])
    alpha = np.array(rgba)[:, :, 3]
    local_ink = ink_mask[
        comp["y"]:comp["y"] + comp["h"],
        comp["x"]:comp["x"] + comp["w"],
    ]

    rgb = np.array(rgba)[:, :, :3]

    ink_alpha = alpha[local_ink > 0]
    repaired_rgb = rgb[(local_ink > 0) & (alpha > 0)]

    assert np.count_nonzero(ink_alpha == 255) >= ink_alpha.size * 0.95
    assert np.max(repaired_rgb) < 230


def test_narrow_text_gap_does_not_split_component_grouping(tmp_path: Path) -> None:
    import numpy as np

    img = np.full((80, 160, 3), 245, dtype=np.uint8)
    mask = np.zeros((80, 160), dtype=np.uint8)
    mask[20:60, 20:140] = 255
    mask[20:60, 76:84] = 0
    img[20:60, 20:140] = [160, 20, 20]

    components = split_components(img, mask, tmp_path, min_area=20)

    assert len(components) == 1


def test_background_repair_does_not_blur_unmasked_pixels() -> None:
    import numpy as np

    img = np.full((60, 80, 3), 240, dtype=np.uint8)
    img[:, ::2] = [220, 230, 240]
    img[:, 1::2] = [250, 245, 235]
    exclude_mask = np.zeros((60, 80), dtype=np.uint8)
    exclude_mask[20:40, 30:50] = 255
    img[exclude_mask > 0] = [120, 80, 20]

    repaired = _original_based_background(
        img, exclude_mask, np.array([240.0, 240.0, 240.0])
    )

    assert np.array_equal(repaired[exclude_mask == 0], img[exclude_mask == 0])


def test_clean_background_restores_trusted_text_dilation_after_component_inpaint() -> None:
    import numpy as np

    source = np.zeros((20, 30, 3), dtype=np.uint8)
    component_mask = np.full((20, 30), 255, dtype=np.uint8)
    text_mask = np.zeros((20, 30), dtype=np.uint8)
    text_mask[6:14, 8:22] = 255
    text_removal = bg_model.build_removal_mask([], text_mask) > 0
    text_halo = text_removal & (text_mask == 0)
    text_clean = np.full_like(source, 23)
    text_clean[text_removal] = [41, 73, 109]

    def gray_inpaint(image, mask):
        return np.full_like(image, 127)

    cleaned = bg_model.build_clean_background(
        source,
        [component_mask],
        text_mask,
        large_inpainter=gray_inpaint,
        text_clean_image=text_clean,
        text_restore_mask=None,
    )

    assert np.any(text_halo)
    np.testing.assert_array_equal(cleaned[text_halo], text_clean[text_halo])
    np.testing.assert_array_equal(cleaned[text_removal], text_clean[text_removal])
    assert np.all(cleaned[~text_removal] == 127)


def test_clean_background_restores_independent_text_exclusion_mask() -> None:
    import numpy as np

    source = np.zeros((20, 30, 3), dtype=np.uint8)
    component_mask = np.full(source.shape[:2], 255, dtype=np.uint8)
    cleanup_mask = np.zeros(source.shape[:2], dtype=np.uint8)
    cleanup_mask[8:12, 13:17] = 255
    text_restore_mask = np.zeros(source.shape[:2], dtype=np.uint8)
    text_restore_mask[5:15, 7:23] = 255
    restore_region = bg_model.build_removal_mask([], text_restore_mask) > 0
    text_clean = np.full_like(source, 37)

    cleaned = bg_model.build_clean_background(
        source,
        [component_mask],
        cleanup_mask,
        large_inpainter=lambda image, mask: np.full_like(image, 127),
        text_clean_image=text_clean,
        text_restore_mask=text_restore_mask,
    )

    np.testing.assert_array_equal(cleaned[restore_region], text_clean[restore_region])
    assert np.all(cleaned[~restore_region] == 127)


def test_clean_background_rejects_mismatched_trusted_text_clean_shape() -> None:
    import numpy as np

    source = np.zeros((12, 16, 3), dtype=np.uint8)
    text_mask = np.zeros(source.shape[:2], dtype=np.uint8)

    with pytest.raises(ValueError, match="text-clean image must match"):
        bg_model.build_clean_background(
            source,
            [],
            text_mask,
            text_clean_image=np.zeros((11, 16, 3), dtype=np.uint8),
        )


def test_clean_background_rejects_mismatched_text_restore_mask_shape() -> None:
    import numpy as np

    source = np.zeros((12, 16, 3), dtype=np.uint8)
    text_mask = np.zeros(source.shape[:2], dtype=np.uint8)

    with pytest.raises(ValueError, match="text restore mask must match"):
        bg_model.build_clean_background(
            source,
            [],
            text_mask,
            text_clean_image=source.copy(),
            text_restore_mask=np.zeros((11, 16), dtype=np.uint8),
        )


def test_clean_background_without_trusted_text_clean_keeps_inpaint_result() -> None:
    import numpy as np

    source = np.zeros((12, 16, 3), dtype=np.uint8)
    component_mask = np.full(source.shape[:2], 255, dtype=np.uint8)
    text_mask = np.zeros(source.shape[:2], dtype=np.uint8)
    text_mask[5:7, 7:9] = 255

    cleaned = bg_model.build_clean_background(
        source,
        [component_mask],
        text_mask,
        large_inpainter=lambda image, mask: np.full_like(image, 127),
        text_clean_image=None,
        text_restore_mask=np.zeros((1, 1), dtype=np.uint8),
    )

    assert np.all(cleaned == 127)


def test_process_image_restores_raw_text_mask_after_each_component_inpaint(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((20, 30, 3), 240, dtype=np.uint8)
    image_path = tmp_path / "source.png"
    Image.fromarray(source, mode="RGB").save(image_path)
    work_dir = tmp_path / "work"
    work_dir.mkdir()
    raw_text_mask = np.zeros(source.shape[:2], dtype=np.uint8)
    raw_text_mask[4:16, 5:25] = 255
    text_mask_path = work_dir / "source-text-mask.png"
    Image.fromarray(raw_text_mask, mode="L").save(text_mask_path)
    cleanup_mask = np.zeros(source.shape[:2], dtype=np.uint8)
    cleanup_mask[8:12, 13:17] = 255
    text_clean = np.full_like(source, 229)
    element_mask = np.zeros(source.shape[:2], dtype=bool)
    element_mask[2:18, 3:27] = True
    element = types.SimpleNamespace(
        mask=element_mask,
        semantic_mask=element_mask.copy(),
    )
    build_calls = []

    monkeypatch.setattr(
        image_to_ppt,
        "_build_text_ink_mask",
        lambda *args, **kwargs: np.zeros(source.shape[:2], dtype=np.uint8),
    )

    def fake_cleanup(image, text_mask, text_items):
        np.testing.assert_array_equal(text_mask, raw_text_mask)
        return cleanup_mask.copy()

    monkeypatch.setattr(image_to_ppt, "_build_text_cleanup_mask", fake_cleanup)
    monkeypatch.setattr(
        image_to_ppt,
        "_repair_text_background",
        lambda image, mask, **kwargs: text_clean,
    )
    monkeypatch.setattr(image_to_ppt, "generate_object_proposals", lambda *args: [])
    monkeypatch.setattr(
        image_to_ppt,
        "generate_prompted_mask_candidates",
        lambda *args: [],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "generate_mask_candidates",
        lambda *args, **kwargs: [],
    )
    monkeypatch.setattr(image_to_ppt, "filter_prompt_free_candidates", lambda *args: [])
    monkeypatch.setattr(
        image_to_ppt,
        "combine_residual_candidates",
        lambda **kwargs: ([], 1),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "resolve_visual_elements",
        lambda *args: [element],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "recheck_visual_element_holes",
        lambda *args: None,
    )

    def fake_build(
        image,
        element_masks,
        positional_cleanup_mask,
        large_inpainter=None,
        text_clean_image=None,
        text_restore_mask=None,
    ):
        build_calls.append(
            (
                positional_cleanup_mask.copy(),
                None if text_restore_mask is None else text_restore_mask.copy(),
                text_clean_image,
            )
        )
        return source.copy()

    monkeypatch.setattr(image_to_ppt, "build_clean_background", fake_build)
    monkeypatch.setattr(
        image_to_ppt,
        "export_visual_components",
        lambda *args, **kwargs: [],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        lambda background: (background.copy(), 0, 0, "identity"),
    )

    image_to_ppt._process_image(
        image_path,
        work_dir,
        object_detector=object(),
        mask_generator=object(),
        lang="en",
        text_analysis={
            "items": [
                {
                    "box": [5, 4, 20, 12],
                    "text": "Editable",
                    "color": "#202020",
                }
            ],
            "mask_path": str(text_mask_path),
        },
        defer_quality=True,
    )

    assert len(build_calls) == 3
    for positional_cleanup, restore_mask, trusted_image in build_calls:
        np.testing.assert_array_equal(positional_cleanup, cleanup_mask)
        np.testing.assert_array_equal(restore_mask, raw_text_mask)
        assert trusted_image is text_clean


def test_white_text_on_colored_bar_is_cleaned_to_bar_color() -> None:
    import cv2
    import numpy as np

    bar_color = np.array([0, 75, 180], dtype=np.uint8)
    img = np.full((90, 220, 3), 245, dtype=np.uint8)
    img[20:60, 10:210] = bar_color
    cv2.putText(
        img,
        "HEADER",
        (25, 48),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.1,
        (255, 255, 255),
        2,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((90, 220), dtype=np.uint8)
    text_mask[18:62, 20:160] = 255

    cleaned = _fill_text_regions(img, text_mask)
    bar_region = np.zeros(text_mask.shape, dtype=bool)
    bar_region[20:60, 10:210] = True
    original_text_pixels = (text_mask > 0) & bar_region & np.all(img > 220, axis=2)
    cleaned_text_pixels = cleaned[original_text_pixels]

    assert np.count_nonzero(np.all(cleaned_text_pixels > 220, axis=1)) == 0
    assert np.mean(np.abs(cleaned_text_pixels.astype(int) - bar_color.astype(int))) < 20


def test_text_ink_mask_ignores_colored_bar_background() -> None:
    import cv2
    import numpy as np

    bar_color = np.array([0, 75, 180], dtype=np.uint8)
    img = np.full((90, 220, 3), 245, dtype=np.uint8)
    img[20:60, 10:210] = bar_color
    cv2.putText(
        img,
        "HEADER",
        (25, 48),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.1,
        (255, 255, 255),
        2,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((90, 220), dtype=np.uint8)
    text_mask[18:62, 20:160] = 255

    ink_mask = _build_text_ink_mask(img, text_mask)
    blue_background = (text_mask > 0) & np.all(img == bar_color, axis=2)
    white_text = (text_mask > 0) & np.all(img > 220, axis=2)

    assert np.count_nonzero(ink_mask[blue_background]) < np.count_nonzero(blue_background) * 0.10
    assert np.count_nonzero(ink_mask[white_text]) >= 800


def test_large_foreground_region_is_replaced_in_background() -> None:
    import numpy as np

    img = np.full((100, 140, 3), 245, dtype=np.uint8)
    img[20:90, 50:130] = [80, 140, 210]
    fg_mask = np.zeros((100, 140), dtype=np.uint8)
    fg_mask[20:90, 50:130] = 255

    repaired = _original_based_background(
        img,
        fg_mask,
        np.array([245.0, 245.0, 245.0]),
        fg_mask=fg_mask,
    )

    assert not np.array_equal(repaired[fg_mask > 0], img[fg_mask > 0])
    assert np.max(np.abs(repaired[fg_mask > 0].astype(int) - 245)) <= 8


def test_load_rgb_handles_unicode_windows_paths(tmp_path: Path) -> None:
    import numpy as np

    image_path = tmp_path / "qmd+文本向量化.png"
    Image.new("RGB", (12, 8), (10, 20, 30)).save(image_path)

    img = _load_rgb(str(image_path))

    assert img.shape == (8, 12, 3)
    assert np.array_equal(img[0, 0], np.array([10, 20, 30], dtype=np.uint8))


def test_small_component_repair_mask_includes_shadow_halo() -> None:
    import numpy as np

    fg_mask = np.zeros((80, 100), dtype=np.uint8)
    fg_mask[30:50, 40:60] = 255

    repair_mask = _build_component_repair_mask(fg_mask)

    assert repair_mask[29, 39] == 255
    assert repair_mask[50, 60] == 255
    assert repair_mask[10, 10] == 0


def test_large_component_alpha_fills_internal_light_holes(tmp_path: Path) -> None:
    import numpy as np

    img = np.full((120, 180, 3), 245, dtype=np.uint8)
    img[20:110, 40:170] = [120, 180, 220]

    mask = np.zeros((120, 180), dtype=np.uint8)
    mask[20:110, 40:170] = 255
    mask[35:65, 70:130] = 0

    components = split_components(img, mask, tmp_path, min_area=20)
    comp = components[0]
    alpha = np.array(Image.open(comp["path"]).convert("RGBA"))[:, :, 3]

    local_x = 90 - comp["x"]
    local_y = 45 - comp["y"]
    assert alpha[local_y, local_x] == 255


def test_east_asian_font_is_written_as_typeface_child() -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    run = box.text_frame.paragraphs[0].add_run()

    _set_run_font(run, "Microsoft YaHei")

    xml = run._r.xml
    assert '<a:ea typeface="Microsoft YaHei"/>' in xml
    assert '<a:latin typeface="Microsoft YaHei"/>' in xml


def test_large_chinese_title_uses_sans_serif_with_detected_weight() -> None:
    text = "辛亥革命的烽火岁月"

    assert _select_font(text, 80.0) == "Microsoft YaHei"
    assert not _should_force_regular_weight(text, 80.0)
    assert _adjust_font_size(text, 100.0) == 88.0
    assert _adjust_font_size("红色全景资源创意展示", 20.0) == 20.0
def test_visual_compare_qa_writes_artifacts(tmp_path: Path) -> None:
    source = tmp_path / "source.png"
    preview = tmp_path / "preview.png"
    out_dir = tmp_path / "qa"
    Image.new("RGB", (20, 10), (255, 255, 255)).save(source)
    Image.new("RGB", (20, 10), (245, 245, 245)).save(preview)

    report = write_visual_compare(source, preview, out_dir)

    assert report["preview_size"] == [20, 10]
    assert report["mean_abs_diff_0_255"] == 10.0
    assert (out_dir / "side_by_side.png").exists()
    assert (out_dir / "blend.png").exists()
    assert (out_dir / "diff_heatmap.png").exists()
    assert (out_dir / "report.json").exists()


def test_assemble_psd_requires_aspose_license(tmp_path: Path, monkeypatch) -> None:
    bg_path = tmp_path / "bg.png"
    Image.new("RGB", (20, 10), "white").save(bg_path)
    out_path = tmp_path / "out.psd"
    monkeypatch.delenv("ASPOSE_PSD_LICENSE", raising=False)

    try:
        assemble_psd(
            background_path=bg_path,
            components=[],
            text_items=[],
            img_width=20,
            img_height=10,
            output_path=out_path,
        )
    except AsposePsdLicenseError as exc:
        assert "ASPOSE_PSD_LICENSE" in str(exc)
    else:
        raise AssertionError("assemble_psd should require an Aspose.PSD license")

    assert not out_path.exists()


def test_resolve_psd_outputs_writes_one_file_per_input(tmp_path: Path) -> None:
    one = tmp_path / "one.png"
    two = tmp_path / "two.png"
    one.write_bytes(b"x")
    two.write_bytes(b"x")

    assert _resolve_output_paths([one], None) == [one.with_suffix(".psd")]
    assert _resolve_output_paths([one], tmp_path / "custom.psd") == [
        tmp_path / "custom.psd"
    ]
    assert _resolve_output_paths([one, two], tmp_path / "psd_out") == [
        tmp_path / "psd_out" / "one.psd",
        tmp_path / "psd_out" / "two.psd",
    ]


def test_assemble_psd_writes_real_text_layers_with_aspose_backend(
    tmp_path: Path, monkeypatch
) -> None:
    bg_path = tmp_path / "bg.png"
    Image.new("RGB", (20, 10), "white").save(bg_path)
    license_path = tmp_path / "Aspose.PSD.lic"
    license_path.write_text("fake", encoding="utf-8")
    out_path = tmp_path / "out.psd"
    calls = []

    class FakeLicense:
        def set_license(self, path):
            calls.append(("license", str(path)))

    class FakeColor:
        @staticmethod
        def from_argb(a, r, g, b):
            return (a, r, g, b)

    class FakeRectangle:
        def __init__(self, x, y, w, h):
            self.x = x
            self.y = y
            self.w = w
            self.h = h

    class FakeLayer:
        def __init__(self, stream):
            self.stream = stream
            self.display_name = ""

    class FakeTextLayer:
        def __init__(self):
            self.display_name = ""
            self.text_color = None
            self.font_size = None

    class FakePsdImage:
        def __init__(self, width, height):
            self.width = width
            self.height = height
            self.layers = []

        def add_layer(self, layer):
            calls.append(("pixel", layer.display_name))

        def add_text_layer(self, text, rect):
            calls.append(("text", text, rect.x, rect.y, rect.w, rect.h))
            return FakeTextLayer()

        def save(self, path):
            Path(path).write_bytes(b"psd")

        def dispose(self):
            calls.append(("dispose",))

    aspose_mod = types.ModuleType("aspose")
    psd_mod = types.ModuleType("aspose.psd")
    psd_mod.License = FakeLicense
    psd_mod.Color = FakeColor
    psd_mod.Rectangle = FakeRectangle
    fileformats_mod = types.ModuleType("aspose.psd.fileformats")
    psd_file_mod = types.ModuleType("aspose.psd.fileformats.psd")
    psd_file_mod.PsdImage = FakePsdImage
    layers_mod = types.ModuleType("aspose.psd.fileformats.psd.layers")
    layers_mod.Layer = FakeLayer

    monkeypatch.setenv("ASPOSE_PSD_LICENSE", str(license_path))
    monkeypatch.setitem(sys.modules, "aspose", aspose_mod)
    monkeypatch.setitem(sys.modules, "aspose.psd", psd_mod)
    monkeypatch.setitem(sys.modules, "aspose.psd.fileformats", fileformats_mod)
    monkeypatch.setitem(sys.modules, "aspose.psd.fileformats.psd", psd_file_mod)
    monkeypatch.setitem(sys.modules, "aspose.psd.fileformats.psd.layers", layers_mod)

    assemble_psd(
        background_path=bg_path,
        components=[],
        text_items=[{
            "box": [2, 3, 10, 4],
            "text": "Hello",
            "font_size": 12,
            "color": "#112233",
        }],
        img_width=20,
        img_height=10,
        output_path=out_path,
    )

    assert ("license", str(license_path)) in calls
    assert ("pixel", "Background") in calls
    assert ("text", "Hello", 2, 3, 10, 4) in calls
    assert out_path.exists()


def test_prepare_multiple_images_separates_ocr_and_visual_stages(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    image_paths = [tmp_path / "first.png", tmp_path / "second.png"]
    for image_path in image_paths:
        Image.new("RGB", (20, 30), "white").save(image_path)

    events = []
    detector = object()
    sam = object()

    def fake_detect(image_path, lang):
        events.append(f"ocr:{Path(image_path).stem}")
        return [], np.zeros((30, 20), dtype=np.uint8)

    def fake_process(
        image_path,
        work_dir,
        object_detector,
        mask_generator,
        lang,
        *,
        text_analysis,
        defer_quality,
    ):
        assert object_detector is detector
        assert mask_generator is sam
        assert text_analysis["items"] == []
        assert Path(text_analysis["mask_path"]).name == "source-text-mask.png"
        assert defer_quality is True
        return {
            "original_image_path": str(image_path),
            "_work_dir": str(work_dir),
            "components": [],
        }

    def fake_finalize(slide_data, lang):
        events.append(f"finalize:{Path(slide_data['original_image_path']).stem}")
        return slide_data

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(
        image_to_ppt,
        "close_ocr_engines",
        lambda: events.append("close-ocr"),
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "create_object_detector",
        lambda: events.append("load-dino") or detector,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "create_sam_generator",
        lambda checkpoint: events.append("load-sam") or sam,
    )
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(image_to_ppt, "_process_image", fake_process)
    monkeypatch.setattr(
        image_to_ppt,
        "_release_visual_resources",
        lambda: events.append("release-visual"),
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_finalize_slide_quality",
        fake_finalize,
        raising=False,
    )

    slides_data = image_to_ppt._prepare_multiple_images(image_paths, "en")

    assert [Path(slide["original_image_path"]).stem for slide in slides_data] == [
        "first",
        "second",
    ]
    assert events == [
        "ocr:first",
        "ocr:second",
        "close-ocr",
        "load-dino",
        "load-sam",
        "release-visual",
        "finalize:first",
        "finalize:second",
        "close-ocr",
    ]


def test_prepare_single_image_reuses_staged_resource_pipeline(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    events = []
    detector = object()
    sam = object()

    def fake_detect(path, lang):
        events.append("source-ocr")
        return [], np.zeros((30, 20), dtype=np.uint8)

    def fake_process(
        path,
        work_dir,
        object_detector,
        mask_generator,
        lang,
        *,
        text_analysis=None,
        defer_quality=False,
    ):
        assert object_detector is detector
        assert mask_generator is sam
        assert text_analysis is not None
        assert defer_quality is True
        events.append("visual")
        return {
            "background_original_path": str(work_dir / "background-original.png"),
            "background_widescreen_path": str(work_dir / "background-16x9.png"),
            "components": [],
            "text_items": [],
            "img_width": 20,
            "img_height": 30,
            "canvas_width": 53,
            "canvas_height": 30,
            "content_offset_x": 16,
            "content_offset_y": 0,
            "original_image_path": str(path),
            "_work_dir": str(work_dir),
            "_text_mask_path": text_analysis["mask_path"],
            "_element_mask_paths": [],
        }

    def fake_finalize(slide_data, lang):
        events.append("finalize")
        slide_data.pop("_work_dir")
        slide_data.pop("_text_mask_path")
        slide_data.pop("_element_mask_paths")
        return slide_data

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(
        image_to_ppt,
        "close_ocr_engines",
        lambda: events.append("close-ocr"),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "create_object_detector",
        lambda: events.append("load-dino") or detector,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "create_sam_generator",
        lambda checkpoint: events.append("load-sam") or sam,
    )
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(image_to_ppt, "_process_image", fake_process)
    monkeypatch.setattr(
        image_to_ppt,
        "_release_visual_resources",
        lambda: events.append("release-visual"),
    )
    monkeypatch.setattr(image_to_ppt, "_finalize_slide_quality", fake_finalize)

    work_root = tmp_path / "run" / "work"
    slide_data, work_dir = image_to_ppt._prepare_single_image(
        image_path,
        "en",
        _work_root=work_root,
    )

    assert "_work_dir" not in slide_data
    assert work_dir == work_root.resolve() / "page_001"
    assert work_dir == Path(slide_data["background_original_path"]).parent
    assert events == [
        "source-ocr",
        "close-ocr",
        "load-dino",
        "load-sam",
        "visual",
        "release-visual",
        "finalize",
        "close-ocr",
    ]


def test_prepare_multiple_images_uses_isolated_run_page_directories(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    image_paths = [tmp_path / "first.png", tmp_path / "second.png"]
    for image_path in image_paths:
        Image.new("RGB", (20, 30), "white").save(image_path)
    work_root = tmp_path / "run" / "work"
    seen_work_dirs = []

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: (
            [],
            np.zeros((30, 20), dtype=np.uint8),
        ),
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(
        image_to_ppt,
        "create_sam_generator",
        lambda checkpoint: object(),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "resolve_sam_checkpoint",
        lambda: Path("sam.pt"),
    )

    def fake_process(
        image_path,
        work_dir,
        *args,
        text_analysis,
        **kwargs,
    ):
        marker = work_dir / "marker.txt"
        marker.write_text(Path(image_path).stem, encoding="utf-8")
        seen_work_dirs.append(work_dir)
        return {
            "background_original_path": str(work_dir / "background.png"),
            "original_image_path": str(image_path),
            "components": [],
        }

    monkeypatch.setattr(image_to_ppt, "_process_image", fake_process)
    monkeypatch.setattr(image_to_ppt, "_release_visual_resources", lambda: None)
    monkeypatch.setattr(
        image_to_ppt,
        "_finalize_slide_quality",
        lambda slide_data, lang: slide_data,
    )

    image_to_ppt._prepare_multiple_images(
        image_paths,
        "en",
        _work_root=work_root,
    )

    assert seen_work_dirs == [
        work_root.resolve() / "page_001",
        work_root.resolve() / "page_002",
    ]
    assert (seen_work_dirs[0] / "marker.txt").read_text(encoding="utf-8") == "first"
    assert (seen_work_dirs[1] / "marker.txt").read_text(encoding="utf-8") == "second"


def test_prepare_multiple_images_refuses_existing_run_page_directory(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    page_dir = tmp_path / "work" / "page_001"
    page_dir.mkdir(parents=True)
    sentinel = page_dir / "sentinel.txt"
    sentinel.write_text("keep", encoding="utf-8")
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: pytest.fail("existing page reached OCR"),
    )

    with pytest.raises(FileExistsError):
        image_to_ppt._prepare_multiple_images(
            [image_path],
            "en",
            _work_root=tmp_path / "work",
        )

    assert sentinel.read_text(encoding="utf-8") == "keep"


def test_prepare_multiple_images_keeps_system_temp_default(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    created = tmp_path / "system-temp"
    calls = []

    def fake_mkdtemp(*, prefix: str) -> str:
        calls.append(prefix)
        created.mkdir()
        return str(created)

    monkeypatch.setattr(image_to_ppt.tempfile, "mkdtemp", fake_mkdtemp)

    assert image_to_ppt._work_directory(None, 0) == created.resolve()
    assert calls == ["img2ppt_0_"]


@pytest.mark.parametrize(
    ("function_name", "batched"),
    [
        ("convert", False),
        ("convert_variants", False),
        ("convert_batch", True),
        ("convert_batch_variants", True),
    ],
)
def test_conversion_entrypoints_forward_internal_work_root(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    function_name: str,
    batched: bool,
) -> None:
    image_path = tmp_path / "source.png"
    work_root = tmp_path / "run" / "work"
    slide_data = {
        "background_original_path": str(tmp_path / "background.png"),
        "background_widescreen_path": str(tmp_path / "background-wide.png"),
        "components": [],
        "text_items": [],
        "img_width": 20,
        "img_height": 30,
        "original_image_path": str(image_path),
    }
    seen_options = []

    def fake_single(
        image_path,
        lang,
        _work_root=None,
        _resource_isolation=False,
    ):
        seen_options.append((_work_root, _resource_isolation))
        return slide_data, work_root / "page_001"

    def fake_multiple(
        image_paths,
        lang,
        _work_root=None,
        _resource_isolation=False,
    ):
        seen_options.append((_work_root, _resource_isolation))
        return [slide_data]

    monkeypatch.setattr(image_to_ppt, "_prepare_single_image", fake_single)
    monkeypatch.setattr(image_to_ppt, "_prepare_multiple_images", fake_multiple)
    monkeypatch.setattr(
        image_to_ppt,
        "_assemble_prepared_slide",
        lambda data, output, add_reference, slide_size: str(output),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "assemble_pptx_multi",
        lambda **kwargs: str(kwargs["output_path"]),
    )

    converter = getattr(image_to_ppt, function_name)
    source = [image_path] if batched else image_path
    kwargs = {
        "output_path": tmp_path / "deck.pptx",
        "_work_root": work_root,
        "_resource_isolation": True,
    }
    if function_name == "convert_batch_variants":
        kwargs["combine_original"] = True
    converter(source, **kwargs)

    assert seen_options == [(work_root, True)]


def test_convert_variants_accepts_finalized_single_slide_without_private_work_dir(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    image_path = tmp_path / "source.png"
    work_dir = tmp_path / "work"
    slide_data = {
        "background_original_path": str(work_dir / "background-original.png"),
        "background_widescreen_path": str(work_dir / "background-16x9.png"),
        "components": [],
        "text_items": [],
        "img_width": 20,
        "img_height": 30,
        "canvas_width": 53,
        "canvas_height": 30,
        "content_offset_x": 16,
        "content_offset_y": 0,
        "original_image_path": str(image_path),
    }
    assembled_sizes = []

    monkeypatch.setattr(
        image_to_ppt,
        "_prepare_multiple_images",
        lambda image_paths, lang, _work_root=None: [slide_data],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_assemble_prepared_slide",
        lambda data, output, add_reference, slide_size: (
            assembled_sizes.append(slide_size) or str(output)
        ),
    )

    result = image_to_ppt.convert_variants(image_path, tmp_path / "deck.pptx")

    assert result == {
        "original": str(tmp_path / "deck_original.pptx"),
        "16:9": str(tmp_path / "deck_16x9.pptx"),
    }
    assert assembled_sizes == ["original", "16:9"]


def test_release_lama_model_drops_cached_instance() -> None:
    sentinel = object()
    lama_inpaint._MODEL = sentinel

    lama_inpaint.release_model()

    assert lama_inpaint._MODEL is None


@pytest.mark.parametrize(
    ("failure_stage", "expected_events"),
    [
        (
            "source-ocr",
            ["ocr:first", "ocr:second", "close-ocr"],
        ),
        (
            "load-sam",
            [
                "ocr:first",
                "ocr:second",
                "close-ocr",
                "load-dino",
                "load-sam",
                "release-visual",
            ],
        ),
        (
            "visual",
            [
                "ocr:first",
                "ocr:second",
                "close-ocr",
                "load-dino",
                "load-sam",
                "visual:first",
                "release-visual",
            ],
        ),
        (
            "finalize",
            [
                "ocr:first",
                "ocr:second",
                "close-ocr",
                "load-dino",
                "load-sam",
                "visual:first",
                "visual:second",
                "release-visual",
                "finalize:first",
                "close-ocr",
            ],
        ),
    ],
)
def test_staged_resource_cleanup_preserves_original_failure(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    failure_stage: str,
    expected_events: list[str],
) -> None:
    import numpy as np

    image_paths = [tmp_path / "first.png", tmp_path / "second.png"]
    for image_path in image_paths:
        Image.new("RGB", (20, 30), "white").save(image_path)

    events = []
    failure = RuntimeError(f"{failure_stage} failed")

    def fake_detect(image_path, lang):
        name = Path(image_path).stem
        events.append(f"ocr:{name}")
        if failure_stage == "source-ocr" and name == "second":
            raise failure
        return [], np.zeros((30, 20), dtype=np.uint8)

    def fake_create_sam(checkpoint):
        events.append("load-sam")
        if failure_stage == "load-sam":
            raise failure
        return object()

    def fake_process(
        image_path,
        work_dir,
        object_detector,
        mask_generator,
        lang,
        **kwargs,
    ):
        name = Path(image_path).stem
        events.append(f"visual:{name}")
        if failure_stage == "visual":
            raise failure
        return {
            "original_image_path": str(image_path),
            "_work_dir": str(work_dir),
            "components": [],
        }

    def fake_finalize(slide_data, lang):
        name = Path(slide_data["original_image_path"]).stem
        events.append(f"finalize:{name}")
        if failure_stage == "finalize":
            raise failure
        return slide_data

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(
        image_to_ppt,
        "close_ocr_engines",
        lambda: events.append("close-ocr"),
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "create_object_detector",
        lambda: events.append("load-dino") or object(),
    )
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", fake_create_sam)
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(image_to_ppt, "_process_image", fake_process)
    monkeypatch.setattr(
        image_to_ppt,
        "_release_visual_resources",
        lambda: events.append("release-visual"),
        raising=False,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_finalize_slide_quality",
        fake_finalize,
        raising=False,
    )

    with pytest.raises(RuntimeError) as caught:
        image_to_ppt._prepare_multiple_images(image_paths, "en")

    assert caught.value is failure
    assert events == expected_events


def test_staged_resource_cleanup_releases_models_before_cache_cleanup(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import gc
    import weakref
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    references = {}
    release_state = []
    failure = RuntimeError("visual failed")

    class Model:
        pass

    def fake_create_detector():
        model = Model()
        references["detector"] = weakref.ref(model)
        return model

    def fake_create_sam(checkpoint):
        model = Model()
        references["sam"] = weakref.ref(model)
        return model

    def fake_process(*args, **kwargs):
        raise failure

    def fake_release():
        gc.collect()
        release_state.append(
            (references["detector"]() is None, references["sam"]() is None)
        )

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((30, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", fake_create_detector)
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", fake_create_sam)
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(image_to_ppt, "_process_image", fake_process)
    monkeypatch.setattr(image_to_ppt, "_release_visual_resources", fake_release)
    monkeypatch.setattr(
        image_to_ppt,
        "_finalize_slide_quality",
        lambda *args: pytest.fail("finalization must not run"),
    )

    with pytest.raises(RuntimeError) as caught:
        image_to_ppt._prepare_multiple_images([image_path], "en")

    assert caught.value is failure
    assert release_state == [(True, True)]


@pytest.mark.parametrize("failure_type", [RuntimeError, KeyboardInterrupt])
def test_staged_resource_cleanup_preserves_visual_failure_when_release_fails(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    caplog: pytest.LogCaptureFixture,
    failure_type: type[BaseException],
) -> None:
    import logging
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    failure = failure_type("visual failed")
    cleanup_failure = RuntimeError("visual cleanup failed")

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((30, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", lambda path: object())
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(
        image_to_ppt,
        "_process_image",
        lambda *args, **kwargs: (_ for _ in ()).throw(failure),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_release_visual_resources",
        lambda: (_ for _ in ()).throw(cleanup_failure),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_finalize_slide_quality",
        lambda *args: pytest.fail("finalization must not run"),
    )

    with caplog.at_level(logging.ERROR, logger=image_to_ppt.__name__):
        with pytest.raises(failure_type) as caught:
            image_to_ppt._prepare_multiple_images([image_path], "en")

    assert caught.value is failure
    assert any("cleanup failed" in record.message for record in caplog.records)


def test_staged_resource_cleanup_propagates_release_failure_without_primary_error(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    cleanup_failure = RuntimeError("visual cleanup failed")
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((30, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", lambda path: object())
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(
        image_to_ppt,
        "_process_image",
        lambda image_path, work_dir, *args, **kwargs: {
            "original_image_path": str(image_path),
            "_work_dir": str(work_dir),
            "components": [],
        },
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_release_visual_resources",
        lambda: (_ for _ in ()).throw(cleanup_failure),
    )

    with pytest.raises(RuntimeError) as caught:
        image_to_ppt._prepare_multiple_images([image_path], "en")

    assert caught.value is cleanup_failure


def test_staged_resource_cleanup_preserves_caller_exception_traceback_boundary(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import gc
    import traceback as traceback_module
    import weakref

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    primary_failure = RuntimeError("source OCR failed")
    cleanup_failure = RuntimeError("OCR close failed")
    references = {}
    caller_state = {}

    class Marker:
        pass

    def raise_caller_failure():
        marker = Marker()
        references["marker"] = weakref.ref(marker)
        raise ValueError("caller's handled error")

    def fake_detect(*args, **kwargs):
        raise primary_failure

    def fake_close():
        raise cleanup_failure

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", fake_close)

    try:
        raise_caller_failure()
    except ValueError as caller_exception:
        caller_state["exception"] = caller_exception
        with pytest.raises(RuntimeError) as caught:
            image_to_ppt._prepare_multiple_images([image_path], "en")

    gc.collect()
    caller_exception = caller_state["exception"]
    caller_frames = [
        frame
        for frame, _ in traceback_module.walk_tb(caller_exception.__traceback__)
        if frame.f_code.co_name == "raise_caller_failure"
    ]

    assert caught.value is primary_failure
    assert primary_failure.__context__ is caller_exception
    assert caller_frames
    assert references["marker"]() is not None
    assert caller_frames[0].f_locals["marker"] is references["marker"]()


def test_staged_resource_cleanup_log_does_not_retain_cleanup_resources(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    caplog: pytest.LogCaptureFixture,
) -> None:
    import gc
    import logging
    import weakref
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    references = {}
    failure = RuntimeError("visual failed")

    class Engine:
        pass

    def fake_release():
        engine = Engine()
        references["engine"] = weakref.ref(engine)
        raise RuntimeError("visual cleanup failed")

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((30, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", lambda path: object())
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(
        image_to_ppt,
        "_process_image",
        lambda *args, **kwargs: (_ for _ in ()).throw(failure),
    )
    monkeypatch.setattr(image_to_ppt, "_release_visual_resources", fake_release)

    with caplog.at_level(logging.ERROR, logger=image_to_ppt.__name__):
        with pytest.raises(RuntimeError) as caught:
            image_to_ppt._prepare_multiple_images([image_path], "en")

    gc.collect()
    assert caught.value is failure
    assert references["engine"]() is None
    assert any("cleanup failed" in record.message for record in caplog.records)


def test_staged_resource_cleanup_clears_nested_lama_cause_frames(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import gc
    import traceback as traceback_module
    import weakref
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    cause_failure = RuntimeError("LaMa backend failed")

    class FailingLama:
        def __call__(self, image, mask):
            raise cause_failure

    lama_inpaint._MODEL = FailingLama()
    model_reference = weakref.ref(lama_inpaint._MODEL)

    def fake_process(*args, **kwargs):
        source = np.zeros((8, 8, 3), dtype=np.uint8)
        mask = np.ones((8, 8), dtype=np.uint8) * 255
        return lama_inpaint.inpaint_large_mask(source, mask)

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((30, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", lambda path: object())
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(image_to_ppt, "_process_image", fake_process)

    try:
        with pytest.raises(lama_inpaint.LargeMaskInpaintError) as caught:
            image_to_ppt._prepare_multiple_images([image_path], "en")

        gc.collect()
        assert caught.value.__cause__ is cause_failure
        assert caught.value.__context__ is cause_failure
        assert "__call__" in [
            frame.name
            for frame in traceback_module.extract_tb(cause_failure.__traceback__)
        ]
        assert "inpaint_large_mask" in [
            frame.name
            for frame in traceback_module.extract_tb(caught.value.__traceback__)
        ]
        assert model_reference() is None
    finally:
        lama_inpaint.release_model()


def test_staged_resource_cleanup_clears_suppressed_cleanup_cause_frames(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import gc
    import traceback as traceback_module
    import weakref
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    primary_failure = RuntimeError("visual failed")
    cleanup_chain = {}
    references = {}

    class FailingCleanup:
        def __call__(self):
            raise RuntimeError("cleanup backend failed")

    def fake_cleanup():
        cleanup_model = FailingCleanup()
        references["model"] = weakref.ref(cleanup_model)
        try:
            cleanup_model()
        except RuntimeError as cause:
            wrapped = RuntimeError("cleanup wrapper failed")
            cleanup_chain["wrapped"] = wrapped
            cleanup_chain["cause"] = cause
            raise wrapped from cause

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((30, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", lambda path: object())
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(
        image_to_ppt,
        "_process_image",
        lambda *args, **kwargs: (_ for _ in ()).throw(primary_failure),
    )
    monkeypatch.setattr(image_to_ppt, "_release_visual_resources", fake_cleanup)

    with pytest.raises(RuntimeError) as caught:
        image_to_ppt._prepare_multiple_images([image_path], "en")

    gc.collect()
    assert caught.value is primary_failure
    assert cleanup_chain["wrapped"].__cause__ is cleanup_chain["cause"]
    assert cleanup_chain["wrapped"].__context__ is cleanup_chain["cause"]
    assert "__call__" in [
        frame.name
        for frame in traceback_module.extract_tb(
            cleanup_chain["cause"].__traceback__
        )
    ]
    assert "fake_cleanup" in [
        frame.name
        for frame in traceback_module.extract_tb(
            cleanup_chain["wrapped"].__traceback__
        )
    ]
    assert references["model"]() is None


@pytest.mark.parametrize("close_fails", [False, True])
def test_staged_resource_cleanup_source_ocr_interrupt_releases_engine(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    caplog: pytest.LogCaptureFixture,
    close_fails: bool,
) -> None:
    import gc
    import logging
    import weakref

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    cache = {}
    references = {}
    failure = KeyboardInterrupt("source OCR interrupted")
    cleanup_failure = RuntimeError("OCR close failed")

    class Engine:
        pass

    def fake_detect(*args, **kwargs):
        engine = Engine()
        references["engine"] = weakref.ref(engine)
        cache["engine"] = engine
        raise failure

    def fake_close():
        cache.clear()
        if close_fails:
            raise cleanup_failure

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", fake_close)

    with caplog.at_level(logging.ERROR, logger=image_to_ppt.__name__):
        with pytest.raises(KeyboardInterrupt) as caught:
            image_to_ppt._prepare_multiple_images([image_path], "en")

    gc.collect()
    assert caught.value is failure
    assert references["engine"]() is None
    if close_fails:
        assert any("cleanup failed" in record.message for record in caplog.records)


@pytest.mark.parametrize("close_fails", [False, True])
def test_staged_resource_cleanup_final_ocr_interrupt_releases_engine(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
    caplog: pytest.LogCaptureFixture,
    close_fails: bool,
) -> None:
    import gc
    import logging
    import weakref
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    cache = {}
    references = {}
    detect_calls = 0
    close_calls = 0
    failure = KeyboardInterrupt("final OCR interrupted")
    cleanup_failure = RuntimeError("OCR close failed")

    class Engine:
        pass

    def fake_detect(*args, **kwargs):
        nonlocal detect_calls
        detect_calls += 1
        if detect_calls == 1:
            return [], np.zeros((30, 20), dtype=np.uint8)
        engine = Engine()
        references["engine"] = weakref.ref(engine)
        cache["engine"] = engine
        raise failure

    def fake_close():
        nonlocal close_calls
        close_calls += 1
        cache.clear()
        if close_fails and close_calls == 2:
            raise cleanup_failure

    def fake_process(image_path, work_dir, *args, text_analysis, **kwargs):
        return {
            "background_original_path": str(image_path),
            "components": [],
            "text_items": [],
            "original_image_path": str(image_path),
            "_work_dir": str(work_dir),
            "_text_mask_path": text_analysis["mask_path"],
            "_element_mask_paths": [],
        }

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", fake_close)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", lambda path: object())
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam.pt"))
    monkeypatch.setattr(image_to_ppt, "_process_image", fake_process)
    monkeypatch.setattr(image_to_ppt, "_release_visual_resources", lambda: None)

    with caplog.at_level(logging.ERROR, logger=image_to_ppt.__name__):
        with pytest.raises(KeyboardInterrupt) as caught:
            image_to_ppt._prepare_multiple_images([image_path], "en")

    gc.collect()
    assert caught.value is failure
    assert references["engine"]() is None
    if close_fails:
        assert any("cleanup failed" in record.message for record in caplog.records)


def test_staged_resource_cleanup_propagates_ocr_close_failure_without_primary_error(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    cleanup_failure = RuntimeError("OCR close failed")
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((30, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "close_ocr_engines",
        lambda: (_ for _ in ()).throw(cleanup_failure),
    )

    with pytest.raises(RuntimeError) as caught:
        image_to_ppt._prepare_multiple_images([image_path], "en")

    assert caught.value is cleanup_failure


def test_deferred_mask_paths_release_source_masks_before_visual_loading(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import gc
    import weakref
    import numpy as np

    image_paths = [tmp_path / "first.png", tmp_path / "second.png"]
    for image_path in image_paths:
        Image.new("RGB", (20, 30), "white").save(image_path)
    mask_references = []
    release_state = []
    failure = RuntimeError("stop after OCR")

    def fake_detect(*args, **kwargs):
        mask = np.zeros((30, 20), dtype=np.uint8)
        mask_references.append(weakref.ref(mask))
        return [], mask

    def fake_create_detector():
        gc.collect()
        release_state.append([reference() is None for reference in mask_references])
        raise failure

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", fake_create_detector)
    monkeypatch.setattr(image_to_ppt, "_release_visual_resources", lambda: None)

    with pytest.raises(RuntimeError) as caught:
        image_to_ppt._prepare_multiple_images(image_paths, "en")

    assert caught.value is failure
    assert release_state == [[True, True]]


def test_deferred_mask_paths_release_source_mask_when_png_save_fails(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import gc
    import weakref
    import numpy as np

    image_path = tmp_path / "source.png"
    Image.new("RGB", (20, 30), "white").save(image_path)
    references = {}
    failure = RuntimeError("mask save failed")

    def fake_detect(*args, **kwargs):
        mask = np.zeros((30, 20), dtype=np.uint8)
        references["mask"] = weakref.ref(mask)
        return [], mask

    class FailingImage:
        def save(self, path):
            raise failure

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(
        image_to_ppt.Image,
        "fromarray",
        lambda *args, **kwargs: FailingImage(),
    )

    with pytest.raises(RuntimeError) as caught:
        image_to_ppt._prepare_multiple_images([image_path], "en")

    gc.collect()
    assert caught.value is failure
    assert references["mask"]() is None


def test_deferred_mask_paths_are_lossless_and_hold_no_arrays(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((30, 20, 3), 40, dtype=np.uint8)
    image_path = tmp_path / "source.png"
    Image.fromarray(source, mode="RGB").save(image_path)
    work_dir = tmp_path / "work"
    text_mask = np.zeros((30, 20), dtype=np.uint8)
    text_mask[1:4, 2:7] = 255
    text_mask_path = work_dir / "source-text-mask.png"
    work_dir.mkdir()
    Image.fromarray(text_mask, mode="L").save(text_mask_path)
    masks = []
    for x1, x2 in ((1, 5), (11, 17)):
        mask = np.zeros((30, 20), dtype=bool)
        mask[5:25, x1:x2] = True
        masks.append(mask)
    elements = [
        types.SimpleNamespace(mask=mask, semantic_mask=mask.copy())
        for mask in masks
    ]

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: pytest.fail("source OCR must not repeat"),
    )
    monkeypatch.setattr(image_to_ppt, "generate_object_proposals", lambda *args: [])
    monkeypatch.setattr(
        image_to_ppt,
        "generate_prompted_mask_candidates",
        lambda *args: [],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "generate_mask_candidates",
        lambda *args, **kwargs: [],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "filter_prompt_free_candidates",
        lambda *args: [],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "combine_residual_candidates",
        lambda **kwargs: ([], 0),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "resolve_visual_elements",
        lambda candidates: elements,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "recheck_visual_element_holes",
        lambda *args: None,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "build_clean_background",
        lambda *args, **kwargs: source.copy(),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "export_visual_components",
        lambda *args, **kwargs: [],
    )
    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        lambda background: (background.copy(), 0, 0, "identity"),
    )

    slide_data = image_to_ppt._process_image(
        image_path,
        work_dir,
        object_detector=object(),
        mask_generator=object(),
        lang="en",
        text_analysis={"items": [], "mask_path": str(text_mask_path)},
        defer_quality=True,
    )

    assert slide_data["_text_mask_path"] == str(text_mask_path.resolve())
    assert "_element_masks" not in slide_data
    assert "_semantic_masks" not in slide_data
    for key in ("_element_mask_paths", "_semantic_mask_paths"):
        assert len(slide_data[key]) == 2
        assert all(
            Path(mask_path).is_absolute() and Path(mask_path).is_file()
            for mask_path in slide_data[key]
        )
    assert not any(isinstance(value, np.ndarray) for value in slide_data.values())
    for key in ("_element_mask_paths", "_semantic_mask_paths"):
        for mask_path, expected in zip(slide_data[key], masks):
            with Image.open(mask_path) as saved:
                assert saved.mode == "L"
                actual = np.asarray(saved)
            assert set(np.unique(actual)) <= {0, 255}
            assert np.array_equal(actual > 0, expected)


def test_deferred_mask_paths_preserve_quality_error_message(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((30, 20, 3), 40, dtype=np.uint8)
    image_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_mask_path = tmp_path / "source-text-mask.png"
    Image.fromarray(source, mode="RGB").save(image_path)
    Image.fromarray(source, mode="RGB").save(background_path)
    Image.fromarray(np.zeros((30, 20), dtype=np.uint8), mode="L").save(
        text_mask_path
    )
    slide_data = {
        "background_original_path": str(background_path),
        "components": [],
        "text_items": [],
        "original_image_path": str(image_path),
        "_work_dir": str(tmp_path),
        "_text_mask_path": str(text_mask_path),
        "_element_mask_paths": [],
    }
    quality = {"mae": 1.2345, "p95": 9.8765}
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((30, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "visual_difference",
        lambda *args: quality,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "write_segmentation_diagnostics",
        lambda *args, **kwargs: None,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "require_visual_quality",
        lambda value: (_ for _ in ()).throw(
            visual_segment.VisualSegmentationError("quality failed")
        ),
    )

    with pytest.raises(visual_segment.VisualSegmentationError) as caught:
        image_to_ppt._finalize_slide_quality(slide_data, "en")

    assert str(caught.value) == (
        "quality failed; mae=1.234, p95=9.877, "
        f"diagnostics={(tmp_path / 'diagnostics').resolve()}"
    )
    assert slide_data["_text_mask_path"] == str(text_mask_path)
    assert slide_data["_element_mask_paths"] == []


def test_visual_difference_flags_sparse_visible_artifacts() -> None:
    import numpy as np

    source = np.zeros((100, 100, 3), dtype=np.uint8)
    reconstructed = source.copy()
    reconstructed[:20, :10] = 12
    metrics = visual_segment.visual_difference(
        source,
        reconstructed,
        np.zeros((100, 100), dtype=np.uint8),
    )

    assert metrics["p99"] == pytest.approx(12.0)
    assert metrics["changed_ratio"] == pytest.approx(0.02)
    assert metrics["largest_artifact_ratio"] == pytest.approx(0.02)
    assert visual_segment.needs_text_only_fallback(metrics)


def test_text_cleanup_mask_covers_glyphs_without_masking_the_whole_ocr_box() -> None:
    import cv2
    import numpy as np

    source = np.full((80, 240, 3), 255, dtype=np.uint8)
    cv2.putText(
        source,
        "Editable",
        (18, 57),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.6,
        (20, 20, 20),
        3,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((80, 240), dtype=np.uint8)
    text_mask[5:72, 5:230] = 255

    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [{"box": [5, 5, 225, 67], "text": "Editable", "color": "#141414"}],
    )

    visible_glyph = np.mean(source, axis=2) < 220
    assert np.all(cleanup[visible_glyph] == 255)
    assert cleanup[8, 220] == 0
    assert np.count_nonzero(cleanup) < np.count_nonzero(text_mask) * 0.5


def test_text_cleanup_mask_keeps_sparse_low_variance_glyphs() -> None:
    import cv2
    import numpy as np

    source = np.full((30, 40, 3), [19, 95, 167], dtype=np.uint8)
    cv2.putText(
        source,
        "3.1",
        (2, 21),
        cv2.FONT_HERSHEY_SIMPLEX,
        0.45,
        (16, 61, 104),
        1,
        cv2.LINE_AA,
    )
    text_mask = np.full((30, 40), 255, dtype=np.uint8)

    assert float(np.std(cv2.cvtColor(source, cv2.COLOR_RGB2GRAY))) < 8.0
    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [{"box": [0, 0, 40, 30], "text": "3.1", "color": "#103d68"}],
    )

    glyphs = cv2.cvtColor(source, cv2.COLOR_RGB2GRAY) < 70
    assert np.count_nonzero(cleanup[glyphs]) >= np.count_nonzero(glyphs) * 0.9
    assert np.count_nonzero(cleanup) < np.count_nonzero(text_mask) * 0.6


def test_text_cleanup_preserves_non_glyph_background_pixels() -> None:
    import cv2
    import numpy as np

    source = np.full((80, 240, 3), 247, dtype=np.uint8)
    cv2.putText(
        source,
        "No ghost",
        (15, 56),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.5,
        (30, 30, 30),
        3,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((80, 240), dtype=np.uint8)
    text_mask[4:72, 4:232] = 255
    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [{"box": [4, 4, 228, 68], "text": "No ghost", "color": "#1e1e1e"}],
    )

    repaired = image_to_ppt._repair_text_background(
        source,
        cleanup,
        large_inpainter=lambda image, mask: np.full_like(image, 247),
    )

    assert np.array_equal(repaired[cleanup == 0], source[cleanup == 0])
    assert float(np.mean(np.abs(repaired.astype(float) - 247))) < 1.0


def test_text_cleanup_escalates_when_lightweight_repair_retains_glyph_edges(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import numpy as np

    source = np.full((80, 160, 3), 255, dtype=np.uint8)
    source[25:55, 50:110] = 20
    cleanup = np.zeros((80, 160), dtype=np.uint8)
    cleanup[20:60, 45:115] = 255
    calls = []

    monkeypatch.setattr(
        image_to_ppt,
        "repair_masked_background",
        lambda image, mask, large_inpainter=None: image.copy(),
    )

    def fake_large(image, mask):
        calls.append(mask.copy())
        repaired = image.copy()
        repaired[mask > 0] = 255
        return repaired

    repaired = image_to_ppt._repair_text_background(
        source,
        cleanup,
        large_inpainter=fake_large,
    )

    assert len(calls) == 1
    assert np.all(repaired[25:55, 50:110] == 255)
    assert np.array_equal(repaired[cleanup == 0], source[cleanup == 0])


def test_text_cleanup_uses_local_background_model_before_large_inpaint(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import cv2
    import numpy as np

    x = np.linspace(238, 252, 180, dtype=np.uint8)
    source = np.repeat(x[None, :, None], 100, axis=0)
    source = np.repeat(source, 3, axis=2)
    source[38:63, 58:122] = 20
    cleanup = np.zeros((100, 180), dtype=np.uint8)
    cleanup[35:66, 55:125] = 255

    monkeypatch.setattr(
        image_to_ppt,
        "repair_masked_background",
        lambda image, mask, large_inpainter=None: image.copy(),
    )

    repaired = image_to_ppt._repair_text_background(
        source,
        cleanup,
        text_items=[
            {
                "box": [55, 35, 70, 31],
                "text": "ghost",
                "color": "#141414",
            }
        ],
        large_inpainter=lambda *args: pytest.fail(
            "local model should avoid large inpaint"
        ),
    )

    residual = visual_segment.background_residual_metrics(
        source,
        repaired,
        cleanup,
    )
    changed = np.any(repaired != source, axis=2)
    bounded = cv2.dilate(
        cleanup,
        np.ones((13, 13), dtype=np.uint8),
        iterations=1,
    ) > 0
    assert not visual_segment.has_background_residual(residual)
    assert not np.any(changed & ~bounded)


def test_text_cleanup_removes_blurred_antialiasing_without_feather_ghosts(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import cv2
    import numpy as np

    alpha = np.zeros((140, 500), dtype=np.uint8)
    cv2.putText(
        alpha,
        "Ghost text",
        (20, 95),
        cv2.FONT_HERSHEY_SIMPLEX,
        2.4,
        255,
        3,
        cv2.LINE_AA,
    )
    alpha = cv2.GaussianBlur(alpha, (0, 0), sigmaX=2).astype(float) / 255
    source = (
        255 * (1 - alpha[:, :, None]) + 5 * alpha[:, :, None]
    ).astype(np.uint8)
    source = np.repeat(source, 3, axis=2)
    text_mask = np.zeros((140, 500), dtype=np.uint8)
    text_mask[15:115, 10:480] = 255
    text_items = [
        {
            "box": [10, 15, 470, 100],
            "text": "Ghost text",
            "color": "#050505",
            "font_size": 33.0,
        }
    ]
    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        text_items,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "repair_masked_background",
        lambda image, mask, large_inpainter=None: image.copy(),
    )

    repaired = image_to_ppt._repair_text_background(
        source,
        cleanup,
        text_items=text_items,
        large_inpainter=lambda *args: pytest.fail(
            "blurred glyph edges should be handled by the local model"
        ),
    )

    residual = visual_segment.background_residual_metrics(
        source,
        repaired,
        cleanup,
    )
    assert residual["retained_edge_ratio"] < 0.02
    assert int(repaired.min()) >= 240
    assert np.array_equal(repaired[cleanup == 0], source[cleanup == 0])


def test_text_cleanup_reconstructs_colored_button_without_text_imprint() -> None:
    import cv2
    import numpy as np

    height, width = 120, 360
    clean = np.full((height, width, 3), 255, dtype=np.uint8)
    button = np.zeros((height, width), dtype=np.uint8)
    cv2.rectangle(button, (80, 25), (280, 95), 255, -1)
    cv2.circle(button, (80, 60), 35, 255, -1)
    cv2.circle(button, (280, 60), 35, 255, -1)
    y_grid, x_grid = np.indices((height, width))
    gradient = np.zeros_like(clean)
    gradient[:, :, 0] = np.clip(
        245 + (x_grid - 60) * 0.04,
        0,
        255,
    )
    gradient[:, :, 1] = np.clip(
        35 + (x_grid - 60) * 0.08 + (y_grid - 25) * 0.2,
        0,
        255,
    )
    gradient[:, :, 2] = 15
    clean[button > 0] = gradient[button > 0]
    source = clean.copy()
    cv2.putText(
        source,
        "BAD",
        (125, 72),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.5,
        (250, 250, 250),
        3,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((height, width), dtype=np.uint8)
    text_mask[30:88, 110:250] = 255
    text_items = [
        {
            "box": [110, 30, 140, 58],
            "text": "BAD",
            "color": "#fafafa",
            "font_size": 22.0,
        }
    ]
    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        text_items,
    )

    repaired = image_to_ppt._repair_text_with_local_planes(
        source,
        cleanup,
        text_items,
    )

    error = np.abs(repaired.astype(float) - clean.astype(float))
    assert float(np.mean(error[cleanup > 0])) < 3.0
    assert float(np.percentile(error[cleanup > 0], 95)) < 8.0
    assert np.array_equal(repaired[cleanup == 0], source[cleanup == 0])


def test_text_cleanup_mask_includes_antialiasing_just_outside_ocr_box() -> None:
    import numpy as np

    source = np.full((30, 50, 3), 250, dtype=np.uint8)
    source[10:20, 20:24] = 20
    source[20:24, 21:23] = 20
    source[16:24, 28:31] = 20
    text_mask = np.zeros((30, 50), dtype=np.uint8)
    text_mask[8:20, 17:31] = 255

    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [{"box": [17, 8, 14, 12], "text": "I", "color": "#141414"}],
    )

    assert np.all(cleanup[20:24, 21:23] == 255)
    assert np.all(cleanup[16:24, 28:31] == 255)


def test_text_cleanup_mask_covers_wide_antialiasing_connected_to_glyph() -> None:
    import numpy as np

    source = np.full((40, 60, 3), 250, dtype=np.uint8)
    source[8:30, 16:34] = 80
    source[12:26, 20:30] = 20
    text_mask = np.zeros((40, 60), dtype=np.uint8)
    text_mask[6:32, 14:36] = 255

    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [{"box": [14, 6, 22, 26], "text": "I", "color": "#141414"}],
    )

    assert np.all(cleanup[8:30, 16:34] == 255)
    assert cleanup[2, 8] == 0


def test_text_cleanup_mask_covers_multiple_text_colors_in_one_ocr_box() -> None:
    import cv2
    import numpy as np

    source = np.full((70, 300, 3), 255, dtype=np.uint8)
    cv2.putText(
        source,
        "RED",
        (10, 48),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.2,
        (220, 30, 20),
        2,
        cv2.LINE_AA,
    )
    cv2.putText(
        source,
        "black",
        (100, 48),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.2,
        (20, 20, 20),
        2,
        cv2.LINE_AA,
    )
    text_mask = np.zeros((70, 300), dtype=np.uint8)
    text_mask[5:62, 5:285] = 255

    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [
            {
                "box": [5, 5, 280, 57],
                "text": "RED black",
                "color": "#dc1e14",
                "font_size": 18.0,
            }
        ],
    )

    red_ink = (source[:, :, 0] > 100) & (source[:, :, 1] < 180)
    black_ink = np.mean(source, axis=2) < 160
    assert np.all(cleanup[red_ink] == 255)
    assert np.all(cleanup[black_ink] == 255)


def test_text_cleanup_mask_covers_secondary_color_touching_ocr_box_edge() -> None:
    import cv2
    import numpy as np

    source = np.full((70, 320, 3), 255, dtype=np.uint8)
    cv2.putText(
        source, "RED", (10, 45), cv2.FONT_HERSHEY_SIMPLEX, 1.0,
        (230, 45, 20), 2, cv2.LINE_AA,
    )
    cv2.putText(
        source, "black", (100, 45), cv2.FONT_HERSHEY_SIMPLEX, 1.0,
        (30, 30, 30), 2, cv2.LINE_AA,
    )
    text_mask = np.zeros(source.shape[:2], dtype=np.uint8)
    text_mask[5:46, 5:280] = 255

    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [{
            "box": [5, 5, 275, 41],
            "text": "RED black",
            "color": "#e62d14",
            "font_size": 18.0,
        }],
    )

    black_ink = (
        np.all(source < 80, axis=2)
        & (np.indices(source.shape[:2])[1] >= 100)
    )
    assert np.mean(cleanup[black_ink] > 0) >= 0.98


def test_text_cleanup_mask_preserves_long_graphic_line_near_text() -> None:
    import cv2
    import numpy as np

    source = np.full((80, 240, 3), 255, dtype=np.uint8)
    cv2.putText(
        source,
        "Text",
        (30, 38),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.0,
        (20, 20, 20),
        2,
        cv2.LINE_AA,
    )
    cv2.line(source, (10, 44), (230, 44), (30, 120, 220), 2)
    text_mask = np.zeros((80, 240), dtype=np.uint8)
    text_mask[10:48, 20:150] = 255

    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [
            {
                "box": [20, 10, 130, 38],
                "text": "Text",
                "color": "#141414",
                "font_size": 14.0,
            }
        ],
    )

    text_ink = np.mean(source, axis=2) < 80
    line = np.zeros(cleanup.shape, dtype=bool)
    line[43:46, 10:231] = True
    assert np.all(cleanup[text_ink & ~line] == 255)
    assert np.count_nonzero(cleanup[line]) < 12


def test_text_cleanup_mask_preserves_colored_card_edge_at_text_box_boundary() -> None:
    import cv2
    import numpy as np

    source = np.full((70, 150, 3), (240, 249, 240), dtype=np.uint8)
    cv2.line(source, (20, 8), (20, 62), (25, 110, 50), 3)
    cv2.putText(
        source,
        "Label",
        (24, 43),
        cv2.FONT_HERSHEY_SIMPLEX,
        0.8,
        (45, 47, 46),
        2,
        cv2.LINE_AA,
    )
    text_mask = np.zeros(source.shape[:2], dtype=np.uint8)
    text_mask[14:52, 20:105] = 255

    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [{"box": [20, 14, 85, 38], "text": "Label", "color": "#2d2f2e"}],
    )

    glyph = (
        (source[:, :, 0] < 100)
        & (np.max(source, axis=2) - np.min(source, axis=2) < 20)
        & (np.indices(source.shape[:2])[1] > 22)
    )
    card_edge = np.zeros(cleanup.shape, dtype=bool)
    card_edge[8:63, 18:23] = True
    assert np.count_nonzero(cleanup[glyph]) >= np.count_nonzero(glyph) * 0.95
    assert np.count_nonzero(cleanup[card_edge]) == 0


def test_low_confidence_short_square_ocr_candidate_is_treated_as_icon() -> None:
    items = [
        {
            "box": [905, 333, 69, 74],
            "text": "вт",
            "font_size": 30.1,
            "confidence": 0.859,
        },
        {
            "box": [385, 142, 87, 51],
            "text": "优点",
            "font_size": 21.1,
            "confidence": 0.99,
        },
    ]

    assert image_to_ppt._filter_probable_icon_text_items(items) == [items[1]]


def test_probable_icon_filter_removes_discarded_item_from_ocr_mask() -> None:
    import numpy as np

    mask = np.zeros((120, 180), dtype=np.uint8)
    mask[20:50, 20:50] = 255
    mask[60:85, 70:150] = 255
    items = [
        {
            "box": [20, 20, 30, 30],
            "text": "x",
            "confidence": 0.6,
        },
        {
            "box": [70, 60, 80, 25],
            "text": "Editable",
            "confidence": 0.99,
        },
    ]

    filtered, filtered_mask = (
        image_to_ppt._filter_probable_icon_text_analysis(items, mask)
    )

    assert filtered == [items[1]]
    assert np.count_nonzero(filtered_mask[20:50, 20:50]) == 0
    assert np.all(filtered_mask[60:85, 70:150] == 255)


def test_overlapping_ocr_prefixes_keep_only_the_complete_text() -> None:
    items = [
        {"box": [1046, 451, 413, 34], "text": "3.大仓库成本高：规划阶段依赖文", "confidence": 0.98},
        {"box": [1047, 455, 506, 28], "text": "3.大仓库成本高：规划阶段依赖文件遍历，", "confidence": 0.97},
        {"box": [733, 802, 81, 48], "text": "一句", "confidence": 0.99},
        {"box": [738, 803, 182, 46], "text": "一句话总结", "confidence": 0.98},
        {"box": [100, 100, 120, 30], "text": "独立文本", "confidence": 0.99},
    ]

    assert image_to_ppt._deduplicate_overlapping_text_items(items) == [
        items[1], items[3], items[4]
    ]


def test_text_cleanup_recovers_detached_same_color_glyph_outside_box_boundary(
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    import numpy as np

    source = np.full((60, 50, 3), 250, dtype=np.uint8)
    source[22:38, 24:31] = 20
    source[22:38, 33:37] = 20
    text_mask = np.zeros((60, 50), dtype=np.uint8)
    text_mask[10:50, 15:31] = 255
    detected_ink = np.zeros((60, 50), dtype=np.uint8)
    detected_ink[22:38, 24:31] = 255
    monkeypatch.setattr(
        image_to_ppt,
        "_build_text_ink_mask",
        lambda *args, **kwargs: detected_ink,
    )

    cleanup = image_to_ppt._build_text_cleanup_mask(
        source,
        text_mask,
        [{"box": [15, 10, 16, 40], "text": "II", "color": "#141414"}],
    )

    assert np.all(cleanup[22:38, 33:37] == 255)


def test_disconnected_text_glyphs_do_not_trigger_large_mask_inpainting() -> None:
    import numpy as np

    mask = np.zeros((200, 400), dtype=np.uint8)
    for y in range(10, 190, 20):
        for x in range(10, 390, 20):
            mask[y:y + 5, x:x + 8] = 255

    assert np.mean(mask > 0) > 0.08
    assert not bg_model.needs_large_mask_inpaint(mask)


def test_background_residual_metrics_detect_faint_component_imprint() -> None:
    import numpy as np

    source = np.full((120, 180, 3), 255, dtype=np.uint8)
    source[30:90, 50:130] = 20
    removal = np.zeros((120, 180), dtype=np.uint8)
    removal[26:94, 46:134] = 255
    clean = np.full_like(source, 255)
    imprinted = clean.copy()
    imprinted[30:90, 50:130] = 235

    clean_metrics = visual_segment.background_residual_metrics(
        source, clean, removal
    )
    imprint_metrics = visual_segment.background_residual_metrics(
        source, imprinted, removal
    )

    assert not visual_segment.has_background_residual(clean_metrics)
    assert visual_segment.has_background_residual(imprint_metrics)
    assert imprint_metrics["retained_edge_ratio"] > 0.8


def test_background_residual_metrics_detect_low_contrast_source_imprint() -> None:
    import numpy as np

    source = np.full((120, 180, 3), 255, dtype=np.uint8)
    source[30:90, 50:130] = 243
    removal = np.zeros((120, 180), dtype=np.uint8)
    removal[26:94, 46:134] = 255
    imprinted = np.full_like(source, 255)
    imprinted[30:90, 50:130] = 249

    metrics = visual_segment.background_residual_metrics(
        source,
        imprinted,
        removal,
    )

    assert visual_segment.has_background_residual(metrics)
    assert metrics["retained_edge_ratio"] > 0.8


def test_final_quality_excludes_only_editable_text_glyphs(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((20, 20, 3), 240, dtype=np.uint8)
    source[8:12, 9:11] = 20
    image_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_mask_path = tmp_path / "source-text-mask.png"
    Image.fromarray(source).save(image_path)
    Image.fromarray(source).save(background_path)
    text_mask = np.zeros((20, 20), dtype=np.uint8)
    text_mask[4:16, 4:16] = 255
    Image.fromarray(text_mask).save(text_mask_path)
    captured_masks = []

    def fake_difference(source, reconstructed, mask):
        captured_masks.append(mask.copy())
        return {"mae": 0.0, "p95": 0.0}

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((20, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(image_to_ppt, "visual_difference", fake_difference)
    monkeypatch.setattr(
        image_to_ppt,
        "write_segmentation_diagnostics",
        lambda *args, **kwargs: None,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "require_visual_quality",
        lambda quality: None,
    )
    slide_data = {
        "background_original_path": str(background_path),
        "components": [],
        "text_items": [
            {"box": [4, 4, 12, 12], "text": "I", "color": "#141414"}
        ],
        "original_image_path": str(image_path),
        "_work_dir": str(tmp_path),
        "_text_mask_path": str(text_mask_path),
        "_element_mask_paths": [],
    }

    semantic_mask_paths = [str(tmp_path / "semantic-mask.png")]
    slide_data["_semantic_mask_paths"] = semantic_mask_paths

    result = image_to_ppt._finalize_slide_quality(slide_data, "en")

    assert captured_masks[0][0, 0] == 0
    assert captured_masks[0][9, 9] == 255
    assert captured_masks[0][5, 15] == 0
    assert "_semantic_mask_paths" not in result
    assert slide_data["_semantic_mask_paths"] == semantic_mask_paths


def test_bold_estimation_handles_light_text_on_dark_background() -> None:
    import cv2
    import numpy as np

    regular = np.full((80, 260, 3), 40, dtype=np.uint8)
    bold = regular.copy()
    cv2.putText(
        regular,
        "Embedding",
        (8, 58),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.8,
        (245, 245, 245),
        2,
        cv2.LINE_AA,
    )
    cv2.putText(
        bold,
        "Embedding",
        (8, 58),
        cv2.FONT_HERSHEY_SIMPLEX,
        1.8,
        (245, 245, 245),
        6,
        cv2.LINE_AA,
    )

    assert not text_detect._estimate_bold(regular)
    assert text_detect._estimate_bold(bold)


def test_text_only_fallback_api_is_removed():
    import inspect

    assert not hasattr(image_to_ppt, "_apply_text_only_fallback")
    assert "_allow_text_only_fallback" not in inspect.signature(
        image_to_ppt._finalize_slide_quality
    ).parameters


def test_final_visual_difference_is_total_gate_not_text_only_success(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((20, 20, 3), 40, dtype=np.uint8)
    image_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_clean_path = tmp_path / "text-clean.png"
    text_mask_path = tmp_path / "source-text-mask.png"
    component_path = tmp_path / "component.png"
    Image.fromarray(source).save(image_path)
    Image.fromarray(source).save(background_path)
    Image.fromarray(source).save(text_clean_path)
    Image.fromarray(np.zeros((20, 20), dtype=np.uint8)).save(text_mask_path)
    Image.fromarray(np.zeros((2, 2, 4), dtype=np.uint8)).save(component_path)
    qualities = [
        {
            "mae": 1.5,
            "p95": 0.0,
            "p99": 12.0,
            "changed_ratio": 0.02,
        },
        {
            "mae": 0.0,
            "p95": 0.0,
            "p99": 0.0,
            "changed_ratio": 0.0,
        },
    ]
    checked = []
    widescreen_kwargs = []
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((20, 20), dtype=np.uint8)),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "visual_difference",
        lambda *args: qualities.pop(0),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "write_segmentation_diagnostics",
        lambda *args, **kwargs: None,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "require_visual_quality",
        lambda quality: checked.append(quality),
    )
    def fake_widescreen(image, **kwargs):
        widescreen_kwargs.append(kwargs)
        return image.copy(), 0, 0, "identity"

    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        fake_widescreen,
    )
    slide_data = {
        "background_original_path": str(background_path),
        "background_widescreen_path": str(background_path),
        "background_path": str(background_path),
        "components": [
            {
                "path": str(component_path),
                "x": 0,
                "y": 0,
                "w": 2,
                "h": 2,
            }
        ],
        "text_items": [{"box": [1, 1, 3, 3], "text": "x"}],
        "original_image_path": str(image_path),
        "_work_dir": str(tmp_path),
        "_text_clean_path": str(text_clean_path),
        "_text_mask_path": str(text_mask_path),
        "_element_mask_paths": [],
    }

    with pytest.raises(
        image_to_ppt.VisualSegmentationError,
        match="visible_visual_artifacts",
    ):
        image_to_ppt._finalize_slide_quality(
            slide_data,
            "en",
            _resource_isolation=True,
        )
    assert slide_data["components"]
    assert "quality_fallback" not in slide_data
    assert checked == []
    assert widescreen_kwargs == []


def test_final_quality_rejects_when_component_raster_text_remains(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((40, 80, 3), 245, dtype=np.uint8)
    component = np.zeros((20, 30, 4), dtype=np.uint8)
    component[:, :, :3] = 40
    component[:, :, 3] = 255
    image_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_clean_path = tmp_path / "text-clean.png"
    component_path = tmp_path / "component.png"
    text_mask_path = tmp_path / "source-text-mask.png"
    Image.fromarray(source).save(image_path)
    Image.fromarray(source).save(background_path)
    Image.fromarray(source).save(text_clean_path)
    Image.fromarray(component, mode="RGBA").save(component_path)
    Image.fromarray(np.zeros((40, 80), dtype=np.uint8)).save(text_mask_path)
    raster_item = {
        "box": [1, 1, 10, 8],
        "text": "raster",
        "color": "#282828",
    }
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: (
            [raster_item],
            np.zeros((40, 80), dtype=np.uint8),
        ),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "repair_exported_component_text",
        lambda *args, **kwargs: None,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        lambda image, **kwargs: (image.copy(), 0, 0, "identity"),
    )
    slide_data = {
        "background_original_path": str(background_path),
        "background_widescreen_path": str(background_path),
        "background_path": str(background_path),
        "components": [
            {
                "path": str(component_path),
                "x": 0,
                "y": 0,
                "w": 30,
                "h": 20,
            }
        ],
        "text_items": [],
        "original_image_path": str(image_path),
        "_work_dir": str(tmp_path),
        "_text_clean_path": str(text_clean_path),
        "_text_mask_path": str(text_mask_path),
        "_element_mask_paths": [],
    }

    with pytest.raises(
        image_to_ppt.VisualSegmentationError,
        match="component_raster_text",
    ):
        image_to_ppt._finalize_slide_quality(slide_data, "en")
    assert slide_data["components"]
    assert "quality_fallback" not in slide_data


def test_final_quality_reports_component_text_overlap_without_text_only_success(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((80, 160, 3), 255, dtype=np.uint8)
    source[30:50, 45:105] = 20
    imprinted = np.full_like(source, 255)
    imprinted[30:50, 45:105] = 252
    text_clean = np.full_like(source, 255)
    component = np.zeros((30, 80, 4), dtype=np.uint8)
    component[:, :, :3] = 255
    component[:, :, 3] = 255
    text_mask = np.zeros((80, 160), dtype=np.uint8)
    text_mask[30:50, 45:105] = 255
    element_mask = np.zeros((80, 160), dtype=np.uint8)
    element_mask[25:55, 35:115] = 255

    image_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_clean_path = tmp_path / "text-clean.png"
    component_path = tmp_path / "component.png"
    text_mask_path = tmp_path / "source-text-mask.png"
    element_mask_path = tmp_path / "element-mask.png"
    Image.fromarray(source).save(image_path)
    Image.fromarray(imprinted).save(background_path)
    Image.fromarray(text_clean).save(text_clean_path)
    Image.fromarray(component, mode="RGBA").save(component_path)
    Image.fromarray(text_mask).save(text_mask_path)
    Image.fromarray(element_mask).save(element_mask_path)

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((80, 160), dtype=np.uint8)),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        lambda image, **kwargs: (image.copy(), 0, 0, "identity"),
    )
    slide_data = {
        "background_original_path": str(background_path),
        "background_widescreen_path": str(background_path),
        "background_path": str(background_path),
        "components": [
            {
                "path": str(component_path),
                "x": 35,
                "y": 25,
                "w": 80,
                "h": 30,
            }
        ],
        "text_items": [
            {
                "box": [40, 25, 70, 35],
                "text": "editable",
                "font_size": 20,
                "color": "#141414",
            }
        ],
        "original_image_path": str(image_path),
        "_work_dir": str(tmp_path),
        "_text_clean_path": str(text_clean_path),
        "_text_mask_path": str(text_mask_path),
        "_element_mask_paths": [str(element_mask_path)],
    }

    with pytest.raises(
        image_to_ppt.VisualSegmentationError,
        match="component_text_overlap:component_0001",
    ):
        image_to_ppt._finalize_slide_quality(slide_data, "en")
    assert slide_data["components"]
    assert "quality_fallback" not in slide_data


def test_final_quality_rejects_when_clean_background_keeps_component_imprint(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((120, 180, 3), 255, dtype=np.uint8)
    source[30:90, 50:130] = 20
    clean_background = np.full_like(source, 255)
    clean_background[30:90, 50:130] = 235
    component = np.zeros((60, 80, 4), dtype=np.uint8)
    component[:, :, :3] = 20
    component[:, :, 3] = 255
    element_mask = np.zeros((120, 180), dtype=np.uint8)
    element_mask[30:90, 50:130] = 255

    image_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    component_path = tmp_path / "component.png"
    text_mask_path = tmp_path / "source-text-mask.png"
    element_mask_path = tmp_path / "element-mask.png"
    Image.fromarray(source).save(image_path)
    Image.fromarray(clean_background).save(background_path)
    Image.fromarray(component, mode="RGBA").save(component_path)
    Image.fromarray(np.zeros((120, 180), dtype=np.uint8)).save(text_mask_path)
    Image.fromarray(element_mask).save(element_mask_path)

    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((120, 180), dtype=np.uint8)),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        lambda image, **kwargs: (image.copy(), 0, 0, "identity"),
    )
    captured_residual = {}
    real_has_background_residual = image_to_ppt.has_background_residual

    def capture_background_residual(metrics):
        captured_residual.update(metrics)
        return real_has_background_residual(metrics)

    monkeypatch.setattr(
        image_to_ppt,
        "has_background_residual",
        capture_background_residual,
    )
    slide_data = {
        "background_original_path": str(background_path),
        "background_widescreen_path": str(background_path),
        "background_path": str(background_path),
        "components": [
            {
                "path": str(component_path),
                "x": 50,
                "y": 30,
                "w": 80,
                "h": 60,
            }
        ],
        "text_items": [],
        "original_image_path": str(image_path),
        "_work_dir": str(tmp_path),
        "_text_mask_path": str(text_mask_path),
        "_element_mask_paths": [str(element_mask_path)],
    }

    with pytest.raises(
        image_to_ppt.VisualSegmentationError,
        match="background_residual",
    ):
        image_to_ppt._finalize_slide_quality(slide_data, "en")
    assert slide_data["components"]
    assert captured_residual["retained_edge_ratio"] > 0.8
    assert "quality_fallback" not in slide_data


def test_final_quality_rejects_text_fallback_that_keeps_glyph_imprint(
    monkeypatch: pytest.MonkeyPatch,
    tmp_path: Path,
) -> None:
    import numpy as np

    source = np.full((120, 180, 3), 255, dtype=np.uint8)
    source[40:80, 60:120] = 20
    imprinted = np.full_like(source, 255)
    imprinted[40:80, 60:120] = 235
    text_mask = np.zeros((120, 180), dtype=np.uint8)
    text_mask[35:85, 55:125] = 255

    image_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_clean_path = tmp_path / "text-clean.png"
    text_mask_path = tmp_path / "source-text-mask.png"
    Image.fromarray(source).save(image_path)
    Image.fromarray(imprinted).save(background_path)
    Image.fromarray(imprinted).save(text_clean_path)
    Image.fromarray(text_mask).save(text_mask_path)
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: ([], np.zeros((120, 180), dtype=np.uint8)),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "build_widescreen_background",
        lambda image, **kwargs: (image.copy(), 0, 0, "identity"),
    )
    slide_data = {
        "background_original_path": str(background_path),
        "background_widescreen_path": str(background_path),
        "background_path": str(background_path),
        "components": [],
        "text_items": [
            {
                "box": [55, 35, 70, 50],
                "text": "ghost",
                "color": "#141414",
            }
        ],
        "original_image_path": str(image_path),
        "_work_dir": str(tmp_path),
        "_text_clean_path": str(text_clean_path),
        "_text_mask_path": str(text_mask_path),
        "_element_mask_paths": [],
    }

    with pytest.raises(
        image_to_ppt.VisualSegmentationError,
        match="background_residual",
    ):
        image_to_ppt._finalize_slide_quality(slide_data, "en")


def test_persist_visual_masks_normalizes_binary_l_images(tmp_path: Path) -> None:
    masks = [
        image_to_ppt.np.asarray([[0, 1], [2, 0]], dtype=image_to_ppt.np.uint8),
    ]

    paths = image_to_ppt._persist_visual_masks(tmp_path, "semantic-masks", masks)

    with Image.open(paths[0]) as stored:
        assert stored.mode == "L"
        assert image_to_ppt.np.array_equal(
            image_to_ppt.np.asarray(stored),
            image_to_ppt.np.asarray([[0, 255], [255, 0]], dtype=image_to_ppt.np.uint8),
        )


def _prepare_component_layers_fixture(
    tmp_path: Path,
    monkeypatch,
    *,
    resource_isolation: bool = False,
    before_visual_worker=None,
    before_prepare=None,
) -> tuple[dict, Path]:
    source_path = tmp_path / "outside-source.png"
    Image.new("RGB", (16, 10), "white").save(source_path)
    work_dir = tmp_path / "prepared"
    calls = {"ocr": 0, "visual": 0}

    def fake_detect(path, *, lang, **kwargs):
        calls["ocr"] += 1
        assert Path(path).resolve().is_relative_to(work_dir.resolve())
        assert lang == "en"
        assert bool(kwargs.get("isolated")) is resource_isolation
        return (
            [{
                "box": [1, 1, 3, 2],
                "text": "A",
                "font_size": 12.0,
                "color": "#000000",
                "bold": False,
                "font": "Arial",
                "align": 1,
                "confidence": 0.99,
            }],
            image_to_ppt.np.zeros((10, 16), dtype=image_to_ppt.np.uint8),
        )

    def fake_slide(path: Path, target: Path, text_analysis: dict) -> dict:
        calls["visual"] += 1
        assert path.resolve().is_relative_to(target.resolve())
        assert Path(text_analysis["mask_path"]).is_file()
        components_dir = target / "components"
        masks_dir = target / "element-masks"
        semantic_masks_dir = target / "semantic-masks"
        components_dir.mkdir()
        masks_dir.mkdir()
        semantic_masks_dir.mkdir()
        component_paths = []
        mask_paths = []
        semantic_mask_paths = []
        for index, color in enumerate(("red", "blue")):
            component_path = components_dir / f"component_{index:04d}.png"
            Image.new("RGBA", (4, 4), color).save(component_path)
            component_paths.append(component_path)
            mask_path = masks_dir / f"{index:04d}.png"
            semantic_mask_path = semantic_masks_dir / f"{index:04d}.png"
            element_mask = image_to_ppt.np.zeros((10, 16), dtype=image_to_ppt.np.uint8)
            semantic_mask = image_to_ppt.np.zeros((10, 16), dtype=image_to_ppt.np.uint8)
            if index == 0:
                element_mask[1:5, 1:5] = 255
                semantic_mask[0:6, 0:6] = 255
            else:
                element_mask[5:9, 10:15] = 255
                semantic_mask[4:10, 9:16] = 255
            Image.fromarray(element_mask, mode="L").save(mask_path)
            Image.fromarray(semantic_mask, mode="L").save(semantic_mask_path)
            mask_paths.append(mask_path)
            semantic_mask_paths.append(semantic_mask_path)
        foreground_evidence_path = target / "foreground-evidence-mask.png"
        foreground_evidence = image_to_ppt.np.zeros(
            (10, 16), dtype=image_to_ppt.np.uint8
        )
        for path in semantic_mask_paths:
            with Image.open(path) as stored_mask:
                foreground_evidence |= image_to_ppt.np.asarray(
                    stored_mask.convert("L")
                )
        Image.fromarray(foreground_evidence, mode="L").save(
            foreground_evidence_path
        )
        background_path = target / "background-original.png"
        Image.new("RGB", (16, 10), "white").save(background_path)
        Image.new("L", (16, 10), 0).save(target / "background-removal-mask.png")
        Image.new("RGB", (16, 10), "black").save(
            target / "background-difference.png"
        )
        text_clean_path = target / "text-clean.png"
        Image.new("RGB", (16, 10), "white").save(text_clean_path)
        return {
            "background_path": str(background_path),
            "background_original_path": str(background_path),
            "background_widescreen_path": str(background_path),
            "components": [
                {
                    "path": str(component_path),
                    "x": index * 4,
                    "y": 0,
                    "w": 4,
                    "h": 4,
                    "area": 16,
                    "z_index": index,
                }
                for index, component_path in enumerate(component_paths)
            ],
            "text_items": text_analysis["items"],
            "img_width": 16,
            "img_height": 10,
            "canvas_width": 16,
            "canvas_height": 10,
            "content_offset_x": 0,
            "content_offset_y": 0,
            "widescreen_background_method": "identity",
            "original_image_path": str(path),
            "_work_dir": str(target),
            "_text_mask_path": text_analysis["mask_path"],
            "_text_clean_path": str(text_clean_path),
            "_element_mask_paths": [str(mask) for mask in mask_paths],
            "_semantic_mask_paths": [str(mask) for mask in semantic_mask_paths],
            "_foreground_evidence_mask_path": str(foreground_evidence_path),
        }

    def fake_process(
        path,
        target,
        object_detector,
        mask_generator,
        lang,
        *,
        text_analysis,
        defer_quality,
        **kwargs,
    ):
        assert defer_quality is True
        return fake_slide(Path(path), Path(target), text_analysis)

    def fake_process_isolated(path, target, lang, text_analysis):
        if before_visual_worker is not None:
            before_visual_worker()
        return fake_slide(Path(path), Path(target), text_analysis)

    monkeypatch.setattr(image_to_ppt, "detect_text", fake_detect)
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", lambda path: object())
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam"))
    monkeypatch.setattr(image_to_ppt, "_process_image", fake_process)
    monkeypatch.setattr(image_to_ppt, "_process_image_isolated", fake_process_isolated)
    monkeypatch.setattr(image_to_ppt, "_release_visual_resources", lambda: None)
    monkeypatch.setattr(
        image_to_ppt,
        "_finalize_slide_quality",
        lambda *args, **kwargs: pytest.fail("prepare must not finalize quality"),
    )
    if before_prepare is not None:
        before_prepare(work_dir)
    prepared = image_to_ppt.prepare_component_layers(
        source_path,
        work_dir,
        lang="en",
        resource_isolation=resource_isolation,
    )
    assert calls == {"ocr": 1, "visual": 1}
    return prepared, work_dir


def _accepted_component_layers(prepared: dict) -> dict:
    return {
        "components": prepared["components"],
        "element_masks": prepared["_element_mask_paths"],
    }


def _write_prepared_manifest(state_path: Path, manifest: dict) -> None:
    state_path.write_text(json.dumps(manifest), encoding="utf-8")
    state_hash = hashlib.sha256(state_path.read_bytes()).hexdigest()
    state_path.with_name("prepared_page.sha256").write_bytes(
        f"{state_hash}\n".encode("ascii"),
    )


def _replace_prepared_mask_asset(
    prepared: dict,
    work_dir: Path,
    asset_name: str,
    *,
    size: tuple[int, int],
    color: int,
) -> tuple[Path, str]:
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    record = manifest["assets"][asset_name][0]
    asset_path = work_dir / record["path"]
    Image.new("L", size, color).save(asset_path)
    record["sha256"] = hashlib.sha256(asset_path.read_bytes()).hexdigest()
    _write_prepared_manifest(state_path, manifest)
    return state_path, record["path"]


@pytest.mark.parametrize("write_fails", [False, True])
def test_prepare_component_layers_isolated_releases_cleanup_arrays_before_worker(
    tmp_path: Path,
    monkeypatch,
    write_fails: bool,
) -> None:
    references = {}
    gc_calls = []

    def fake_load_rgb(path):
        source = image_to_ppt.np.full((10, 16, 3), 255, dtype=image_to_ppt.np.uint8)
        references["source_image"] = weakref.ref(source)
        return source

    def fake_cleanup_mask(source_image, stored_mask, text_items):
        references["stored_mask"] = weakref.ref(stored_mask)
        removal_mask = image_to_ppt.np.zeros((10, 16), dtype=image_to_ppt.np.uint8)
        references["removal_mask"] = weakref.ref(removal_mask)
        return removal_mask

    def fake_repair(source_image, removal_mask, **kwargs):
        text_clean = source_image.copy()
        references["text_clean"] = weakref.ref(text_clean)
        return text_clean

    def fake_save(path, image):
        if write_fails:
            raise RuntimeError("text-clean write failed")
        Image.fromarray(image).save(path)

    def assert_released_before_worker():
        assert gc_calls == ["collect"]
        assert set(references) == {
            "source_image",
            "stored_mask",
            "removal_mask",
            "text_clean",
        }
        assert all(reference() is None for reference in references.values())

    monkeypatch.setattr(image_to_ppt, "_load_rgb", fake_load_rgb)
    monkeypatch.setattr(image_to_ppt, "_build_text_cleanup_mask", fake_cleanup_mask)
    monkeypatch.setattr(image_to_ppt, "_repair_text_background", fake_repair)
    monkeypatch.setattr(image_to_ppt, "_save_rgb", fake_save)
    monkeypatch.setattr(
        image_to_ppt,
        "_isolated_large_inpainter",
        lambda work_dir: object(),
    )
    monkeypatch.setattr(
        image_to_ppt.gc,
        "collect",
        lambda: gc_calls.append("collect") or 0,
    )

    if write_fails:
        with pytest.raises(RuntimeError, match="text-clean write failed"):
            _prepare_component_layers_fixture(
                tmp_path,
                monkeypatch,
                resource_isolation=True,
            )
        assert gc_calls == ["collect"]
        assert all(reference() is None for reference in references.values())
    else:
        _prepare_component_layers_fixture(
            tmp_path,
            monkeypatch,
            resource_isolation=True,
            before_visual_worker=assert_released_before_worker,
        )


def test_prepare_component_layers_persists_initial_components_without_quality(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)

    assert prepared["phase"] == "initial_layers"
    assert prepared["initial_component_count"] == 2
    assert len(prepared["components"]) == 2
    assert Path(prepared["state_path"]) == (work_dir / "prepared_page.json").resolve()
    manifest = json.loads(Path(prepared["state_path"]).read_text(encoding="utf-8"))
    assert manifest["schema_version"] == 5
    cleanup_path = Path(prepared["_text_cleanup_mask_path"])
    assert cleanup_path == (work_dir / "text-clean-removal-mask.png").resolve()
    assert manifest["assets"]["text_cleanup_mask"]["path"] == cleanup_path.name
    assert manifest["initial_diagnostics"] == []
    assert manifest["phase"] == "initial_layers"
    assert len(manifest["assets"]["semantic_masks"]) == 2
    assert Path(prepared["_foreground_evidence_mask_path"]).is_file()
    assert manifest["assets"]["foreground_evidence_mask"]["path"] == (
        "foreground-evidence-mask.png"
    )
    sidecar_path = work_dir / "prepared_page.sha256"
    assert sidecar_path.read_text(encoding="ascii") == (
        hashlib.sha256(Path(prepared["state_path"]).read_bytes()).hexdigest()
        + "\n"
    )
    assert not (work_dir / ".prepared_page.json.tmp").exists()

    def assert_asset(record: dict) -> None:
        path = Path(record["path"])
        assert not path.is_absolute()
        assert ".." not in path.parts
        absolute = work_dir / path
        assert hashlib.sha256(absolute.read_bytes()).hexdigest() == record["sha256"]

    for key, value in manifest["assets"].items():
        if value is None:
            continue
        if isinstance(value, list):
            for record in value:
                assert_asset(record)
        else:
            assert_asset(value)
    for component in manifest["components"]:
            assert_asset(component["asset"])


def test_load_component_layers_rejects_foreground_evidence_not_matching_semantics(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    evidence_path = work_dir / manifest["assets"]["foreground_evidence_mask"]["path"]
    Image.new("L", (16, 10), 0).save(evidence_path)
    manifest["assets"]["foreground_evidence_mask"]["sha256"] = hashlib.sha256(
        evidence_path.read_bytes()
    ).hexdigest()
    _write_prepared_manifest(state_path, manifest)

    with pytest.raises(
        ValueError,
        match="foreground evidence does not match semantic masks",
    ):
        image_to_ppt.load_component_layers(state_path)


def test_prepare_component_layers_authenticates_text_cleanup_mask(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(
        tmp_path, monkeypatch, resource_isolation=True,
    )

    cleanup_path = Path(prepared["_text_cleanup_mask_path"])
    assert cleanup_path == (work_dir / "text-clean-removal-mask.png").resolve()
    manifest = json.loads(Path(prepared["state_path"]).read_text(encoding="utf-8"))
    record = manifest["assets"]["text_cleanup_mask"]
    assert record["path"] == "text-clean-removal-mask.png"
    assert record["sha256"] == hashlib.sha256(cleanup_path.read_bytes()).hexdigest()


def test_prepare_component_layers_overwrites_stale_text_cleanup_mask(
    tmp_path: Path,
    monkeypatch,
) -> None:
    def write_stale_mask(work_dir: Path) -> None:
        work_dir.mkdir()
        Image.new("L", (16, 10), 77).save(
            work_dir / "text-clean-removal-mask.png"
        )

    monkeypatch.setattr(
        image_to_ppt,
        "_build_text_cleanup_mask",
        lambda source, mask, items: image_to_ppt.np.full(
            (10, 16), 255, dtype=image_to_ppt.np.uint8
        ),
    )
    prepared, _ = _prepare_component_layers_fixture(
        tmp_path,
        monkeypatch,
        before_prepare=write_stale_mask,
    )

    with Image.open(prepared["_text_cleanup_mask_path"]) as cleanup:
        assert set(cleanup.getdata()) == {255}


def test_prepare_component_layers_does_not_authorize_state_replaced_after_publish(
    tmp_path: Path,
    monkeypatch,
) -> None:
    real_atomic_write = image_to_ppt._atomic_write_prepared_text
    replaced = False

    def replace_after_state_publish(work_dir, path, content, **kwargs):
        nonlocal replaced
        real_atomic_write(work_dir, path, content, **kwargs)
        if path.name == "prepared_page.json" and not replaced:
            manifest = json.loads(path.read_text(encoding="utf-8"))
            manifest["resource_isolation"] = True
            path.write_text(
                json.dumps(manifest, ensure_ascii=False, indent=2),
                encoding="utf-8",
            )
            replaced = True

    monkeypatch.setattr(
        image_to_ppt,
        "_atomic_write_prepared_text",
        replace_after_state_publish,
    )

    with pytest.raises(ValueError, match="state sha256"):
        _prepare_component_layers_fixture(tmp_path, monkeypatch)

    assert replaced is True


def test_prepare_component_layers_rejects_preexisting_asset_reparse_point(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source_path = tmp_path / "source.png"
    Image.new("RGB", (4, 4), "red").save(source_path)
    work_dir = tmp_path / "prepared"
    work_dir.mkdir()
    prepared_target = work_dir / "source-image.png"
    Image.new("RGB", (4, 4), "blue").save(prepared_target)
    original_target = prepared_target.read_bytes()
    real_lstat = image_to_ppt.os.lstat

    def fake_lstat(path):
        status = real_lstat(path)
        if Path(path) == prepared_target:
            return types.SimpleNamespace(
                st_mode=status.st_mode,
                st_file_attributes=getattr(
                    image_to_ppt.stat,
                    "FILE_ATTRIBUTE_REPARSE_POINT",
                    0x400,
                ),
            )
        return status

    monkeypatch.setattr(image_to_ppt.os, "lstat", fake_lstat)
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: pytest.fail("OCR must not run"),
    )

    with pytest.raises(ValueError, match="link or reparse"):
        image_to_ppt.prepare_component_layers(
            source_path,
            work_dir,
            lang="en",
            resource_isolation=False,
        )

    assert prepared_target.read_bytes() == original_target


def test_prepare_component_layers_rejects_preexisting_source_hardlink(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source_path = tmp_path / "source.png"
    external_path = tmp_path / "external.png"
    Image.new("RGB", (4, 4), "red").save(source_path)
    Image.new("RGB", (4, 4), "blue").save(external_path)
    original_external = external_path.read_bytes()
    work_dir = tmp_path / "prepared"
    work_dir.mkdir()
    os.link(external_path, work_dir / "source-image.png")
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: pytest.fail("OCR must not run"),
    )

    with pytest.raises(ValueError, match="hardlink|single-link"):
        image_to_ppt.prepare_component_layers(
            source_path,
            work_dir,
            lang="en",
            resource_isolation=False,
        )

    assert external_path.read_bytes() == original_external


def test_prepare_component_layers_does_not_create_through_linked_parent(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "source.png"
    Image.new("RGB", (4, 4), "red").save(source_path)
    external = tmp_path / "external"
    external.mkdir()
    linked_parent = tmp_path / "linked"
    try:
        linked_parent.symlink_to(external, target_is_directory=True)
    except OSError as error:
        if os.name != "nt":
            pytest.skip(f"directory links are unavailable: {error}")
        junction = subprocess.run(
            ["cmd", "/c", "mklink", "/J", str(linked_parent), str(external)],
            capture_output=True,
            text=True,
            check=False,
        )
        if junction.returncode:
            pytest.skip(f"directory links are unavailable: {error}")

    with pytest.raises(ValueError, match="link|resolve"):
        image_to_ppt.prepare_component_layers(
            source_path,
            linked_parent / "new-work-dir",
            lang="en",
            resource_isolation=False,
        )

    assert not (external / "new-work-dir").exists()


def test_prepare_component_layers_ocr_cleanup_does_not_mask_primary_error(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source_path = tmp_path / "source.png"
    Image.new("RGB", (4, 4), "red").save(source_path)
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            RuntimeError("OCR primary failure")
        ),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "close_ocr_engines",
        lambda: (_ for _ in ()).throw(OSError("OCR cleanup failure")),
    )

    with pytest.raises(RuntimeError, match="OCR primary failure"):
        image_to_ppt.prepare_component_layers(
            source_path,
            tmp_path / "prepared",
            lang="en",
            resource_isolation=False,
        )


def test_prepare_component_layers_visual_cleanup_does_not_mask_primary_error(
    tmp_path: Path,
    monkeypatch,
) -> None:
    source_path = tmp_path / "source.png"
    Image.new("RGB", (4, 4), "red").save(source_path)
    monkeypatch.setattr(
        image_to_ppt,
        "detect_text",
        lambda *args, **kwargs: (
            [],
            image_to_ppt.np.zeros((4, 4), dtype=image_to_ppt.np.uint8),
        ),
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(image_to_ppt, "create_object_detector", lambda: object())
    monkeypatch.setattr(image_to_ppt, "create_sam_generator", lambda path: object())
    monkeypatch.setattr(image_to_ppt, "resolve_sam_checkpoint", lambda: Path("sam"))
    monkeypatch.setattr(
        image_to_ppt,
        "_process_image",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            RuntimeError("visual primary failure")
        ),
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_release_visual_resources",
        lambda: (_ for _ in ()).throw(OSError("visual cleanup failure")),
    )

    with pytest.raises(RuntimeError, match="visual primary failure"):
        image_to_ppt.prepare_component_layers(
            source_path,
            tmp_path / "prepared",
            lang="en",
            resource_isolation=False,
        )


def test_load_component_layers_recovers_absolute_owned_paths(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = prepared["state_path"]
    del prepared

    restored = image_to_ppt.load_component_layers(state_path)

    assert restored["phase"] == "initial_layers"
    assert restored["initial_component_count"] == 2
    path_values = [
        restored["original_image_path"],
        restored["background_original_path"],
        restored["background_widescreen_path"],
        restored["background_removal_mask_path"],
        restored["background_difference_path"],
        restored["_text_mask_path"],
        restored["_text_clean_path"],
        *restored["_element_mask_paths"],
        *restored["_semantic_mask_paths"],
        *(component["path"] for component in restored["components"]),
    ]
    assert all(Path(value).is_absolute() for value in path_values)
    assert all(Path(value).is_relative_to(work_dir.resolve()) for value in path_values)


def test_load_component_layers_v2_rejects_mismatched_mask_counts(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    manifest["assets"]["semantic_masks"].pop()
    _write_prepared_manifest(state_path, manifest)

    with pytest.raises(ValueError, match="mask counts"):
        image_to_ppt.load_component_layers(state_path)


def test_load_component_layers_v2_rejects_child_outside_parent(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    child_record = manifest["assets"]["element_masks"][0]
    child_path = work_dir / child_record["path"]
    with Image.open(child_path) as stored_child:
        child = image_to_ppt.np.asarray(stored_child.convert("L")).copy()
    child[9, 15] = 255
    Image.fromarray(child, mode="L").save(child_path)
    child_record["sha256"] = hashlib.sha256(child_path.read_bytes()).hexdigest()
    _write_prepared_manifest(state_path, manifest)

    with pytest.raises(ValueError, match="inside its parent"):
        image_to_ppt.load_component_layers(state_path)


@pytest.mark.parametrize("asset_name", ["element_masks", "semantic_masks"])
def test_load_component_layers_v2_rejects_empty_masks_with_one_read(
    tmp_path: Path,
    monkeypatch,
    asset_name: str,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path, target_path = _replace_prepared_mask_asset(
        prepared,
        work_dir,
        asset_name,
        size=(16, 10),
        color=0,
    )
    real_read = image_to_ppt._read_prepared_asset_bytes
    target_reads = 0

    def track_target_reads(asset_work_dir, record, label):
        nonlocal target_reads
        if record["path"] == target_path:
            target_reads += 1
        return real_read(asset_work_dir, record, label)

    monkeypatch.setattr(
        image_to_ppt,
        "_read_prepared_asset_bytes",
        track_target_reads,
    )

    with pytest.raises(ValueError, match="non-empty"):
        image_to_ppt.load_component_layers(state_path)

    assert target_reads == 1


@pytest.mark.parametrize("asset_name", ["element_masks", "semantic_masks"])
def test_load_component_layers_v2_rejects_wrong_size_before_convert(
    tmp_path: Path,
    monkeypatch,
    asset_name: str,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path, _ = _replace_prepared_mask_asset(
        prepared,
        work_dir,
        asset_name,
        size=(17, 10),
        color=255,
    )
    real_convert = Image.Image.convert

    def reject_wrong_size_convert(image, *args, **kwargs):
        if image.size != (16, 10):
            pytest.fail("wrong-size mask reached convert")
        return real_convert(image, *args, **kwargs)

    monkeypatch.setattr(Image.Image, "convert", reject_wrong_size_convert)

    with pytest.raises(ValueError, match="dimensions"):
        image_to_ppt.load_component_layers(state_path)


def test_load_component_layers_v2_validates_authenticated_mask_bytes(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    target_path = manifest["assets"]["element_masks"][0]["path"]
    replacement_path = work_dir / "replacement-child.png"
    replacement = image_to_ppt.np.zeros((10, 16), dtype=image_to_ppt.np.uint8)
    replacement[9, 15] = 255
    Image.fromarray(replacement, mode="L").save(replacement_path)
    real_read = image_to_ppt._read_prepared_asset_bytes
    target_reads = 0

    def replace_after_authenticated_read(asset_work_dir, record, label):
        nonlocal target_reads
        owned, content = real_read(asset_work_dir, record, label)
        if record["path"] == target_path:
            target_reads += 1
            if target_reads == 1:
                os.replace(replacement_path, owned)
        return owned, content

    monkeypatch.setattr(
        image_to_ppt,
        "_read_prepared_asset_bytes",
        replace_after_authenticated_read,
    )

    restored = image_to_ppt.load_component_layers(state_path)

    assert target_reads == 1
    assert restored["_element_mask_paths"][0] == str(work_dir / target_path)


def test_load_component_layers_reads_v1_without_fabricating_semantic_masks(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    assert manifest["schema_version"] == 5
    manifest["schema_version"] = 1
    manifest.pop("initial_diagnostics")
    manifest["assets"].pop("semantic_masks")
    manifest["assets"].pop("text_cleanup_mask")
    manifest["assets"].pop("foreground_evidence_mask")
    _write_prepared_manifest(state_path, manifest)

    restored = image_to_ppt.load_component_layers(state_path)

    assert len(restored["_element_mask_paths"]) == 2
    assert "_semantic_mask_paths" not in restored


def test_load_component_layers_does_not_create_missing_work_directory(
    tmp_path: Path,
) -> None:
    state_path = tmp_path / "missing-work" / "prepared_page.json"

    with pytest.raises(ValueError, match="work directory.*does not exist"):
        image_to_ppt.load_component_layers(state_path)

    assert not state_path.parent.exists()


def test_load_component_layers_rejects_json_changed_without_sidecar_update(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    manifest["phase"] = "changed"
    state_path.write_text(json.dumps(manifest), encoding="utf-8")

    with pytest.raises(ValueError, match="state sha256"):
        image_to_ppt.load_component_layers(state_path)


def test_load_component_layers_never_parses_state_modified_after_hash(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    original_resource_isolation = prepared["_resource_isolation"]
    replacement_manifest = json.loads(state_path.read_text(encoding="utf-8"))
    replacement_manifest["resource_isolation"] = not original_resource_isolation
    replacement_bytes = json.dumps(replacement_manifest).encode("utf-8")
    real_open = Path.open
    modified = False

    def modify_state() -> None:
        nonlocal modified
        if modified:
            return
        modified = True
        with real_open(state_path, "wb") as destination:
            destination.write(replacement_bytes)

    class StateReader:
        def __init__(self, source) -> None:
            self.source = source

        def __enter__(self):
            self.source.__enter__()
            return self

        def __exit__(self, *args):
            return self.source.__exit__(*args)

        def __getattr__(self, name):
            return getattr(self.source, name)

        def read(self, *args, **kwargs):
            data = self.source.read(*args, **kwargs)
            requested_size = args[0] if args else -1
            if requested_size == -1 or data == b"":
                modify_state()
            return data

    def controlled_open(path, *args, **kwargs):
        source = real_open(path, *args, **kwargs)
        mode = args[0] if args else kwargs.get("mode", "r")
        if Path(path) == state_path and mode == "rb" and not modified:
            return StateReader(source)
        return source

    monkeypatch.setattr(Path, "open", controlled_open)

    try:
        restored = image_to_ppt.load_component_layers(state_path)
    except ValueError:
        pass
    else:
        assert restored["_resource_isolation"] is original_resource_isolation

    assert modified is True


@pytest.mark.parametrize("sidecar", ["", "0" * 64, "A" * 64 + "\n"])
def test_load_component_layers_rejects_invalid_state_sidecar(
    tmp_path: Path,
    monkeypatch,
    sidecar: str,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    (work_dir / "prepared_page.sha256").write_text(sidecar, encoding="ascii")

    with pytest.raises(ValueError, match="sidecar"):
        image_to_ppt.load_component_layers(prepared["state_path"])


@pytest.mark.parametrize(
    "tamper",
    [
        lambda manifest: manifest.update({"schema_version": True}),
        lambda manifest: manifest["dimensions"].update({"img_width": True}),
        lambda manifest: manifest["components"][0]["metadata"].update({"x": True}),
    ],
)
def test_load_component_layers_rejects_bool_for_integer_fields(
    tmp_path: Path,
    monkeypatch,
    tamper,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    tamper(manifest)
    _write_prepared_manifest(state_path, manifest)

    with pytest.raises(ValueError, match="schema_version|dimension|metadata"):
        image_to_ppt.load_component_layers(state_path)


@pytest.mark.parametrize("target_name", ["state", "sidecar", "source"])
def test_load_component_layers_rejects_hardlinked_owned_file(
    tmp_path: Path,
    monkeypatch,
    target_name: str,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    target = (
        Path(prepared["state_path"])
        if target_name == "state"
        else work_dir / "prepared_page.sha256"
        if target_name == "sidecar"
        else Path(prepared["original_image_path"])
    )
    os.link(target, work_dir / f"{target_name}-hardlink")

    with pytest.raises(ValueError, match="hardlink|single-link"):
        image_to_ppt.load_component_layers(prepared["state_path"])


@pytest.mark.parametrize(
    "asset_name,index",
    [
        ("source_image", None),
        ("ocr_mask", None),
        ("element_masks", 0),
        ("semantic_masks", 0),
        ("foreground_evidence_mask", None),
        ("background_original", None),
        ("background_removal_mask", None),
        ("background_difference", None),
    ],
)
def test_load_component_layers_rejects_tampered_assets(
    tmp_path: Path,
    monkeypatch,
    asset_name: str,
    index: int | None,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    manifest = json.loads(Path(prepared["state_path"]).read_text(encoding="utf-8"))
    record = manifest["assets"][asset_name]
    if index is not None:
        record = record[index]
    (work_dir / record["path"]).write_bytes(b"tampered")

    with pytest.raises(ValueError, match="sha256"):
        image_to_ppt.load_component_layers(prepared["state_path"])


def test_load_component_layers_rejects_tampered_component_rgba(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    manifest = json.loads(Path(prepared["state_path"]).read_text(encoding="utf-8"))
    component = manifest["components"][0]["asset"]
    (work_dir / component["path"]).write_bytes(b"tampered")

    with pytest.raises(ValueError, match="sha256"):
        image_to_ppt.load_component_layers(prepared["state_path"])


@pytest.mark.parametrize("malicious_path", ["../outside.png", "C:/outside.png"])
def test_load_component_layers_rejects_json_asset_path_escape(
    tmp_path: Path,
    monkeypatch,
    malicious_path: str,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    manifest["assets"]["source_image"]["path"] = malicious_path
    _write_prepared_manifest(state_path, manifest)

    with pytest.raises(ValueError, match="asset path"):
        image_to_ppt.load_component_layers(state_path)


@pytest.mark.parametrize(
    "tamper,match",
    [
        (
            lambda manifest: manifest["components"][0]["metadata"].update(
                {"unexpected_path": "C:/outside-secret.png"}
            ),
            "component metadata",
        ),
        (
            lambda manifest: manifest["text_items"][0].update(
                {"box": [-1, 1, 3, 2]}
            ),
            "text item",
        ),
        (
            lambda manifest: manifest["dimensions"].update({"img_width": 0}),
            "dimension",
        ),
    ],
)
def test_load_component_layers_rejects_invalid_recovery_fields(
    tmp_path: Path,
    monkeypatch,
    tamper,
    match: str,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    state_path = Path(prepared["state_path"])
    manifest = json.loads(state_path.read_text(encoding="utf-8"))
    tamper(manifest)
    _write_prepared_manifest(state_path, manifest)

    with pytest.raises(ValueError, match=match):
        image_to_ppt.load_component_layers(state_path)


def test_agent_managed_finalize_rejects_component_text_overlap_without_mutation(
    tmp_path: Path,
    monkeypatch,
) -> None:
    real_finalize = image_to_ppt._finalize_slide_quality
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    component_records = [dict(component) for component in prepared["components"]]
    component_hashes = [
        hashlib.sha256(Path(component["path"]).read_bytes()).hexdigest()
        for component in prepared["components"]
    ]
    monkeypatch.setattr(
        image_to_ppt,
        "_has_component_text_overlap",
        lambda masks, text_mask: True,
    )
    monkeypatch.setattr(image_to_ppt, "_finalize_slide_quality", real_finalize)

    with pytest.raises(
        image_to_ppt.VisualSegmentationError,
        match="component_text_overlap",
    ):
        image_to_ppt.finalize_component_layers(
            prepared,
            _accepted_component_layers(prepared),
            lang="en",
        )

    assert prepared["components"] == component_records
    assert [
        hashlib.sha256(Path(component["path"]).read_bytes()).hexdigest()
        for component in prepared["components"]
    ] == component_hashes


def test_agent_managed_finalize_stages_components_before_quality_failure(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    original_path = Path(prepared["components"][0]["path"])
    original_hash = hashlib.sha256(original_path.read_bytes()).hexdigest()
    cleanup_calls = []
    monkeypatch.setattr(
        image_to_ppt,
        "_has_component_text_overlap",
        lambda masks, text_mask: False,
    )

    def fake_finalize(slide_data, lang, **kwargs):
        assert kwargs == {
            "_resource_isolation": False,
        }
        staged_path = Path(slide_data["components"][0]["path"])
        assert staged_path != original_path
        staged_dir = staged_path.parent
        staged_asset_paths = [
            slide_data["original_image_path"],
            slide_data["background_path"],
            slide_data["background_original_path"],
            slide_data["background_widescreen_path"],
            slide_data["background_removal_mask_path"],
            slide_data["background_difference_path"],
            slide_data["_text_mask_path"],
            slide_data["_text_clean_path"],
            *slide_data["_element_mask_paths"],
            *(component["path"] for component in slide_data["components"]),
        ]
        assert all(Path(path).parent == staged_dir for path in staged_asset_paths)
        staged_path.write_bytes(b"mutated during quality")
        raise image_to_ppt.VisualSegmentationError(
            "agent-managed quality failed: background_residual"
        )

    monkeypatch.setattr(image_to_ppt, "_finalize_slide_quality", fake_finalize)
    monkeypatch.setattr(
        image_to_ppt,
        "close_ocr_engines",
        lambda: cleanup_calls.append("close"),
    )

    with pytest.raises(
        image_to_ppt.VisualSegmentationError,
        match="background_residual",
    ):
        image_to_ppt.finalize_component_layers(
            prepared,
            _accepted_component_layers(prepared),
            lang="en",
        )

    assert prepared["initial_component_count"] == 2
    assert len(prepared["components"]) == 2
    assert hashlib.sha256(original_path.read_bytes()).hexdigest() == original_hash
    assert cleanup_calls == ["close"]
    assert not list(Path(prepared["_work_dir"]).glob("quality-components-*"))


def test_agent_managed_finalize_success_keeps_staging_without_duplicate_loads(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    prepared_components = [dict(component) for component in prepared["components"]]
    prepared_paths = {
        Path(path)
        for path in (
            prepared["original_image_path"],
            prepared["background_original_path"],
            prepared["background_widescreen_path"],
            prepared["background_removal_mask_path"],
            prepared["background_difference_path"],
            prepared["_text_mask_path"],
            prepared["_text_clean_path"],
            *prepared["_element_mask_paths"],
            *prepared["_semantic_mask_paths"],
            *(component["path"] for component in prepared["components"]),
        )
    }
    prepared_hashes = {
        path: hashlib.sha256(path.read_bytes()).hexdigest()
        for path in prepared_paths
    }
    load_calls = []
    open_calls = []
    cleanup_calls = []
    staged_paths = []
    real_load_rgb = image_to_ppt._load_rgb
    real_image_open = image_to_ppt.Image.open

    def tracked_load_rgb(path):
        load_calls.append(Path(path))
        return real_load_rgb(path)

    def tracked_image_open(path, *args, **kwargs):
        open_calls.append(path)
        return real_image_open(path, *args, **kwargs)

    def fake_finalize(slide_data, lang, **kwargs):
        staged_paths.extend(
            Path(path)
            for path in (
                slide_data["original_image_path"],
                slide_data["background_path"],
                slide_data["background_original_path"],
                slide_data["background_widescreen_path"],
                slide_data["background_removal_mask_path"],
                slide_data["background_difference_path"],
                slide_data["_text_mask_path"],
                slide_data["_text_clean_path"],
                *slide_data["_element_mask_paths"],
                *(component["path"] for component in slide_data["components"]),
            )
        )
        return slide_data

    monkeypatch.setattr(image_to_ppt, "_load_rgb", tracked_load_rgb)
    monkeypatch.setattr(image_to_ppt.Image, "open", tracked_image_open)
    monkeypatch.setattr(
        image_to_ppt,
        "_has_component_text_overlap",
        lambda masks, text_mask: False,
    )
    monkeypatch.setattr(image_to_ppt, "_finalize_slide_quality", fake_finalize)
    monkeypatch.setattr(
        image_to_ppt,
        "close_ocr_engines",
        lambda: cleanup_calls.append("close"),
    )

    result = image_to_ppt.finalize_component_layers(
        prepared,
        _accepted_component_layers(prepared),
        lang="en",
    )

    assert result["phase"] == "quality_accepted"
    assert result["initial_component_count"] == len(prepared_components)
    assert prepared["components"] == prepared_components
    assert cleanup_calls == ["close"]
    assert load_calls == []
    assert len(open_calls) == (
        len(prepared["_element_mask_paths"])
        + len(prepared["_semantic_mask_paths"])
        + 1
    )
    assert staged_paths
    staging_dirs = {path.parent for path in staged_paths}
    assert len(staging_dirs) == 1
    assert next(iter(staging_dirs)).name.startswith("quality-components-")
    assert all(path.is_file() for path in staged_paths)
    assert all(Path(component["path"]).is_file() for component in result["components"])
    assert {
        path: hashlib.sha256(path.read_bytes()).hexdigest()
        for path in prepared_paths
    } == prepared_hashes


@pytest.mark.parametrize(
    "accepted_factory",
    [
        lambda prepared: None,
        lambda prepared: prepared["components"],
        lambda prepared: {
            "components": prepared["components"],
            "_element_mask_paths": prepared["_element_mask_paths"],
        },
        lambda prepared: {"components": prepared["components"]},
    ],
)
def test_agent_managed_finalize_rejects_legacy_accepted_shapes(
    tmp_path: Path,
    monkeypatch,
    accepted_factory,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    monkeypatch.setattr(
        image_to_ppt,
        "_has_component_text_overlap",
        lambda masks, text_mask: False,
    )

    with pytest.raises(ValueError, match="accepted"):
        image_to_ppt.finalize_component_layers(
            prepared,
            accepted_factory(prepared),
            lang="en",
        )


def test_agent_managed_finalize_rejects_noncurrent_owned_component(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    replacement = work_dir / "replacement.png"
    replacement.write_bytes(Path(prepared["components"][0]["path"]).read_bytes())
    accepted = _accepted_component_layers(prepared)
    accepted["components"] = [
        {**accepted["components"][0], "path": str(replacement)},
        accepted["components"][1],
    ]

    with pytest.raises(ValueError, match="current prepared state"):
        image_to_ppt.finalize_component_layers(prepared, accepted, lang="en")


def test_agent_managed_finalize_rejects_noncurrent_owned_element_mask(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    replacement = work_dir / "replacement-mask.png"
    replacement.write_bytes(Path(prepared["_element_mask_paths"][0]).read_bytes())
    accepted = _accepted_component_layers(prepared)
    accepted["element_masks"] = [
        str(replacement),
        *accepted["element_masks"][1:],
    ]

    with pytest.raises(ValueError, match="current prepared state"):
        image_to_ppt.finalize_component_layers(prepared, accepted, lang="en")


def test_agent_managed_finalize_reloads_state_before_consuming_assets(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    accepted = _accepted_component_layers(prepared)
    Path(prepared["components"][0]["path"]).write_bytes(b"replaced")

    with pytest.raises(ValueError, match="sha256"):
        image_to_ppt.finalize_component_layers(prepared, accepted, lang="en")


@pytest.mark.parametrize(
    "target_label,target_path,mode,size,color",
    [
        (
            "source_image",
            lambda prepared: prepared["original_image_path"],
            "RGB",
            (16, 10),
            "green",
        ),
        (
            "background_original",
            lambda prepared: prepared["background_original_path"],
            "RGB",
            (16, 10),
            "green",
        ),
        (
            "ocr_mask",
            lambda prepared: prepared["_text_mask_path"],
            "L",
            (16, 10),
            255,
        ),
        (
            "text_clean",
            lambda prepared: prepared["_text_clean_path"],
            "RGB",
            (16, 10),
            "green",
        ),
        (
            "element mask",
            lambda prepared: prepared["_element_mask_paths"][0],
            "L",
            (16, 10),
            127,
        ),
        (
            "component RGBA",
            lambda prepared: prepared["components"][0]["path"],
            "RGBA",
            (4, 4),
            "green",
        ),
    ],
)
def test_agent_managed_finalize_rejects_asset_replaced_after_load_validation(
    tmp_path: Path,
    monkeypatch,
    target_label: str,
    target_path,
    mode: str,
    size: tuple[int, int],
    color,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    if target_label == "background_original":
        state_path = Path(prepared["state_path"])
        manifest = json.loads(state_path.read_text(encoding="utf-8"))
        original_record = manifest["assets"]["background_original"]
        widescreen_path = work_dir / "background-widescreen.png"
        widescreen_path.write_bytes(
            (work_dir / original_record["path"]).read_bytes()
        )
        manifest["assets"]["background_widescreen"] = {
            "path": widescreen_path.relative_to(work_dir).as_posix(),
            "sha256": hashlib.sha256(widescreen_path.read_bytes()).hexdigest(),
        }
        _write_prepared_manifest(state_path, manifest)
        prepared = image_to_ppt.load_component_layers(state_path)

    owned_target = Path(target_path(prepared))
    replacement_path = work_dir / f"replacement-{target_label.replace(' ', '-')}.png"
    Image.new(mode, size, color).save(replacement_path)
    real_read_asset = image_to_ppt._read_prepared_asset_bytes
    replaced = False

    def replace_after_validation(asset_work_dir, record, label):
        nonlocal replaced
        loaded_path, content = real_read_asset(asset_work_dir, record, label)
        if label == target_label and loaded_path == owned_target and not replaced:
            replaced = True
            os.replace(replacement_path, owned_target)
        return loaded_path, content

    monkeypatch.setattr(
        image_to_ppt,
        "_read_prepared_asset_bytes",
        replace_after_validation,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_has_component_text_overlap",
        lambda masks, text_mask: False,
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(
        image_to_ppt,
        "_finalize_slide_quality",
        lambda slide_data, lang, **kwargs: slide_data,
    )

    with pytest.raises(ValueError, match="sha256|changed while being read"):
        image_to_ppt.finalize_component_layers(
            prepared,
            _accepted_component_layers(prepared),
            lang="en",
        )

    assert replaced is True
    assert not list(work_dir.glob("quality-components-*"))


def test_agent_managed_finalize_rejects_tampered_prepared_dict(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    accepted = _accepted_component_layers(prepared)
    prepared["img_width"] += 1

    with pytest.raises(ValueError, match="fresh prepared state"):
        image_to_ppt.finalize_component_layers(prepared, accepted, lang="en")


def test_agent_managed_finalize_cleans_staging_when_result_becomes_empty(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    monkeypatch.setattr(
        image_to_ppt,
        "_has_component_text_overlap",
        lambda masks, text_mask: False,
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)

    def fake_finalize(slide_data, lang, **kwargs):
        slide_data["components"] = []
        return slide_data

    monkeypatch.setattr(image_to_ppt, "_finalize_slide_quality", fake_finalize)

    with pytest.raises(
        image_to_ppt.VisualSegmentationError,
        match="became empty",
    ):
        image_to_ppt.finalize_component_layers(
            prepared,
            _accepted_component_layers(prepared),
            lang="en",
        )

    assert not list(Path(prepared["_work_dir"]).glob("quality-components-*"))


def test_agent_managed_finalize_cleans_staging_when_component_snapshot_fails(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, _ = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    monkeypatch.setattr(
        image_to_ppt,
        "_has_component_text_overlap",
        lambda masks, text_mask: False,
    )
    original_snapshot = image_to_ppt._snapshot_prepared_asset
    staged_snapshot_count = 0

    def fail_second_component_snapshot(
        work_dir,
        record,
        label,
        staged_dir,
        name,
    ):
        nonlocal staged_snapshot_count
        if name.startswith("component_"):
            staged_snapshot_count += 1
            if staged_snapshot_count == 2:
                raise OSError("staged snapshot failed")
        return original_snapshot(work_dir, record, label, staged_dir, name)

    monkeypatch.setattr(
        image_to_ppt,
        "_snapshot_prepared_asset",
        fail_second_component_snapshot,
    )

    with pytest.raises(OSError, match="staged snapshot failed"):
        image_to_ppt.finalize_component_layers(
            prepared,
            _accepted_component_layers(prepared),
            lang="en",
        )

    assert staged_snapshot_count == 2
    assert not list(Path(prepared["_work_dir"]).glob("quality-components-*"))


def test_agent_managed_finalize_does_not_overwrite_hardlink_occupying_snapshot(
    tmp_path: Path,
    monkeypatch,
) -> None:
    prepared, work_dir = _prepare_component_layers_fixture(tmp_path, monkeypatch)
    external_path = tmp_path / "external-snapshot.bin"
    external_content = b"external content must remain unchanged"
    external_path.write_bytes(external_content)
    real_read_asset = image_to_ppt._read_prepared_asset_bytes
    occupied = False

    def occupy_source_snapshot(asset_work_dir, record, label):
        nonlocal occupied
        owned, content = real_read_asset(asset_work_dir, record, label)
        if label == "source image" and not occupied:
            staged_dirs = list(work_dir.glob("quality-components-*"))
            assert len(staged_dirs) == 1
            os.link(
                external_path,
                staged_dirs[0] / f"source-image{owned.suffix}",
            )
            occupied = True
        return owned, content

    monkeypatch.setattr(
        image_to_ppt,
        "_read_prepared_asset_bytes",
        occupy_source_snapshot,
    )
    monkeypatch.setattr(
        image_to_ppt,
        "_has_component_text_overlap",
        lambda masks, text_mask: False,
    )
    monkeypatch.setattr(image_to_ppt, "close_ocr_engines", lambda: None)
    monkeypatch.setattr(
        image_to_ppt,
        "_finalize_slide_quality",
        lambda slide_data, lang, **kwargs: slide_data,
    )
    exclusive_error = None

    try:
        image_to_ppt.finalize_component_layers(
            prepared,
            _accepted_component_layers(prepared),
            lang="en",
        )
    except FileExistsError as error:
        exclusive_error = error

    assert occupied is True
    assert external_path.read_bytes() == external_content
    assert exclusive_error is not None
    assert not list(work_dir.glob("quality-components-*"))


@pytest.mark.parametrize(
    "content_kind",
    ["card", "person_outline", "dense_lines", "gradient_object"],
)
def test_component_tree_preserves_intact_parent_and_exports_only_active_child(
    tmp_path: Path,
    content_kind: str,
) -> None:
    import numpy as np

    parent = np.zeros((24, 32), dtype=bool)
    if content_kind == "card":
        parent[3:21, 4:28] = True
    elif content_kind == "person_outline":
        parent[3:9, 13:19] = True
        parent[8:20, 10:22] = True
    elif content_kind == "dense_lines":
        parent[4:20:3, 3:29] = True
        parent[3:21, 5:29:4] = True
    else:
        yy, xx = np.ogrid[:24, :32]
        parent = ((yy - 12) ** 2) / 64 + ((xx - 16) ** 2) / 144 <= 1
    child = parent.copy()
    first_y, first_x = np.argwhere(parent)[0]
    child[first_y, first_x] = False
    element = visual_segment.VisualElement(
        mask=child,
        semantic_mask=parent,
        z_index=3,
        score=0.9,
        source="synthetic",
    )

    layers = visual_segment.build_component_mask_layers([element])
    image = np.full((24, 32, 3), 180, dtype=np.uint8)
    exported = fg_extract.export_component_tree(
        image,
        layers,
        tmp_path / content_kind,
        text_mask=np.zeros((24, 32), dtype=np.uint8),
    )

    from image2editable.component_contracts import validate_component_graph

    assert validate_component_graph(exported["graph"]) is exported["graph"]
    parent_node, child_node = exported["graph"]["nodes"]
    assert parent_node["kind"] == "parent"
    assert parent_node["state"] == "inactive"
    assert child_node["kind"] == "child"
    assert child_node["state"] == "pending"
    assert child_node["parent_id"] == parent_node["id"]
    assert len(exported["components"]) == 1
    assert Path(exported["components"][0]["path"]).is_file()
    assert (tmp_path / content_kind / "component-graph.json").is_file()
    saved_parent = np.asarray(
        Image.open(tmp_path / content_kind / parent_node["mask"])
    ) > 0
    saved_child = np.asarray(
        Image.open(tmp_path / content_kind / child_node["mask"])
    ) > 0
    assert np.array_equal(saved_parent, parent)
    assert np.array_equal(saved_child, child)
    assert np.all(~saved_child | saved_parent)
    assert np.count_nonzero(saved_parent) > np.count_nonzero(saved_child)


def test_component_mask_layers_reject_overlapping_active_children() -> None:
    import numpy as np

    first = np.zeros((10, 10), dtype=bool)
    first[1:6, 1:6] = True
    second = np.zeros((10, 10), dtype=bool)
    second[5:9, 5:9] = True
    elements = [
        visual_segment.VisualElement(first, 0, 0.9, "synthetic"),
        visual_segment.VisualElement(second, 1, 0.9, "synthetic"),
    ]

    with pytest.raises(
        visual_segment.VisualSegmentationError,
        match="overlapping visual ownership",
    ):
        visual_segment.build_component_mask_layers(elements)


def test_visual_mask_validation_streams_without_numpy_sum(monkeypatch) -> None:
    import numpy as np

    first = np.zeros((20, 20), dtype=bool)
    first[1:5, 1:5] = True
    second = np.zeros((20, 20), dtype=bool)
    second[10:15, 10:15] = True
    monkeypatch.setattr(
        visual_segment.np,
        "sum",
        lambda *args, **kwargs: pytest.fail("mask stack sum is not allowed"),
    )

    visual_segment.validate_visual_masks([first, second])


def test_component_mask_layers_share_valid_bool_inputs_without_mutation() -> None:
    import numpy as np

    child = np.zeros((12, 12), dtype=bool)
    child[3:9, 3:9] = True
    parent = child.copy()
    parent[2, 5] = True
    originals = child.copy(), parent.copy()
    element = visual_segment.VisualElement(
        child,
        0,
        0.9,
        "synthetic",
        semantic_mask=parent,
    )

    layer = visual_segment.build_component_mask_layers([element])[0]

    assert np.shares_memory(layer["child_mask"], child)
    assert np.shares_memory(layer["parent_mask"], parent)
    assert np.array_equal(child, originals[0])
    assert np.array_equal(parent, originals[1])


@pytest.mark.parametrize(
    "invalid",
    [
        [["mask"]],
        [[float("nan")]],
        [[-1]],
    ],
)
def test_visual_mask_validation_rejects_invalid_values(invalid) -> None:
    with pytest.raises(ValueError, match="mask"):
        visual_segment.validate_visual_masks([invalid])


def _persist_component_graph_masks(
    graph_root: Path,
    graph: dict,
    masks: dict[str, object],
) -> None:
    import numpy as np

    for node in graph["nodes"]:
        mask = np.asarray(masks[node["id"]], dtype=bool)
        mask_path = graph_root / node["mask"]
        mask_path.parent.mkdir(parents=True, exist_ok=True)
        Image.fromarray(mask.astype(np.uint8) * 255, mode="L").save(mask_path)
        node["mask_sha256"] = hashlib.sha256(mask_path.read_bytes()).hexdigest()
        ys, xs = np.nonzero(mask)
        node["bbox"] = [
            int(xs.min()),
            int(ys.min()),
            int(xs.max()) + 1,
            int(ys.max()) + 1,
        ]


def test_component_graph_export_includes_frozen_but_not_failed_nodes(
    tmp_path: Path,
) -> None:
    import numpy as np

    from image2editable.component_contracts import validate_component_graph

    frozen_mask = np.zeros((12, 12), dtype=bool)
    frozen_mask[1:5, 1:5] = True
    failed_mask = np.zeros((12, 12), dtype=bool)
    failed_mask[7:10, 7:10] = True
    graph = {
        "nodes": [
            {
                "id": "component_frozen",
                "kind": "parent",
                "parent_id": None,
                "state": "frozen",
                "mask": "masks/component_frozen.png",
                "mask_sha256": "a" * 64,
                "bbox": [1, 1, 5, 5],
                "z_index": 4,
                "text_ids": [],
            },
            {
                "id": "component_failed",
                "kind": "parent",
                "parent_id": None,
                "state": "failed",
                "mask": "masks/component_failed.png",
                "mask_sha256": "b" * 64,
                "bbox": [7, 7, 10, 10],
                "z_index": 5,
                "text_ids": [],
            },
        ]
    }
    graph_root = tmp_path / "graph"
    _persist_component_graph_masks(
        graph_root,
        graph,
        {
            "component_frozen": frozen_mask,
            "component_failed": failed_mask,
        },
    )
    validate_component_graph(graph)

    components = fg_extract.export_component_graph(
        np.full((12, 12, 3), 180, dtype=np.uint8),
        graph,
        graph_root,
        tmp_path / "components",
        text_mask=np.zeros((12, 12), dtype=np.uint8),
        foreground_mask=frozen_mask,
    )

    assert len(components) == 1
    assert components[0]["component_id"] == "component_frozen"
    assert components[0]["z_index"] == 4


def test_component_graph_export_rejects_duplicate_active_ownership(
    tmp_path: Path,
) -> None:
    import numpy as np

    mask_a = np.zeros((10, 10), dtype=bool)
    mask_a[1:6, 1:6] = True
    mask_b = np.zeros((10, 10), dtype=bool)
    mask_b[5:9, 5:9] = True

    def node(component_id: str, z_index: int) -> dict:
        return {
            "id": component_id,
            "kind": "parent",
            "parent_id": None,
            "state": "frozen",
            "mask": f"masks/{component_id}.png",
            "mask_sha256": str(z_index + 1) * 64,
            "bbox": [1, 1, 9, 9],
            "z_index": z_index,
            "text_ids": [],
        }

    graph = {"nodes": [node("component_a", 0), node("component_b", 1)]}
    graph_root = tmp_path / "graph"
    _persist_component_graph_masks(
        graph_root,
        graph,
        {"component_a": mask_a, "component_b": mask_b},
    )

    with pytest.raises(
        fg_extract.ComponentExtractionError,
        match="ownership",
    ):
        fg_extract.export_component_graph(
            np.full((10, 10, 3), 180, dtype=np.uint8),
            graph,
            graph_root,
            tmp_path / "components",
            text_mask=np.zeros((10, 10), dtype=np.uint8),
            foreground_mask=mask_a | mask_b,
        )


@pytest.mark.parametrize("corruption", ["hash", "bbox", "extra"])
def test_component_graph_export_rejects_unbound_mask_assets(
    tmp_path: Path,
    corruption: str,
) -> None:
    import numpy as np

    mask = np.zeros((10, 10), dtype=bool)
    mask[2:7, 3:8] = True
    graph = {
        "nodes": [{
            "id": "component_0001",
            "kind": "parent",
            "parent_id": None,
            "state": "frozen",
            "mask": "masks/component_0001.png",
            "mask_sha256": "0" * 64,
            "bbox": [3, 2, 8, 7],
            "z_index": 0,
            "text_ids": [],
        }]
    }
    graph_root = tmp_path / "graph"
    _persist_component_graph_masks(
        graph_root,
        graph,
        {"component_0001": mask},
    )
    if corruption == "hash":
        graph["nodes"][0]["mask_sha256"] = "f" * 64
    elif corruption == "bbox":
        graph["nodes"][0]["bbox"] = [2, 2, 8, 7]
    else:
        Image.fromarray(mask.astype(np.uint8) * 255, mode="L").save(
            graph_root / "masks" / "undeclared.png"
        )

    with pytest.raises(ValueError, match="mask|bbox|undeclared"):
        fg_extract.export_component_graph(
            np.full((10, 10, 3), 180, dtype=np.uint8),
            graph,
            graph_root,
            tmp_path / "components",
            text_mask=np.zeros((10, 10), dtype=np.uint8),
            foreground_mask=mask,
        )
    assert not (tmp_path / "components").exists()


def test_component_graph_export_sorts_by_z_index_and_cleans_embedded_text(
    tmp_path: Path,
) -> None:
    import numpy as np

    low = np.zeros((16, 20), dtype=bool)
    low[2:7, 2:8] = True
    high = np.zeros_like(low)
    high[9:14, 11:18] = True
    text_mask = np.zeros_like(low, dtype=np.uint8)
    text_mask[3:5, 4:6] = 255
    image = np.full((16, 20, 3), 240, dtype=np.uint8)
    image[4, 2:8] = 0
    image[3:5, 4:6] = [220, 20, 20]
    text_clean = image.copy()
    text_clean[3:5, 4:6] = 240
    text_clean[4, 2:8] = 0
    text_items = [{"box": [4, 3, 2, 2], "color": "#DC1414"}]

    def node(component_id: str, z_index: int) -> dict:
        return {
            "id": component_id,
            "kind": "parent",
            "parent_id": None,
            "state": "frozen",
            "mask": f"masks/{component_id}.png",
            "mask_sha256": "0" * 64,
            "bbox": [1, 1, 2, 2],
            "z_index": z_index,
            "text_ids": [],
        }

    graph = {"nodes": [node("high", 9), node("low", 1)]}
    graph_root = tmp_path / "graph"
    _persist_component_graph_masks(
        graph_root,
        graph,
        {"high": high, "low": low},
    )
    components = fg_extract.export_component_graph(
        image,
        graph,
        graph_root,
        tmp_path / "components",
        text_mask=text_mask,
        foreground_mask=low | high,
        text_items=text_items,
        text_clean_image=text_clean,
    )

    assert [item["component_id"] for item in components] == ["low", "high"]
    assert [item["z_index"] for item in components] == [1, 9]
    low_component = components[0]
    with Image.open(low_component["path"]) as exported:
        rgba = np.asarray(exported.convert("RGBA"))
    local_y = 4 - low_component["y"]
    local_x = 3 - low_component["x"]
    assert np.all(rgba[local_y, local_x, :3] < 20)
    assert rgba[local_y, local_x, 3] > 0


def test_component_graph_export_requires_reliable_text_clean_assets(
    tmp_path: Path,
) -> None:
    import numpy as np

    mask = np.zeros((8, 8), dtype=bool)
    mask[1:7, 1:7] = True
    text_mask = np.zeros((8, 8), dtype=np.uint8)
    text_mask[3:5, 3:5] = 255
    graph = {"nodes": [{
        "id": "component_0001",
        "kind": "parent",
        "parent_id": None,
        "state": "frozen",
        "mask": "masks/component_0001.png",
        "mask_sha256": "0" * 64,
        "bbox": [1, 1, 7, 7],
        "z_index": 0,
        "text_ids": [],
    }]}
    graph_root = tmp_path / "graph"
    _persist_component_graph_masks(
        graph_root,
        graph,
        {"component_0001": mask},
    )

    with pytest.raises(
        fg_extract.ComponentExtractionError,
        match="text_items.*text_clean_image",
    ):
        fg_extract.export_component_graph(
            np.full((8, 8, 3), 180, dtype=np.uint8),
            graph,
            graph_root,
            tmp_path / "components",
            text_mask=text_mask,
            foreground_mask=mask,
        )
    assert not (tmp_path / "components").exists()


def test_component_graph_export_rejects_missing_foreground(tmp_path: Path) -> None:
    import numpy as np

    mask = np.zeros((10, 10), dtype=bool)
    mask[2:5, 2:5] = True
    foreground = mask.copy()
    foreground[7, 7] = True
    graph = {"nodes": [{
        "id": "component_0001",
        "kind": "parent",
        "parent_id": None,
        "state": "frozen",
        "mask": "masks/component_0001.png",
        "mask_sha256": "0" * 64,
        "bbox": [2, 2, 5, 5],
        "z_index": 0,
        "text_ids": [],
    }]}
    graph_root = tmp_path / "graph"
    _persist_component_graph_masks(
        graph_root,
        graph,
        {"component_0001": mask},
    )

    with pytest.raises(fg_extract.ComponentExtractionError, match="missing"):
        fg_extract.export_component_graph(
            np.full((10, 10, 3), 180, dtype=np.uint8),
            graph,
            graph_root,
            tmp_path / "components",
            text_mask=np.zeros((10, 10), dtype=np.uint8),
            foreground_mask=foreground,
        )
    assert not (tmp_path / "components").exists()


def test_component_tree_failure_leaves_no_partial_or_old_assets(
    tmp_path: Path,
) -> None:
    import numpy as np

    output_dir = tmp_path / "tree"
    output_dir.mkdir()
    (output_dir / "old.png").write_bytes(b"old")
    mask = np.zeros((8, 8), dtype=bool)
    mask[2:6, 2:6] = True
    layers = [{"parent_mask": mask, "child_mask": mask, "z_index": 0}]

    with pytest.raises(FileExistsError):
        fg_extract.export_component_tree(
            np.full((8, 8, 3), 180, dtype=np.uint8),
            layers,
            output_dir,
            text_mask=np.zeros((8, 8), dtype=np.uint8),
        )
    assert (output_dir / "old.png").read_bytes() == b"old"
    assert not list(tmp_path.glob(".tree-staging-*"))


def test_component_tree_cleans_staging_when_component_export_fails(
    tmp_path: Path,
    monkeypatch,
) -> None:
    import numpy as np

    mask = np.zeros((8, 8), dtype=bool)
    mask[2:6, 2:6] = True
    layers = [{"parent_mask": mask, "child_mask": mask, "z_index": 0}]

    def fail_after_partial_write(*args, **kwargs):
        output_dir = Path(args[2])
        output_dir.mkdir(parents=True, exist_ok=True)
        (output_dir / "partial.png").write_bytes(b"partial")
        raise RuntimeError("forced component failure")

    monkeypatch.setattr(
        fg_extract,
        "export_visual_components",
        fail_after_partial_write,
    )
    with pytest.raises(RuntimeError, match="forced component failure"):
        fg_extract.export_component_tree(
            np.full((8, 8, 3), 180, dtype=np.uint8),
            layers,
            tmp_path / "tree",
            text_mask=np.zeros((8, 8), dtype=np.uint8),
        )
    assert not (tmp_path / "tree").exists()
    assert not list(tmp_path.glob(".tree-staging-*"))


@pytest.mark.parametrize("invalid", [[[-1]], [[float("nan")]], [["mask"]]])
def test_component_tree_rejects_invalid_mask_values(
    tmp_path: Path,
    invalid,
) -> None:
    import numpy as np

    with pytest.raises(ValueError, match="mask"):
        fg_extract.export_component_tree(
            np.full((1, 1, 3), 180, dtype=np.uint8),
            [{"parent_mask": invalid, "child_mask": [[1]], "z_index": 0}],
            tmp_path / "tree",
            text_mask=np.zeros((1, 1), dtype=np.uint8),
        )
    assert not (tmp_path / "tree").exists()
    assert not list(tmp_path.glob(".tree-staging-*"))


def test_text_hole_repairs_use_one_owner_map_and_skip_empty_text() -> None:
    import numpy as np

    first = np.zeros((12, 12), dtype=bool)
    first[1:5, 1:5] = True
    second = np.zeros((12, 12), dtype=bool)
    second[7:11, 7:11] = True
    empty = np.zeros((12, 12), dtype=np.uint8)

    assert fg_extract._assign_text_hole_repairs(
        [first, second],
        empty,
        empty,
    ) is None

    text = empty.copy()
    text[4:8, 4:8] = 255
    owners = fg_extract._assign_text_hole_repairs(
        [first, second],
        text,
        text,
    )
    assert isinstance(owners, np.ndarray)
    assert owners.shape == text.shape
    assert owners.dtype == np.uint32


def test_component_graph_export_rejects_hardlinked_mask(tmp_path: Path) -> None:
    import numpy as np

    mask = np.zeros((8, 8), dtype=bool)
    mask[2:6, 2:6] = True
    graph = {"nodes": [{
        "id": "component_0001",
        "kind": "parent",
        "parent_id": None,
        "state": "frozen",
        "mask": "masks/component_0001.png",
        "mask_sha256": "0" * 64,
        "bbox": [2, 2, 6, 6],
        "z_index": 0,
        "text_ids": [],
    }]}
    graph_root = tmp_path / "graph"
    _persist_component_graph_masks(
        graph_root,
        graph,
        {"component_0001": mask},
    )
    mask_path = graph_root / graph["nodes"][0]["mask"]
    external = tmp_path / "external.png"
    mask_path.replace(external)
    try:
        os.link(external, mask_path)
    except OSError as error:
        pytest.skip(f"hardlink is unavailable: {error}")

    with pytest.raises(ValueError, match="single-link"):
        fg_extract.export_component_graph(
            np.full((8, 8, 3), 180, dtype=np.uint8),
            graph,
            graph_root,
            tmp_path / "components",
            text_mask=np.zeros((8, 8), dtype=np.uint8),
            foreground_mask=mask,
        )
    assert external.is_file()
    assert not (tmp_path / "components").exists()


def test_component_graph_export_rejects_output_inside_mask_evidence(
    tmp_path: Path,
) -> None:
    import numpy as np

    mask = np.zeros((8, 8), dtype=bool)
    mask[2:6, 2:6] = True
    graph = {"nodes": [{
        "id": "component_0001",
        "kind": "parent",
        "parent_id": None,
        "state": "frozen",
        "mask": "masks/component_0001.png",
        "mask_sha256": "0" * 64,
        "bbox": [2, 2, 6, 6],
        "z_index": 0,
        "text_ids": [],
    }]}
    graph_root = tmp_path / "graph"
    _persist_component_graph_masks(
        graph_root,
        graph,
        {"component_0001": mask},
    )

    with pytest.raises(ValueError, match="graph masks"):
        fg_extract.export_component_graph(
            np.full((8, 8, 3), 180, dtype=np.uint8),
            graph,
            graph_root,
            graph_root / "masks" / "generated",
            text_mask=np.zeros((8, 8), dtype=np.uint8),
            foreground_mask=mask,
        )
    assert not (graph_root / "masks" / "generated").exists()


def test_skill_component_graph_modules_run_outside_repository(tmp_path: Path) -> None:
    skill_root = Path(__file__).parents[1] / "skills" / "image-to-ppt"
    command = (
        "from scripts import fg_extract, component_contracts, component_quality; "
        "api = fg_extract._component_graph_api(); "
        "assert component_contracts.COMPONENT_KINDS; "
        "assert api[1] is component_contracts.validate_component_graph; "
        "assert api[2] is component_quality.validate_pixel_ownership"
    )
    environment = os.environ.copy()
    environment["PYTHONPATH"] = str(skill_root)
    completed = subprocess.run(
        [sys.executable, "-c", command],
        cwd=tmp_path,
        env=environment,
        text=True,
        capture_output=True,
    )
    assert completed.returncode == 0, completed.stderr
