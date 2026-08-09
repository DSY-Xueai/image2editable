from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from scripts.ppt_assemble import assemble_pptx, assemble_pptx_multi


def _background(tmp_path: Path) -> Path:
    path = tmp_path / "background.png"
    Image.new("RGB", (400, 300), "white").save(path)
    return path


def _native_rectangle(z_index: int = 1) -> dict:
    return {
        "object_id": "shape_1",
        "route": "native_shape",
        "z_index": z_index,
        "bbox": [40, 60, 160, 140],
        "shape": {
            "shape_type": "rectangle",
            "fill_rgb": [30, 90, 180],
        },
    }


def test_native_shape_is_not_written_as_picture(tmp_path: Path) -> None:
    output = tmp_path / "native.pptx"

    assemble_pptx(
        _background(tmp_path),
        [],
        [],
        400,
        300,
        output,
        slide_size="original",
        visual_elements=[_native_rectangle()],
    )

    presentation = Presentation(output)
    content = list(presentation.slides[0].shapes)[1:]
    assert len(content) == 1
    assert content[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert content[0].name == "image2editable:shape_1"


def test_visual_elements_follow_z_order_and_text_stays_on_top(tmp_path: Path) -> None:
    component_path = tmp_path / "component.png"
    Image.new("RGBA", (80, 60), (200, 30, 20, 255)).save(component_path)
    output = tmp_path / "layers.pptx"
    raster = {
        "object_id": "raster_1",
        "route": "raster_component",
        "z_index": 1,
        "component": {
            "path": str(component_path),
            "x": 10,
            "y": 10,
            "w": 80,
            "h": 60,
        },
    }

    assemble_pptx(
        _background(tmp_path),
        [],
        [{"box": [20, 20, 100, 30], "text": "top"}],
        400,
        300,
        output,
        slide_size="original",
        visual_elements=[_native_rectangle(z_index=2), raster],
    )

    presentation = Presentation(output)
    content = list(presentation.slides[0].shapes)[1:]
    assert [shape.shape_type for shape in content] == [
        MSO_SHAPE_TYPE.PICTURE,
        MSO_SHAPE_TYPE.AUTO_SHAPE,
        MSO_SHAPE_TYPE.TEXT_BOX,
    ]


def test_native_line_uses_measured_endpoints(tmp_path: Path) -> None:
    output = tmp_path / "line.pptx"
    element = {
        "object_id": "line_1",
        "route": "native_shape",
        "z_index": 1,
        "bbox": [20, 20, 180, 100],
        "shape": {
            "shape_type": "line",
            "fill_rgb": [10, 20, 30],
            "line_start": [20, 100],
            "line_end": [180, 20],
            "line_width": 5,
        },
    }

    assemble_pptx(
        _background(tmp_path),
        [],
        [],
        400,
        300,
        output,
        slide_size="original",
        visual_elements=[element],
    )

    presentation = Presentation(output)
    line = list(presentation.slides[0].shapes)[1]
    assert line.shape_type == MSO_SHAPE_TYPE.LINE
    assert line.name == "image2editable:line_1"


def test_multi_slide_uses_visual_elements_when_present(tmp_path: Path) -> None:
    output = tmp_path / "multi.pptx"
    background = _background(tmp_path)

    assemble_pptx_multi(
        [
            {
                "background_path": str(background),
                "background_original_path": str(background),
                "components": [],
                "visual_elements": [_native_rectangle()],
                "text_items": [],
                "img_width": 400,
                "img_height": 300,
            }
        ],
        output,
        slide_size="original",
    )

    presentation = Presentation(output)
    content = list(presentation.slides[0].shapes)[1:]
    assert [shape.shape_type for shape in content] == [MSO_SHAPE_TYPE.AUTO_SHAPE]


def test_unsupported_visual_route_fails(tmp_path: Path) -> None:
    import pytest

    with pytest.raises(ValueError, match="Unsupported visual route"):
        assemble_pptx(
            _background(tmp_path),
            [],
            [],
            400,
            300,
            tmp_path / "bad.pptx",
            slide_size="original",
            visual_elements=[{"route": "svg", "z_index": 1}],
        )
