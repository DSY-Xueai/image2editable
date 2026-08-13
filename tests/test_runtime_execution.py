from __future__ import annotations

import copy
import hashlib
import json
from datetime import datetime
import multiprocessing
import os
import re
import shutil
import subprocess
import types
from pathlib import Path, PurePosixPath
from typing import Any

import numpy as np
import pytest
from PIL import Image, ImageDraw
from pypdf import PdfWriter
from pptx import Presentation

from image2editable import component_repair, host_agent, legacy, runtime
from image2editable.component_contracts import MAX_REPAIR_ROUNDS
from image2editable.contracts import PageStatus, RunStatus, SCHEMA_VERSION
from image2editable.execution import ExecutionLease
from image2editable.pptx_input import prepare_pptx_job
from image2editable.resources import safe_default_policy
from image2editable.store import RunStore


def _component_source(path: Path) -> None:
    image = Image.new("RGB", (32, 32), "black")
    for y in range(8, 24):
        for x in range(8, 24):
            image.putpixel((x, y), (255, 255, 255))
    image.save(path)


def _install_component_e2e_boundaries(
    monkeypatch: pytest.MonkeyPatch,
    *,
    baked_background_pages: set[str] | None = None,
    component_count: int = 1,
) -> tuple[dict[str, Any], list[str], list[str]]:
    prepared_pages: dict[str, dict[str, Any]] = {}
    initial_calls: list[str] = []
    assembly_calls: list[str] = []
    baked_background_pages = baked_background_pages or set()

    class BoundaryImageModule:
        @staticmethod
        def prepare_component_layers(source, work_dir, **kwargs):
            work = Path(work_dir)
            page_id = work.parent.parent.name
            initial_calls.append(page_id)
            work.mkdir(parents=True)
            with Image.open(source) as opened:
                source_rgb = opened.convert("RGB")
                width, height = source_rgb.size
            background = work / "background.png"
            difference = work / "difference.png"
            text_mask = work / "text-mask.png"
            if page_id in baked_background_pages:
                source_rgb.save(background)
            else:
                Image.new("RGB", (width, height), "black").save(background)
            Image.new("RGB", (width, height), "black").save(difference)
            Image.new("L", (width, height), 0).save(text_mask)
            component_masks = []
            components = []
            foreground_evidence = np.zeros((height, width), dtype=np.uint8)
            for index in range(component_count):
                if component_count == 1:
                    left, top = width // 4, height // 4
                    right, bottom = width * 3 // 4, height * 3 // 4
                else:
                    left, top = width * index // component_count, 0
                    right, bottom = width * (index + 1) // component_count, height
                component_mask = work / f"component-mask-{index:04d}.png"
                mask = Image.new("L", (width, height), 0)
                ImageDraw.Draw(mask).rectangle(
                    (left, top, right - 1, bottom - 1), fill=255
                )
                mask.save(component_mask)
                foreground_evidence[top:bottom, left:right] = 255
                component_masks.append(str(component_mask))
                component = work / f"component-{index:04d}.png"
                source_rgb.crop((left, top, right, bottom)).convert("RGBA").save(component)
                components.append({
                    "path": str(component), "x": left, "y": top,
                    "w": right - left, "h": bottom - top, "z_index": index,
                })
            state_path = work / "prepared-page.json"
            state_path.write_text("{}", encoding="utf-8")
            foreground_evidence_path = work / "foreground-evidence-mask.png"
            Image.fromarray(foreground_evidence, mode="L").save(
                foreground_evidence_path
            )
            prepared = {
                "_prepared_schema_version": 5,
                "state_path": str(state_path),
                "initial_component_count": component_count,
                "original_image_path": str(Path(source).resolve()),
                "background_original_path": str(background),
                "background_difference_path": str(difference),
                "_text_mask_path": str(text_mask),
                "_element_mask_paths": component_masks,
                "_foreground_evidence_mask_path": str(foreground_evidence_path),
                "components": components,
                "img_width": width, "img_height": height,
                "canvas_width": width, "canvas_height": height,
                "content_offset_x": 0, "content_offset_y": 0,
                "text_items": [],
            }
            prepared_pages[page_id] = prepared
            return prepared

        @staticmethod
        def load_component_layers(state_path):
            return prepared_pages[Path(state_path).parents[1].parent.name]

        @staticmethod
        def _assemble_prepared_slide(slide, output, *args):
            assembly_calls.append("single")
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(output)
            return str(output)

        @staticmethod
        def assemble_pptx_multi(slides, output, **kwargs):
            assembly_calls.append("multi")
            presentation = Presentation()
            for _ in slides:
                presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(output)
            return str(output)

    real_import_module = legacy.importlib.import_module
    monkeypatch.setattr(
        legacy.importlib,
        "import_module",
        lambda name: (
            BoundaryImageModule
            if name == "image_to_ppt"
            else real_import_module(name)
        ),
    )
    monkeypatch.setattr(
        "scripts.sam_worker.run_component_prompt_worker",
        lambda *args, **kwargs: (_ for _ in ()).throw(
            AssertionError("SAM must not run for accept/fallback plans")
        ),
    )
    return prepared_pages, initial_calls, assembly_calls


def _record_current_component_plan(
    run_dir: Path, plan_path: Path, actions: list[dict[str, Any]]
) -> None:
    store = RunStore.open(run_dir)
    summary = store.read_json("run_summary.json")
    page_id = summary["current_page"]
    state = store.read_json(
        f"pages/{page_id}/reconstruction/component_state.json"
    )
    request_ref = state["current_round"]["request_ref"]
    plan = {
        "schema_version": 1,
        "kind": "component_plan",
        "page_id": page_id,
        "provider": "host",
        "repair_round": state["repair_round"],
        "request_sha256": request_ref["sha256"],
        "actions": actions,
    }
    plan_path.write_text(json.dumps(plan), encoding="utf-8")
    runtime.record_host_agent_plan(run_dir, plan_path)


def _accept_action(component_id: str = "component_0001") -> dict[str, Any]:
    return {
        "action": "accept",
        "object_ids": [component_id],
        "parameters": {},
        "confidence": 0.95,
        "evidence": ["component boundary is complete"],
    }


def _assert_round_presentation_manifest(
    run_dir: Path, page_id: str, repair_round: int
) -> None:
    round_dir = (
        run_dir / "pages" / page_id / "reconstruction" / "agent"
        / f"round-{repair_round:02d}"
    )
    request = json.loads(
        (round_dir / "component_agent_request.json").read_text(encoding="utf-8")
    )
    assert "presentation-manifest.json" in request["evidence"]
    manifest = json.loads(
        (round_dir / "presentation-manifest.json").read_text(encoding="utf-8")
    )
    assert manifest["source_sha256"] == request["source_sha256"]
    assert manifest["graph_sha256"] == request["graph_sha256"]
    assert set(manifest["components"][0]) == {
        "component_id", "rgba", "ownership_mask",
        "presentation_alpha_mask", "generated_underlay_mask", "metrics",
    }
    for component in manifest["components"]:
        for name in (
            "rgba", "ownership_mask", "presentation_alpha_mask",
            "generated_underlay_mask",
        ):
            reference = component[name]
            assert set(reference) == {"path", "sha256"}
            pure = PurePosixPath(reference["path"])
            assert not pure.is_absolute()
            assert ".." not in pure.parts
            asset = run_dir / Path(*pure.parts)
            assert asset.is_file()
            assert asset.resolve().is_relative_to(run_dir.resolve())
            assert hashlib.sha256(asset.read_bytes()).hexdigest() == reference["sha256"]


def _build_test_presentation_manifest(
    run_root: Path,
    *,
    source_path: Path,
    text_clean_path: Path,
    graph: dict,
    graph_dir: Path,
    output_dir: Path,
    text_mask_path: Path | None = None,
) -> tuple[Path, str]:
    graph_path = graph_dir / "component-graph.json"
    graph_path.write_text(json.dumps(graph), encoding="utf-8")
    manifest_path = legacy._build_presentation_assets(
        types.SimpleNamespace(root=run_root),
        source_path=source_path,
        text_clean_path=text_clean_path,
        text_mask_path=text_mask_path,
        graph_path=graph_path,
        output_dir=output_dir,
    )
    return manifest_path, hashlib.sha256(graph_path.read_bytes()).hexdigest()


def _manifest_asset_path(
    run_root: Path, manifest: dict, component_index: int, field: str
) -> Path:
    reference = manifest["components"][component_index][field]
    return run_root / Path(*PurePosixPath(reference["path"]).parts)


def test_presentation_rgba_zeroes_only_transparent_rgb(tmp_path: Path) -> None:
    from scripts.component_underlay import build_presentation_layer

    source_path = tmp_path / "source.png"
    mask_path = tmp_path / "mask.png"
    source = np.arange(6 * 4 * 3, dtype=np.uint8).reshape((4, 6, 3))
    ownership = np.zeros((4, 6), dtype=bool)
    ownership[:, :3] = True
    Image.fromarray(source, mode="RGB").save(source_path)
    Image.fromarray(ownership.astype(np.uint8) * 255, mode="L").save(mask_path)
    graph = {"nodes": [{
        "id": "component", "kind": "parent", "parent_id": None,
        "state": "pending", "mask": mask_path.name,
        "mask_sha256": hashlib.sha256(mask_path.read_bytes()).hexdigest(),
        "bbox": [0, 0, 3, 4], "z_index": 0, "text_ids": [],
    }]}
    manifest_path, _ = _build_test_presentation_manifest(
        tmp_path,
        source_path=source_path,
        text_clean_path=source_path,
        graph=graph,
        graph_dir=tmp_path,
        output_dir=tmp_path,
    )
    expected = build_presentation_layer(
        source_rgb=source,
        text_clean_rgb=source,
        ownership_mask=ownership,
        semantic_mask=ownership,
        higher_layer_mask=np.zeros_like(ownership),
        text_mask=np.zeros_like(ownership),
    )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    with Image.open(_manifest_asset_path(tmp_path, manifest, 0, "rgba")) as image:
        rgba = np.asarray(image.convert("RGBA")).copy()

    alpha = expected["presentation_alpha_mask"]
    assert np.array_equal(rgba[:, :, 3] == 255, alpha)
    assert np.all(rgba[~alpha, :3] == 0)
    assert np.array_equal(rgba[alpha, :3], expected["rgb"][alpha])


def test_presentation_assets_assign_text_hole_to_colored_shape(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "source.png"
    clean_path = tmp_path / "text-clean.png"
    visual_path = tmp_path / "visual.png"
    text_path = tmp_path / "text.png"
    source = np.full((20, 40, 3), 255, dtype=np.uint8)
    source[4:16, 4:36] = (20, 160, 60)
    source[8:12, 16:24] = 255
    clean = source.copy()
    clean[8:12, 16:24] = (20, 160, 60)
    visual = np.zeros((20, 40), dtype=np.uint8)
    visual[4:16, 4:36] = 255
    visual[8:12, 16:24] = 0
    text = np.zeros((20, 40), dtype=np.uint8)
    text[8:12, 16:24] = 255
    Image.fromarray(source).save(source_path)
    Image.fromarray(clean).save(clean_path)
    Image.fromarray(visual).save(visual_path)
    Image.fromarray(text).save(text_path)
    graph = {"nodes": [{
        "id": "component", "kind": "parent", "parent_id": None,
        "state": "pending", "mask": visual_path.name,
        "mask_sha256": hashlib.sha256(visual_path.read_bytes()).hexdigest(),
        "bbox": [4, 4, 36, 16], "z_index": 0, "text_ids": [],
    }, {
        "id": "text", "kind": "text", "parent_id": None,
        "state": "frozen", "mask": text_path.name,
        "mask_sha256": hashlib.sha256(text_path.read_bytes()).hexdigest(),
        "bbox": [16, 8, 24, 12], "z_index": 1, "text_ids": [],
    }]}

    manifest_path, _ = _build_test_presentation_manifest(
        tmp_path, source_path=source_path, text_clean_path=clean_path,
        graph=graph, graph_dir=tmp_path, output_dir=tmp_path,
    )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    with Image.open(_manifest_asset_path(
        tmp_path, manifest, 0, "ownership_mask"
    )) as image:
        ownership = np.asarray(image.convert("L")) > 0
    with Image.open(_manifest_asset_path(
        tmp_path, manifest, 0, "rgba"
    )) as image:
        rgba = np.asarray(image.convert("RGBA"))

    assert not np.any(ownership[8:12, 16:24])
    assert np.all(rgba[8:12, 16:24, :3] == (20, 160, 60))


def test_presentation_assets_publish_atomically_and_retry_after_save_failure(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    source_path = tmp_path / "source.png"
    mask_path = tmp_path / "mask.png"
    Image.new("RGB", (4, 4), "red").save(source_path)
    Image.new("L", (4, 4), 255).save(mask_path)
    graph = {"nodes": [{
        "id": "component", "kind": "parent", "parent_id": None,
        "state": "pending", "mask": mask_path.name,
        "mask_sha256": hashlib.sha256(mask_path.read_bytes()).hexdigest(),
        "bbox": [0, 0, 4, 4], "z_index": 0, "text_ids": [],
    }]}
    graph_path = tmp_path / "component-graph.json"
    graph_path.write_text(json.dumps(graph), encoding="utf-8")
    real_save = Image.Image.save
    save_count = 0

    def fail_second_save(image, *args, **kwargs):
        nonlocal save_count
        save_count += 1
        if save_count == 2:
            raise RuntimeError("controlled presentation save failure")
        return real_save(image, *args, **kwargs)

    monkeypatch.setattr(Image.Image, "save", fail_second_save)
    with pytest.raises(RuntimeError, match="controlled presentation save failure"):
        legacy._build_presentation_assets(
            types.SimpleNamespace(root=tmp_path),
            source_path=source_path,
            text_clean_path=source_path,
            graph_path=graph_path,
            output_dir=tmp_path,
        )

    final_dir = tmp_path / "presentation-assets"
    assert not final_dir.exists()
    assert not list(tmp_path.glob(".presentation-assets.tmp-*"))

    monkeypatch.setattr(Image.Image, "save", real_save)
    real_sha256_file = legacy.sha256_file

    def reject_staging_hash(path):
        assert not Path(path).parent.name.startswith(".presentation-assets.tmp-")
        return real_sha256_file(path)

    monkeypatch.setattr(legacy, "sha256_file", reject_staging_hash)
    manifest_path = legacy._build_presentation_assets(
        types.SimpleNamespace(root=tmp_path),
        source_path=source_path,
        text_clean_path=source_path,
        graph_path=graph_path,
        output_dir=tmp_path,
    )
    assert manifest_path == final_dir / "presentation-manifest.json"
    assert manifest_path.is_file()

    conflict_output = tmp_path / "conflict"
    conflict_output.mkdir()
    conflict_final = conflict_output / "presentation-assets"
    conflict_final.mkdir()
    with pytest.raises(RuntimeError, match="presentation assets already published"):
        legacy._build_presentation_assets(
            types.SimpleNamespace(root=tmp_path),
            source_path=source_path,
            text_clean_path=source_path,
            graph_path=graph_path,
            output_dir=conflict_output,
        )
    assert list(conflict_final.iterdir()) == []
    assert not list(conflict_output.glob(".presentation-assets.tmp-*"))


def test_presentation_assets_keep_fully_occluded_component_transparent(
    tmp_path: Path,
) -> None:
    source_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_mask_path = tmp_path / "text-mask.png"
    mask_path = tmp_path / "mask.png"
    Image.new("RGB", (4, 4), "red").save(source_path)
    Image.new("RGB", (4, 4), "black").save(background_path)
    Image.new("L", (4, 4), 0).save(text_mask_path)
    Image.new("L", (4, 4), 255).save(mask_path)
    mask_sha256 = hashlib.sha256(mask_path.read_bytes()).hexdigest()
    graph = {"nodes": [
        {
            "id": "lower", "kind": "parent", "parent_id": None,
            "state": "pending", "mask": mask_path.name,
            "mask_sha256": mask_sha256, "bbox": [0, 0, 4, 4],
            "z_index": 0, "text_ids": [],
        },
        {
            "id": "upper", "kind": "parent", "parent_id": None,
            "state": "pending", "mask": mask_path.name,
            "mask_sha256": mask_sha256, "bbox": [0, 0, 4, 4],
            "z_index": 1, "text_ids": [],
        },
    ]}
    manifest_path, _ = _build_test_presentation_manifest(
        tmp_path,
        source_path=source_path,
        text_clean_path=source_path,
        graph=graph,
        graph_dir=tmp_path,
        output_dir=tmp_path,
    )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

    assert [item["component_id"] for item in manifest["components"]] == [
        "lower", "upper",
    ]
    for field in (
        "rgba", "ownership_mask", "presentation_alpha_mask",
        "generated_underlay_mask",
    ):
        with Image.open(_manifest_asset_path(tmp_path, manifest, 0, field)) as image:
            assert not np.any(np.asarray(image))
    assert all(
        value == 0.0 for value in manifest["components"][0]["metrics"].values()
    )
    evidence_dir = tmp_path / "evidence"
    evidence_dir.mkdir()
    evidence = legacy._render_component_evidence(
        source_path=source_path,
        graph=graph,
        text_mask_path=text_mask_path,
        background_path=background_path,
        presentation_manifest_path=manifest_path,
        run_root=tmp_path,
        reconstruction=tmp_path,
        graph_sha256=hashlib.sha256(
            (tmp_path / "component-graph.json").read_bytes()
        ).hexdigest(),
        output_dir=evidence_dir,
        text_items=[],
    )
    with Image.open(evidence["component-isolation.png"]) as isolation:
        assert np.any(np.asarray(isolation))


def test_presentation_assets_allow_empty_active_graph(tmp_path: Path) -> None:
    source_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_mask_path = tmp_path / "text-mask.png"
    mask_path = tmp_path / "mask.png"
    Image.new("RGB", (4, 4), "red").save(source_path)
    Image.new("RGB", (4, 4), "white").save(background_path)
    Image.new("L", (4, 4), 0).save(text_mask_path)
    Image.new("L", (4, 4), 255).save(mask_path)
    graph = {"nodes": [{
        "id": "inactive", "kind": "parent", "parent_id": None,
        "state": "inactive", "mask": mask_path.name,
        "mask_sha256": hashlib.sha256(mask_path.read_bytes()).hexdigest(),
        "bbox": [0, 0, 4, 4], "z_index": 0, "text_ids": [],
    }]}
    manifest_path, graph_sha256 = _build_test_presentation_manifest(
        tmp_path,
        source_path=source_path,
        text_clean_path=source_path,
        graph=graph,
        graph_dir=tmp_path,
        output_dir=tmp_path,
    )
    assert json.loads(manifest_path.read_text(encoding="utf-8"))["components"] == []
    evidence = legacy._render_component_evidence(
        source_path=source_path,
        graph=graph,
        text_mask_path=text_mask_path,
        background_path=background_path,
        presentation_manifest_path=manifest_path,
        run_root=tmp_path,
        reconstruction=tmp_path,
        graph_sha256=graph_sha256,
        output_dir=tmp_path,
        text_items=[],
    )
    with Image.open(evidence["reconstructed.png"]) as reconstructed:
        assert reconstructed.getpixel((0, 0)) == (255, 255, 255)


def test_presentation_asset_decode_rejects_post_validation_replacement(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    source, graph, manifest_path, graph_sha256 = (
        _presentation_asset_validation_case(tmp_path)
    )
    real_validate = legacy._validate_presentation_manifest

    def replace_after_validation(*args, **kwargs):
        manifest = real_validate(*args, **kwargs)
        rgba_path = _manifest_asset_path(tmp_path, manifest, 0, "rgba")
        Image.new("RGBA", (4, 4), (0, 255, 0, 255)).save(rgba_path)
        return manifest

    monkeypatch.setattr(
        legacy, "_validate_presentation_manifest", replace_after_validation
    )
    with pytest.raises(RuntimeError, match="presentation asset hash mismatch"):
        for _ in legacy._load_presentation_assets(
            run_root=tmp_path,
            reconstruction=tmp_path,
            manifest_path=manifest_path,
            source_sha256=hashlib.sha256(source.read_bytes()).hexdigest(),
            graph_sha256=graph_sha256,
            graph=graph,
            page_size=(4, 4),
        ):
            pass


def test_presentation_higher_masks_use_one_reverse_z_accumulator(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    import scripts.component_underlay as component_underlay

    source_path = tmp_path / "source.png"
    Image.new("RGB", (4, 1), "red").save(source_path)
    specs = (("mid_a", 1, 0), ("low", 0, 1), ("top", 2, 2), ("mid_b", 1, 3))
    nodes = []
    for component_id, z_index, x in specs:
        mask_path = tmp_path / f"{component_id}.png"
        mask = np.zeros((1, 4), dtype=np.uint8)
        mask[0, x] = 255
        Image.fromarray(mask, mode="L").save(mask_path)
        nodes.append({
            "id": component_id, "kind": "parent", "parent_id": None,
            "state": "pending", "mask": mask_path.name,
            "mask_sha256": hashlib.sha256(mask_path.read_bytes()).hexdigest(),
            "bbox": [x, 0, x + 1, 1], "z_index": z_index, "text_ids": [],
        })
    captured = {}

    def capture_layer(**kwargs):
        ownership = kwargs["ownership_mask"]
        x = int(np.flatnonzero(ownership)[0])
        captured[x] = kwargs["higher_layer_mask"].copy()
        return {
            "rgb": kwargs["source_rgb"].copy(),
            "ownership_mask": ownership.copy(),
            "presentation_alpha_mask": ownership.copy(),
            "generated_underlay_mask": np.zeros_like(ownership),
            "metrics": {
                "boundary_color_mae": 0.0,
                "gradient_jump_p95": 0.0,
                "added_high_frequency_pixels": 0.0,
            },
        }

    real_zeros = np.zeros
    zero_calls = []

    def counted_zeros(*args, **kwargs):
        zero_calls.append((args, kwargs))
        return real_zeros(*args, **kwargs)

    monkeypatch.setattr(component_underlay, "build_presentation_layer", capture_layer)
    monkeypatch.setattr(np, "zeros", counted_zeros)
    manifest_path, _ = _build_test_presentation_manifest(
        tmp_path,
        source_path=source_path,
        text_clean_path=source_path,
        graph={"nodes": nodes},
        graph_dir=tmp_path,
        output_dir=tmp_path,
    )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))

    assert [item["component_id"] for item in manifest["components"]] == [
        item[0] for item in specs
    ]
    assert len(zero_calls) == 3
    assert not np.any(captured[2])
    assert np.array_equal(captured[0], np.array([[False, False, True, False]]))
    assert np.array_equal(captured[3], np.array([[False, False, True, False]]))
    assert np.array_equal(captured[1], np.array([[True, False, True, True]]))


def test_presentation_loader_and_compositor_consume_layers_incrementally(
    tmp_path: Path,
) -> None:
    import gc
    import weakref

    source, graph, manifest_path, graph_sha256 = (
        _presentation_asset_validation_case(tmp_path)
    )
    loaded = legacy._load_presentation_assets(
        run_root=tmp_path,
        reconstruction=tmp_path,
        manifest_path=manifest_path,
        source_sha256=hashlib.sha256(source.read_bytes()).hexdigest(),
        graph_sha256=graph_sha256,
        graph=graph,
        page_size=(4, 4),
    )
    assert iter(loaded) is loaded
    assert next(loaded)["component_id"] == "component"
    with pytest.raises(StopIteration):
        next(loaded)

    class Layer(dict):
        pass

    first_reference = None

    def layers():
        nonlocal first_reference
        first = Layer(component_id="first", rgba=np.zeros((1, 1, 4), dtype=np.uint8))
        first_reference = weakref.ref(first)
        yield first
        del first
        gc.collect()
        assert first_reference() is None
        yield Layer(component_id="second", rgba=np.zeros((1, 1, 4), dtype=np.uint8))
        yield Layer(component_id="third", rgba=np.zeros((1, 1, 4), dtype=np.uint8))

    stream_graph = {"nodes": [
        {"id": component_id, "kind": "parent", "state": "pending", "z_index": index}
        for index, component_id in enumerate(("first", "second", "third"))
    ]}
    background = Image.new("RGB", (1, 1), "black")
    try:
        composited = legacy._composite_presentation_layers(
            background, stream_graph, layers()
        )
        composited.close()
    finally:
        background.close()


def _local_receipt(tmp_path: Path) -> dict[str, object]:
    return {
        "schema_version": 1,
        "model_id": "test/model",
        "requested_revision": "test",
        "resolved_revision": "a" * 40,
        "stability": "experimental",
        "snapshot_path": str((tmp_path / "model-snapshot").resolve()),
        "files": [
            {
                "path": "config.json",
                "size": 2,
                "sha256": "b" * 64,
            }
        ],
    }


def _write_mock_component_state(store: RunStore, page_id: str) -> None:
    reference = {"path": "mock-artifact.json", "sha256": "0" * 64}
    store.write_json(
        f"pages/{page_id}/reconstruction/component_state.json",
        {
            "schema_version": 1, "page_id": page_id, "provider": "host",
            "source_sha256": "1" * 64, "initial_component_count": 0,
            "quality_gate_version": 1, "revision": 1,
            "phase": "awaiting_plan", "status": "active",
            "repair_round": 1, "plan_count": 0, "stop_reason": None,
            "graph_ref": reference,
            "current_round": {
                "round": 1, "request_ref": reference,
                "plan_ref": None, "execution_ref": None, "quality_ref": None,
            },
            "frozen": {}, "candidate_ids": [], "failed_ids": [],
            "fallback": {"status": "none", "parent_ids": []},
            "last_normalized_plan_sha256": None, "result_ref": None,
            "delivery_checks": {"pptx_reopen": "unknown"},
            "updated_at": "now", "round_history": [], "parent_assets": {},
            "fallback_graph_ref": None, "fallback_quality_ref": None,
            "fallback_input_refs": None,
        },
    )


def test_image_component_plan_e2e_pauses_then_assembles_once(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _component_source(source)
    _, initial_calls, assembly_calls = _install_component_e2e_boundaries(monkeypatch)
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", slide_size="16:9",
        agent_provider="host",
    )

    waiting = runtime.run_job(run_dir)

    assert waiting["status"] == "awaiting_agent"
    assert waiting["pending_components"] == 1
    assert waiting["frozen_components"] == 0
    assert initial_calls == ["page_001"]
    assert assembly_calls == []
    status = runtime.get_status(run_dir)
    assert status["current_page"] == "page_001"
    assert status["repair_round"] == 1
    assert status["pending_components"] == 1
    assert Path(status["diagnostics"]).is_absolute()

    _record_current_component_plan(
        run_dir, tmp_path / "image-plan.json", [_accept_action()]
    )
    completed = runtime.run_job(run_dir)

    assert completed["status"] == "completed"
    assert completed["quality_gate_version"] == runtime.COMPONENT_QUALITY_GATE_VERSION
    assert initial_calls == ["page_001"]
    assert assembly_calls == ["single"]
    result = RunStore.open(run_dir).read_json(
        "pages/page_001/reconstruction/component_result.json"
    )
    assert result["status"] == "ready_for_assembly"
    assert result["final_component_ids"] == ["component_0001"]


def test_image_component_retry_passes_source_image_once_to_sam_worker(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _component_source(source)
    _install_component_e2e_boundaries(monkeypatch)
    calls: list[dict[str, Any]] = []

    def fake_sam_worker(image, *, box, positive, negative, work_dir):
        calls.append({
            "shape": image.shape,
            "box": box,
            "positive": positive,
            "negative": negative,
            "work_dir": Path(work_dir),
        })
        mask = np.zeros(image.shape[:2], dtype=bool)
        mask[8:24, 8:24] = True
        return mask

    monkeypatch.setattr(
        "scripts.sam_worker.run_component_prompt_worker", fake_sam_worker
    )
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", slide_size="16:9",
        agent_provider="host",
    )
    assert runtime.run_job(run_dir)["status"] == "awaiting_agent"
    _record_current_component_plan(
        run_dir,
        tmp_path / "retry-plan.json",
        [{
            "action": "retry_with_box",
            "object_ids": ["component_0001"],
            "parameters": {"box": [0.25, 0.25, 0.75, 0.75]},
            "confidence": 0.95,
            "evidence": ["retry the incomplete component boundary"],
        }],
    )

    runtime.run_job(run_dir)

    assert calls == [{
        "shape": (32, 32, 3),
        "box": [8.0, 8.0, 24.0, 24.0],
        "positive": [],
        "negative": [],
        "work_dir": run_dir / "pages/page_001/reconstruction",
    }]


def test_local_provider_runs_complete_without_host_receipt(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _component_source(source)
    _, initial_calls, assembly_calls = _install_component_e2e_boundaries(monkeypatch)
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        slide_size="16:9",
        agent_provider="local",
    )
    expected_service = object()
    plans = []
    monkeypatch.setattr(runtime, "_local_service_config", lambda: expected_service)

    def fake_local_agent(request_path, *, service_config):
        request_path = Path(request_path)
        request = json.loads(request_path.read_text(encoding="utf-8"))
        plan = {
            "schema_version": 1,
            "kind": "component_plan",
            "page_id": request["page_id"],
            "provider": "local",
            "repair_round": request["repair_round"],
            "request_sha256": hashlib.sha256(request_path.read_bytes()).hexdigest(),
            "actions": [_accept_action()],
        }
        plans.append(plan)
        assert service_config is expected_service
        return plan

    monkeypatch.setattr(runtime, "_run_local_service_agent", fake_local_agent)

    completed = runtime.run_job(run_dir)

    assert completed["status"] == "completed"
    assert initial_calls == ["page_001"]
    assert assembly_calls == ["single"]
    assert [plan["provider"] for plan in plans] == ["local"]
    assert not (run_dir / "host_capabilities.json").exists()
    assert not (run_dir / "host-challenge").exists()
    state = RunStore.open(run_dir).read_json(
        "pages/page_001/reconstruction/component_state.json"
    )
    assert state["provider"] == "local"
    assert state["plan_count"] == 1


def test_local_provider_warning_fails_without_fake_output(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _component_source(source)
    _install_component_e2e_boundaries(
        monkeypatch,
        baked_background_pages={"page_001"},
    )
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        slide_size="16:9",
        agent_provider="local",
    )
    monkeypatch.setattr(runtime, "_local_service_config", lambda: object())
    rounds = []

    def fake_local_agent(request_path, **kwargs):
        request_path = Path(request_path)
        request = json.loads(request_path.read_text(encoding="utf-8"))
        repair_round = request["repair_round"]
        rounds.append(repair_round)
        return {
            "schema_version": 1,
            "kind": "component_plan",
            "page_id": request["page_id"],
            "provider": "local",
            "repair_round": request["repair_round"],
            "request_sha256": hashlib.sha256(request_path.read_bytes()).hexdigest(),
            "actions": [_accept_action()] if request["repair_round"] == 1 else [],
        }

    monkeypatch.setattr(runtime, "_run_local_service_agent", fake_local_agent)

    with pytest.raises(RuntimeError, match="editable reconstruction incomplete"):
        runtime.run_job(run_dir)

    assert rounds == [1, 2]
    assert RunStore.open(run_dir).read_json("run_state.json")["status"] == "failed"
    state = RunStore.open(run_dir).read_json(
        "pages/page_001/reconstruction/component_state.json"
    )
    assert state["plan_count"] == 2
    assert state["status"] == "preserved_with_warning"
    assert not (run_dir / "final/output.pptx").exists()
    assert not (run_dir / "host_capabilities.json").exists()
    round_two = run_dir / "pages/page_001/reconstruction/agent/round-02"
    assert (round_two / "numbered-masks.png").read_bytes() != (
        round_two / "source.png"
    ).read_bytes()
    assert (round_two / "ownership.png").read_bytes() != (
        round_two / "source.png"
    ).read_bytes()
    request = json.loads(
        (round_two / "component_agent_request.json").read_text(encoding="utf-8")
    )
    unexplained = round_two / "unexplained-mask.png"
    assert unexplained.read_bytes() == (
        run_dir / "pages/page_001/reconstruction/execution-01/unexplained-mask.png"
    ).read_bytes()
    assert request["evidence"]["unexplained-mask.png"]["sha256"] == (
        hashlib.sha256(unexplained.read_bytes()).hexdigest()
    )


def test_host_component_rounds_publish_hash_bound_presentation_manifests(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _component_source(source)
    _install_component_e2e_boundaries(
        monkeypatch,
        baked_background_pages={"page_001"},
    )
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        slide_size="16:9",
        agent_provider="host",
    )

    first_wait = runtime.run_job(run_dir)
    assert first_wait["status"] == "awaiting_agent"
    _assert_round_presentation_manifest(run_dir, "page_001", 1)

    _record_current_component_plan(
        run_dir,
        tmp_path / "round-1-plan.json",
        [_accept_action()],
    )
    second_wait = runtime.run_job(run_dir)

    assert second_wait["status"] == "awaiting_agent"
    assert second_wait["repair_round"] == 2
    _assert_round_presentation_manifest(run_dir, "page_001", 1)
    _assert_round_presentation_manifest(run_dir, "page_001", 2)


def test_next_round_disk_reserve_fails_before_evidence_publication(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _component_source(source)
    _install_component_e2e_boundaries(
        monkeypatch,
        baked_background_pages={"page_001"},
    )
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        slide_size="16:9",
        agent_provider="local",
    )
    monkeypatch.setattr(runtime, "_local_service_config", lambda: object())

    def fake_local_agent(request_path, **kwargs):
        request_path = Path(request_path)
        request = json.loads(request_path.read_text(encoding="utf-8"))
        return {
            "schema_version": 1,
            "kind": "component_plan",
            "page_id": request["page_id"],
            "provider": "local",
            "repair_round": request["repair_round"],
            "request_sha256": hashlib.sha256(request_path.read_bytes()).hexdigest(),
            "actions": [_accept_action()],
        }

    monkeypatch.setattr(runtime, "_run_local_service_agent", fake_local_agent)
    real_reserve = legacy._ensure_component_disk_reserve
    reserve_calls = 0

    def fail_next_publication(*args, **kwargs):
        nonlocal reserve_calls
        reserve_calls += 1
        if reserve_calls == 3:
            raise RuntimeError("component page disk reserve is insufficient")
        return real_reserve(*args, **kwargs)

    monkeypatch.setattr(
        legacy,
        "_ensure_component_disk_reserve",
        fail_next_publication,
    )

    with pytest.raises(RuntimeError, match="disk reserve"):
        runtime.run_job(run_dir)

    reconstruction = run_dir / "pages/page_001/reconstruction"
    assert reserve_calls == 3
    assert not (reconstruction / "evidence-round-02").exists()


def test_local_provider_stops_when_page_quality_does_not_improve(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _component_source(source)
    _install_component_e2e_boundaries(
        monkeypatch,
        baked_background_pages={"page_001"},
    )
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        slide_size="16:9",
        agent_provider="local",
    )
    monkeypatch.setattr(runtime, "_local_service_config", lambda: object())
    rounds = []

    def fake_local_agent(request_path, **kwargs):
        request_path = Path(request_path)
        request = json.loads(request_path.read_text(encoding="utf-8"))
        repair_round = request["repair_round"]
        rounds.append(repair_round)
        return {
            "schema_version": 1,
            "kind": "component_plan",
            "page_id": request["page_id"],
            "provider": "local",
            "repair_round": request["repair_round"],
            "request_sha256": hashlib.sha256(request_path.read_bytes()).hexdigest(),
            "actions": [
                {
                    "action": "expand",
                    "object_ids": ["component_0001"],
                    "parameters": {"margin_ratio": repair_round / 100},
                    "confidence": 0.95,
                    "evidence": ["bounded distinct repair attempt"],
                }
            ],
        }

    monkeypatch.setattr(runtime, "_run_local_service_agent", fake_local_agent)

    with pytest.raises(RuntimeError, match="editable reconstruction incomplete"):
        runtime.run_job(run_dir)

    assert rounds == [1, 2]
    state = RunStore.open(run_dir).read_json(
        "pages/page_001/reconstruction/component_state.json"
    )
    assert state["plan_count"] == 2
    assert state["stop_reason"] == "no_quality_improvement"
    assert state["status"] == "preserved_with_warning"
    assert not (run_dir / "pages/page_001/reconstruction/agent/round-03").exists()
    assert not (run_dir / "final/output.pptx").exists()


def test_local_missing_service_configuration_stops_before_heavy_page_initialization(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _component_source(source)
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        agent_provider="local",
    )
    initialized = False

    def missing_service():
        raise RuntimeError("Local model service is not configured")

    def unexpected_initialize(*args, **kwargs):
        nonlocal initialized
        initialized = True
        raise AssertionError("heavy page initialization must not start")

    monkeypatch.setattr(runtime, "_local_service_config", missing_service)
    monkeypatch.setattr(runtime, "initialize_legacy_page", unexpected_initialize)

    with pytest.raises(RuntimeError, match="not configured"):
        runtime.run_job(run_dir)

    assert initialized is False
    assert not (run_dir / "host_capabilities.json").exists()


def test_pdf_component_plan_e2e_is_serial_and_falls_back_before_one_assembly(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    import image2editable.pdf_input as pdf_input

    source = tmp_path / "source.pdf"
    writer = PdfWriter()
    writer.add_blank_page(width=320, height=180)
    writer.add_blank_page(width=320, height=180)
    with source.open("wb") as stream:
        writer.write(stream)

    def render_pdf(source_path, outputs, *, profile):
        records = []
        for index, output in enumerate(outputs):
            target = Path(output)
            target.parent.mkdir(parents=True, exist_ok=True)
            _component_source(target)
            records.append({
                "page_index": index, "page_number": index + 1,
                "width_pt": 320.0, "height_pt": 180.0, "rotation": 0,
                "media_box": [0.0, 0.0, 320.0, 180.0],
                "crop_box": [0.0, 0.0, 320.0, 180.0], "profile": profile,
                "target_dpi": 144, "effective_dpi": 144.0, "scale": 2.0,
                "reasons": [], "pixel_width": 32, "pixel_height": 32,
                "sha256": hashlib.sha256(target.read_bytes()).hexdigest(),
                "renderer": "test-boundary", "renderer_version": "1",
            })
        return records

    monkeypatch.setattr(pdf_input, "render_pdf_document", render_pdf)
    _, initial_calls, assembly_calls = _install_component_e2e_boundaries(
        monkeypatch, baked_background_pages={"page_002"}
    )
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", slide_size="16:9",
        agent_provider="host",
    )

    first_wait = runtime.run_job(run_dir)
    assert first_wait["current_page"] == "page_001"
    assert initial_calls == ["page_001"]
    assert assembly_calls == []

    _record_current_component_plan(
        run_dir, tmp_path / "page-1-plan.json", [_accept_action()]
    )
    second_wait = runtime.run_job(run_dir)
    assert second_wait["current_page"] == "page_002"
    assert second_wait["repair_round"] == 1
    assert initial_calls == ["page_001", "page_002"]
    assert assembly_calls == []

    _record_current_component_plan(
        run_dir, tmp_path / "page-2-round-1.json", [_accept_action()]
    )
    third_wait = runtime.run_job(run_dir)
    assert third_wait["current_page"] == "page_002"
    assert third_wait["repair_round"] == 2
    assert initial_calls == ["page_001", "page_002"]
    assert assembly_calls == []

    _record_current_component_plan(
        run_dir, tmp_path / "page-2-round-2.json", []
    )
    with pytest.raises(RuntimeError, match="editable reconstruction incomplete"):
        runtime.run_job(run_dir)

    assert initial_calls == ["page_001", "page_002"]
    assert assembly_calls == []
    store = RunStore.open(run_dir)
    page_1_state = store.read_json(
        "pages/page_001/reconstruction/component_state.json"
    )
    page_2_state = store.read_json("pages/page_002/reconstruction/component_state.json")
    assert (
        page_1_state["quality_gate_version"]
        == page_2_state["quality_gate_version"]
    )
    assert page_2_state["status"] == "preserved_with_warning"
    assert page_2_state["fallback"]["status"] == "warning"
    assert store.read_json("run_state.json")["status"] == "failed"
    assert not (run_dir / "final/output.pptx").exists()


def test_image_host_pauses_without_assembly_or_completed_summary(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", agent_provider="host"
    )
    calls = {"initialize": 0, "assemble": 0}

    def initialize(store: RunStore, page_id: str, *, _lease=None) -> dict:
        calls["initialize"] += 1
        reconstruction = store.root / "pages" / page_id / "reconstruction"
        reconstruction.mkdir(parents=True, exist_ok=True)
        _write_mock_component_state(store, page_id)
        return {"status": "request_published", "page_id": page_id}

    monkeypatch.setattr(runtime, "initialize_legacy_page", initialize, raising=False)
    monkeypatch.setattr(
        runtime, "advance_legacy_page",
        lambda store, page_id, **kwargs: {"status": "awaiting_agent", "page_id": page_id},
        raising=False,
    )

    def assemble(store: RunStore) -> dict:
        calls["assemble"] += 1
        raise AssertionError("assembly must not run while Host is waiting")

    monkeypatch.setattr(runtime, "assemble_legacy_results", assemble, raising=False)
    monkeypatch.setattr(
        runtime, "execute_legacy",
        lambda store: (_ for _ in ()).throw(
            AssertionError("one-shot legacy execution must not run")
        ),
    )

    summary = runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert summary["status"] == "awaiting_agent"
    assert summary["provider"] == "host"
    assert summary["quality_gate_version"] == runtime.COMPONENT_QUALITY_GATE_VERSION
    assert summary["current_page"] == "page_001"
    assert store.read_json("run_state.json")["status"] == "awaiting_agent"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "awaiting_agent"
    assert not (run_dir / "final" / "output.pptx").exists()
    assert calls == {"initialize": 1, "assemble": 0}


def test_image_resume_does_not_repeat_initialized_layers(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", agent_provider="host"
    )
    calls = {"initialize": 0, "advance": 0, "assemble": 0}

    def initialize(store: RunStore, page_id: str, *, _lease=None) -> dict:
        calls["initialize"] += 1
        reconstruction = store.root / "pages" / page_id / "reconstruction"
        reconstruction.mkdir(parents=True, exist_ok=True)
        _write_mock_component_state(store, page_id)
        return {"status": "request_published", "page_id": page_id}

    def advance(store: RunStore, page_id: str, *, _lease=None) -> dict:
        calls["advance"] += 1
        status = "awaiting_agent" if calls["advance"] == 1 else "ready_for_assembly"
        return {"status": status, "page_id": page_id}

    def assemble(store: RunStore) -> dict:
        calls["assemble"] += 1
        output = store.root / "final" / "output.pptx"
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_bytes(b"pptx")
        return {"pptx": str(output)}

    monkeypatch.setattr(runtime, "initialize_legacy_page", initialize, raising=False)
    monkeypatch.setattr(runtime, "advance_legacy_page", advance, raising=False)
    monkeypatch.setattr(runtime, "assemble_legacy_results", assemble, raising=False)
    monkeypatch.setattr(
        runtime, "execute_legacy",
        lambda store: (_ for _ in ()).throw(
            AssertionError("one-shot legacy execution must not run")
        ),
    )
    monkeypatch.setattr(runtime, "_validate_completed_pptx_output", lambda *args: None)

    assert runtime.run_job(run_dir)["status"] == "awaiting_agent"
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.PREPARED)
    store.transition_page("page_001", PageStatus.PROCESSING)
    summary = runtime.run_job(run_dir)

    assert summary["status"] == "completed"
    assert summary["quality_gate_version"] == runtime.COMPONENT_QUALITY_GATE_VERSION
    assert calls == {"initialize": 1, "advance": 2, "assemble": 1}


def test_legacy_page_initialization_is_idempotent_without_rerunning_models(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    calls = {"prepare": 0, "request": 0, "state": 0}
    reconstruction = run_dir / "pages" / "page_001" / "reconstruction"

    class FakeImageModule:
        @staticmethod
        def prepare_component_layers(*args, **kwargs):
            calls["prepare"] += 1
            return {"state_path": str(reconstruction / "initial/prepared-page.json"),
                    "initial_component_count": 2}

    monkeypatch.setattr(
        legacy.importlib, "import_module", lambda name: FakeImageModule
    )
    monkeypatch.setattr(
        legacy, "_build_initial_page_session",
        lambda *args: {"page_id": "page_001"}, raising=False,
    )

    def request(session: dict, *, repair_round: int) -> Path:
        calls["request"] += 1
        path = reconstruction / "agent/round-01/component_agent_request.json"
        path.parent.mkdir(parents=True, exist_ok=True)
        path.write_text("{}", encoding="utf-8")
        return path

    def initialize(store: RunStore, page_id: str, **kwargs) -> dict:
        calls["state"] += 1
        store.write_json(
            f"pages/{page_id}/reconstruction/component_state.json", {"ready": True}
        )
        return {"phase": "request_published"}

    monkeypatch.setattr(legacy, "build_component_agent_request", request, raising=False)
    monkeypatch.setattr(
        legacy, "initialize_component_repair_state", initialize, raising=False
    )

    with ExecutionLease(store.root / "execution.lock", run_root=store.root) as lease:
        first = legacy.initialize_legacy_page(store, "page_001", _lease=lease)
        second = legacy.initialize_legacy_page(store, "page_001", _lease=lease)

    assert first["status"] == "initialized"
    assert second["status"] == "already_initialized"
    assert calls == {"prepare": 1, "request": 1, "state": 1}


@pytest.mark.parametrize("prepared_schema_version", [1, 5])
def test_initial_page_session_uses_versioned_foreground_evidence(
    tmp_path: Path,
    prepared_schema_version: int,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    prepared_root = reconstruction / "initial"
    prepared_root.mkdir(parents=True)
    background = prepared_root / "background.png"
    difference = prepared_root / "difference.png"
    text_mask = prepared_root / "text-mask.png"
    component = prepared_root / "component.png"
    component_mask = prepared_root / "component-mask.png"
    Image.new("RGB", (20, 10), (1, 2, 3)).save(source)
    Image.new("RGB", (20, 10), "black").save(background)
    Image.new("RGB", (20, 10), "black").save(difference)
    Image.new("L", (20, 10), 0).save(text_mask)
    with Image.open(text_mask) as image:
        image.putpixel((10, 5), 255)
        image.save(text_mask)
    Image.new("RGBA", (4, 3), (10, 20, 30, 255)).save(component)
    Image.new("L", (20, 10), 0).save(component_mask)
    mask = Image.open(component_mask)
    mask.putpixel((2, 1), 255)
    mask.save(component_mask)
    prepared = {
        "_prepared_schema_version": prepared_schema_version,
        "state_path": str(prepared_root / "prepared-page.json"),
        "initial_component_count": 1,
        "original_image_path": str(source),
        "background_original_path": str(background),
        "background_difference_path": str(difference),
        "_text_mask_path": str(text_mask),
        "_element_mask_paths": [str(component_mask)],
        "components": [{
            "path": str(component), "x": 2, "y": 1, "w": 1, "h": 1,
            "z_index": 0,
        }],
        "text_items": [{"text": "T", "box": [10, 4, 3, 2]}],
    }
    if prepared_schema_version == 5:
        foreground = prepared_root / "foreground-evidence-mask.png"
        Image.new("L", (20, 10), 0).save(foreground)
        with Image.open(foreground) as image:
            image.putpixel((2, 1), 255)
            image.save(foreground)
        prepared["_foreground_evidence_mask_path"] = str(foreground)

    session = legacy._build_initial_page_session(
        store, "page_001", prepared, reconstruction
    )

    assert session["provider"] == "host"
    expected_evidence = set(legacy.EVIDENCE_NAMES)
    if prepared_schema_version < 5:
        expected_evidence.remove("unexplained-mask.png")
    assert set(session["evidence"]) == expected_evidence
    graph = json.loads(Path(session["evidence"]["component-graph.json"]).read_text())
    visual, text = graph["nodes"]
    assert visual["kind"] == "parent"
    assert visual["state"] == "pending"
    assert visual["bbox"] == [2, 1, 3, 2]
    assert text["id"] == "text_0001"
    assert text["kind"] == "text"
    assert text["state"] == "frozen"
    assert text["bbox"] == [10, 4, 13, 6]
    graph_root = Path(session["evidence"]["component-graph.json"]).parent
    for node in graph["nodes"]:
        mask_path = graph_root / node["mask"]
        assert hashlib.sha256(mask_path.read_bytes()).hexdigest() == node["mask_sha256"]
    request_path = legacy.build_component_agent_request(session, repair_round=1)
    request = json.loads(request_path.read_text(encoding="utf-8"))
    assert request["candidate_ids"] == ["component_0001"]
    assert request["frozen_ids"] == ["text_0001"]
    if prepared_schema_version == 5:
        record = request["evidence"]["unexplained-mask.png"]
        unexplained = request_path.parent / record["path"]
        assert unexplained.read_bytes() == foreground.read_bytes()
        assert record["sha256"] == hashlib.sha256(unexplained.read_bytes()).hexdigest()
    else:
        assert "unexplained-mask.png" not in request["evidence"]
    evidence = session["evidence"]
    for name in ("numbered-masks.png", "ocr-overlay.png", "ownership.png"):
        assert Path(evidence[name]).read_bytes() != Path(evidence["source.png"]).read_bytes()
    with Image.open(evidence["reconstructed.png"]) as reconstructed:
        assert reconstructed.getpixel((0, 0)) == (0, 0, 0)
        assert reconstructed.getpixel((2, 1)) == (1, 2, 3)


def test_initial_page_session_renders_cjk_ocr_evidence(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    background = tmp_path / "background.png"
    difference = tmp_path / "difference.png"
    text_mask = tmp_path / "text-mask.png"
    for path in (source, background, difference):
        Image.new("RGB", (20, 10), "white").save(path)
    Image.new("L", (20, 10), 0).save(text_mask)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(parents=True)
    labels = []
    original_text = ImageDraw.ImageDraw.text

    def draw_ascii_label(draw, xy, text, *args, **kwargs):
        text.encode("ascii")
        labels.append(text)
        return original_text(draw, xy, text, *args, **kwargs)

    monkeypatch.setattr(ImageDraw.ImageDraw, "text", draw_ascii_label)

    session = legacy._build_initial_page_session(
        RunStore.open(run_dir),
        "page_001",
        {
            "original_image_path": str(source),
            "background_original_path": str(background),
            "background_difference_path": str(difference),
            "_text_mask_path": str(text_mask),
            "_element_mask_paths": [],
            "components": [],
            "text_items": [{"text": "虚拟机", "box": [1, 1, 5, 3]}],
        },
        reconstruction,
    )

    assert Path(session["evidence"]["ocr-overlay.png"]).is_file()
    assert labels[-1] == "text_0001"
    assert all(label.isascii() for label in labels)
    report = json.loads(
        Path(session["evidence"]["quality-report.json"]).read_text(encoding="utf-8")
    )
    assert report["text_items"][0]["text"] == "虚拟机"


def test_initial_page_session_v2_builds_parent_child_graph_and_excludes_text(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (20, 10), (1, 2, 3)).save(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    prepared_root = reconstruction / "initial"
    prepared_root.mkdir(parents=True)
    background = prepared_root / "background.png"
    difference = prepared_root / "difference.png"
    text_mask = prepared_root / "text-mask.png"
    component = prepared_root / "component.png"
    component_mask = prepared_root / "component-mask.png"
    semantic_mask = prepared_root / "semantic-mask.png"
    Image.new("RGB", (20, 10), "black").save(background)
    Image.new("RGB", (20, 10), "black").save(difference)
    mask = Image.new("L", (20, 10), 0)
    mask.putpixel((3, 2), 255)
    mask.save(text_mask)
    mask.close()
    Image.new("RGBA", (4, 3), (10, 20, 30, 255)).save(component)
    mask = Image.new("L", (20, 10), 0)
    mask.putpixel((2, 1), 255)
    mask.putpixel((3, 2), 255)
    mask.save(component_mask)
    mask.close()
    mask = Image.new("L", (20, 10), 0)
    ImageDraw.Draw(mask).rectangle((1, 0, 4, 3), fill=255)
    mask.save(semantic_mask)
    mask.close()
    prepared = {
        "state_path": str(prepared_root / "prepared-page.json"),
        "initial_component_count": 1,
        "original_image_path": str(source),
        "background_original_path": str(background),
        "background_difference_path": str(difference),
        "_text_mask_path": str(text_mask),
        "_element_mask_paths": [str(component_mask)],
        "_semantic_mask_paths": [str(semantic_mask)],
        "components": [{
            "path": str(component), "x": 2, "y": 1, "w": 2, "h": 2,
            "z_index": 0,
        }],
        "text_items": [{"text": "T", "box": [3, 2, 1, 1]}],
    }
    reserve_node_counts = []
    monkeypatch.setattr(
        legacy,
        "_ensure_component_disk_reserve",
        lambda *args, **kwargs: reserve_node_counts.append(kwargs["node_count"]),
    )

    session = legacy._build_initial_page_session(
        store, "page_001", prepared, reconstruction
    )

    assert reserve_node_counts == [3]
    with Image.open(session["evidence"]["reconstructed.png"]) as reconstructed:
        assert reconstructed.getpixel((2, 1)) == (1, 2, 3)
        assert max(
            abs(value - expected)
            for value, expected in zip(
                reconstructed.getpixel((3, 2)), (1, 2, 3)
            )
        ) <= 1
    graph = json.loads(Path(session["evidence"]["component-graph.json"]).read_text())
    parent, child, text = graph["nodes"]
    assert parent["id"] == "parent_0001"
    assert parent["kind"] == "parent"
    assert parent["state"] == "inactive"
    assert parent["parent_id"] is None
    assert parent["bbox"] == [1, 0, 5, 4]
    assert child["id"] == "component_0001"
    assert child["kind"] == "child"
    assert child["state"] == "pending"
    assert child["parent_id"] == "parent_0001"
    assert child["bbox"] == [2, 1, 4, 3]
    assert text["id"] == "text_0001"
    assert text["kind"] == "text"
    assert text["state"] == "frozen"
    graph_root = Path(session["evidence"]["component-graph.json"]).parent
    assert (graph_root / parent["mask"]).read_bytes() == semantic_mask.read_bytes()
    assert (graph_root / child["mask"]).read_bytes() == component_mask.read_bytes()
    for node in graph["nodes"]:
        mask_path = graph_root / node["mask"]
        assert hashlib.sha256(mask_path.read_bytes()).hexdigest() == node["mask_sha256"]
    request_path = legacy.build_component_agent_request(session, repair_round=1)
    request = json.loads(request_path.read_text(encoding="utf-8"))
    assert request["candidate_ids"] == ["component_0001"]
    assert request["frozen_ids"] == ["text_0001"]


def test_initial_page_session_v2_requires_matching_semantic_mask_count(
    tmp_path: Path
) -> None:
    source = tmp_path / "source.png"
    background = tmp_path / "background.png"
    difference = tmp_path / "difference.png"
    text_mask = tmp_path / "text-mask.png"
    for path in (source, background, difference):
        Image.new("RGB", (20, 10), "black").save(path)
    Image.new("L", (20, 10), 0).save(text_mask)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(parents=True)

    with pytest.raises(ValueError, match="semantic mask count"):
        legacy._build_initial_page_session(
            RunStore.open(run_dir),
            "page_001",
            {
                "original_image_path": str(source),
                "background_original_path": str(background),
                "background_difference_path": str(difference),
                "_text_mask_path": str(text_mask),
                "_element_mask_paths": [],
                "_semantic_mask_paths": [str(text_mask)],
                "components": [],
                "text_items": [],
            },
            reconstruction,
        )


def test_initial_page_session_closes_grayscale_when_bbox_raises(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (20, 10), "white").save(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    reconstruction = run_dir / "pages/page_001/reconstruction"
    prepared_root = reconstruction / "initial"
    prepared_root.mkdir(parents=True)
    component_mask = prepared_root / "component-mask.png"
    semantic_mask = prepared_root / "semantic-mask.png"
    Image.new("L", (20, 10), 255).save(component_mask)
    Image.new("L", (20, 10), 255).save(semantic_mask)
    converted_images = []
    closed_image_ids = set()
    real_convert = Image.Image.convert
    real_close = Image.Image.close

    def tracked_convert(image, *args, **kwargs):
        converted = real_convert(image, *args, **kwargs)
        converted_images.append(converted)
        return converted

    def tracked_close(image):
        closed_image_ids.add(id(image))
        real_close(image)

    def failing_getbbox(image):
        raise RuntimeError("controlled bbox failure")

    monkeypatch.setattr(legacy, "_ensure_component_disk_reserve", lambda *a, **k: None)
    monkeypatch.setattr(Image.Image, "convert", tracked_convert)
    monkeypatch.setattr(Image.Image, "close", tracked_close)
    monkeypatch.setattr(Image.Image, "getbbox", failing_getbbox)

    with pytest.raises(RuntimeError, match="controlled bbox failure"):
        legacy._build_initial_page_session(
            RunStore.open(run_dir),
            "page_001",
            {
                "original_image_path": str(source),
                "_element_mask_paths": [str(component_mask)],
                "_semantic_mask_paths": [str(semantic_mask)],
                "components": [{"z_index": 0}],
                "text_items": [],
            },
            reconstruction,
        )

    assert (reconstruction / "evidence-source/masks/parent_0001.png").is_file()
    assert len(converted_images) == 1
    assert id(converted_images[0]) in closed_image_ids


def test_initial_page_session_closes_text_mask_when_save_raises(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    Image.new("RGB", (20, 10), "white").save(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    reconstruction = run_dir / "pages/page_001/reconstruction"
    prepared_root = reconstruction / "initial"
    prepared_root.mkdir(parents=True)
    component_mask = prepared_root / "component-mask.png"
    Image.new("L", (20, 10), 255).save(component_mask)
    created_images = []
    closed_image_ids = set()
    real_new = Image.new
    real_close = Image.Image.close

    def tracked_new(*args, **kwargs):
        image = real_new(*args, **kwargs)
        created_images.append(image)
        return image

    def tracked_close(image):
        closed_image_ids.add(id(image))
        real_close(image)

    def failing_save(image, *args, **kwargs):
        assert image.getpixel((2, 1)) == 255
        raise RuntimeError("controlled text mask save failure")

    monkeypatch.setattr(legacy, "_ensure_component_disk_reserve", lambda *a, **k: None)
    monkeypatch.setattr(Image, "new", tracked_new)
    monkeypatch.setattr(Image.Image, "close", tracked_close)
    monkeypatch.setattr(Image.Image, "save", failing_save)

    with pytest.raises(RuntimeError, match="controlled text mask save failure"):
        legacy._build_initial_page_session(
            RunStore.open(run_dir),
            "page_001",
            {
                "original_image_path": str(source),
                "_element_mask_paths": [str(component_mask)],
                "components": [{"z_index": 0}],
                "text_items": [{"text": "T", "box": [2, 1, 1, 1]}],
            },
            reconstruction,
        )

    assert len(created_images) == 1
    assert id(created_images[0]) in closed_image_ids


def test_component_evidence_closes_loaded_images_when_mask_validation_fails(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    background = tmp_path / "background.png"
    text_mask = tmp_path / "text-mask.png"
    mask_path = tmp_path / "component-mask.png"
    Image.new("RGB", (20, 10), "white").save(source)
    Image.new("RGB", (20, 10), "black").save(background)
    Image.new("L", (20, 10), 0).save(text_mask)
    Image.new("L", (20, 10), 255).save(mask_path)
    graph = {"nodes": [{
        "id": "component_0001", "kind": "parent", "parent_id": None,
        "state": "pending", "mask": mask_path.name,
        "mask_sha256": hashlib.sha256(mask_path.read_bytes()).hexdigest(),
        "bbox": [0, 0, 20, 10], "z_index": 0, "text_ids": [],
    }]}
    manifest_path, graph_sha256 = _build_test_presentation_manifest(
        tmp_path,
        source_path=source,
        text_clean_path=source,
        graph=graph,
        graph_dir=tmp_path,
        output_dir=tmp_path,
    )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    rgba_path = tmp_path / manifest["components"][0]["rgba"]["path"]
    rgba_path.write_bytes(rgba_path.read_bytes() + b"changed")
    copied_images = []
    closed_image_ids = set()
    real_copy = Image.Image.copy
    real_close = Image.Image.close

    def tracked_copy(image):
        copied = real_copy(image)
        copied_images.append(copied)
        return copied

    def tracked_close(image):
        closed_image_ids.add(id(image))
        real_close(image)

    monkeypatch.setattr(Image.Image, "copy", tracked_copy)
    monkeypatch.setattr(Image.Image, "close", tracked_close)

    with pytest.raises(RuntimeError, match="presentation asset hash mismatch"):
        legacy._render_component_evidence(
            source_path=source,
            graph=graph,
            text_mask_path=text_mask,
            background_path=background,
            presentation_manifest_path=manifest_path,
            run_root=tmp_path,
            reconstruction=tmp_path,
            graph_sha256=graph_sha256,
            output_dir=tmp_path,
            text_items=[],
        )

    assert copied_images
    missing = [id(image) for image in copied_images if id(image) not in closed_image_ids]
    assert not missing, (len(copied_images), len(closed_image_ids), missing)


def test_component_evidence_closes_node_images_before_rendering_next_node(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    background = tmp_path / "background.png"
    text_mask = tmp_path / "text-mask.png"
    Image.new("RGB", (20, 10), "white").save(source)
    Image.new("RGB", (20, 10), "black").save(background)
    Image.new("L", (20, 10), 0).save(text_mask)
    nodes = []
    for index in range(2):
        mask_path = tmp_path / f"component-mask-{index}.png"
        mask = Image.new("L", (20, 10), 0)
        ImageDraw.Draw(mask).rectangle(
            (index * 10, 0, index * 10 + 9, 9), fill=255
        )
        mask.save(mask_path)
        mask.close()
        nodes.append({
            "id": f"component_{index + 1:04d}", "kind": "parent",
            "parent_id": None, "state": "pending", "mask": mask_path.name,
            "mask_sha256": hashlib.sha256(mask_path.read_bytes()).hexdigest(),
            "bbox": [0, 0, 20, 10], "z_index": index, "text_ids": [],
        })
    graph = {"nodes": nodes}
    manifest_path, graph_sha256 = _build_test_presentation_manifest(
        tmp_path,
        source_path=source,
        text_clean_path=source,
        graph=graph,
        graph_dir=tmp_path,
        output_dir=tmp_path,
    )
    copied_images = []
    closed_image_ids = set()
    real_convert = Image.Image.convert
    real_close = Image.Image.close

    def tracked_convert(image, *args, **kwargs):
        copied = real_convert(image, *args, **kwargs)
        copied_images.append(copied)
        return copied

    def tracked_close(image):
        closed_image_ids.add(id(image))
        real_close(image)

    monkeypatch.setattr(Image.Image, "convert", tracked_convert)
    monkeypatch.setattr(Image.Image, "close", tracked_close)

    legacy._render_component_evidence(
        source_path=source,
        graph=graph,
        text_mask_path=text_mask,
        background_path=background,
        presentation_manifest_path=manifest_path,
        run_root=tmp_path,
        reconstruction=tmp_path,
        graph_sha256=graph_sha256,
        output_dir=tmp_path,
        text_items=[],
    )

    assert copied_images
    assert all(id(image) in closed_image_ids for image in copied_images)


def test_component_evidence_uses_exact_manifest_rgba_for_every_render(
    tmp_path: Path
) -> None:
    size = (200, 100)
    overlap = (190, 90)
    non_text = (0, 0)
    source_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_mask_path = tmp_path / "text-mask.png"
    parent_mask_path = tmp_path / "parent-mask.png"
    child_mask_path = tmp_path / "child-mask.png"
    Image.new("RGB", size, (1, 2, 3)).save(source_path)
    Image.new("RGB", size, "black").save(background_path)
    text_clean_path = tmp_path / "text-clean.png"
    Image.new("RGB", size, (30, 40, 50)).save(text_clean_path)
    mask = Image.new("L", size, 0)
    mask.putpixel(overlap, 255)
    mask.save(text_mask_path)
    mask.close()
    mask = Image.new("L", size, 0)
    mask.putpixel(non_text, 255)
    mask.putpixel(overlap, 255)
    mask.save(child_mask_path)
    mask.save(parent_mask_path)
    mask.close()
    graph = {"nodes": [
        {
            "id": "parent_0001", "kind": "parent", "parent_id": None,
            "state": "inactive", "mask": parent_mask_path.name,
            "mask_sha256": hashlib.sha256(parent_mask_path.read_bytes()).hexdigest(),
            "bbox": [0, 0, 191, 91], "z_index": 0, "text_ids": [],
        },
        {
            "id": "component_0001", "kind": "child",
            "parent_id": "parent_0001", "state": "pending",
            "mask": child_mask_path.name,
            "mask_sha256": hashlib.sha256(child_mask_path.read_bytes()).hexdigest(),
            "bbox": [0, 0, 191, 91], "z_index": 0, "text_ids": [],
        },
    ]}
    initial_dir = tmp_path / "initial-evidence"
    later_dir = tmp_path / "later-evidence"
    initial_dir.mkdir()
    later_dir.mkdir()
    manifest_path, graph_sha256 = _build_test_presentation_manifest(
        tmp_path,
        source_path=source_path,
        text_clean_path=text_clean_path,
        text_mask_path=text_mask_path,
        graph=graph,
        graph_dir=tmp_path,
        output_dir=initial_dir,
    )

    initial = legacy._render_component_evidence(
        source_path=source_path,
        graph=graph,
        text_mask_path=text_mask_path,
        background_path=background_path,
        presentation_manifest_path=manifest_path,
        run_root=tmp_path,
        reconstruction=tmp_path,
        graph_sha256=graph_sha256,
        output_dir=initial_dir,
        text_items=[],
    )
    later = legacy._render_component_evidence(
        source_path=source_path,
        graph=graph,
        text_mask_path=text_mask_path,
        background_path=background_path,
        presentation_manifest_path=manifest_path,
        run_root=tmp_path,
        reconstruction=tmp_path,
        graph_sha256=graph_sha256,
        output_dir=later_dir,
        text_items=[],
    )

    with Image.open(initial["numbered-masks.png"]) as numbered:
        assert numbered.getpixel(overlap) == (1, 2, 3)
    with Image.open(initial["ownership.png"]) as ownership:
        assert ownership.getpixel(overlap) == (24, 24, 24)
    with Image.open(initial["reconstructed.png"]) as reconstructed:
        assert reconstructed.getpixel(overlap) == (30, 40, 50)
    with Image.open(later["reconstructed.png"]) as reconstructed:
        assert reconstructed.getpixel(non_text) == (1, 2, 3)
        assert reconstructed.getpixel(overlap) == (30, 40, 50)
    with Image.open(initial["component-isolation.png"]) as isolation:
        assert isolation.mode == "RGBA"
        pixels = np.asarray(isolation)
        opaque_rgb = pixels[pixels[:, :, 3] == 255, :3]
        assert np.any(np.all(opaque_rgb == (1, 2, 3), axis=1))
        assert np.any(np.all(opaque_rgb == (30, 40, 50), axis=1))


def test_component_evidence_uses_distinct_z_index_colors_for_four_v2_children(
    tmp_path: Path
) -> None:
    size = (400, 100)
    source_path = tmp_path / "source.png"
    background_path = tmp_path / "background.png"
    text_mask_path = tmp_path / "text-mask.png"
    Image.new("RGB", size, "white").save(source_path)
    Image.new("RGB", size, "black").save(background_path)
    Image.new("L", size, 0).save(text_mask_path)
    nodes = []
    sample_points = []
    for index in range(4):
        parent_id = f"parent_{index + 1:04d}"
        component_id = f"component_{index + 1:04d}"
        parent_mask_path = tmp_path / f"parent-mask-{index}.png"
        child_mask_path = tmp_path / f"child-mask-{index}.png"
        left = index * 100 + 5
        mask = Image.new("L", size, 0)
        ImageDraw.Draw(mask).rectangle((left, 5, left + 80, 85), fill=255)
        mask.save(parent_mask_path)
        mask.save(child_mask_path)
        mask.close()
        nodes.extend((
            {
                "id": parent_id, "kind": "parent", "parent_id": None,
                "state": "inactive", "mask": parent_mask_path.name,
                "mask_sha256": hashlib.sha256(
                    parent_mask_path.read_bytes()
                ).hexdigest(),
                "bbox": [left, 5, left + 81, 86], "z_index": index,
                "text_ids": [],
            },
            {
                "id": component_id, "kind": "child", "parent_id": parent_id,
                "state": "pending", "mask": child_mask_path.name,
                "mask_sha256": hashlib.sha256(
                    child_mask_path.read_bytes()
                ).hexdigest(),
                "bbox": [left, 5, left + 81, 86], "z_index": index,
                "text_ids": [],
            },
        ))
        sample_points.append((left + 5, 75))
    graph = {"nodes": nodes}
    manifest_path, graph_sha256 = _build_test_presentation_manifest(
        tmp_path,
        source_path=source_path,
        text_clean_path=source_path,
        graph=graph,
        graph_dir=tmp_path,
        output_dir=tmp_path,
    )

    evidence = legacy._render_component_evidence(
        source_path=source_path,
        graph=graph,
        text_mask_path=text_mask_path,
        background_path=background_path,
        presentation_manifest_path=manifest_path,
        run_root=tmp_path,
        reconstruction=tmp_path,
        graph_sha256=graph_sha256,
        output_dir=tmp_path,
        text_items=[],
    )

    with Image.open(evidence["ownership.png"]) as ownership:
        colors = [ownership.getpixel(point) for point in sample_points]
    assert colors == [
        (255, 80, 80),
        (70, 180, 255),
        (90, 220, 120),
        (255, 190, 60),
    ]


def _presentation_asset_validation_case(
    tmp_path: Path,
) -> tuple[Path, dict, Path, str]:
    source = tmp_path / "source.png"
    mask = tmp_path / "mask.png"
    Image.new("RGB", (4, 4), "red").save(source)
    Image.new("L", (4, 4), 255).save(mask)
    graph = {"nodes": [{
        "id": "component", "kind": "parent", "parent_id": None,
        "state": "pending", "mask": mask.name,
        "mask_sha256": hashlib.sha256(mask.read_bytes()).hexdigest(),
        "bbox": [0, 0, 4, 4], "z_index": 0, "text_ids": [],
    }]}
    manifest_path, graph_sha256 = _build_test_presentation_manifest(
        tmp_path,
        source_path=source,
        text_clean_path=source,
        graph=graph,
        graph_dir=tmp_path,
        output_dir=tmp_path,
    )
    return source, graph, manifest_path, graph_sha256


def _set_manifest_mask_pixel(
    tmp_path: Path,
    manifest_path: Path,
    field: str,
    value: int,
) -> None:
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    asset_path = tmp_path / manifest["components"][0][field]["path"]
    with Image.open(asset_path) as opened:
        mask = opened.convert("L")
    mask.putpixel((0, 0), value)
    mask.save(asset_path)
    mask.close()
    manifest["components"][0][field]["sha256"] = hashlib.sha256(
        asset_path.read_bytes()
    ).hexdigest()
    manifest_path.write_text(json.dumps(manifest), encoding="utf-8")


def _load_test_presentation_assets(
    tmp_path: Path,
    source: Path,
    graph: dict,
    manifest_path: Path,
    graph_sha256: str,
) -> None:
    for _ in legacy._load_presentation_assets(
        run_root=tmp_path,
        reconstruction=tmp_path,
        manifest_path=manifest_path,
        source_sha256=hashlib.sha256(source.read_bytes()).hexdigest(),
        graph_sha256=graph_sha256,
        graph=graph,
        page_size=(4, 4),
    ):
        pass


def test_presentation_assets_reject_non_binary_rgba_alpha(tmp_path: Path) -> None:
    source, graph, manifest_path, graph_sha256 = (
        _presentation_asset_validation_case(tmp_path)
    )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    rgba_path = tmp_path / manifest["components"][0]["rgba"]["path"]
    with Image.open(rgba_path) as opened:
        rgba = opened.convert("RGBA")
    rgba.putpixel((0, 0), (*rgba.getpixel((0, 0))[:3], 1))
    rgba.save(rgba_path)
    rgba.close()
    manifest["components"][0]["rgba"]["sha256"] = hashlib.sha256(
        rgba_path.read_bytes()
    ).hexdigest()
    manifest_path.write_text(json.dumps(manifest), encoding="utf-8")

    with pytest.raises(
        ValueError,
        match="presentation RGBA alpha does not match alpha mask",
    ):
        _load_test_presentation_assets(
            tmp_path, source, graph, manifest_path, graph_sha256
        )


def test_presentation_assets_reject_non_binary_mask(tmp_path: Path) -> None:
    source, graph, manifest_path, graph_sha256 = (
        _presentation_asset_validation_case(tmp_path)
    )
    _set_manifest_mask_pixel(
        tmp_path, manifest_path, "ownership_mask", 127
    )

    with pytest.raises(
        ValueError,
        match="presentation asset masks must be binary",
    ):
        _load_test_presentation_assets(
            tmp_path, source, graph, manifest_path, graph_sha256
        )


def test_presentation_assets_reject_ownership_generated_overlap(
    tmp_path: Path,
) -> None:
    source, graph, manifest_path, graph_sha256 = (
        _presentation_asset_validation_case(tmp_path)
    )
    _set_manifest_mask_pixel(
        tmp_path, manifest_path, "generated_underlay_mask", 255
    )

    with pytest.raises(
        ValueError,
        match="presentation ownership and generated underlay masks overlap",
    ):
        _load_test_presentation_assets(
            tmp_path, source, graph, manifest_path, graph_sha256
        )


def test_presentation_assets_reject_alpha_mask_union_mismatch(
    tmp_path: Path,
) -> None:
    source, graph, manifest_path, graph_sha256 = (
        _presentation_asset_validation_case(tmp_path)
    )
    _set_manifest_mask_pixel(
        tmp_path, manifest_path, "ownership_mask", 0
    )

    with pytest.raises(
        ValueError,
        match="presentation asset masks do not match RGBA alpha",
    ):
        _load_test_presentation_assets(
            tmp_path, source, graph, manifest_path, graph_sha256
        )


def test_quality_reconstruction_uses_text_clean_pixels_without_alpha_holes(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    background = tmp_path / "background.png"
    text_mask = tmp_path / "text-mask.png"
    mask_dir = tmp_path / "graph/masks"
    mask_dir.mkdir(parents=True)
    component_mask = mask_dir / "component.png"
    nested_mask = mask_dir / "nested.png"
    Image.new("RGB", (4, 4), "red").save(source)
    Image.new("RGB", (4, 4), "white").save(background)
    Image.new("L", (4, 4), 255).save(component_mask)
    nested = Image.new("L", (4, 4), 0)
    for x in (2, 3):
        for y in (2, 3):
            nested.putpixel((x, y), 255)
    nested.save(nested_mask)
    nested.close()
    mask = Image.new("L", (4, 4), 0)
    mask.putpixel((1, 1), 255)
    mask.save(text_mask)
    mask.close()
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    foreground_evidence = run_dir / "input/foreground-evidence-mask.png"
    Image.new("L", (4, 4), 255).save(foreground_evidence)
    prepared = {
        "_prepared_schema_version": 5,
        "original_image_path": str(source),
        "background_original_path": str(background),
        "_text_mask_path": str(text_mask),
        "_foreground_evidence_mask_path": str(foreground_evidence),
        "text_items": [{"text": "editable"}],
    }

    class FakeImageModule:
        @staticmethod
        def load_component_layers(path):
            return prepared

    monkeypatch.setattr(
        legacy.importlib, "import_module", lambda name: FakeImageModule,
    )
    output_dir = run_dir / "pages/page_001/reconstruction/quality"
    output_dir.mkdir(parents=True)
    graph = {"nodes": [{
        "id": "component", "kind": "parent", "parent_id": None,
        "state": "pending_gate", "mask": "masks/component.png",
        "mask_sha256": hashlib.sha256(component_mask.read_bytes()).hexdigest(),
        "bbox": [0, 0, 4, 4], "z_index": 0, "text_ids": [],
    }, {
        "id": "nested", "kind": "parent", "parent_id": None,
        "state": "pending_gate", "mask": "masks/nested.png",
        "mask_sha256": hashlib.sha256(nested_mask.read_bytes()).hexdigest(),
        "bbox": [2, 2, 3, 3], "z_index": 1, "text_ids": [],
    }]}
    original_masks = {
        component_mask: component_mask.read_bytes(),
        nested_mask: nested_mask.read_bytes(),
    }
    original_graph = json.loads(json.dumps(graph))
    (mask_dir.parent / "component-graph.json").write_text(
        json.dumps(graph), encoding="utf-8"
    )

    refs = legacy._quality_assets(
        store, "page_001", graph, mask_dir.parent, output_dir,
    )

    assert refs["foreground_evidence"] == {
        "path": "input/foreground-evidence-mask.png",
        "sha256": hashlib.sha256(foreground_evidence.read_bytes()).hexdigest(),
    }
    with Image.open(run_dir / refs["reconstructed"]["path"]) as reconstructed:
        assert reconstructed.getpixel((0, 0)) == (255, 0, 0)
        assert reconstructed.getpixel((1, 1)) == (255, 0, 0)
    native = json.loads(
        (run_dir / refs["native_check"]["path"]).read_text(encoding="utf-8")
    )
    assert native["contained_parent_pairs"] == [["component", "nested"]]
    assert graph == original_graph
    assert all(path.read_bytes() == payload for path, payload in original_masks.items())


def test_quality_assets_remove_suppressed_ocr_and_restore_source_visual(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    text_clean = tmp_path / "text-clean.png"
    background = tmp_path / "background.png"
    text_mask = tmp_path / "text-mask.png"
    graph_dir = tmp_path / "graph"
    mask_dir = graph_dir / "masks"
    mask_dir.mkdir(parents=True)
    source_pixels = np.full((4, 6, 3), 255, dtype=np.uint8)
    source_pixels[1:3, 1:4] = (47, 111, 237)
    Image.fromarray(source_pixels, mode="RGB").save(source)
    Image.new("RGB", (6, 4), "white").save(text_clean)
    Image.new("RGB", (6, 4), "white").save(background)
    mistaken_region = np.zeros((4, 6), dtype=np.uint8)
    mistaken_region[1:3, 1:4] = 255
    Image.fromarray(mistaken_region, mode="L").save(text_mask)
    visual_mask = mask_dir / "visual.png"
    frozen_text_mask = mask_dir / "text_0001.png"
    Image.fromarray(mistaken_region, mode="L").save(visual_mask)
    Image.fromarray(mistaken_region, mode="L").save(frozen_text_mask)
    graph = {"nodes": [{
        "id": "visual", "kind": "parent", "parent_id": None,
        "state": "pending_gate", "mask": "masks/visual.png",
        "mask_sha256": hashlib.sha256(visual_mask.read_bytes()).hexdigest(),
        "bbox": [1, 1, 4, 3], "z_index": 0, "text_ids": [],
    }, {
        "id": "text_0001", "kind": "text", "parent_id": None,
        "state": "inactive", "mask": "masks/text_0001.png",
        "mask_sha256": hashlib.sha256(frozen_text_mask.read_bytes()).hexdigest(),
        "bbox": [1, 1, 4, 3], "z_index": 1, "text_ids": [],
    }]}
    (graph_dir / "component-graph.json").write_text(
        json.dumps(graph), encoding="utf-8",
    )
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    prepared = {
        "original_image_path": str(source),
        "background_original_path": str(background),
        "_text_clean_path": str(text_clean),
        "_text_mask_path": str(text_mask),
        "_text_cleanup_mask_path": str(text_mask),
        "text_items": [{"text": "mistaken", "box": [1, 1, 3, 2]}],
    }

    class FakeImageModule:
        @staticmethod
        def load_component_layers(path):
            return prepared

        @staticmethod
        def _interpolate_text_item_boxes(image, items, padding):
            return image

    monkeypatch.setattr(
        legacy.importlib, "import_module", lambda name: FakeImageModule,
    )
    original_effective_text_context = legacy._effective_text_context
    refine_text_clean_flags = []

    def record_refinement_policy(**kwargs):
        refine_text_clean_flags.append(kwargs["refine_text_clean"])
        return original_effective_text_context(**kwargs)

    monkeypatch.setattr(
        legacy, "_effective_text_context", record_refinement_policy,
    )
    output_dir = run_dir / "pages/page_001/reconstruction/quality"
    output_dir.mkdir(parents=True)

    refs = legacy._quality_assets(
        store, "page_001", graph, graph_dir, output_dir,
    )

    assert refine_text_clean_flags == [True]
    native = json.loads(
        (run_dir / refs["native_check"]["path"]).read_text(encoding="utf-8")
    )
    assert native["text_items"] == []
    with Image.open(run_dir / refs["text_mask"]["path"]) as effective_mask:
        assert effective_mask.getbbox() is None
    with Image.open(run_dir / refs["reconstructed"]["path"]) as reconstructed:
        assert reconstructed.getpixel((2, 1)) == (47, 111, 237)


def test_filtered_text_items_keep_their_original_component_ids() -> None:
    items = [{
        "_component_id": "text_0002",
        "text": "keep",
        "box": [2, 1, 3, 2],
    }]
    assert legacy._component_text_items(items, (8, 6)) == [{
        "id": "text_0002",
        "text": "keep",
        "box": [2, 1, 5, 3],
    }]


def test_component_text_items_reject_non_ascii_numeric_id() -> None:
    with pytest.raises(ValueError, match="component text id is invalid"):
        legacy._component_text_items(
            [{"_component_id": "text_１２３", "text": "x", "box": [0, 0, 1, 1]}],
            (2, 2),
        )


def test_frozen_presentation_records_are_reused_without_changing_graph_binding() -> None:
    old_frozen = {
        "component_id": "frozen",
        "rgba": {"path": "old-rgba.png", "sha256": "a" * 64},
        "ownership_mask": {"path": "old-own.png", "sha256": "b" * 64},
        "presentation_alpha_mask": {"path": "old-alpha.png", "sha256": "c" * 64},
        "generated_underlay_mask": {"path": "old-underlay.png", "sha256": "d" * 64},
        "metrics": {},
    }
    new_frozen = {
        **old_frozen,
        "rgba": {"path": "new-rgba.png", "sha256": "e" * 64},
    }
    current = {
        "schema_version": 1,
        "source_sha256": "f" * 64,
        "graph_sha256": "1" * 64,
        "components": [new_frozen, {"component_id": "pending"}],
    }
    previous = {
        **current,
        "graph_sha256": "0" * 64,
        "components": [old_frozen],
    }

    reused = legacy._reuse_frozen_presentation_records(
        current, previous, {"frozen"}
    )

    assert reused["graph_sha256"] == "1" * 64
    assert reused["components"] == [old_frozen, {"component_id": "pending"}]


def test_text_region_is_assigned_to_the_best_component_owner() -> None:
    text = np.zeros((8, 10), dtype=bool)
    text[2:6, 3:7] = True
    left = np.zeros_like(text)
    right = np.zeros_like(text)
    left[2:6, 3:5] = True
    right[2:4, 6:7] = True

    assigned = legacy._assign_text_regions_to_component_masks(
        [left, right], text
    )

    assert np.all(assigned[0][text])
    assert np.count_nonzero(assigned[1] & text) == 0
    assert np.count_nonzero(left & text) == 8


def test_item_text_region_fills_a_containing_sparse_component() -> None:
    text = np.zeros((12, 20), dtype=bool)
    text[4:8, 5:7] = True
    text[4:8, 9:11] = True
    text[4:8, 13:15] = True
    component = np.zeros_like(text)
    component[2:10, 2] = True
    component[2:10, 17] = True
    component[2, 2:18] = True
    component[9, 2:18] = True
    component[4:8, 5:8] = True

    assigned = legacy._assign_text_regions_to_component_masks(
        [component], text, [{"box": [5, 4, 10, 4]}]
    )

    assert np.all(assigned[0][2:10, 3:17])


def test_item_text_region_fills_text_hole_backed_by_surrounding_component() -> None:
    text = np.zeros((12, 20), dtype=bool)
    text[4:8, 5:7] = True
    text[4:8, 9:11] = True
    text[4:8, 13:15] = True
    component = np.zeros_like(text)
    component[2:10, 2:18] = True
    component[text] = False

    assigned = legacy._assign_text_regions_to_component_masks(
        [component], text, [{"box": [5, 4, 10, 4]}]
    )

    assert np.all(assigned[0][2:10, 3:17])


def test_item_text_region_does_not_expand_beyond_component_silhouette() -> None:
    text = np.zeros((100, 120), dtype=bool)
    text[30:60, 50:70] = True
    component = np.zeros_like(text)
    component[20:80, 20:100] = True
    component[text] = False
    silhouette = np.zeros_like(text)
    silhouette[20:80, 20:100] = True

    assigned = legacy._assign_text_regions_to_component_masks(
        [component], text, [{"box": [45, 23, 30, 51]}]
    )

    assert np.all(assigned[0][text])
    assert not np.any(assigned[0] & ~silhouette)


def test_text_halo_scales_beyond_the_bounded_pixel_repair() -> None:
    assert legacy._text_item_repair_padding_px(40) == 4
    assert legacy._text_item_halo_px(40) == 12


def test_item_text_region_does_not_steal_an_existing_nested_visual() -> None:
    text = np.zeros((12, 20), dtype=bool)
    text[4:8, 5:7] = True
    text[4:8, 9:11] = True
    text[4:8, 13:15] = True
    container = np.zeros_like(text)
    container[2:10, 2] = True
    container[2:10, 17] = True
    container[2, 2:18] = True
    container[9, 2:18] = True
    container[4:8, 5:8] = True
    nested = np.zeros_like(text)
    nested[5:7, 11:13] = True

    assigned = legacy._assign_text_regions_to_component_masks(
        [container, nested], text, [{"box": [5, 4, 10, 4]}]
    )

    assert np.all(assigned[0][4:8, 5:15] | nested[4:8, 5:15])
    assert np.array_equal(assigned[1], nested)
    assert not np.any(assigned[0] & assigned[1])


def test_item_text_region_does_not_expand_a_partial_icon_owner() -> None:
    text = np.zeros((12, 20), dtype=bool)
    text[4:8, 5:15] = True
    icon = np.zeros_like(text)
    icon[4:8, 5:8] = True

    assigned = legacy._assign_text_regions_to_component_masks(
        [icon], text, [{"box": [5, 4, 10, 4]}]
    )

    assert np.array_equal(assigned[0], icon)


def _background_box_mask(
    shape: tuple[int, int], box: tuple[int, int, int, int]
) -> np.ndarray:
    mask = np.zeros(shape, dtype=bool)
    left, top, right, bottom = box
    mask[top:bottom, left:right] = True
    return mask


def _write_background_action_graph(
    graph_dir: Path, masks: dict[str, np.ndarray]
) -> dict:
    nodes = []
    for z_index, (component_id, mask) in enumerate(masks.items()):
        path = graph_dir / "masks" / f"{component_id}.png"
        Image.fromarray(mask.astype(np.uint8) * 255, mode="L").save(path)
        ys, xs = np.where(mask)
        nodes.append({
            "id": component_id,
            "kind": "text" if component_id.startswith("text_") else "parent",
            "parent_id": None,
            "state": "frozen",
            "mask": f"masks/{component_id}.png",
            "mask_sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
            "bbox": [
                int(xs.min()), int(ys.min()),
                int(xs.max()) + 1, int(ys.max()) + 1,
            ],
            "z_index": z_index,
            "text_ids": [],
        })
    return {"nodes": nodes}


def test_rebuild_canvas_background_consumes_every_repair_request(
    tmp_path: Path,
) -> None:
    shape = (80, 120)
    source_pixels = np.full((*shape, 3), 240, dtype=np.uint8)
    source_pixels[8:24, 8:36] = 20
    source_pixels[40:64, 80:112] = 40
    source = tmp_path / "source.png"
    current = tmp_path / "current.png"
    Image.fromarray(source_pixels).save(source)
    Image.fromarray(source_pixels).save(current)
    graph_dir = tmp_path / "graph"
    (graph_dir / "masks").mkdir(parents=True)
    left = _background_box_mask(shape, (8, 8, 36, 24))
    right = _background_box_mask(shape, (80, 40, 112, 64))
    graph = _write_background_action_graph(
        graph_dir,
        {"text_left": left, "component_right": right},
    )
    text_mask = tmp_path / "text-mask.png"
    Image.fromarray(left.astype(np.uint8) * 255, mode="L").save(text_mask)
    output = tmp_path / "rebuilt.png"

    legacy._rebuild_canvas_background(
        source_path=source,
        current_background_path=current,
        repair_requests=[
            ({"text_left"}, 0.01),
            ({"component_right"}, 0.02),
        ],
        graph=graph,
        graph_dir=graph_dir,
        text_mask_path=text_mask,
        output_path=output,
    )

    actual = np.asarray(Image.open(output).convert("RGB"))
    assert not np.array_equal(actual[8:24, 8:36], source_pixels[8:24, 8:36])
    assert not np.array_equal(actual[40:64, 80:112], source_pixels[40:64, 80:112])
    assert np.array_equal(actual[:4, :4], source_pixels[:4, :4])


def test_execute_legacy_round_aggregates_background_actions(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    from scripts import sam_worker

    run_dir = tmp_path / "run"
    reconstruction = run_dir / "pages/page_001/reconstruction"
    round_dir = reconstruction / "round-01"
    round_dir.mkdir(parents=True)
    store = RunStore(run_dir)

    source = round_dir / "source.png"
    background = reconstruction / "background.png"
    text_mask = reconstruction / "text-mask.png"
    for path in (source, background):
        Image.new("RGB", (12, 8), "white").save(path)
    Image.new("L", (12, 8), 0).save(text_mask)
    quality = round_dir / "quality-report.json"
    quality.write_text("{}", encoding="utf-8")
    presentation_manifest = round_dir / "presentation-manifest.json"
    presentation_manifest.write_text("{}", encoding="utf-8")
    graph_path = round_dir / "component-graph.json"
    graph_path.write_text(json.dumps({"nodes": []}), encoding="utf-8")

    def reference(path: Path) -> dict[str, str]:
        return {
            "path": path.relative_to(run_dir).as_posix(),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }

    request = round_dir / "request.json"
    request.write_text(json.dumps({
        "evidence": {
            "source.png": {"path": "source.png"},
            "quality-report.json": {"path": "quality-report.json"},
            "presentation-manifest.json": {"path": "presentation-manifest.json"},
        },
    }), encoding="utf-8")
    plan = round_dir / "plan.json"
    plan.write_text(json.dumps({"actions": [
        {
            "action": "rebuild_background",
            "object_ids": ["text_left"],
            "parameters": {"margin_ratio": 0.01},
        },
        {
            "action": "rebuild_background",
            "object_ids": ["component_right"],
            "parameters": {"margin_ratio": 0.02},
        },
    ]}), encoding="utf-8")
    store.write_json("pages/page_001/reconstruction/component_state.json", {
        "repair_round": 1,
        "provider": "host",
        "graph_ref": reference(graph_path),
        "current_round": {
            "request_ref": reference(request),
            "plan_ref": reference(plan),
        },
        "frozen": {},
    })

    def execute_round(pixels, graph, actions, **kwargs):
        prompt = {
            "component_id": "component",
            "box": [1, 1, 4, 4],
            "positive": [],
            "negative": [],
        }
        with pytest.raises(RuntimeError, match="mask count"):
            kwargs["sam_batch_runner"](image=pixels, prompts=[prompt])
        output_dir = kwargs["output_dir"]
        output_dir.mkdir(parents=True)
        (output_dir / "component-graph.json").write_text(
            json.dumps(graph), encoding="utf-8"
        )
        return graph

    captured: dict[str, Any] = {}

    def rebuild_background(**kwargs):
        captured["repair_requests"] = (
            kwargs["repair_requests"]
            if "repair_requests" in kwargs
            else [(kwargs["component_ids"], kwargs["margin_ratio"])]
        )
        return kwargs["current_background_path"]

    fake_image_module = types.SimpleNamespace(
        load_component_layers=lambda path: {
            "_text_clean_path": str(source),
            "_text_cleanup_mask_path": str(text_mask),
            "_text_mask_path": str(text_mask),
            "background_original_path": str(background),
            "text_items": [],
        }
    )
    monkeypatch.setattr(legacy, "execute_component_action_round", execute_round)
    monkeypatch.setattr(
        sam_worker,
        "run_component_prompt_batch_worker",
        lambda *args, **kwargs: [
            np.ones((8, 12), dtype=bool),
            np.ones((8, 12), dtype=bool),
        ],
    )
    monkeypatch.setattr(legacy, "_ensure_component_disk_reserve", lambda *a, **k: None)
    monkeypatch.setattr(legacy, "_effective_text_context", lambda **kwargs: (
        [], np.zeros((8, 12), dtype=bool), np.full((8, 12, 3), 255, dtype=np.uint8)
    ))
    monkeypatch.setattr(legacy, "_rebuild_canvas_background", rebuild_background)
    monkeypatch.setattr(legacy, "_quality_assets", lambda *a, **k: {})
    monkeypatch.setattr(legacy, "record_component_execution", lambda *a, **k: None)
    monkeypatch.setattr(
        legacy.importlib, "import_module", lambda name: fake_image_module
    )

    legacy._execute_legacy_round(store, "page_001", object())

    assert captured["repair_requests"] == [
        ({"text_left"}, 0.01),
        ({"component_right"}, 0.02),
    ]


def test_agent_can_rebuild_uniform_canvas_under_active_components(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    current = tmp_path / "current.png"
    text_mask = tmp_path / "text-mask.png"
    graph_dir = tmp_path / "graph"
    masks = graph_dir / "masks"
    masks.mkdir(parents=True)
    component_mask = masks / "component.png"
    source_image = Image.new("RGB", (40, 30), "white")
    ImageDraw.Draw(source_image).rectangle((12, 8, 27, 21), fill="black")
    source_image.save(source)
    source_image.close()
    current_image = Image.new("RGB", (40, 30), "white")
    ImageDraw.Draw(current_image).rectangle((10, 6, 29, 23), fill=(180, 180, 180))
    current_image.save(current)
    current_image.close()
    mask = Image.new("L", (40, 30), 0)
    ImageDraw.Draw(mask).rectangle((12, 8, 27, 21), fill=255)
    mask.save(component_mask)
    Image.new("L", (40, 30), 0).save(text_mask)
    graph = {"nodes": [{
        "id": "component", "kind": "parent", "parent_id": None,
        "state": "pending", "mask": "masks/component.png",
        "mask_sha256": hashlib.sha256(component_mask.read_bytes()).hexdigest(),
        "bbox": [12, 8, 28, 22], "z_index": 0, "text_ids": [],
    }]}
    output = tmp_path / "rebuilt.png"

    legacy._rebuild_canvas_background(
        source_path=source, current_background_path=current,
        repair_requests=[({"component"}, 0.1)],
        graph=graph, graph_dir=graph_dir, text_mask_path=text_mask,
        output_path=output,
    )

    with Image.open(output) as rebuilt:
        assert rebuilt.getpixel((20, 15)) == (255, 255, 255)
        assert rebuilt.getpixel((0, 0)) == (255, 255, 255)


def test_background_rebuild_cleans_discarded_component_residual(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    current = tmp_path / "current.png"
    text_mask = tmp_path / "text-mask.png"
    graph_dir = tmp_path / "graph"
    masks = graph_dir / "masks"
    masks.mkdir(parents=True)
    component_mask = masks / "discarded.png"
    Image.new("RGB", (40, 30), "white").save(source)
    current_image = Image.new("RGB", (40, 30), "white")
    ImageDraw.Draw(current_image).rectangle((12, 8, 27, 21), fill=(180, 180, 180))
    current_image.save(current)
    mask = Image.new("L", (40, 30), 0)
    ImageDraw.Draw(mask).rectangle((12, 8, 27, 21), fill=255)
    mask.save(component_mask)
    Image.new("L", (40, 30), 0).save(text_mask)
    graph = {"nodes": [{
        "id": "discarded", "kind": "parent", "parent_id": None,
        "state": "inactive", "mask": "masks/discarded.png",
        "mask_sha256": hashlib.sha256(component_mask.read_bytes()).hexdigest(),
        "bbox": [12, 8, 28, 22], "z_index": 0, "text_ids": [],
    }]}
    output = tmp_path / "rebuilt.png"

    legacy._rebuild_canvas_background(
        source_path=source, current_background_path=current,
        repair_requests=[({"discarded"}, 0.01)],
        graph=graph, graph_dir=graph_dir, text_mask_path=text_mask,
        output_path=output,
    )

    with Image.open(output) as rebuilt:
        assert rebuilt.getpixel((20, 15)) == (255, 255, 255)


def test_background_rebuild_restores_structure_and_clears_only_selected_visual(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    current = tmp_path / "current.png"
    restored = tmp_path / "text-clean.png"
    text_mask = tmp_path / "text-mask.png"
    graph_dir = tmp_path / "graph"
    masks = graph_dir / "masks"
    masks.mkdir(parents=True)
    Image.new("RGB", (40, 30), "white").save(source)
    current_image = Image.new("RGB", (40, 30), "white")
    ImageDraw.Draw(current_image).rectangle(
        (28, 8, 34, 14), fill=(210, 210, 210)
    )
    current_image.save(current)
    clean = Image.new("RGB", (40, 30), "white")
    draw = ImageDraw.Draw(clean)
    draw.line((4, 20, 35, 20), fill="blue", width=1)
    draw.rectangle((12, 8, 18, 14), fill="red")
    draw.rectangle((28, 8, 34, 14), fill="green")
    clean.save(restored)
    Image.new("L", (40, 30), 0).save(text_mask)
    selected_mask = masks / "selected.png"
    mask = Image.new("L", (40, 30), 0)
    ImageDraw.Draw(mask).rectangle((12, 8, 18, 14), fill=255)
    mask.save(selected_mask)
    other_mask = masks / "other.png"
    mask = Image.new("L", (40, 30), 0)
    ImageDraw.Draw(mask).rectangle((28, 8, 34, 14), fill=255)
    mask.save(other_mask)
    graph = {"nodes": [{
        "id": "selected", "kind": "child", "parent_id": "parent",
        "state": "frozen", "mask": "masks/selected.png",
        "mask_sha256": hashlib.sha256(selected_mask.read_bytes()).hexdigest(),
        "bbox": [12, 8, 19, 15], "z_index": 0, "text_ids": [],
    }, {
        "id": "other", "kind": "child", "parent_id": "parent",
        "state": "frozen", "mask": "masks/other.png",
        "mask_sha256": hashlib.sha256(other_mask.read_bytes()).hexdigest(),
        "bbox": [28, 8, 35, 15], "z_index": 1, "text_ids": [],
    }]}
    output = tmp_path / "rebuilt.png"
    from scripts import component_underlay

    choose_visual_fill = component_underlay._choose_visual_fill
    allow_original_values = []

    def record_fill_policy(**kwargs):
        allow_original_values.append(kwargs.get("allow_original"))
        return choose_visual_fill(**kwargs)

    monkeypatch.setattr(
        component_underlay, "_choose_visual_fill", record_fill_policy
    )

    legacy._rebuild_canvas_background(
        source_path=source, current_background_path=current,
        restore_background_path=restored,
        repair_requests=[({"selected"}, 0.01)],
        graph=graph, graph_dir=graph_dir, text_mask_path=text_mask,
        output_path=output,
    )

    with Image.open(output) as rebuilt:
        assert rebuilt.getpixel((15, 11)) == (255, 255, 255)
        assert rebuilt.getpixel((31, 11)) == (255, 255, 255)
        assert rebuilt.getpixel((20, 20)) == (0, 0, 255)
    assert allow_original_values == [False]


def test_background_rebuild_preserves_local_tinted_surface(tmp_path: Path) -> None:
    source = tmp_path / "source.png"
    current = tmp_path / "current.png"
    restored = tmp_path / "text-clean.png"
    text_mask = tmp_path / "text-mask.png"
    graph_dir = tmp_path / "graph"
    masks = graph_dir / "masks"
    masks.mkdir(parents=True)
    component_mask = masks / "component.png"
    source_image = Image.new("RGB", (60, 40), "white")
    draw = ImageDraw.Draw(source_image)
    draw.rectangle((8, 4, 51, 35), fill=(235, 248, 237))
    draw.rectangle((25, 15, 34, 24), fill=(20, 120, 40))
    source_image.save(source)
    source_image.save(current)
    source_image.save(restored)
    mask = Image.new("L", (60, 40), 0)
    ImageDraw.Draw(mask).rectangle((25, 15, 34, 24), fill=255)
    mask.save(component_mask)
    Image.new("L", (60, 40), 0).save(text_mask)
    graph = {"nodes": [{
        "id": "component", "kind": "parent", "parent_id": None,
        "state": "pending", "mask": "masks/component.png",
        "mask_sha256": hashlib.sha256(component_mask.read_bytes()).hexdigest(),
        "bbox": [25, 15, 35, 25], "z_index": 0, "text_ids": [],
    }]}
    output = tmp_path / "rebuilt.png"

    legacy._rebuild_canvas_background(
        source_path=source, current_background_path=current,
        restore_background_path=restored,
        repair_requests=[({"component"}, 0.025)],
        graph=graph, graph_dir=graph_dir, text_mask_path=text_mask,
        output_path=output,
    )

    with Image.open(output) as rebuilt:
        pixel = rebuilt.getpixel((29, 19))
        assert max(abs(pixel[index] - value) for index, value in enumerate((235, 248, 237))) <= 6


def test_background_rebuild_does_not_expand_text_box_into_neighbor_content(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    current = tmp_path / "current.png"
    text_mask = tmp_path / "text-mask.png"
    graph_dir = tmp_path / "graph"
    graph_dir.mkdir()
    source_image = Image.new("RGB", (40, 30), "white")
    ImageDraw.Draw(source_image).rectangle((16, 12, 18, 14), fill="blue")
    source_image.save(source)
    source_image.save(current)
    mask = Image.new("L", (40, 30), 0)
    ImageDraw.Draw(mask).rectangle((12, 12, 14, 14), fill=255)
    mask.save(text_mask)
    output = tmp_path / "rebuilt.png"

    legacy._rebuild_canvas_background(
        source_path=source, current_background_path=current,
        repair_requests=[],
        graph={"nodes": []}, graph_dir=graph_dir, text_mask_path=text_mask,
        output_path=output,
    )

    with Image.open(output) as rebuilt:
        assert rebuilt.getpixel((16, 13)) == (0, 0, 255)


def test_background_rebuild_supports_nonuniform_canvas_border(tmp_path: Path) -> None:
    source = tmp_path / "source.png"
    current = tmp_path / "current.png"
    text_mask = tmp_path / "text-mask.png"
    graph_dir = tmp_path / "graph"
    masks = graph_dir / "masks"
    masks.mkdir(parents=True)
    component_mask = masks / "component.png"
    source_image = Image.new("RGB", (40, 30), "white")
    draw = ImageDraw.Draw(source_image)
    draw.rectangle((0, 0, 39, 2), fill="red")
    draw.rectangle((0, 27, 39, 29), fill="blue")
    source_image.save(source)
    Image.new("RGB", (40, 30), "white").save(current)
    mask = Image.new("L", (40, 30), 0)
    ImageDraw.Draw(mask).rectangle((12, 8, 27, 21), fill=255)
    mask.save(component_mask)
    Image.new("L", (40, 30), 0).save(text_mask)
    graph = {"nodes": [{
        "id": "component", "kind": "parent", "parent_id": None,
        "state": "pending", "mask": "masks/component.png",
        "mask_sha256": hashlib.sha256(component_mask.read_bytes()).hexdigest(),
        "bbox": [12, 8, 28, 22], "z_index": 0, "text_ids": [],
    }]}

    output = tmp_path / "rebuilt.png"
    legacy._rebuild_canvas_background(
        source_path=source, current_background_path=current,
        repair_requests=[({"component"}, 0.01)],
        graph=graph, graph_dir=graph_dir, text_mask_path=text_mask,
        output_path=output,
    )

    with Image.open(output) as rebuilt:
        assert rebuilt.getpixel((20, 15)) == (255, 255, 255)


def test_initial_page_session_reserves_disk_before_writing_evidence(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    background = tmp_path / "background.png"
    difference = tmp_path / "difference.png"
    text_mask = tmp_path / "text-mask.png"
    for path in (source, background, difference):
        Image.new("RGB", (20, 10), "black").save(path)
    Image.new("L", (20, 10), 0).save(text_mask)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    reconstruction = run_dir / "pages/page_001/reconstruction"
    monkeypatch.setattr(
        legacy.shutil,
        "disk_usage",
        lambda path: types.SimpleNamespace(free=1),
    )

    with pytest.raises(RuntimeError, match="disk reserve"):
        legacy._build_initial_page_session(
            RunStore.open(run_dir),
            "page_001",
            {
                "original_image_path": str(source),
                "background_original_path": str(background),
                "background_difference_path": str(difference),
                "_text_mask_path": str(text_mask),
                "_element_mask_paths": [],
                "components": [],
                "text_items": [],
            },
            reconstruction,
        )

    assert not (reconstruction / "evidence-source").exists()


def test_accepted_presentation_assembly_preserves_bound_underlay_and_text_clean_rgb(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", slide_size="16:9"
    )
    store = RunStore.open(run_dir)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    assets = reconstruction / "accepted"
    masks = assets / "masks"
    presentation_assets = assets / "presentation-assets"
    masks.mkdir(parents=True)
    presentation_assets.mkdir()
    accepted_source = assets / "source.png"
    background = assets / "background.png"
    reconstructed = assets / "reconstructed.png"
    text_mask = assets / "text-mask.png"
    native = assets / "native.json"
    source_pixels = np.full((4, 4, 3), 255, dtype=np.uint8)
    source_pixels[1:3, 1:3] = (0, 0, 255)
    source_pixels[0, 0] = (255, 0, 0)
    Image.fromarray(source_pixels, mode="RGB").save(accepted_source)
    Image.new("RGB", (4, 4), "white").save(background)
    Image.new("RGB", (4, 4), "red").save(reconstructed)
    Image.new("L", (4, 4), 0).save(text_mask)
    with Image.open(text_mask) as image:
        image.putpixel((0, 0), 255)
        image.save(text_mask)
    native.write_text("{}", encoding="utf-8")
    base_mask_path = masks / "base.png"
    child_mask_path = masks / "child.png"
    text_node_mask_path = masks / "text.png"
    Image.new("L", (4, 4), 255).save(base_mask_path)
    child_mask = np.zeros((4, 4), dtype=np.uint8)
    child_mask[1:3, 1:3] = 255
    Image.fromarray(child_mask, mode="L").save(child_mask_path)
    text_node_mask = np.zeros((4, 4), dtype=np.uint8)
    text_node_mask[0, 0] = 255
    Image.fromarray(text_node_mask, mode="L").save(text_node_mask_path)
    graph = {
        "nodes": [
            {
                "id": "base", "kind": "parent", "parent_id": None,
                "state": "frozen", "mask": "masks/base.png",
                "mask_sha256": hashlib.sha256(
                    base_mask_path.read_bytes()
                ).hexdigest(),
                "bbox": [0, 0, 4, 4], "z_index": 0,
                "text_ids": ["text_1"],
            },
            {
                "id": "child", "kind": "parent", "parent_id": None,
                "state": "frozen", "mask": "masks/child.png",
                "mask_sha256": hashlib.sha256(
                    child_mask_path.read_bytes()
                ).hexdigest(),
                "bbox": [1, 1, 3, 3], "z_index": 1, "text_ids": [],
            },
            {
                "id": "text_1", "kind": "text", "parent_id": None,
                "state": "frozen", "mask": "masks/text.png",
                "mask_sha256": hashlib.sha256(
                    text_node_mask_path.read_bytes()
                ).hexdigest(),
                "bbox": [0, 0, 1, 1], "z_index": 2, "text_ids": [],
            },
        ]
    }
    graph_path = assets / "component-graph.json"
    graph_path.write_text(json.dumps(graph), encoding="utf-8")

    def ref(path: Path) -> dict:
        return {
            "path": path.resolve().relative_to(run_dir.resolve()).as_posix(),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }

    ownership_base = np.full((4, 4), 255, dtype=np.uint8)
    ownership_base[1:3, 1:3] = 0
    ownership_base[0, 0] = 0
    generated_base = np.zeros((4, 4), dtype=np.uint8)
    generated_base[1:3, 1:3] = 255
    generated_base[0, 0] = 255
    alpha_base = np.full((4, 4), 255, dtype=np.uint8)
    base_rgba = np.zeros((4, 4, 4), dtype=np.uint8)
    base_rgba[:, :, :3] = (255, 255, 255)
    base_rgba[1:3, 1:3, :3] = (0, 255, 0)
    base_rgba[:, :, 3] = alpha_base
    child_rgba = np.zeros((4, 4, 4), dtype=np.uint8)
    child_rgba[1:3, 1:3, :3] = source_pixels[1:3, 1:3]
    child_rgba[1:3, 1:3, 3] = 255

    def presentation_component(
        component_id: str,
        index: int,
        rgba: np.ndarray,
        ownership: np.ndarray,
        alpha: np.ndarray,
        generated: np.ndarray,
    ) -> dict:
        references = {}
        for name, array, mode in (
            ("rgba", rgba, "RGBA"),
            ("ownership_mask", ownership, "L"),
            ("presentation_alpha_mask", alpha, "L"),
            ("generated_underlay_mask", generated, "L"),
        ):
            path = presentation_assets / f"{index:04d}-{name}.png"
            Image.fromarray(array, mode=mode).save(path)
            references[name] = ref(path)
        return {
            "component_id": component_id,
            **references,
            "metrics": {
                "boundary_color_mae": 0.0,
                "gradient_jump_p95": 0.0,
                "added_high_frequency_pixels": 0.0,
            },
        }

    presentation_manifest = presentation_assets / "presentation-manifest.json"
    presentation_manifest.write_text(json.dumps({
        "schema_version": 1,
        "source_sha256": ref(accepted_source)["sha256"],
        "graph_sha256": ref(graph_path)["sha256"],
        "components": [
            presentation_component(
                "base", 1, base_rgba, ownership_base,
                alpha_base, generated_base,
            ),
            presentation_component(
                "child", 2, child_rgba, child_mask, child_mask,
                np.zeros((4, 4), dtype=np.uint8),
            ),
        ],
    }), encoding="utf-8")

    result = {
        "schema_version": 1, "page_id": "page_001",
        "status": "ready_for_assembly", "provider": "host",
        "repair_rounds": 1, "initial_component_count": 2,
        "final_component_ids": ["base", "child"], "graph_ref": ref(graph_path),
        "accepted_graph_sha256": ref(graph_path)["sha256"],
        "round_history": [], "fallback": {"status": "none", "parent_ids": []},
        "accepted_asset_refs": {
            "source": ref(accepted_source), "background": ref(background),
            "reconstructed": ref(reconstructed), "text_mask": ref(text_mask),
            "native_check": ref(native),
            "presentation_manifest": ref(presentation_manifest),
        },
        "delivery_checks": {"pptx_reopen": "unknown"},
    }
    slide = legacy._accepted_slide_data(
        store,
        reconstruction,
        {"text_items": [{"text": "editable"}]},
        result,
    )

    assert [item["component_id"] for item in slide["components"]] == [
        "base", "child",
    ]
    with Image.open(slide["components"][0]["path"]) as component:
        rgba = np.asarray(component.convert("RGBA"))
    assert np.all(rgba[1:3, 1:3, 3] == 255)
    assert np.all(rgba[1:3, 1:3, :3] == (0, 255, 0))
    assert not np.array_equal(rgba[1:3, 1:3, :3], source_pixels[1:3, 1:3])
    assert tuple(rgba[0, 0, :3]) == (255, 255, 255)
    assert not np.array_equal(rgba[0, 0, :3], source_pixels[0, 0])


def _accepted_presentation_case(tmp_path: Path) -> tuple[
    RunStore, Path, dict, Path, Path
]:
    run_root = tmp_path / "run"
    reconstruction = run_root / "pages/page_001/reconstruction"
    accepted = reconstruction / "accepted"
    masks = accepted / "masks"
    masks.mkdir(parents=True)
    source = accepted / "source.png"
    background = accepted / "background.png"
    reconstructed = accepted / "reconstructed.png"
    text_mask = accepted / "text-mask.png"
    foreground_evidence = accepted / "foreground-evidence-mask.png"
    native = accepted / "native.json"
    for path, color in ((source, "white"), (background, "black"),
                        (reconstructed, "red")):
        Image.new("RGB", (4, 4), color).save(path)
    Image.new("L", (4, 4), 0).save(text_mask)
    Image.new("L", (4, 4), 255).save(foreground_evidence)
    native.write_text("{}", encoding="utf-8")
    mask = masks / "component_0001.png"
    Image.new("L", (4, 4), 255).save(mask)
    graph = {"nodes": [{
        "id": "component_0001", "kind": "parent", "parent_id": None,
        "state": "frozen", "mask": "masks/component_0001.png",
        "mask_sha256": hashlib.sha256(mask.read_bytes()).hexdigest(),
        "bbox": [0, 0, 4, 4], "z_index": 0, "text_ids": [],
    }]}
    graph_path = accepted / "component-graph.json"
    graph_path.write_text(json.dumps(graph), encoding="utf-8")
    store = RunStore(run_root)

    def ref(path: Path) -> dict[str, str]:
        return {
            "path": path.resolve().relative_to(run_root.resolve()).as_posix(),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }

    presentation_output = accepted / "presentation"
    presentation_output.mkdir()
    manifest_path, _ = _build_test_presentation_manifest(
        run_root,
        source_path=source,
        text_clean_path=source,
        graph=graph,
        graph_dir=accepted,
        output_dir=presentation_output,
    )
    result = {
        "accepted_asset_refs": {
            name: ref(path) for name, path in {
                "source": source, "background": background,
                "reconstructed": reconstructed, "text_mask": text_mask,
                "foreground_evidence": foreground_evidence,
                "native_check": native,
            }.items()
        },
        "graph_ref": ref(graph_path),
        "accepted_graph_sha256": ref(graph_path)["sha256"],
        "final_component_ids": ["component_0001"],
    }
    result["accepted_asset_refs"]["presentation_manifest"] = ref(manifest_path)
    return store, reconstruction, result, manifest_path, reconstructed


def test_accepted_slide_uses_effective_text_items_from_result(tmp_path: Path) -> None:
    store, reconstruction, result, _, _ = _accepted_presentation_case(tmp_path)
    result["text_items"] = []

    slide = legacy._accepted_slide_data(
        store,
        reconstruction,
        {"text_items": [{"text": "mistaken", "box": [1, 1, 2, 2]}]},
        result,
    )

    assert slide["text_items"] == []


def _accepted_assembly_job(tmp_path: Path) -> tuple[RunStore, Path, Path]:
    import image_to_ppt

    store, reconstruction, result, _, _ = _accepted_presentation_case(tmp_path)
    output = tmp_path / "accepted-output.pptx"
    initial = reconstruction / "initial"
    initial.mkdir()
    prepared_source = initial / "source.png"
    prepared_background = initial / "background.png"
    prepared_text_mask = initial / "text-mask.png"
    prepared_foreground = initial / "foreground-evidence-mask.png"
    prepared_removal_mask = initial / "removal-mask.png"
    prepared_difference = initial / "difference.png"
    Image.new("RGB", (16, 9), "white").save(prepared_source)
    Image.new("RGB", (16, 9), "white").save(prepared_background)
    Image.new("L", (16, 9), 0).save(prepared_text_mask)
    Image.new("L", (16, 9), 0).save(prepared_foreground)
    Image.new("L", (16, 9), 0).save(prepared_removal_mask)
    Image.new("RGB", (16, 9), "black").save(prepared_difference)
    image_to_ppt._write_prepared_page({
        "img_width": 16,
        "img_height": 9,
        "canvas_width": 16,
        "canvas_height": 9,
        "content_offset_x": 0,
        "content_offset_y": 0,
        "widescreen_background_method": "identity",
        "original_image_path": str(prepared_source),
        "background_original_path": str(prepared_background),
        "background_widescreen_path": str(prepared_background),
        "background_removal_mask_path": str(prepared_removal_mask),
        "background_difference_path": str(prepared_difference),
        "_text_mask_path": str(prepared_text_mask),
        "_foreground_evidence_mask_path": str(prepared_foreground),
        "_element_mask_paths": [],
        "_semantic_mask_paths": [],
        "_resource_isolation": False,
        "_initial_diagnostics": [],
        "components": [],
        "text_items": [{
            "box": [0, 2, 8, 4],
            "text": "editable",
            "font_size": 10.0,
            "color": "#000000",
            "bold": False,
            "font": "Arial",
            "align": 1,
            "confidence": 1.0,
        }],
    }, initial)
    result["schema_version"] = 1
    result["page_id"] = "page_001"
    result["status"] = "ready_for_assembly"
    result_path = reconstruction / "component_result.json"
    result_path.write_text(json.dumps(result), encoding="utf-8")

    def ref(path: Path) -> dict[str, str]:
        return {
            "path": path.relative_to(store.root).as_posix(),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }

    store.write_json("job_manifest.json", {
        "pages": ["page_001"],
        "options": {"output_path": str(output), "slide_size": "16:9"},
        "input": {},
    })
    store.write_json("pages/page_001/reconstruction/component_state.json", {
        "status": "ready_for_assembly", "result_ref": ref(result_path),
    })
    return store, reconstruction, output


def _publish_native_route_sidecar(
    store: RunStore, reconstruction: Path, *, object_id: str
) -> dict:
    result_path = reconstruction / "component_result.json"
    result = json.loads(result_path.read_text(encoding="utf-8"))
    graph_path = store.root / result["graph_ref"]["path"]
    graph = json.loads(graph_path.read_text(encoding="utf-8"))
    node = next(item for item in graph["nodes"] if item["id"] == object_id)
    mask_path = graph_path.parent / node["mask"]

    def ref(path: Path) -> dict[str, str]:
        return {
            "path": path.relative_to(store.root).as_posix(),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }

    source_ref = result["accepted_asset_refs"]["source"]
    route_dir = reconstruction / "route"
    route_dir.mkdir()
    ir = {
        "schema_version": 1,
        "page_id": "page_001",
        "canvas": {"width": 4, "height": 4},
        "objects": [{
            "id": object_id,
            "bbox": node["bbox"],
            "z_index": node["z_index"],
            "source_refs": [source_ref],
            "mask_ref": ref(mask_path),
            "relations": [],
            "candidate_representations": [
                {
                    "kind": "raster_component",
                    "confidence": 1.0,
                    "payload": {"asset_ref": source_ref},
                    "evidence_refs": [],
                    "required_qa_checks": [],
                },
                {
                    "kind": "native_shape",
                    "confidence": 1.0,
                    "payload": {
                        "shape_type": "rectangle",
                        "fill_rgb": [255, 0, 0],
                    },
                    "evidence_refs": [],
                    "required_qa_checks": ["render_difference"],
                },
            ],
        }],
    }
    ir_path = route_dir / "reconstruction-ir.json"
    ir_path.write_text(json.dumps(ir), encoding="utf-8")
    from image2editable.reconstruction_contracts import reconstruction_ir_sha256

    plan = {
        "schema_version": 1,
        "page_id": "page_001",
        "ir_sha256": reconstruction_ir_sha256(ir),
        "adapter": "pptx",
        "routes": [{
            "object_id": object_id,
            "selected_route": "native_shape",
            "fallback_route": "raster_component",
            "candidate_confidence": 1.0,
            "evidence_refs": [],
            "qa_requirements": ["render_difference"],
        }],
    }
    plan_path = route_dir / "reconstruction-plan.json"
    plan_path.write_text(json.dumps(plan), encoding="utf-8")
    qa = {
        "schema_version": 1,
        "renderer": {"renderer": "powerpoint", "available": True},
        "initial": {"accepted": True},
        "fallback": None,
        "final_plan_sha256": hashlib.sha256(
            (
                json.dumps(
                    plan,
                    ensure_ascii=False,
                    indent=2,
                    sort_keys=True,
                )
                + "\n"
            ).encode("utf-8")
        ).hexdigest(),
    }
    qa_path = route_dir / "render-qa.json"
    qa_path.write_text(json.dumps(qa), encoding="utf-8")
    route_result = {
        "schema_version": 1,
        "page_id": "page_001",
        "status": "native_accepted",
        "component_result_sha256": hashlib.sha256(
            result_path.read_bytes()
        ).hexdigest(),
        "ir_ref": ref(ir_path),
        "plan_ref": ref(plan_path),
        "qa_ref": ref(qa_path),
        "reason": None,
    }
    route_result_path = route_dir / "route_result.json"
    route_result_path.write_text(json.dumps(route_result), encoding="utf-8")
    return ref(route_result_path)


def test_accepted_presentation_pptx_e2e_cleans_temporary_assets(
    tmp_path: Path,
) -> None:
    store, reconstruction, output = _accepted_assembly_job(tmp_path)

    outputs = legacy.assemble_legacy_results(store)

    reopened = Presentation(outputs["16:9"])
    assert len(reopened.slides) == 1
    assert any(
        getattr(shape, "text", "") == "editable"
        for shape in reopened.slides[0].shapes
    )
    assert any(shape.shape_type == 13 for shape in reopened.slides[0].shapes)
    assert output.is_file()


def test_accepted_presentation_consumes_published_native_route(tmp_path: Path) -> None:
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    store, reconstruction, _ = _accepted_assembly_job(tmp_path)
    expected_route_ref = _publish_native_route_sidecar(
        store, reconstruction, object_id="component_0001"
    )

    outputs = legacy.assemble_legacy_results(store)

    presentation = Presentation(outputs["16:9"])
    content = list(presentation.slides[0].shapes)[1:]
    assert content[0].shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    assert content[0].name == "image2editable:component_0001"
    delivery = store.read_json(
        "pages/page_001/reconstruction/component_delivery.json"
    )
    assert delivery["route_result"] == expected_route_ref


def test_runtime_finalizes_ready_route_with_target_capabilities(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    store = RunStore(tmp_path / "run")
    source = store.root / "pages/page_001/source.png"
    source.parent.mkdir(parents=True)
    Image.new("RGB", (4, 4), "white").save(source)
    result_path = store.root / "pages/page_001/reconstruction/component_result.json"
    result_path.parent.mkdir()
    result_path.write_text("{}", encoding="utf-8")

    def ref(path: Path) -> dict[str, str]:
        return {
            "path": path.relative_to(store.root).as_posix(),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }

    store.write_json("pages/page_001/page_request.json", {
        "schema_version": 1,
        "source": source.relative_to(store.root).as_posix(),
        "sha256": ref(source)["sha256"],
    })
    store.write_json("pages/page_001/reconstruction/component_state.json", {
        "status": "ready_for_assembly",
        "result_ref": ref(result_path),
    })
    captured = {}

    def finalize(context, *, renderer, policy):
        captured.update(context=context, renderer=renderer, policy=policy)
        return {"status": "raster_fallback"}

    monkeypatch.setattr(runtime, "_finalize_reconstruction_route", finalize)
    monkeypatch.setattr(
        runtime,
        "_discover_powerpoint_renderer",
        lambda: (_ for _ in ()).throw(
            AssertionError("PSD route must not discover PowerPoint")
        ),
    )

    runtime._finalize_reconstruction_routes(
        store, {"output_format": "psd"}, ["page_001"]
    )

    assert captured["context"].component_result_path == result_path
    assert captured["context"].capabilities == frozenset(
        {"editable_text", "raster_component"}
    )
    assert captured["policy"]["native_shape_enabled"] is False
    assert captured["renderer"].available() is False


def test_accepted_presentation_psd_uses_final_agent_layers(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    store, reconstruction, _ = _accepted_assembly_job(tmp_path)
    output = tmp_path / "accepted-output.psd"
    manifest = store.read_json("job_manifest.json")
    manifest["output_format"] = "psd"
    manifest["options"]["output_path"] = str(output)
    store.write_json("job_manifest.json", manifest)
    captured = {}

    def fake_assemble_psd(**kwargs: object) -> str:
        captured.update(kwargs)
        captured["background_exists"] = Path(kwargs["background_path"]).is_file()
        Path(kwargs["output_path"]).write_bytes(b"psd")
        return str(kwargs["output_path"])

    monkeypatch.setattr(legacy, "assemble_psd", fake_assemble_psd, raising=False)

    outputs = legacy.assemble_legacy_results(store)

    assert outputs == {"page_001": str(output)}
    assert captured["background_exists"] is True
    assert len(captured["components"]) == 1
    assert captured["text_items"][0]["text"] == "editable"
    assert captured["img_width"] == 16
    assert captured["img_height"] == 9
    assert output.read_bytes() == b"psd"
    assert not list(reconstruction.glob("assembly-assets-*"))


def test_psd_assembly_rejects_preserved_warning_page(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    store, reconstruction, _ = _accepted_assembly_job(tmp_path)
    output = tmp_path / "accepted-output.psd"
    manifest = store.read_json("job_manifest.json")
    manifest["output_format"] = "psd"
    manifest["options"]["output_path"] = str(output)
    store.write_json("job_manifest.json", manifest)
    store.write_json(
        "pages/page_001/reconstruction/component_state.json",
        {"status": "preserved_with_warning"},
    )
    monkeypatch.setattr(
        legacy,
        "assemble_psd",
        lambda **kwargs: pytest.fail("PSD assembler must not run"),
        raising=False,
    )

    with pytest.raises(RuntimeError, match="quality gate"):
        legacy.assemble_legacy_results(store)

    assert not output.exists()
    assert not list(reconstruction.glob("assembly-assets-*"))


def test_psd_recovery_tracks_only_the_declared_output(tmp_path: Path) -> None:
    output = (tmp_path / "output.psd").resolve()
    manifest = {
        "output_format": "psd",
        "pages": ["page_001"],
        "input": {"type": "images", "items": []},
        "options": {"output_path": str(output), "slide_size": "both"},
    }

    assert runtime._expected_legacy_output_entries(manifest, "images") == [output]


def test_psd_batch_targets_disambiguate_duplicate_source_names(tmp_path: Path) -> None:
    manifest = {
        "pages": ["page_001", "page_002", "page_003"],
        "input": {
            "type": "images",
            "items": [
                {"original_path": str(tmp_path / "first" / "slide.png")},
                {"original_path": str(tmp_path / "second" / "slide.png")},
                {"original_path": str(tmp_path / "third" / "other.png")},
            ],
        },
    }
    output_dir = tmp_path / "psd"

    assert legacy._legacy_psd_output_targets(manifest, output_dir) == {
        "page_001": output_dir / "001_slide.psd",
        "page_002": output_dir / "002_slide.psd",
        "page_003": output_dir / "other.psd",
    }


def test_suppressed_text_does_not_reappear_in_assembled_pptx(
    tmp_path: Path,
) -> None:
    store, reconstruction, _ = _accepted_assembly_job(tmp_path)
    result_path = reconstruction / "component_result.json"
    result = json.loads(result_path.read_text(encoding="utf-8"))
    result["text_items"] = []
    result_path.write_text(json.dumps(result), encoding="utf-8")
    state = store.read_json(
        "pages/page_001/reconstruction/component_state.json"
    )
    state["result_ref"]["sha256"] = hashlib.sha256(
        result_path.read_bytes()
    ).hexdigest()
    store.write_json(
        "pages/page_001/reconstruction/component_state.json", state
    )

    outputs = legacy.assemble_legacy_results(store)

    reopened = Presentation(outputs["16:9"])
    assert not any(
        getattr(shape, "text", "") == "editable"
        for shape in reopened.slides[0].shapes
    )
    assert any(shape.shape_type == 13 for shape in reopened.slides[0].shapes)


def test_accepted_presentation_assembly_failure_cleans_temporary_assets(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    store, reconstruction, output = _accepted_assembly_job(tmp_path)

    class FailingImageModule:
        @staticmethod
        def load_component_layers(path):
            return {"text_items": [{"text": "editable"}]}

        @staticmethod
        def _assemble_prepared_slide(*args, **kwargs):
            raise RuntimeError("controlled assembly failure")

    monkeypatch.setattr(
        legacy.importlib, "import_module", lambda name: FailingImageModule
    )

    with pytest.raises(RuntimeError, match="controlled assembly failure"):
        legacy.assemble_legacy_results(store)

    assert not output.exists()
    assert not list(reconstruction.glob("assembly-assets-*"))


@pytest.mark.parametrize(
    "mutation", ["missing_manifest", "manifest_content", "component_id", "rgba"],
)
def test_presentation_tamper_is_rejected(tmp_path: Path, mutation: str) -> None:
    store, reconstruction, result, manifest_path, _ = (
        _accepted_presentation_case(tmp_path)
    )
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    if mutation == "missing_manifest":
        result["accepted_asset_refs"].pop("presentation_manifest")
    elif mutation == "manifest_content":
        manifest["components"][0]["metrics"]["boundary_color_mae"] = 1.0
        manifest_path.write_text(json.dumps(manifest), encoding="utf-8")
    elif mutation == "component_id":
        manifest["components"][0]["component_id"] = "replacement"
        manifest_path.write_text(json.dumps(manifest), encoding="utf-8")
        result["accepted_asset_refs"]["presentation_manifest"]["sha256"] = (
            hashlib.sha256(manifest_path.read_bytes()).hexdigest()
        )
    else:
        rgba_ref = manifest["components"][0]["rgba"]
        rgba_path = store.root / Path(*PurePosixPath(rgba_ref["path"]).parts)
        Image.new("RGBA", (4, 4), (0, 255, 0, 255)).save(rgba_path)

    with pytest.raises((ValueError, RuntimeError), match="presentation"):
        legacy._accepted_slide_data(
            store, reconstruction, {"text_items": []}, result
        )
    assert not list(reconstruction.glob("assembly-assets-*"))


def test_presentation_tamper_between_ref_check_and_loader_is_rejected(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    store, reconstruction, result, manifest_path, _ = (
        _accepted_presentation_case(tmp_path)
    )
    real_loader = legacy._load_presentation_assets

    def replace_then_load(**kwargs):
        manifest_path.write_bytes(manifest_path.read_bytes() + b"replacement")
        return real_loader(**kwargs)

    monkeypatch.setattr(legacy, "_load_presentation_assets", replace_then_load)

    with pytest.raises((ValueError, RuntimeError), match="presentation"):
        legacy._accepted_slide_data(
            store, reconstruction, {"text_items": []}, result
        )
    assert not list(reconstruction.glob("assembly-assets-*"))


def test_unused_accepted_asset_reference_structure_is_still_validated(
    tmp_path: Path,
) -> None:
    store, reconstruction, result, _, _ = _accepted_presentation_case(tmp_path)
    result["accepted_asset_refs"]["reconstructed"]["path"] = "../outside.png"

    with pytest.raises(ValueError, match="reference"):
        legacy._accepted_slide_data(
            store, reconstruction, {"text_items": []}, result
        )


@pytest.mark.parametrize("component_id", ["../escape", "..\\escape"])
def test_accepted_presentation_component_id_cannot_escape_asset_directory(
    tmp_path: Path, component_id: str,
) -> None:
    store, reconstruction, result, manifest_path, _ = (
        _accepted_presentation_case(tmp_path)
    )
    graph_path = store.root / result["graph_ref"]["path"]
    graph = json.loads(graph_path.read_text(encoding="utf-8"))
    graph["nodes"][0]["id"] = component_id
    graph_path.write_text(json.dumps(graph), encoding="utf-8")
    graph_sha256 = hashlib.sha256(graph_path.read_bytes()).hexdigest()
    result["graph_ref"]["sha256"] = graph_sha256
    result["accepted_graph_sha256"] = graph_sha256
    result["final_component_ids"] = [component_id]
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    manifest["graph_sha256"] = graph_sha256
    manifest["components"][0]["component_id"] = component_id
    manifest_path.write_text(json.dumps(manifest), encoding="utf-8")
    result["accepted_asset_refs"]["presentation_manifest"]["sha256"] = (
        hashlib.sha256(manifest_path.read_bytes()).hexdigest()
    )

    slide = legacy._accepted_slide_data(
        store, reconstruction, {"text_items": []}, result
    )

    component_path = Path(slide["components"][0]["path"]).resolve()
    assert component_path.parent.name.startswith("assembly-assets-")
    assert component_path.is_relative_to(reconstruction.resolve())


def test_warning_image_page_assembly_refuses_fake_editable_output(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", slide_size="16:9"
    )
    store = RunStore.open(run_dir)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(exist_ok=True)
    store.write_json(
        "pages/page_001/reconstruction/component_state.json",
        {"status": "preserved_with_warning"},
    )
    class FakeImageModule:
        @staticmethod
        def load_component_layers(path):
            return {
                "img_width": 12, "img_height": 8,
                "text_items": [{"text": "must not survive"}],
                "components": [{"path": "must-not-survive"}],
            }

        @staticmethod
        def _assemble_prepared_slide(slide_data, output_path, *args):
            raise AssertionError("warning image must not enter the assembler")

    monkeypatch.setattr(
        legacy.importlib, "import_module", lambda name: FakeImageModule
    )

    with pytest.raises(RuntimeError, match="editable reconstruction incomplete"):
        legacy.assemble_legacy_results(store)

    assert not (run_dir / "final/output.pptx").exists()


def test_legacy_assembly_refuses_existing_output_before_assembler_call(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    output = tmp_path / "existing.pptx"
    _image(source)
    output.write_bytes(b"owner data")
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", output_path=output,
        slide_size="16:9",
    )
    store = RunStore.open(run_dir)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(exist_ok=True)
    store.write_json(
        "pages/page_001/reconstruction/component_state.json",
        {"status": "preserved_with_warning"},
    )

    class FakeImageModule:
        @staticmethod
        def load_component_layers(path):
            return {"img_width": 12, "img_height": 8}

        @staticmethod
        def _assemble_prepared_slide(*args, **kwargs):
            raise AssertionError("assembler must not run")

    monkeypatch.setattr(
        legacy.importlib, "import_module", lambda name: FakeImageModule
    )

    with pytest.raises(RuntimeError, match="Refusing to overwrite"):
        legacy.assemble_legacy_results(store)
    assert output.read_bytes() == b"owner data"
    assert not list(reconstruction.glob("assembly-assets-*"))


def test_warning_image_with_multiple_variants_publishes_no_outputs(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run", slide_size="both")
    store = RunStore.open(run_dir)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(exist_ok=True)
    store.write_json(
        "pages/page_001/reconstruction/component_state.json",
        {"status": "preserved_with_warning"},
    )
    calls = 0

    class FakeImageModule:
        @staticmethod
        def load_component_layers(path):
            return {"img_width": 12, "img_height": 8}

        @staticmethod
        def _assemble_prepared_slide(slide_data, output_path, *args):
            nonlocal calls
            calls += 1
            if calls == 2:
                raise RuntimeError("second variant failed")
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(output_path)

    monkeypatch.setattr(
        legacy.importlib, "import_module", lambda name: FakeImageModule
    )

    with pytest.raises(RuntimeError, match="editable reconstruction incomplete"):
        legacy.assemble_legacy_results(store)
    assert calls == 0
    assert not (run_dir / "final/output_original.pptx").exists()
    assert not (run_dir / "final/output_16x9.pptx").exists()


def test_host_agent_next_times_out_explicitly_while_execution_lease_is_held(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_run(RunStatus.AWAITING_AGENT)

    ticks = iter([0.0, 31.0])
    monkeypatch.setattr(host_agent.time, "monotonic", lambda: next(ticks))
    monkeypatch.setattr(host_agent.time, "sleep", lambda seconds: None)
    with ExecutionLease(run_dir / "execution.lock", run_root=run_dir):
        with pytest.raises(RuntimeError, match="Timed out waiting for Run execution lock"):
            runtime.next_host_agent_item(run_dir)


def test_parent_fallback_reuses_previous_background_and_frozen_assets(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch,
) -> None:
    run_dir = tmp_path / "run"
    page_dir = run_dir / "pages/page_001"
    graph_dir = page_dir / "reconstruction/execution-05"
    (graph_dir / "masks").mkdir(parents=True)
    source = page_dir / "source.png"
    _image(source)
    store = RunStore(run_dir)

    def reference(path: Path) -> dict[str, str]:
        return {
            "path": path.relative_to(run_dir).as_posix(),
            "sha256": hashlib.sha256(path.read_bytes()).hexdigest(),
        }

    store.write_json("pages/page_001/page_request.json", {
        "schema_version": 1,
        "source": "pages/page_001/source.png",
        "sha256": hashlib.sha256(source.read_bytes()).hexdigest(),
    })
    mask = graph_dir / "masks/parent_0001.png"
    Image.new("L", (12, 8), 255).save(mask)
    graph_path = graph_dir / "component-graph.json"
    graph_path.write_text(json.dumps({"nodes": [{
        "id": "parent_0001", "mask": "masks/parent_0001.png",
        "mask_sha256": hashlib.sha256(mask.read_bytes()).hexdigest(),
        "bbox": [0, 0, 12, 8], "state": "pending_gate",
    }, {
        "id": "component_0001", "mask": "masks/parent_0001.png",
        "mask_sha256": hashlib.sha256(mask.read_bytes()).hexdigest(),
        "bbox": [0, 0, 12, 8], "state": "pending",
    }]}), encoding="utf-8")
    background = graph_dir / "background.png"
    _image(background, (4, 5, 6))
    manifest = graph_dir / "presentation-manifest.json"
    manifest.write_text("{}", encoding="utf-8")
    quality = graph_dir / "quality-report.json"
    quality.write_text(json.dumps({"input_refs": {
        "background": reference(background),
        "presentation_manifest": reference(manifest),
    }}), encoding="utf-8")
    store.write_json(
        "pages/page_001/reconstruction/component_state.json",
        {
            "repair_round": 5,
            "graph_ref": reference(graph_path),
            "current_round": {"quality_ref": None},
            "fallback": {"parent_ids": ["parent_0001"]},
            "failed_ids": ["component_0001", "parent_0001"],
            "parent_assets": {"parent_0001": reference(mask)},
            "frozen": {"component_0004": "frozen-hash"},
        },
    )

    def execute(pixels, graph, actions, **kwargs):
        output_dir = kwargs["output_dir"]
        (output_dir / "masks").mkdir(parents=True)
        shutil.copy2(mask, output_dir / "masks/parent_0001.png")
        return copy.deepcopy(graph)

    captured = {}
    captured_actions = []
    def capture_execute(pixels, graph, actions, **kwargs):
        captured_actions.extend(actions)
        return execute(pixels, graph, actions, **kwargs)
    monkeypatch.setattr(legacy, "execute_component_action_round", capture_execute)
    monkeypatch.setattr(legacy, "_ensure_component_disk_reserve", lambda *args, **kwargs: None)
    monkeypatch.setattr(
        legacy, "_quality_assets",
        lambda *args, **kwargs: captured.update(kwargs) or {},
    )
    monkeypatch.setattr(legacy, "record_parent_fallback_execution", lambda *args, **kwargs: None)

    legacy._execute_legacy_parent_fallback(store, "page_001", object())

    assert captured["background_path_override"] == background
    assert captured["frozen_manifest_path"] == manifest
    assert captured["frozen_component_ids"] == {"component_0004"}
    assert [action["action"] for action in captured_actions] == [
        "discard", "collapse_to_parent",
    ]


def _image(path: Path, color: tuple[int, int, int] = (1, 2, 3)) -> None:
    Image.new("RGB", (12, 8), color).save(path)


def _pptx(path: Path, slide_count: int = 2) -> None:
    presentation = Presentation()
    for _ in range(slide_count):
        presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(path)


def _mock_legacy_completion(
    monkeypatch: pytest.MonkeyPatch, assemble
) -> None:
    def initialize(store: RunStore, page_id: str, **kwargs) -> dict:
        reconstruction = store.root / "pages" / page_id / "reconstruction"
        reconstruction.mkdir(parents=True, exist_ok=True)
        (reconstruction / "component_state.json").write_text(
            "{}", encoding="utf-8"
        )
        return {"status": "initialized", "page_id": page_id}

    monkeypatch.setattr(runtime, "initialize_legacy_page", initialize)
    monkeypatch.setattr(
        runtime, "advance_legacy_page",
        lambda store, page_id, **kwargs: {
            "status": "ready_for_assembly", "page_id": page_id
        },
    )
    monkeypatch.setattr(runtime, "assemble_legacy_results", assemble)


def _run_synchronized(
    run_dir: str,
    barrier: object,
    release: object,
    results: object,
) -> None:
    real_open = runtime.RunStore.open
    first_state_read = True

    def synchronized_open(root: str | Path) -> RunStore:
        nonlocal first_state_read
        store = real_open(root)
        read_json = store.read_json

        def synchronized_read(relative: str | Path) -> dict[str, Any]:
            nonlocal first_state_read
            document = read_json(relative)
            if str(relative) == "run_state.json" and first_state_read:
                first_state_read = False
                barrier.wait(10)
            return document

        store.read_json = synchronized_read
        return store

    def execute(store: RunStore) -> dict[str, str]:
        release.wait(10)
        return {"pptx": str(store.root / "final" / "output.pptx")}

    def initialize(store: RunStore, page_id: str, **kwargs) -> dict:
        reconstruction = store.root / "pages" / page_id / "reconstruction"
        reconstruction.mkdir(parents=True, exist_ok=True)
        (reconstruction / "component_state.json").write_text(
            "{}", encoding="utf-8"
        )
        return {"status": "initialized", "page_id": page_id}

    runtime.RunStore.open = synchronized_open
    runtime.initialize_legacy_page = initialize
    runtime.advance_legacy_page = lambda store, page_id, **kwargs: {
        "status": "ready_for_assembly", "page_id": page_id
    }
    runtime.assemble_legacy_results = execute
    try:
        results.put(("ok", runtime.run_job(run_dir)))
    except Exception as error:
        results.put(("error", type(error).__name__, str(error)))


def test_run_job_prepared_race_has_only_one_executor(tmp_path: Path) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    context = multiprocessing.get_context("spawn")
    barrier = context.Barrier(2)
    release = context.Event()
    results = context.Queue()
    processes = [
        context.Process(
            target=_run_synchronized,
            args=(str(run_dir), barrier, release, results),
        )
        for _ in range(2)
    ]
    for process in processes:
        process.start()
    try:
        outcomes = [results.get(timeout=10)]
        release.set()
        outcomes.append(results.get(timeout=10))
    finally:
        release.set()
        for process in processes:
            process.join(10)
            if process.is_alive():
                process.terminate()
                process.join(10)

    assert sorted(item[0] for item in outcomes) == ["error", "ok"]
    assert "already executing" in next(
        item[2] for item in outcomes if item[0] == "error"
    )
    assert all(process.exitcode == 0 for process in processes)


def test_run_job_writes_execution_metadata_while_lease_is_held(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    captured = {}

    def execute(store: RunStore) -> dict[str, str]:
        captured.update(store.read_json("execution.json"))
        with pytest.raises(RuntimeError, match="already executing"):
            with ExecutionLease(store.root / "execution.lock"):
                pass
        return {"pptx": str(store.root / "final" / "output.pptx")}

    _mock_legacy_completion(monkeypatch, execute)

    runtime.run_job(run_dir)

    assert captured["schema_version"] == SCHEMA_VERSION
    assert re.fullmatch(r"[0-9a-f]{32}", captured["token"])
    assert captured["pid"] == os.getpid()
    assert datetime.fromisoformat(captured["started_at"].replace("Z", "+00:00")).tzinfo
    assert captured["input_type"] == "images"


def test_recover_orphaned_image_run_resets_pages_and_cleans_owned_dirs(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    for name in ("work", "final"):
        directory = run_dir / name
        directory.mkdir()
        (directory / "partial.bin").write_bytes(b"partial")

    status = runtime.recover_job(run_dir)

    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "pending"
    assert not (run_dir / "work").exists()
    assert not (run_dir / "final").exists()


def test_recover_finalizing_image_run_resets_validated_page(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    store.transition_page("page_001", PageStatus.VALIDATED)
    store.transition_run(RunStatus.FINALIZING)

    status = runtime.recover_job(run_dir)

    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "pending"


def test_recover_orphaned_pptx_run_keeps_pages_analyzed(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)

    status = runtime.recover_job(run_dir)

    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "analyzed"


def test_recover_rejects_active_execution_lease(tmp_path: Path) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)

    with ExecutionLease(run_dir / "execution.lock"):
        with pytest.raises(RuntimeError, match="already executing"):
            runtime.recover_job(run_dir)


def test_recover_rejects_existing_external_output_without_state_change(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    output = tmp_path / "external.pptx"
    _image(source)
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        output_path=output,
    )
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    output.write_bytes(b"user")
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="external output"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert output.read_bytes() == b"user"


def test_recover_treats_linked_external_parent_as_external(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    final = run_dir / "final"
    final.mkdir()
    external_parent = tmp_path / "external"
    try:
        external_parent.symlink_to(final, target_is_directory=True)
    except OSError as error:
        if os.name != "nt":
            pytest.skip(f"directory links are unavailable: {error}")
        junction = subprocess.run(
            [
                "cmd",
                "/c",
                "mklink",
                "/J",
                str(external_parent),
                str(final),
            ],
            capture_output=True,
            text=True,
            check=False,
        )
        if junction.returncode:
            pytest.skip(f"directory links are unavailable: {error}")
    output = external_parent / "output.pptx"
    output.write_bytes(b"user")
    manifest = store.read_json("job_manifest.json")
    manifest["options"]["output_path"] = str(output.absolute())
    store.write_json("job_manifest.json", manifest)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="external output"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert output.read_bytes() == b"user"


def test_recover_rejects_existing_external_single_image_variant(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    output = tmp_path / "external.pptx"
    _image(source)
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        output_path=output,
        slide_size="both",
    )
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    variant = tmp_path / "external_original.pptx"
    variant.write_bytes(b"user")
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="external output"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert variant.read_bytes() == b"user"


def test_recover_rejects_existing_external_batch_variant_directory(
    tmp_path: Path,
) -> None:
    sources = [tmp_path / "first.png", tmp_path / "second.png"]
    for source in sources:
        _image(source)
    output = tmp_path / "external.pptx"
    run_dir = runtime.prepare_job(
        sources,
        run_dir=tmp_path / "run",
        output_path=output,
        slide_size="both",
    )
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    for page_id in ("page_001", "page_002"):
        store.transition_page(page_id, PageStatus.PROCESSING)
    variant_dir = tmp_path / "external_original"
    variant_dir.mkdir()
    sentinel = variant_dir / "first_original.pptx"
    sentinel.write_bytes(b"user")
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="external output"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert sentinel.read_bytes() == b"user"


def test_recover_rejects_existing_external_original_batch_directory(
    tmp_path: Path,
) -> None:
    sources = [tmp_path / "first.png", tmp_path / "second.png"]
    for source in sources:
        _image(source)
    output = tmp_path / "external.pptx"
    run_dir = runtime.prepare_job(
        sources,
        run_dir=tmp_path / "run",
        output_path=output,
        slide_size="original",
    )
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    for page_id in ("page_001", "page_002"):
        store.transition_page(page_id, PageStatus.PROCESSING)
    variant_dir = tmp_path / "external_original"
    variant_dir.mkdir()
    sentinel = variant_dir / "first_original.pptx"
    sentinel.write_bytes(b"user")
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="external output"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert sentinel.read_bytes() == b"user"


def test_recover_rejects_existing_pptx_output_without_state_change(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    output = run_dir / "final" / "output.pptx"
    output.parent.mkdir()
    output.write_bytes(b"unknown owner")
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="PPTX.*output"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert output.read_bytes() == b"unknown owner"


@pytest.mark.parametrize("status", [RunStatus.PREPARED, RunStatus.FAILED])
def test_recover_rejects_non_orphan_status_without_state_change(
    tmp_path: Path, status: RunStatus
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / status.value)
    store = RunStore.open(run_dir)
    if status is RunStatus.FAILED:
        store.transition_run(status)
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="running or finalizing"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before


def test_recover_rejects_linked_cleanup_path_without_state_change(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    external = tmp_path / "external"
    external.mkdir()
    sentinel = external / "sentinel.txt"
    sentinel.write_text("keep", encoding="utf-8")
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    try:
        (run_dir / "work").symlink_to(external, target_is_directory=True)
    except OSError as error:
        pytest.skip(f"Cannot create symlink: {error}")
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="work"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert sentinel.read_text(encoding="utf-8") == "keep"


def test_recover_rejects_non_directory_cleanup_path_without_state_change(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    work = run_dir / "work"
    work.write_bytes(b"keep")
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="work"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert work.read_bytes() == b"keep"


def test_recover_rejects_pptx_preserved_page_before_cleanup(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PRESERVED)
    work = run_dir / "work"
    work.mkdir()
    sentinel = work / "sentinel"
    sentinel.write_bytes(b"keep")
    before = runtime.get_status(run_dir)

    with pytest.raises(RuntimeError, match="PPTX.*blocked"):
        runtime.recover_job(run_dir)

    assert runtime.get_status(run_dir) == before
    assert sentinel.read_bytes() == b"keep"


def test_recover_job_is_exported() -> None:
    import image2editable

    assert image2editable.recover_job is runtime.recover_job


def test_run_job_does_not_acquire_lease_for_non_prepared_run(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)

    with ExecutionLease(run_dir / "execution.lock"):
        with pytest.raises(RuntimeError, match="current status is running"):
            runtime.run_job(run_dir)


def test_run_job_preserves_pptx_without_calling_legacy(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source)
    source_bytes = source.read_bytes()
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")

    def unexpected_legacy(store: RunStore) -> dict[str, Any]:
        raise AssertionError("PPTX run entered legacy execution")

    monkeypatch.setattr(runtime, "execute_legacy", unexpected_legacy)

    summary = runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    output = Path(summary["outputs"]["pptx"])

    assert summary["status"] == "completed"
    assert summary["pages"] == 2
    assert summary["resource_policy"] == safe_default_policy()
    assert "_output_identity" not in summary
    assert output.read_bytes() == source_bytes
    assert store.read_json("run_summary.json") == summary
    assert store.read_json("run_state.json")["status"] == "completed"
    assert {
        page["status"]
        for page in store.read_json("page_jobs.json")["pages"].values()
    } == {"preserved"}


def test_host_pptx_completion_never_reads_local_model_state(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    monkeypatch.setattr(
        runtime,
        "_local_model_summary",
        lambda store: (_ for _ in ()).throw(
            AssertionError("Host read Local model state")
        ),
    )

    assert runtime.run_job(run_dir)["status"] == "completed"


def test_local_pptx_completed_summary_with_model_is_idempotent(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        agent_provider="local",
    )
    runtime._bind_local_model_receipt(
        RunStore.open(run_dir),
        _local_receipt(tmp_path),
    )
    completed = runtime.run_job(run_dir)

    def unexpected_execute(store: RunStore) -> dict[str, object]:
        raise AssertionError("completed Local PPTX executed again")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    assert completed["agent_model"]["provider"] == "local"
    assert runtime.run_job(run_dir) == completed


def test_run_job_executes_agent_approved_shadow_plan(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    image = tmp_path / "slide.png"
    _image(image)
    source = tmp_path / "source.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image),
        0,
        0,
        presentation.slide_width,
        presentation.slide_height,
    )
    presentation.save(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    candidate = runtime.next_candidate(run_dir)["candidate"]
    runtime.record_decision(
        run_dir,
        page_id="page_001",
        object_id=candidate["source_shape_id"],
        decision="replace",
        confidence=0.99,
        category="full_slide_screenshot",
        evidence=["complete slide layout"],
    )
    _install_component_e2e_boundaries(monkeypatch)
    calls = []

    def fake_execute(store: RunStore, plans) -> dict[str, Any]:
        calls.append(plans)
        output = store.root / "final/output.pptx"
        output.parent.mkdir(parents=True, exist_ok=True)
        input_path = store.root / "input/original.pptx"
        output.write_bytes(input_path.read_bytes())
        digest = runtime.sha256_file(output)
        status = output.lstat()
        return {
            "schema_version": SCHEMA_VERSION,
            "status": "completed",
            "pages": 1,
            "preserved_objects": 0,
            "pending_candidates": 0,
            "replaced_pages": 1,
            "preserved_with_warning_pages": 0,
            "page_results": [
                {
                    "schema_version": SCHEMA_VERSION,
                    "page_id": "page_001",
                    "status": "replaced",
                }
            ],
            "warnings": [],
            "outputs": {"pptx": str(output)},
            "input_sha256": runtime.sha256_file(input_path),
            "output_sha256": digest,
            "_output_identity": {
                "version": 1,
                "path": str(output),
                "dev": status.st_dev,
                "ino": status.st_ino,
                "mode": status.st_mode,
                "size": status.st_size,
                "mtime_ns": status.st_mtime_ns,
                "sha256": digest,
            },
        }

    monkeypatch.setattr(
        runtime,
        "execute_pptx_shadow",
        fake_execute,
        raising=False,
    )

    summary = runtime.run_job(run_dir)

    # An approved page request first enters the durable component-repair
    # boundary.  Until the Host has supplied an accepted result, shadow
    # execution and final output publication are forbidden.
    assert summary["status"] == "awaiting_agent"
    assert calls == []
    store = RunStore.open(run_dir)
    assert (
        store.read_json("page_jobs.json")["pages"]["page_001"]["status"]
        == "awaiting_agent"
    )
    assert not (run_dir / "final" / "output.pptx").exists()


def test_mixed_pptx_warning_output_is_recovery_not_reconstruction_success(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    image = tmp_path / "slide.png"
    _image(image)
    source = tmp_path / "source.pptx"
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(
        str(image), 0, 0, presentation.slide_width, presentation.slide_height
    )
    presentation.save(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    monkeypatch.setattr(
        runtime,
        "shadow_replacement_plans",
        lambda store, manifest: [{"page_id": "page_001"}],
    )

    def warning_recovery(store: RunStore, plans: list[dict]) -> dict[str, Any]:
        output = store.root / "final/output.pptx"
        output.parent.mkdir(parents=True)
        source_copy = store.root / manifest["input"]["source"]
        output.write_bytes(source_copy.read_bytes())
        digest = runtime.sha256_file(output)
        status = output.lstat()
        return {
            "schema_version": SCHEMA_VERSION,
            "status": "completed",
            "pages": 1,
            "preserved_objects": manifest["input"]["object_count"],
            "pending_candidates": manifest["input"]["candidate_count"],
            "replaced_pages": 0,
            "preserved_with_warning_pages": 1,
            "page_results": [{
                "schema_version": SCHEMA_VERSION,
                "page_id": "page_001",
                "status": "preserved_with_warning",
                "warning": "editable reconstruction incomplete",
                "error_type": "RuntimeError",
            }],
            "warnings": ["editable reconstruction incomplete"],
            "outputs": {"pptx": str(output)},
            "input_sha256": manifest["input"]["sha256"],
            "output_sha256": digest,
            "_output_identity": {
                "version": 1,
                "path": str(output),
                "dev": status.st_dev,
                "ino": status.st_ino,
                "mode": status.st_mode,
                "size": status.st_size,
                "mtime_ns": status.st_mtime_ns,
                "sha256": digest,
            },
        }

    monkeypatch.setattr(runtime, "execute_pptx_shadow", warning_recovery)

    summary = runtime.run_job(run_dir)

    assert summary["replaced_pages"] == 0
    assert summary["preserved_with_warning_pages"] == 1
    assert Path(summary["outputs"]["pptx"]).is_file()
    assert summary["page_results"][0]["status"] == "preserved_with_warning"
    assert (
        store.read_json("page_jobs.json")["pages"]["page_001"]["status"]
        == "preserved_with_warning"
    )


def test_completed_pptx_run_is_idempotent_without_recopy(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    completed = runtime.run_job(run_dir)
    output = Path(completed["outputs"]["pptx"])
    before = output.stat()

    def unexpected_execute(store: RunStore) -> dict[str, object]:
        raise AssertionError("completed PPTX run copied output again")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    assert runtime.run_job(run_dir) == completed
    assert output.stat().st_mtime_ns == before.st_mtime_ns
    assert output.stat().st_size == before.st_size


@pytest.mark.parametrize(
    "damage",
    [
        "native_missing",
        "candidates_missing",
        "native_hash",
        "candidates_hash",
        "metadata",
        "objects",
        "candidates",
        "count",
        "manifest_record",
    ],
)
def test_completed_pptx_run_revalidates_bound_inventories_without_mutating_state(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    damage: str,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    native_relative = "pages/page_001/native_objects.json"
    candidates_relative = "pages/page_001/screenshot_candidates.json"
    native_path = run_dir / native_relative
    candidates_path = run_dir / candidates_relative
    manifest = store.read_json("job_manifest.json")
    record = manifest["input"]["inventories"][0]

    if damage == "native_missing":
        native_path.unlink()
    elif damage == "candidates_missing":
        candidates_path.unlink()
    elif damage == "native_hash":
        native_path.write_bytes(native_path.read_bytes() + b" ")
    elif damage == "candidates_hash":
        candidates_path.write_bytes(candidates_path.read_bytes() + b" ")
    elif damage == "metadata":
        native = store.read_json(native_relative)
        native["slide_part"] = "ppt/slides/slide999.xml"
        store.write_json(native_relative, native)
        record["native_objects_sha256"] = runtime.sha256_file(native_path)
        store.write_json("job_manifest.json", manifest)
    elif damage == "objects":
        native = store.read_json(native_relative)
        native["objects"] = [{"action": "invalid"}]
        store.write_json(native_relative, native)
        record["native_objects_sha256"] = runtime.sha256_file(native_path)
        store.write_json("job_manifest.json", manifest)
    elif damage == "candidates":
        candidates = store.read_json(candidates_relative)
        candidates["candidates"] = [{"action": "candidate"}]
        store.write_json(candidates_relative, candidates)
        record["screenshot_candidates_sha256"] = runtime.sha256_file(
            candidates_path
        )
        store.write_json("job_manifest.json", manifest)
    elif damage == "count":
        native = store.read_json(native_relative)
        native["objects"] = [{"action": "preserve"}]
        store.write_json(native_relative, native)
        record["native_objects_sha256"] = runtime.sha256_file(native_path)
        store.write_json("job_manifest.json", manifest)
    else:
        manifest["input"]["inventories"].pop()
        store.write_json("job_manifest.json", manifest)

    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    before_summary = store.read_json("run_summary.json")

    def unexpected_execute(_store: RunStore) -> dict[str, object]:
        raise AssertionError("completed PPTX run executed again")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises((RuntimeError, ValueError), match="PPTX"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert store.read_json("run_summary.json") == before_summary


def test_pptx_run_revalidates_inventory_before_running_state_transition(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    native_path = run_dir / "pages/page_001/native_objects.json"
    native_path.write_bytes(native_path.read_bytes() + b" ")
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")

    def unexpected_execute(_store: RunStore) -> dict[str, object]:
        raise AssertionError("invalid PPTX inventory executed")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises(RuntimeError, match="PPTX inventory hash"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages


def test_completed_pptx_inventory_validation_uses_runtime_manifest_snapshot(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    native_path = run_dir / "pages/page_001/native_objects.json"
    native_path.write_bytes(native_path.read_bytes() + b" ")
    replacement = copy.deepcopy(manifest)
    replacement["input"]["inventories"][0][
        "native_objects_sha256"
    ] = runtime.sha256_file(native_path)
    store.write_json("job_manifest.json", replacement)
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    before_summary = store.read_json("run_summary.json")

    monkeypatch.setattr(
        runtime,
        "_manifest_input",
        lambda _store: (manifest, "pptx"),
    )

    with pytest.raises(RuntimeError, match="PPTX inventory hash"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert store.read_json("run_summary.json") == before_summary


def test_pptx_execution_reuses_pretransition_manifest_inventory_snapshot(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    native_path = run_dir / "pages/page_001/native_objects.json"
    original_transition_run = RunStore.transition_run
    replaced = False

    def replace_manifest_after_running(
        self: RunStore, target: RunStatus
    ) -> dict[str, Any]:
        nonlocal replaced
        result = original_transition_run(self, target)
        if target is RunStatus.RUNNING and not replaced:
            replaced = True
            native_path.write_bytes(native_path.read_bytes() + b" ")
            replacement = copy.deepcopy(manifest)
            replacement["input"]["inventories"][0][
                "native_objects_sha256"
            ] = runtime.sha256_file(native_path)
            self.write_json("job_manifest.json", replacement)
        return result

    monkeypatch.setattr(RunStore, "transition_run", replace_manifest_after_running)

    with pytest.raises(RuntimeError, match="PPTX inventory hash"):
        runtime.run_job(run_dir)

    assert replaced is True
    assert store.read_json("run_state.json")["status"] == "failed"
    assert not (run_dir / "final/output.pptx").exists()


def test_completed_pptx_rejects_inventory_replaced_after_trusted_read(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    native_path = run_dir / "pages/page_001/native_objects.json"
    replacement = tmp_path / "replacement.json"
    replacement.write_bytes(b'{"tampered":true}')
    original_open = Path.open
    replaced = False

    class ReplaceAfterRead:
        def __init__(self, file):
            self.file = file

        def __enter__(self):
            return self.file.__enter__()

        def __exit__(self, *args):
            nonlocal replaced
            result = self.file.__exit__(*args)
            runtime.os.replace(replacement, native_path)
            replaced = True
            return result

    def replace_inventory(path: Path, *args, **kwargs):
        file = original_open(path, *args, **kwargs)
        if path == native_path and args == ("rb",):
            return ReplaceAfterRead(file)
        return file

    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    before_summary = store.read_json("run_summary.json")
    monkeypatch.setattr(Path, "open", replace_inventory)

    with pytest.raises(RuntimeError, match="changed during verification"):
        runtime.run_job(run_dir)

    assert replaced is True
    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert store.read_json("run_summary.json") == before_summary


@pytest.mark.parametrize("damage", ["missing", "bad_bytes", "directory"])
def test_completed_pptx_run_rejects_invalid_output_entry_without_mutating_state(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    damage: str,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    completed = runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    output = Path(completed["outputs"]["pptx"])
    output.unlink()
    if damage == "bad_bytes":
        output.write_bytes(b"corrupt")
    elif damage == "directory":
        output.mkdir()
    before_run = store.read_json("run_state.json")
    before_summary = store.read_json("run_summary.json")

    def unexpected_execute(_store: RunStore) -> dict[str, object]:
        raise AssertionError("completed PPTX run executed again")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises(RuntimeError, match="PPTX completed output"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("run_summary.json") == before_summary


def test_completed_pptx_run_rejects_output_symlink_without_mutating_state(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    completed = runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    output = Path(completed["outputs"]["pptx"])
    output.unlink()
    try:
        output.symlink_to(source)
    except OSError as error:
        pytest.skip(f"symlinks unavailable: {error}")
    before_run = store.read_json("run_state.json")
    before_summary = store.read_json("run_summary.json")

    def unexpected_execute(_store: RunStore) -> dict[str, object]:
        raise AssertionError("completed PPTX run executed again")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises(RuntimeError, match="PPTX completed output"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("run_summary.json") == before_summary


def test_completed_pptx_run_rejects_output_replaced_during_hash(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    completed = runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    output = Path(completed["outputs"]["pptx"])
    replacement = tmp_path / "replacement.pptx"
    replacement.write_bytes(output.read_bytes())
    before_run = store.read_json("run_state.json")
    before_summary = store.read_json("run_summary.json")
    original_sha256_file = runtime.sha256_file

    def replace_after_hash(path: Path) -> str:
        digest = original_sha256_file(path)
        if Path(path) == output:
            runtime.os.replace(replacement, output)
        return digest

    def unexpected_execute(_store: RunStore) -> dict[str, object]:
        raise AssertionError("completed PPTX run executed again")

    monkeypatch.setattr(runtime, "sha256_file", replace_after_hash)
    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises(RuntimeError, match="changed during verification"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("run_summary.json") == before_summary


def test_pptx_run_rejects_legacy_unbound_inventories_before_execution(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    manifest["input"].pop("inventories", None)
    store.write_json("job_manifest.json", manifest)
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")

    def unexpected_execute(_store: RunStore) -> dict[str, object]:
        raise AssertionError("unbound PPTX run executed")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises(RuntimeError, match="inventor"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages


@pytest.mark.parametrize("invalid_slide_count", [True, 1.0, -1, 2])
def test_pptx_manifest_slide_count_is_validated_before_execution(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    invalid_slide_count: object,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    manifest["input"]["slide_count"] = invalid_slide_count
    store.write_json("job_manifest.json", manifest)
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    called = False

    def malicious_execute(_store: RunStore) -> dict[str, object]:
        nonlocal called
        called = True
        return {
            "schema_version": SCHEMA_VERSION,
            "status": "completed",
            "pages": invalid_slide_count,
        }

    monkeypatch.setattr(runtime, "execute_pptx_preserve", malicious_execute)

    with pytest.raises(RuntimeError, match="slide_count"):
        runtime.run_job(run_dir)

    assert called is False
    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("schema_version", True),
        ("status", "failed"),
        ("pages", True),
        ("preserved_objects", 1),
        ("pending_candidates", 1),
        ("warnings", ["unexpected"]),
        ("outputs", {"pptx": "wrong"}),
        ("input_sha256", "0" * 64),
        ("output_sha256", "0" * 64),
        (
            "resource_policy",
            {**safe_default_policy(), "heavy_page_concurrency": True},
        ),
        ("_output_identity", {}),
        ("unknown_public", True),
    ],
)
def test_completed_pptx_summary_is_revalidated_against_manifest(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    field: str,
    value: object,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    completed = runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    summary = dict(completed)
    summary[field] = value
    store.write_json("run_summary.json", summary)
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")

    def unexpected_execute(_store: RunStore) -> dict[str, object]:
        raise AssertionError("completed PPTX run executed again")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises(RuntimeError, match="PPTX execution summary"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert store.read_json("run_summary.json") == summary


def test_pptx_execution_failure_records_analyzed_pages_and_run(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")

    def fail_execute(store: RunStore) -> dict[str, object]:
        raise RuntimeError("preserve failed")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", fail_execute)

    with pytest.raises(RuntimeError, match="preserve failed"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    assert {
        page["status"]
        for page in store.read_json("page_jobs.json")["pages"].values()
    } == {"failed"}
    assert store.read_json("run_summary.json") == {
        "schema_version": SCHEMA_VERSION,
        "status": "failed",
        "error": {"type": "RuntimeError", "message": "preserve failed"},
        "outputs": {},
    }


def test_pptx_run_never_overwrites_preexisting_output(tmp_path: Path) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "output.pptx"
    _pptx(source, slide_count=1)
    output.write_bytes(b"existing")
    run_dir = runtime.prepare_job(
        source,
        run_dir=tmp_path / "run",
        output_path=output,
    )

    with pytest.raises(FileExistsError, match="already exists"):
        runtime.run_job(run_dir)

    assert output.read_bytes() == b"existing"
    store = RunStore.open(run_dir)
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    with pytest.raises(RuntimeError, match="blocked"):
        runtime.retry_page(run_dir, "page_001")
    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert output.read_bytes() == b"existing"


@pytest.mark.parametrize(
    "failure_point",
    ["finalizing", "summary", "completed"],
)
def test_pptx_post_publish_failure_compensates_and_retries(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    failure_point: str,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    source_bytes = source.read_bytes()
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_transition_run = RunStore.transition_run
    original_write_json = RunStore.write_json
    injected = False

    def fail_transition(
        self: RunStore, target: RunStatus
    ) -> dict[str, Any]:
        nonlocal injected
        should_fail = (
            failure_point == "finalizing" and target is RunStatus.FINALIZING
        ) or (
            failure_point == "completed" and target is RunStatus.COMPLETED
        )
        if should_fail and not injected:
            injected = True
            raise OSError(f"{failure_point} state write failed")
        return original_transition_run(self, target)

    def fail_summary_write(
        self: RunStore, relative: str | Path, document: dict[str, Any]
    ) -> None:
        nonlocal injected
        if (
            failure_point == "summary"
            and Path(relative) == Path("run_summary.json")
            and document.get("status") == "completed"
            and not injected
        ):
            injected = True
            raise OSError("summary write failed")
        original_write_json(self, relative, document)

    monkeypatch.setattr(RunStore, "transition_run", fail_transition)
    monkeypatch.setattr(RunStore, "write_json", fail_summary_write)

    with pytest.raises(OSError, match="write failed"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert not output.exists()
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "failed"
    assert store.read_json("run_summary.json")["status"] == "failed"

    retried = runtime.retry_page(run_dir, "page_001")
    assert retried["run"]["status"] == "prepared"
    assert retried["pages"]["pages"]["page_001"]["status"] == "analyzed"
    assert runtime.run_job(run_dir)["status"] == "completed"
    assert output.read_bytes() == source_bytes


def test_pptx_compensation_does_not_delete_concurrent_replacement(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_transition_run = RunStore.transition_run

    def replace_then_fail(
        self: RunStore, target: RunStatus
    ) -> dict[str, Any]:
        if target is RunStatus.FINALIZING:
            output.unlink()
            output.write_bytes(b"concurrent replacement")
            raise OSError("finalizing state write failed")
        return original_transition_run(self, target)

    monkeypatch.setattr(RunStore, "transition_run", replace_then_fail)

    with pytest.raises(OSError, match="finalizing") as error:
        runtime.run_job(run_dir)

    assert error.value.__cause__ is not None
    assert "safely" in str(error.value.__cause__)
    assert output.read_bytes() == b"concurrent replacement"
    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "preserved"
    with pytest.raises(RuntimeError, match="blocked"):
        runtime.retry_page(run_dir, "page_001")
    assert store.read_json("run_state.json")["status"] == "failed"
    assert output.read_bytes() == b"concurrent replacement"


def test_pptx_compensation_isolates_path_before_deleting(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_transition_run = RunStore.transition_run
    original_replace = runtime.os.replace
    replacement_injected = False

    def replace_before_isolation(source_path, destination_path):
        nonlocal replacement_injected
        if Path(source_path) == output and not replacement_injected:
            replacement_injected = True
            output.unlink()
            output.write_bytes(b"last-moment replacement")
        return original_replace(source_path, destination_path)

    def fail_finalizing(
        self: RunStore, target: RunStatus
    ) -> dict[str, Any]:
        if target is RunStatus.FINALIZING:
            raise OSError("finalizing state write failed")
        return original_transition_run(self, target)

    monkeypatch.setattr(runtime.os, "replace", replace_before_isolation)
    monkeypatch.setattr(RunStore, "transition_run", fail_finalizing)

    with pytest.raises(OSError, match="finalizing"):
        runtime.run_job(run_dir)

    assert replacement_injected is True
    assert output.read_bytes() == b"last-moment replacement"
    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    with pytest.raises(RuntimeError, match="blocked"):
        runtime.retry_page(run_dir, "page_001")
    assert store.read_json("run_state.json")["status"] == "failed"


def test_pptx_forged_summary_hash_cleans_token_owned_output(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_execute = runtime.execute_pptx_preserve

    def return_wrong_output_hash(store: RunStore) -> dict[str, object]:
        summary = original_execute(store)
        summary["output_sha256"] = "0" * 64
        return summary

    monkeypatch.setattr(
        runtime, "execute_pptx_preserve", return_wrong_output_hash
    )

    with pytest.raises(RuntimeError, match="hash does not match"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert not output.exists()
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "failed"
    assert store.read_json("run_summary.json")["status"] == "failed"
    assert "retry_blocked" not in store.read_json("run_summary.json")
    monkeypatch.setattr(runtime, "execute_pptx_preserve", original_execute)
    assert runtime.retry_page(run_dir, "page_001")["run"]["status"] == "prepared"
    assert runtime.run_job(run_dir)["status"] == "completed"


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("schema_version", True),
        ("status", "failed"),
        ("pages", True),
        ("preserved_objects", True),
        ("pending_candidates", 0.0),
        ("warnings", ["unexpected"]),
        ("outputs", {"pptx": "wrong"}),
        ("input_sha256", "0" * 64),
    ],
)
def test_pptx_invalid_summary_is_cleaned_and_retryable_when_token_owned(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    field: str,
    value: object,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_execute = runtime.execute_pptx_preserve

    def forge_summary(store: RunStore) -> dict[str, object]:
        summary = original_execute(store)
        summary[field] = value
        return summary

    monkeypatch.setattr(runtime, "execute_pptx_preserve", forge_summary)

    with pytest.raises(RuntimeError, match="PPTX execution summary"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert not output.exists()
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "failed"
    assert "retry_blocked" not in store.read_json("run_summary.json")
    monkeypatch.setattr(runtime, "execute_pptx_preserve", original_execute)
    assert runtime.retry_page(run_dir, "page_001")["run"]["status"] == "prepared"
    assert runtime.run_job(run_dir)["status"] == "completed"


def test_pptx_self_consistent_wrong_bytes_are_rejected_and_cleaned(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    from image2editable.pptx_input import _publish_pptx_no_clobber

    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_execute = runtime.execute_pptx_preserve

    def publish_wrong_bytes(store: RunStore) -> dict[str, object]:
        manifest = store.read_json("job_manifest.json")
        temporary = output.parent / ".malicious.tmp"
        output.parent.mkdir(parents=True, exist_ok=True)
        temporary.write_bytes(b"not-a-pptx")
        token = _publish_pptx_no_clobber(temporary, output)
        temporary.unlink()
        digest = runtime.sha256_file(output)
        return {
            "schema_version": SCHEMA_VERSION,
            "status": "completed",
            "pages": 1,
            "preserved_objects": manifest["input"]["object_count"],
            "pending_candidates": manifest["input"]["candidate_count"],
            "warnings": [],
            "outputs": {"pptx": str(output)},
            "input_sha256": digest,
            "output_sha256": digest,
            "_output_identity": token,
        }

    monkeypatch.setattr(runtime, "execute_pptx_preserve", publish_wrong_bytes)

    with pytest.raises(RuntimeError, match="manifest"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert not output.exists()
    assert "retry_blocked" not in store.read_json("run_summary.json")
    monkeypatch.setattr(runtime, "execute_pptx_preserve", original_execute)
    assert runtime.retry_page(run_dir, "page_001")["run"]["status"] == "prepared"
    assert runtime.run_job(run_dir)["status"] == "completed"


@pytest.mark.parametrize("invalid_sha256", [True, "A" * 64, "0" * 63])
def test_pptx_invalid_manifest_sha256_is_rejected_before_execution(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    invalid_sha256: object,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    manifest["input"]["sha256"] = invalid_sha256
    store.write_json("job_manifest.json", manifest)
    before_state = store.read_json("run_state.json")

    def unexpected_execute(_store: RunStore) -> dict[str, object]:
        raise AssertionError("invalid manifest executed")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises(RuntimeError, match="manifest.*sha256"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_state
    assert not (run_dir / "final" / "output.pptx").exists()


def test_pptx_executor_cannot_claim_another_output_path(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    user_output = tmp_path / "user-owned.pptx"
    _pptx(source, slide_count=1)
    user_output.write_bytes(b"user")
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    status = user_output.lstat()

    def claim_user_output(_store: RunStore) -> dict[str, object]:
        digest = runtime.sha256_file(user_output)
        return {
            "schema_version": SCHEMA_VERSION,
            "status": "completed",
            "pages": 1,
            "outputs": {"pptx": str(user_output)},
            "output_sha256": digest,
            "_output_identity": {
                "version": 1,
                "path": str(user_output),
                "dev": status.st_dev,
                "ino": status.st_ino,
                "mode": status.st_mode,
                "size": status.st_size,
                "mtime_ns": status.st_mtime_ns,
                "sha256": digest,
            },
        }

    monkeypatch.setattr(runtime, "execute_pptx_preserve", claim_user_output)

    with pytest.raises(RuntimeError, match="expected output path"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert user_output.read_bytes() == b"user"
    assert store.read_json("run_summary.json")["retry_blocked"] is True
    with pytest.raises(RuntimeError, match="blocked"):
        runtime.retry_page(run_dir, "page_001")
    assert user_output.read_bytes() == b"user"


def test_pptx_executor_cannot_claim_preexisting_expected_output(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    output = tmp_path / "expected.pptx"
    _pptx(source, slide_count=1)
    output.write_bytes(b"user")
    run_dir = runtime.prepare_job(
        source, run_dir=tmp_path / "run", output_path=output
    )
    status = output.lstat()

    def claim_preexisting(_store: RunStore) -> dict[str, object]:
        digest = runtime.sha256_file(output)
        return {
            "schema_version": SCHEMA_VERSION,
            "status": "completed",
            "pages": 1,
            "outputs": {"pptx": str(output)},
            "output_sha256": digest,
            "_output_identity": {
                "version": 1,
                "path": str(output),
                "dev": status.st_dev,
                "ino": status.st_ino,
                "mode": status.st_mode,
                "size": status.st_size,
                "mtime_ns": status.st_mtime_ns,
                "sha256": digest,
            },
        }

    monkeypatch.setattr(runtime, "execute_pptx_preserve", claim_preexisting)

    with pytest.raises(RuntimeError, match="already existed"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert output.read_bytes() == b"user"
    assert store.read_json("run_summary.json")["retry_blocked"] is True


def test_pptx_same_bytes_replacement_before_return_is_not_owned(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    replacement = tmp_path / "replacement.pptx"
    original_execute = runtime.execute_pptx_preserve

    def replace_before_return(store: RunStore) -> dict[str, object]:
        summary = original_execute(store)
        replacement.write_bytes(output.read_bytes())
        runtime.os.replace(replacement, output)
        return summary

    monkeypatch.setattr(
        runtime, "execute_pptx_preserve", replace_before_return
    )

    with pytest.raises(RuntimeError, match="identity token"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert output.read_bytes() == source.read_bytes()
    assert store.read_json("run_summary.json")["retry_blocked"] is True
    with pytest.raises(RuntimeError, match="blocked"):
        runtime.retry_page(run_dir, "page_001")
    assert output.read_bytes() == source.read_bytes()


@pytest.mark.parametrize(
    "mutation",
    [
        "path",
        "identity",
        "hash",
        "missing",
        "absent",
        "malformed",
        "version_bool",
        "version_float",
        "future",
    ],
)
def test_pptx_forged_or_unknown_identity_token_is_fail_closed(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    mutation: str,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_execute = runtime.execute_pptx_preserve

    def forge_token(store: RunStore) -> dict[str, object]:
        summary = original_execute(store)
        token = dict(summary["_output_identity"])
        if mutation == "path":
            token["path"] = str(tmp_path / "other.pptx")
        elif mutation == "identity":
            token["ino"] += 1
        elif mutation == "hash":
            token["sha256"] = "0" * 64
        elif mutation == "missing":
            token.pop("ino")
        elif mutation == "absent":
            summary.pop("_output_identity")
            return summary
        elif mutation == "malformed":
            summary["_output_identity"] = ["not", "an", "object"]
            return summary
        elif mutation == "version_bool":
            token["version"] = True
        elif mutation == "version_float":
            token["version"] = 1.0
        else:
            token["future"] = True
        summary["_output_identity"] = token
        return summary

    monkeypatch.setattr(runtime, "execute_pptx_preserve", forge_token)

    with pytest.raises(RuntimeError, match="identity token"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert output.is_file()
    assert store.read_json("run_summary.json")["retry_blocked"] is True


def test_pptx_execute_post_publish_error_blocks_retry(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_execute = runtime.execute_pptx_preserve

    def publish_then_fail(store: RunStore) -> dict[str, object]:
        original_execute(store)
        raise OSError("post-publish cleanup failed")

    monkeypatch.setattr(
        runtime, "execute_pptx_preserve", publish_then_fail
    )

    with pytest.raises(OSError, match="post-publish"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert output.is_file()
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "preserved"
    assert store.read_json("run_summary.json")["retry_blocked"] is True
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    with pytest.raises(RuntimeError, match="blocked"):
        runtime.retry_page(run_dir, "page_001")
    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert output.is_file()


def test_pptx_compensation_failure_blocks_retry_before_pages_preserved(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_transition_pages = runtime._transition_pages
    original_replace = runtime.os.replace
    replacement_injected = False

    def fail_preserved_transition(
        store: RunStore, page_ids: list[str], target: Any
    ) -> None:
        if target.value == "preserved":
            raise OSError("preserved page write failed")
        original_transition_pages(store, page_ids, target)

    def replace_before_isolation(source_path, destination_path):
        nonlocal replacement_injected
        if Path(source_path) == output and not replacement_injected:
            replacement_injected = True
            output.unlink()
            output.write_bytes(b"concurrent replacement")
        return original_replace(source_path, destination_path)

    monkeypatch.setattr(
        runtime, "_transition_pages", fail_preserved_transition
    )
    monkeypatch.setattr(runtime.os, "replace", replace_before_isolation)

    with pytest.raises(OSError, match="preserved page") as error:
        runtime.run_job(run_dir)

    assert error.value.__cause__ is not None
    assert output.read_bytes() == b"concurrent replacement"
    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "failed"
    assert store.read_json("run_summary.json")["retry_blocked"] is True
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    with pytest.raises(RuntimeError, match="blocked"):
        runtime.retry_page(run_dir, "page_001")
    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages


def test_pptx_compensation_recovers_one_shot_page_snapshot_write_failure(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_transition_run = RunStore.transition_run
    original_write_json = RunStore.write_json
    finalizing_failed = False
    snapshot_failed = False

    def fail_finalizing_once(
        self: RunStore, target: RunStatus
    ) -> dict[str, Any]:
        nonlocal finalizing_failed
        if target is RunStatus.FINALIZING and not finalizing_failed:
            finalizing_failed = True
            raise OSError("finalizing state write failed")
        return original_transition_run(self, target)

    def fail_snapshot_once(
        self: RunStore, relative: str | Path, document: dict[str, Any]
    ) -> None:
        nonlocal snapshot_failed
        if (
            Path(relative) == Path("page_jobs.json")
            and {
                page["status"] for page in document["pages"].values()
            }
            == {"analyzed"}
            and not snapshot_failed
        ):
            snapshot_failed = True
            raise OSError("page snapshot write failed")
        original_write_json(self, relative, document)

    monkeypatch.setattr(RunStore, "transition_run", fail_finalizing_once)
    monkeypatch.setattr(RunStore, "write_json", fail_snapshot_once)

    with pytest.raises(OSError, match="finalizing") as error:
        runtime.run_job(run_dir)

    assert isinstance(error.value.__cause__, OSError)
    assert str(error.value.__cause__) == "page snapshot write failed"
    store = RunStore.open(run_dir)
    assert not output.exists()
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "failed"
    assert store.read_json("run_summary.json")["status"] == "failed"

    retried = runtime.retry_page(run_dir, "page_001")
    assert retried["run"]["status"] == "prepared"
    assert retried["pages"]["pages"]["page_001"]["status"] == "analyzed"
    assert runtime.run_job(run_dir)["status"] == "completed"


def test_pptx_compensation_recovers_completed_post_write_failure(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    output = run_dir / "final" / "output.pptx"
    original_transition_run = RunStore.transition_run
    completed_failed = False

    def persist_completed_then_fail(
        self: RunStore, target: RunStatus
    ) -> dict[str, Any]:
        nonlocal completed_failed
        result = original_transition_run(self, target)
        if target is RunStatus.COMPLETED and not completed_failed:
            completed_failed = True
            raise OSError("completed state post-write failure")
        return result

    monkeypatch.setattr(
        RunStore, "transition_run", persist_completed_then_fail
    )

    with pytest.raises(OSError, match="post-write"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert not output.exists()
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "failed"
    assert store.read_json("run_summary.json")["status"] == "failed"

    retried = runtime.retry_page(run_dir, "page_001")
    assert retried["run"]["status"] == "prepared"
    assert retried["pages"]["pages"]["page_001"]["status"] == "analyzed"
    assert runtime.run_job(run_dir)["status"] == "completed"


@pytest.mark.parametrize("mutation", ["manifest_pages", "page_jobs", "status"])
def test_pptx_run_rejects_inconsistent_pages_before_execution(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    mutation: str,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source)
    run_dir = prepare_pptx_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    if mutation == "manifest_pages":
        manifest = store.read_json("job_manifest.json")
        manifest["pages"] = ["page_001", "page_003"]
        store.write_json("job_manifest.json", manifest)
    else:
        page_jobs = store.read_json("page_jobs.json")
        if mutation == "page_jobs":
            page_jobs["pages"]["page_003"] = page_jobs["pages"].pop("page_002")
        else:
            page_jobs["pages"]["page_001"]["status"] = "pending"
        store.write_json("page_jobs.json", page_jobs)
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")

    def unexpected_execute(store: RunStore) -> dict[str, object]:
        raise AssertionError("inconsistent PPTX run executed")

    monkeypatch.setattr(runtime, "execute_pptx_preserve", unexpected_execute)

    with pytest.raises(RuntimeError, match="PPTX"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert not (run_dir / "final" / "output.pptx").exists()


def test_pptx_page_validation_accepts_sorted_json_order_for_1000_pages() -> None:
    page_ids = [f"page_{index:03d}" for index in range(1, 1001)]
    page_jobs = {
        "pages": {
            page_id: {"status": "analyzed"}
            for page_id in sorted(page_ids)
        }
    }

    assert runtime._pptx_page_ids(
        {"pages": page_ids},
        page_jobs,
        runtime.PageStatus.ANALYZED,
    ) == page_ids


def test_retry_pptx_run_restores_analyzed_pages(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    original_execute = runtime.execute_pptx_preserve
    calls = 0

    def fail_once(store: RunStore) -> dict[str, object]:
        nonlocal calls
        calls += 1
        if calls == 1:
            raise RuntimeError("preserve failed")
        return original_execute(store)

    monkeypatch.setattr(runtime, "execute_pptx_preserve", fail_once)
    with pytest.raises(RuntimeError, match="preserve failed"):
        runtime.run_job(run_dir)
    work_root = run_dir / "work"
    work_root.mkdir()
    (work_root / "legacy.txt").write_text("stale", encoding="utf-8")

    status = runtime.retry_page(run_dir, "page_001")

    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "analyzed"
    assert not work_root.exists()
    assert runtime.run_job(run_dir)["status"] == "completed"
    assert calls == 2


def test_run_job_rejects_unknown_input_type_before_execution(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    manifest["input"]["type"] = "unknown"
    store.write_json("job_manifest.json", manifest)

    def unexpected_legacy(store: RunStore) -> dict[str, Any]:
        raise AssertionError("unknown input entered legacy execution")

    monkeypatch.setattr(runtime, "execute_legacy", unexpected_legacy)

    with pytest.raises(RuntimeError, match="Unsupported input type"):
        runtime.run_job(run_dir)

    assert store.read_json("run_state.json")["status"] == "prepared"


def test_run_job_completes_and_writes_summary_and_page_results(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    _image(first)
    _image(second, (4, 5, 6))
    run_dir = runtime.prepare_job([first, second], run_dir=tmp_path / "run")
    outputs = {
        "16:9": str((tmp_path / "wide.pptx").resolve()),
        "original": [str((tmp_path / "first.pptx").resolve())],
    }
    _mock_legacy_completion(monkeypatch, lambda store: outputs)

    summary = runtime.run_job(run_dir)
    store = RunStore.open(run_dir)

    assert summary == {
        "schema_version": SCHEMA_VERSION,
        "status": "completed",
        "pages": 2,
        "outputs": outputs,
        "resource_policy": safe_default_policy(),
        "quality_gate_version": runtime.COMPONENT_QUALITY_GATE_VERSION,
    }
    assert store.read_json("run_summary.json") == summary
    assert store.read_json("run_state.json")["status"] == "completed"
    for page_id in ("page_001", "page_002"):
        assert store.read_json(f"pages/{page_id}/page_result.json") == {
            "schema_version": SCHEMA_VERSION,
            "page_id": page_id,
            "status": "validated",
            "outputs": outputs,
        }
        assert (
            store.read_json("page_jobs.json")["pages"][page_id]["status"]
            == "validated"
        )


def test_run_job_records_execution_failure_for_run_and_all_pages(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    _image(first)
    _image(second)
    run_dir = runtime.prepare_job([first, second], run_dir=tmp_path / "run")

    def fail_execute(store: RunStore) -> dict[str, Any]:
        raise RuntimeError("model unavailable")

    _mock_legacy_completion(monkeypatch, fail_execute)

    with pytest.raises(RuntimeError, match="model unavailable"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    assert {
        page["status"]
        for page in store.read_json("page_jobs.json")["pages"].values()
    } == {"failed"}
    assert store.read_json("run_summary.json") == {
        "schema_version": SCHEMA_VERSION,
        "status": "failed",
        "error": {"type": "RuntimeError", "message": "model unavailable"},
        "outputs": {},
    }


def test_page_result_write_failure_records_failed_run_and_pages(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    _image(first)
    _image(second)
    run_dir = runtime.prepare_job([first, second], run_dir=tmp_path / "run")
    _mock_legacy_completion(
        monkeypatch, lambda store: {"16:9": str(tmp_path / "output.pptx")}
    )
    original_write_json = RunStore.write_json

    def fail_second_page_result(
        self: RunStore, relative: str | Path, document: dict[str, Any]
    ) -> None:
        if Path(relative) == Path("pages/page_002/page_result.json"):
            raise OSError("page result write failed")
        original_write_json(self, relative, document)

    monkeypatch.setattr(RunStore, "write_json", fail_second_page_result)

    with pytest.raises(OSError, match="page result write failed"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    assert {
        page["status"]
        for page in store.read_json("page_jobs.json")["pages"].values()
    } == {"failed"}
    assert store.read_json("run_summary.json") == {
        "schema_version": SCHEMA_VERSION,
        "status": "failed",
        "error": {"type": "OSError", "message": "page result write failed"},
        "outputs": {},
    }


def test_completed_transition_failure_can_retry_the_entire_batch(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    _mock_legacy_completion(
        monkeypatch, lambda store: {"16:9": str(tmp_path / "output.pptx")}
    )
    original_transition_run = RunStore.transition_run

    def fail_completed(self: RunStore, target: RunStatus) -> dict[str, Any]:
        if target is RunStatus.COMPLETED:
            raise OSError("completed state write failed")
        return original_transition_run(self, target)

    monkeypatch.setattr(RunStore, "transition_run", fail_completed)

    with pytest.raises(OSError, match="completed state write failed"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    assert (
        store.read_json("page_jobs.json")["pages"]["page_001"]["status"]
        == "failed"
    )
    assert store.read_json("run_summary.json") == {
        "schema_version": SCHEMA_VERSION,
        "status": "failed",
        "error": {"type": "OSError", "message": "completed state write failed"},
        "outputs": {},
    }

    retried = runtime.retry_page(run_dir, "page_001")
    assert retried["run"]["status"] == "prepared"
    assert retried["pages"]["pages"]["page_001"]["status"] == "pending"

    monkeypatch.setattr(RunStore, "transition_run", original_transition_run)
    completed = runtime.run_job(run_dir)

    assert completed["status"] == "completed"
    assert runtime.get_status(run_dir)["run"]["status"] == "completed"


def test_success_validates_pages_with_one_page_jobs_write(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    _image(first)
    _image(second)
    run_dir = runtime.prepare_job([first, second], run_dir=tmp_path / "run")
    _mock_legacy_completion(monkeypatch, lambda store: {})
    original_write_json = RunStore.write_json
    validated_writes = 0

    def count_validated_write(
        self: RunStore, relative: str | Path, document: dict[str, Any]
    ) -> None:
        nonlocal validated_writes
        if (
            Path(relative) == Path("page_jobs.json")
            and "validated"
            in {
                page["status"] for page in document["pages"].values()
            }
        ):
            validated_writes += 1
        original_write_json(self, relative, document)

    monkeypatch.setattr(RunStore, "write_json", count_validated_write)

    runtime.run_job(run_dir)

    assert validated_writes == 2


def test_cleanup_error_is_cause_and_does_not_stop_later_failure_records(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    original_transition_pages = runtime._transition_pages

    def fail_execute(store: RunStore) -> dict[str, Any]:
        raise RuntimeError("execution failed")

    def fail_failed_pages(
        store: RunStore, page_ids: list[str], target: Any
    ) -> None:
        if target.value == "failed":
            raise OSError("page cleanup failed")
        original_transition_pages(store, page_ids, target)

    _mock_legacy_completion(monkeypatch, fail_execute)
    monkeypatch.setattr(runtime, "_transition_pages", fail_failed_pages)

    with pytest.raises(RuntimeError, match="execution failed") as error:
        runtime.run_job(run_dir)

    assert isinstance(error.value.__cause__, OSError)
    assert str(error.value.__cause__) == "page cleanup failed"
    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("run_summary.json")["status"] == "failed"
    assert (
        store.read_json("page_jobs.json")["pages"]["page_001"]["status"]
        == "validated"
    )

    monkeypatch.setattr(runtime, "_transition_pages", original_transition_pages)
    retried = runtime.retry_page(run_dir, "page_001")

    assert retried["run"]["status"] == "prepared"
    assert retried["pages"]["pages"]["page_001"]["status"] == "pending"


def test_retry_recovers_failed_batch_left_running_by_one_run_write_failure(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    original_transition_run = RunStore.transition_run
    failed_transition = False
    execute_calls = 0

    def fail_execute_once(store: RunStore) -> dict[str, Any]:
        nonlocal execute_calls
        execute_calls += 1
        if execute_calls == 1:
            raise RuntimeError("execution failed")
        return {}

    def fail_run_failed_once(
        self: RunStore, target: RunStatus
    ) -> dict[str, Any]:
        nonlocal failed_transition
        if target is RunStatus.FAILED and not failed_transition:
            failed_transition = True
            raise OSError("run failed write failed")
        return original_transition_run(self, target)

    _mock_legacy_completion(monkeypatch, fail_execute_once)
    monkeypatch.setattr(RunStore, "transition_run", fail_run_failed_once)

    with pytest.raises(RuntimeError, match="execution failed") as error:
        runtime.run_job(run_dir)

    assert isinstance(error.value.__cause__, OSError)
    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "running"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "failed"
    assert store.read_json("run_summary.json")["status"] == "failed"

    retried = runtime.retry_page(run_dir, "page_001")
    assert retried["run"]["status"] == "prepared"
    assert retried["pages"]["pages"]["page_001"]["status"] == "pending"

    completed = runtime.run_job(run_dir)
    assert completed["status"] == "completed"
    assert execute_calls == 2


def test_retry_page_resets_the_entire_failed_batch(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    _image(first)
    _image(second)
    run_dir = runtime.prepare_job([first, second], run_dir=tmp_path / "run")

    def fail_execute(store: RunStore) -> dict[str, Any]:
        raise RuntimeError("failed")

    _mock_legacy_completion(monkeypatch, fail_execute)
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)

    status = runtime.retry_page(run_dir, "page_001")

    assert status["run"]["status"] == "prepared"
    assert {
        page["status"] for page in status["pages"]["pages"].values()
    } == {"pending"}


@pytest.mark.parametrize(
    ("analyzed", "expected"),
    [(False, "pending"), (True, "analyzed")],
)
def test_reset_page_jobs_allows_reconstruction_warning_retry(
    analyzed: bool, expected: str,
) -> None:
    page_jobs = {
        "schema_version": 1,
        "pages": {
            "page_001": {
                "schema_version": 1,
                "status": "preserved_with_warning",
                "updated_at": "2026-01-01T00:00:00Z",
            },
        },
    }

    assert runtime._reset_page_jobs(page_jobs, analyzed=analyzed) is True
    assert page_jobs["pages"]["page_001"]["status"] == expected


def test_retry_warning_removes_stale_page_reconstruction(tmp_path: Path) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    store.transition_page("page_001", PageStatus.PRESERVED_WITH_WARNING)
    store.transition_run(RunStatus.FAILED)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(parents=True)
    (reconstruction / "component_state.json").write_text(
        "stale", encoding="utf-8"
    )

    status = runtime.retry_page(run_dir, "page_001")

    assert not reconstruction.exists()
    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "pending"


def test_retry_warning_during_agent_wait_preserves_other_page(
    tmp_path: Path,
) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    _image(first)
    _image(second)
    run_dir = runtime.prepare_job([first, second], run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    store.transition_page("page_001", PageStatus.PRESERVED_WITH_WARNING)
    store.transition_page("page_002", PageStatus.PROCESSING)
    store.transition_page("page_002", PageStatus.AWAITING_AGENT)
    store.transition_run(RunStatus.AWAITING_AGENT)
    first_reconstruction = run_dir / "pages/page_001/reconstruction"
    second_reconstruction = run_dir / "pages/page_002/reconstruction"
    first_reconstruction.mkdir(parents=True)
    second_reconstruction.mkdir(parents=True)
    (first_reconstruction / "stale.txt").write_text("stale", encoding="utf-8")
    (second_reconstruction / "keep.txt").write_text("keep", encoding="utf-8")

    status = runtime.retry_page(run_dir, "page_001")

    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "pending"
    assert status["pages"]["pages"]["page_002"]["status"] == "awaiting_agent"
    assert not first_reconstruction.exists()
    assert (second_reconstruction / "keep.txt").read_text(encoding="utf-8") == "keep"


@pytest.mark.parametrize(
    ("first_status", "ready_after_assembly_warning", "fail_after_cleanup"),
    [
        ("replaced", False, False),
        ("preserved_with_warning", True, False),
        ("replaced", False, True),
    ],
)
def test_retry_completed_pptx_warning_preserves_validated_other_page(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    first_status: str,
    ready_after_assembly_warning: bool,
    fail_after_cleanup: bool,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=2)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    output = run_dir / "final/output.pptx"
    output.parent.mkdir(parents=True)
    output.write_bytes(b"owned output")
    output_hash = runtime.sha256_file(output)
    store.write_json("run_state.json", {
        "schema_version": SCHEMA_VERSION,
        "status": "completed",
        "updated_at": runtime.utc_now(),
    })
    page_jobs = store.read_json("page_jobs.json")
    page_jobs["pages"]["page_001"].update({
        "status": first_status, "updated_at": runtime.utc_now(),
    })
    page_jobs["pages"]["page_002"].update({
        "status": "preserved_with_warning", "updated_at": runtime.utc_now(),
    })
    store.write_json("page_jobs.json", page_jobs)
    store.write_json("run_summary.json", {
        "schema_version": SCHEMA_VERSION,
        "status": "completed",
        "outputs": {"pptx": str(output)},
        "output_sha256": output_hash,
        "page_results": [
            {"page_id": "page_001", "status": first_status},
            {"page_id": "page_002", "status": "preserved_with_warning"},
        ],
        "input_sha256": manifest["input"]["sha256"],
    })
    first_reconstruction = run_dir / "pages/page_001/reconstruction"
    second_reconstruction = run_dir / "pages/page_002/reconstruction"
    first_reconstruction.mkdir(parents=True)
    second_reconstruction.mkdir(parents=True)
    if ready_after_assembly_warning:
        _write_mock_component_state(store, "page_001")
        component_state = store.read_json(
            "pages/page_001/reconstruction/component_state.json"
        )
        component_state.update({
            "phase": "ready_for_assembly",
            "status": "ready_for_assembly",
            "result_ref": {
                "path": "mock-artifact.json",
                "sha256": "0" * 64,
            },
        })
        store.write_json(
            "pages/page_001/reconstruction/component_state.json",
            component_state,
        )
    (first_reconstruction / "keep.txt").write_text("keep", encoding="utf-8")
    (first_reconstruction / "donor.pptx").write_bytes(b"stale donor")
    (second_reconstruction / "stale.txt").write_text("stale", encoding="utf-8")
    first_plan = run_dir / "host-component-plan-page_001-01-first.json"
    second_plan = run_dir / "host-component-plan-page_002-05-stale.json"
    first_plan.write_text("keep", encoding="utf-8")
    second_plan.write_text("stale", encoding="utf-8")

    if fail_after_cleanup:
        original_unlink = Path.unlink
        failed = False

        def fail_isolated_unlink(path: Path, *args, **kwargs) -> None:
            nonlocal failed
            if path.name.startswith(".output.pptx.recovery-") and not failed:
                failed = True
                raise OSError("isolated output cleanup failed")
            original_unlink(path, *args, **kwargs)

        monkeypatch.setattr(Path, "unlink", fail_isolated_unlink)
        with pytest.raises(OSError, match="isolated output cleanup failed"):
            runtime.retry_page(run_dir, "page_002")
        assert output.read_bytes() == b"owned output"
        assert (second_reconstruction / "stale.txt").read_text(
            encoding="utf-8"
        ) == "stale"
        assert store.read_json("run_state.json")["status"] == "completed"
        return

    status = runtime.retry_page(run_dir, "page_002")

    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "validated"
    assert status["pages"]["pages"]["page_002"]["status"] == "analyzed"
    assert not output.exists()
    assert (first_reconstruction / "keep.txt").read_text(encoding="utf-8") == "keep"
    assert (first_reconstruction / "donor.pptx").read_bytes() == b"stale donor"
    assert not second_reconstruction.exists()
    assert first_plan.read_text(encoding="utf-8") == "keep"
    assert not second_plan.exists()


def test_retry_failed_page_preserves_other_ready_reconstruction(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=2)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.write_json("run_state.json", {
        "schema_version": SCHEMA_VERSION,
        "status": "failed",
        "updated_at": runtime.utc_now(),
    })
    page_jobs = store.read_json("page_jobs.json")
    for page in page_jobs["pages"].values():
        page.update({"status": "failed", "updated_at": runtime.utc_now()})
    store.write_json("page_jobs.json", page_jobs)
    store.write_json("run_summary.json", {
        "schema_version": SCHEMA_VERSION,
        "status": "failed",
    })
    first_reconstruction = run_dir / "pages/page_001/reconstruction"
    second_reconstruction = run_dir / "pages/page_002/reconstruction"
    first_reconstruction.mkdir(parents=True)
    second_reconstruction.mkdir(parents=True)
    (first_reconstruction / "stale.txt").write_text("stale", encoding="utf-8")
    _write_mock_component_state(store, "page_002")
    second_state = store.read_json(
        "pages/page_002/reconstruction/component_state.json"
    )
    second_state.update({
        "phase": "ready_for_assembly",
        "status": "ready_for_assembly",
        "result_ref": {
            "path": "mock-artifact.json",
            "sha256": "0" * 64,
        },
    })
    store.write_json(
        "pages/page_002/reconstruction/component_state.json", second_state
    )
    (second_reconstruction / "donor.pptx").write_bytes(b"validated donor")

    status = runtime.retry_page(run_dir, "page_001")

    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "analyzed"
    assert status["pages"]["pages"]["page_002"]["status"] == "validated"
    assert not first_reconstruction.exists()
    assert (second_reconstruction / "donor.pptx").read_bytes() == b"validated donor"


@pytest.mark.parametrize("stop_reason", ["round_limit", "no_quality_improvement"])
def test_retry_completed_component_warning_resumes_component_state(
    tmp_path: Path,
    stop_reason: str,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    output = run_dir / "final/output.pptx"
    output.parent.mkdir(parents=True)
    output.write_bytes(b"owned output")
    output_hash = runtime.sha256_file(output)
    store.write_json("run_state.json", {
        "schema_version": SCHEMA_VERSION, "status": "completed",
        "updated_at": runtime.utc_now(),
    })
    page_jobs = store.read_json("page_jobs.json")
    page_jobs["pages"]["page_001"].update({
        "status": "preserved_with_warning", "updated_at": runtime.utc_now(),
    })
    store.write_json("page_jobs.json", page_jobs)
    store.write_json("run_summary.json", {
        "schema_version": SCHEMA_VERSION, "status": "completed",
        "outputs": {"pptx": str(output)}, "output_sha256": output_hash,
        "page_results": [{
            "page_id": "page_001", "status": "preserved_with_warning",
        }],
        "input_sha256": manifest["input"]["sha256"],
    })
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(parents=True)
    _write_mock_component_state(store, "page_001")
    state = store.read_json(
        "pages/page_001/reconstruction/component_state.json"
    )
    reference = state["graph_ref"]
    graph_path = store.root / reference["path"]
    graph_path.write_text(json.dumps({
        "nodes": [{
            "id": "candidate", "kind": "parent", "parent_id": None,
            "state": "pending", "mask": "masks/candidate.png",
            "mask_sha256": "2" * 64, "bbox": [0, 0, 1, 1],
            "z_index": 0, "text_ids": [],
        }],
    }), encoding="utf-8")
    reference["sha256"] = runtime.sha256_file(graph_path)
    state.update({
        "phase": "preserved_with_warning",
        "status": "preserved_with_warning",
        "stop_reason": stop_reason,
        "plan_count": 1,
        "candidate_ids": ["candidate"],
        "failed_ids": ["candidate"],
        "fallback": {"status": "warning", "parent_ids": []},
        "current_round": {
            "round": 1, "request_ref": reference, "plan_ref": reference,
            "execution_ref": reference, "quality_ref": reference,
        },
        "round_history": [{
            "round": 1, "plan_sha256": "0" * 64,
            "normalized_plan_sha256": "0" * 64,
            "execution_sha256": "0" * 64, "quality_sha256": "0" * 64,
            "frozen_ids": [], "failed_ids": ["candidate"],
        }],
    })
    if stop_reason == "round_limit":
        state["repair_round"] = MAX_REPAIR_ROUNDS
        state["plan_count"] = MAX_REPAIR_ROUNDS
        state["current_round"]["round"] = MAX_REPAIR_ROUNDS
        state["round_history"] = [
            {
                **state["round_history"][0],
                "round": repair_round,
            }
            for repair_round in range(1, MAX_REPAIR_ROUNDS + 1)
        ]
    store.write_json(
        "pages/page_001/reconstruction/component_state.json", state
    )

    if stop_reason == "round_limit":
        with pytest.raises(RuntimeError, match="non-resumable repair boundary"):
            runtime.retry_page(run_dir, "page_001")
        assert output.read_bytes() == b"owned output"
        assert reconstruction.exists()
        assert store.read_json("run_state.json")["status"] == "completed"
        assert store.read_json(
            "pages/page_001/reconstruction/component_state.json"
        ) == state
        return

    status = runtime.retry_page(run_dir, "page_001")

    resumed = store.read_json(
        "pages/page_001/reconstruction/component_state.json"
    )
    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "analyzed"
    assert resumed["phase"] == "freeze_committed"
    assert resumed["status"] == "active"
    assert reconstruction.exists()
    assert not output.exists()
    assert component_repair.advance_component_repair(store, "page_001") == {
        "status": "needs_next_round",
        "page_id": "page_001",
        "repair_round": 2,
        "candidate_ids": ["candidate"],
        "page_violations": [],
    }


def test_retry_page_removes_work_before_resetting_state(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    _mock_legacy_completion(
        monkeypatch,
        lambda store: (_ for _ in ()).throw(RuntimeError("failed")),
    )
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)
    work_root = run_dir / "work"
    work_root.mkdir()
    (work_root / "diagnostic.txt").write_text("keep until retry", encoding="utf-8")

    status = runtime.retry_page(run_dir, "page_001")

    assert not work_root.exists()
    assert status["run"]["status"] == "prepared"
    assert status["pages"]["pages"]["page_001"]["status"] == "pending"


def test_retry_work_cleanup_failure_preserves_all_state(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    _mock_legacy_completion(
        monkeypatch,
        lambda store: (_ for _ in ()).throw(RuntimeError("failed")),
    )
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    work_root = run_dir / "work"
    work_root.mkdir()
    (work_root / "diagnostic.txt").write_text("keep", encoding="utf-8")
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    before_summary = store.read_json("run_summary.json")

    def fail_cleanup(path: Path, expected_identity: tuple[int, int]) -> None:
        assert path == work_root.resolve()
        raise OSError("work cleanup failed")

    monkeypatch.setattr(runtime, "_safe_rmtree", fail_cleanup)

    with pytest.raises(OSError, match="work cleanup failed"):
        runtime.retry_page(run_dir, "page_001")

    assert work_root.is_dir()
    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert store.read_json("run_summary.json") == before_summary


def test_retry_rejects_work_symlink_without_deleting_external_target(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    _mock_legacy_completion(
        monkeypatch,
        lambda store: (_ for _ in ()).throw(RuntimeError("failed")),
    )
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    external = tmp_path / "external"
    external.mkdir()
    sentinel = external / "sentinel.txt"
    sentinel.write_text("outside", encoding="utf-8")
    work_root = run_dir / "work"
    try:
        work_root.symlink_to(external, target_is_directory=True)
    except OSError as error:
        pytest.skip(f"Cannot create symlink: {error}")
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    before_summary = store.read_json("run_summary.json")

    with pytest.raises(RuntimeError, match="work"):
        runtime.retry_page(run_dir, "page_001")

    assert sentinel.read_text(encoding="utf-8") == "outside"
    assert work_root.is_symlink()
    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages
    assert store.read_json("run_summary.json") == before_summary


@pytest.mark.parametrize("module", [legacy, runtime])
def test_work_safety_detects_windows_reparse_attribute(
    module: Any,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    reparse_flag = 0x400
    monkeypatch.setattr(
        module.stat,
        "FILE_ATTRIBUTE_REPARSE_POINT",
        reparse_flag,
        raising=False,
    )
    directory_mode = module.stat.S_IFDIR

    assert module._is_link_or_reparse(
        types.SimpleNamespace(
            st_mode=directory_mode,
            st_file_attributes=reparse_flag,
        )
    )
    assert not module._is_link_or_reparse(
        types.SimpleNamespace(
            st_mode=directory_mode,
            st_file_attributes=0,
        )
    )


def test_retry_page_writes_page_jobs_once(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    _image(first)
    _image(second)
    run_dir = runtime.prepare_job([first, second], run_dir=tmp_path / "run")

    def fail_execute(store: RunStore) -> dict[str, Any]:
        raise RuntimeError("failed")

    _mock_legacy_completion(monkeypatch, fail_execute)
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)

    original_write_json = RunStore.write_json
    page_jobs_writes = 0

    def count_page_jobs_write(
        self: RunStore, relative: str | Path, document: dict[str, Any]
    ) -> None:
        nonlocal page_jobs_writes
        if Path(relative) == Path("page_jobs.json"):
            page_jobs_writes += 1
        original_write_json(self, relative, document)

    monkeypatch.setattr(RunStore, "write_json", count_page_jobs_write)

    runtime.retry_page(run_dir, "page_001")

    assert page_jobs_writes == 1


def test_retry_page_write_failure_preserves_page_jobs(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    first = tmp_path / "first.png"
    second = tmp_path / "second.png"
    _image(first)
    _image(second)
    run_dir = runtime.prepare_job([first, second], run_dir=tmp_path / "run")

    def fail_execute(store: RunStore) -> dict[str, Any]:
        raise RuntimeError("failed")

    _mock_legacy_completion(monkeypatch, fail_execute)
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)

    store = RunStore.open(run_dir)
    original_page_jobs = store.read_json("page_jobs.json")
    original_write_json = RunStore.write_json
    page_jobs_writes = 0

    def fail_page_jobs_write(
        self: RunStore, relative: str | Path, document: dict[str, Any]
    ) -> None:
        nonlocal page_jobs_writes
        if Path(relative) == Path("page_jobs.json"):
            page_jobs_writes += 1
            raise OSError("page jobs write failed")
        original_write_json(self, relative, document)

    monkeypatch.setattr(RunStore, "write_json", fail_page_jobs_write)

    with pytest.raises(OSError, match="page jobs write failed"):
        runtime.retry_page(run_dir, "page_001")

    assert page_jobs_writes == 1
    assert store.read_json("page_jobs.json") == original_page_jobs
    assert store.read_json("run_state.json")["status"] == "failed"

    monkeypatch.setattr(RunStore, "write_json", original_write_json)

    retried = runtime.retry_page(run_dir, "page_001")
    repeated = runtime.retry_page(run_dir, "page_001")

    assert retried["run"]["status"] == "prepared"
    assert {
        page["status"] for page in retried["pages"]["pages"].values()
    } == {"pending"}
    assert repeated == retried


def test_pptx_retry_writes_pages_before_run_and_recovers_run_write_failure(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    monkeypatch.setattr(
        runtime,
        "execute_pptx_preserve",
        lambda store: (_ for _ in ()).throw(RuntimeError("failed")),
    )
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)

    original_transition_run = RunStore.transition_run
    failed = False

    def fail_prepared_once(
        self: RunStore, target: RunStatus
    ) -> dict[str, Any]:
        nonlocal failed
        if target is RunStatus.PREPARED and not failed:
            failed = True
            raise OSError("prepared state write failed")
        return original_transition_run(self, target)

    monkeypatch.setattr(RunStore, "transition_run", fail_prepared_once)

    with pytest.raises(OSError, match="prepared state write failed"):
        runtime.retry_page(run_dir, "page_001")

    store = RunStore.open(run_dir)
    assert store.read_json("run_state.json")["status"] == "failed"
    assert store.read_json("page_jobs.json")["pages"]["page_001"]["status"] == "analyzed"
    retried = runtime.retry_page(run_dir, "page_001")
    assert retried["run"]["status"] == "prepared"
    assert retried["pages"]["pages"]["page_001"]["status"] == "analyzed"


def test_pptx_retry_removes_stale_reconstruction_donor(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    monkeypatch.setattr(
        runtime,
        "execute_pptx_preserve",
        lambda store: (_ for _ in ()).throw(RuntimeError("failed")),
    )
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir()
    (reconstruction / "donor.pptx").write_bytes(b"stale")

    runtime.retry_page(run_dir, "page_001")

    assert not reconstruction.exists()


def test_pptx_recover_removes_stale_reconstruction_donor(
    tmp_path: Path,
) -> None:
    source = tmp_path / "source.pptx"
    _pptx(source, slide_count=1)
    run_dir = runtime.prepare_job(source, run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    store.transition_page("page_001", PageStatus.PROCESSING)
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir()
    (reconstruction / "donor.pptx").write_bytes(b"stale")

    runtime.recover_job(run_dir)

    assert not reconstruction.exists()


def test_retry_page_rejects_unknown_or_nonfailed_page(tmp_path: Path) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")

    with pytest.raises(KeyError, match="Unknown page_id"):
        runtime.retry_page(run_dir, "missing")
    with pytest.raises(RuntimeError, match="not failed"):
        runtime.retry_page(run_dir, "page_001")


def test_run_job_rejects_non_prepared_run_without_changing_state(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    store.transition_run(RunStatus.RUNNING)
    before_run = store.read_json("run_state.json")
    before_pages = store.read_json("page_jobs.json")
    called = False

    def fake_execute(store: RunStore) -> dict[str, Any]:
        nonlocal called
        called = True
        return {}

    monkeypatch.setattr(runtime, "execute_legacy", fake_execute)

    with pytest.raises(RuntimeError, match="must be prepared"):
        runtime.run_job(run_dir)

    assert not called
    assert store.read_json("run_state.json") == before_run
    assert store.read_json("page_jobs.json") == before_pages


def test_run_job_returns_existing_completed_summary(tmp_path: Path, monkeypatch) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    _mock_legacy_completion(monkeypatch, lambda store: {})
    completed = runtime.run_job(run_dir)

    def unexpected_execute(store: RunStore) -> dict[str, Any]:
        raise AssertionError("completed run executed again")

    monkeypatch.setattr(runtime, "execute_legacy", unexpected_execute)

    assert runtime.run_job(run_dir) == completed


def test_run_job_validates_existing_completed_summary(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    _mock_legacy_completion(monkeypatch, lambda store: {})
    runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    summary = store.read_json("run_summary.json")
    summary["schema_version"] = 2
    store.write_json("run_summary.json", summary)

    with pytest.raises(ValueError, match="Unsupported schema_version"):
        runtime.run_job(run_dir)


def test_retry_validates_existing_failed_summary(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")

    def fail(store: RunStore) -> dict[str, Any]:
        raise RuntimeError("failed")

    _mock_legacy_completion(monkeypatch, fail)
    with pytest.raises(RuntimeError, match="failed"):
        runtime.run_job(run_dir)
    store = RunStore.open(run_dir)
    summary = store.read_json("run_summary.json")
    summary["schema_version"] = 2
    store.write_json("run_summary.json", summary)

    with pytest.raises(ValueError, match="Unsupported schema_version"):
        runtime.retry_page(run_dir, "page_001")


@pytest.mark.parametrize(
    ("image_count", "slide_size", "function_name", "expected_extra"),
    [
        (1, "both", "convert_variants", {}),
        (1, "original", "convert", {"slide_size": "original"}),
        (1, "16:9", "convert", {"slide_size": "16:9"}),
        (2, "both", "convert_batch_variants", {}),
        (
            2,
            "original",
            "convert_batch_variants",
            {"include_widescreen": False},
        ),
        (2, "16:9", "convert_batch", {}),
    ],
)
def test_execute_legacy_dispatches_with_real_signatures(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    capsys: pytest.CaptureFixture[str],
    image_count: int,
    slide_size: str,
    function_name: str,
    expected_extra: dict[str, Any],
) -> None:
    sources = []
    for index in range(image_count):
        source = tmp_path / f"source-{index}.png"
        _image(source)
        sources.append(source)
    output_path = tmp_path / "chosen.pptx"
    run_dir = runtime.prepare_job(
        sources,
        run_dir=tmp_path / "run",
        output_path=output_path,
        slide_size=slide_size,
        lang="en",
    )
    calls: list[tuple[str, tuple[Any, ...], dict[str, Any]]] = []

    def record(name: str):
        def converter(*args: Any, **kwargs: Any) -> Any:
            print("legacy progress")
            calls.append((name, args, kwargs))
            if name in {"convert_variants", "convert_batch_variants"}:
                return {"original": "original.pptx", "16:9": "wide.pptx"}
            return "single.pptx"

        return converter

    fake_module = types.SimpleNamespace(
        convert=record("convert"),
        convert_variants=record("convert_variants"),
        convert_batch=record("convert_batch"),
        convert_batch_variants=record("convert_batch_variants"),
    )
    monkeypatch.setattr(legacy.importlib, "import_module", lambda name: fake_module)

    result = legacy.execute_legacy(RunStore.open(run_dir))
    captured = capsys.readouterr()

    assert len(calls) == 1
    assert captured.out == ""
    assert "legacy progress" in captured.err
    name, args, kwargs = calls[0]
    assert name == function_name
    copied_sources = [
        (Path(run_dir) / "input" / f"{index:03d}_source-{index - 1}.png").resolve()
        for index in range(1, image_count + 1)
    ]
    assert args[0] == (copied_sources[0] if image_count == 1 else copied_sources)
    assert kwargs == {
        "output_path": str(output_path.resolve()),
        "lang": "en",
        "_work_root": (run_dir / "work").resolve(),
        "_resource_isolation": True,
        **expected_extra,
    }
    assert not (run_dir / "work").exists()
    assert set(result) == ({"original", "16:9"} if "variants" in name else {slide_size})


@pytest.mark.parametrize(
    ("input_document", "slide_size", "expected_name", "expected_extra"),
    [
        (
            {"type": "pdf", "page_ratios_equal": True, "page_aspect_ratio": 2.0},
            "both",
            "convert_batch_variants",
            {"combine_original": True, "original_aspect_ratio": 2.0},
        ),
        (
            {"type": "pdf", "page_ratios_equal": True, "page_aspect_ratio": 2.0},
            "original",
            "convert_batch_variants",
            {
                "include_widescreen": False,
                "combine_original": True,
                "original_aspect_ratio": 2.0,
            },
        ),
        (
            {"type": "pdf", "page_ratios_equal": True, "page_aspect_ratio": 2.0},
            "16:9",
            "convert_batch",
            {},
        ),
        (
            {"type": "pdf", "page_ratios_equal": False, "page_aspect_ratio": None},
            "both",
            "convert_batch_variants",
            {},
        ),
        (
            {"type": "pdf", "page_ratios_equal": True},
            "both",
            "convert_batch_variants",
            {"combine_original": True},
        ),
        ({"type": "images"}, "both", "convert_batch_variants", {}),
    ],
)
def test_execute_legacy_combines_only_equal_ratio_pdf_originals(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    input_document: dict[str, Any],
    slide_size: str,
    expected_name: str,
    expected_extra: dict[str, Any],
) -> None:
    sources = [tmp_path / "first.png", tmp_path / "second.png"]
    for source in sources:
        _image(source)
    run_dir = runtime.prepare_job(
        sources,
        run_dir=tmp_path / "run",
        slide_size=slide_size,
    )
    store = RunStore.open(run_dir)
    manifest = store.read_json("job_manifest.json")
    manifest["input"] = input_document
    store.write_json("job_manifest.json", manifest)
    calls: list[tuple[str, dict[str, Any]]] = []

    def record(name: str):
        def converter(*args: Any, **kwargs: Any) -> Any:
            calls.append((name, kwargs))
            return {"original": "original.pptx", "16:9": "wide.pptx"}

        return converter

    fake_module = types.SimpleNamespace(
        convert_batch=record("convert_batch"),
        convert_batch_variants=record("convert_batch_variants"),
    )
    monkeypatch.setattr(legacy.importlib, "import_module", lambda name: fake_module)

    legacy.execute_legacy(store)

    assert calls == [
        (
            expected_name,
            {
                    "output_path": str((run_dir / "final" / "output.pptx").resolve()),
                    "lang": "ch",
                    "_work_root": (run_dir / "work").resolve(),
                    "_resource_isolation": True,
                    **expected_extra,
            },
        )
    ]
    assert not (run_dir / "work").exists()


def test_execute_legacy_uses_default_output_path(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(
        [source], run_dir=tmp_path / "run", slide_size="16:9"
    )
    captured: dict[str, Any] = {}

    def convert_image(
        image_path: Path,
        output_path: str,
        lang: str,
        slide_size: str,
        _work_root: Path,
        _resource_isolation: bool,
    ) -> str:
        captured.update(
            output_path=output_path,
            lang=lang,
            slide_size=slide_size,
            _work_root=_work_root,
            _resource_isolation=_resource_isolation,
        )
        return output_path

    fake_module = types.SimpleNamespace(convert=convert_image)
    monkeypatch.setattr(legacy.importlib, "import_module", lambda name: fake_module)

    legacy.execute_legacy(RunStore.open(run_dir))

    assert captured == {
        "output_path": str((run_dir / "final" / "output.pptx").resolve()),
        "lang": "ch",
        "slide_size": "16:9",
        "_work_root": (run_dir / "work").resolve(),
        "_resource_isolation": True,
    }
    assert not (run_dir / "work").exists()


def test_execute_legacy_accepts_preexisting_empty_work_directory(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(
        [source],
        run_dir=tmp_path / "run",
        slide_size="16:9",
    )
    work_root = run_dir / "work"
    work_root.mkdir()
    seen_roots = []

    def convert_image(*args: Any, _work_root: Path, **kwargs: Any) -> str:
        seen_roots.append(_work_root)
        return str(tmp_path / "output.pptx")

    monkeypatch.setattr(
        legacy.importlib,
        "import_module",
        lambda name: types.SimpleNamespace(convert=convert_image),
    )

    legacy.execute_legacy(RunStore.open(run_dir))

    assert seen_roots == [work_root.resolve()]
    assert not work_root.exists()


@pytest.mark.skipif(os.name != "nt", reason="Windows directory handle semantics")
@pytest.mark.parametrize("replacement", ["root", "nested"])
def test_execute_legacy_cleanup_rejects_directory_replacement_during_enumeration(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    replacement: str,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(
        [source],
        run_dir=tmp_path / "run",
        slide_size="16:9",
    )
    work_root = run_dir / "work"
    external = tmp_path / "external"
    external.mkdir()
    sentinel = external / "sentinel.txt"
    sentinel.write_text("outside", encoding="utf-8")
    target = work_root if replacement == "root" else work_root / "nested"
    displaced = tmp_path / "displaced"
    converted = False

    def convert_image(*args: Any, _work_root: Path, **kwargs: Any) -> str:
        nonlocal converted
        target.mkdir(parents=True, exist_ok=True)
        (target / "owned.txt").write_text("owned", encoding="utf-8")
        converted = True
        return str(tmp_path / "output.pptx")

    monkeypatch.setattr(
        legacy.importlib,
        "import_module",
        lambda name: types.SimpleNamespace(convert=convert_image),
    )
    original_entries = getattr(legacy, "_windows_entries", None)
    attempted = False

    def replace_directory_during_enumeration(
        kernel32: Any,
        handle: Any,
        path: Path,
        status: Any,
    ) -> Any:
        nonlocal attempted
        assert original_entries is not None
        entries = original_entries(kernel32, handle, path, status)
        if converted and not attempted and path == target:
            attempted = True
            target.rename(displaced)
            external.rename(target)
        return entries

    monkeypatch.setattr(
        legacy,
        "_windows_entries",
        replace_directory_during_enumeration,
        raising=False,
    )

    with pytest.raises(OSError):
        legacy.execute_legacy(RunStore.open(run_dir))

    assert attempted
    assert sentinel.read_text(encoding="utf-8") == "outside"
    assert external.is_dir()


@pytest.mark.skipif(os.name != "nt", reason="Windows directory handle semantics")
@pytest.mark.parametrize("entry_kind", ["directory", "file"])
def test_execute_legacy_cleanup_rejects_entry_replacement_after_parent_yield(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    entry_kind: str,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job(
        [source],
        run_dir=tmp_path / "run",
        slide_size="16:9",
    )
    work_root = run_dir / "work"
    target = work_root / ("nested" if entry_kind == "directory" else "owned.txt")
    displaced = tmp_path / "displaced"
    external = tmp_path / "external"
    external.mkdir()
    sentinel = external / "sentinel.txt"
    sentinel.write_text("outside", encoding="utf-8")
    converted = False
    swapped = False

    def convert_image(*args: Any, _work_root: Path, **kwargs: Any) -> str:
        nonlocal converted
        work_root.mkdir(exist_ok=True)
        if entry_kind == "directory":
            target.mkdir()
            (target / "owned.txt").write_text("owned", encoding="utf-8")
        else:
            target.write_text("owned", encoding="utf-8")
        converted = True
        return str(tmp_path / "output.pptx")

    monkeypatch.setattr(
        legacy.importlib,
        "import_module",
        lambda name: types.SimpleNamespace(convert=convert_image),
    )
    original_entries = getattr(legacy, "_windows_entries", None)

    def replace_after_parent_enumeration(
        kernel32: Any,
        handle: Any,
        path: Path,
        status: Any,
    ) -> Any:
        nonlocal swapped
        assert original_entries is not None
        entries = original_entries(kernel32, handle, path, status)
        if converted and not swapped and path == work_root:
            target.rename(displaced)
            if entry_kind == "directory":
                external.rename(target)
            else:
                sentinel.rename(target)
            swapped = True
        return entries

    monkeypatch.setattr(
        legacy,
        "_windows_entries",
        replace_after_parent_enumeration,
        raising=False,
    )

    with pytest.raises(RuntimeError, match="changed"):
        legacy.execute_legacy(RunStore.open(run_dir))

    assert swapped
    if entry_kind == "directory":
        target.rename(external)
    else:
        target.rename(sentinel)
    assert sentinel.read_text(encoding="utf-8") == "outside"


@pytest.mark.parametrize("case", ["nonempty", "file"])
def test_execute_legacy_rejects_unsafe_existing_work_before_import(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    case: str,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    work_root = run_dir / "work"
    if case == "nonempty":
        work_root.mkdir()
        (work_root / "sentinel.txt").write_text("keep", encoding="utf-8")
    else:
        work_root.write_text("keep", encoding="utf-8")

    def unexpected_import(name: str) -> Any:
        raise AssertionError("legacy module imported before work validation")

    monkeypatch.setattr(legacy.importlib, "import_module", unexpected_import)

    with pytest.raises(RuntimeError, match="work"):
        legacy.execute_legacy(RunStore.open(run_dir))


def test_execute_legacy_rejects_work_symlink_before_import_without_external_delete(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    external = tmp_path / "external"
    external.mkdir()
    sentinel = external / "sentinel.txt"
    sentinel.write_text("outside", encoding="utf-8")
    work_root = run_dir / "work"
    try:
        work_root.symlink_to(external, target_is_directory=True)
    except OSError as error:
        pytest.skip(f"Cannot create symlink: {error}")

    def unexpected_import(name: str) -> Any:
        raise AssertionError("legacy module imported before work validation")

    monkeypatch.setattr(legacy.importlib, "import_module", unexpected_import)

    with pytest.raises(RuntimeError, match="work"):
        legacy.execute_legacy(RunStore.open(run_dir))

    assert sentinel.read_text(encoding="utf-8") == "outside"
    assert work_root.is_symlink()


def test_legacy_failure_retains_work_and_records_absolute_diagnostics(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    diagnostic_name = "failure.txt"

    def fail_initialization(store, page_id, **kwargs):
        page_root = store.root / "pages" / page_id / "reconstruction"
        page_root.mkdir(parents=True)
        (page_root / diagnostic_name).write_text("details", encoding="utf-8")
        raise RuntimeError("conversion failed")

    monkeypatch.setattr(runtime, "initialize_legacy_page", fail_initialization)

    with pytest.raises(RuntimeError, match="conversion failed"):
        runtime.run_job(run_dir)

    diagnostics = (run_dir / "pages/page_001/reconstruction").resolve()
    assert (diagnostics / diagnostic_name).read_text(
        encoding="utf-8"
    ) == "details"
    assert RunStore.open(run_dir).read_json("run_summary.json") == {
        "schema_version": SCHEMA_VERSION,
        "status": "failed",
        "error": {"type": "RuntimeError", "message": "conversion failed"},
        "outputs": {},
        "diagnostics": str(diagnostics),
    }


def test_legacy_cleanup_failure_after_conversion_records_failed_run_and_diagnostics(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    diagnostic_name = "conversion.txt"
    converted = False

    def convert_image(*args: Any, _work_root: Path, **kwargs: Any) -> Any:
        nonlocal converted
        page_root = _work_root / "page_001"
        page_root.mkdir()
        (page_root / diagnostic_name).write_text("complete", encoding="utf-8")
        converted = True
        return {"16:9": str(tmp_path / "output.pptx")}

    def fail_cleanup(*args: Any, **kwargs: Any) -> None:
        assert converted
        raise OSError("work cleanup failed")

    monkeypatch.setattr(
        legacy.importlib,
        "import_module",
        lambda name: types.SimpleNamespace(convert_variants=convert_image),
    )
    monkeypatch.setattr(legacy, "_safe_rmtree", fail_cleanup, raising=False)

    with pytest.raises(OSError, match="work cleanup failed"):
        legacy.execute_legacy(RunStore.open(run_dir))

    work_root = (run_dir / "work").resolve()
    assert (work_root / "page_001" / diagnostic_name).read_text(
        encoding="utf-8"
    ) == "complete"


@pytest.mark.parametrize(
    ("case", "reason"),
    [
        ("outside", "outside"),
        ("missing", "not a file"),
        ("directory", "not a file"),
        ("sha256", "sha256 mismatch"),
    ],
)
def test_execute_legacy_rejects_unsafe_or_changed_source_before_import(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    case: str,
    reason: str,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    request_path = Path("pages/page_001/page_request.json")
    request = store.read_json(request_path)
    if case == "outside":
        request["source"] = "../source.png"
    elif case == "missing":
        request["source"] = "input/missing.png"
    elif case == "directory":
        request["source"] = "input"
    else:
        request["sha256"] = "0" * 64
    store.write_json(request_path, request)

    def unexpected_import(name: str) -> Any:
        raise AssertionError("legacy module imported before source validation")

    monkeypatch.setattr(legacy.importlib, "import_module", unexpected_import)

    with pytest.raises(ValueError, match=rf"page_001.*{reason}"):
        legacy.execute_legacy(store)


@pytest.mark.parametrize("document_name", ["job_manifest", "page_request"])
def test_execute_legacy_validates_consumed_document_versions_before_import(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    document_name: str,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    run_dir = runtime.prepare_job([source], run_dir=tmp_path / "run")
    store = RunStore.open(run_dir)
    relative = (
        Path("job_manifest.json")
        if document_name == "job_manifest"
        else Path("pages/page_001/page_request.json")
    )
    document = store.read_json(relative)
    document["schema_version"] = 2
    store.write_json(relative, document)

    def unexpected_import(name: str) -> Any:
        raise AssertionError("legacy module imported before schema validation")

    monkeypatch.setattr(legacy.importlib, "import_module", unexpected_import)

    with pytest.raises(ValueError, match="Unsupported schema_version"):
        legacy.execute_legacy(store)


def test_absolute_outputs_recurses_without_changing_other_values(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    monkeypatch.chdir(tmp_path)
    value = {
        "single": "one.pptx",
        "nested": ["two.pptx", {"empty": None, "count": 2}],
    }

    assert legacy._absolute_outputs(value) == {
        "single": str((tmp_path / "one.pptx").resolve()),
        "nested": [
            str((tmp_path / "two.pptx").resolve()),
            {"empty": None, "count": 2},
        ],
    }


def test_convert_prepares_then_runs(monkeypatch: pytest.MonkeyPatch) -> None:
    calls: list[tuple[str, Any]] = []
    prepared = Path("prepared-run")

    def fake_prepare(inputs: list[str], **kwargs: Any) -> Path:
        calls.append(("prepare", (inputs, kwargs)))
        return prepared

    def fake_run(run_dir: Path) -> dict[str, Any]:
        calls.append(("run", run_dir))
        return {"status": "completed"}

    monkeypatch.setattr(runtime, "prepare_job", fake_prepare)
    monkeypatch.setattr(runtime, "run_job", fake_run)

    result = runtime.convert(
        ["source.png"],
        run_dir="run",
        output_path="output.pptx",
        slide_size="original",
        lang="en",
    )

    assert calls == [
        (
            "prepare",
            (
                ["source.png"],
                {
                    "run_dir": "run",
                    "output_path": "output.pptx",
                    "slide_size": "original",
                    "lang": "en",
                    "agent_provider": "host",
                },
            ),
        ),
        ("run", prepared),
    ]
    assert result == {"status": "completed"}


@pytest.mark.parametrize("as_string", [False, True])
def test_convert_accepts_one_path_directly(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    as_string: bool,
) -> None:
    source = tmp_path / "source.png"
    _image(source)
    value = str(source) if as_string else source
    _mock_legacy_completion(monkeypatch, lambda store: {})

    summary = runtime.convert(value, run_dir=tmp_path / "run")

    assert summary["status"] == "completed"
    assert summary["pages"] == 1
