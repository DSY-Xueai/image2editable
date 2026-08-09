from pathlib import Path

from PIL import Image
from pptx import Presentation
import pytest

from image2editable.powerpoint_renderer import (
    PowerPointRenderer,
    RendererUnavailable,
)


class _FakeSlide:
    def __init__(self, *, fail: bool = False) -> None:
        self.fail = fail

    def Export(self, output: str, format_name: str, width: int, height: int) -> None:
        assert format_name == "PNG"
        if self.fail:
            raise RuntimeError("export failed")
        Image.new("RGB", (width, height), "white").save(output)


class _FakePresentation:
    def __init__(self, *, fail: bool = False) -> None:
        self.closed = False
        self.slide = _FakeSlide(fail=fail)

    def Slides(self, page_number: int) -> _FakeSlide:
        assert page_number == 1
        return self.slide

    def Close(self) -> None:
        self.closed = True


class _FakePresentations:
    def __init__(self, presentation: _FakePresentation) -> None:
        self.presentation = presentation
        self.open_arguments = None

    def Open(self, path: str, **kwargs) -> _FakePresentation:
        self.open_arguments = (path, kwargs)
        return self.presentation


class _FakeApplication:
    Version = "16.0"

    def __init__(self, presentation: _FakePresentation) -> None:
        self.Presentations = _FakePresentations(presentation)
        self.quit_called = False

    def Quit(self) -> None:
        self.quit_called = True


class _FakeDispatch:
    def __init__(self, *, fail: bool = False) -> None:
        self.presentation = _FakePresentation(fail=fail)
        self.application = _FakeApplication(self.presentation)

    def __call__(self, prog_id: str) -> _FakeApplication:
        assert prog_id == "PowerPoint.Application"
        return self.application


def test_renderer_reports_unavailable_without_pywin32(tmp_path: Path) -> None:
    renderer = PowerPointRenderer(dispatch_factory=None)

    assert renderer.available() is False
    with pytest.raises(RendererUnavailable):
        renderer.render_page(
            tmp_path / "input.pptx",
            1,
            tmp_path / "page.png",
            width=1600,
            height=900,
        )


def test_renderer_closes_presentation_and_quits_application(tmp_path: Path) -> None:
    fake_dispatch = _FakeDispatch()
    renderer = PowerPointRenderer(dispatch_factory=fake_dispatch)
    output = tmp_path / "page.png"

    result = renderer.render_page(
        tmp_path / "input.pptx", 1, output, width=1600, height=900
    )

    assert result == {
        "renderer": "powerpoint",
        "version": "16.0",
        "width": 1600,
        "height": 900,
        "path": str(output.resolve()),
    }
    assert fake_dispatch.presentation.closed is True
    assert fake_dispatch.application.quit_called is True
    assert fake_dispatch.application.Presentations.open_arguments[1] == {
        "ReadOnly": True,
        "Untitled": False,
        "WithWindow": False,
    }


def test_renderer_cleans_up_when_export_fails(tmp_path: Path) -> None:
    fake_dispatch = _FakeDispatch(fail=True)
    renderer = PowerPointRenderer(dispatch_factory=fake_dispatch)

    with pytest.raises(RuntimeError, match="export failed"):
        renderer.render_page(
            tmp_path / "input.pptx",
            1,
            tmp_path / "page.png",
            width=1600,
            height=900,
        )

    assert fake_dispatch.presentation.closed is True
    assert fake_dispatch.application.quit_called is True


def test_renderer_rejects_wrong_export_size(tmp_path: Path) -> None:
    fake_dispatch = _FakeDispatch()
    fake_dispatch.presentation.slide.Export = lambda output, *_: Image.new(
        "RGB", (10, 10), "white"
    ).save(output)
    renderer = PowerPointRenderer(dispatch_factory=fake_dispatch)

    with pytest.raises(RuntimeError, match="dimensions"):
        renderer.render_page(
            tmp_path / "input.pptx",
            1,
            tmp_path / "page.png",
            width=1600,
            height=900,
        )


def test_renderer_does_not_accept_stale_output(tmp_path: Path) -> None:
    fake_dispatch = _FakeDispatch()
    fake_dispatch.presentation.slide.Export = lambda *_: None
    output = tmp_path / "page.png"
    Image.new("RGB", (1600, 900), "white").save(output)
    renderer = PowerPointRenderer(dispatch_factory=fake_dispatch)

    with pytest.raises(RuntimeError, match="did not produce"):
        renderer.render_page(
            tmp_path / "input.pptx", 1, output, width=1600, height=900
        )


@pytest.mark.powerpoint
def test_real_powerpoint_exports_one_slide(tmp_path: Path) -> None:
    renderer = PowerPointRenderer.discover()
    if not renderer.available():
        pytest.skip("PowerPoint COM dependencies are unavailable")
    source = tmp_path / "source.pptx"
    presentation = Presentation()
    presentation.slides.add_slide(presentation.slide_layouts[6])
    presentation.save(source)

    try:
        result = renderer.render_page(
            source, 1, tmp_path / "page.png", width=400, height=300
        )
    except RendererUnavailable as error:
        pytest.skip(str(error))

    assert (result["width"], result["height"]) == (400, 300)
