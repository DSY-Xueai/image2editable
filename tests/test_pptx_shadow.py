from __future__ import annotations

import posixpath
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches
import pytest

from image2editable.pptx_input import scan_pptx
from image2editable.pptx_shadow import patch_slide_background


P = "http://schemas.openxmlformats.org/presentationml/2006/main"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
PR = "http://schemas.openxmlformats.org/package/2006/relationships"


def test_patch_slide_background_transplants_reconstruction_only(tmp_path):
    source = _source_with_background_and_native_text(tmp_path)
    donor = _reconstruction_donor(tmp_path)
    output = tmp_path / "output.pptx"
    source_before = source.read_bytes()
    source_members = _members(source)
    native_before = _object_by_name(scan_pptx(source), "Native overlay")

    result = patch_slide_background(
        source,
        donor,
        output,
        slide_part="ppt/slides/slide1.xml",
    )

    assert source.read_bytes() == source_before
    assert result == {
        "source_shape_id": "background",
        "slide_part": "ppt/slides/slide1.xml",
        "imported_shapes": 3,
        "imported_media": 2,
    }
    inventory = scan_pptx(output)
    names = [item["name"] for item in inventory["slides"][0]["objects"]]
    assert names[-1] == "Native overlay"
    assert "Slide Background Image" not in names
    assert (
        _object_by_name(inventory, "Native overlay")["xml_c14n_sha256"]
        == native_before["xml_c14n_sha256"]
    )

    changed = {
        "[Content_Types].xml",
        "ppt/slides/slide1.xml",
        "ppt/slides/_rels/slide1.xml.rels",
    }
    output_members = _members(output)
    for name, contents in source_members.items():
        if name not in changed:
            assert output_members[name] == contents
    assert len(output_members) == len(source_members) + 2

    reopened = Presentation(output)
    assert [
        shape.text for shape in reopened.slides[0].shapes if shape.has_text_frame
    ] == [
        "editable donor text",
        "native text",
    ]


def test_patch_slide_background_assigns_unique_shape_ids_and_valid_media_targets(
    tmp_path,
):
    source = _source_with_background_and_native_text(tmp_path)
    donor = _reconstruction_donor(tmp_path)
    output = tmp_path / "output.pptx"

    patch_slide_background(
        source,
        donor,
        output,
        slide_part="ppt/slides/slide1.xml",
    )

    with zipfile.ZipFile(output) as archive:
        slide = ET.fromstring(archive.read("ppt/slides/slide1.xml"))
        shape_ids = [item.get("id") for item in slide.findall(f".//{{{P}}}cNvPr")]
        assert len(shape_ids) == len(set(shape_ids))
        relationships = ET.fromstring(archive.read("ppt/slides/_rels/slide1.xml.rels"))
        names = set(archive.namelist())
        for relationship in relationships.findall(f"{{{PR}}}Relationship"):
            if relationship.get("Type", "").endswith("/image"):
                target = relationship.get("Target")
                assert _resolve_slide_target(target) in names


def test_patch_slide_background_replaces_full_slide_picture_shape(tmp_path):
    source = _source_with_picture_and_native_text(tmp_path)
    donor = _reconstruction_donor(tmp_path)
    output = tmp_path / "output.pptx"
    source_before = source.read_bytes()
    before = scan_pptx(source)
    picture = _object_by_name(before, "Screenshot")
    native_before = _object_by_name(before, "Native overlay")

    result = patch_slide_background(
        source,
        donor,
        output,
        slide_part="ppt/slides/slide1.xml",
        source_shape_id=picture["shape_id"],
    )

    assert source.read_bytes() == source_before
    assert result["source_shape_id"] == picture["shape_id"]
    inventory = scan_pptx(output)
    names = [item["name"] for item in inventory["slides"][0]["objects"]]
    assert "Screenshot" not in names
    assert names[-1] == "Native overlay"
    assert (
        _object_by_name(inventory, "Native overlay")["xml_c14n_sha256"]
        == native_before["xml_c14n_sha256"]
    )
    assert all(
        item["shape_id"] != picture["shape_id"]
        for item in inventory["slides"][0]["objects"]
    )
    assert (
        Presentation(output)
        .slides[0]
        .notes_slide.notes_text_frame.text.endswith("keep this note")
    )


def test_patch_slide_background_maps_donor_to_original_picture_bounds(tmp_path):
    source = _source_with_picture_and_native_text(
        tmp_path,
        left=0.5,
        top=0.25,
        width=9,
        height=4.5,
    )
    donor = _reconstruction_donor(tmp_path)
    output = tmp_path / "output.pptx"
    before = scan_pptx(source)
    picture = _object_by_name(before, "Screenshot")

    patch_slide_background(
        source,
        donor,
        output,
        slide_part="ppt/slides/slide1.xml",
        source_shape_id=picture["shape_id"],
    )

    after = scan_pptx(output)
    imported_background = max(
        (item for item in after["slides"][0]["objects"] if item["type"] == "picture"),
        key=lambda item: item["cx"] * item["cy"],
    )
    assert (
        imported_background["x"],
        imported_background["y"],
        imported_background["cx"],
        imported_background["cy"],
    ) == (
        picture["x"],
        picture["y"],
        picture["cx"],
        picture["cy"],
    )


def test_patch_slide_background_rejects_picture_with_connector_reference(
    tmp_path,
):
    source = _source_with_picture_and_connector(tmp_path)
    donor = _reconstruction_donor(tmp_path)
    output = tmp_path / "output.pptx"
    picture = _object_by_name(scan_pptx(source), "Screenshot")

    with pytest.raises(ValueError, match="referenced"):
        patch_slide_background(
            source,
            donor,
            output,
            slide_part="ppt/slides/slide1.xml",
            source_shape_id=picture["shape_id"],
        )

    assert not output.exists()


def test_patch_slide_background_refuses_existing_output(tmp_path):
    source = _source_with_background_and_native_text(tmp_path)
    donor = _reconstruction_donor(tmp_path)
    output = tmp_path / "output.pptx"
    output.write_bytes(b"user")

    try:
        patch_slide_background(
            source,
            donor,
            output,
            slide_part="ppt/slides/slide1.xml",
        )
    except FileExistsError:
        pass
    else:
        raise AssertionError("existing output was overwritten")

    assert output.read_bytes() == b"user"


def _source_with_background_and_native_text(tmp_path: Path) -> Path:
    source = _source_with_picture_and_native_text(tmp_path)
    _move_first_picture_to_background(source)
    return source


def _source_with_picture_and_native_text(
    tmp_path: Path,
    *,
    left: float = 0,
    top: float = 0,
    width: float = 10,
    height: float = 5,
) -> Path:
    source_image = tmp_path / "source.png"
    Image.new("RGB", (400, 200), "navy").save(source_image)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    screenshot = slide.shapes.add_picture(
        str(source_image),
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    screenshot.name = "Screenshot"
    native = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1))
    native.name = "Native overlay"
    native.text = "native text"
    slide.notes_slide.notes_text_frame.text = "keep this note"
    source = tmp_path / "source.pptx"
    presentation.save(source)
    return source


def _source_with_picture_and_connector(tmp_path: Path) -> Path:
    source_image = tmp_path / "source.png"
    Image.new("RGB", (400, 200), "navy").save(source_image)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    screenshot = slide.shapes.add_picture(
        str(source_image), 0, 0, Inches(10), Inches(5)
    )
    screenshot.name = "Screenshot"
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(1),
        Inches(1),
        Inches(2),
        Inches(2),
    )
    connection = OxmlElement("a:stCxn")
    connection.set("id", str(screenshot.shape_id))
    connection.set("idx", "0")
    connector._element.nvCxnSpPr.cNvCxnSpPr.append(connection)
    source = tmp_path / "source-with-connector.pptx"
    presentation.save(source)
    return source


def _reconstruction_donor(tmp_path: Path) -> Path:
    background = tmp_path / "clean-background.png"
    component = tmp_path / "component.png"
    Image.new("RGB", (400, 200), "white").save(background)
    Image.new("RGBA", (80, 60), (220, 30, 30, 255)).save(component)
    presentation = Presentation()
    presentation.slide_width = Inches(10)
    presentation.slide_height = Inches(5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_picture(str(background), 0, 0, Inches(10), Inches(5))
    slide.shapes.add_picture(
        str(component), Inches(2), Inches(2), Inches(2), Inches(1.5)
    )
    text = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(3), Inches(1))
    text.text = "editable donor text"
    donor = tmp_path / "donor.pptx"
    presentation.save(donor)
    return donor


def _move_first_picture_to_background(source: Path) -> None:
    p = f"{{{P}}}"
    a = f"{{{A}}}"

    def transform(contents: bytes) -> bytes:
        root = ET.fromstring(contents)
        common = root.find(f"{p}cSld")
        tree = common.find(f"{p}spTree")
        picture = tree.find(f"{p}pic")
        tree.remove(picture)
        background = ET.Element(f"{p}bg")
        properties = ET.SubElement(background, f"{p}bgPr")
        fill = picture.find(f"{p}blipFill")
        fill.tag = f"{a}blipFill"
        properties.append(fill)
        common.insert(0, background)
        return ET.tostring(root, encoding="utf-8", xml_declaration=True)

    _replace_member(source, "ppt/slides/slide1.xml", transform)


def _replace_member(path: Path, name: str, transform) -> None:
    temporary = path.with_suffix(".tmp")
    with zipfile.ZipFile(path) as old, zipfile.ZipFile(temporary, "w") as new:
        for info in old.infolist():
            contents = old.read(info.filename)
            if info.filename == name:
                contents = transform(contents)
            new.writestr(info, contents)
    temporary.replace(path)


def _members(path: Path) -> dict[str, bytes]:
    with zipfile.ZipFile(path) as archive:
        return {name: archive.read(name) for name in archive.namelist()}


def _object_by_name(inventory: dict, name: str) -> dict:
    return next(
        item for item in inventory["slides"][0]["objects"] if item["name"] == name
    )


def _resolve_slide_target(target: str) -> str:
    return posixpath.normpath(posixpath.join("ppt/slides", target))
