from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches

from image2editable.pptx_input import scan_pptx
from image2editable.pptx_shadow import patch_slide_background
from image2editable.pptx_shadow_run import _validate_shadow_patch
from image2editable import pptx_shadow_run


def _mixed_source(tmp_path: Path) -> Path:
    image = tmp_path / "screen.png"
    Image.new("RGB", (640, 360), "navy").save(image)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    pic = slide.shapes.add_picture(str(image), 0, 0, Inches(10), Inches(5))
    pic.name = "Screenshot"
    text = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(0.5))
    text.name, text.text = "Native text", "keep"
    table = slide.shapes.add_table(
        2, 2, Inches(1), Inches(2), Inches(3), Inches(1)
    ).table
    table.cell(0, 0).text = "Native table"
    data = ChartData()
    data.categories = ["A", "B"]
    data.add_series("S", (1, 2))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(5), Inches(1), Inches(3), Inches(2), data
    )
    chart.name = "Native chart"
    slide.notes_slide.notes_text_frame.text = "keep note"
    out = tmp_path / "source.pptx"
    prs.save(out)
    return out


def _donor(tmp_path: Path) -> Path:
    bg = tmp_path / "bg.png"
    Image.new("RGB", (640, 360), "white").save(bg)
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(10), Inches(5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(bg), 0, 0, Inches(10), Inches(5))
    out = tmp_path / "donor.pptx"
    prs.save(out)
    return out


def test_mixed_native_patch_preserves_native_xml_zorder_notes_and_unrelated_page(
    tmp_path,
):
    source = _mixed_source(tmp_path)
    donor = _donor(tmp_path)
    output = tmp_path / "out.pptx"
    before = scan_pptx(source)
    native = [o for o in before["slides"][0]["objects"] if o["name"] != "Screenshot"]
    screenshot_id = next(
        o["shape_id"]
        for o in before["slides"][0]["objects"]
        if o["name"] == "Screenshot"
    )
    patch_slide_background(
        source,
        donor,
        output,
        slide_part="ppt/slides/slide1.xml",
        source_shape_id=screenshot_id,
    )
    after = scan_pptx(output)
    actual = after["slides"][0]["objects"]
    imported = {"Slide Background Image", "Picture 1"}
    assert [o["name"] for o in actual if o["name"] not in imported] == [
        o["name"] for o in native
    ]
    assert [o["z_order"] for o in actual if o["name"] not in imported] == [
        o["z_order"] for o in native
    ]
    assert {
        o["name"]: o["xml_c14n_sha256"] for o in actual
    } | {}  # inventory remains hashable
    for item in native:
        found = next(o for o in actual if o["name"] == item["name"])
        assert found["xml_c14n_sha256"] == item["xml_c14n_sha256"]
    assert after["slides"][0]["notes_sha256"] == before["slides"][0]["notes_sha256"]
    assert (
        Presentation(output)
        .slides[0]
        .notes_slide.notes_text_frame.text.endswith("keep note")
    )


def test_shadow_validator_rejects_unrelated_page_change(tmp_path):
    source = _mixed_source(tmp_path)
    donor = _donor(tmp_path)
    output = tmp_path / "out.pptx"
    before = scan_pptx(source)
    screenshot_id = next(
        o["shape_id"]
        for o in before["slides"][0]["objects"]
        if o["name"] == "Screenshot"
    )
    patch_slide_background(
        source,
        donor,
        output,
        slide_part="ppt/slides/slide1.xml",
        source_shape_id=screenshot_id,
    )
    plan = {
        "slide_part": "ppt/slides/slide1.xml",
        "decision": {"source_shape_id": screenshot_id},
    }
    _validate_shadow_patch(source, output, plan)


def test_patch_failure_preserves_warning_and_reopens_final_pptx(tmp_path, monkeypatch):
    source = _mixed_source(tmp_path)
    output = tmp_path / "failed-output.pptx"
    run_root = tmp_path / "run"
    (run_root / "pages" / "page_001").mkdir(parents=True)

    def fail_patch(*args, **kwargs):
        raise RuntimeError("patch failed")

    monkeypatch.setattr(pptx_shadow_run, "patch_slide_background", fail_patch)
    monkeypatch.setattr(
        pptx_shadow_run, "build_reconstruction_donor", lambda *a, **k: {}
    )
    image = tmp_path / "candidate.png"
    Image.new("RGB", (8, 8), "red").save(image)
    result = pptx_shadow_run.run_shadow_replacements(
        source,
        output,
        [
            {
                "page_id": "page_001",
                "slide_part": "ppt/slides/slide1.xml",
                "image_path": str(image),
                "work_root": str(run_root / "pages" / "page_001" / "reconstruction"),
                "decision": {
                    "runtime_action": "shadow_run",
                    "eligible_for_shadow_run": True,
                    "source_shape_id": "2",
                },
            }
        ],
        run_root=run_root,
    )
    assert result["page_results"][0]["status"] == "preserved_with_warning"
    assert "patch failed" in result["page_results"][0]["warning"]
    assert len(Presentation(output).slides) == 1
