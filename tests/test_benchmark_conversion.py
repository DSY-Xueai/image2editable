import hashlib
import io
import json
import math
import os
import re
import subprocess
import sys
import warnings
import zipfile
import zlib
from dataclasses import replace
from pathlib import Path

import pytest
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches

import scripts.benchmark_conversion as benchmark


CASE_SPECS = tuple(
    (f"image_{index}", "image", f"image-{index}.png", 1) for index in range(8)
) + (
    ("document", "pdf", "document.pdf", 3),
    ("mixed", "pptx", "mixed.pptx", 3),
)


def _canonical_sha256(value: object) -> str:
    payload = json.dumps(
        value,
        ensure_ascii=False,
        sort_keys=True,
        separators=(",", ":"),
        allow_nan=False,
    ).encode("utf-8")
    return hashlib.sha256(payload).hexdigest()


def _valid_manifest(root: Path) -> dict[str, object]:
    cases = []
    for identifier, kind, filename, pages in CASE_SPECS:
        payload = f"asset:{identifier}".encode()
        (root / filename).write_bytes(payload)
        cases.append(
            {
                "id": identifier,
                "kind": kind,
                "path": filename,
                "pages": pages,
                "bytes": len(payload),
                "sha256": hashlib.sha256(payload).hexdigest(),
            }
        )
    routes = [
        {"id": "images", "cases": [case[0] for case in CASE_SPECS[:8]], "pages": 8},
        {"id": "pdf", "cases": ["document"], "pages": 3},
        {"id": "mixed_pptx", "cases": ["mixed"], "pages": 3},
    ]
    manifest: dict[str, object] = {
        "schema_version": 1,
        "cases": cases,
        "routes": routes,
    }
    manifest["corpus_sha256"] = _canonical_sha256(manifest)
    return manifest


def _write_manifest(
    root: Path,
    manifest: dict[str, object],
    *,
    refresh_corpus_sha: bool = True,
) -> Path:
    if refresh_corpus_sha:
        manifest["corpus_sha256"] = _canonical_sha256(
            {
                "schema_version": manifest.get("schema_version"),
                "cases": manifest.get("cases"),
                "routes": manifest.get("routes"),
            }
        )
    path = root / "manifest.json"
    path.write_text(json.dumps(manifest), encoding="utf-8")
    return path


def _assert_invalid(path: Path) -> None:
    with pytest.raises(benchmark.BenchmarkError) as raised:
        benchmark.load_manifest(path)
    assert str(raised.value) == "invalid_corpus"


def test_module_exposes_benchmark_error() -> None:
    assert issubclass(benchmark.BenchmarkError, RuntimeError)


def _assert_benchmark_readme_contract(readme: str) -> None:
    visible = re.sub(r"<!--.*?-->", "", readme, flags=re.DOTALL)
    expected_headings = [
        "# 转换基准",
        "## 公开语料",
        "## 环境前置",
        "## 运行",
        "## 通过标准",
        "## 安全报告",
        "## 结果解释",
    ]
    headings = [line for line in visible.splitlines() if line.startswith("#")]
    assert headings == expected_headings
    parts = re.split(r"(?m)^## ([^\n]+)\n", visible)
    assert parts[0].strip() == "# 转换基准"
    assert parts[1::2] == [heading.removeprefix("## ") for heading in expected_headings[1:]]
    sections = dict(zip(parts[1::2], parts[2::2]))

    doctor = "image2editable doctor --agent-local"
    runner = (
        "python scripts/benchmark_conversion.py --corpus benchmark/corpus "
        "--output-dir benchmark-results"
    )
    environment_commands = [
        block.strip()
        for block in re.findall(
            r"(?ms)^```bash\n(.*?)^```$", sections["环境前置"]
        )
    ]
    run_commands = [
        block.strip()
        for block in re.findall(r"(?ms)^```bash\n(.*?)^```$", sections["运行"])
    ]
    assert environment_commands == [doctor]
    assert run_commands == [runner]
    assert visible.index(doctor) < visible.index(runner)

    assert "10 个输入、14 页、3 条 routes" in sections["公开语料"]
    assert "8 张图片、3 页 PDF、3 页 mixed PPTX" in sections["公开语料"]
    assert "维护者用来验证转换质量" in sections["环境前置"]
    assert "benchmark/private/" not in visible
    assert "只有 `ready=true` 才运行真实 benchmark" in sections["环境前置"]
    assert "runner 不会自动下载模型" in sections["环境前置"]
    assert "输出目录必须尚不存在" in sections["运行"]

    assert (
        "`passed` 必须同时满足：3 routes、14 pages、0 failed_routes、"
        "0 warning_pages" in sections["通过标准"]
    )
    assert "所有必须重建页都通过可编辑结构门禁" in sections["通过标准"]
    assert (
        "`preserved_with_warning`、缺页、损坏输出、整页单图或不可见组件绕过"
        "均判定失败" in sections["通过标准"]
    )

    assert (
        "不包含任何绝对路径、URL、密钥、stderr 或异常正文"
        in sections["安全报告"]
    )
    assert "耗时只是本机事实，不代表其他机器或输入" in sections["结果解释"]


def test_benchmark_readme_documents_strict_safe_execution_contract() -> None:
    root = Path(__file__).resolve().parents[1]
    readme = (root / "benchmark" / "README.md").read_text(encoding="utf-8")
    _assert_benchmark_readme_contract(readme)


def test_benchmark_readme_contract_rejects_html_comment_content() -> None:
    root = Path(__file__).resolve().parents[1]
    readme = (root / "benchmark" / "README.md").read_text(encoding="utf-8")

    with pytest.raises(AssertionError):
        _assert_benchmark_readme_contract(f"<!--\n{readme}\n-->")


def test_benchmark_readme_contract_rejects_misplaced_commands() -> None:
    root = Path(__file__).resolve().parents[1]
    readme = (root / "benchmark" / "README.md").read_text(encoding="utf-8")

    doctor_block = "```bash\nimage2editable doctor --agent-local\n```"
    runner_block = (
        "```bash\npython scripts/benchmark_conversion.py --corpus benchmark/corpus "
        "--output-dir benchmark-results\n```"
    )
    misplaced_commands = readme.replace(doctor_block, "").replace(runner_block, "")
    misplaced_commands = misplaced_commands.replace(
        "## 公开语料\n",
        "## 公开语料\n\n"
        "```bash\n"
        "image2editable doctor --agent-local\n"
        "python scripts/benchmark_conversion.py --corpus benchmark/corpus "
        "--output-dir benchmark-results\n"
        "```\n",
        1,
    )
    with pytest.raises(AssertionError):
        _assert_benchmark_readme_contract(misplaced_commands)


def test_loads_tracked_manifest_into_three_ordered_routes() -> None:
    root = Path(__file__).resolve().parents[1] / "benchmark" / "corpus"

    manifest = benchmark.load_manifest(root / "manifest.json")

    assert manifest.root == root.resolve()
    assert manifest.corpus_sha256 == "5aec7fb2cf751e42fe0c1c51c49bac613af85fae1727ba66de517a7a0c1718a0"
    assert [route.identifier for route in manifest.routes] == [
        "images",
        "pdf",
        "mixed_pptx",
    ]
    assert [route.expected_pages for route in manifest.routes] == [8, 3, 3]
    assert sum(route.expected_pages for route in manifest.routes) == 14
    assert [path.name for path in manifest.routes[0].sources] == [
        "01-zh-courseware.png",
        "02-typography.png",
        "03-flowchart.png",
        "04-table-chart.png",
        "05-photo-overlay.png",
        "06-transparency-shadow.png",
        "07-compressed.jpg",
        "08-portrait.png",
    ]


def test_loads_tracked_manifest_from_documented_corpus_directory() -> None:
    root = Path(__file__).resolve().parents[1] / "benchmark" / "corpus"

    manifest = benchmark.load_manifest(root)

    assert manifest.root == root.resolve()
    assert [route.identifier for route in manifest.routes] == [
        "images",
        "pdf",
        "mixed_pptx",
    ]


@pytest.mark.parametrize(
    "mutation",
    [
        lambda manifest: manifest.update(extra=True),
        lambda manifest: manifest.update(schema_version=2),
        lambda manifest: manifest.update(schema_version=True),
        lambda manifest: manifest["cases"][0].update(extra=True),
        lambda manifest: manifest["cases"][0].update(id=""),
        lambda manifest: manifest["cases"][0].update(kind="video"),
        lambda manifest: manifest["cases"][0].update(pages=True),
        lambda manifest: manifest["cases"][0].update(pages=0),
        lambda manifest: manifest["cases"][0].update(bytes=True),
        lambda manifest: manifest["cases"][0].update(bytes=0),
        lambda manifest: manifest["cases"][0].update(sha256="A" * 64),
    ],
)
def test_rejects_invalid_exact_schema_and_strict_types(
    tmp_path: Path, mutation: object
) -> None:
    manifest = _valid_manifest(tmp_path)
    mutation(manifest)

    _assert_invalid(_write_manifest(tmp_path, manifest))


@pytest.mark.parametrize(
    "unsafe_path",
    [
        "../secret.png",
        "/absolute.png",
        "C:/absolute.png",
        "\\\\server\\share\\x.png",
        ".",
        "",
        "nested/file.png",
        "nested\\file.png",
    ],
)
def test_rejects_unsafe_asset_paths(tmp_path: Path, unsafe_path: str) -> None:
    manifest = _valid_manifest(tmp_path)
    manifest["cases"][0]["path"] = unsafe_path

    _assert_invalid(_write_manifest(tmp_path, manifest))


@pytest.mark.parametrize(
    "unsafe_path",
    [
        "bad<name.png",
        "bad>name.png",
        'bad"name.png',
        "bad:name.png",
        "bad|name.png",
        "bad?name.png",
        "bad*name.png",
        "bad\x00name.png",
        "bad\x1fname.png",
        "bad\x7fname.png",
        "trailing.",
        "trailing ",
        "CON",
        "CONIN$",
        "CONOUT$",
        "CON .txt",
        "prn.txt",
        "AUX",
        "nul.png",
        "NUL .txt",
        "COM1",
        "COM1 .txt",
        "com9.log",
        "COM¹",
        "com².txt",
        "LPT1",
        "lpt9.txt",
        "LPT³.log",
    ],
)
def test_rejects_windows_unsafe_filename_before_reading_assets(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    unsafe_path: str,
) -> None:
    manifest = _valid_manifest(tmp_path)
    manifest["cases"][0]["path"] = unsafe_path
    manifest_path = _write_manifest(tmp_path, manifest)
    asset_reads: list[Path] = []
    original_read = benchmark._read_regular_file

    def tracked_read(path: Path, *args: object, **kwargs: object) -> bytes:
        if path != manifest_path:
            asset_reads.append(path)
        return original_read(path, *args, **kwargs)

    monkeypatch.setattr(benchmark, "_read_regular_file", tracked_read)

    _assert_invalid(manifest_path)
    assert asset_reads == []


@pytest.mark.parametrize("field", ["id", "path"])
def test_rejects_duplicate_case_identity(tmp_path: Path, field: str) -> None:
    manifest = _valid_manifest(tmp_path)
    manifest["cases"][1][field] = manifest["cases"][0][field]

    _assert_invalid(_write_manifest(tmp_path, manifest))


def test_rejects_casefolded_duplicate_path_before_reading_assets(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    manifest = _valid_manifest(tmp_path)
    first, second = manifest["cases"][:2]
    second["path"] = first["path"].upper()
    second["bytes"] = first["bytes"]
    second["sha256"] = first["sha256"]
    manifest_path = _write_manifest(tmp_path, manifest)
    asset_reads: list[Path] = []
    original_read = benchmark._read_regular_file

    def tracked_read(path: Path, *args: object, **kwargs: object) -> bytes:
        if path != manifest_path:
            asset_reads.append(path)
        return original_read(path, *args, **kwargs)

    monkeypatch.setattr(benchmark, "_read_regular_file", tracked_read)

    _assert_invalid(manifest_path)
    assert asset_reads == []


@pytest.mark.skipif(os.name != "nt", reason="requires Windows filesystem semantics")
def test_rejects_real_windows_case_alias(tmp_path: Path) -> None:
    manifest = _valid_manifest(tmp_path)
    first, second = manifest["cases"][:2]
    alias = first["path"].upper()
    if not (tmp_path / alias).exists():
        pytest.skip("filesystem is case-sensitive")
    second["path"] = alias
    second["bytes"] = first["bytes"]
    second["sha256"] = first["sha256"]

    _assert_invalid(_write_manifest(tmp_path, manifest))


@pytest.mark.skipif(os.name != "nt", reason="requires NTFS alternate data streams")
def test_rejects_real_windows_alternate_data_stream(tmp_path: Path) -> None:
    manifest = _valid_manifest(tmp_path)
    source = tmp_path / manifest["cases"][0]["path"]
    stream_name = f"{source.name}:asset"
    stream_path = tmp_path / stream_name
    payload = b"alternate-data"
    try:
        stream_path.write_bytes(payload)
    except OSError as error:
        pytest.skip(f"alternate data streams unavailable: {error.__class__.__name__}")
    second = manifest["cases"][1]
    second["path"] = stream_name
    second["bytes"] = len(payload)
    second["sha256"] = hashlib.sha256(payload).hexdigest()

    _assert_invalid(_write_manifest(tmp_path, manifest))


@pytest.mark.parametrize("problem", ["size", "sha256"])
def test_rejects_asset_size_or_hash_mismatch(tmp_path: Path, problem: str) -> None:
    manifest = _valid_manifest(tmp_path)
    if problem == "size":
        manifest["cases"][0]["bytes"] += 1
    else:
        manifest["cases"][0]["sha256"] = "0" * 64

    _assert_invalid(_write_manifest(tmp_path, manifest))


def test_rejects_asset_symlink_without_leaking_target_path(tmp_path: Path) -> None:
    manifest = _valid_manifest(tmp_path)
    target = tmp_path / "private-target"
    target.write_bytes(b"secret")
    source = tmp_path / manifest["cases"][0]["path"]
    source.unlink()
    try:
        source.symlink_to(target)
    except OSError as error:
        pytest.skip(f"symlink unavailable: {error.__class__.__name__}")

    _assert_invalid(_write_manifest(tmp_path, manifest))


def test_rejects_multiply_linked_asset(tmp_path: Path) -> None:
    manifest = _valid_manifest(tmp_path)
    source = tmp_path / manifest["cases"][0]["path"]
    os.link(source, tmp_path / "second-name.png")

    _assert_invalid(_write_manifest(tmp_path, manifest))


def test_file_identity_changes_when_link_count_changes(tmp_path: Path) -> None:
    source = tmp_path / "asset.png"
    source.write_bytes(b"asset")
    before = source.lstat()

    os.link(source, tmp_path / "second-name.png")

    assert benchmark._identity(source.lstat()) != benchmark._identity(before)


def test_rejects_manifest_symlink(tmp_path: Path) -> None:
    manifest = _valid_manifest(tmp_path)
    real_manifest = _write_manifest(tmp_path, manifest)
    linked_manifest = tmp_path / "linked.json"
    try:
        linked_manifest.symlink_to(real_manifest)
    except OSError as error:
        pytest.skip(f"symlink unavailable: {error.__class__.__name__}")
    _assert_invalid(linked_manifest)


def test_rejects_oversize_manifest(tmp_path: Path) -> None:
    manifest = _valid_manifest(tmp_path)
    real_manifest = _write_manifest(tmp_path, manifest)
    real_manifest.write_bytes(b" " * (64 * 1024 + 1))

    _assert_invalid(real_manifest)


@pytest.mark.parametrize(
    "mutation",
    [
        lambda manifest: manifest["routes"][0].update(extra=True),
        lambda manifest: manifest["routes"][0].update(id="pdf"),
        lambda manifest: manifest["routes"][0].update(pages=True),
        lambda manifest: manifest["routes"][0].update(pages=7),
        lambda manifest: manifest["routes"][0]["cases"].__setitem__(0, "missing"),
        lambda manifest: manifest["routes"][1]["cases"].append("image_0"),
        lambda manifest: manifest["routes"].reverse(),
    ],
)
def test_rejects_invalid_route_contract(tmp_path: Path, mutation: object) -> None:
    manifest = _valid_manifest(tmp_path)
    mutation(manifest)

    _assert_invalid(_write_manifest(tmp_path, manifest))


def test_rejects_wrong_corpus_sha256(tmp_path: Path) -> None:
    manifest = _valid_manifest(tmp_path)
    manifest["corpus_sha256"] = "0" * 64

    _assert_invalid(_write_manifest(tmp_path, manifest, refresh_corpus_sha=False))


def test_rejects_more_than_twelve_mib_of_assets(tmp_path: Path) -> None:
    manifest = _valid_manifest(tmp_path)
    payload = b"x" * (12 * 1024 * 1024)
    first = manifest["cases"][0]
    (tmp_path / first["path"]).write_bytes(payload)
    first["bytes"] = len(payload)
    first["sha256"] = hashlib.sha256(payload).hexdigest()

    _assert_invalid(_write_manifest(tmp_path, manifest))


def _completed_process(
    command: list[str],
    *,
    returncode: int = 0,
    stdout: str = "",
    stderr: str = "",
) -> subprocess.CompletedProcess[str]:
    return subprocess.CompletedProcess(
        args=command,
        returncode=returncode,
        stdout=stdout,
        stderr=stderr,
    )


def _loaded_manifest(tmp_path: Path) -> benchmark.CorpusManifest:
    corpus = tmp_path / "corpus"
    corpus.mkdir()
    return benchmark.load_manifest(_write_manifest(corpus, _valid_manifest(corpus)))


def test_run_cli_uses_isolated_interpreter_and_clean_environment(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    calls: list[tuple[list[str], dict[str, object]]] = []
    monkeypatch.setenv("PYTHONPATH", "untrusted")

    def fake_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        calls.append((command, kwargs))
        return _completed_process(command, stdout="ok")

    monkeypatch.setattr(subprocess, "run", fake_run)

    completed = benchmark._run_cli(
        ["doctor", "--agent-local"], cwd=tmp_path, timeout=180.0
    )

    assert completed.stdout == "ok"
    assert len(calls) == 1
    command, options = calls[0]
    assert command[:5] == [sys.executable, "-I", "-B", "-m", "image2editable"]
    assert command[5:] == ["doctor", "--agent-local"]
    assert options["cwd"] == tmp_path
    assert "PYTHONPATH" not in options["env"]
    assert options["capture_output"] is True
    assert options["text"] is True
    assert options["check"] is False
    assert options["timeout"] == 180.0


@pytest.mark.parametrize(
    ("doctor_result", "raises_timeout"),
    [
        ({"returncode": 4, "stdout": '{"ready":true,"checks":{}}'}, False),
        ({"returncode": 0, "stdout": '{"ready":false,"checks":{}}'}, False),
        ({"returncode": 0, "stdout": "not-json"}, False),
        ({"returncode": 0, "stdout": "{}\n{}"}, False),
        ({"returncode": 0, "stdout": '{"ready":true,"checks":[]}'}, False),
        (
            {
                "returncode": 0,
                "stdout": '{"ready":true,"checks":{},"extra":true}',
            },
            False,
        ),
        (
            {
                "returncode": 0,
                "stdout": '{"ready":false,"ready":true,"checks":{}}',
            },
            False,
        ),
        (
            {
                "returncode": 0,
                "stdout": '{"ready":true,"checks":{"probe":NaN}}',
            },
            False,
        ),
        (
            {
                "returncode": 0,
                "stdout": '{"ready":true,"checks":{"probe":Infinity}}',
            },
            False,
        ),
        (
            {
                "returncode": 0,
                "stdout": '{"ready":true,"checks":{"probe":-Infinity}}',
            },
            False,
        ),
        ({}, True),
    ],
)
def test_doctor_failure_prevents_all_conversions(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    doctor_result: dict[str, object],
    raises_timeout: bool,
) -> None:
    manifest = _loaded_manifest(tmp_path)
    result_root = tmp_path / "results"
    commands: list[list[str]] = []

    def fake_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        commands.append(command)
        if raises_timeout:
            raise subprocess.TimeoutExpired(command, kwargs["timeout"])
        return _completed_process(command, **doctor_result)

    monkeypatch.setattr(subprocess, "run", fake_run)

    with pytest.raises(benchmark.BenchmarkError) as raised:
        benchmark.execute_routes(manifest, result_root)

    assert str(raised.value) == "doctor_not_ready"
    assert [command[5:] for command in commands] == [["doctor", "--agent-local"]]


def test_execute_routes_calls_doctor_once_then_three_exact_conversions(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    manifest = _loaded_manifest(tmp_path)
    result_root = tmp_path / "results"
    calls: list[tuple[list[str], dict[str, object]]] = []
    monkeypatch.setenv("PYTHONPATH", "untrusted")

    def fake_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        calls.append((command, kwargs))
        if command[5] == "doctor":
            return _completed_process(
                command, stdout='{"ready":true,"checks":{"agent_local":"ok"}}\n'
            )
        run_path = Path(command[command.index("--run-dir") + 1])
        output_path = Path(command[command.index("--output") + 1])
        assert not run_path.exists()
        assert not output_path.exists()
        return _completed_process(command, stdout=f"converted:{command[6]}")

    monkeypatch.setattr(subprocess, "run", fake_run)

    executions = benchmark.execute_routes(manifest, result_root)

    assert result_root.is_dir()
    assert len(calls) == 4
    assert [command[5] for command, _ in calls] == [
        "doctor",
        "convert",
        "convert",
        "convert",
    ]
    assert calls[0][0][5:] == ["doctor", "--agent-local"]
    assert calls[0][1]["timeout"] == 180.0
    for index, (route, execution) in enumerate(zip(manifest.routes, executions), start=1):
        command, options = calls[index]
        run_path = result_root.resolve() / route.identifier / "run"
        output_path = result_root.resolve() / route.identifier / "output.pptx"
        assert command == [
            sys.executable,
            "-I",
            "-B",
            "-m",
            "image2editable",
            "convert",
            *(str(source) for source in route.sources),
            "--run-dir",
            str(run_path),
            "--output",
            str(output_path),
            "--slide-size",
            "16:9",
            "--agent-provider",
            "local",
        ]
        assert options["cwd"] == result_root.resolve()
        assert options["timeout"] is None
        assert "PYTHONPATH" not in options["env"]
        assert execution.route is route
        assert execution.run_root == run_path
        assert execution.output_path == output_path
        assert execution.returncode == 0
        assert execution.stdout.startswith("converted:")
        assert math.isfinite(execution.duration_seconds)
        assert execution.duration_seconds >= 0


def test_nonzero_conversion_is_recorded_and_does_not_stop_later_routes(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    manifest = _loaded_manifest(tmp_path)
    result_root = tmp_path / "results"
    calls: list[list[str]] = []

    def fake_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        calls.append(command)
        if command[5] == "doctor":
            return _completed_process(command, stdout='{"ready":true,"checks":{}}')
        if len(calls) == 2:
            return _completed_process(
                command, returncode=9, stdout="failed-output", stderr="secret-stderr"
            )
        return _completed_process(command, stdout="later-output")

    monkeypatch.setattr(subprocess, "run", fake_run)

    executions = benchmark.execute_routes(manifest, result_root)

    assert [execution.route.identifier for execution in executions] == [
        "images",
        "pdf",
        "mixed_pptx",
    ]
    assert [execution.returncode for execution in executions] == [9, 0, 0]
    assert [execution.stdout for execution in executions] == [
        "failed-output",
        "later-output",
        "later-output",
    ]
    assert not hasattr(executions[0], "stderr")
    assert len(calls) == 4


def test_nonfinite_clock_duration_is_normalized(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    manifest = _loaded_manifest(tmp_path)
    clock_values = iter([0.0, math.inf, 1.0, 2.0, 3.0, 4.0])

    def fake_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        if command[5] == "doctor":
            return _completed_process(command, stdout='{"ready":true,"checks":{}}')
        return _completed_process(command)

    monkeypatch.setattr(subprocess, "run", fake_run)
    monkeypatch.setattr(benchmark.time, "perf_counter", lambda: next(clock_values))

    executions = benchmark.execute_routes(manifest, tmp_path / "results")

    assert executions[0].duration_seconds == 0.0
    assert all(
        math.isfinite(execution.duration_seconds) for execution in executions
    )


def test_conversion_exception_is_safely_recorded_and_later_routes_continue(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    manifest = _loaded_manifest(tmp_path)
    calls: list[list[str]] = []

    def fake_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        calls.append(command)
        if command[5] == "doctor":
            return _completed_process(command, stdout='{"ready":true,"checks":{}}')
        if len(calls) == 2:
            raise OSError("sk-secret failed at C:\\private\\input")
        return _completed_process(command, returncode=8, stdout="later failed")

    monkeypatch.setattr(subprocess, "run", fake_run)

    executions = benchmark.execute_routes(manifest, tmp_path / "results")

    assert [command[5] for command in calls] == [
        "doctor",
        "convert",
        "convert",
        "convert",
    ]
    assert [execution.returncode for execution in executions] == [-1, 8, 8]
    assert executions[0].stdout == ""


def test_existing_result_root_fails_before_subprocess(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    manifest = _loaded_manifest(tmp_path)
    result_root = tmp_path / "results"
    result_root.mkdir()

    def unexpected_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        raise AssertionError("subprocess must not run")

    monkeypatch.setattr(subprocess, "run", unexpected_run)

    with pytest.raises(benchmark.BenchmarkError) as raised:
        benchmark.execute_routes(manifest, result_root)

    assert str(raised.value) == "invalid_output"


def test_execute_revalidates_assets_before_output_or_subprocess(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    manifest = _loaded_manifest(tmp_path)
    source = manifest.routes[0].sources[0]
    original = source.read_bytes()
    replacement = b"x" * len(original)
    assert replacement != original
    source.write_bytes(replacement)
    result_root = tmp_path / "results"
    calls: list[list[str]] = []

    def fake_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        calls.append(command)
        if command[5] == "doctor":
            return _completed_process(command, stdout='{"ready":true,"checks":{}}')
        return _completed_process(command)

    monkeypatch.setattr(subprocess, "run", fake_run)

    with pytest.raises(benchmark.BenchmarkError) as raised:
        benchmark.execute_routes(manifest, result_root)

    assert str(raised.value) == "invalid_corpus"
    assert calls == []
    assert not result_root.exists()


def test_dangling_result_entry_fails_before_subprocess(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    manifest = _loaded_manifest(tmp_path)
    result_root = tmp_path / "dangling-results"
    try:
        result_root.symlink_to(tmp_path / "missing-target", target_is_directory=True)
    except OSError as error:
        pytest.skip(f"directory symlink unavailable: {error.__class__.__name__}")

    def unexpected_run(
        command: list[str], **kwargs: object
    ) -> subprocess.CompletedProcess[str]:
        raise AssertionError("subprocess must not run")

    monkeypatch.setattr(subprocess, "run", unexpected_run)

    with pytest.raises(benchmark.BenchmarkError) as raised:
        benchmark.execute_routes(manifest, result_root)

    assert str(raised.value) == "invalid_output"


PERFORMANCE_FIELDS = {
    "model_loads",
    "stage_runs",
    "stage_duration_ms",
    "worker_runs",
    "worker_duration_ms",
    "inference_runs",
    "inference_operations",
    "inference_duration_ms",
    "agent_runs",
    "agent_image_count",
    "agent_total_bytes",
    "agent_duration_ms",
}


def _performance(pages: int) -> dict[str, object]:
    values = {}
    for index in range(1, pages + 1):
        page = {field: {} for field in PERFORMANCE_FIELDS}
        page.update(
            agent_runs=1,
            agent_image_count=1,
            agent_total_bytes=10,
            agent_duration_ms=2,
        )
        page["model_loads"] = {"layout-model": 1}
        page["stage_runs"] = {"reconstruct": 1}
        page["stage_duration_ms"] = {"reconstruct": 2}
        page["worker_runs"] = {"layout-model": 1}
        page["worker_duration_ms"] = {"layout-model": 2}
        page["inference_runs"] = {"layout-model": 1}
        page["inference_operations"] = {"layout-model": 1}
        page["inference_duration_ms"] = {"layout-model": 2}
        values[f"page_{index:03d}"] = page
    return {"pages": values}


def _reported_performance(pages: int) -> dict[str, object]:
    performance = _performance(pages)
    for page in performance["pages"].values():
        for field in PERFORMANCE_FIELDS - {
            "agent_runs",
            "agent_image_count",
            "agent_total_bytes",
            "agent_duration_ms",
        }:
            page[field] = sum(page[field].values())
    return performance


def _write_editable_pptx(
    path: Path,
    pages: int,
    *,
    mode: str = "text",
) -> bytes:
    presentation = Presentation()
    presentation.slide_height = Inches(7.5)
    presentation.slide_width = (
        Inches(10)
        if mode == "aspect_4_3"
        else Inches(15)
        if mode == "aspect_2_1"
        else presentation.slide_height * 16 // 9
    )
    image_path = path.with_suffix(".png")
    transparent_path = path.with_name(f"{path.stem}-transparent.png")
    sparse_path = path.with_name(f"{path.stem}-sparse.png")
    Image.new("RGB", (40, 40), "navy").save(image_path)
    Image.new("RGBA", (40, 40), (0, 0, 0, 0)).save(transparent_path)
    sparse_image = Image.new("RGBA", (40, 40), (0, 0, 0, 0))
    sparse_image.putpixel((0, 0), (0, 0, 0, 255))
    sparse_image.save(sparse_path)
    for index in range(pages):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        page_mode = mode
        if mode in {"aspect_4_3", "aspect_2_1"}:
            page_mode = "text"
        if mode == "mixed_valid":
            page_mode = "native" if index == 0 else "text"
        if page_mode == "text":
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            textbox.text = "editable"
        elif page_mode == "native":
            slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
            )
        elif page_mode == "partial_gradient_native":
            slide.shapes.add_picture(
                str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
            )
            native = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(2)
            )
            native.fill.solid()
            native.line.fill.background()
            solid_fill = next(
                child
                for child in native._element.spPr
                if child.tag.endswith("}solidFill")
            )
            gradient_fill = OxmlElement("a:gradFill")
            stops = OxmlElement("a:gsLst")
            transparent_stop = OxmlElement("a:gs")
            transparent_stop.set("pos", "0")
            transparent_color = OxmlElement("a:srgbClr")
            transparent_color.set("val", "FF0000")
            alpha = OxmlElement("a:alpha")
            alpha.set("val", "0")
            transparent_color.append(alpha)
            transparent_stop.append(transparent_color)
            opaque_stop = OxmlElement("a:gs")
            opaque_stop.set("pos", "100000")
            opaque_color = OxmlElement("a:srgbClr")
            opaque_color.set("val", "FF0000")
            opaque_stop.append(opaque_color)
            stops.extend((transparent_stop, opaque_stop))
            gradient_fill.append(stops)
            solid_fill.getparent().replace(solid_fill, gradient_fill)
        elif page_mode == "full_picture":
            slide.shapes.add_picture(
                str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
            )
        elif page_mode == "small_picture":
            slide.shapes.add_picture(
                str(image_path), Inches(1), Inches(1), Inches(2), Inches(2)
            )
        elif page_mode == "pictures":
            slide.shapes.add_picture(
                str(image_path), Inches(1), Inches(1), Inches(2), Inches(2)
            )
            slide.shapes.add_picture(
                str(image_path), Inches(4), Inches(1), Inches(2), Inches(2)
            )
        elif page_mode in {"partial_picture_alpha", "partial_picture_alpha_repl"}:
            slide.shapes.add_picture(
                str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
            )
            picture = slide.shapes.add_picture(
                str(image_path), Inches(1), Inches(1), Inches(2), Inches(2)
            )
            alpha = OxmlElement(
                "a:alphaRepl"
                if page_mode == "partial_picture_alpha_repl"
                else "a:alphaModFix"
            )
            alpha.set("a" if page_mode == "partial_picture_alpha_repl" else "amt", "50000")
            picture._element.blipFill.blip.append(alpha)
        elif page_mode in {
            "group_text",
            "group_offslide",
            "group_zero_width",
            "group_zero_height",
        }:
            if page_mode != "group_text":
                slide.shapes.add_picture(
                    str(image_path),
                    0,
                    0,
                    presentation.slide_width,
                    presentation.slide_height,
                )
            group = slide.shapes.add_group_shape()
            group.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1)).text = (
                "nested editable"
            )
            if page_mode == "group_offslide":
                group.left = presentation.slide_width + Inches(1)
            elif page_mode == "group_zero_width":
                group.width = 0
            elif page_mode == "group_zero_height":
                group.height = 0
        elif page_mode in {
            "group_flip_h_offslide",
            "group_flip_v_offslide",
            "group_flip_h_visible",
            "nested_group_flip_h_offslide",
        }:
            slide.shapes.add_picture(
                str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
            )
            group = slide.shapes.add_group_shape()
            text_shapes = group.shapes
            if page_mode == "nested_group_flip_h_offslide":
                text_shapes = group.shapes.add_group_shape().shapes
            text_shapes.add_textbox(0, 0, Inches(1), Inches(0.5)).text = "nested editable"
            anchor = text_shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(3), Inches(1.5), Inches(1), Inches(0.5)
            )
            next(
                node
                for node in anchor._element.iter()
                if node.tag.endswith("}cNvPr")
            ).set("hidden", "1")
            group.width = Inches(4)
            group.height = Inches(2)
            if page_mode == "group_flip_v_offslide":
                group.top = presentation.slide_height - Inches(0.5)
                group._element.grpSpPr.xfrm.flipV = True
            else:
                group.left = (
                    Inches(1)
                    if page_mode == "group_flip_h_visible"
                    else presentation.slide_width - Inches(1)
                )
                group._element.grpSpPr.xfrm.flipH = True
        elif page_mode in {"empty_text", "group_empty_text", "invisible_shape"}:
            slide.shapes.add_picture(
                str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
            )
            shapes = slide.shapes
            if page_mode == "group_empty_text":
                shapes = slide.shapes.add_group_shape().shapes
            shape = shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(1))
            if page_mode == "invisible_shape":
                shape = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(2), Inches(1)
                )
                shape.fill.background()
                shape.line.fill.background()
        elif page_mode in {
            "outside_text",
            "zero_text",
            "outside_picture",
            "zero_picture",
        }:
            slide.shapes.add_picture(
                str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
            )
            if page_mode.endswith("text"):
                textbox = slide.shapes.add_textbox(
                    presentation.slide_width + Inches(1)
                    if page_mode == "outside_text"
                    else Inches(1),
                    Inches(1),
                    Inches(2) if page_mode == "outside_text" else 0,
                    Inches(1),
                )
                textbox.text = "editable but invisible"
            else:
                picture = slide.shapes.add_picture(
                    str(image_path), Inches(1), Inches(1), Inches(2), Inches(2)
                )
                if page_mode == "outside_picture":
                    picture.left = presentation.slide_width + Inches(1)
                else:
                    picture.width = 0
        elif page_mode in {
            "tiny_text",
            "tiny_picture",
            "tiny_native",
            "hidden_text",
            "hidden_picture",
            "group_hidden_text",
            "transparent_text",
            "transparent_picture",
            "shape_transparent_picture",
            "shape_alpha_repl_transparent_picture",
            "shape_unknown_alpha_picture",
            "shape_invalid_alpha_picture",
            "shape_extra_alpha_field_picture",
            "sparse_picture",
            "transparent_native",
            "transparent_native_leading_zero",
        }:
            slide.shapes.add_picture(
                str(image_path), 0, 0, presentation.slide_width, presentation.slide_height
            )
            if page_mode.endswith("text"):
                if page_mode == "group_hidden_text":
                    group = slide.shapes.add_group_shape()
                    textbox = group.shapes.add_textbox(
                        Inches(1), Inches(1), Inches(2), Inches(1)
                    )
                    hidden_shape = group
                else:
                    textbox = slide.shapes.add_textbox(
                        Inches(1),
                        Inches(1),
                        1 if page_mode == "tiny_text" else Inches(2),
                        1 if page_mode == "tiny_text" else Inches(1),
                    )
                    hidden_shape = textbox
                textbox.text = "editable"
                if page_mode in {"hidden_text", "group_hidden_text"}:
                    next(
                        node
                        for node in hidden_shape._element.iter()
                        if node.tag.endswith("}cNvPr")
                    ).set("hidden", "1")
                elif page_mode == "transparent_text":
                    run = textbox.text_frame.paragraphs[0].runs[0]
                    properties = run._r.get_or_add_rPr()
                    fill = OxmlElement("a:solidFill")
                    color = OxmlElement("a:srgbClr")
                    color.set("val", "000000")
                    alpha = OxmlElement("a:alpha")
                    alpha.set("val", "0")
                    color.append(alpha)
                    fill.append(color)
                    properties.append(fill)
            elif page_mode in {
                "tiny_native",
                "transparent_native",
                "transparent_native_leading_zero",
            }:
                native = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE,
                    Inches(1),
                    Inches(1),
                    1 if page_mode == "tiny_native" else Inches(2),
                    1 if page_mode == "tiny_native" else Inches(2),
                )
                if page_mode.startswith("transparent_native"):
                    native.fill.solid()
                    native.fill.fore_color.rgb = RGBColor(0, 0, 0)
                    native.line.fill.solid()
                    native.line.fill.fore_color.rgb = RGBColor(0, 0, 0)
                    alpha_value = (
                        "00000"
                        if page_mode == "transparent_native_leading_zero"
                        else "0"
                    )
                    for fill in (
                        node
                        for node in native._element.iter()
                        if node.tag.endswith("}solidFill")
                    ):
                        color = next(iter(fill))
                        alpha = OxmlElement("a:alpha")
                        alpha.set("val", alpha_value)
                        color.append(alpha)
            else:
                picture = slide.shapes.add_picture(
                    str(
                        transparent_path
                        if page_mode == "transparent_picture"
                        else sparse_path
                        if page_mode == "sparse_picture"
                        else image_path
                    ),
                    Inches(1),
                    Inches(1),
                    1 if page_mode == "tiny_picture" else Inches(2),
                    1 if page_mode == "tiny_picture" else Inches(2),
                )
                if page_mode == "hidden_picture":
                    next(
                        node
                        for node in picture._element.iter()
                        if node.tag.endswith("}cNvPr")
                    ).set("hidden", "1")
                elif page_mode in {
                    "shape_transparent_picture",
                    "shape_alpha_repl_transparent_picture",
                    "shape_unknown_alpha_picture",
                    "shape_invalid_alpha_picture",
                    "shape_extra_alpha_field_picture",
                }:
                    if page_mode == "shape_unknown_alpha_picture":
                        alpha = OxmlElement("a:alphaBiLevel")
                        alpha.set("thresh", "50000")
                    elif page_mode == "shape_invalid_alpha_picture":
                        alpha = OxmlElement("a:alphaModFix")
                        alpha.set("amt", "NaN")
                    elif page_mode == "shape_extra_alpha_field_picture":
                        alpha = OxmlElement("a:alphaModFix")
                        alpha.set("amt", "50000")
                        alpha.set("unexpected", "1")
                    else:
                        alpha = OxmlElement(
                            "a:alphaRepl"
                            if page_mode == "shape_alpha_repl_transparent_picture"
                            else "a:alphaModFix"
                        )
                        alpha.set(
                            "a"
                            if page_mode == "shape_alpha_repl_transparent_picture"
                            else "amt",
                            "00000",
                        )
                    picture._element.blipFill.blip.append(alpha)
        else:
            raise AssertionError(page_mode)
    stream = io.BytesIO()
    presentation.save(stream)
    payload = stream.getvalue()
    path.write_bytes(payload)
    return payload


def _zip_bytes(names: list[str], payload: bytes = b"content") -> bytes:
    stream = io.BytesIO()
    with warnings.catch_warnings():
        warnings.simplefilter("ignore", UserWarning)
        with zipfile.ZipFile(stream, "w", zipfile.ZIP_DEFLATED) as archive:
            for name in names:
                archive.writestr(name, payload)
    return stream.getvalue()


def _patch_zip_central_sizes(
    payload: bytes, *, compressed: list[int], uncompressed: list[int]
) -> bytes:
    patched = bytearray(payload)
    position = 0
    for compressed_size, uncompressed_size in zip(compressed, uncompressed):
        position = patched.find(b"PK\x01\x02", position)
        assert position >= 0
        patched[position + 20 : position + 24] = compressed_size.to_bytes(4, "little")
        patched[position + 24 : position + 28] = uncompressed_size.to_bytes(4, "little")
        position += 4
    return bytes(patched)


def _patch_zip_central_encrypted(payload: bytes) -> bytes:
    patched = bytearray(payload)
    position = patched.find(b"PK\x01\x02")
    assert position >= 0
    flags = int.from_bytes(patched[position + 8 : position + 10], "little") | 1
    patched[position + 8 : position + 10] = flags.to_bytes(2, "little")
    return bytes(patched)


def _png_with_dimensions(width: int, height: int) -> bytes:
    stream = io.BytesIO()
    Image.new("RGBA", (1, 1), (0, 0, 0, 0)).save(stream, format="PNG")
    payload = bytearray(stream.getvalue())
    payload[16:20] = width.to_bytes(4, "big")
    payload[20:24] = height.to_bytes(4, "big")
    payload[29:33] = zlib.crc32(payload[12:29]).to_bytes(4, "big")
    return bytes(payload)


def _replace_first_pptx_media(path: Path, payload: bytes) -> None:
    with zipfile.ZipFile(io.BytesIO(path.read_bytes())) as source:
        members = [(member, source.read(member)) for member in source.infolist()]
    output = io.BytesIO()
    replaced = False
    with zipfile.ZipFile(output, "w") as target:
        for member, content in members:
            if not replaced and member.filename.startswith("ppt/media/"):
                content = payload
                replaced = True
            target.writestr(member, content)
    assert replaced
    path.write_bytes(output.getvalue())


def _execution(
    root: Path,
    route_id: str,
    pages: int,
    *,
    mode: str = "text",
    returncode: int = 0,
    stdout: str | None = None,
    statuses: list[str] | None = None,
) -> benchmark.RouteExecution:
    base = root / route_id
    run_root = base / "run"
    output_path = base / "output.pptx"
    run_root.mkdir(parents=True)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    payload = _write_editable_pptx(output_path, pages, mode=mode)
    digest = hashlib.sha256(payload).hexdigest()
    page_ids = [f"page_{index:03d}" for index in range(1, pages + 1)]
    summary: dict[str, object] = {
        "schema_version": 1,
        "status": "completed",
        "pages": pages,
        "outputs": {
            "pptx" if route_id == "mixed_pptx" else "16:9": str(output_path)
        },
        "performance": _performance(pages),
    }
    if route_id == "mixed_pptx":
        statuses = statuses or ["preserved", "replaced", "replaced"]
        summary.update(
            page_results=[
                {"schema_version": 1, "page_id": page_id, "status": status}
                for page_id, status in zip(page_ids, statuses)
            ],
            preserved_with_warning_pages=statuses.count("preserved_with_warning"),
            warnings=[],
            output_sha256=digest,
        )
    else:
        for page_id in page_ids:
            page_root = run_root / "pages" / page_id
            page_root.mkdir(parents=True)
            (page_root / "page_result.json").write_text(
                json.dumps(
                    {
                        "schema_version": 1,
                        "page_id": page_id,
                        "status": "validated",
                    }
                ),
                encoding="utf-8",
            )
    return benchmark.RouteExecution(
        route=benchmark.Route(route_id, (), pages),
        run_root=run_root,
        output_path=output_path,
        returncode=returncode,
        stdout=json.dumps(summary) if stdout is None else stdout,
        duration_seconds=0.125,
    )


def test_build_report_has_fixed_safe_exact_schema_and_totals(tmp_path: Path) -> None:
    executions = (
        _execution(tmp_path, "images", 8),
        _execution(tmp_path, "pdf", 3, mode="group_text"),
        _execution(tmp_path, "mixed_pptx", 3, mode="mixed_valid"),
    )
    manifest = benchmark.CorpusManifest(
        tmp_path, "a" * 64, (), tuple(execution.route for execution in executions)
    )

    report = benchmark.build_report(manifest, executions)
    report_path = tmp_path / "benchmark-report.json"
    benchmark.write_report(report_path, report)
    payload = report_path.read_bytes()
    stored = json.loads(payload)

    assert set(stored) == {
        "schema_version",
        "status",
        "corpus_sha256",
        "environment",
        "routes",
        "totals",
    }
    assert set(stored["environment"]) == {
        "python",
        "platform",
        "device_interface",
    }
    assert [set(route) for route in stored["routes"]] == [
        {
            "id",
            "kind",
            "input_count",
            "pages",
            "duration_ms",
            "status",
            "error_type",
            "warning_pages",
            "output_sha256",
            "performance",
        }
    ] * 3
    assert stored["status"] == "passed"
    assert [route["kind"] for route in stored["routes"]] == ["image", "pdf", "pptx"]
    assert [route["input_count"] for route in stored["routes"]] == [8, 1, 1]
    assert stored["routes"][0]["performance"] == _reported_performance(8)
    assert stored["totals"] == {
        "routes": 3,
        "inputs": 10,
        "pages": 14,
        "duration_ms": 375,
        "failed_routes": 0,
        "warning_pages": 0,
    }
    assert payload.endswith(b"\n")
    assert b"\r\n" not in payload
    assert str(tmp_path).encode() not in payload


@pytest.mark.parametrize(
    "stdout",
    [
        "not-json",
        "{}\n{}",
        '{"status":"completed","status":"completed"}',
        '{"value":NaN}',
        '{"value":Infinity}',
        " " * (4 * 1024 * 1024 + 1),
    ],
    ids=["syntax", "multiple", "duplicate", "nan", "infinity", "oversize"],
)
def test_zero_exit_rejects_noncanonical_summary_json(
    tmp_path: Path, stdout: str
) -> None:
    execution = _execution(tmp_path, "images", 8, stdout=stdout)

    result = benchmark.evaluate_execution(execution)

    assert result["status"] == "failed"
    assert result["error_type"] == "invalid_summary"


def test_nonzero_exit_precedes_invalid_summary_and_does_not_leak(tmp_path: Path) -> None:
    secret = "sk-secret conversion failed at C:\\private\\asset and https://secret"
    execution = _execution(
        tmp_path, "images", 8, returncode=9, stdout=secret
    )

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "conversion_failed"
    assert secret not in json.dumps(result)


@pytest.mark.parametrize("problem", ["missing", "corrupt", "hash", "pages"])
def test_output_integrity_failures_are_output_invalid(
    tmp_path: Path, problem: str
) -> None:
    execution = _execution(tmp_path, "images", 8)
    summary = json.loads(execution.stdout)
    if problem == "missing":
        execution.output_path.unlink()
    elif problem == "corrupt":
        execution.output_path.write_bytes(b"not a pptx")
        summary.pop("output_sha256", None)
    elif problem == "hash":
        summary["output_sha256"] = "0" * 64
    else:
        _write_editable_pptx(execution.output_path, 7)
        summary.pop("output_sha256", None)
    execution = replace(execution, stdout=json.dumps(summary))

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "output_invalid"


def test_pptx_archive_preflight_accepts_real_corpus_and_normal_output(
    tmp_path: Path,
) -> None:
    corpus = (
        Path(__file__).resolve().parents[1] / "benchmark" / "corpus" / "10-mixed.pptx"
    )
    generated = tmp_path / "normal.pptx"
    _write_editable_pptx(generated, 3)

    benchmark._validate_pptx_archive(corpus.read_bytes())
    benchmark._validate_pptx_archive(generated.read_bytes())


@pytest.mark.parametrize(
    "payload",
    [
        b"not-a-zip",
        _zip_bytes([f"member-{index}.xml" for index in range(2049)]),
        _patch_zip_central_sizes(
            _zip_bytes(["large.xml"]),
            compressed=[64 * 1024 * 1024 + 1],
            uncompressed=[64 * 1024 * 1024 + 1],
        ),
        _patch_zip_central_sizes(
            _zip_bytes([f"part-{index}.xml" for index in range(5)]),
            compressed=[60 * 1024 * 1024] * 5,
            uncompressed=[60 * 1024 * 1024] * 5,
        ),
        _zip_bytes(["ratio.xml"], b"0" * (2 * 1024 * 1024)),
        _zip_bytes(["duplicate.xml", "duplicate.xml"]),
        _patch_zip_central_encrypted(_zip_bytes(["encrypted.xml"])),
    ],
    ids=[
        "bad-zip",
        "member-count",
        "single-size",
        "total-size",
        "compression-ratio",
        "duplicate",
        "encrypted",
    ],
)
def test_pptx_archive_preflight_rejects_unsafe_metadata(payload: bytes) -> None:
    with pytest.raises(ValueError):
        benchmark._validate_pptx_archive(payload)


@pytest.mark.parametrize(
    "name",
    [
        "/absolute.xml",
        "C:/drive.xml",
        ".",
        "..",
        "a/./part.xml",
        "a/../part.xml",
        "a\\part.xml",
    ],
)
def test_pptx_archive_preflight_rejects_dangerous_member_name(name: str) -> None:
    payload = _zip_bytes([name])
    if "\\" in name:
        payload = payload.replace(b"a/part.xml", b"a\\part.xml")
    with pytest.raises(ValueError):
        benchmark._validate_pptx_archive(payload)


def test_unsafe_archive_is_rejected_before_python_pptx(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    execution = _execution(tmp_path, "images", 8)
    execution.output_path.write_bytes(
        _zip_bytes(["ratio.xml"], b"0" * (2 * 1024 * 1024))
    )
    presentation_called = False

    def unexpected_presentation(stream: object) -> object:
        nonlocal presentation_called
        presentation_called = True
        raise RuntimeError

    monkeypatch.setattr("pptx.Presentation", unexpected_presentation)

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "output_invalid"
    assert presentation_called is False


@pytest.mark.parametrize(("route_id", "pages"), [("images", 8), ("pdf", 3)])
def test_legacy_product_summary_without_output_sha256_uses_bound_digest(
    tmp_path: Path, route_id: str, pages: int
) -> None:
    execution = _execution(tmp_path, route_id, pages)

    result = benchmark.evaluate_execution(execution)

    assert result["status"] == "passed"
    assert result["output_sha256"] == hashlib.sha256(
        execution.output_path.read_bytes()
    ).hexdigest()


@pytest.mark.parametrize(("route_id", "pages"), [("images", 8), ("pdf", 3)])
def test_legacy_product_summary_rejects_null_output_sha256(
    tmp_path: Path, route_id: str, pages: int
) -> None:
    execution = _execution(tmp_path, route_id, pages)
    summary = json.loads(execution.stdout)
    summary["output_sha256"] = None

    result = benchmark.evaluate_execution(
        replace(execution, stdout=json.dumps(summary))
    )

    assert result["error_type"] == "output_invalid"


def test_mixed_product_summary_requires_output_sha256(tmp_path: Path) -> None:
    execution = _execution(tmp_path, "mixed_pptx", 3, mode="mixed_valid")
    summary = json.loads(execution.stdout)
    summary.pop("output_sha256")

    result = benchmark.evaluate_execution(
        replace(execution, stdout=json.dumps(summary))
    )

    assert result["error_type"] == "output_invalid"


@pytest.mark.parametrize(
    ("route_id", "pages", "wrong_key"),
    [("images", 8, "pptx"), ("pdf", 3, "pptx"), ("mixed_pptx", 3, "16:9")],
)
def test_route_rejects_wrong_product_output_key(
    tmp_path: Path, route_id: str, pages: int, wrong_key: str
) -> None:
    execution = _execution(
        tmp_path,
        route_id,
        pages,
        mode="mixed_valid" if route_id == "mixed_pptx" else "text",
    )
    summary = json.loads(execution.stdout)
    summary["outputs"] = {wrong_key: str(execution.output_path)}

    result = benchmark.evaluate_execution(
        replace(execution, stdout=json.dumps(summary))
    )

    assert result["error_type"] == "invalid_summary"


@pytest.mark.parametrize(("route_id", "pages"), [("images", 8), ("pdf", 3)])
@pytest.mark.parametrize("mode", ["aspect_4_3", "aspect_2_1"])
def test_legacy_route_rejects_non_widescreen_output(
    tmp_path: Path, route_id: str, pages: int, mode: str
) -> None:
    execution = _execution(tmp_path, route_id, pages, mode=mode)

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "output_invalid"


def test_mixed_route_does_not_require_widescreen_output(tmp_path: Path) -> None:
    execution = _execution(tmp_path, "mixed_pptx", 3, mode="aspect_4_3")

    result = benchmark.evaluate_execution(execution)

    assert result["status"] == "passed"


def test_missing_page_result_is_invalid_summary(tmp_path: Path) -> None:
    execution = _execution(tmp_path, "pdf", 3)
    (execution.run_root / "pages" / "page_002" / "page_result.json").unlink()

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "invalid_summary"


@pytest.mark.parametrize("payload", ['{"schema_version":1,"page_id":"page_001","page_id":"page_001","status":"validated"}', '{"schema_version":1,"page_id":"page_001","status":"validated","score":NaN}'])
def test_page_result_rejects_duplicate_or_nonfinite_json(
    tmp_path: Path, payload: str
) -> None:
    execution = _execution(tmp_path, "images", 8)
    path = execution.run_root / "pages" / "page_001" / "page_result.json"
    path.write_text(payload, encoding="utf-8")

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "invalid_summary"


def test_runtime_warning_fallback_precedes_editability(tmp_path: Path) -> None:
    execution = _execution(tmp_path, "images", 8, mode="full_picture")
    path = execution.run_root / "pages" / "page_001" / "page_result.json"
    page = json.loads(path.read_text(encoding="utf-8"))
    page["warnings"] = ["secret warning"]
    path.write_text(json.dumps(page), encoding="utf-8")

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "warning_fallback"
    assert "secret warning" not in json.dumps(result)


@pytest.mark.parametrize("mode", ["full_picture", "small_picture"])
def test_single_picture_reconstructed_page_is_flattened_output(
    tmp_path: Path, mode: str
) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode=mode)

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "flattened_output"


@pytest.mark.parametrize("mode", ["empty_text", "group_empty_text", "invisible_shape"])
def test_nonvisual_native_shape_cannot_hide_flattened_output(
    tmp_path: Path, mode: str
) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode=mode)

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "flattened_output"


@pytest.mark.parametrize(
    "mode", ["group_offslide", "group_zero_width", "group_zero_height"]
)
def test_group_transform_cannot_make_local_child_look_visible(
    tmp_path: Path, mode: str
) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode=mode)

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "flattened_output"


@pytest.mark.parametrize(
    "mode",
    [
        "group_flip_h_offslide",
        "group_flip_v_offslide",
        "nested_group_flip_h_offslide",
    ],
)
def test_group_flip_cannot_make_local_child_look_visible(
    tmp_path: Path, mode: str
) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode=mode)

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "flattened_output"


def test_group_flip_keeps_page_local_text_editable(tmp_path: Path) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode="group_flip_h_visible")

    result = benchmark.evaluate_execution(execution)

    assert result["status"] == "passed"


@pytest.mark.parametrize(
    "mode", ["outside_text", "zero_text", "outside_picture", "zero_picture"]
)
def test_off_slide_or_zero_size_content_cannot_hide_flattened_output(
    tmp_path: Path, mode: str
) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode=mode)

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "flattened_output"


@pytest.mark.parametrize(
    "mode",
    [
        "tiny_text",
        "tiny_picture",
        "tiny_native",
        "hidden_text",
        "hidden_picture",
        "group_hidden_text",
        "transparent_text",
        "transparent_picture",
        "shape_transparent_picture",
        "shape_alpha_repl_transparent_picture",
        "shape_unknown_alpha_picture",
        "shape_invalid_alpha_picture",
        "shape_extra_alpha_field_picture",
        "sparse_picture",
        "transparent_native",
        "transparent_native_leading_zero",
    ],
)
def test_tiny_hidden_or_transparent_content_cannot_hide_flattened_output(
    tmp_path: Path, mode: str
) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode=mode)

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "flattened_output"


@pytest.mark.parametrize("dimensions", [(10_000, 10_000), (20_000, 10_000)])
def test_picture_pixel_bomb_is_output_invalid_without_warning_escape(
    tmp_path: Path, dimensions: tuple[int, int]
) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode="pictures")
    _replace_first_pptx_media(
        execution.output_path, _png_with_dimensions(*dimensions)
    )

    with warnings.catch_warnings(record=True) as caught:
        warnings.simplefilter("always")
        result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "output_invalid"
    assert not any(
        issubclass(warning.category, Image.DecompressionBombWarning)
        for warning in caught
    )


@pytest.mark.parametrize(
    "mode",
    [
        "text",
        "native",
        "pictures",
        "group_text",
        "partial_gradient_native",
        "partial_picture_alpha",
        "partial_picture_alpha_repl",
    ],
)
def test_real_editable_shapes_pass_structure_gate(tmp_path: Path, mode: str) -> None:
    execution = _execution(tmp_path, "pdf", 3, mode=mode)

    result = benchmark.evaluate_execution(execution)

    assert result["status"] == "passed"
    assert result["error_type"] is None


@pytest.mark.parametrize(
    ("statuses", "expected"),
    [
        (["preserved", "replaced", "replaced"], None),
        (["replaced", "replaced", "replaced"], None),
        (["preserved", "preserved", "replaced"], "warning_fallback"),
        (["preserved", "replaced", "preserved"], "warning_fallback"),
        (["preserved_with_warning", "replaced", "replaced"], "warning_fallback"),
    ],
)
def test_mixed_route_requires_replacement_of_pages_two_and_three(
    tmp_path: Path, statuses: list[str], expected: str | None
) -> None:
    execution = _execution(
        tmp_path, "mixed_pptx", 3, mode="mixed_valid", statuses=statuses
    )

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == expected


@pytest.mark.parametrize("problem", ["missing", "bool", "negative", "huge", "identifier", "page"])
def test_performance_contract_is_required_and_strict(
    tmp_path: Path, problem: str
) -> None:
    execution = _execution(tmp_path, "pdf", 3)
    summary = json.loads(execution.stdout)
    if problem == "missing":
        summary.pop("performance")
    elif problem == "page":
        summary["performance"]["pages"].pop("page_003")
    else:
        page = summary["performance"]["pages"]["page_001"]
        if problem == "bool":
            page["agent_runs"] = True
        elif problem == "negative":
            page["agent_runs"] = -1
        elif problem == "huge":
            page["agent_runs"] = 1_000_000_000_001
        else:
            page["model_loads"] = {"bad identifier!": 1}
    execution = replace(execution, stdout=json.dumps(summary))

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "invalid_summary"


def test_write_report_is_exclusive(tmp_path: Path) -> None:
    path = tmp_path / "benchmark-report.json"
    path.write_text("existing", encoding="utf-8")

    with pytest.raises(FileExistsError):
        benchmark.write_report(path, {"status": "passed"})

    assert path.read_text(encoding="utf-8") == "existing"


def test_multiply_linked_output_is_output_invalid(tmp_path: Path) -> None:
    execution = _execution(tmp_path, "pdf", 3)
    os.link(execution.output_path, execution.output_path.with_name("alias.pptx"))

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "output_invalid"


@pytest.mark.parametrize("problem", ["map_bool", "map_negative", "map_huge", "extra_field"])
def test_performance_maps_and_exact_page_fields_are_strict(
    tmp_path: Path, problem: str
) -> None:
    execution = _execution(tmp_path, "pdf", 3)
    summary = json.loads(execution.stdout)
    page = summary["performance"]["pages"]["page_001"]
    if problem == "extra_field":
        page["invented"] = 0
    else:
        page["model_loads"]["layout-model"] = {
            "map_bool": True,
            "map_negative": -1,
            "map_huge": 1_000_000_000_001,
        }[problem]
    execution = replace(execution, stdout=json.dumps(summary))

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "invalid_summary"


@pytest.mark.parametrize("location", ["summary", "page"])
def test_mixed_warning_fields_are_warning_fallback(
    tmp_path: Path, location: str
) -> None:
    execution = _execution(tmp_path, "mixed_pptx", 3, mode="mixed_valid")
    summary = json.loads(execution.stdout)
    if location == "summary":
        summary["warnings"] = ["private warning"]
        summary["preserved_with_warning_pages"] = 1
        summary["page_results"][0]["status"] = "preserved_with_warning"
    else:
        summary["page_results"][1]["warning"] = "private warning"
    execution = replace(execution, stdout=json.dumps(summary))

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "warning_fallback"
    assert "private warning" not in json.dumps(result)


def test_mixed_structured_page_warning_is_warning_fallback(tmp_path: Path) -> None:
    execution = _execution(tmp_path, "mixed_pptx", 3, mode="mixed_valid")
    summary = json.loads(execution.stdout)
    summary["page_results"][1]["warning"] = ["structured warning"]
    execution = replace(execution, stdout=json.dumps(summary))

    result = benchmark.evaluate_execution(execution)

    assert result["error_type"] == "warning_fallback"


def test_untrusted_summary_fields_never_reach_report(tmp_path: Path) -> None:
    executions = [
        _execution(tmp_path, "images", 8),
        _execution(tmp_path, "pdf", 3),
        _execution(tmp_path, "mixed_pptx", 3, mode="mixed_valid"),
    ]
    injected = json.loads(executions[0].stdout)
    injected.update(
        error="sk-secret",
        diagnostics="C:\\private\\asset",
        url="https://private.example/token",
    )
    executions[0] = replace(executions[0], stdout=json.dumps(injected))
    manifest = benchmark.CorpusManifest(
        tmp_path, "b" * 64, (), tuple(item.route for item in executions)
    )

    report = benchmark.build_report(manifest, tuple(executions))
    encoded = json.dumps(report)

    assert "sk-secret" not in encoded
    assert "C:\\\\private" not in encoded
    assert "https://" not in encoded


def test_performance_identifier_secret_is_aggregated_without_copying_key(
    tmp_path: Path,
) -> None:
    executions = [
        _execution(tmp_path, "images", 8),
        _execution(tmp_path, "pdf", 3),
        _execution(tmp_path, "mixed_pptx", 3, mode="mixed_valid"),
    ]
    summary = json.loads(executions[0].stdout)
    metrics = summary["performance"]["pages"]["page_001"]["model_loads"]
    metrics["sk-proj-secret-token"] = 4
    executions[0] = replace(executions[0], stdout=json.dumps(summary))
    manifest = benchmark.CorpusManifest(
        tmp_path, "d" * 64, (), tuple(item.route for item in executions)
    )

    report = benchmark.build_report(manifest, tuple(executions))
    encoded = json.dumps(report)

    assert report["routes"][0]["performance"]["pages"]["page_001"][
        "model_loads"
    ] == 5
    assert "sk-proj-secret-token" not in encoded


@pytest.mark.parametrize(
    ("warning_count", "statuses", "expected_error"),
    [
        (-1, ["preserved", "replaced", "replaced"], "invalid_summary"),
        (4, ["preserved", "replaced", "replaced"], "invalid_summary"),
        (1, ["preserved", "replaced", "replaced"], "invalid_summary"),
        (0, ["preserved_with_warning", "replaced", "replaced"], "invalid_summary"),
    ],
)
def test_mixed_warning_count_is_bounded_and_matches_page_results(
    tmp_path: Path,
    warning_count: int,
    statuses: list[str],
    expected_error: str,
) -> None:
    execution = _execution(
        tmp_path, "mixed_pptx", 3, mode="mixed_valid", statuses=statuses
    )
    summary = json.loads(execution.stdout)
    summary["preserved_with_warning_pages"] = warning_count

    result = benchmark.evaluate_execution(
        replace(execution, stdout=json.dumps(summary))
    )

    assert result["error_type"] == expected_error


def test_mixed_warning_pages_are_deduplicated_and_bounded(tmp_path: Path) -> None:
    execution = _execution(
        tmp_path,
        "mixed_pptx",
        3,
        mode="mixed_valid",
        statuses=["preserved", "preserved_with_warning", "replaced"],
    )
    summary = json.loads(execution.stdout)
    summary["preserved_with_warning_pages"] = 1
    summary["warnings"] = [f"warning {index}" for index in range(10)]

    result = benchmark.evaluate_execution(
        replace(execution, stdout=json.dumps(summary))
    )

    assert result["error_type"] == "warning_fallback"
    assert result["warning_pages"] == 3


def _manifest_for_executions(
    tmp_path: Path, executions: tuple[benchmark.RouteExecution, ...]
) -> benchmark.CorpusManifest:
    return benchmark.CorpusManifest(
        tmp_path, "c" * 64, (), tuple(item.route for item in executions)
    )


@pytest.mark.parametrize(("returncode", "expected_exit"), [(7, 1), (0, 0)])
def test_main_writes_report_with_matching_exit_code(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    returncode: int,
    expected_exit: int,
) -> None:
    data_root = tmp_path / "data"
    executions = (
        _execution(
            data_root,
            "images",
            8,
            returncode=returncode,
            stdout="sk-secret" if returncode else None,
        ),
        _execution(data_root, "pdf", 3),
        _execution(data_root, "mixed_pptx", 3, mode="mixed_valid"),
    )
    manifest = _manifest_for_executions(tmp_path, executions)
    output_root = tmp_path / "result"
    monkeypatch.setattr(benchmark, "load_manifest", lambda path: manifest)

    def fake_execute(
        received: benchmark.CorpusManifest, result: Path
    ) -> tuple[benchmark.RouteExecution, ...]:
        assert received is manifest
        result.mkdir()
        return executions

    monkeypatch.setattr(benchmark, "execute_routes", fake_execute)

    exit_code = benchmark.main(
        ["--corpus", str(tmp_path / "manifest.json"), "--output-dir", str(output_root)]
    )
    report = json.loads((output_root / "benchmark-report.json").read_text("utf-8"))

    assert exit_code == expected_exit
    assert report["status"] == ("failed" if returncode else "passed")
    assert report["routes"][0]["error_type"] == (
        "conversion_failed" if returncode else None
    )
    assert report["routes"][1]["status"] == "passed"
    assert "sk-secret" not in json.dumps(report)


def test_main_doctor_failure_reports_all_routes(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    executions = (
        _execution(tmp_path / "data", "images", 8),
        _execution(tmp_path / "data", "pdf", 3),
        _execution(tmp_path / "data", "mixed_pptx", 3, mode="mixed_valid"),
    )
    manifest = _manifest_for_executions(tmp_path, executions)
    output_root = tmp_path / "result"
    monkeypatch.setattr(benchmark, "load_manifest", lambda path: manifest)

    def failed_doctor(received: benchmark.CorpusManifest, result: Path) -> object:
        result.mkdir()
        raise benchmark.BenchmarkError("doctor_not_ready")

    monkeypatch.setattr(benchmark, "execute_routes", failed_doctor)

    exit_code = benchmark.main(
        ["--corpus", str(tmp_path / "manifest.json"), "--output-dir", str(output_root)]
    )
    report = json.loads((output_root / "benchmark-report.json").read_text("utf-8"))

    assert exit_code == 1
    assert [route["error_type"] for route in report["routes"]] == [
        "doctor_not_ready"
    ] * 3
    assert report["totals"]["failed_routes"] == 3


def test_main_report_exclusive_failure_cannot_return_success(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch, capsys: pytest.CaptureFixture[str]
) -> None:
    data_root = tmp_path / "data"
    executions = (
        _execution(data_root, "images", 8),
        _execution(data_root, "pdf", 3),
        _execution(data_root, "mixed_pptx", 3, mode="mixed_valid"),
    )
    manifest = _manifest_for_executions(tmp_path, executions)
    output_root = tmp_path / "result"
    monkeypatch.setattr(benchmark, "load_manifest", lambda path: manifest)

    def fake_execute(
        received: benchmark.CorpusManifest, result: Path
    ) -> tuple[benchmark.RouteExecution, ...]:
        result.mkdir()
        (result / "benchmark-report.json").write_text("existing", encoding="utf-8")
        return executions

    monkeypatch.setattr(benchmark, "execute_routes", fake_execute)

    exit_code = benchmark.main(
        ["--corpus", str(tmp_path / "manifest.json"), "--output-dir", str(output_root)]
    )

    assert exit_code == 1
    assert capsys.readouterr().err == "benchmark_failed\n"
    assert (output_root / "benchmark-report.json").read_text("utf-8") == "existing"


def test_cli_invalid_corpus_has_fixed_exit_code_and_stderr(tmp_path: Path) -> None:
    completed = subprocess.run(
        [
            sys.executable,
            "scripts/benchmark_conversion.py",
            "--corpus",
            str(tmp_path / "missing.json"),
            "--output-dir",
            str(tmp_path / "result"),
        ],
        cwd=Path(__file__).resolve().parents[1],
        text=True,
        capture_output=True,
        check=False,
    )

    assert completed.returncode == 1
    assert completed.stdout == ""
    assert completed.stderr == "invalid_corpus\n"
    assert not (tmp_path / "result").exists()
