from __future__ import annotations

import copy
import hashlib
import importlib
import inspect
import json
import os
import re
import subprocess
from datetime import datetime
from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from PIL import Image, ImageChops, UnidentifiedImageError
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader, PdfWriter


ROOT = Path(__file__).resolve().parents[1]
RELEASE_ROOT = ROOT / "benchmarks" / "release"
MANIFEST_PATH = RELEASE_ROOT / "manifest.json"
CORE_MANIFEST_PATH = RELEASE_ROOT / "core-v0.2-manifest.json"

ROOT_FIELDS = {"schema_version", "cases", "categories"}
CASE_FIELDS = {
    "id",
    "kind",
    "path",
    "sha256",
    "source",
    "license",
    "page_count",
    "categories",
    "agent_provider",
    "pptx_mode",
    "expected_pages",
}
PAGE_FIELDS = {
    "page_id",
    "expected_status",
    "min_visual_components",
    "min_text_boxes",
    "max_unexplained_pixels",
    "max_quality_violations",
}

CASE_SPECS = [
    (
        "image-bilingual-dashboard",
        "image",
        "inputs/01-bilingual-dashboard.png",
        1,
        "bilingual_dashboard",
        None,
        12,
        8,
    ),
    (
        "image-dense-parameter-comparison",
        "image",
        "inputs/02-dense-parameter-comparison.png",
        1,
        "dense_parameter_comparison",
        None,
        15,
        12,
    ),
    (
        "image-profile-cards",
        "image",
        "inputs/03-profile-cards.png",
        1,
        "profile_cards",
        None,
        10,
        8,
    ),
    (
        "image-four-stage-timeline",
        "image",
        "inputs/04-four-stage-timeline.png",
        1,
        "four_stage_timeline",
        None,
        10,
        6,
    ),
    (
        "image-combo-chart",
        "image",
        "inputs/05-combo-chart.png",
        1,
        "combo_chart",
        None,
        8,
        5,
    ),
    (
        "image-flowchart",
        "image",
        "inputs/06-flowchart.png",
        1,
        "flowchart",
        None,
        10,
        6,
    ),
    (
        "image-icon-matrix",
        "image",
        "inputs/07-icon-matrix.png",
        1,
        "icon_matrix",
        None,
        12,
        4,
    ),
    (
        "image-light-text-gradient",
        "image",
        "inputs/08-light-text-gradient.png",
        1,
        "light_text_gradient",
        None,
        3,
        4,
    ),
    (
        "image-thin-line-network",
        "image",
        "inputs/09-thin-line-network.png",
        1,
        "thin_line_network",
        None,
        12,
        5,
    ),
    (
        "image-tiny-element-table",
        "image",
        "inputs/10-tiny-element-table.png",
        1,
        "tiny_element_table",
        None,
        15,
        10,
    ),
    (
        "image-dark-poster",
        "image",
        "inputs/11-dark-poster.png",
        1,
        "dark_poster",
        None,
        5,
        5,
    ),
    (
        "image-non-16-9-infographic",
        "image",
        "inputs/12-non-16-9-infographic.png",
        1,
        "non_16_9_infographic",
        None,
        12,
        8,
    ),
    (
        "pdf-mixed-page-sizes",
        "pdf",
        "inputs/13-mixed-page-sizes.pdf",
        2,
        "pdf_mixed_page_sizes",
        None,
        4,
        11,
    ),
    (
        "pdf-rotated-page",
        "pdf",
        "inputs/14-rotated-page.pdf",
        2,
        "pdf_rotated_page",
        None,
        4,
        11,
    ),
    ("pdf-high-dpi", "pdf", "inputs/15-high-dpi.pdf", 2, "pdf_high_dpi", None, 6, 4),
    (
        "pptx-image-only",
        "pptx",
        "inputs/16-image-only.pptx",
        4,
        "pptx_image_only",
        "image_only",
        6,
        4,
    ),
    (
        "pptx-mixed-native",
        "pptx",
        "inputs/17-mixed-native.pptx",
        4,
        "pptx_mixed_native",
        "mixed_native",
        6,
        4,
    ),
    (
        "pptx-mixed-screenshot-candidates",
        "pptx",
        "inputs/18-mixed-screenshot-candidates.pptx",
        4,
        "pptx_mixed_screenshot_candidates",
        "mixed_screenshot_candidates",
        6,
        4,
    ),
]

EXPECTED_CATEGORIES = [spec[4] for spec in CASE_SPECS]
CORE_CASE_IDS = (
    "image-bilingual-dashboard",
    "image-combo-chart",
    "image-flowchart",
    "image-icon-matrix",
    "image-thin-line-network",
    "image-tiny-element-table",
    "image-dark-poster",
    "image-non-16-9-infographic",
    "pdf-rotated-page",
    "pptx-mixed-screenshot-candidates",
)
EXPECTED_README = """# v0.2 核心 14 页 benchmark

这里保存 v0.2 发布门禁使用的固定语料、manifest 和审核过的 Host plans。核心语料由 10 个完整 case、14 页组成：8 张图片、2 页旋转 PDF 和 4 页 mixed screenshot candidates PPTX。多页文件会整份测试，不截取其中几页。

仓库还保留了用于补充覆盖的生成语料，记录在 `manifest.json` 中；它不属于 v0.2 核心通过条件，也不会被写进核心完成率。

## 覆盖内容

核心图片覆盖中英双语仪表盘、柱线组合图、流程图、图标矩阵、细线网络图、小元素表格、深色海报和非 16:9 信息图。PDF 用来验证旋转页面，PPTX 用来验证 native objects 与 screenshot candidates 混合存在时的处理。

输入由项目脚本生成，manifest 记录 `source=project-generated` 和 `license=CC0-1.0`。每个输入的 bytes 与 SHA-256 都已固定，运行时会重新校验。

## Release Gate 如何运行

真实模型 benchmark 不放进普通 push CI，需要在受保护的 Release Gate 中显式开启。GitHub-hosted Windows 会把 10 个 case 分成 5 组并行运行；每个 case 仍执行 3 次独立重复。分片只缩短等待时间，不减少测试次数，也不放宽任何质量门禁。

五份分片报告会由独立的聚合步骤重新校验。v0.2 的正式参考环境固定为 GitHub-hosted Windows、AMD64、Python 3.12 和 CPU。只有 manifest、依赖约束、运行环境、10 个 case、30 次尝试和 42 个累计页面全部一致，且性能没有超过同环境基线 15%，才会生成 `report_kind: official`、`status: passed` 的正式报告。单个分片不能代表 benchmark 通过。

当 GitHub-hosted 环境产生的 request/graph hash 与已有 plans 不一致时，可以先运行 diagnostic。runner 会先选择与当前 request/graph 完整匹配的 plan；没有 exact match 时，只能在 action 通过当前 request/graph 完整契约校验的旧 plan 中选择唯一一个，并更新绑定 hash。无论哪种情况，原有 decision、actions、parameters、confidence 和 evidence 都必须保持不变；三次重复得到的绑定也必须完全一致。它会继续执行相同的页面与质量检查，但报告只会是 `report_kind: diagnostic`，不能算作正式通过。

性能基线只能通过 `baseline-candidate` 从五份成功的三次重复 diagnostic 报告生成。该命令会重新检查固定的 CPU 环境、完整 case 和页数、manifest、依赖约束及性能数据；probe、失败报告、CUDA 报告、缺失或重复的 case 都不能生成候选基线。候选文件仍需审核并提交为 `BASELINE.json`，它本身不会把 diagnostic 变成正式通过。

如果只需要定位一次耗时较长的失败，可以把 diagnostic repeats 选为 `1`。这会运行相同的真实模型、页面检查和质量门禁，但只生成 `report_kind: probe` 的单次诊断报告，不计算性能，也不生成候选 plan。问题修复后仍需把 repeats 选回 `3`，完成三次严格 diagnostic。

Release Gate 只上传 JSON 报告，以及 diagnostic 产生的候选 plan JSON。模型文件、模型缓存、输入副本、运行 workspace 和生成的 PPTX 都不会上传为 benchmark 工件，也不会提交到 Git。

## 严格通过标准

每页必须达到 manifest 中的最小非文本视觉组件数 `min_visual_components` 和最小原生文本框数 `min_text_boxes`。两类对象独立计数；已经转为原生文本的 OCR 内容不能继续留在 raster 中重复计数。

所有页面还必须满足预期状态，并同时达到 0 warning、0 fallback、0 unexplained pixels 和 0 quality violations。缺少质量文件、plan 不匹配、任一重复失败或性能回归都会使门禁失败。

## 字体与重新生成

生成器使用仓库中的 Google Fonts `Noto Sans SC` variable TTF：`fonts/NotoSansSC[wght].ttf`。字体按 `fonts/OFL.txt` 中的 SIL Open Font License 1.1 分发，固定来源为 `google/fonts` commit `e1118da94a8cb00cf6d06cdac9ef13eb1e5c6ab7`。字体本身不是 CC0；由它渲染出的 benchmark 输入按 CC0-1.0 发布。

可以运行 `python scripts/build_release_corpus.py <output-root>` 在一个不存在的目录中重新生成语料。同一环境的两次 fresh generation 要求 PNG RGB 像素一致；不同平台只承诺格式、尺寸、页数和对象 inventory 等明确语义一致，不承诺 PNG 像素或 PDF/PPTX 字节完全相同。
"""

EXPECTED_INPUT_NAMES = {Path(spec[2]).name for spec in CASE_SPECS}
EXPECTED_FONT_PATH = RELEASE_ROOT / "fonts" / "NotoSansSC[wght].ttf"
EXPECTED_FONT_SHA256 = "a3041811a78c361b1de50f953c805e0244951c21c5bd412f7232ef0d899af0da"
EXPECTED_IMAGE_SIZES = {
    "01-bilingual-dashboard.png": (1600, 900),
    "02-dense-parameter-comparison.png": (1600, 900),
    "03-profile-cards.png": (1600, 900),
    "04-four-stage-timeline.png": (1600, 900),
    "05-combo-chart.png": (1600, 900),
    "06-flowchart.png": (1600, 900),
    "07-icon-matrix.png": (1600, 900),
    "08-light-text-gradient.png": (1600, 900),
    "09-thin-line-network.png": (1600, 900),
    "10-tiny-element-table.png": (1600, 900),
    "11-dark-poster.png": (1600, 900),
    "12-non-16-9-infographic.png": (1000, 1400),
}


@pytest.fixture(scope="module")
def fresh_release_corpora(tmp_path_factory: pytest.TempPathFactory) -> tuple[Path, Path]:
    builder = importlib.import_module("scripts.build_release_corpus")
    parent = tmp_path_factory.mktemp("release-corpus")
    first = parent / "first"
    second = parent / "second"

    builder.build(first)
    builder.build(second)

    return first, second


def _pptx_inventory(path: Path) -> tuple[tuple[tuple[int, str], ...], ...]:
    deck = Presentation(path)
    return tuple(
        tuple((int(shape.shape_type), shape.text if shape.has_text_frame else "") for shape in slide.shapes)
        for slide in deck.slides
    )


def _pdf_inventory(path: Path) -> tuple[tuple[float, float, int], ...]:
    return tuple(
        (
            round(float(page.mediabox.width), 3),
            round(float(page.mediabox.height), 3),
            page.rotation,
        )
        for page in PdfReader(path).pages
    )


def _assert_same_rgb_pixels(first: Image.Image, second: Image.Image) -> None:
    assert ImageChops.difference(
        first.convert("RGB"), second.convert("RGB")
    ).getbbox() is None


def _reject_duplicate_keys(pairs: list[tuple[str, object]]) -> dict[str, object]:
    value = {}
    for key, item in pairs:
        if key in value:
            raise ValueError(f"duplicate JSON key: {key}")
        value[key] = item
    return value


def _manifest() -> dict[str, object]:
    return json.loads(
        MANIFEST_PATH.read_text(encoding="utf-8"),
        object_pairs_hook=_reject_duplicate_keys,
    )


def _core_manifest() -> dict[str, object]:
    return json.loads(
        CORE_MANIFEST_PATH.read_text(encoding="utf-8"),
        object_pairs_hook=_reject_duplicate_keys,
    )


def _assert_exact_fields(manifest: dict[str, object]) -> None:
    assert set(manifest) == ROOT_FIELDS
    for case in manifest["cases"]:
        assert set(case) == CASE_FIELDS
        for page in case["expected_pages"]:
            assert set(page) == PAGE_FIELDS


def _assert_numeric_contract(manifest: dict[str, object]) -> None:
    assert type(manifest["schema_version"]) is int
    assert manifest["schema_version"] == 2
    for case in manifest["cases"]:
        assert type(case["page_count"]) is int
        assert case["page_count"] > 0
        for page in case["expected_pages"]:
            for field in (
                "min_visual_components",
                "min_text_boxes",
                "max_unexplained_pixels",
                "max_quality_violations",
            ):
                assert type(page[field]) is int
                assert page[field] >= 0


def _assert_release_input(case: dict[str, object], release_root: Path) -> None:
    path = release_root / case["path"]
    assert path.is_file(), path
    digest = hashlib.sha256(path.read_bytes()).hexdigest()
    assert case["sha256"] == digest
    assert digest != "0" * 64

    if case["kind"] == "image":
        with Image.open(path) as image:
            assert image.format == "PNG"
            assert image.n_frames == case["page_count"]
            image.verify()
    elif case["kind"] == "pdf":
        assert len(PdfReader(path).pages) == case["page_count"]
    elif case["kind"] == "pptx":
        assert len(Presentation(path).slides) == case["page_count"]
    else:
        raise AssertionError(f"unsupported release input kind: {case['kind']}")


def test_release_manifest_has_exact_root_schema_and_categories() -> None:
    manifest = _manifest()

    _assert_exact_fields(manifest)
    _assert_numeric_contract(manifest)
    assert manifest["schema_version"] == 2
    assert manifest["categories"] == EXPECTED_CATEGORIES
    assert len(manifest["cases"]) == 18
    assert len({case["id"] for case in manifest["cases"]}) == 18


def test_release_manifest_defines_eighteen_inputs_and_thirty_pages() -> None:
    manifest = _manifest()
    cases = manifest["cases"]

    assert [
        (
            case["id"],
            case["kind"],
            case["path"],
            case["page_count"],
            case["categories"][0],
            case["pptx_mode"],
        )
        for case in cases
    ] == [spec[:6] for spec in CASE_SPECS]
    assert {
        kind: sum(case["kind"] == kind for case in cases)
        for kind in ("image", "pdf", "pptx")
    } == {
        "image": 12,
        "pdf": 3,
        "pptx": 3,
    }
    assert sum(case["page_count"] for case in cases) == 30


def test_core_v020_manifest_is_exact_subset_of_extended_corpus() -> None:
    extended = _manifest()
    core = _core_manifest()
    selected = {case["id"]: case for case in extended["cases"]}

    _assert_exact_fields(core)
    _assert_numeric_contract(core)
    assert [case["id"] for case in core["cases"]] == list(CORE_CASE_IDS)
    assert core["cases"] == [selected[case_id] for case_id in CORE_CASE_IDS]
    assert core["categories"] == [
        case["categories"][0] for case in core["cases"]
    ]
    assert len(core["cases"]) == 10
    assert sum(case["page_count"] for case in core["cases"]) == 14
    kinds = [case["kind"] for case in core["cases"]]
    assert {kind: kinds.count(kind) for kind in ("image", "pdf", "pptx")} == {
        "image": 8,
        "pdf": 1,
        "pptx": 1,
    }


def test_core_v020_pptx_contract_matches_authored_page_routes_and_objects() -> None:
    core = _core_manifest()
    case = next(
        item
        for item in core["cases"]
        if item["id"] == "pptx-mixed-screenshot-candidates"
    )

    assert [
        (
            page["expected_status"],
            page["min_visual_components"],
            page["min_text_boxes"],
        )
        for page in case["expected_pages"]
    ] == [
        ("replaced", 4, 39),
        ("replaced", 24, 14),
        ("replaced", 3, 4),
        ("preserved", 1, 1),
    ]


def test_core_v020_pptx_replay_plans_match_authored_requests() -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = [
        json.loads(path.read_text(encoding="utf-8"))
        for path in sorted(
            (RELEASE_ROOT / "plans").glob("pptx-mixed-screenshot-candidates--*.json")
        )
    ]
    candidate_plans = [plan for plan in plans if plan.get("kind") == "candidate_decision"]
    component_plans = [plan for plan in plans if plan.get("kind") == "component_plan"]

    assert len(candidate_plans) == 3
    assert all(runner._valid_candidate_plan(plan) for plan in candidate_plans)
    assert [
        (plan["page_id"], plan["image_sha256"], plan["source_object_sha256"])
        for plan in candidate_plans
    ] == [
        ("page_001", "594d8f1eeeb94746e8aaa5730a3fa291855423f2eaea48ba30b1eeb577ed4e9d", "3362dd457e443599faf0e3f1254e3c5ff1a570e3c3362e5d3e09e77ba0b010a5"),
        ("page_002", "cc2d8482b48fd32f9920cf0092da88ad16d345821db0d0f98b11103b09871900", "7fa090a9f03eef6ec086214c8c9639197369193d850c72d51faa241b617c5c18"),
        ("page_003", "65eafdda6be280feceee856cda45b76a68e2e9cf1f82646abed7c1dadacab798", "0030b760342c417d22d0b35e2ece11c7adfc1c65dde72018bd93ed61f8e63c02"),
    ]
    assert len(component_plans) == 11
    assert all(runner._valid_component_plan(plan) for plan in component_plans)
    assert [
        (plan["page_id"], plan["repair_round"], plan["request_sha256"], plan["graph_sha256"])
        for plan in component_plans
    ] == [
        ("page_001", 1, "24185f07f11855780141dcd851bb0d96a256194bc2abfe8bcd23996f9c160bb0", "047df8ef508cf23eb7280963f2ceea816b252b2d0a9e1d9fb077426d6d608628"),
        ("page_001", 2, "6202fe0d9f9f6c31dd59d47fca39660b9a2b0b6c4b4cd02fa9c3bbfc1b7efe20", "e630d52184ec4c498e095e82432f563253490550d10bdca11a76516439d746b3"),
        ("page_001", 3, "bcba8d9f50dd6932d8696395bdc2f2c9182e0bb4cf3caba5184a4c5ab6eaad50", "fbf0a49079c99395d6ef1f9ff2b8374d0a517de64f73ddcaa00209e9d97ba194"),
        ("page_001", 4, "7a726142d6b973eb46084d978edc9f62d87a0d9dd0a7c298962f3eba17051abc", "59e852e5a394bdff979dabe9f80928f1cec7f1ae8fdef83c53b2c66371c55175"),
        ("page_002", 1, "d124520930be78fa42a275a10604426a73c555f6e8fe6d2f230bbbce8c511ca0", "14c6305d6e1f3a3ec966f4b218186d5f0a48aa11872171b5b60c43431639756a"),
        ("page_002", 2, "b280a53626cbdd9fe331a31b5b30988ae02ddcf9ef9c7adb784d853fadf39a7c", "37aec06e3ce7d98602a2d4f85050be559f767a684dae2e64c7be098706eef7da"),
        ("page_002", 3, "4c0fecf5705cda34bbad48eca8b223f6f8c09fee314561ab8112b0e84078e382", "ddc1b99bdea467ac116df5e4144a12271fc8507ca84d8fa82c422dabf2d33090"),
        ("page_003", 1, "2edd9fcdfe62eeac56afbc0f299a076b6f604d1d0baaef755ec3837ead44af59", "c6f7ccf4c5ce928242841297b5f76cc059f417e4c08a8064d9998a84449dd4f6"),
        ("page_003", 2, "f16861801c987e00427c18b9a7b2f939ca62421638ae7f515be2ef4874960474", "eabb6e008ab8fb528d08bfe3c8dd6babe309cb08cb3cfb4c25fd959ad4e45e5d"),
        ("page_003", 3, "aefc7f8b2ea9264ac205807e4c9c0bd85bbf9d88b92de23b1a735c767e6cc5a7", "0bbc6904d38b036700dbaa2b6d3f118e0e02c4984b4c6a020ddf8593cfb666f9"),
        ("page_003", 4, "41a4627b8c909193d70f2b1ad3f60f4fc71216a12d62d4ecd986fed9034cc929", "d4b580ca78872dcc6e1c62ae13640f0e25063e566ae3736c214bbae3c6c927f1"),
    ]
def test_release_manifest_cases_and_pages_use_exact_strict_fields() -> None:
    cases = _manifest()["cases"]
    real_statuses = {
        "pptx-mixed-screenshot-candidates": [
            "replaced",
            "replaced",
            "replaced",
            "preserved",
        ]
    }
    real_thresholds = {
        "pptx-mixed-screenshot-candidates": [
            (4, 39),
            (24, 14),
            (3, 4),
            (1, 1),
        ]
    }

    for case, spec in zip(cases, CASE_SPECS, strict=True):
        identifier, _, _, page_count, category, _, min_visual, min_text_boxes = spec
        assert re.fullmatch(r"[0-9a-f]{64}", case["sha256"])
        assert case["source"] == "project-generated"
        assert case["license"] == "CC0-1.0"
        assert case["categories"] == [category]
        assert case["agent_provider"] == "host"
        assert len(case["expected_pages"]) == page_count

        for page_number, page in enumerate(case["expected_pages"], start=1):
            expected_status = real_statuses.get(
                identifier, ["validated"] * page_count
            )[page_number - 1]
            page_min_visual, page_min_text = real_thresholds.get(
                identifier, [(min_visual, min_text_boxes)] * page_count
            )[page_number - 1]
            assert page == {
                "page_id": f"{identifier}-page-{page_number:03d}",
                "expected_status": expected_status,
                "min_visual_components": page_min_visual,
                "min_text_boxes": page_min_text,
                "max_unexplained_pixels": 0,
                "max_quality_violations": 0,
            }


def test_release_manifest_pptx_modes_are_exact_and_only_apply_to_pptx() -> None:
    cases = _manifest()["cases"]
    pptx_cases = [case for case in cases if case["kind"] == "pptx"]

    assert {case["pptx_mode"] for case in pptx_cases} == {
        "image_only",
        "mixed_native",
        "mixed_screenshot_candidates",
    }
    assert all(case["pptx_mode"] is None for case in cases if case["kind"] != "pptx")


@pytest.mark.parametrize(
    ("needle", "duplicate"),
    [
        ('  "schema_version": 2,', '  "schema_version": 2,\n  "schema_version": 2,'),
        (
            '      "id": "image-bilingual-dashboard",',
            '      "id": "image-bilingual-dashboard",\n'
            '      "id": "image-bilingual-dashboard",',
        ),
        (
            '          "page_id": "image-bilingual-dashboard-page-001",',
            '          "page_id": "image-bilingual-dashboard-page-001",\n'
            '          "page_id": "image-bilingual-dashboard-page-001",',
        ),
    ],
)
def test_release_manifest_rejects_duplicate_keys_in_every_object(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    needle: str,
    duplicate: str,
) -> None:
    original = MANIFEST_PATH.read_text(encoding="utf-8")
    mutant = original.replace(needle, duplicate, 1)
    assert mutant != original
    mutant_path = tmp_path / "manifest.json"
    mutant_path.write_text(mutant, encoding="utf-8")
    monkeypatch.setitem(globals(), "MANIFEST_PATH", mutant_path)

    with pytest.raises(ValueError, match="duplicate JSON key"):
        _manifest()


@pytest.mark.parametrize(
    ("field", "value"),
    [
        ("schema_version", True),
        ("page_count", True),
        ("page_count", 0),
        ("min_visual_components", True),
        ("min_visual_components", -1),
        ("min_text_boxes", True),
        ("min_text_boxes", -1),
        ("max_unexplained_pixels", False),
        ("max_unexplained_pixels", -1),
        ("max_quality_violations", False),
        ("max_quality_violations", -1),
    ],
)
def test_release_manifest_rejects_bool_and_out_of_range_integers(
    field: str, value: object
) -> None:
    mutant = copy.deepcopy(_manifest())
    if field == "schema_version":
        mutant[field] = value
    elif field == "page_count":
        mutant["cases"][0][field] = value
    else:
        mutant["cases"][0]["expected_pages"][0][field] = value

    with pytest.raises(AssertionError):
        _assert_numeric_contract(mutant)


def test_release_manifest_object_order_is_not_schema() -> None:
    mutant = copy.deepcopy(_manifest())
    mutant["cases"][0]["expected_pages"][0] = dict(
        reversed(mutant["cases"][0]["expected_pages"][0].items())
    )
    mutant["cases"][0] = dict(reversed(mutant["cases"][0].items()))
    mutant = dict(reversed(mutant.items()))

    _assert_exact_fields(mutant)


def test_release_readme_documents_phase_one_and_future_strict_run() -> None:
    readme = (RELEASE_ROOT / "README.md").read_text(encoding="utf-8")

    assert readme == EXPECTED_README


def test_release_binary_inputs_have_git_attributes() -> None:
    attributes = (ROOT / ".gitattributes").read_text(encoding="utf-8").splitlines()
    ignores = (ROOT / ".gitignore").read_text(encoding="utf-8").splitlines()

    assert "benchmarks/release/inputs/*.png binary" in attributes
    assert "benchmarks/release/inputs/*.pdf binary" in attributes
    assert "benchmarks/release/inputs/*.pptx binary" in attributes
    assert "benchmarks/release/fonts/*.ttf binary" in attributes
    assert "!benchmarks/release/inputs/*.png" in ignores
    assert "!benchmarks/release/inputs/*.pdf" in ignores
    assert "!benchmarks/release/inputs/*.pptx" in ignores


def test_rgb_pixel_comparison_rejects_red_vs_blue_mutation() -> None:
    red = Image.new("RGB", (2, 2), "red")
    blue = Image.new("RGB", (2, 2), "blue")

    with pytest.raises(AssertionError):
        _assert_same_rgb_pixels(red, blue)


def test_release_builder_uses_only_bundled_regular_noto_sans_sc() -> None:
    builder = importlib.import_module("scripts.build_release_corpus")

    assert builder.FONT_PATH == EXPECTED_FONT_PATH
    assert builder.FONT_VARIATION == "Regular"
    assert builder.FONT_PATH.is_file()
    assert builder.FONT_PATH.with_name("OFL.txt").is_file()
    assert hashlib.sha256(builder.FONT_PATH.read_bytes()).hexdigest() == EXPECTED_FONT_SHA256
    source = inspect.getsource(builder._font)
    assert "FONT_PATH" in source
    assert "Windows\\Fonts" not in source
    assert "DejaVuSans" not in source
    assert "load_default" not in source


def test_release_builder_fails_fast_when_bundled_font_is_missing(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    builder = importlib.import_module("scripts.build_release_corpus")
    monkeypatch.setattr(builder, "FONT_PATH", tmp_path / "missing.ttf")

    with pytest.raises(FileNotFoundError):
        builder._font(24)


def test_bundled_font_has_distinct_cjk_glyphs_without_notdef() -> None:
    builder = importlib.import_module("scripts.build_release_corpus")
    font = builder._font(48)

    def fingerprint(character: str) -> tuple[tuple[int, int, int, int], bytes]:
        return font.getbbox(character), bytes(font.getmask(character))

    notdef = fingerprint("\U0010ffff")
    replacement = fingerprint("\ufffd")
    glyphs = [fingerprint(character) for character in "运营仪表盘"]

    assert all(glyph not in {notdef, replacement} for glyph in glyphs)
    assert len(set(glyphs)) == len(glyphs)


def test_release_builder_requires_a_new_output_root(tmp_path: Path) -> None:
    builder = importlib.import_module("scripts.build_release_corpus")
    existing = tmp_path / "existing"
    existing.mkdir()

    with pytest.raises(FileExistsError):
        builder.build(existing)


def test_release_builder_main_accepts_one_positional_output(tmp_path: Path) -> None:
    builder = importlib.import_module("scripts.build_release_corpus")
    output = tmp_path / "from-main"

    assert builder.main([str(output)]) is None
    assert {path.name for path in output.iterdir()} == EXPECTED_INPUT_NAMES


def test_release_builder_creates_exact_eighteen_inputs_twice(
    fresh_release_corpora: tuple[Path, Path],
) -> None:
    first, second = fresh_release_corpora

    assert {path.name for path in first.iterdir()} == EXPECTED_INPUT_NAMES
    assert {path.name for path in second.iterdir()} == EXPECTED_INPUT_NAMES


def test_fresh_release_pngs_are_valid_and_pixel_equivalent(
    fresh_release_corpora: tuple[Path, Path],
) -> None:
    first, second = fresh_release_corpora

    for name, expected_size in EXPECTED_IMAGE_SIZES.items():
        with (
            Image.open(first / name) as first_image,
            Image.open(second / name) as second_image,
            Image.open(RELEASE_ROOT / "inputs" / name) as canonical_image,
        ):
            assert first_image.format == "PNG"
            assert second_image.format == "PNG"
            assert canonical_image.format == "PNG"
            assert first_image.size == expected_size
            assert second_image.size == expected_size
            assert canonical_image.size == expected_size
            _assert_same_rgb_pixels(first_image, second_image)


def test_fresh_release_pdfs_have_two_real_pages_and_equivalent_geometry(
    fresh_release_corpora: tuple[Path, Path],
) -> None:
    first, second = fresh_release_corpora

    for name in ("13-mixed-page-sizes.pdf", "14-rotated-page.pdf", "15-high-dpi.pdf"):
        assert _pdf_inventory(first / name) == _pdf_inventory(second / name)
        assert _pdf_inventory(first / name) == _pdf_inventory(
            RELEASE_ROOT / "inputs" / name
        )
        assert len(_pdf_inventory(first / name)) == 2

    mixed_sizes = _pdf_inventory(first / "13-mixed-page-sizes.pdf")
    assert mixed_sizes[0][:2] != mixed_sizes[1][:2]
    assert _pdf_inventory(first / "14-rotated-page.pdf")[1][2] == 90

    high_dpi_reader = PdfReader(first / "15-high-dpi.pdf")
    image_sizes = [
        (int(image.image.width), int(image.image.height))
        for page in high_dpi_reader.pages
        for image in page.images
    ]
    assert len(image_sizes) == 2
    assert all(max(size) >= 2400 for size in image_sizes)


def test_fresh_release_pptx_modes_have_expected_object_inventory(
    fresh_release_corpora: tuple[Path, Path],
) -> None:
    first, second = fresh_release_corpora
    names = (
        "16-image-only.pptx",
        "17-mixed-native.pptx",
        "18-mixed-screenshot-candidates.pptx",
    )

    for name in names:
        deck = Presentation(first / name)
        assert len(deck.slides) == 4
        assert deck.slide_width * 9 == deck.slide_height * 16
        assert deck.core_properties.created == datetime(2020, 1, 1)
        assert deck.core_properties.modified == datetime(2020, 1, 1)
        assert _pptx_inventory(first / name) == _pptx_inventory(second / name)
        assert _pptx_inventory(first / name) == _pptx_inventory(
            RELEASE_ROOT / "inputs" / name
        )
        with ZipFile(first / name) as archive:
            members = archive.infolist()
            assert [member.filename for member in members] == sorted(
                member.filename for member in members
            )
            assert all(member.date_time == (1980, 1, 1, 0, 0, 0) for member in members)
            assert all(member.compress_type == ZIP_DEFLATED for member in members)

    image_only = Presentation(first / "16-image-only.pptx")
    assert all(
        len(slide.shapes) == 1
        and slide.shapes[0].shape_type == MSO_SHAPE_TYPE.PICTURE
        for slide in image_only.slides
    )

    mixed_native = Presentation(first / "17-mixed-native.pptx")
    native_shapes = [shape for slide in mixed_native.slides for shape in slide.shapes]
    assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in native_shapes)
    assert sum(shape.has_text_frame for shape in native_shapes) >= 8
    assert sum(shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for shape in native_shapes) >= 8

    screenshot_candidates = Presentation(first / "18-mixed-screenshot-candidates.pptx")
    for slide in screenshot_candidates.slides:
        assert any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in slide.shapes)
        assert any(shape.has_text_frame and shape.text for shape in slide.shapes)
        assert any(shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE for shape in slide.shapes)


def test_release_input_rejects_invalid_image_with_matching_hash(tmp_path: Path) -> None:
    case = copy.deepcopy(_manifest()["cases"][0])
    case["path"] = "invalid.png"
    path = tmp_path / case["path"]
    path.write_bytes(b"not an image")
    case["sha256"] = hashlib.sha256(path.read_bytes()).hexdigest()

    with pytest.raises(UnidentifiedImageError):
        _assert_release_input(case, tmp_path)


def test_release_input_rejects_valid_image_with_wrong_hash(tmp_path: Path) -> None:
    case = copy.deepcopy(_manifest()["cases"][0])
    case["path"] = "wrong-hash.png"
    path = tmp_path / case["path"]
    Image.new("RGB", (2, 2), "red").save(path)
    case["sha256"] = "f" * 64

    with pytest.raises(AssertionError):
        _assert_release_input(case, tmp_path)


def test_release_input_rejects_jpeg_renamed_as_png(tmp_path: Path) -> None:
    case = copy.deepcopy(_manifest()["cases"][0])
    case["path"] = "renamed.png"
    path = tmp_path / case["path"]
    Image.new("RGB", (2, 2), "red").save(path, format="JPEG")
    case["sha256"] = hashlib.sha256(path.read_bytes()).hexdigest()
    with Image.open(path) as image:
        assert image.format == "JPEG"

    with pytest.raises(AssertionError):
        _assert_release_input(case, tmp_path)


def test_release_input_rejects_multiframe_apng(tmp_path: Path) -> None:
    case = copy.deepcopy(_manifest()["cases"][0])
    case["path"] = "two-frame.png"
    path = tmp_path / case["path"]
    first = Image.new("RGBA", (2, 2), "red")
    second = Image.new("RGBA", (2, 2), "blue")
    first.save(
        path,
        format="PNG",
        save_all=True,
        append_images=[second],
        duration=100,
        loop=0,
    )
    case["sha256"] = hashlib.sha256(path.read_bytes()).hexdigest()
    with Image.open(path) as image:
        assert image.format == "PNG"
        assert image.n_frames == 2

    with pytest.raises(AssertionError):
        _assert_release_input(case, tmp_path)


def test_release_input_rejects_image_that_opens_but_fails_verify(
    tmp_path: Path,
) -> None:
    case = copy.deepcopy(_manifest()["cases"][0])
    case["path"] = "truncated.png"
    path = tmp_path / case["path"]
    Image.new("RGB", (2, 2), "red").save(path)
    path.write_bytes(path.read_bytes()[:-10])
    case["sha256"] = hashlib.sha256(path.read_bytes()).hexdigest()

    with pytest.raises((OSError, SyntaxError)):
        _assert_release_input(case, tmp_path)


def test_release_input_rejects_wrong_pdf_page_count_with_matching_hash(
    tmp_path: Path,
) -> None:
    case = copy.deepcopy(_manifest()["cases"][12])
    case["path"] = "one-page.pdf"
    path = tmp_path / case["path"]
    writer = PdfWriter()
    writer.add_blank_page(width=100, height=100)
    with path.open("wb") as stream:
        writer.write(stream)
    case["sha256"] = hashlib.sha256(path.read_bytes()).hexdigest()

    with pytest.raises(AssertionError):
        _assert_release_input(case, tmp_path)


def test_release_input_rejects_wrong_pptx_slide_count_with_matching_hash(
    tmp_path: Path,
) -> None:
    case = copy.deepcopy(_manifest()["cases"][15])
    case["path"] = "one-slide.pptx"
    path = tmp_path / case["path"]
    deck = Presentation()
    deck.slides.add_slide(deck.slide_layouts[6])
    deck.save(path)
    case["sha256"] = hashlib.sha256(path.read_bytes()).hexdigest()

    with pytest.raises(AssertionError):
        _assert_release_input(case, tmp_path)


def test_release_inputs_exist_and_match_manifest_sha256() -> None:
    cases = _manifest()["cases"]

    for case in cases:
        _assert_release_input(case, RELEASE_ROOT)


def _write_benchmark_plan(path: Path, value: dict[str, object]) -> None:
    path.write_text(json.dumps(value, sort_keys=True), encoding="utf-8")


def _candidate_response() -> dict[str, object]:
    return {
        "candidate": {
            "candidate_id": "candidate_001",
            "page_id": "page_001",
            "source_shape_id": "2",
            "source_object_sha256": "1" * 64,
            "image_sha256": "2" * 64,
        }
    }


def _candidate_plan() -> dict[str, object]:
    return {
        "schema_version": 1,
        "kind": "candidate_decision",
        "page_id": "page_001",
        "candidate_id": "candidate_001",
        "source_shape_id": "2",
        "source_object_sha256": "1" * 64,
        "image_sha256": "2" * 64,
        "decision": "replace",
        "confidence": 0.99,
        "category": "full_slide_screenshot",
        "evidence": ["full-slide raster bound to the candidate hashes"],
    }


def _component_request() -> dict[str, object]:
    return {
        "kind": "component_request",
        "page_id": "page_001",
        "provider": "host",
        "repair_round": 1,
        "request_sha256": "3" * 64,
        "graph_sha256": "4" * 64,
    }


def _bound_component_request(run_dir: Path) -> tuple[dict[str, object], dict[str, object]]:
    request_dir = run_dir / "pages/page_001/reconstruction/agent/round-01"
    request_dir.mkdir(parents=True)
    graph_document = {
        "nodes": [
            {
                "id": "component_0001",
                "kind": "child",
                "parent_id": None,
                "state": "pending",
                "mask": "masks/component_0001.png",
                "mask_sha256": "1" * 64,
                "bbox": [0, 0, 5, 5],
                "z_index": 0,
                "text_ids": [],
            }
        ]
    }
    graph = (
        json.dumps(graph_document, ensure_ascii=False, indent=2, sort_keys=True).encode(
            "utf-8"
        )
        + b"\n"
    )
    graph_sha256 = hashlib.sha256(graph).hexdigest()
    (request_dir / "component-graph.json").write_bytes(graph)
    evidence_names = (
        "source.png",
        "numbered-masks.png",
        "ocr-overlay.png",
        "component-isolation.png",
        "ownership.png",
        "reconstructed.png",
        "difference.png",
        "unexplained-mask.png",
        "component-graph.json",
        "quality-report.json",
        "presentation-manifest.json",
    )
    request = {
        "schema_version": 1,
        "page_id": "page_001",
        "provider": "host",
        "repair_round": 1,
        "source_sha256": "0" * 64,
        "graph_sha256": graph_sha256,
        "candidate_ids": ["component_0001"],
        "frozen_ids": [],
        "evidence": {
            name: {
                "path": name,
                "sha256": graph_sha256 if name == "component-graph.json" else "0" * 64,
            }
            for name in evidence_names
        },
        "review_evidence": [
            name
            for name in evidence_names
            if name not in {"component-graph.json", "presentation-manifest.json"}
        ],
    }
    payload = (
        json.dumps(request, ensure_ascii=False, indent=2, sort_keys=True).encode("utf-8")
        + b"\n"
    )
    request_path = request_dir / "component_agent_request.json"
    request_path.write_bytes(payload)
    response = {
        "kind": "component_request",
        "page_id": "page_001",
        "provider": "host",
        "repair_round": 1,
        "request_sha256": hashlib.sha256(payload).hexdigest(),
        "request_path": str(request_path.resolve()),
    }
    plan = {
        **_component_plan(),
        "request_sha256": response["request_sha256"],
        "graph_sha256": graph_sha256,
    }
    return response, plan


def _component_plan() -> dict[str, object]:
    return {
        "schema_version": 1,
        "kind": "component_plan",
        "page_id": "page_001",
        "provider": "host",
        "repair_round": 1,
        "request_sha256": "3" * 64,
        "graph_sha256": "4" * 64,
        "actions": [
            {
                "action": "accept",
                "object_ids": ["component_0001"],
                "parameters": {},
                "confidence": 0.99,
                "evidence": ["component boundary matches the source"],
            }
        ],
    }


def _runner_case() -> dict[str, object]:
    return {
        "id": "pptx-mixed-screenshot-candidates",
        "kind": "pptx",
        "path": "inputs/18-mixed-screenshot-candidates.pptx",
        "agent_provider": "host",
    }


def _install_runner_plans(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    *,
    component_plan: dict[str, object] | None = None,
) -> Path:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = tmp_path / "plans"
    plans.mkdir()
    _write_benchmark_plan(
        plans / "pptx-mixed-screenshot-candidates--candidate.json",
        _candidate_plan(),
    )
    if component_plan is not None:
        _write_benchmark_plan(
            plans / "pptx-mixed-screenshot-candidates--component.json",
            component_plan,
        )
    monkeypatch.setattr(runner, "PLAN_ROOT", plans)
    return plans


def test_release_runner_locks_the_complete_host_protocol(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = _install_runner_plans(tmp_path, monkeypatch)
    calls: list[list[str]] = []
    run_next_count = 0
    run_execute_count = 0
    agent_next_count = 0
    expected_component_plan: dict[str, object] | None = None
    observed_plans: list[tuple[str, bool]] = []

    def fake_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        nonlocal agent_next_count, expected_component_plan, run_next_count, run_execute_count
        assert cwd == tmp_path / "workspace"
        calls.append(arguments)
        command = arguments[:2]
        if arguments[0] == "prepare":
            run_dir = Path(arguments[arguments.index("--run-dir") + 1])
            return subprocess.CompletedProcess(
                arguments,
                0,
                json.dumps({"run_dir": str(run_dir.resolve()), "status": "prepared"}),
                "",
            )
        if command == ["run", "next"]:
            run_next_count += 1
            value = _candidate_response() if run_next_count == 1 else {"candidate": None}
        elif command == ["decision", "record"]:
            value = {"status": "recorded"}
        elif command == ["run", "execute"]:
            run_execute_count += 1
            value = (
                {"status": "awaiting_agent"}
                if run_execute_count == 1
                else {
                    "status": "completed",
                    "page_results": [
                        {"page_id": "page_001", "status": "validated"}
                    ],
                }
            )
        elif command == ["agent", "next"]:
            agent_next_count += 1
            if agent_next_count == 1:
                challenge = (
                    tmp_path
                    / "workspace/pptx-mixed-screenshot-candidates/run/host-challenge/challenge.png"
                )
                challenge.parent.mkdir(parents=True)
                image = Image.new("RGB", (240, 120), "white")
                for left in (24, 97, 170):
                    for x in range(left, left + 45):
                        for y in range(20, 65):
                            image.putpixel((x, y), (217, 72, 95))
                image.save(challenge)
                value = {
                    "kind": "capability_handshake",
                    "challenge_id": "5" * 64,
                    "image_path": str(challenge.resolve()),
                    "required_capabilities": [
                        "vision",
                        "local_file_read",
                        "tool_use",
                        "structured_json",
                    ],
                }
            else:
                run_dir = Path(arguments[2])
                value, expected_component_plan = _bound_component_request(run_dir)
                _write_benchmark_plan(
                    plans / "pptx-mixed-screenshot-candidates--component.json",
                    {
                        **expected_component_plan,
                        "request_sha256": "8" * 64,
                        "graph_sha256": "9" * 64,
                    },
                )
        elif command == ["agent", "record"]:
            submitted = json.loads(
                Path(arguments[arguments.index("--plan") + 1]).read_text(
                    encoding="utf-8"
                )
            )
            if submitted["kind"] == "host_capability_response":
                assert submitted == {
                    "schema_version": 1,
                    "kind": "host_capability_response",
                    "challenge_id": "5" * 64,
                    "observed": {"shape": "square", "color": "#d9485f", "count": 3},
                }
                value = {
                    "capabilities": [
                        "vision",
                        "local_file_read",
                        "tool_use",
                        "structured_json",
                    ],
                    "status": "capabilities_recorded",
                }
            else:
                assert expected_component_plan is not None
                assert submitted == {
                    key: value
                    for key, value in expected_component_plan.items()
                    if key != "graph_sha256"
                }
                plan_path = Path(arguments[2]) / "recorded-component-plan.json"
                plan_path.write_bytes(
                    Path(arguments[arguments.index("--plan") + 1]).read_bytes()
                )
                value = {
                    "plan_path": str(plan_path.resolve()),
                    "recovered": False,
                    "status": "recorded",
                }
        else:
            raise AssertionError(arguments)
        return subprocess.CompletedProcess(arguments, 0, json.dumps(value), "")

    result = runner.run_case(
        _runner_case(),
        workspace=tmp_path / "workspace",
        command=fake_command,
        allow_compatible_component_bindings=True,
        plan_observer=lambda filename, plan, rebound: observed_plans.append(
            (filename, rebound)
        ),
    )

    run_dir = str((tmp_path / "workspace" / "pptx-mixed-screenshot-candidates" / "run").resolve())
    input_path = str((RELEASE_ROOT / "inputs/18-mixed-screenshot-candidates.pptx").resolve())
    assert calls == [
        [
            "prepare",
            input_path,
            "--run-dir",
            run_dir,
            "--output",
            str((tmp_path / "workspace" / "pptx-mixed-screenshot-candidates" / "output.pptx").resolve()),
            "--slide-size",
            "original",
            "--agent-provider",
            "host",
        ],
        ["run", "next", run_dir],
        [
            "decision",
            "record",
            run_dir,
            "--page",
            "page_001",
            "--object",
            "2",
            "--decision",
            "replace",
            "--confidence",
            "0.99",
            "--category",
            "full_slide_screenshot",
            "--evidence",
            "full-slide raster bound to the candidate hashes",
        ],
        ["run", "next", run_dir],
        ["run", "execute", run_dir],
        ["agent", "next", run_dir],
        ["agent", "record", run_dir, "--plan", calls[6][-1]],
        ["agent", "next", run_dir],
        ["agent", "record", run_dir, "--plan", calls[8][-1]],
        ["run", "execute", run_dir],
    ]
    assert result.case_id == "pptx-mixed-screenshot-candidates"
    assert result.run_dir == run_dir
    assert result.pages == [{"page_id": "page_001", "status": "validated"}]
    assert type(result.duration_ms) is int and result.duration_ms >= 0
    assert observed_plans == [
        ("pptx-mixed-screenshot-candidates--component.json", True),
    ]


def test_release_runner_rejects_invalid_capability_png_size(tmp_path: Path) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    image_path = tmp_path / "challenge.png"
    Image.new("RGB", (241, 120), "white").save(image_path)

    with pytest.raises(runner.BenchmarkFailure, match="invalid_response"):
        runner._observe_capability_challenge(image_path, tmp_path)


def test_release_runner_rejects_out_of_bounds_capability_png(tmp_path: Path) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    run_dir = tmp_path / "run"
    run_dir.mkdir()
    image_path = tmp_path / "challenge.png"
    image = Image.new("RGB", (240, 120), "white")
    for left in (24, 97, 170):
        for x in range(left, left + 45):
            for y in range(20, 65):
                image.putpixel((x, y), (217, 72, 95))
    image.save(image_path)

    with pytest.raises(runner.BenchmarkFailure, match="invalid_response"):
        runner._observe_capability_challenge(image_path, run_dir)


def test_release_runner_requires_bound_component_request_path(tmp_path: Path) -> None:
    runner = importlib.import_module("scripts.release_benchmark")

    with pytest.raises(runner.BenchmarkFailure, match="invalid_response"):
        runner._component_binding(_component_request(), tmp_path)


def test_release_runner_preserves_preexisting_host_plan_file(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    workspace = tmp_path / "workspace"
    run_dir = workspace / "case/run"
    run_dir.mkdir(parents=True)
    sentinel = workspace / ".host-plan-001.json"
    sentinel.write_text("owner data", encoding="utf-8")
    submitted: list[Path] = []

    def fake_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        path = Path(arguments[arguments.index("--plan") + 1])
        assert path.parent == run_dir.parent
        assert path != sentinel
        assert path.is_file()
        submitted.append(path)
        artifact = run_dir / "recorded-plan.json"
        artifact.write_bytes(path.read_bytes())
        response = {
            "plan_path": str(artifact.resolve()),
            "recovered": False,
            "status": "recorded",
        }
        return subprocess.CompletedProcess(arguments, 0, json.dumps(response), "")

    runner._record_host_document(
        fake_command,
        {"kind": "component_plan"},
        run_dir,
        workspace,
        1,
    )

    assert sentinel.read_text(encoding="utf-8") == "owner data"
    assert len(submitted) == 1
    assert not submitted[0].exists()


def test_release_runner_rejects_replaced_host_submission_without_deleting_it(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    workspace = tmp_path / "workspace"
    run_dir = workspace / "case/run"
    run_dir.mkdir(parents=True)
    replacement = b'{"kind":"replacement"}\n'
    submitted: list[Path] = []
    original_metadata: dict[Path, os.stat_result] = {}
    original_lstat = Path.lstat
    original_read_regular_file = runner._read_regular_file

    def reused_identity_lstat(path: Path) -> os.stat_result:
        if path in original_metadata:
            return original_metadata[path]
        return original_lstat(path)

    def read_replacement(
        path: Path, limit: int, *, require_single_link: bool
    ) -> bytes:
        if path in original_metadata:
            return path.read_bytes()
        return original_read_regular_file(
            path, limit, require_single_link=require_single_link
        )

    monkeypatch.setattr(Path, "lstat", reused_identity_lstat)
    monkeypatch.setattr(runner, "_read_regular_file", read_replacement)

    def replacing_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        path = Path(arguments[arguments.index("--plan") + 1])
        original_metadata[path] = original_lstat(path)
        submitted.append(path)
        path.unlink()
        path.write_bytes(replacement)
        artifact = run_dir / "recorded-plan.json"
        artifact.write_bytes(replacement)
        response = {
            "plan_path": str(artifact.resolve()),
            "recovered": False,
            "status": "recorded",
        }
        return subprocess.CompletedProcess(arguments, 0, json.dumps(response), "")

    with pytest.raises(runner.BenchmarkFailure, match="invalid_plan"):
        runner._record_host_document(
            replacing_command,
            {"kind": "component_plan"},
            run_dir,
            workspace,
            1,
        )

    assert len(submitted) == 1
    assert submitted[0].read_bytes() == replacement


def test_release_runner_rejects_mismatched_recorded_host_plan(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    workspace = tmp_path / "workspace"
    run_dir = workspace / "case/run"
    run_dir.mkdir(parents=True)

    def mismatched_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        artifact = run_dir / "recorded-plan.json"
        artifact.write_text('{"kind":"different"}\n', encoding="utf-8")
        response = {
            "plan_path": str(artifact.resolve()),
            "recovered": False,
            "status": "recorded",
        }
        return subprocess.CompletedProcess(arguments, 0, json.dumps(response), "")

    with pytest.raises(runner.BenchmarkFailure, match="invalid_plan"):
        runner._record_host_document(
            mismatched_command,
            {"kind": "component_plan"},
            run_dir,
            workspace,
            1,
        )


@pytest.mark.parametrize(
    ("returncode", "stdout"),
    [(9, "secret stdout"), (0, "not json"), (0, '{"candidate":NaN}')],
)
def test_release_runner_normalizes_command_and_json_failures_without_leaks(
    tmp_path: Path,
    returncode: int,
    stdout: str,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")

    def broken_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        return subprocess.CompletedProcess(arguments, returncode, stdout, "private stderr")

    with pytest.raises(runner.BenchmarkFailure) as caught:
        runner.run_case(
            _runner_case(), workspace=tmp_path / "workspace", command=broken_command
        )

    message = str(caught.value)
    assert message in {"command_failed", "invalid_json"}
    assert "secret stdout" not in message
    assert "private stderr" not in message
    assert str(tmp_path.resolve()) not in message


def test_release_runner_reports_safe_command_failure_coordinates(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    secret = str((tmp_path / "private-input.png").resolve())
    stderr = (
        "Traceback (most recent call last):\n"
        f'  File "{secret}", line 7, in private_function\n'
        "  File \"C:\\runner\\site-packages\\image2editable\\"
        "component_repair.py\", line 1112, in advance_round\n"
        f"RuntimeError: failed while reading {secret}\n"
    )

    def broken_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        return subprocess.CompletedProcess(arguments, 9, "secret stdout", stderr)

    with pytest.raises(runner.BenchmarkFailure, match="command_failed") as caught:
        runner._call(
            broken_command,
            ["run", "execute", str(tmp_path / "run")],
            cwd=tmp_path,
        )

    assert caught.value.details == {
        "stage": "run_execute",
        "returncode": 9,
        "exception_type": "RuntimeError",
        "frame": {
            "module": "image2editable.component_repair",
            "line": 1112,
        },
    }
    serialized = json.dumps(caught.value.details)
    assert secret not in serialized
    assert "failed while reading" not in serialized
    assert "secret stdout" not in serialized


def test_release_runner_rejects_forged_stderr_diagnostic_names(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    secret = "MODEL_SECRET"
    stderr = (
        f'  File "C:\\private\\image2editable\\{secret}.py", '
        f"line 7, in {secret}\n"
        f"Prompt{secret}Error: private\n"
    )

    def broken_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        return subprocess.CompletedProcess(arguments, 9, "", stderr)

    with pytest.raises(runner.BenchmarkFailure, match="command_failed") as caught:
        runner._call(
            broken_command,
            ["run", "execute", str(tmp_path / "run")],
            cwd=tmp_path,
        )

    assert caught.value.details == {"stage": "run_execute", "returncode": 9}
    assert secret not in json.dumps(caught.value.details)


def test_release_runner_reports_safe_invocation_and_invalid_json_stages(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")

    def invocation_failure(arguments: list[str], *, cwd: Path) -> None:
        raise FileNotFoundError(str(tmp_path / "private-command.exe"))

    with pytest.raises(runner.BenchmarkFailure, match="command_failed") as caught:
        runner._call(
            invocation_failure,
            ["prepare", str(tmp_path / "private-input.png")],
            cwd=tmp_path,
        )
    assert caught.value.details == {
        "stage": "prepare",
        "exception_type": "FileNotFoundError",
    }

    def invalid_json(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        return subprocess.CompletedProcess(arguments, 0, "private not-json", "")

    with pytest.raises(runner.BenchmarkFailure, match="invalid_json") as caught:
        runner._call(
            invalid_json,
            ["agent", "next", str(tmp_path / "run")],
            cwd=tmp_path,
        )
    assert caught.value.details == {"stage": "agent_next", "returncode": 0}


def test_release_runner_rejects_missing_component_plan(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    _install_runner_plans(tmp_path, monkeypatch)
    responses = iter(
        [
            {"status": "prepared"},
            _candidate_response(),
            {"status": "recorded"},
            {"candidate": None},
            {"status": "awaiting_agent"},
            None,
        ]
    )

    def fake_command(arguments: list[str], *, cwd: Path) -> subprocess.CompletedProcess[str]:
        value = next(responses)
        if arguments[0] == "prepare":
            value["run_dir"] = str(
                Path(arguments[arguments.index("--run-dir") + 1]).resolve()
            )
        elif arguments[:2] == ["agent", "next"]:
            value, _ = _bound_component_request(Path(arguments[2]))
        return subprocess.CompletedProcess(arguments, 0, json.dumps(value), "")

    with pytest.raises(runner.BenchmarkFailure, match="missing_plan"):
        runner.run_case(
            _runner_case(), workspace=tmp_path / "workspace", command=fake_command
        )


@pytest.mark.parametrize(
    ("mutation", "message"),
    [
        ({"request_sha256": "5" * 64}, "stale_plan"),
        ({"graph_sha256": "5" * 64}, "stale_plan"),
        ({"page_id": "page_999"}, "mismatched_plan"),
        ({"repair_round": 2}, "mismatched_plan"),
    ],
)
def test_release_runner_rejects_stale_or_mismatched_component_plan(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    mutation: dict[str, object],
    message: str,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plan = {**_component_plan(), **mutation}
    _install_runner_plans(tmp_path, monkeypatch, component_plan=plan)

    with pytest.raises(runner.BenchmarkFailure, match=message):
        runner._select_component_plan("pptx-mixed-screenshot-candidates", _component_request())


def test_release_runner_reports_safe_component_identity_mismatch(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    _install_runner_plans(
        tmp_path,
        monkeypatch,
        component_plan={**_component_plan(), "page_id": "page_999"},
    )

    with pytest.raises(runner.BenchmarkFailure, match="mismatched_plan") as caught:
        runner._select_component_plan(
            "pptx-mixed-screenshot-candidates", _component_request()
        )

    assert caught.value.details == {
        "stage": "component_plan_identity",
        "page_id": "page_001",
        "repair_round": 1,
        "expected_identities": [{"page_id": "page_999", "repair_round": 1}],
    }


def test_release_runner_reports_safe_candidate_identity_mismatch(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = _install_runner_plans(tmp_path, monkeypatch)
    _write_benchmark_plan(
        plans / "pptx-mixed-screenshot-candidates--candidate.json",
        {**_candidate_plan(), "source_shape_id": "999"},
    )

    with pytest.raises(runner.BenchmarkFailure, match="mismatched_plan") as caught:
        runner._select_candidate_plan(
            "pptx-mixed-screenshot-candidates",
            _candidate_response()["candidate"],
        )

    assert caught.value.details == {
        "stage": "candidate_plan_identity",
        "page_id": "page_001",
        "candidate_id": "candidate_001",
        "source_shape_id": "2",
        "expected_identities": [
            {
                "page_id": "page_001",
                "candidate_id": "candidate_001",
                "source_shape_id": "999",
            }
        ],
    }


@pytest.mark.parametrize(
    ("kind", "mutation"),
    [
        ("candidate", {"page_id": "C:\\private\\input.png"}),
        ("candidate", {"candidate_id": "https://private.example/secret"}),
        ("candidate", {"source_shape_id": "document text"}),
        ("component", {"page_id": "C:\\private\\input.png"}),
    ],
)
def test_release_runner_rejects_untrusted_protocol_identifiers_without_echo(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    kind: str,
    mutation: dict[str, object],
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    _install_runner_plans(tmp_path, monkeypatch, component_plan=_component_plan())
    value = (
        {**_candidate_response()["candidate"], **mutation}
        if kind == "candidate"
        else {**_component_request(), **mutation}
    )
    select = (
        runner._select_candidate_plan
        if kind == "candidate"
        else runner._select_component_plan
    )

    with pytest.raises(runner.BenchmarkFailure, match="invalid_response") as caught:
        select("pptx-mixed-screenshot-candidates", value)

    serialized = json.dumps(caught.value.details)
    assert next(iter(mutation.values())) not in serialized


def test_release_runner_reports_safe_component_stale_plan_bindings(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    stale = {**_component_plan(), "request_sha256": "5" * 64}
    _install_runner_plans(tmp_path, monkeypatch, component_plan=stale)

    with pytest.raises(runner.BenchmarkFailure, match="stale_plan") as caught:
        runner._select_component_plan(
            "pptx-mixed-screenshot-candidates", _component_request()
        )

    assert caught.value.details == {
        "stage": "component_plan",
        "page_id": "page_001",
        "repair_round": 1,
        "actual_request_sha256": "3" * 64,
        "actual_graph_sha256": "4" * 64,
        "expected_bindings": [
            {"request_sha256": "5" * 64, "graph_sha256": "4" * 64}
        ],
    }


def test_release_runner_reports_safe_candidate_stale_plan_bindings(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    stale = {**_candidate_plan(), "source_object_sha256": "5" * 64}
    plans = _install_runner_plans(tmp_path, monkeypatch)
    _write_benchmark_plan(
        plans / "pptx-mixed-screenshot-candidates--candidate.json", stale
    )

    with pytest.raises(runner.BenchmarkFailure, match="stale_plan") as caught:
        runner._select_candidate_plan(
            "pptx-mixed-screenshot-candidates", _candidate_response()["candidate"]
        )

    assert caught.value.details == {
        "stage": "candidate_plan",
        "page_id": "page_001",
        "candidate_id": "candidate_001",
        "source_shape_id": "2",
        "actual_source_object_sha256": "1" * 64,
        "actual_image_sha256": "2" * 64,
        "expected_bindings": [
            {"source_object_sha256": "5" * 64, "image_sha256": "2" * 64}
        ],
    }


def test_release_runner_official_component_rebinding_keeps_candidate_binding_exact(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = _install_runner_plans(tmp_path, monkeypatch)
    stale = {**_candidate_plan(), "source_object_sha256": "5" * 64}
    _write_benchmark_plan(
        plans / "pptx-mixed-screenshot-candidates--candidate.json", stale
    )

    def fake_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        if arguments[0] == "prepare":
            run_dir = Path(arguments[arguments.index("--run-dir") + 1])
            run_dir.mkdir()
            value = {"run_dir": str(run_dir.resolve()), "status": "prepared"}
        elif arguments[:2] == ["run", "next"]:
            value = _candidate_response()
        else:
            raise AssertionError(arguments)
        return subprocess.CompletedProcess(arguments, 0, json.dumps(value), "")

    with pytest.raises(runner.BenchmarkFailure, match="stale_plan") as caught:
        runner.run_case(
            _runner_case(),
            workspace=tmp_path / "workspace",
            command=fake_command,
            allow_compatible_component_bindings=True,
        )

    assert caught.value.details is not None
    assert caught.value.details["stage"] == "candidate_plan"


def test_release_runner_reports_repeated_component_request_identity(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = _install_runner_plans(tmp_path, monkeypatch)
    execute_count = 0
    bound_response: dict[str, object] | None = None

    def fake_command(
        arguments: list[str], *, cwd: Path
    ) -> subprocess.CompletedProcess[str]:
        nonlocal bound_response, execute_count
        command = arguments[:2]
        if arguments[0] == "prepare":
            run_dir = Path(arguments[arguments.index("--run-dir") + 1])
            value = {"run_dir": str(run_dir.resolve()), "status": "prepared"}
        elif command == ["run", "execute"]:
            execute_count += 1
            value = {"status": "awaiting_agent"}
        elif command == ["agent", "next"]:
            if bound_response is None:
                bound_response, plan = _bound_component_request(Path(arguments[2]))
                plan_path = plans / "pptx-mixed-screenshot-candidates--component.json"
                _write_benchmark_plan(plan_path, plan)
            value = dict(bound_response)
        elif command == ["agent", "record"]:
            submitted = Path(arguments[arguments.index("--plan") + 1])
            artifact = Path(arguments[2]) / "recorded-component-plan.json"
            artifact.write_bytes(submitted.read_bytes())
            value = {
                "plan_path": str(artifact.resolve()),
                "recovered": False,
                "status": "recorded",
            }
        else:
            raise AssertionError(arguments)
        return subprocess.CompletedProcess(arguments, 0, json.dumps(value), "")

    case = {**_runner_case(), "kind": "image"}
    with pytest.raises(runner.BenchmarkFailure, match="stale_plan") as caught:
        runner.run_case(
            case,
            workspace=tmp_path / "workspace",
            command=fake_command,
            allow_stale_bindings=True,
        )

    assert execute_count == 2
    assert caught.value.details is not None
    assert set(caught.value.details) == {
        "stage",
        "page_id",
        "repair_round",
        "request_sha256",
        "graph_sha256",
    }
    assert caught.value.details["stage"] == "component_request_repeated"
    assert caught.value.details["page_id"] == "page_001"
    assert caught.value.details["repair_round"] == 1
    assert re.fullmatch(r"[0-9a-f]{64}", caught.value.details["request_sha256"])
    assert re.fullmatch(r"[0-9a-f]{64}", caught.value.details["graph_sha256"])


def test_release_runner_diagnostic_rebinds_only_component_hashes(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    stale = {**_component_plan(), "request_sha256": "5" * 64}
    _install_runner_plans(tmp_path, monkeypatch, component_plan=stale)
    run_dir = tmp_path / "run"
    response, _ = _bound_component_request(run_dir)
    request = runner._component_binding(response, run_dir)

    selection = runner._resolve_component_plan(
        "pptx-mixed-screenshot-candidates",
        request,
        allow_stale_binding=True,
        graph=request["_component_graph"],
    )

    assert selection.filename == (
        "pptx-mixed-screenshot-candidates--component.json"
    )
    assert selection.plan == stale
    assert selection.rebound_plan == {
        **stale,
        "request_sha256": request["request_sha256"],
        "graph_sha256": request["graph_sha256"],
    }
    assert selection.rebound_plan["actions"] == stale["actions"]


def test_release_runner_rejects_stale_component_without_complete_contract(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    stale = {**_component_plan(), "request_sha256": "5" * 64}
    _install_runner_plans(tmp_path, monkeypatch, component_plan=stale)
    request = {
        **_component_request(),
        "_component_request": {"page_id": "page_001"},
        "_component_graph": {"nodes": []},
    }

    with pytest.raises(runner.BenchmarkFailure, match="incompatible_plan"):
        runner._resolve_component_plan(
            "pptx-mixed-screenshot-candidates",
            request,
            allow_stale_binding=True,
            graph=request["_component_graph"],
        )


def test_release_runner_diagnostic_rebinds_only_candidate_hashes(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    stale = {**_candidate_plan(), "source_object_sha256": "5" * 64}
    plans = _install_runner_plans(tmp_path, monkeypatch)
    _write_benchmark_plan(
        plans / "pptx-mixed-screenshot-candidates--candidate.json", stale
    )

    selection = runner._resolve_candidate_plan(
        "pptx-mixed-screenshot-candidates",
        _candidate_response()["candidate"],
        allow_stale_binding=True,
    )

    assert selection.filename == "pptx-mixed-screenshot-candidates--candidate.json"
    assert selection.plan == stale
    assert selection.rebound_plan == {
        **stale,
        "source_object_sha256": "1" * 64,
        "image_sha256": "2" * 64,
    }
    assert selection.rebound_plan["decision"] == stale["decision"]


def test_release_runner_diagnostic_rejects_duplicate_identity_plans(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = _install_runner_plans(
        tmp_path,
        monkeypatch,
        component_plan={**_component_plan(), "request_sha256": "5" * 64},
    )
    _write_benchmark_plan(
        plans / "pptx-mixed-screenshot-candidates--component-copy.json",
        {**_component_plan(), "request_sha256": "6" * 64},
    )
    run_dir = tmp_path / "run"
    response, _ = _bound_component_request(run_dir)
    request = runner._component_binding(response, run_dir)

    with pytest.raises(runner.BenchmarkFailure, match="duplicate_plan"):
        runner._resolve_component_plan(
            "pptx-mixed-screenshot-candidates",
            request,
            allow_stale_binding=True,
            graph=request["_component_graph"],
        )


def test_release_runner_uses_graph_selector_for_compatible_plan_families(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    selected_filename = "pptx-mixed-screenshot-candidates--component.json"
    alternate_filename = "pptx-mixed-screenshot-candidates--component-copy.json"
    plans = _install_runner_plans(
        tmp_path,
        monkeypatch,
        component_plan={**_component_plan(), "request_sha256": "5" * 64},
    )
    _write_benchmark_plan(
        plans / alternate_filename,
        {**_component_plan(), "request_sha256": "6" * 64},
    )
    monkeypatch.setattr(
        runner,
        "_COMPONENT_PLAN_FAMILIES",
        (
            (
                "component_0001",
                "component_0002",
                frozenset({selected_filename}),
            ),
            (
                "component_0002",
                "component_0001",
                frozenset({alternate_filename}),
            ),
        ),
    )
    monkeypatch.setattr(
        runner,
        "_COMPONENT_PLAN_ROUND_SIGNATURES",
        {
            1: (
                ((0, 0, 5, 5), "child", "pending"),
                ((10, 20, 15, 25), "child", "inactive"),
            )
        },
    )
    run_dir = tmp_path / "run"
    response, _ = _bound_component_request(run_dir)
    request = runner._component_binding(response, run_dir)
    request["_component_graph"]["nodes"].append(
        {
            "id": "component_0002",
            "kind": "child",
            "state": "inactive",
            "bbox": [10, 20, 15, 25],
        }
    )

    selection = runner._resolve_component_plan(
        "pptx-mixed-screenshot-candidates",
        request,
        allow_stale_binding=True,
        graph=request["_component_graph"],
    )

    assert selection.filename == selected_filename
    assert selection.rebound_plan == {
        **selection.plan,
        "request_sha256": request["request_sha256"],
        "graph_sha256": request["graph_sha256"],
    }


@pytest.mark.parametrize("mutation", ["unreviewed_bbox", "unregistered_plan"])
def test_release_runner_rejects_unreviewed_compatible_plan_family(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    mutation: str,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    selected_filename = "pptx-mixed-screenshot-candidates--component.json"
    alternate_filename = "pptx-mixed-screenshot-candidates--component-copy.json"
    plans = _install_runner_plans(
        tmp_path,
        monkeypatch,
        component_plan={**_component_plan(), "request_sha256": "5" * 64},
    )
    _write_benchmark_plan(
        plans / alternate_filename,
        {**_component_plan(), "request_sha256": "6" * 64},
    )
    monkeypatch.setattr(
        runner,
        "_COMPONENT_PLAN_FAMILIES",
        (
            (
                "component_0001",
                "component_0002",
                frozenset({selected_filename}),
            ),
            (
                "component_0002",
                "component_0001",
                frozenset({alternate_filename}),
            ),
        ),
    )
    monkeypatch.setattr(
        runner,
        "_COMPONENT_PLAN_ROUND_SIGNATURES",
        {
            1: (
                ((0, 0, 5, 5), "child", "pending"),
                ((10, 20, 15, 25), "child", "inactive"),
            )
        },
    )
    run_dir = tmp_path / "run"
    response, _ = _bound_component_request(run_dir)
    request = runner._component_binding(response, run_dir)
    request["_component_graph"]["nodes"].append(
        {
            "id": "component_0002",
            "kind": "child",
            "state": "inactive",
            "bbox": [10, 20, 15, 25],
        }
    )
    if mutation == "unreviewed_bbox":
        request["_component_graph"]["nodes"][0]["bbox"] = [0, 5, 5, 10]
    else:
        _write_benchmark_plan(
            plans / "pptx-mixed-screenshot-candidates--component-third.json",
            {**_component_plan(), "request_sha256": "7" * 64},
        )

    with pytest.raises(runner.BenchmarkFailure, match="duplicate_plan"):
        runner._resolve_component_plan(
            "pptx-mixed-screenshot-candidates",
            request,
            allow_stale_binding=True,
            graph=request["_component_graph"],
        )


@pytest.mark.parametrize(
    ("top_id", "bottom_id", "expected_filenames"),
    [
        (
            "component_0003",
            "component_0004",
            (
                "image-non-16-9-infographic--component-round-01-20260823.json",
                "image-non-16-9-infographic--component-round-02-20260823.json",
                "image-non-16-9-infographic--component-round-03-20260823.json",
                "image-non-16-9-infographic--component-round-04.json",
            ),
        ),
        (
            "component_0004",
            "component_0003",
            (
                "image-non-16-9-infographic--component-round-01.json",
                "image-non-16-9-infographic--component-round-02.json",
                "image-non-16-9-infographic--component-round-03.json",
                "image-non-16-9-infographic--component-round-04-bca.json",
            ),
        ),
    ],
)
def test_release_infographic_plan_families_match_reviewed_graph_layouts(
    top_id: str,
    bottom_id: str,
    expected_filenames: tuple[str, ...],
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    case_entries = runner._case_plan_entries("image-non-16-9-infographic")
    signatures = (
        (
            ((61, 316, 940, 540), "child", "pending"),
            ((62, 887, 939, 1110), "child", "pending"),
        ),
        (
            ((61, 316, 940, 540), "child", "pending"),
            ((62, 887, 939, 1110), "child", "inactive"),
        ),
        (
            ((61, 316, 940, 540), "parent", "pending"),
            ((62, 887, 939, 1110), "child", "inactive"),
        ),
        (
            ((60, 315, 941, 541), "parent", "pending"),
            ((62, 887, 939, 1110), "child", "inactive"),
        ),
    )
    for repair_round, expected_filename in enumerate(expected_filenames, start=1):
        entries = [
            (filename, plan)
            for filename, plan in case_entries
            if plan.get("kind") == "component_plan"
            and plan.get("page_id") == "page_001"
            and plan.get("repair_round") == repair_round
        ]
        top, bottom = signatures[repair_round - 1]
        selection = runner._select_component_plan_family(
            entries,
            {
                "nodes": [
                    {
                        "id": top_id,
                        "bbox": list(top[0]),
                        "kind": top[1],
                        "state": top[2],
                    },
                    {
                        "id": bottom_id,
                        "bbox": list(bottom[0]),
                        "kind": bottom[1],
                        "state": bottom[2],
                    },
                ]
            },
        )

        assert selection is not None
        assert selection[0] == expected_filename


def test_release_runner_diagnostic_prefers_exact_component_binding(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = _install_runner_plans(
        tmp_path,
        monkeypatch,
        component_plan={**_component_plan(), "request_sha256": "5" * 64,
                        "graph_sha256": "6" * 64},
    )
    exact = {**_component_plan()}
    _write_benchmark_plan(
        plans / "pptx-mixed-screenshot-candidates--component-exact.json",
        exact,
    )

    selection = runner._resolve_component_plan(
        "pptx-mixed-screenshot-candidates",
        _component_request(),
        allow_stale_binding=True,
    )

    assert selection.filename == (
        "pptx-mixed-screenshot-candidates--component-exact.json"
    )
    assert selection.plan == exact
    assert selection.rebound_plan is None


def test_release_runner_diagnostic_rejects_incompatible_stale_component_plan(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    _install_runner_plans(
        tmp_path,
        monkeypatch,
        component_plan={
            **_component_plan(),
            "request_sha256": "5" * 64,
            "graph_sha256": "6" * 64,
            "actions": [
                {
                    "action": "accept",
                    "object_ids": ["component_0002"],
                    "parameters": {},
                    "confidence": 0.99,
                    "evidence": ["stale plan targets an unknown component"],
                }
            ],
        },
    )
    run_dir = tmp_path / "run"
    response, _ = _bound_component_request(run_dir)
    request = runner._component_binding(response, run_dir)

    with pytest.raises(runner.BenchmarkFailure, match="incompatible_plan"):
        runner._resolve_component_plan(
            "pptx-mixed-screenshot-candidates",
            request,
            allow_stale_binding=True,
            graph=request["_component_graph"],
        )


def test_release_runner_diagnostic_rejects_invalid_rebound_hashes(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    _install_runner_plans(
        tmp_path,
        monkeypatch,
        component_plan={**_component_plan(), "request_sha256": "5" * 64},
    )
    request = {**_component_request(), "request_sha256": "not-a-sha256"}

    with pytest.raises(runner.BenchmarkFailure, match="invalid_response"):
        runner._resolve_component_plan(
            "pptx-mixed-screenshot-candidates",
            request,
            allow_stale_binding=True,
        )


def test_release_runner_rejects_duplicate_and_hardlinked_plans(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    plans = _install_runner_plans(
        tmp_path, monkeypatch, component_plan=_component_plan()
    )
    duplicate = plans / "pptx-mixed-screenshot-candidates--component-copy.json"
    duplicate.write_bytes(
        (plans / "pptx-mixed-screenshot-candidates--component.json").read_bytes()
    )
    with pytest.raises(runner.BenchmarkFailure, match="duplicate_plan"):
        runner._select_component_plan("pptx-mixed-screenshot-candidates", _component_request())

    duplicate.unlink()
    hardlink = plans / "pptx-mixed-screenshot-candidates--component-hardlink.json"
    hardlink.hardlink_to(plans / "pptx-mixed-screenshot-candidates--component.json")
    with pytest.raises(runner.BenchmarkFailure, match="invalid_plan"):
        runner._select_component_plan("pptx-mixed-screenshot-candidates", _component_request())


def _write_valid_release_quality(reconstruction: Path) -> None:
    graph_sha256 = "a" * 64
    (reconstruction / "component_result.json").write_text(
        json.dumps(
            {
                "final_component_ids": [],
                "text_items": [],
                "repair_rounds": 1,
                "accepted_graph_sha256": graph_sha256,
                "warning": None,
                "fallback": {"status": "none", "parent_ids": []},
            }
        ),
        encoding="utf-8",
    )
    quality_path = reconstruction / "execution-01/component-quality.json"
    quality_path.parent.mkdir()
    quality_path.write_text(
        json.dumps(
            {
                "page_id": "page_001",
                "repair_round": 1,
                "input_graph_sha256": graph_sha256,
                "report": {
                    "visual_metrics": {"unexplained_visual_pixels": 0},
                    "violations": ["pptx_reopen_unknown"],
                    "component_reports": [],
                },
            }
        ),
        encoding="utf-8",
    )


def _batch_manifest(tmp_path: Path, *, warning: bool = False) -> Path:
    root = tmp_path / "release"
    source = root / "inputs/01.png"
    source.parent.mkdir(parents=True)
    source.write_bytes(b"benchmark-input")
    case = {
        "id": "image-one",
        "kind": "image",
        "path": "inputs/01.png",
        "sha256": hashlib.sha256(source.read_bytes()).hexdigest(),
        "page_count": 1,
        "expected_pages": [
            {
                "page_id": "page_001",
                "expected_status": "validated",
                "min_visual_components": 0,
                "min_text_boxes": 0,
                "max_unexplained_pixels": 0,
                "max_quality_violations": 0,
            }
        ],
        "agent_provider": "host",
    }
    if warning:
        case["expected_pages"][0]["warning"] = True
    manifest = root / "manifest.json"
    manifest.write_text(json.dumps({"schema_version": 2, "cases": [case]}), encoding="utf-8")
    return manifest


def test_release_performance_summary_uses_three_passed_repeats() -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    attempts = [
        {"case_id": "a", "repeat": 1, "status": "passed", "duration_ms": 300},
        {"case_id": "b", "repeat": 1, "status": "passed", "duration_ms": 30},
        {"case_id": "a", "repeat": 2, "status": "passed", "duration_ms": 100},
        {"case_id": "b", "repeat": 2, "status": "passed", "duration_ms": 10},
        {"case_id": "a", "repeat": 3, "status": "passed", "duration_ms": 200},
        {"case_id": "b", "repeat": 3, "status": "passed", "duration_ms": 20},
    ]

    assert runner.aggregate_performance(attempts) == {
        "repeat_total_duration_ms": [330, 110, 220],
        "median_total_duration_ms": 220,
        "case_median_duration_ms": {"a": 200, "b": 20},
    }


def _release_baseline() -> dict[str, object]:
    return {
        "schema_version": 1,
        "benchmark": "v0.2-core-14-page",
        "manifest_sha256": "a" * 64,
        "constraints_sha256": "b" * 64,
        "environment": {
            "os": "Windows",
            "architecture": "AMD64",
            "python": "3.12",
            "device": "cuda",
        },
        "median_total_duration_ms": 100,
        "case_median_duration_ms": {"image-one": 100},
    }


def _release_report(duration_ms: int = 100) -> dict[str, object]:
    return {
        "schema_version": 1,
        "status": "passed",
        "manifest_sha256": "a" * 64,
        "repeat": 3,
        "totals": {"cases": 1, "pages": 3, "failed_attempts": 0},
        "attempts": [],
        "performance": {
            "repeat_total_duration_ms": [duration_ms] * 3,
            "median_total_duration_ms": duration_ms,
            "case_median_duration_ms": {"image-one": duration_ms},
        },
    }


def test_release_baseline_comparison_uses_same_environment_and_fifteen_percent_limit() -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    baseline = _release_baseline()
    keywords = {
        "constraints_sha256": "b" * 64,
        "environment": baseline["environment"],
    }

    assert runner.compare_baseline(_release_report(115), baseline, **keywords) == {
        "status": "passed",
        "median_total_duration_ms": 115,
        "limit_ms": 115,
    }
    assert runner.compare_baseline(_release_report(116), baseline, **keywords) == {
        "status": "regressed",
        "median_total_duration_ms": 116,
        "limit_ms": 115,
    }


def test_release_baseline_does_not_compare_different_device_or_inputs() -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    baseline = _release_baseline()
    report = _release_report()
    report["manifest_sha256"] = "c" * 64
    environment = {**baseline["environment"], "device": "cpu"}

    assert runner.compare_baseline(
        report,
        baseline,
        constraints_sha256="d" * 64,
        environment=environment,
    ) == {
        "status": "not_comparable",
        "reasons": ["manifest_sha256", "constraints_sha256", "environment"],
    }


@pytest.mark.parametrize(
    "mutation",
    [
        {"status": "failed"},
        {"repeat": 2},
        {"performance": {"median_total_duration_ms": 100}},
    ],
)
def test_release_baseline_rejects_non_strict_report(mutation: dict[str, object]) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    report = {**_release_report(), **mutation}
    baseline = _release_baseline()

    with pytest.raises(runner.BenchmarkFailure, match="invalid_performance_result"):
        runner.compare_baseline(
            report,
            baseline,
            constraints_sha256="b" * 64,
            environment=baseline["environment"],
        )


@pytest.mark.parametrize(
    "attempts",
    [
        [],
        [
            {
                "case_id": "a",
                "repeat": 1,
                "status": "failed",
                "duration_ms": 1,
            }
        ],
        [
            {
                "case_id": "a",
                "repeat": True,
                "status": "passed",
                "duration_ms": 1,
            }
        ],
        [
            {
                "case_id": "a",
                "repeat": 1,
                "status": "passed",
                "duration_ms": -1,
            }
        ],
        [
            {
                "case_id": "a",
                "repeat": 1,
                "status": "passed",
                "duration_ms": 1,
            },
            {
                "case_id": "a",
                "repeat": 1,
                "status": "passed",
                "duration_ms": 2,
            },
        ],
        [
            {
                "case_id": "a",
                "repeat": repeat,
                "status": "passed",
                "duration_ms": repeat,
            }
            for repeat in (1, 2)
        ],
    ],
)
def test_release_performance_summary_rejects_incomplete_or_failed_attempts(
    attempts: list[dict[str, object]],
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")

    with pytest.raises(runner.BenchmarkFailure, match="invalid_performance_result"):
        runner.aggregate_performance(attempts)


def test_release_runner_manifest_repeats_three_times_and_writes_report(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    calls: list[Path] = []

    def fake_case(case: dict[str, object], *, workspace: Path, command: object) -> runner.BenchmarkCaseResult:
        calls.append(workspace)
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        _write_valid_release_quality(page / "reconstruction")
        return runner.BenchmarkCaseResult(
            case_id=str(case["id"]),
            run_dir=str(run_dir),
            pages=[{"page_id": "page_001", "status": "validated"}],
            duration_ms=7,
        )

    report_path = tmp_path / "results/report.json"
    report = runner.run_manifest(
        manifest,
        workspace=tmp_path / "runs",
        report_path=report_path,
        case_runner=fake_case,
    )
    assert report["status"] == "passed"
    assert report["manifest_sha256"] == runner.manifest_sha256(manifest)
    assert report["repeat"] == 3
    assert report["totals"] == {"cases": 1, "pages": 3, "failed_attempts": 0}
    assert report["performance"] == {
        "repeat_total_duration_ms": [7, 7, 7],
        "median_total_duration_ms": 7,
        "case_median_duration_ms": {"image-one": 7},
    }
    assert len(calls) == 3
    assert len(set(calls)) == 3
    assert json.loads(report_path.read_text(encoding="utf-8"))["status"] == "passed"


def test_release_runner_manifest_persists_only_structured_failure_details(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    details = {
        "stage": "component_plan",
        "page_id": "page_001",
        "repair_round": 1,
        "actual_request_sha256": "3" * 64,
        "actual_graph_sha256": "4" * 64,
        "expected_bindings": [
            {"request_sha256": "5" * 64, "graph_sha256": "6" * 64}
        ],
    }

    def stale_case(*_: object, **__: object) -> runner.BenchmarkCaseResult:
        raise runner.BenchmarkFailure("stale_plan", details)

    report = runner.run_manifest(
        manifest,
        workspace=tmp_path / "runs",
        report_path=tmp_path / "report.json",
        case_runner=stale_case,
    )

    assert report["status"] == "failed"
    assert [attempt["diagnostics"] for attempt in report["attempts"]] == [
        details,
        details,
        details,
    ]


def test_release_diagnostic_writes_stable_rebound_plans_without_official_pass(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    rebound = {
        **_component_plan(),
        "request_sha256": "7" * 64,
        "graph_sha256": "8" * 64,
    }
    environment = {
        "os": "Windows",
        "architecture": "AMD64",
        "python": "3.12",
        "device": "cpu",
    }
    monkeypatch.setattr(runner, "benchmark_environment", lambda: environment)
    constraints = tmp_path / "constraints.txt"
    constraints.write_text("package==1\n", encoding="utf-8")
    calls: list[tuple[bool, Path]] = []

    def diagnostic_case(
        case: dict[str, object],
        *,
        workspace: Path,
        command: object,
        allow_stale_bindings: bool,
        plan_observer: object,
    ) -> runner.BenchmarkCaseResult:
        calls.append((allow_stale_bindings, workspace))
        plan_observer("image-one--component-round-01.json", rebound)
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        _write_valid_release_quality(page / "reconstruction")
        return runner.BenchmarkCaseResult(
            str(case["id"]),
            str(run_dir),
            [{"page_id": "page_001", "status": "validated"}],
            9,
        )

    plans_output = tmp_path / "rebound-plans"
    report = runner.run_diagnostic_manifest(
        manifest,
        case_ids=["image-one"],
        workspace=tmp_path / "runs",
        report_path=tmp_path / "report.json",
        plans_output=plans_output,
        constraints_path=constraints,
        case_runner=diagnostic_case,
    )

    assert report["report_kind"] == "diagnostic"
    assert report["status"] == "diagnostic_complete"
    assert report["status"] != "passed"
    assert report["repeat"] == 3
    assert report["selected_case_ids"] == ["image-one"]
    assert report["constraints_sha256"] == runner.canonical_text_sha256(constraints)
    assert report["environment"] == environment
    assert report["totals"] == {
        "cases": 1,
        "attempts": 3,
        "pages": 3,
        "failed_attempts": 0,
    }
    assert report["performance"] == {
        "repeat_total_duration_ms": [9, 9, 9],
        "median_total_duration_ms": 9,
        "case_median_duration_ms": {"image-one": 9},
    }
    assert len(calls) == 3
    assert all(allow_stale is True for allow_stale, _ in calls)
    assert len({workspace for _, workspace in calls}) == 3
    stored = json.loads(
        (plans_output / "image-one--component-round-01.json").read_text(
            encoding="utf-8"
        )
    )
    assert stored == rebound


def test_release_diagnostic_sorts_selected_case_ids_for_strict_validation(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    payload = json.loads(manifest.read_text(encoding="utf-8"))
    second_case = dict(payload["cases"][0])
    second_case["id"] = "image-alpha"
    payload["cases"].append(second_case)
    manifest.write_text(json.dumps(payload), encoding="utf-8")
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)

    def diagnostic_case(
        case: dict[str, object],
        *,
        workspace: Path,
        command: object,
        allow_stale_bindings: bool,
        plan_observer: object,
    ) -> runner.BenchmarkCaseResult:
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        _write_valid_release_quality(page / "reconstruction")
        return runner.BenchmarkCaseResult(
            str(case["id"]),
            str(run_dir),
            [{"page_id": "page_001", "status": "validated"}],
            9,
        )

    report = runner.run_diagnostic_manifest(
        manifest,
        case_ids=["image-one", "image-alpha"],
        workspace=tmp_path / "runs",
        report_path=tmp_path / "report.json",
        plans_output=tmp_path / "plans",
        case_runner=diagnostic_case,
    )

    assert report["selected_case_ids"] == ["image-alpha", "image-one"]
    attempts = report["attempts"]
    assert {(attempt["case_id"], attempt["repeat"]) for attempt in attempts} == {
        (case_id, repeat)
        for case_id in ("image-one", "image-alpha")
        for repeat in (1, 2, 3)
    }


def test_release_probe_runs_once_without_performance_or_plan_output(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    constraints = tmp_path / "constraints.txt"
    constraints.write_text("package==1\n", encoding="utf-8")
    calls: list[Path] = []

    def probe_case(
        case: dict[str, object],
        *,
        workspace: Path,
        command: object,
        allow_stale_bindings: bool,
        plan_observer: object,
    ) -> runner.BenchmarkCaseResult:
        assert allow_stale_bindings is True
        calls.append(workspace)
        plan_observer(
            "image-one--component-round-01.json",
            {
                **_component_plan(),
                "request_sha256": "7" * 64,
                "graph_sha256": "8" * 64,
            },
        )
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        _write_valid_release_quality(page / "reconstruction")
        return runner.BenchmarkCaseResult(
            str(case["id"]),
            str(run_dir),
            [{"page_id": "page_001", "status": "validated"}],
            9,
        )

    plans_output = tmp_path / "probe-plans"
    report = runner.run_diagnostic_manifest(
        manifest,
        case_ids=["image-one"],
        workspace=tmp_path / "runs",
        report_path=tmp_path / "probe.json",
        plans_output=plans_output,
        constraints_path=constraints,
        repeat=1,
        case_runner=probe_case,
    )

    assert report["report_kind"] == "probe"
    assert report["status"] == "probe_complete"
    assert report["repeat"] == 1
    assert report["totals"] == {
        "cases": 1,
        "attempts": 1,
        "pages": 1,
        "failed_attempts": 0,
    }
    assert "performance" not in report
    assert calls == [tmp_path / "runs/repeat-01"]
    assert not plans_output.exists()

    with pytest.raises(runner.BenchmarkFailure, match="invalid_repeat"):
        runner.run_shard_manifest(
            manifest,
            case_ids=["image-one"],
            workspace=tmp_path / "official-runs",
            report_path=tmp_path / "official.json",
            constraints_path=constraints,
            repeat=1,
            case_runner=probe_case,
        )

    with pytest.raises(runner.BenchmarkFailure, match="invalid_repeat"):
        runner.run_diagnostic_manifest(
            manifest,
            case_ids=["image-one"],
            workspace=tmp_path / "invalid-runs",
            report_path=tmp_path / "invalid.json",
            plans_output=tmp_path / "invalid-plans",
            constraints_path=constraints,
            repeat=2,
            case_runner=probe_case,
        )


def test_release_diagnostic_rejects_unstable_rebound_plan_across_repeats(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    attempt = 0

    def unstable_case(
        case: dict[str, object],
        *,
        workspace: Path,
        command: object,
        allow_stale_bindings: bool,
        plan_observer: object,
    ) -> runner.BenchmarkCaseResult:
        nonlocal attempt
        attempt += 1
        filename = (
            "image-one--component-round-01.json"
            if attempt == 1
            else "image-one--component-round-01-alt.json"
        )
        plan_observer(
            filename,
            {**_component_plan(), "request_sha256": str(attempt) * 64},
        )
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        _write_valid_release_quality(page / "reconstruction")
        return runner.BenchmarkCaseResult(
            str(case["id"]),
            str(run_dir),
            [{"page_id": "page_001", "status": "validated"}],
            9,
        )

    plans_output = tmp_path / "rebound-plans"
    report = runner.run_diagnostic_manifest(
        manifest,
        case_ids=["image-one"],
        workspace=tmp_path / "runs",
        report_path=tmp_path / "report.json",
        plans_output=plans_output,
        case_runner=unstable_case,
    )

    assert report["status"] == "failed"
    assert report["totals"]["failed_attempts"] == 2
    assert report["attempts"][1]["error_type"] == "unstable_plan_binding"
    assert not plans_output.exists()


def test_release_diagnostic_checks_exact_bindings_without_exporting_them(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    monkeypatch.setattr(
        runner,
        "benchmark_environment",
        lambda: {
            "os": "Windows",
            "architecture": "AMD64",
            "python": "3.12",
            "device": "cpu",
        },
    )

    def exact_case(
        case: dict[str, object],
        *,
        workspace: Path,
        command: object,
        allow_stale_bindings: bool,
        plan_observer: object,
    ) -> runner.BenchmarkCaseResult:
        plan_observer("image-one--component-round-01.json", _component_plan(), False)
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        _write_valid_release_quality(page / "reconstruction")
        return runner.BenchmarkCaseResult(
            str(case["id"]),
            str(run_dir),
            [{"page_id": "page_001", "status": "validated"}],
            9,
        )

    plans_output = tmp_path / "plans"
    report = runner.run_diagnostic_manifest(
        manifest,
        case_ids=["image-one"],
        workspace=tmp_path / "runs",
        report_path=tmp_path / "report.json",
        plans_output=plans_output,
        case_runner=exact_case,
    )

    assert report["status"] == "diagnostic_complete"
    assert plans_output.is_dir()
    assert list(plans_output.iterdir()) == []


def test_release_diagnostic_rejects_existing_plans_output_before_running(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    plans_output = tmp_path / "plans"
    plans_output.mkdir()

    with pytest.raises(runner.BenchmarkFailure, match="invalid_workspace"):
        runner.run_diagnostic_manifest(
            manifest,
            case_ids=["image-one"],
            workspace=tmp_path / "runs",
            report_path=tmp_path / "report.json",
            plans_output=plans_output,
        )

    assert not (tmp_path / "report.json").exists()


def test_release_official_shard_runs_only_selected_cases_and_keeps_manifest_identity(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    environment = {
        "os": "Windows",
        "architecture": "AMD64",
        "python": "3.12",
        "device": "cpu",
    }
    monkeypatch.setattr(runner, "benchmark_environment", lambda: environment)
    constraints = tmp_path / "constraints.txt"
    constraints.write_text("package==1\n", encoding="utf-8")

    calls: list[bool] = []

    def shard_case(
        case: dict[str, object],
        *,
        workspace: Path,
        command: object,
        allow_compatible_component_bindings: bool,
        plan_observer: object,
    ) -> runner.BenchmarkCaseResult:
        calls.append(allow_compatible_component_bindings)
        plan_observer(
            "image-one--component-round-01.json", _component_plan(), False
        )
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        _write_valid_release_quality(page / "reconstruction")
        return runner.BenchmarkCaseResult(
            str(case["id"]),
            str(run_dir),
            [{"page_id": "page_001", "status": "validated"}],
            11,
        )

    report = runner.run_shard_manifest(
        manifest,
        case_ids=["image-one"],
        workspace=tmp_path / "runs",
        report_path=tmp_path / "shard.json",
        constraints_path=constraints,
        case_runner=shard_case,
    )

    assert report["report_kind"] == "shard"
    assert report["status"] == "passed"
    assert report["manifest_sha256"] == runner.manifest_sha256(manifest)
    assert report["constraints_sha256"] == runner.canonical_text_sha256(constraints)
    assert report["environment"] == environment
    assert report["selected_case_ids"] == ["image-one"]
    assert report["totals"] == {
        "cases": 1,
        "attempts": 3,
        "pages": 3,
        "failed_attempts": 0,
    }
    assert calls == [True, True, True]


@pytest.mark.parametrize(
    ("variation", "expected_failed_attempts"),
    [("binding", 2), ("missing", 1)],
)
def test_release_official_shard_rejects_unstable_component_plans(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
    variation: str,
    expected_failed_attempts: int,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)
    attempt = 0

    def shard_case(
        case: dict[str, object],
        *,
        workspace: Path,
        command: object,
        allow_compatible_component_bindings: bool,
        plan_observer: object,
    ) -> runner.BenchmarkCaseResult:
        nonlocal attempt
        assert allow_compatible_component_bindings is True
        attempt += 1
        if variation != "missing" or attempt != 2:
            plan = _component_plan()
            if variation == "binding":
                plan = {**plan, "request_sha256": str(attempt) * 64}
            plan_observer(
                "image-one--component-round-01.json", plan, variation == "binding"
            )
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        _write_valid_release_quality(page / "reconstruction")
        return runner.BenchmarkCaseResult(
            str(case["id"]),
            str(run_dir),
            [{"page_id": "page_001", "status": "validated"}],
            11,
        )

    report = runner.run_shard_manifest(
        manifest,
        case_ids=["image-one"],
        workspace=tmp_path / "runs",
        report_path=tmp_path / "shard.json",
        case_runner=shard_case,
    )

    assert report["status"] == "failed"
    failed = [attempt for attempt in report["attempts"] if attempt["status"] == "failed"]
    assert len(failed) == expected_failed_attempts
    assert all(attempt["error_type"] == "unstable_plan_binding" for attempt in failed)


_CORE_SHARDS = [
    ["pptx-mixed-screenshot-candidates"],
    ["pdf-rotated-page", "image-thin-line-network"],
    ["image-icon-matrix", "image-dark-poster"],
    ["image-flowchart", "image-combo-chart"],
    [
        "image-bilingual-dashboard",
        "image-tiny-element-table",
        "image-non-16-9-infographic",
    ],
]


def _write_core_shard_reports(
    tmp_path: Path,
    runner: object,
) -> tuple[list[Path], Path, Path, dict[str, object]]:
    manifest = _core_manifest()
    cases = {case["id"]: case for case in manifest["cases"]}
    manifest_hash = runner.manifest_sha256(CORE_MANIFEST_PATH)
    constraints = tmp_path / "constraints.txt"
    constraints.write_text("package==1\n", encoding="utf-8")
    constraints_hash = runner.canonical_text_sha256(constraints)
    environment = {
        "os": "Windows",
        "architecture": "AMD64",
        "python": "3.12",
        "device": "cpu",
    }
    paths = []
    for shard_index, case_ids in enumerate(_CORE_SHARDS, start=1):
        attempts = []
        for repeat in (1, 2, 3):
            for case_id in case_ids:
                case = cases[case_id]
                attempts.append(
                    {
                        "case_id": case_id,
                        "repeat": repeat,
                        "status": "passed",
                        "duration_ms": 10,
                        "pages": [
                            {
                                "page_id": f"page_{page_number:03d}",
                                "status": page["expected_status"],
                            }
                            for page_number, page in enumerate(
                                case["expected_pages"], start=1
                            )
                        ],
                    }
                )
        report = {
            "schema_version": 1,
            "report_kind": "shard",
            "status": "passed",
            "manifest_sha256": manifest_hash,
            "constraints_sha256": constraints_hash,
            "environment": environment,
            "repeat": 3,
            "selected_case_ids": sorted(case_ids),
            "attempts": attempts,
            "totals": {
                "cases": len(case_ids),
                "attempts": len(attempts),
                "pages": sum(cases[case_id]["page_count"] for case_id in case_ids)
                * 3,
                "failed_attempts": 0,
            },
            "performance": runner.aggregate_performance(attempts),
        }
        path = tmp_path / f"shard-{shard_index}.json"
        path.write_text(json.dumps(report), encoding="utf-8")
        paths.append(path)
    baseline = {
        "schema_version": 1,
        "benchmark": "v0.2-core-14-page",
        "manifest_sha256": manifest_hash,
        "constraints_sha256": constraints_hash,
        "environment": environment,
        "median_total_duration_ms": 1000,
        "case_median_duration_ms": {case_id: 100 for case_id in cases},
    }
    baseline_path = tmp_path / "baseline.json"
    baseline_path.write_text(json.dumps(baseline), encoding="utf-8")
    return paths, constraints, baseline_path, baseline


def _write_core_diagnostic_reports(
    tmp_path: Path,
    runner: object,
) -> tuple[list[Path], Path]:
    paths, constraints, _, _ = _write_core_shard_reports(tmp_path, runner)
    for path in paths:
        report = json.loads(path.read_text(encoding="utf-8"))
        report["report_kind"] = "diagnostic"
        report["status"] = "diagnostic_complete"
        path.write_text(json.dumps(report), encoding="utf-8")
    return paths, constraints


def test_release_baseline_candidate_uses_only_complete_cpu_diagnostics(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    reports, constraints = _write_core_diagnostic_reports(tmp_path, runner)

    candidate = runner.create_baseline_candidate(
        CORE_MANIFEST_PATH,
        shard_report_paths=reports,
        constraints_path=constraints,
        report_path=tmp_path / "baseline-candidate.json",
    )

    assert candidate == {
        "schema_version": 1,
        "benchmark": "v0.2-core-14-page",
        "manifest_sha256": runner.manifest_sha256(CORE_MANIFEST_PATH),
        "constraints_sha256": runner.canonical_text_sha256(constraints),
        "environment": {
            "os": "Windows",
            "architecture": "AMD64",
            "python": "3.12",
            "device": "cpu",
        },
        "median_total_duration_ms": 100,
        "case_median_duration_ms": {
            case["id"]: 10 for case in _core_manifest()["cases"]
        },
    }


@pytest.mark.parametrize(
    "mutation",
    [
        "probe",
        "failed_diagnostic",
        "cuda",
        "wrong_manifest",
        "wrong_constraints",
        "wrong_totals",
        "wrong_performance",
        "missing_report",
        "duplicate_case_coverage",
        "existing_output",
    ],
)
def test_release_baseline_candidate_rejects_non_official_evidence(
    tmp_path: Path, mutation: str
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    reports, constraints = _write_core_diagnostic_reports(tmp_path, runner)
    payload = json.loads(reports[0].read_text(encoding="utf-8"))
    report_path = tmp_path / "baseline-candidate.json"
    if mutation == "probe":
        payload["report_kind"] = "probe"
        payload["status"] = "probe_complete"
    elif mutation == "failed_diagnostic":
        payload["status"] = "failed"
    elif mutation == "cuda":
        payload["environment"] = {**payload["environment"], "device": "cuda"}
    elif mutation == "wrong_manifest":
        payload["manifest_sha256"] = "f" * 64
    elif mutation == "wrong_constraints":
        payload["constraints_sha256"] = "f" * 64
    elif mutation == "wrong_totals":
        payload["totals"]["pages"] += 1
    elif mutation == "wrong_performance":
        payload["performance"]["median_total_duration_ms"] += 1
    elif mutation == "duplicate_case_coverage":
        other = json.loads(reports[1].read_text(encoding="utf-8"))
        other["selected_case_ids"] = payload["selected_case_ids"]
        reports[1].write_text(json.dumps(other), encoding="utf-8")
    elif mutation == "missing_report":
        reports.pop()
    elif mutation == "existing_output":
        report_path.write_text("{}", encoding="utf-8")
    if mutation not in {"missing_report", "existing_output"}:
        reports[0].write_text(json.dumps(payload), encoding="utf-8")

    with pytest.raises(runner.BenchmarkFailure):
        runner.create_baseline_candidate(
            CORE_MANIFEST_PATH,
            shard_report_paths=reports,
            constraints_path=constraints,
            report_path=report_path,
        )


def test_release_aggregate_accepts_exact_five_shard_core_evidence(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    reports, constraints, baseline_path, baseline = _write_core_shard_reports(
        tmp_path, runner
    )

    report = runner.aggregate_shard_reports(
        CORE_MANIFEST_PATH,
        shard_report_paths=reports,
        constraints_path=constraints,
        baseline_path=baseline_path,
        report_path=tmp_path / "official.json",
    )

    assert report["report_kind"] == "official"
    assert report["status"] == "passed"
    assert report["manifest_sha256"] == baseline["manifest_sha256"]
    assert report["totals"] == {
        "cases": 10,
        "attempts": 30,
        "pages": 42,
        "failed_attempts": 0,
    }
    assert report["performance_comparison"]["status"] == "passed"


@pytest.mark.parametrize(
    "mutation",
    [
        "wrong_manifest",
        "wrong_constraints",
        "mixed_environment",
        "wrong_kind",
        "failed_attempt",
        "duplicate_repeat",
        "wrong_totals",
        "missing_shard",
        "duplicate_case_coverage",
    ],
)
def test_release_aggregate_rejects_incomplete_or_mixed_shard_evidence(
    tmp_path: Path, mutation: str
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    reports, constraints, baseline_path, _ = _write_core_shard_reports(tmp_path, runner)
    payload = json.loads(reports[0].read_text(encoding="utf-8"))
    if mutation == "wrong_manifest":
        payload["manifest_sha256"] = "f" * 64
    elif mutation == "wrong_constraints":
        payload["constraints_sha256"] = "f" * 64
    elif mutation == "mixed_environment":
        payload["environment"] = {**payload["environment"], "device": "cuda"}
    elif mutation == "wrong_kind":
        payload["report_kind"] = "diagnostic"
    elif mutation == "failed_attempt":
        payload["attempts"][0]["status"] = "failed"
    elif mutation == "duplicate_repeat":
        payload["attempts"][1]["repeat"] = payload["attempts"][0]["repeat"]
    elif mutation == "wrong_totals":
        payload["totals"]["pages"] += 1
    elif mutation == "duplicate_case_coverage":
        other = json.loads(reports[1].read_text(encoding="utf-8"))
        other["selected_case_ids"] = payload["selected_case_ids"]
        reports[1].write_text(json.dumps(other), encoding="utf-8")
    if mutation != "missing_shard":
        reports[0].write_text(json.dumps(payload), encoding="utf-8")
    else:
        reports.pop()

    with pytest.raises(runner.BenchmarkFailure):
        runner.aggregate_shard_reports(
            CORE_MANIFEST_PATH,
            shard_report_paths=reports,
            constraints_path=constraints,
            baseline_path=baseline_path,
            report_path=tmp_path / "official.json",
        )


def test_release_runner_manifest_identity_is_line_ending_independent(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    lf = tmp_path / "manifest-lf.json"
    crlf = tmp_path / "manifest-crlf.json"
    payload = CORE_MANIFEST_PATH.read_bytes().replace(b"\r\n", b"\n")
    lf.write_bytes(payload)
    crlf.write_bytes(payload.replace(b"\n", b"\r\n"))

    assert runner.manifest_sha256(lf) == runner.manifest_sha256(crlf)


def test_release_runner_maps_namespaced_manifest_page_to_local_run_page(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    run_dir = tmp_path / "run"
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(parents=True)
    _write_valid_release_quality(reconstruction)
    case = {
        "expected_pages": [
            {
                "page_id": "image-one-page-001",
                "expected_status": "replaced",
                "min_visual_components": 0,
                "min_text_boxes": 0,
                "max_unexplained_pixels": 0,
                "max_quality_violations": 0,
            }
        ]
    }
    result = runner.BenchmarkCaseResult(
        "image-one",
        str(run_dir),
        [{"page_id": "page_001", "status": "replaced"}],
        1,
    )

    runner._validate_batch_case(case, result)


def test_release_runner_accepts_strict_legacy_validated_page_status(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    run_dir = tmp_path / "run"
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(parents=True)
    _write_valid_release_quality(reconstruction)
    case = {
        "expected_pages": [
            {
                "page_id": "image-one-page-001",
                "expected_status": "validated",
                "min_visual_components": 0,
                "min_text_boxes": 0,
                "max_unexplained_pixels": 0,
                "max_quality_violations": 0,
            }
        ]
    }
    result = runner.BenchmarkCaseResult(
        "image-one",
        str(run_dir),
        [{"page_id": "page_001", "status": "validated"}],
        1,
    )

    runner._validate_batch_case(case, result)


def test_release_runner_rejects_replaced_page_when_validated_is_expected(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    case = {
        "expected_pages": [
            {
                "page_id": "image-one-page-001",
                "expected_status": "validated",
                "min_visual_components": 0,
                "min_text_boxes": 0,
                "max_unexplained_pixels": 0,
                "max_quality_violations": 0,
            }
        ]
    }
    result = runner.BenchmarkCaseResult(
        "image-one",
        str(tmp_path / "run"),
        [{"page_id": "page_001", "status": "replaced"}],
        1,
    )

    with pytest.raises(runner.BenchmarkFailure, match="invalid_page_result"):
        runner._validate_batch_case(case, result)


def test_release_runner_enforces_visual_components_separately_from_text(
    tmp_path: Path,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    run_dir = tmp_path / "run"
    reconstruction = run_dir / "pages/page_001/reconstruction"
    reconstruction.mkdir(parents=True)
    result_path = reconstruction / "component_result.json"
    result_path.write_text(
        json.dumps(
            {
                "final_component_ids": ["visual_001", "visual_002", "visual_003"],
                "text_items": [
                    {"_component_id": f"text_{index:03d}"}
                    for index in range(1, 12)
                ],
                "warning": None,
                "fallback": {"status": "none", "parent_ids": []},
            }
        ),
        encoding="utf-8",
    )
    case = {
        "expected_pages": [
            {
                "page_id": "pdf-page-001",
                "expected_status": "replaced",
                "min_visual_components": 4,
                "min_text_boxes": 11,
                "max_unexplained_pixels": 0,
                "max_quality_violations": 0,
            }
        ]
    }
    result = runner.BenchmarkCaseResult(
        "pdf",
        str(run_dir),
        [{"page_id": "page_001", "status": "replaced"}],
        1,
    )

    with pytest.raises(runner.BenchmarkFailure, match="quality_gate"):
        runner._validate_batch_case(case, result)

    duplicate_visual = ["visual_001"] * 4
    text_items = [
        {"_component_id": f"text_{index:03d}"} for index in range(1, 12)
    ]
    result_path.write_text(
        json.dumps(
            {
                "final_component_ids": duplicate_visual,
                "text_items": text_items,
                "warning": None,
                "fallback": {"status": "none", "parent_ids": []},
            }
        ),
        encoding="utf-8",
    )
    with pytest.raises(runner.BenchmarkFailure, match="invalid_quality_result"):
        runner._validate_batch_case(case, result)

    unique_visual = [
        "visual_001",
        "visual_002",
        "visual_003",
        "visual_004",
    ]
    result_path.write_text(
        json.dumps(
            {
                "final_component_ids": unique_visual,
                "text_items": [{"_component_id": "text_001"}] * 11,
                "warning": None,
                "fallback": {"status": "none", "parent_ids": []},
            }
        ),
        encoding="utf-8",
    )
    with pytest.raises(runner.BenchmarkFailure, match="invalid_quality_result"):
        runner._validate_batch_case(case, result)

    component_result = {
        "final_component_ids": unique_visual,
        "text_items": text_items,
        "warning": None,
        "fallback": {"status": "none", "parent_ids": []},
    }
    result_path.write_text(json.dumps(component_result), encoding="utf-8")
    with pytest.raises(runner.BenchmarkFailure, match="invalid_quality_result"):
        runner._validate_batch_case(case, result)

    component_result.update(
        {"repair_rounds": 1, "accepted_graph_sha256": "a" * 64}
    )
    result_path.write_text(json.dumps(component_result), encoding="utf-8")
    quality_path = reconstruction / "execution-01/component-quality.json"
    quality_path.parent.mkdir()
    quality = {
        "page_id": "page_001",
        "repair_round": 1,
        "input_graph_sha256": "a" * 64,
        "report": {
            "visual_metrics": {"unexplained_visual_pixels": 0},
            "violations": ["pptx_reopen_unknown"],
            "component_reports": [
                {"component_id": identifier, "violations": []}
                for identifier in unique_visual
            ],
        },
    }
    quality_path.write_text(json.dumps(quality), encoding="utf-8")
    runner._validate_batch_case(case, result)

    quality_path.unlink()
    with pytest.raises(runner.BenchmarkFailure, match="invalid_quality_result"):
        runner._validate_batch_case(case, result)

    quality_path.write_text(json.dumps(quality), encoding="utf-8")
    del component_result["warning"]
    result_path.write_text(json.dumps(component_result), encoding="utf-8")
    with pytest.raises(runner.BenchmarkFailure, match="invalid_quality_result"):
        runner._validate_batch_case(case, result)

    component_result["warning"] = None
    del component_result["fallback"]
    result_path.write_text(json.dumps(component_result), encoding="utf-8")
    with pytest.raises(runner.BenchmarkFailure, match="invalid_quality_result"):
        runner._validate_batch_case(case, result)

    result_path.unlink()
    with pytest.raises(runner.BenchmarkFailure, match="invalid_quality_result"):
        runner._validate_batch_case(case, result)


def test_release_runner_requires_preserved_pptx_page_and_related_parts_to_match(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    release_root = tmp_path / "release"
    source = release_root / "inputs/source.pptx"
    run_dir = tmp_path / "runs/pptx/run"
    output = run_dir.parent / "output.pptx"
    source.parent.mkdir(parents=True)
    output.parent.mkdir(parents=True)
    slide = (
        b'<p:sld xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
        b'xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        b'<p:cSld><p:spTree><p:nvGrpSpPr/><p:grpSpPr/>'
        b'<p:pic/><p:sp><p:txBody><a:p><a:r><a:t>04</a:t></a:r></a:p></p:txBody></p:sp>'
        b'</p:spTree></p:cSld></p:sld>'
    )
    relationships = (
        b'<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
        b'<Relationship Id="rId1" Type="image" Target="../media/image1.png"/>'
        b'</Relationships>'
    )

    def write_pptx(path: Path, media: bytes) -> None:
        with ZipFile(path, "w", ZIP_DEFLATED) as archive:
            archive.writestr("ppt/slides/slide1.xml", slide)
            archive.writestr("ppt/slides/_rels/slide1.xml.rels", relationships)
            archive.writestr("ppt/media/image1.png", media)

    write_pptx(source, b"same-image")
    write_pptx(output, b"same-image")
    monkeypatch.setattr(runner, "RELEASE_ROOT", release_root)
    case = {
        "kind": "pptx",
        "path": "inputs/source.pptx",
        "expected_pages": [
            {
                "page_id": "pptx-page-001",
                "expected_status": "preserved",
                "min_visual_components": 1,
                "min_text_boxes": 1,
                "max_unexplained_pixels": 0,
                "max_quality_violations": 0,
            }
        ],
    }
    result = runner.BenchmarkCaseResult(
        "pptx", str(run_dir), [{"page_id": "page_001", "status": "preserved"}], 1
    )

    runner._validate_batch_case(case, result)

    write_pptx(output, b"changed-image")
    with pytest.raises(runner.BenchmarkFailure, match="quality_gate"):
        runner._validate_batch_case(case, result)


def test_release_loader_rejects_float_schema_version(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _manifest()
    manifest["schema_version"] = 2.0
    path = tmp_path / "manifest.json"
    path.write_text(json.dumps(manifest), encoding="utf-8")
    monkeypatch.setattr(runner, "RELEASE_ROOT", RELEASE_ROOT)

    with pytest.raises(runner.BenchmarkFailure, match="invalid_manifest"):
        runner._load_batch_cases(path)


def test_release_runner_manifest_fails_closed_on_warning_and_invalid_repeat(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path, warning=True)
    monkeypatch.setattr(runner, "RELEASE_ROOT", manifest.parent)

    def warning_case(case: dict[str, object], *, workspace: Path, command: object) -> runner.BenchmarkCaseResult:
        run_dir = workspace / str(case["id"]) / "run"
        page = run_dir / "pages/page_001"
        (page / "reconstruction").mkdir(parents=True)
        (page / "page_result.json").write_text(
            '{"page_id":"page_001","status":"validated"}', encoding="utf-8"
        )
        (page / "reconstruction/component_result.json").write_text(
            '{"warning":"fallback","fallback":{"status":"none"}}', encoding="utf-8"
        )
        return runner.BenchmarkCaseResult(str(case["id"]), str(run_dir), [{"page_id":"page_001","status":"validated"}], 1)

    report = runner.run_manifest(
        manifest,
        workspace=tmp_path / "runs",
        report_path=tmp_path / "report.json",
        case_runner=warning_case,
    )
    assert report["status"] == "failed"
    assert report["totals"]["failed_attempts"] == 3
    assert "performance" not in report
    with pytest.raises(runner.BenchmarkFailure, match="invalid_repeat"):
        runner.run_manifest(
            manifest,
            workspace=tmp_path / "other",
            report_path=tmp_path / "other.json",
            repeat=2,
            case_runner=warning_case,
        )


def test_release_runner_cli_routes_explicit_modes(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    captured: list[tuple[str, Path, dict[str, object]]] = []

    def fake_shard(path: Path, **kwargs: object) -> dict[str, object]:
        captured.append(("official-shard", path, kwargs))
        return {"status": "passed"}

    def fake_diagnostic(path: Path, **kwargs: object) -> dict[str, object]:
        captured.append(("diagnostic-shard", path, kwargs))
        return {"status": "diagnostic_complete"}

    def fake_aggregate(path: Path, **kwargs: object) -> dict[str, object]:
        captured.append(("aggregate", path, kwargs))
        return {"status": "passed"}

    def fake_baseline_candidate(path: Path, **kwargs: object) -> dict[str, object]:
        captured.append(("baseline-candidate", path, kwargs))
        return {"status": "created"}

    monkeypatch.setattr(runner, "run_shard_manifest", fake_shard)
    monkeypatch.setattr(runner, "run_diagnostic_manifest", fake_diagnostic)
    monkeypatch.setattr(runner, "aggregate_shard_reports", fake_aggregate)
    monkeypatch.setattr(runner, "create_baseline_candidate", fake_baseline_candidate)
    assert (
        runner.main(
            [
                "--mode",
                "official-shard",
                "--manifest",
                str(manifest),
                "--workspace",
                str(tmp_path / "runs"),
                "--report",
                str(tmp_path / "report.json"),
                "--case-id",
                "image-one",
            ]
        )
        == 0
    )
    assert (
        runner.main(
            [
                "--mode",
                "diagnostic-shard",
                "--manifest",
                str(manifest),
                "--workspace",
                str(tmp_path / "diagnostic-runs"),
                "--report",
                str(tmp_path / "diagnostic.json"),
                "--plans-output",
                str(tmp_path / "plans"),
                "--case-id",
                "image-one",
            ]
        )
        == 0
    )
    shard_reports = [tmp_path / f"shard-{index}.json" for index in range(5)]
    aggregate_arguments = [
        "--mode",
        "aggregate",
        "--manifest",
        str(manifest),
        "--report",
        str(tmp_path / "official.json"),
        "--baseline",
        str(tmp_path / "baseline.json"),
        "--constraints",
        str(tmp_path / "constraints.txt"),
    ]
    for shard_report in shard_reports:
        aggregate_arguments.extend(["--shard-report", str(shard_report)])
    assert runner.main(aggregate_arguments) == 0
    candidate_arguments = [
        "--mode",
        "baseline-candidate",
        "--manifest",
        str(manifest),
        "--report",
        str(tmp_path / "baseline-candidate.json"),
        "--constraints",
        str(tmp_path / "constraints.txt"),
    ]
    for shard_report in shard_reports:
        candidate_arguments.extend(["--shard-report", str(shard_report)])
    assert runner.main(candidate_arguments) == 0
    assert [mode for mode, _, _ in captured] == [
        "official-shard",
        "diagnostic-shard",
        "aggregate",
        "baseline-candidate",
    ]
    assert all(path == manifest for _, path, _ in captured)
    assert captured[0][2]["case_ids"] == ["image-one"]
    assert captured[0][2]["repeat"] == 3
    assert captured[1][2]["plans_output"] == tmp_path / "plans"
    assert captured[2][2]["shard_report_paths"] == shard_reports
    assert captured[3][2]["shard_report_paths"] == shard_reports


def test_release_runner_cli_routes_one_repeat_probe(
    tmp_path: Path, monkeypatch: pytest.MonkeyPatch
) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    manifest = _batch_manifest(tmp_path)
    captured: list[dict[str, object]] = []

    def fake_probe(path: Path, **kwargs: object) -> dict[str, object]:
        assert path == manifest
        captured.append(kwargs)
        return {"status": "probe_complete"}

    monkeypatch.setattr(runner, "run_diagnostic_manifest", fake_probe)
    assert (
        runner.main(
            [
                "--mode",
                "diagnostic-shard",
                "--manifest",
                str(manifest),
                "--workspace",
                str(tmp_path / "probe-runs"),
                "--report",
                str(tmp_path / "probe.json"),
                "--plans-output",
                str(tmp_path / "plans"),
                "--case-id",
                "image-one",
                "--repeat",
                "1",
            ]
        )
        == 0
    )
    assert captured[0]["repeat"] == 1


def test_release_runner_cli_requires_mode_specific_arguments(tmp_path: Path) -> None:
    runner = importlib.import_module("scripts.release_benchmark")
    with pytest.raises(SystemExit):
        runner.main(
            [
                "--manifest",
                str(tmp_path / "manifest.json"),
                "--report",
                str(tmp_path / "report.json"),
            ]
        )
    with pytest.raises(SystemExit):
        runner.main(
            [
                "--mode",
                "diagnostic-shard",
                "--manifest",
                str(tmp_path / "manifest.json"),
                "--workspace",
                str(tmp_path / "runs"),
                "--report",
                str(tmp_path / "report.json"),
                "--case-id",
                "image-one",
            ]
        )
